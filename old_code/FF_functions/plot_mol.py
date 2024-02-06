#!/bin/env python                                                                                                                                                              
# Author: Brett Savoie (brettsavoie@gmail.com)

import sys,argparse,os,ast,re,fnmatch,matplotlib,subprocess,shutil,itertools,datetime
matplotlib.use('Agg') # Needed for cluster image generation
from scipy.optimize import curve_fit,minimize
from scipy.spatial.distance import cdist
from scipy.interpolate import UnivariateSpline
import numpy as np #from numpy import *
from numpy.linalg import norm
from scipy.linalg import lstsq
import time as Time
#from pylab import *
from copy import deepcopy
from scipy import interpolate
from math import acos
from matplotlib import pyplot as plt
# Powerpoint generation dependencies
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches,Pt

# Add TAFFI Lib to path
sys.path.append('/'.join(os.path.abspath(__file__).split('/')[:-2])+'/Lib')
from extract_charges import *
from transify import *
from adjacency import *
from file_parsers import *
from id_types import *
from kekule import *
import random

def main(argv):

    parser = argparse.ArgumentParser(description='This program collects force field parameters from the output of Orca calculations. This script should be used in conjunction with '+\
                                                 'the paramgen.py program, which automatically generates the Orca jobs for submission. After the quantum chemistry completes, this '+\
                                                 'script collates the information and generates a database file that can be used in conjunction with polygen to generate the MD job.')

    #required (positional) arguments                                                                                                  
    parser.add_argument('xyzfile', help = 'A quoted list of the *.xyz files to be included in the sim, or a folder holding a previous md cycles. If a folder is supplied then the FF parameters and last configuration '+\
                                              'are used to start the new job along with the most recently fit VDW parameters.')

    parser.add_argument('-o', dest='outfile', default='',
                        help = 'output pdb filename, default: name of xyz file')

    # Make parse inputs
    args=parser.parse_args(argv)
   
    Elements,Geometry = xyz_parse(args.xyzfile)
    print(Elements)
    Adj_mat = Table_generator(Elements,Geometry)
    print(Adj_mat)
    Hybridizations = Hybridization_finder(Elements,Adj_mat)
    gens =2
    Atom_types = id_types(Elements,Adj_mat,gens)
    print(Atom_types)
    geo_2D = kekule(Elements,Atom_types,Geometry,Adj_mat)
    fig = plt.figure(figsize=(12,10))
    ax = plt.subplot(111)
    linewidth = 5
    fontsize = 50
    for count_i,i in enumerate(Adj_mat):
       for count_j,j in enumerate(i):
           if count_j > count_i:
               if j == 1:
                   ax.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

    for count_i,i in enumerate(geo_2D):
       ax.text(i[0], i[1], Elements[count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0})
    ax.axis('image')
    ax.axis('off')
    plot_name = 'test'
    plt.savefig(plot_name, dpi=300,bbox_inches='tight')
    plt.close(fig)
    quit()

# This function collects the modes being parsed, the location of the relvant scan data for each mode, and the relvant atom labeling for each scan geometry.
# The function also performs some basic checks to ensure that everything is in place for the parse. 
def parse_run_dirs(base_name=None,parse_modes=[],FF_dict=[],modes_from_FF=[],xyz_fit=[],undefined_modes=[],gens=2):    

   #### WE ARE IN INTRA FOLDER    

    # Change to the user supplied directory (if supplied)
    current_dir = os.getcwd()        
    if base_name != None:
        os.chdir(base_name) 

    # If mode_files are present then the script will determine fitting requirements based on these. 
    # If no *.modes files are present then the script will automatically determine the modes being
    # fit based upon the structure of the base_name folder.



    ####### ORIGINAL driver.py : when we don't have out_inchi.modes yet   

    #mode_files = [ '../'+i  for i in os.listdir('..') if os.path.isdir(i) == False and fnmatch.fnmatch(i,"*.modes") ]



    ####### CURRENT driver.py: when we have out_inchi.modes Not sure this is gonna work for higher generation
    ####### Not sure if when we have inchi_geoopt folder instead of geoopt folder will cause a problem or not 
    inchikey = (os.getcwd()).split("/")[-2]
    if(len(inchikey)) != 27:
      print("ERROR: not an inchikey folder")
      quit()
    #inchi_folders = ['../../'+i for i in os.listdir('../../') if os.path.isdir(i) == True]
    #print(inchi_folders)
    #mode_files = ['../../'+i+'/'+j for i in inchi_folders for j in os.listdir(i) if os.path.isdir(j) == False and fnmatch.fnmatch(j,"*_"+inchikey+".modes")]
    mode_files = [ '../../'+i+'/'+j  for i in os.listdir('../..') if os.path.isdir('../../'+i) == True for j in os.listdir('../../'+i) if os.path.isdir('../../'+i+'/'+j) == False and fnmatch.fnmatch(j,"*_"+inchikey+".modes") ]
    #print(mode_files)




    
    
    # If no *.modes files were discovered then default to parsing based upon the directory structure
    if len(mode_files) == 0:

        if os.path.isdir("geoopt") == False:
            print("ERROR in parse_run_dirs: No modes file were found and the geoopt folder is missing (needed for the fallback protocol). Exiting...")
            quit()
        else:
            print("ERROR: this function is not completed yet :<")
            quit()
            geo_folders = ["geoopt"]
            geo_atomtypes = [ get_scan_atomtypes('bonds_angles/bonds_angles_master.xyz') for i in geo_folders ]
            master_files = [ 'bonds_angles/bonds_angles_master.xyz' for i in geo_folders ]

        if os.path.isdir("bonds") == False:
            bond_folders = []
        else:
            bond_folders = [ 'bonds/'+i for i in os.listdir('bonds') if os.path.isdir('bonds/'+i) ]
            bond_types   = [ tuple(i.split('_')) for i in os.listdir('bonds') if os.path.isdir('bonds/'+i) ] 
            bond_atomtypes = [ get_scan_atomtypes(i) for i in bond_folders ]

        if os.path.isdir("angles") == False:
            angle_folders = []
        else:
            angle_folders = [ 'angles/'+i for i in os.listdir('angles') if os.path.isdir('angles/'+i) ]
            angle_types   = [ tuple(i.split('_')) for i in os.listdir('angles') if os.path.isdir('angles/'+i) ] 
            angle_atomtypes = [ get_scan_atomtypes(i) for i in angle_folders ]

        if os.path.isdir("dihedrals") == False:
            dihedral_folders = []
        else:
            dihedral_folders = [ 'dihedrals/'+i for i in os.listdir('dihedrals') if os.path.isdir('dihedrals/'+i) ]
            print(dihedral_folders)
            dihedral_types   = [ tuple(i.split('_')) for i in os.listdir('dihedrals') if os.path.isdir('dihedrals/'+i) ] 
            dihedral_atomtypes = [ get_scan_atomtypes(i.split("/")[0]) for i in dihedral_folders ]

        if os.path.isdir("dihedrals_harmonic") == False:
            harm_dihedral_folders = []
        else:
            harm_dihedral_folders = [ 'dihedrals_harmonic/'+i for i in os.listdir('dihedrals_harmonic') if os.path.isdir('dihedrals_harmonic/'+i) ]
            harm_dihedral_types   = [ tuple(i.split('_')) for i in os.listdir('dihedrals_harmonic') if os.path.isdir('dihedrals_harmonic/'+i) ] 
            harm_dihedral_atomtypes = [ get_scan_atomtypes(i) for i in harm_dihedral_folders ]

    # Parse protocols based upon the contents of the supplied *.modes files
    else:

        # Populate the folder and atomtype lists based on the contents of the discovered *.modes files
        bond_types,bond_atoms,bond_folders,bond_atomtypes,\
        angle_types,angle_atoms,angle_folders,angle_atomtypes,\
        dihedral_types,dihedral_atoms,dihedral_folders,dihedral_atomtypes,\
        harm_dihedral_types,harm_dihedral_atoms,harm_dihedral_folders,harm_dihedral_atomtypes,\
        bond_folders_frag,angle_folders_frag,dihedral_folders_frag,harm_dihedral_folders_frag = parse_mode_files(mode_files,parse_modes)
      

        # Grab the geo_files based upon the mode_files
        summed_folders = []
        add_folders = []
        main_inchi = mode_files[0].split('_')[-1].split('.')[-2] 
        if "bonds" in parse_modes: 
            summed_folders += bond_folders
        if "angles" in parse_modes: 
            summed_folders += angle_folders
        if "harm_dihedrals" in parse_modes: 
            summed_folders += harm_dihedral_folders
        if "dihedrals" in parse_modes: 
            summed_folders += dihedral_folders
        geo_folders = natural_sort(set([ i.split('/')[0]+'/geoopt' for i in summed_folders ]))
        if "bonds" in parse_modes: 
            other_inchi = [ i  for i in bond_folders_frag  if i != main_inchi]
            add_folders += [ '../../'+i+'/Intra/'+j for i in other_inchi for j in geo_folders if fnmatch.fnmatch(j,'*bonds*')]
        if "angles" in parse_modes: 
            other_inchi = [ i  for i in angle_folders_frag  if i != main_inchi]
            add_folders += [ '../../'+i+'/Intra/'+j for i in other_inchi for j in geo_folders if fnmatch.fnmatch(j,'*angles*')]
        if "harm_dihedrals" in parse_modes: 
            other_inchi = [ i  for i in harm_dihedral_folders_frag  if i != main_inchi]
            add_folders += [ '../../'+i+'/Intra/'+j for i in other_inchi for j in geo_folders if fnmatch.fnmatch(j,'*harm_dihedral*')]
        if "dihedrals" in parse_modes: 
            other_inchi = [ i  for i in dihedral_folders_frag  if i != main_inchi]
            add_folders += [ '../../'+i+'/Intra/'+j for i in other_inchi for j in geo_folders if fnmatch.fnmatch(j,'*dihedrals*')]

        geo_folders = natural_sort(set([ i.split('/')[0]+'/geoopt' for i in summed_folders ] + add_folders))
        # geo_folders has to include other frag's geo, but master.xyz would only be in the main inchi
        # that's why there's old_geo_folders and geo_folders
        # this may need further revision
        
     
          

            
        old_geo_folders = natural_sort(set([ i.split('/')[0]+'/geoopt' for i in summed_folders ]))
        #print("old")
        #print(old_geo_folders)

            

        # Grab geo_atomtypes (in the *master.xyz file)
        #geo_atomtypes = [ get_scan_atomtypes('bonds_angles/bonds_angles_master.xyz') for i in geo_folders ]


        # Save a master geometry for each fragment
        #master_files = [ next( '/'.join(i.split('/')[:-1])+'/'+j for j in os.listdir('/'.join(i.split('/')[:-1])) if fnmatch.fnmatch(j,"*master.xyz") ) for i in geo_folders ]
        #master_files = [ '/'.join(i.split('/')[:-1])+'/'+j for i in old_geo_folders for j in os.listdir('/'.join(i.split('/')[:-1])) if fnmatch.fnmatch(j,"*master.xyz") ]
        master_files=[]
        for k in geo_folders:
            chain = k.split('/')
            if(len(chain)==6):
               master_files.append('/'.join(chain[:-1])+'/'+chain[-2]+'_'+chain[-4]+'_master.xyz')
            elif(len(chain)==2):
               master_files.append(chain[0]+'/'+chain[0]+'_'+main_inchi+'_master.xyz')

               
        #print("master")
        #print(master_files)
        # an example of master_files
        #['../../LCGLNKUTAGEVQW-UHFFFAOYSA-N/Intra/bonds_angles/bonds_angles_LCGLNKUTAGEVQW-UHFFFAOYSA-N_master.xyz', '../../LFQSCWFLJHTTHZ-UHFFFAOYSA-N/Intra/bonds_angles/bonds_angles_LFQSCWFLJHTTHZ-UHFFFAOYSA-N_master.xyz', 'bonds_angles/bonds_angles_XOBKSJJDNFUZPF-UHFFFAOYSA-N_master.xyz']
        print(master_files)
        geo_atomtypes= [get_scan_atomtypes(i) for i in master_files] 
         
        #print(master_files)

    # Check for /REDO folders in the parsed directories XXX THIS IS DEPRECATED XXX
    for count_i,i in enumerate(angle_folders):
        while os.path.isdir(angle_folders[count_i]+'/REDO'):
            angle_folders[count_i] = angle_folders[count_i] + '/REDO'

    # Grab undefined modes from the supplied FF files
    undefined_files = [ '/'.join([dp,files]) for dp,dn,fn in os.walk('.') for files in fn if fnmatch.fnmatch(files,'undefined.db') ]
    _,skip_modes,_,_ = parse_FF_params(undefined_files,["dft"])

    # Generate composite set of modes that should be skipped during the parametrization
    modes_from_FF = set(modes_from_FF)
    skip_modes = set(skip_modes)
    skip_modes.update(modes_from_FF)

    # Remove types/folders that correspond to modes that are already in the modes_from_FF list
    try:
        bond_types,bond_atoms,bond_folders,bond_atomtypes = list(zip(*[ (bond_types[count_i],bond_atoms[count_i],bond_folders[count_i],bond_atomtypes[count_i]) for count_i,i in enumerate(bond_types) if i not in skip_modes ]))
        bond_types,bond_atoms,bond_folders,bond_atomtypes = (list(bond_types),list(bond_atoms),list(bond_folders),list(bond_atomtypes))
    except ValueError:
        bond_types,bond_atoms,bond_folders,bond_atomtypes = [],[],[],[]

    try:
        angle_types,angle_atoms,angle_folders,angle_atomtypes = list(zip(*[ (angle_types[count_i],angle_atoms[count_i],angle_folders[count_i],angle_atomtypes[count_i]) for count_i,i in enumerate(angle_types) if i not in skip_modes ]))
        angle_types,angle_atoms,angle_folders,angle_atomtypes = (list(angle_types),list(angle_atoms),list(angle_folders),list(angle_atomtypes))
    except ValueError:
        angle_types,angle_atoms,angle_folders,angle_atomtypes = [],[],[],[]

    # Only exclude dihedrals if they are part of the undefined set
    # NOTE: For dihedrals sometimes the same fragment is scanned with different coincident dihedrals, requiring the jobs to be included
    try:
        dihedral_types,dihedral_atoms,dihedral_folders,dihedral_atomtypes = \
        list(zip(*[ (dihedral_types[count_i],dihedral_atoms[count_i],dihedral_folders[count_i],dihedral_atomtypes[count_i]) for count_i,i in enumerate(dihedral_types) if i+tuple(["opls"]) not in undefined_modes ]))
        dihedral_types,dihedral_atoms,dihedral_folders,dihedral_atomtypes = (list(dihedral_types),list(dihedral_atoms),list(dihedral_folders),list(dihedral_atomtypes))
    except ValueError:
        dihedral_types,dihedral_atoms,dihedral_folders,dihedral_atomtypes = [],[],[],[]

    try:
        harm_dihedral_types,harm_dihedral_atoms,harm_dihedral_folders,harm_dihedral_atomtypes = \
        list(zip(*[ (harm_dihedral_types[count_i],harm_dihedral_atoms[count_i],harm_dihedral_folders[count_i],harm_dihedral_atomtypes[count_i]) for count_i,i in enumerate(harm_dihedral_types) if i+tuple(["harmonic"]) not in undefined_modes ])) 
        harm_dihedral_types,harm_dihedral_atoms,harm_dihedral_folders,harm_dihedral_atomtypes = (list(harm_dihedral_types),list(harm_dihedral_atoms),list(harm_dihedral_folders),list(harm_dihedral_atomtypes))
    except ValueError:
        harm_dihedral_types,harm_dihedral_atoms,harm_dihedral_folders,harm_dihedral_atomtypes = [],[],[],[]

    # Check for duplicate dihedral types. This occasionally arises when the non-main dihedrals vary between scans.
    # Having duplicate dihedral types will create issues in the dihedral fit function.
    for count_i,i in enumerate(dihedral_types):
        
        if len([ j for j in dihedral_types if j == i ]) == 1:
            continue
        else:
            # Collect the bonds and angles necessary to fit the dihedrals
            atomtypes  = dihedral_atomtypes[count_i]
            geo_folder = dihedral_folders[count_i].split('/')[0]+'/geoopt'
            xyz_file   = next( '/'.join(geo_folder.split('/')[:-1])+'/'+j for j in os.listdir('/'.join(geo_folder.split('/')[:-1])) if fnmatch.fnmatch(j,"*master.xyz") )

            # Parse the modes in the current xyz_fit file
            elements,geo  = xyz_parse(xyz_file)
            adj_mat  = Table_generator(elements,geo,File=xyz_file)
            atomtypes = atomtypes

            # Find the mode types in this fragment and update the fit lists
            Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(adj_mat,atomtypes,return_all=1)            
            shared_dihedrals = [ (Dihedral_types[count_j],Dihedrals[count_j]) for count_j,j in enumerate(Dihedrals) \
                                 if ( j[1] == dihedral_atoms[count_i][1] and j[2] == dihedral_atoms[count_i][2] ) or ( j[2] == dihedral_atoms[count_i][1] and j[1] == dihedral_atoms[count_i][2] ) ]
            new_dihedral = [ j for j in shared_dihedrals if j[0] not in dihedral_types ]
            if len(new_dihedral) > 0:
                dihedral_types[count_i] = new_dihedral[0][0]
                dihedral_atoms[count_i] = new_dihedral[0][1]

    # If xyz_fit option is supplied then only the folders associated with the xyz_fit geometries are retained.
    if len(xyz_fit) > 0:
        print("Sorry this function is incompleted")
        quit()

        # Change to work directory
        os.chdir(current_dir)

        # Loop over the supplied xyz_fit molecules and parse the modes
        fit_bondtypes = []
        fit_angletypes = []
        fit_dihedraltypes = []    
        for count_i,i in enumerate(xyz_fit):

            # On the first iteration, find the number of gens based on the atomtypes being parsed
            if count_i == 0:

                # Run a generation test on a set of all atomtypes involved in the mode scans that were discovered (probably not optimal but it is safe and doesn't cost much)
                unique_types = set([ k for j in (bond_types+angle_types+dihedral_types+harm_dihedral_types) for k in j ])
                gens = 0
                for j in unique_types:
                    tmp_gen = [ count_k for count_k,k in enumerate(j.split("[")) if "]" in k ]
                    if tmp_gen[0] > gens:
                        gens = tmp_gen[0]
                gens = gens - 1

            # Parse the modes in the current xyz_fit file
            fit_elements,fit_geo  = xyz_parse(i)
            fit_adj_mat  = Table_generator(fit_elements,fit_geo,File=i)
            fit_atomtypes = id_types(fit_elements,fit_adj_mat,gens=gens,geo=fit_geo)

            # Find the mode types in this fragment and update the fit lists
            tmp_Bonds,tmp_Angles,tmp_Dihedrals,tmp_One_fives,tmp_Bond_types,tmp_Angle_types,tmp_Dihedral_types,tmp_One_five_types = Find_modes(fit_adj_mat,fit_atomtypes,return_all=1)
            fit_bondtypes     += tmp_Bond_types
            fit_angletypes    += tmp_Angle_types
            fit_dihedraltypes += tmp_Dihedral_types            

        # Return to the base_name directory (if supplied)
        if base_name != None:
            os.chdir(base_name)

        # Create sets of the unique modes in the supplied xyz_fit files
        # fit_bondtypes = set(fit_bondtypes)
        # fit_angletypes = set(fit_angletypes)
        # fit_dihedraltypes = set(fit_dihedraltypes)
        while 1:

            N_dihedrals = len(fit_dihedraltypes)

            # Remove bonds that aren't in the fit_xyz molecules
            keep_ind = [ count_i for count_i,i in enumerate(bond_types) if i in fit_bondtypes ]
    #        keep_ind = [ count_i for count_i,i in enumerate(bond_types) if i in fit_bondtypes or True in [ j == bond_folders[count_i].split("/")[0].split("_")[0] for j in dihedral_keep_strings ] ]
            tmp_bond_types = [ bond_types[i] for i in keep_ind ]
            tmp_bond_folders = [ bond_folders[i] for i in keep_ind ]
            tmp_bond_atomtypes = [ bond_atomtypes[i] for i in keep_ind ]

            # Remove angles that aren't in the fit_xyz molecules
            keep_ind = [ count_i for count_i,i in enumerate(angle_types) if i in fit_angletypes ]
    #        keep_ind = [ count_i for count_i,i in enumerate(angle_types) if i in fit_angletypes or True in [ j == angle_folders[count_i].split("/")[0].split("_")[0] for j in dihedral_keep_strings ] ]
            tmp_angle_types = [ angle_types[i] for i in keep_ind ]
            tmp_angle_folders = [ angle_folders[i] for i in keep_ind ]
            tmp_angle_atomtypes = [ angle_atomtypes[i] for i in keep_ind ]

            # Remove dihedrals that aren't in the fit_xyz molecules
            keep_ind = [ count_i for count_i,i in enumerate(dihedral_types) if i in fit_dihedraltypes ]
            tmp_dihedral_types = [ dihedral_types[i] for i in keep_ind ]
            tmp_dihedral_folders = [ dihedral_folders[i] for i in keep_ind ]
            tmp_dihedral_atomtypes = [ dihedral_atomtypes[i] for i in keep_ind ]

            # Remove harm_dihedrals that aren't in the fit_xyz molecules
            keep_ind = [ count_i for count_i,i in enumerate(harm_dihedral_types) if i in fit_dihedraltypes ]
            tmp_harm_dihedral_types = [ harm_dihedral_types[i] for i in keep_ind ]
            tmp_harm_dihedral_folders = [ harm_dihedral_folders[i] for i in keep_ind ]
            tmp_harm_dihedral_atomtypes = [ harm_dihedral_atomtypes[i] for i in keep_ind ]

            # Collect the bonds and angles necessary to fit the dihedrals
            mode_atomtypes = tmp_bond_atomtypes+tmp_angle_atomtypes+tmp_dihedral_atomtypes+tmp_harm_dihedral_atomtypes
            geo_folders = [ i.split('/')[0]+'/geoopt' for i in (tmp_bond_folders+tmp_angle_folders+tmp_dihedral_folders+tmp_harm_dihedral_folders) ]
            for count_i,i in enumerate(geo_folders):
                xyz_file = next( '/'.join(i.split('/')[:-1])+'/'+j for j in os.listdir('/'.join(i.split('/')[:-1])) if fnmatch.fnmatch(j,"*master.xyz") )

                # Parse the modes in the current xyz_fit file
                fit_elements,fit_geo  = xyz_parse(xyz_file)
                fit_adj_mat  = Table_generator(fit_elements,fit_geo,File=xyz_file)
                fit_atomtypes = mode_atomtypes[count_i]

                # Find the mode types in this fragment and update the fit lists
                tmp2_Bonds,tmp2_Angles,tmp2_Dihedrals,tmp2_One_fives,tmp2_Bond_types,tmp2_Angle_types,tmp2_Dihedral_types,tmp2_One_five_types = Find_modes(fit_adj_mat,fit_atomtypes,return_all=1)            
                fit_bondtypes += tmp2_Bond_types
                fit_angletypes += tmp2_Angle_types
                fit_dihedraltypes += tmp2_Dihedral_types

                fit_bondtypes = list(set(fit_bondtypes))
                fit_angletypes = list(set(fit_angletypes)) 
                fit_dihedraltypes = list(set(fit_dihedraltypes))                

            # Break once convergence is reached
            if N_dihedrals == len(fit_dihedraltypes):
                break

        # Remove bonds that aren't in the fit_xyz molecules
        keep_ind = [ count_i for count_i,i in enumerate(bond_types) if i in fit_bondtypes ]
#        keep_ind = [ count_i for count_i,i in enumerate(bond_types) if i in fit_bondtypes or True in [ j == bond_folders[count_i].split("/")[0].split("_")[0] for j in dihedral_keep_strings ] ]
        bond_types = [ bond_types[i] for i in keep_ind ]
        bond_folders = [ bond_folders[i] for i in keep_ind ]
        bond_atomtypes = [ bond_atomtypes[i] for i in keep_ind ]

        # Remove angles that aren't in the fit_xyz molecules
        keep_ind = [ count_i for count_i,i in enumerate(angle_types) if i in fit_angletypes ]
#        keep_ind = [ count_i for count_i,i in enumerate(angle_types) if i in fit_angletypes or True in [ j == angle_folders[count_i].split("/")[0].split("_")[0] for j in dihedral_keep_strings ] ]
        angle_types = [ angle_types[i] for i in keep_ind ]
        angle_folders = [ angle_folders[i] for i in keep_ind ]
        angle_atomtypes = [ angle_atomtypes[i] for i in keep_ind ]

        # Remove dihedrals that aren't in the fit_xyz molecules
        keep_ind = [ count_i for count_i,i in enumerate(dihedral_types) if i in fit_dihedraltypes ]
        dihedral_types = [ dihedral_types[i] for i in keep_ind ]
        dihedral_folders = [ dihedral_folders[i] for i in keep_ind ]
        dihedral_atomtypes = [ dihedral_atomtypes[i] for i in keep_ind ]

        # Remove harm_dihedrals that aren't in the fit_xyz molecules
        keep_ind = [ count_i for count_i,i in enumerate(harm_dihedral_types) if i in fit_dihedraltypes ]
        harm_dihedral_types = [ harm_dihedral_types[i] for i in keep_ind ]
        harm_dihedral_folders = [ harm_dihedral_folders[i] for i in keep_ind ]
        harm_dihedral_atomtypes = [ harm_dihedral_atomtypes[i] for i in keep_ind ]

        # Grab the geo_files based upon the mode_files
        summed_folders = []
        if "bonds" in parse_modes: summed_folders += bond_folders
        if "angles" in parse_modes: summed_folders += angle_folders
        if "harm_dihedrals" in parse_modes: summed_folders += harm_dihedral_folders
        if "dihedrals" in parse_modes: summed_folders += dihedral_folders
        geo_folders = natural_sort(set([ i.split('/')[0]+'/geoopt' for i in summed_folders ]))

        # Grab geo_atomtypes (in the *master.xyz file)
        #geo_atomtypes = [ get_scan_atomtypes(next( '/'.join(i.split('/')[:-1])+'/'+j for j in os.listdir('/'.join(i.split('/')[:-1])) if fnmatch.fnmatch(j,"*master.xyz") )) for i in geo_folders ]
        geo_atomtypes = [ get_scan_atomtypes(['bonds_angles/bonds_angles_master.xyz']) for i in geo_folders ]

        # Save a master geometry for each fragment
        master_files = [ next( '/'.join(i.split('/')[:-1])+'/'+j for j in os.listdir('/'.join(i.split('/')[:-1])) if fnmatch.fnmatch(j,"*master.xyz") ) for i in geo_folders ]

    # Check that all *_folders exist
    missing  = [ i for i in geo_folders if os.path.isdir(i) == False ]
    if "bonds" in parse_modes: missing += [ i for i in bond_folders if os.path.isdir(i) == False ]
    if "angles" in parse_modes: missing += [ i for i in angle_folders if os.path.isdir(i) == False ]
    if "harm_dihedrals" in parse_modes: missing += [ i for i in harm_dihedral_folders if os.path.isdir(i) == False ]
    if "dihedrals" in parse_modes: missing += [ i for i in dihedral_folders if os.path.isdir(i) == False ]
    if len(missing) > 0:
        print("\nERROR in parse_run_dirs: The following scan folders are missing. Exiting...\n")            
        for i in missing:
            print("\t{}".format(i))
        print("")
        quit()

    # Check that all *_folders contain complete run data
    incomplete  = [ i for i in geo_folders if i not in missing and check_scan_completeness(i) == False ]
    if "bonds" in parse_modes: incomplete += [ i for i in bond_folders if i not in missing and check_scan_completeness(i) == False ]
    if "angles" in parse_modes: incomplete += [ i for i in angle_folders if i not in missing and check_scan_completeness(i) == False ]
    if "harm_dihedrals" in parse_modes: incomplete += [ i for i in harm_dihedral_folders if i not in missing and check_scan_completeness(i) == False ]
    print("current:{}".format(os.getcwd()))
    if "dihedrals" in parse_modes: incomplete += [ i for i in dihedral_folders if i not in missing and check_scan_completeness(i) == False ]
    if len(incomplete) > 0:
        print("\nERROR in parse_run_dirs: The following scan folders are incomplete. Exiting...\n")            
        for i in incomplete:
            print("\t{}".format(i))
        print("")
        quit()
    
    os.chdir(current_dir)

    return master_files,geo_folders,geo_atomtypes,\
           bond_types,bond_folders,bond_atomtypes,\
           angle_types,angle_folders,angle_atomtypes,\
           dihedral_types,dihedral_folders,dihedral_atomtypes,\
           harm_dihedral_types,harm_dihedral_folders,harm_dihedral_atomtypes,\
           bond_folders_frag,angle_folders_frag,dihedral_folders_frag,harm_dihedral_folders_frag

# This function checks the completeness of scan data. (currently only compatible with ORCA)
def check_scan_completeness(folder):

    # If "folder" is actually an output file then the following loop handles the completness determination
    if os.path.isdir(folder) == False:
        complete_flag = 0
        with open(folder,'r') as f:
            for lines in f:
                if "****ORCA TERMINATED NORMALLY****" in lines:
                    complete_flag = 1
        if complete_flag == 0:
            return False
        else:
            return True

    # Check the completeness of any output files in the outermost level of the trajectory (NOTE: the function doesn't expect a file at this level.)
    output_files = [ folder+'/'+j for j in os.listdir(folder) if fnmatch.fnmatch(j,"*.out") ]
    outer_parse = 0
    if len(output_files) > 0:
        outer_parse = 1
        for j in output_files:
            complete_flag = 0
            with open(j,'r') as f:
                for lines in f:
                    if "****ORCA TERMINATED NORMALLY****" in lines:
                        complete_flag = 1
            if complete_flag == 0:
                return False

    # Find the sub_folders within the supplied trajectory
    sub_folders = [ i for i in os.listdir(folder) if os.path.isdir(folder+'/'+i) == True ]

    # Else, the following loop(s) handle the completeness determination for a scan directory
    inner_parse = 0
    for i in sub_folders:
        inner_parse = 1
        output_files = [ folder+'/'+i+'/'+j for j in os.listdir(folder+'/'+i) if fnmatch.fnmatch(j,"*.out") ]
        if len(output_files) == 0:            
            return False
        else:
            for j in output_files:
                complete_flag = 0
                with open(j,'r') as f:
                    for lines in f:
                        if "****ORCA TERMINATED NORMALLY****" in lines:
                            complete_flag = 1
                if complete_flag == 0:
                    return False

    # If a file was actually parsed and didn't no incompleteness was identified (evidenced by the 
    # fact that this point in the function is reached) then return a True
    if inner_parse == 1 or outer_parse == 1:
        return True

    # This else statement catches situations where no output files are discovered in the outer folder and there are no subdirectories
    else:
        return False

# A wrapper for the commands to return the folder that contains the file named "name"
def containing_folder(name):
    folder = "/".join(name.split('/')[:-1])
    if len(folder) == 0:
        return './'
    else:
        return folder

# Grab the modes being scanned from a -*.modes file
def parse_mode_files(mode_files,parse_modes):

    bond_types              = []
    bond_atoms              = []
    bond_atomtypes          = []
    bond_folders            = []
    bond_folders_frag       = []
    angle_types             = []
    angle_atoms             = []
    angle_atomtypes         = []
    angle_folders           = []
    angle_folders_frag      = []
    dihedral_types          = []
    dihedral_atoms          = []
    dihedral_atomtypes      = []
    dihedral_folders        = []
    dihedral_folders_frag   = []
    harm_dihedral_types     = []
    harm_dihedral_atoms     = []
    harm_dihedral_atomtypes = []
    harm_dihedral_folders   = []
    harm_dihedral_folders_frag   = []
    bond_parse              = 0
    angle_parse             = 0
    dihedral_parse          = 0
    harm_dihedral_parse     = 0
    for m in mode_files:
        frag_info = m.split('/')[-2]
        #print(frag_info)
        with open(m,'r') as f:
            for lines in f:
                fields = lines.split()

                # Commmands for parsing bonds are contained in this series of if/elif statements
                if  len(fields) == 2 and fields[0] == "bond" and fields[1] == "start" and "bonds" in parse_modes:
                    bond_parse = 1
                    continue
                elif len(fields) == 2 and fields[0] == "bond" and fields[1] == "end" and "bonds" in parse_modes:
                    bond_parse = 0                    
                    for count_j,j in enumerate(mode_ind):
                        #bond_folders += [next( i for i in os.listdir(containing_folder(m)) if os.path.isdir(i) == True and i.split("_")[0] == ".modes".join(m.split('.modes')[:-1]) and "bonds" in i.split("_") )]
                        bond_folders += ['bonds_angles']
                        bond_folders[-1] = bond_folders[-1] + '/bonds/' +  "_".join([ atom_types[0][i] for i in mode_ind[0] ])
                        bond_folders_frag.append(frag_info)
                        bond_types += [ tuple([ atom_types[count_j][i] for i in j ]) ]
                        bond_atomtypes += [ atom_types[count_j] ]
                        bond_atoms += [j]
                    continue
                elif bond_parse == 1:
                    bond_parse += 1
                    continue
                elif bond_parse == 2:
                    mode_ind = [ [ int(j) for j in i.split('_') ] for i in fields ]
                    atom_types = [ [] for i in mode_ind ]
                    bond_parse += 1
                elif bond_parse == 3:
                    for count_i,i in enumerate(fields):
                        atom_types[count_i] += [i]
                    continue

                # Commmands for parsing angles are contained in this series of if/elif statements
                if  len(fields) == 2 and fields[0] == "angle" and fields[1] == "start" and "angles" in parse_modes:
                    angle_parse = 1
                    continue
                elif len(fields) == 2 and fields[0] == "angle" and fields[1] == "end" and "angles" in parse_modes:
                    angle_parse = 0
                    for count_j,j in enumerate(mode_ind):
                        #angle_folders += [next( i for i in os.listdir(containing_folder(m)) if os.path.isdir(i) == True and i.split("_")[0] == ".modes".join(m.split('.modes')[:-1]) and "angles" in i.split("_") )]
                        angle_folders += ['bonds_angles']
                        angle_folders[-1] = angle_folders[-1] + '/angles/' +  "_".join([ atom_types[0][i] for i in mode_ind[0] ])
                        angle_folders_frag.append(frag_info)
                        angle_types += [ tuple([ atom_types[count_j][i] for i in j ]) ]
                        angle_atomtypes += [ atom_types[count_j] ]
                        angle_atoms += [j]
                    continue
                elif angle_parse == 1:
                    angle_parse += 1
                    continue
                elif angle_parse == 2:
                    mode_ind = [ [ int(j) for j in i.split('_') ] for i in fields ]
                    atom_types = [ [] for i in mode_ind ]
                    angle_parse += 1
                elif angle_parse == 3:
                    for count_i,i in enumerate(fields):
                        atom_types[count_i] += [i]
                    continue

                # Commmands for parsing dihedrals are contained in this series of if/elif statements
                if  len(fields) == 2 and fields[0] == "dihedral" and fields[1] == "start" and "dihedrals" in parse_modes:
                    dihedral_parse = 1
                    continue
                elif len(fields) == 2 and fields[0] == "dihedral" and fields[1] == "end" and "dihedrals" in parse_modes:
                    dihedral_parse = 0
                    for count_j,j in enumerate(mode_ind):
                        #dihedral_folders += [next( i for i in os.listdir(containing_folder(m)) if os.path.isdir(i) == True and i.split("_")[0] == ".modes".join(m.split('.modes')[:-1]) and "dihedrals" in i.split("_") )]
                        dihedral_folders += ['dihedrals']
                        dihedral_folders[-1] = dihedral_folders[-1] + '/dihedrals/' +  "_".join([ atom_types[0][i] for i in mode_ind[0] ])
                        dihedral_folders_frag.append(frag_info)
                        dihedral_types += [ tuple([ atom_types[count_j][i] for i in j ]) ]
                        dihedral_atomtypes += [ atom_types[count_j] ]
                        dihedral_atoms += [j]
                    continue
                elif dihedral_parse == 1:
                    dihedral_parse += 1
                    continue
                elif dihedral_parse == 2:
                    mode_ind = [ [ int(j) for j in i.split('_') ] for i in fields ]
                    atom_types = [ [] for i in mode_ind ]
                    dihedral_parse += 1
                elif dihedral_parse == 3:                    
                    for count_i,i in enumerate(fields):
                        atom_types[count_i] += [i]
                    continue

                # Commmands for parsing harmonic dihedrals are contained in this series of if/elif statements
                if  len(fields) == 2 and fields[0] == "harmonic_dihedral" and fields[1] == "start" and "harm_dihedrals" in parse_modes:
                    harm_dihedral_parse = 1
                    continue
                elif len(fields) == 2 and fields[0] == "harmonic_dihedral" and fields[1] == "end" and "harm_dihedrals" in parse_modes:
                    harm_dihedral_parse = 0
                    for count_j,j in enumerate(mode_ind):
                        #harm_dihedral_folders += [next( i for i in os.listdir(containing_folder(m)) if os.path.isdir(i) == True and i.split("_")[0] == ".modes".join(m.split('.modes')[:-1]) and "dihedrals" in i.split("_") )]
                        harm_dihedral_folders += ['dihedrals']
                        harm_dihedral_folders[-1] = harm_dihedral_folders[-1] + '/dihedrals_harmonic/' +  "_".join([ atom_types[0][i] for i in mode_ind[0] ])
                        harm_dihedral_folders_frag.append(frag_info)
                        harm_dihedral_types += [ tuple([ atom_types[count_j][i] for i in j ]) ]
                        harm_dihedral_atomtypes += [ atom_types[count_j] ]
                        harm_dihedral_atoms += [j]
                    continue
                elif harm_dihedral_parse == 1:
                    harm_dihedral_parse += 1
                    continue
                elif harm_dihedral_parse == 2:
                    mode_ind = [ [ int(j) for j in i.split('_') ] for i in fields ]
                    atom_types = [ [] for i in mode_ind ]
                    harm_dihedral_parse += 1
                elif harm_dihedral_parse == 3:                    
                    for count_i,i in enumerate(fields):
                        atom_types[count_i] += [i]
                    continue

    return bond_types,bond_atoms,bond_folders,bond_atomtypes,\
           angle_types,angle_atoms,angle_folders,angle_atomtypes,\
           dihedral_types,dihedral_atoms,dihedral_folders,dihedral_atomtypes,\
           harm_dihedral_types,harm_dihedral_atoms,harm_dihedral_folders,harm_dihedral_atomtypes,\
           bond_folders_frag,angle_folders_frag,dihedral_folders_frag,harm_dihedral_folders_frag

# Scrap the atomtypes from the scan folder 
# If *.modes files are not supplied then the atomtypes are parsed directly from the scan_folder *.xyz file generated by paramgen.py
def get_scan_atomtypes(folder_name):

    # If the supplied "folder_name" is actually an xyz file then the atomtypes are parsed directly from the file        
    if len(folder_name) > 3 and folder_name[-4:] == ".xyz":
        geo = folder_name    
        if os.path.isfile(geo) == False:
            print("ERROR in get_scan_atomtypes: file {} does not exist. Exiting...".format(geo))
            quit()

    # Else the geometry file is parsed based upon the expected hierarchy of the mode scan(s)
    else:
        # Need modified
        sub =  next( folder_name+'/'+i for i in os.listdir(folder_name) if os.path.isdir(folder_name+'/'+i) == True )
        print(sub)
        #print(folder_name)
        #tmp_nex =  [folder_name+'/'+i for i in os.listdir(folder_name) if os.path.isdir(folder_name+'/'+i) == True]
        #print(tmp_nex)
        geo = [ sub+'/'+i for i in os.listdir(sub) if fnmatch.fnmatch(i,"*"+sub.split('/')[0]+"*.xyz") ]
        geo = str(geo)

    # Parse the atomtypes
    atomtypes = []
    with open(geo,'r') as f:
        for lc,lines in enumerate(f):
            if lc > 1:
                fields = lines.split()
                if len(fields) == 0:
                    break
                elif len(fields) < 5:
                    print("ERROR in get_scan_atomtypes: file {} does not have atomtype information. Exiting...".format(geo))
                    quit()
                else:
                    atomtypes += [fields[4]]
                
    return atomtypes

# Used by curve_fit for bond and angle fits
# x range is a global variable that ensures the x0 value lies within the sampled range
def harmonic(x,k,x0):
    global harm_range
    if x0 < harm_range[0] or x0 > harm_range[1]: return k*(x-x0)**2 * 1.0E10
    return k*(x-x0)**2.0

# Used by curve_fit for bond and angle fits
# x range is a global variable that ensures the x0 value lies within the sampled range
def harmonic_norm(x,k,x0,norm):
    global harm_range
    if x0 < harm_range[0] or x0 > harm_range[1]: return k*(x-x0)**2 * 1.0E10
    return norm+k*(x-x0)**2.0

# Used by curve_fit for bond and angle fits
# x range is a global variable that ensures the x0 value lies within the sampled range
def dihedral_harmonic_norm(x,k,d,n,norm):
    return norm+k*(1.0+d*np.cos(n*x))

# Used by curve_fit for bond and angle fits
# x range is a global variable that ensures the x0 value lies within the sampled range
def dihedral_harmonic_nonorm(x,k,d,n,norm):
    return norm+k*(1.0+d*np.cos(n*x))-norm

# Used by curve_fit for fitting single dihedrals
def opls(x,v1,v2,v3,v4):
    return 0.5*v1*(1+np.cos(x)) + 0.5*v2*(1-np.cos(2.0*x)) + 0.5*v3*(1+np.cos(3.0*x)) + 0.5*v4*(1-np.cos(4.0*x))

def opls_5(x,v0,v1,v2,v3,v4):
    return v0 + 0.5*v1*(1+np.cos(x)) + 0.5*v2*(1-np.cos(2.0*x)) + 0.5*v3*(1+np.cos(3.0*x)) + 0.5*v4*(1-np.cos(4.0*x))

# Used by curve_fit for fitting and arbitrary number of coincident *of the same type*
# Since all dihedrals are of the same type only a single set of Vs is used for the fit.
# The user must pre-assign each dihedrals angles to a global tmp_angles array so that 
# the function can access it.
def opls_many(x,v1,v2,v3,v4):
    result = 0
    for i in range(len(tmp_angles[0,:])):
        result += 0.5*v1*(1+np.cos(tmp_angles[x,i])) + 0.5*v2*(1-np.cos(2.0*tmp_angles[x,i])) + 0.5*v3*(1+np.cos(3.0*tmp_angles[x,i])) + 0.5*v4*(1-np.cos(4.0*tmp_angles[x,i]))
    return result

def tryint(s):
    try:
        return int(s)
    except:
        return s
     
def alphanum_key(s):
    """ Turn a string into a list of string and number chunks.
        "z23a" -> ["z", 23, "a"]
    """
    return [ tryint(c) for c in re.split('([0-9]+)', s) ]

def sort_nicely(l):
    """ Sort the given list in the way that humans expect.
    """
    l.sort(key=alphanum_key)        

# # Description: Initialize VDW_dict based on UFF parameters for the initial guess of the fit.
# def initialize_VDW(atomtypes):

#     # Initialize UFF parameters (tuple corresponds to eps,sigma pairs for each element)
#     # Taken from UFF (Rappe et al. JACS 1992)
#     # Note: LJ parameters in the table are specificed in the eps,r_min form rather than eps,sigma
#     #       the conversion between r_min and sigma is sigma = r_min/2^(1/6)
#     # Note: Units for sigma = angstroms and eps = kcal/mol 
#     UFF_dict = { 1:(0.044,2.5711337005530193),  2:(0.056,2.1043027722474816),  3:(0.025,2.183592758161972),  4:(0.085,2.4455169812952313),\
#                  5:(0.180,3.6375394661670053),  6:(0.105,3.4308509635584463),  7:(0.069,3.260689308393642),  8:(0.060,3.1181455134911875),\
#                  9:(0.050,2.996983287824101),  10:(0.042,2.88918454292912),   11:(0.030,2.657550876212632), 12:(0.111,2.6914050275019648),\
#                 13:(0.505,4.008153332913386),  14:(0.402,3.82640999441276),   15:(0.305,3.694556984127987), 16:(0.274,3.594776327696269),\
#                 17:(0.227,3.5163772404999194), 18:(0.185,3.4459962417668324), 19:(0.035,3.396105913550973), 20:(0.238,3.0281647429590133),\
#                 21:(0.019,2.935511276272418),  22:(0.017,2.828603430095577),  23:(0.016,2.800985569833227), 24:(0.015,2.6931868249382456),\
#                 25:(0.013,2.6379511044135446), 26:(0.013,2.5942970672246677), 27:(0.014,2.558661118499054), 28:(0.015,2.5248069672097215),\
#                 29:(0.005,3.113691019900486),  30:(0.124,2.4615531582217574), 31:(0.415,3.904809081609107), 32:(0.379,3.813046513640652),\
#                 33:(0.309,3.7685015777336357), 34:(0.291,3.746229109780127),  35:(0.251,3.731974730289881), 36:(0.220,3.689211591819145),\
#                 37:(0.040,3.6651573264293558), 38:(0.235,3.2437622327489755), 39:(0.072,2.980056212179435), 40:(0.069,2.78316759547042),\
#                 41:(0.059,2.819694442914174),  42:(0.056,2.7190228877643157), 43:(0.048,2.670914356984738), 44:(0.056,2.6397329018498255),\
#                 45:(0.053,2.6094423454330538), 46:(0.048,2.5827153838888437), 47:(0.036,2.804549164705788), 48:(0.228,2.537279549263686),\
#                 49:(0.599,3.976080979060334),  50:(0.567,3.9128271700723705), 51:(0.449,3.937772334180300), 52:(0.398,3.982317270087316),\
#                 53:(0.339,4.009044231631527),  54:(0.332,3.923517954690054),  55:(0.045,4.024189509839913), 56:(0.364,3.2989979532736764),\
#                 72:(0.072,2.798312873678806),  73:(0.081,2.8241489365048755), 74:(0.067,2.734168165972701), 75:(0.066,2.631714813386562),\
#                 76:(0.037,2.7796040005978586), 77:(0.073,2.5301523595185635), 78:(0.080,2.453535069758495), 79:(0.039,2.9337294788361374),\
#                 80:(0.385,2.4098810325696176), 81:(0.680,3.872736727756055),  82:(0.663,3.828191791849038), 83:(0.518,3.893227398273283),\
#                 84:(0.325,4.195242063722858),  85:(0.284,4.231768911166611),  86:(0.248,4.245132391938716) }

#     # Initialize VDW_dict first guess based on element types and Lorentz-Berthelot mixing rules
#     VDW_dict = {}
#     for count_i,i in enumerate(atomtypes):
#         for count_j,j in enumerate(atomtypes):
#             if count_i < count_j:
#                 continue

#             type_1 = int(i.split('[')[1])
#             type_2 = int(j.split('[')[1])
#             eps    = (UFF_dict[type_1][0]*UFF_dict[type_2][0])**(0.5)
#             sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0
#             if i > j:
#                 VDW_dict[(i,j)] = (eps,sigma)
#             else:
#                 VDW_dict[(j,i)] = (eps,sigma)

#     return VDW_dict


##############################
# Slide generation functions #
##############################
# Adds a slide with title and default blue line
# returns the slide handle
def add_slide(pres,title):

    blank_slide = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide)
    shapes = slide.shapes

    # Create title 
    txBox = slide.shapes.add_textbox(Inches(0.25),Inches(0.1),Inches(10.0),Inches(1.0))
    text_frame = txBox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    font = run.font
    font.name = 'Helvetica'
    font.size = Pt(40)
    font.bold = True

    # Create line 
    left = Inches(1.45)
    top = Inches(-2.5)
    width = Inches(7.1)  # lines are instantiated as diagonal lines. The slide width is 10 in. so np.cos(45)*10 are the necessary width and height.
    height = Inches(7.1) # lines are instantiated as diagonal lines. The slide width is 10 in. so np.cos(45)*10 are the necessary width and height.
    shape = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left, top, width, height)
    shape.rotation = 45.0
    shape.line.color.rgb = RGBColor(0,50,230)
    shape.line.width = Pt(3.0)

    return slide

# Adds a 2x2 array of pictures to a slide.
# Assumes all images have a 4x3 width to height ratio
def add_4(slide,names,h_pad=0.0,v_pad=0.0):
        
    top = 1.5 # Distance of the 2x2 array from the top of the slide in inches.
    left_start = 1 # Distance of the 2x2 array from the left edge of the slide in inches.
    top_start = 1.25
    counter = 0
    for i in range(2):
        for j in range(2):
            try:
                height = Inches(3)
                left = Inches(left_start+(4+h_pad)*j)
                top = Inches(top_start+(3+v_pad)*i)
                pic = slide.shapes.add_picture(names[counter], left, top, height=height)
                counter += 1
            except:
               break

# Adds a 3x3 array of pictures to a slide.
# Assumes all images have a 4x3 width to height ratio
def add_9(slide,names,h_pad=0.0,v_pad=0.0):
        
    top = 1.5 # Distance of the 3x3 array from the top of the slide in inches.
    left_start = 1 # Distance of the 3x3 array from the left edge of the slide in inches.
    top_start = 1.25
    h_pad = 0.0
    v_pad = 0.0
    counter = 0
    for i in range(3):
        for j in range(3):
            try:
                height = Inches(2)
                left = Inches(left_start+(2.67+h_pad)*j)
                top = Inches(top_start+(2+v_pad)*i)
                pic = slide.shapes.add_picture(names[counter], left, top, height=height)
                counter += 1
            except:
               break

# Description: A wrapper function for the commands to parse the elements and atom_types
#              from the *master.xyz file.            
#
# Inputs       filename:   The *master.xyz file
# Returns      geo:        A numpy array holding the optimized geometry
#              elements:   A list of the elements indexed to geo
#              atomtypes:  A list of the atomtypes indexed to geo
#              keep_ind:   A list of indices corresponding to the atoms involved in modes that
#                          were parametrized
#              keep_types: A list of types indexed to keep_ind corresponding to the atomtypes
#                          involved in modes that were parametrized.
def get_info(filename):

    # First find the number of atoms 
    with open(filename,'r') as xyz:
        for _ in xyz:
            N_atoms = int(_.split()[0])
            break

    # Loop over the globally optimized geometry and save it to an array
    geo = np.zeros([N_atoms,3])
    elements = ["X"]*N_atoms
    atomtypes = [0]*N_atoms
    atom_count = 0            
    with open(filename,'r') as xyz:
        for count,lines in enumerate(xyz):

            # Write the geometry and elements to variable
            if count >= 2:
                fields = lines.split()
                geo[atom_count,:] = np.array([float(fields[1]),float(fields[2]),float(fields[3])])
                elements[atom_count] = fields[0]
                atomtypes[atom_count] = fields[4]
                atom_count += 1

    return geo,elements,atomtypes

# Description: A wrapper function for the commands to parse the masses and charges
#              from the optimized geometry located in the first bond being parametrized.
#
# Inputs       elements:   A list of elements (from the *master.xyz file parsed prior)
#              bond_steps: The number of bond configurations used for the parametrization
#                          This is used here to id the optimized (unperturbed) geometry.
#
# Returns      charges:    A list of the CHELPG partial charges for each atom type (indexed to elements).
#              masses:     A list of the masses of each atomtypes (indexed to elements)
def get_q_and_m(geo_folders,geo_atomtypes,force_read_charges=0,zero_opt=0,unique_atomtypes=None,FF_db="",gens=2):

    # Initialize element to atomic number conversion
    atomic2ele = { 1:"H",   2:"He",\
                   3:"Li",  4:"Be",                                                                                                      5: "B",    6:"C",   7: "N",   8: "O",   9: "F",  10:"Ne",\
                  11:"Na", 12:"Mg",                                                                                                     13:"Al",  14:"Si",  15: "P",  16: "S",  17:"Cl",  18:"Ar",\
                  19:"K",  20:"Ca",  21:"Sc",  22:"Ti",  23: "V", 24:"Cr",  25:"Mn",  26:"Fe",  27:"Co",  28:"Ni",  29:"Cu",  30:"Zn",  31:"Ga",  32:"Ge",  33:"As",  34:"Se",  35:"Br",  36:"Kr",\
                  37:"Rb", 38:"Sr",  39: "Y",  40:"Zr",  41:"Nb", 42:"Mo",  43:"Tc",  44:"Ru",  45:"Rh",  46:"Pd",  47:"Ag",  48:"Cd",  49:"In",  50:"Sn",  51:"Sb",  52:"Te",  53: "I",  54:"Xe",\
                  55:"Cs", 56:"Ba",            72:"Hf",  73:"Ta", 74: "W",  75:"Re",  76:"Os",  77:"Ir",  78:"Pt",  79:"Au",  80:"Hg",  81:"Tl",  82:"Pb",  83:"Bi",  84:"Po",  85:"At",  86:"Rn"}

    # Initialize mass_dict and van der waals dictionaries (used as default values in several places).
    mass_dict = {'H':1.00794,'He':4.002602,'Li':6.941,'Be':9.012182,'B':10.811,'C':12.011,'N':14.00674,'O':15.9994,'F':18.9984032,'Ne':20.1797,\
                 'Na':22.989768,'Mg':24.3050,'Al':26.981539,'Si':28.0855,'P':30.973762,'S':32.066,'Cl':35.4527,'Ar':39.948,\
                 'K':39.0983,'Ca':40.078,'Sc':44.955910,'Ti':47.867,'V':50.9415,'Cr':51.9961,'Mn':54.938049,'Fe':55.845,'Co':58.933200,'Ni':58.6934,'Cu':63.546,'Zn':65.39,\
                 'Ga':69.723,'Ge':72.61,'As':74.92159,'Se':78.96,'Br':79.904,'Kr':83.80,\
                 'Rb':85.4678,'Sr':87.62,'Y':88.90585,'Zr':91.224,'Nb':92.90638,'Mo':95.94,'Tc':98.0,'Ru':101.07,'Rh':102.90550,'Pd':106.42,'Ag':107.8682,'Cd':112.411,\
                 'In':114.818,'Sn':118.710,'Sb':121.760,'Te':127.60,'I':126.90447,'Xe':131.29,\
                 'Cs':132.90545,'Ba':137.327,'La':138.9055,'Hf':178.49,'Ta':180.9479,'W':183.84,'Re':186.207,'Os':190.23,'Ir':192.217,'Pt':195.078,'Au':196.96655,'Hg':200.59,\
                 'Tl':204.3833,'Pb':207.2,'Bi':208.98038,'Po':209.0,'At':210.0,'Rn':222.0}

    #########################################################
    #  Process Charges from the Geometry Optimization Runs  #
    #########################################################

    # Check if the force_read option is set, if so then the charges will be read from the output files
    if force_read_charges == 1:
        taffi_flag = 0

    # Check whether the information is present to perform a TAFFI fit
    # if not, the charges are read directly from the quantum chemistry output file
    else:
        taffi_flag = 1
        for count_i,i in enumerate(geo_folders):
            charge_vpot = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"*.vpot") ]
            charge_xyz = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"geoopt.xyz") ]
            if len(charge_vpot) == 0 or len(charge_xyz) == 0:
                taffi_flag = 0
                break

    # Print diagnostic
    if taffi_flag == 1:
        print("Fitting the partial charges using the TAFFI two-step procedure...")
    else:
        print("Reading the partial charges directly from the quantum chemistry output...")

    # Collect charges for atoms that were parametrized: The algorithm uses
    # CHELPG charges from the globally optimized geometry. Several redundant
    charge_list = []
    atomtypes = []
    taffy_flag = 0
    charges = {}
    eq_charges_AA = {}
    eq_charges_UA = {}
    for count_i,i in enumerate(geo_folders):

        # Find the output file of the globally optimized geometry (i.e., the iteration 
        # corresponding to args.bond_steps is the unperturbed globally optimized geometry)
        # and check for *.vpot and *.xyz files
        output = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"*.out") ]
        if len(output) == 0:
            print("ERROR in get_q_and_m: No output for the globally optimized geometry was found. Check that the run completed.")
            quit()
        charge_vpot = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"*.vpot") ]
        charge_xyz = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"geoopt.xyz") ]

        # Parse the charges using hte TAFFI procedure
        if taffi_flag == 1:            
            # parse the total charge from the output file
            qtot = None
            with open(output[0],'r') as f:
                for lines in f:
                    fields = lines.split()
                    if len(fields) == 3 and fields[0] == "Total" and fields[1] == "charge:":
                        qtot = round(float(fields[2]))

            # safety in case something went wrong during the parse
            if qtot is None:
                print("ERROR in get_q_and_m: could not parse the total charge on the molecule from {}. Exiting...".format(output[0]))
                quit()

            # parse the partial charges using the taffi two-step procedure
            if zero_opt == 1:
                tmp_charges_AA = { j:(0.0,0.0) for j in geo_atomtypes[count_i] }
                tmp_charges_AA.update({ j+'-UA':(0.0,0.0) for j in geo_atomtypes[count_i] })
                tmp_charges_UA = tmp_charges_AA
            else:
                tmp_charges_AA,errs = fit_charges(charge_vpot[0],xyz_file=charge_xyz[0],out_file=output[0],qtot=qtot,gens=gens,w_pot=1.0,w_qtot=1.0,w_hyper=0.0,w_dipole=0.1,FF_db=FF_db,UA_opt=False,symmetrize=False,two_step=True)
                tmp_charges_UA,errs = fit_charges(charge_vpot[0],xyz_file=charge_xyz[0],out_file=output[0],qtot=qtot,gens=gens,w_pot=1.0,w_qtot=1.0,w_hyper=0.0,w_dipole=0.1,FF_db=FF_db,UA_opt=True,symmetrize=False,two_step=True)
            
            # append the charges to the appropriate list in the dictionary
            for j in tmp_charges_AA:
                if j not in list(charges.keys()):
                    charges[j] = []
                charges[j] += [tmp_charges_AA[j][0]]

            # append the equilibrium charges to the eq* dictionaries
            eq_charges_AA[i] = np.array([ tmp_charges_AA[j][0] for j in geo_atomtypes[count_i] ])
            #print(tmp_charges_AA.keys())
            #print(i)
            eq_charges_UA[i] = np.array([ tmp_charges_UA[j+'-UA'][0] if return_UA_H(j) == 0 else 0.0 for j in geo_atomtypes[count_i] ])
            
        # If the *.vpot file is absent then the charges are read directly from the quantum chemistry output file
        else:
            charges_flag = 0
            charges_count = 0
            with open(output[0],'r') as f:
                for count,lines in enumerate(f):
                    fields = lines.split()
                    if len(fields) == 2 and fields[0] == "CHELPG" and fields[1] == "Charges":
                        charges_flag = 1
                        continue

                    # Necessary because of a "------" filler line in the Orca CHELPG output
                    if charges_flag == 1:
                        charges_flag = 2
                        continue

                    # Start parsing charges
                    if charges_flag == 2:
                        charge_list += [float(fields[3])]
                        charges_count += 1

                        # Once all charges have been read reset the count and flag.
                        if charges_count == len(geo_atomtypes[count_i]):
                            charges_count = 0
                            charges_flag = 0

                    # Once the end of the geometry optimization job is reached stop reading.
                    if len(fields) == 5 and fields[1] == 'JOB' and fields[2] == 'NUMBER' and fields[3] == '2':
                        break

            # Calculate UA charges for the equiliubrium configuration
            e,g      = xyz_parse(charge_xyz[0]) 
            adj_mat  = Table_generator(e,g,File=charge_xyz[0])
            UA_inds = [ count_j for count_j,j in enumerate(geo_atomtypes[count_i]) if return_UA_H(j) == 1 ]
            UA_charges = np.array(deepcopy(charge_list))
            for j in UA_inds:                
                UA_charges[next( count_k for count_k,k in enumerate(adj_mat[j]) if k == 1 )] += UA_charges[j]
                UA_charges[j] = 0.0

            # append the equilibrium charges to the eq* dictionaries
            eq_charges_AA[i] = np.array(deepcopy(charge_list))
            eq_charges_UA[i] = deepcopy(UA_charges)

        # Append the atomtypes from the current geometry to the total list
        atomtypes += geo_atomtypes[count_i]

    # If the charges were read from the output file then average over like types and initialize a charge dictionary
    if taffi_flag == 0:

        charges = {}
        for i in sorted(set(atomtypes)):
            ind = []
            charge = 0.0

            # find like types and sum charges
            for count_j,j in enumerate(charge_list):
                if atomtypes[count_j] == i:
                    ind += [count_j]
                    charge += j

            # Assign averaged charges to all like types
            for j in ind:
                charge_list[j] = charge/float(len(ind))        
            charges[i] = charge/float(len(ind))

    # Else, calculate the mean of each charge type
    # NOTE: since the fragments may include atomtypes that aren't in the master atomtypes lists the tmp_charges dictionary
    #       is used to only retain the relevant atomtypes
    else:
        tmp_charges = {}
        for i in set(atomtypes):
            tmp_charges[i] = np.mean(charges[i])
        charges = tmp_charges

    #####################
    #  Process Masses   #
    #####################

    # Assign mass to each atom based on element label
    # If unique_atomtypes is supplied then it is used (this is used to ensure that all atomtypes get masses assigned even when only a subset of modes are being parsed; not optimal but safe)
    masses = {}
    if unique_atomtypes is not None:
        for i in set(unique_atomtypes):
            masses[i] = mass_dict[atomic2ele[int(i.split('[')[1].split(']')[0])]]
    else:
        for i in set(atomtypes):
            masses[i] = mass_dict[atomic2ele[int(i.split('[')[1].split(']')[0])]]

    return charges,masses,{ i:atomic2ele[int(i.split('[')[1].split(']')[0])] for i in list(charges.keys()) },eq_charges_AA,eq_charges_UA

# Description: This function reads in the prerequisite all-atom data and 
#              returns the equivalent united-atom (UA) lists included masses,
#              charges, types, elements, and indices
def gen_UA_data(master_files,masses_AA,charges_AA,elements_AA,unique_atomtypes=None):

    # Iterate over input geometry files
    all_types = []
    all_elements = [] 
    all_masses = []
    all_charges = []
    for count_i,i in enumerate(master_files):

        # get the elements, atomtypes, and adjacency matrix of the current geometry
        elements_i,geo_i = xyz_parse(i)        
        adj_mat_i = Table_generator(elements_i,geo_i,File=i)
        atomtypes_i = get_scan_atomtypes(i)   

        # Find atoms connected to hydrogens that were not parametrized (if UA is requested).
        # If a hydrogen is discovered in the atom_indices list, then it is not united to its
        # parent heavy atom. 
        H_ind = [ count_i for count_i,i in enumerate(elements_i) if i == "H" ]

        # Remove hydrogens that aren't attached to sp3 carbon
        del_list = []
        for count_j,j in enumerate(H_ind):
            if len([ a for count_a,a in enumerate(adj_mat_i[j]) if a == 1 and int(atomtypes_i[count_a].split('[')[1].split(']')[0]) == 6 and int(sum(adj_mat_i[count_a])) == 4 ]) == 0:   # includes sp3 provision
                del_list += [count_j]
        H_ind = [ j for count_j,j in enumerate(H_ind) if count_j not in del_list ]            

        # Intitialize a list of UA atom indices
        ind_UA = []

        # Initialize temporary copies of mass and charge lists
        # for calculating the UA charges and masses
        tmp_masses = np.array([ masses_AA[j] for j in atomtypes_i ])
        tmp_charges = np.array([ charges_AA[j] for j in atomtypes_i ])

        # Interate over the rows of the adj_mat
        for count_a,a in enumerate(adj_mat_i):

            # Skip hydrogens and atoms not being fit
            if ( count_a in H_ind ):
                continue
            else:
                ind_UA += [count_a]

            # Find indices of connected hydrogens
            ind = [ count_j for count_j,j in enumerate(a) if ( (j == 1) and (count_j in H_ind) ) ]

            # Add the hydrogen charges and mass to the current atom's charge and mass
            for j in ind:
                tmp_masses[count_a] += tmp_masses[j]
                tmp_charges[count_a] += tmp_charges[j]

        # Create lists of UA types and UA elements, masses and charges with removed hydrogen entries.     
        all_types    += [ atomtypes_i[j]+'-UA' for j in ind_UA ]
        all_elements += [ elements_i[j] for j in ind_UA ]
        all_masses   += [ tmp_masses[j] for j in ind_UA ]
        all_charges  += [ tmp_charges[j] for j in ind_UA ]
    
    # Average over duplicate UA types. 
    atomtypes_UA = sorted(set(all_types))
    masses_UA = {}
    charges_UA = { i:0.0 for i in atomtypes_UA }
    for count_i,i in enumerate(atomtypes_UA):
        charge_counter = 0
        mass_counter = 0
        for count_j,j in enumerate(all_types):

            # Assign mass and element
            if j == i and mass_counter == 0:
                masses_UA[i] = all_masses[count_j]
                mass_counter = 1

            # Add charges to massive nucleus
            if j == i:
                charges_UA[i] += all_charges[count_j]
                charge_counter += 1

        # Average charges
        charges_UA[i] = charges_UA[i]/float(charge_counter)
        
    # If UA charges were included in charges_AA (this happends if the TAFFI fit of partial charges was possible, and is the more accurate scenario),
    # then the UA_charges are populated directly from the charges_AA dict. 
    # XXX This function should eventually be overhauled to assume that the charges_AA dictionary holds the charges_UA data since the alternative is less accurate and inconsistent with TAFFI.
    if False not in [ i in charges_AA for i in list(charges_UA.keys()) ]:
        print("WARNING: an obsolete loop ran. TROUBLESHOOT.")
        quit()
        for i in list(charges_UA.keys()):
            charges_UA[i] = charges_AA[i]

    return charges_UA,masses_UA,{ i:elements_AA[i.split("-UA")[0]] for i in list(charges_UA.keys()) }

# This is a messy wrapper for the bond fitting procedure. In brief, the 
# function looks up all the bond data, calculates harmonic fits, and
# generates some plots and data text files. All parameters are saved to FF_dict
#def parse_bonds(base_name,elements,atomtypes,bond_steps,modes_from_FF,save_folder,qc_types=[]):
def parse_bonds(base_name,bond_types,bond_folders,bond_atomtypes,modes_from_FF,save_folder,qc_types=[]):

    # In python all global variables must be declared in each scope
    global FF_dict,harm_range

    # Save parent directory
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_bond_plots_DFT=[]
    paths_to_bond_plots_MP2=[]

    #####################################
    ###### Process bond directories #####
    #####################################
    for bonds_count,bonds in enumerate(bond_types):

        print("\n{}".format("*"*167))
        print("* {:^163s} *".format("Processing bond {}".format(" ".join(bonds))))
        print("{}".format("*"*167))

        # Change into current bond directory
        os.chdir(bond_folders[bonds_count])

        # Parse bond data
        Geo_dict,bond_atoms,completion_flag = get_bond_angle_data(working_dir+'/'+save_folder+'/'+base_name+'_'+"_".join(bonds)+'_geos.xyz')

        # Avoid the parse is the data is incomplete
        if completion_flag == 0:
            os.chdir(working_dir)
            print("******************************************************************************************")
            print("  Warning: Scan jobs did not run to completion for bond {}...Skipping...".format("_".join(bonds)))
            print("******************************************************************************************")
            continue

#         # Find minimum energies
#         Min_DFT_E = 1000
#         Min_MP2_E = 1000
#         for keys in Geo_dict.keys():
#             if 'dft' in qc_types and Geo_dict[keys]["DFT"] < Min_DFT_E:
#                 Min_DFT_E = Geo_dict[keys]["DFT"]
#             if 'mp2' in qc_types and Geo_dict[keys]["MP2"] < Min_MP2_E:
#                 Min_MP2_E = Geo_dict[keys]["MP2"]

        # Calculate Bond Lengths
        for keys in list(Geo_dict.keys()):
            Geo_dict[keys]["length"] = norm(Geo_dict[keys]["geo"][bond_atoms[1]]-Geo_dict[keys]["geo"][bond_atoms[0]])

        ######################################
        ###### Calculate Fit Parameters ######
        ######################################

        # Collect list of energies and bond lengths
        lengths = []
        DFT_E = []
        MP2_E = []
        for i in range(len(list(Geo_dict.keys()))):
            lengths += [Geo_dict[str(i)]["length"]]
            if 'dft' in qc_types:
#                DFT_E += [Geo_dict[str(i)]["DFT"] - Min_DFT_E]
                DFT_E += [Geo_dict[str(i)]["DFT"]]
            if 'mp2' in qc_types:
#                MP2_E += [Geo_dict[str(i)]["MP2"] - Min_MP2_E]
                MP2_E += [Geo_dict[str(i)]["MP2"]]
        
        # Sort by lengths
        if 'dft' in qc_types and 'mp2' in qc_types:
            together = sorted(zip(lengths,DFT_E,MP2_E))
            lengths  = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            MP2_E    = [ i[2] for i in together ]            
        elif 'dft' in qc_types:
            together = sorted(zip(lengths,DFT_E))
            lengths  = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
        elif 'mp2' in qc_types:
            together = sorted(zip(lengths,MP2_E))
            lengths  = [ i[0] for i in together ]
            MP2_E    = [ i[2] for i in together ]

        # If the minimum is near the edge of the scanned range then automatically perform a reoptimization
        min_ind = DFT_E.index(min(DFT_E)) 
        resub_flag = 0
        if min_ind in [0,1,len(DFT_E)-1,len(DFT_E)-2]:
            print("Minimum lies near the edge of the scanned range...")
            resub_flag = 1

        # Remove edge effects (i.e. edges with the improper derivative)
        # The value of correction_flag is used to toggle a diagnostic print statement
        # NOTE: the delta variable is calculated here in case a resubmission is required
        delta = np.abs(lengths[1]-lengths[0])
        start_ind = 0
        end_ind = len(lengths)-1
        correction_flag = 0
        if resub_flag == 0:
            for count_i,i in enumerate(lengths[:-1]):
                if DFT_E[count_i+1] - DFT_E[count_i] < 0: 
                    start_ind = count_i
                    break
            if start_ind != 0: correction_flag = 1

            for i in range(len(lengths))[::-1]:
                if DFT_E[i]-DFT_E[i-1] > 0: 
                    end_ind = i
                    break
            if end_ind != len(lengths)-1: correction_flag = 1

        # Assembled edited lists and print diagnostic if trimming occured
        if correction_flag == 1:
            print("Data points with improper derivatives were removed to increase the quality of fit...")
        if 'dft' in qc_types: DFT_E = DFT_E[start_ind:end_ind+1]
        if 'mp2' in qc_types: MP2_E = MP2_E[start_ind:end_ind+1]
        lengths = lengths[start_ind:end_ind+1]

        # Remove derivative discontinuities
        # NOTE: the threshold of 1000.0 kcal^2/ang^2 is empirical, it was found to consistently remove most derivative 
        #       discontinuities without removing any continuous data.
        second_d_thresh = 1000.0
        print("DFT_E: {}".format(DFT_E))
        print("lengths: {}".format(lengths))

        # Remove datapoints until derivative discontinuities have been removed
        correction_flag = 0
        clean_flag = 0
        while clean_flag == 0 and len(DFT_E) > 2:

            # Loop over the energies and calculate the second derivatives. 
            # NOTE: the if/else construction is necessary to deal with the endpoints
            clean_flag = 1
            del_list = []            
            for i in range(len(DFT_E))[1:-1]:                

                second_d = ( ( DFT_E[i] - DFT_E[i-1] ) / ( lengths[i] - lengths[i-1] ) - ( DFT_E[i+1] - DFT_E[i] ) / ( lengths[i+1] - lengths[i] ) ) / (lengths[i] - lengths[i-1] ) 
                if second_d > 0 or np.abs(second_d) > second_d_thresh:
                    #print "{} (removing an element; angle: {})".format(second_d,angles[i])
                    correction_flag = 1
                    clean_flag = 0
                    del_list += [i]
                    break
                
            DFT_E = [ i for count_i,i in enumerate(DFT_E) if count_i not in del_list ]
            lengths = [ i for count_i,i in enumerate(lengths) if count_i not in del_list ]

        if correction_flag == 1:
            print("Data points with positive second derivatives or discontinuities were removed to increase the quality of the fit...")
        
        # If more than a third of the data has been discarded then the script generates a reoptimization job using the lowest
        # energy conformation from the scan.
        if len(lengths) <= len(list(Geo_dict.keys()))*4.0/5.0 or resub_flag == 1:
            min_folder = [ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ].index(min([ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ]))
            template_file = str(min_folder) + '/' + bond_folders[bonds_count].split('/')[0] + '_' + '-'.join(bonds) + '_' + str(min_folder) + '.in'
            print("Generating the input files for a reoptimization based on the lowest energy configuration...")
            gen_reoptimization(template=template_file,geo=Geo_dict[str(i)]['geo'],elements=Geo_dict[str(i)]["elements"],delta=delta,N_steps=len(list(Geo_dict.keys())))
            os.chdir(working_dir)
            continue

        # Normalize energies by the minimum value
        if 'dft' in qc_types: DFT_E = np.array(DFT_E)-min(DFT_E)
        if 'mp2' in qc_types: MP2_E = np.array(MP2_E)-min(MP2_E)

        # print diagnostic
        print("constrained atoms: {}".format(" ".join([ str(i) for i in bond_atoms])))
        print("length range: {: <6.4f}-{: <6.4f}".format(min(np.array(lengths)),max(np.array(lengths))))
        print("Calculating harmonic fits...")
        harm_range = [min(lengths),max(lengths)]

        # Use read parameters if available (params holds force constant and theta_0)
        if bonds in modes_from_FF:
            print("Using parameters for this bond from the supplied FF-file(s)...")
            if "dft" in qc_types:
                DFT_params = np.array(FF_dict["bonds"][bonds]["DFT"][1:])
            if "mp2" in qc_types:
                MP2_params = np.array(FF_dict["bonds"][bonds]["MP2"][1:])

        # Calculate parameters (params holds force constant and theta_0)
        else:
            if 'dft' in qc_types:
                DFT_params, DFT_errors = curve_fit(harmonic, lengths, DFT_E,p0=[100.0,Geo_dict[str(int(round(len(lengths)/2.0)))]["length"]],maxfev=10000)
            if 'mp2' in qc_types:
                MP2_params, MP2_errors = curve_fit(harmonic, lengths, MP2_E,p0=[100.0,Geo_dict[str(int(round(len(lengths)/2.0)))]["length"]],maxfev=10000)

        #####################################
        ############# Write Info ############
        #####################################

        # Return to parent/"bonds" folder
        os.chdir(working_dir)

        # Open file and write header
        with open(save_folder+'/'+"_".join(bonds)+'_info.txt', 'w') as f:
            if 'dft' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol Ang^-2\n{:20s} {:< 12.6f} Ang\n".format("k (DFT) =",DFT_params[0],"r_0 (DFT) =",DFT_params[1]))
            if 'mp2' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol Ang^-2\n{:20s} {:< 12.6f} Ang\n".format("k (MP2) =",MP2_params[0],"r_0 (MP2) =",MP2_params[1]))
            if 'dft' in qc_types and 'mp2' in qc_types:
                f.write("{:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s}\n".format("r (Angstroms)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)","MP2_Energy","E-E0 (MP2;kcal/mol)","E-E0 (MP2_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["length"],Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["length"],DFT_params[0],DFT_params[1]),\
                            Geo_dict[str(i)]["MP2"],Geo_dict[str(i)]["MP2"],harmonic(Geo_dict[str(i)]["length"],MP2_params[0],MP2_params[1])))
            elif 'dft' in qc_types: 
                f.write("{:<25s} {:<25s} {:<25s} {:<25s}\n".format("r (Angstroms)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["length"],Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["length"],DFT_params[0],DFT_params[1])))

        #### Generate 2D representation of the molecule KNOWN ISSUE WITH PLOTTING FRAGMENT BASED BONDS AND ANGLES, 2D REPRESENTATION IS TAKEN FROM THE MASTER FRAGMENT
        print("Generating 2D representation of the fit fragment...")
        adj_mat = Table_generator(Geo_dict["0"]["elements"],Geo_dict["0"]["geo"])
        geo_2D = kekule(Geo_dict["0"]["elements"],bond_atomtypes[bonds_count],Geo_dict["0"]["geo"],adj_mat)

        #######################################
        ######### Generate Fit Plots ##########
        #######################################
        print("Rendering fit plots...")

        if 'dft' in qc_types:

            fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
            draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
            if draw_scale > 1.0: draw_scale = 1.0

            # Add bonds
            linewidth = 2.0*draw_scale*5.0
            if linewidth > 2.0:
                linewidth = 2.0
            if linewidth < 1.0:
                linewidth = 1.0
            for count_i,i in enumerate(adj_mat):
                for count_j,j in enumerate(i):
                    if count_j > count_i:
                        if j == 1:
                            ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

            # Add Highlight
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[bond_atoms[0]][0],geo_2D[bond_atoms[1]][0]],[geo_2D[bond_atoms[0]][1],geo_2D[bond_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                            
            # Add atom labels
            fontsize = 20*draw_scale*5.0
            if fontsize > 24:
                fontsize = 24
            if fontsize < 9:
                fontsize = 9            
            for count_i,i in enumerate(geo_2D):
                if count_i in bond_atoms:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                else:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
            ax1.axis('image')
            ax1.axis('off')

            # Plot data
            step = (max(lengths)-min(lengths))/100.0
            xdata = np.arange(min(lengths),max(lengths)+step,step)
            ydata = harmonic(xdata,DFT_params[0],DFT_params[1])
            ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
            ax2.scatter(np.array(lengths),DFT_E,s=120,color=(0.0,0.1,0.8))
            ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
            ax2.set_xlabel(r'$\mathrm{Length (\AA)}$',fontsize=20,labelpad=10)
            ax2.set_xlim([min(lengths), max(lengths)])
            ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
            ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
            y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            ax2.yaxis.set_major_formatter(y_formatter)
            ax2.xaxis.set_major_formatter(x_formatter)
            [j.set_linewidth(3) for j in ax2.spines.values()]
            plt.tight_layout()
            Name = save_folder+'/'+"_".join(bonds)+'-DFT_fit.png'
            plt.savefig(Name, bbox_inches=0,dpi=300)
            plt.savefig(save_folder+'/'+"_".join(bonds)+'-DFT_fit.pdf', bbox_inches=0,dpi=300)
            paths_to_bond_plots_DFT+=[Name]
            plt.close(fig)

        if 'mp2' in qc_types:

            fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
            draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
            draw_scale = 1.0
            # Add bonds
            for count_i,i in enumerate(adj_mat):
                for count_j,j in enumerate(i):
                    if count_j > count_i:
                        if j == 1:
                            ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=2.0*draw_scale))
            # Add Highlight
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[bond_atoms[0]][0],geo_2D[bond_atoms[1]][0]],[geo_2D[bond_atoms[0]][1],geo_2D[bond_atoms[1]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))
                            
            # Add atom labels
            for count_i,i in enumerate(geo_2D):
                if count_i in bond_atoms:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                else:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
            ax1.axis('image')
            ax1.axis('off')

            # Plot Data
            step = (max(lengths)-min(lengths))/100.0
            xdata = np.arange(min(lengths),max(lengths)+step,step)
            ydata = harmonic(xdata,MP2_params[0],MP2_params[1])
            ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
            ax2.scatter(np.array(lengths),MP2_E,s=120,color=(0.0,0.1,0.8))
            ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
            ax2.set_xlabel(r'$\mathrm{Length \, (\AA)}$',fontsize=20,labelpad=10)
            ax2.set_xlim([min(lengths), max(lengths)])
            ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
            ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
            y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            ax2.yaxis.set_major_formatter(y_formatter)
            ax2.xaxis.set_major_formatter(x_formatter)
            [j.set_linewidth(3) for j in ax2.spines.values()]
            plt.tight_layout()

            Name = save_folder+'/'+"_".join(bonds)+'-MP2_fit.png'
            plt.savefig(Name, bbox_inches=0,dpi=300)
            plt.savefig(save_folder+'/'+"_".join(bonds)+'-MP2_fit.pdf', bbox_inches=0,dpi=300)
            paths_to_bond_plots_MP2+=[Name]
            plt.close(fig)

        ##############################################################
        ######## Update FF_dict with DFT and MP2 FF-parameters #######
        ##############################################################
        FF_dict["bonds"][bonds]={}
        if 'dft' in qc_types:
            FF_dict["bonds"][bonds]["DFT"] = ["harmonic"] + DFT_params.tolist()
        if 'mp2' in qc_types:
            FF_dict["bonds"][bonds]["MP2"] = ["harmonic"] + MP2_params.tolist()

    # Return to parent folder after processing all bonds
    os.chdir(working_dir)
    
    return paths_to_bond_plots_DFT, paths_to_bond_plots_MP2

def get_bond_angle_data(xyz_name):

    # Check for a scan run
    dirs = [ dn[0] for dp, dn, filenames in os.walk('.') if 'scan' in dn ]

    # Parse bond data for continuous scan
    if 'scan' in dirs:

        # Find the output file for this scan
        output_file = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ][0]

        # Find the mode being scanned and the number of atoms in the geometry
        con_flag = 0
        geo_flag = 0
        N_atoms  = 0
        N_jobs   = 1
        with open(output_file,'r') as f:
            for lines in f:
                if "%geom scan" in lines.lower(): con_flag = 1;continue
                if con_flag == 1: bond_atoms = [ int(i) for i in lines.split("=")[0].split(">")[1].split()[1:] ]; con_flag = 0
                if "CARTESIAN COORDINATES (ANGSTROEM)" in lines: geo_flag = 1; continue
                if geo_flag == 1:
                    fields = lines.split()
                    if len(fields) == 1: continue
                    if len(fields) == 0: break
                    N_atoms += 1
                if "JOBS TO BE PROCESSED THIS RUN" in lines: N_jobs = int(lines.split()[3])

        # If there are multiple jobs, then that means a preliminary relaxation and/or mode scan is performed prior to the run data
        # so the first job is skipped
        if N_jobs == 1:
            parse_job = 1
        else: 
            parse_job = 0

        # Find the optimized geometries and single point energies
        parse_flag      = -1
        job_count       = 0
        geo_flag        = 0
        completion_flag = 0
        atom_count      = 0
        Geo_dict        = {}        
        with open(output_file,'r') as f:
            for lines in f:

                # Only the last job is parsed, the parse_job flag and job_count flags control appropriate parsing
                if parse_job == 1:
                    if "RELAXED SURFACE SCAN STEP" in lines: 
                        parse_flag += 1
                        Geo_dict[str(parse_flag)] = {}
                        Geo_dict[str(parse_flag)]["geo"] = np.zeros([N_atoms,3])
                        Geo_dict[str(parse_flag)]["elements"] = ["X"]*N_atoms
                    if "CARTESIAN COORDINATES (ANGSTROEM)" in lines: geo_flag = 1; continue
                    if geo_flag == 1:
                        fields = lines.split()
                        if len(fields) == 1: continue
                        Geo_dict[str(parse_flag)]["geo"][atom_count,:] = np.array([float(fields[1]),float(fields[2]),float(fields[3])])
                        Geo_dict[str(parse_flag)]["elements"][atom_count] = fields[0]
                        atom_count += 1
                        if atom_count == N_atoms: 
                            atom_count = 0 
                            geo_flag = 0

                    if "FINAL SINGLE POINT ENERGY" in lines: Geo_dict[str(parse_flag)]["DFT"] = float(lines.split()[4])*Htokcalmol
                    if "MP2 TOTAL ENERGY" in lines: Geo_dict[str(parse_flag)]["MP2"] = float(lines.split()[3])*Htokcalmol

                    if "****ORCA TERMINATED NORMALLY****" in lines: completion_flag = 1

                # Only parse the last job
                elif "JOB NUMBER" in lines: 
                    job_count += 1
                    if job_count == 2:
                        parse_job = 1

        # Proceed to catenation of the optimized geometries if all are completed. While the catenated xyz file is
        # produced, all geometries are also saved to a dictionary called Geo_dict, keyed to their number in the sequence.
        # Geo_dict holds geometries, active angle in degrees, the geometry's DFT energy and the geometry's MP2 energy in the
        # the format: Geo_dict["0"]["geo"] returns geometry as an Nx3 array; Geo_dict["0"]["angle"] returns the angle 
        # in degrees of the active angle in the current geometry ("0"); Geo_dict["0"]["DFT"]/Geo_dict["0"]["MP2"] return the
        # DFT and MP2 energies of the current geometry, respectively. 

        # Open file and initialize dictionary to hold catenated geometries 
        with open(xyz_name, 'w') as f:

            # Loop over all discovered geometries and save their geometries to Geo_dict and
            # catenate to a single file for viewing.
            for i in natural_sort(list(Geo_dict.keys())):
                f.write("{}\n\n".format(len(Geo_dict[i]["elements"])))
                for count_j,j in enumerate(Geo_dict[i]["elements"]):
                    f.write("{:<20s} {:< 12.8f} {:< 12.8f} {:< 12.8f}\n".format(j,Geo_dict[i]["geo"][count_j][0],Geo_dict[i]["geo"][count_j][1],Geo_dict[i]["geo"][count_j][2]))                

    # Parse bond data for parallel scan
    else:
        
        # Find the optimized geometry files, and sort by number.
        geometries = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*geo_opt.xyz") ]
        sort_nicely(geometries)

        # Proceed to catenation of the optimized geometries if all are completed. While the catenated xyz file is
        # produced, all geometries are also saved to a dictionary called Geo_dict, keyed to their number in the sequence.
        # Geo_dict holds geometries, active angle in degrees, the geometry's DFT energy and the geometry's MP2 energy in the
        # the format: Geo_dict["0"]["geo"] returns geometry as an Nx3 array; Geo_dict["0"]["angle"] returns the angle 
        # in degrees of the active angle in the current geometry ("0"); Geo_dict["0"]["DFT"]/Geo_dict["0"]["MP2"] return the
        # DFT and MP2 energies of the current geometry, respectively. 

        # First find the number of atoms 
        with open(geometries[0],'r') as xyz:
            for _ in xyz:
                N_atoms = int(_.split()[0])
                break

        # Open file and initialize dictionary to hold catenated geometries 
        f = open(xyz_name,'w')
        Geo_dict = {}

        # Loop over all discovered geometries and save their geometries to Geo_dict and
        # catenate to a single file for viewing.
        #print(geometries)
        # geometris = ['./0/0_geo_opt.xyz', './1/1_geo_opt.xyz', './2/2_geo_opt.xyz', './3/3_geo_opt.xyz', './4/4_geo_opt.xyz', './5/5_geo_opt.xyz', './6/6_geo_opt.xyz', './7/7_geo_opt.xyz', './8/8_geo_opt.xyz', './9/9_geo_opt.xyz', './10/10_geo_opt.xyz']         


        #### DEBUG: geometries files is the same
        #### geometries files are the output of geoopt ORCA
        #### The ORCA input xyz files have the same atom types 
        
        for count_g,g in enumerate(geometries):
            Geo_dict[str(count_g)] = {}
            Geo_dict[str(count_g)]["geo"] = np.zeros([N_atoms,3])
            Geo_dict[str(count_g)]["elements"] = ["X"]*N_atoms
            atom_count = 0
            with open(g,'r') as xyz:

                for count,lines in enumerate(xyz):
                    f.write(lines)                

                    # Write to dictionary
                    if count >= 2:
                        fields = lines.split()
                        Geo_dict[str(count_g)]["geo"][atom_count,:] = np.array([float(fields[1]),float(fields[2]),float(fields[3])])
                        Geo_dict[str(count_g)]["elements"][atom_count] = fields[0]
                        atom_count += 1

        f.write("\n")
        f.close()

        ######################################
        ########## Collect Energies ##########
        ######################################

        # Find the completed output files, and sort by number.
        outputs = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ]
        sort_nicely(outputs)

        # Check for the completion of all output files
        flag = [0]*len(outputs)
        for count_o,o in enumerate(outputs):
            with open(o,'r') as f:
                for lines in f: 
                    if "****ORCA TERMINATED NORMALLY****" in lines: flag[count_o] = 1
        if 0 in flag: completion_flag = 0
        else: completion_flag = 1

        # First find the atoms that are constrained
        with open(outputs[0],'r') as out:
            flag = 0
            for _ in out:
                fields = _.split()
                if (len(fields) == 4 and fields[2].lower() == "%geom" and fields[3].lower() == "constraints") or (len(fields) == 3 and fields[0] == "|" and fields[2].lower()=="constraints"):
                    flag = 1
                    continue
                if flag == 1:
                    if _.split(">")[1].split()[1] == "B":
                        bond_atoms = [ int(i) for i in _.split(">")[1].split()[2:4] ]
                    if _.split(">")[1].split()[1] == "A":
                        bond_atoms = [ int(i) for i in _.split(">")[1].split()[2:5] ]
                    if _.split(">")[1].split()[1] == "D":
                        bond_atoms = [ int(i) for i in _.split(">")[1].split()[2:6] ]
                    break


        # Find the DFT and MP2 final energies for each output. Note: energies
        # are converted to kcal/mol by default and assume hartree output
        #print(outputs)
        #output = './0/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_0.out', './1/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_1.out', './2/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_2.out', './3/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_3.out', './4/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_4.out', './5/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_5.out', './6/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_6.out', './7/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_7.out', './8/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_8.out', './9/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_9.out', './10/bonds_angles_[6[6[1][1][1]][8[6]][1][1]]-[8[6[6][1][1]][6[1][1][1]]]_10.out']
        for count_o,o in enumerate(outputs):
            with open(o,'r') as out:
                for lines in out:
                    fields = lines.split()

                    if len(fields) == 5 and fields[0].lower() == "final" and fields[1].lower() == "single" and fields[2].lower() == "point" and fields[3].lower() == "energy":
                        Geo_dict[str(count_o)]["DFT"] = float(fields[4])*Htokcalmol
                    if len(fields) == 5 and fields[0].lower() == "mp2" and fields[1].lower() == "total" and fields[2].lower() == "energy:":
                        Geo_dict[str(count_o)]["MP2"] = float(fields[3])*Htokcalmol
                        break
         

    return Geo_dict,bond_atoms,completion_flag

# This is a messy wrapper for the angle fitting procedure. In brief, the 
# function looks up all the angle data, calculates harmonic fits, and
# generates some plots and data text files. All parameters are saved to FF_dict
def parse_angles(base_name,angle_types,angle_folders,angle_atomtypes,modes_from_FF,save_folder,qc_types=[]):

    # In python all global variables must be declared in each scope
    global FF_dict,harm_range

    # Save parent directory
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_angle_plots_DFT=[]
    paths_to_angle_plots_MP2=[]

    ######################################
    ###### Process angle directories #####
    ######################################
    for angles_count,angles in enumerate(angle_types):

        print("\n{}".format("*"*167))
        print("* {:^163s} *".format("Processing angle {}".format(" ".join(angles))))
        print("{}".format("*"*167))

        # Change into current bond directory
        os.chdir(angle_folders[angles_count])

        # Parse angle data
        Geo_dict,angle_atoms,completion_flag = get_bond_angle_data(working_dir+'/'+save_folder+'/'+base_name+'_'+"_".join(angles)+'_geos.xyz')

        # Avoid the parse is the data is incomplete
        if completion_flag == 0:
            os.chdir(working_dir)
            print("******************************************************************************************")
            print("  Warning: Scan jobs did not run to completion for angle {}...Skipping...".format("_".join(angles)))
            print("******************************************************************************************")
            continue

        # Calculate Angles
        for keys in natural_sort(list(Geo_dict.keys())):
            atom_1 = Geo_dict[keys]["geo"][angle_atoms[0]]
            atom_2 = Geo_dict[keys]["geo"][angle_atoms[1]]
            atom_3 = Geo_dict[keys]["geo"][angle_atoms[2]]
            Geo_dict[keys]["angle"] = np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))

        ######################################
        ###### Calculate Fit Parameters ######
        ######################################

        # Collect list of energies and angles
        thetas = []
        DFT_E = []
        MP2_E = []
        for i in range(len(list(Geo_dict.keys()))):
            thetas += [Geo_dict[str(i)]["angle"]]
            if 'dft' in qc_types:
                DFT_E += [Geo_dict[str(i)]["DFT"]]
            if 'mp2' in qc_types:
                MP2_E += [Geo_dict[str(i)]["MP2"]]

        # Check for linear angle
        # (i) a change in gradient and (ii) vicinity to 180.0 are used to check for linearity                
        linearity_flag = 0
        if thetas[1] - thetas[0] >= 0.0:            
            if False in [ (i-thetas[count_i]) > 0.0 for count_i,i in enumerate(thetas[1:]) ] and (180.0 - np.mean(thetas)*180.0/np.pi) < 5.0:
                linearity_flag = 1
                turn = [ (i-thetas[count_i]) > 0.0 for count_i,i in enumerate(thetas[1:]) ].index(False)+1
                thetas[turn:] = [ i+2.0*(np.pi-i) for i in thetas[turn:] ]
        elif thetas[1] - thetas[0] < 0:
            if False in [ (i-thetas[count_i]) < 0.0 for count_i,i in enumerate(thetas[1:]) ] and (180.0 - np.mean(thetas)*180.0/np.pi) < 5.0:
                linearity_flag = 1
                turn = [ (i-thetas[count_i]) < 0.0 for count_i,i in enumerate(thetas[1:]) ].index(False)+1
                thetas[turn:] = [ i+2.0*(np.pi-i) for i in thetas[turn:] ]
        if linearity_flag == 1: print("Angles were unwrapped due to the linearity of the angle.")

        # Sort by angle
        if 'dft' in qc_types and 'mp2' in qc_types:
            together = sorted(zip(thetas,DFT_E,MP2_E))
            thetas = [ i[0] for i in together ]
            DFT_E  = [ i[1] for i in together ]
            MP2_E  = [ i[2] for i in together ]            
        elif 'dft' in qc_types:
            together = sorted(zip(thetas,DFT_E))
            thetas = [ i[0] for i in together ]
            DFT_E  = [ i[1] for i in together ]
        elif 'mp2' in qc_types:
            together = sorted(zip(thetas,MP2_E))
            thetas = [ i[0] for i in together ]
            MP2_E  = [ i[2] for i in together ]

        # If the minimum is near the edge of the scanned range then automatically perform a reoptimization
        min_ind = DFT_E.index(min(DFT_E)) 
        resub_flag = 0
        if min_ind in [0,1,len(DFT_E)-1,len(DFT_E)-2]:
            print("Minimum lies near the edge of the scanned range...")
            resub_flag = 1
        
        # Remove edge effects (i.e. edges with the improper derivative)
        # The value of correction_flag is used to toggle a diagnostic print statement
        delta = np.abs(thetas[1]-thetas[0])*180.0/np.pi
        start_ind = 0
        end_ind = len(thetas) - 1
        correction_flag = 0
        if resub_flag == 0:
            for count_i,i in enumerate(thetas[:-1]):
                if DFT_E[count_i+1] - DFT_E[count_i] < 0: 
                    start_ind = count_i
                    break
            if start_ind != 0: correction_flag = 1

            for i in range(len(thetas))[::-1]:
                if DFT_E[i]-DFT_E[i-1] > 0: 
                    end_ind = i
                    break
            if end_ind != len(thetas)-1: correction_flag = 1
        
        # Assembled edited lists and print diagnostic if trimming occured
        if correction_flag == 1:
            print("Data points with improper derivatives were removed to increase the quality of fit...")
        if 'dft' in qc_types: DFT_E = DFT_E[start_ind:end_ind+1]
        if 'mp2' in qc_types: MP2_E = MP2_E[start_ind:end_ind+1]
        thetas = thetas[start_ind:end_ind+1]

        # # Remove discontinuities in the curve based on simple threshold ( 0.1 kcal/mol per 0.1 degrees )
        # start_ind = 0
        # end_ind = len(thetas)-1
        # mid = float(len(thetas))/2.0
        # correction_flag = 0
        # if resub_flag == 0:
        #     for i in range(1,len(thetas)):
        #         if np.abs((DFT_E[i] - DFT_E[i-1])/(thetas[i] - thetas[i-1])) > 0.01/(0.1*np.pi/180.0):
        #             if i-1 <= mid:
        #                 start_ind = i
        #             else:
        #                 end_ind = i-1
        #     if end_ind != len(thetas)-1: correction_flag = 1

        # # Assembled edited lists and print diagnostic if trimming occured
        # if correction_flag == 1:
        #     print "Data points with derivative discontinuities were removed to increase the quality of fit..."
        # if 'dft' in qc_types: DFT_E = DFT_E[start_ind:end_ind+1]
        # if 'mp2' in qc_types: MP2_E = MP2_E[start_ind:end_ind+1]
        # thetas = thetas[start_ind:end_ind+1]            

        # Remove derivative discontinuities
        # NOTE: the threshold of 500.0 kcal^2/rad is empirical, it was found to consistently remove most derivative 
        #       discontinuities without removing any continuous data.
        second_d_thresh = 500.0

        # Remove datapoints until derivative discontinuities have been removed
        correction_flag = 0
        clean_flag = 0
        while clean_flag == 0 and len(DFT_E) > 2:

            # Loop over the energies and calculate the second derivatives. 
            # NOTE: the if/else construction is necessary to deal with the endpoints
            clean_flag = 1
            del_list = []            
            for i in range(len(DFT_E))[1:-1]:                

                second_d = ( ( DFT_E[i] - DFT_E[i-1] ) / ( thetas[i] - thetas[i-1] ) - ( DFT_E[i+1] - DFT_E[i] ) / ( thetas[i+1] - thetas[i] ) ) / (thetas[i] - thetas[i-1] ) 
                if second_d > 0 or np.abs(second_d) > second_d_thresh:
                    #print "{} (removing an element; angle: {})".format(second_d,angles[i])
                    correction_flag = 1
                    clean_flag = 0
                    del_list += [i]
                    break
                
            DFT_E = [ i for count_i,i in enumerate(DFT_E) if count_i not in del_list ]
            thetas = [ i for count_i,i in enumerate(thetas) if count_i not in del_list ]

        if correction_flag == 1:
            print("Data points with positive second derivatives or discontinuities were removed to increase the quality of the fit...")

        # If more than a third of the data has been discarded then the script generates a reoptimization job using the lowest
        # energy conformation from the scan.
        if len(thetas) <= len(list(Geo_dict.keys()))*4.0/5.0 or resub_flag == 1:
            min_folder = [ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ].index(min([ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ]))
            template_file = str(min_folder) + '/' + angle_folders[angles_count].split('/')[0] + '_' + '-'.join(angles) + '_' + str(min_folder) + '.in'
            print("Generating the input files for a reoptimization based on the lowest energy configuration...")
            gen_reoptimization(template=template_file,geo=Geo_dict[str(i)]['geo'],elements=Geo_dict[str(i)]["elements"],delta=delta,N_steps=len(list(Geo_dict.keys())))
            os.chdir(working_dir)
            continue

        # Normalize energies by the minimum value
        if 'dft' in qc_types: DFT_E = np.array(DFT_E)-min(DFT_E)
        if 'mp2' in qc_types: MP2_E = np.array(MP2_E)-min(MP2_E)

        # Print diagnostic
        print("constrained atoms: {}".format(" ".join([ str(i) for i in angle_atoms])))
        print("theta range: {: <6.4f}-{: <6.4f}".format(min(np.array(thetas))*180.0/np.pi,max(np.array(thetas))*180.0/np.pi))
        print("Calculating harmonic fits...")
        harm_range=[min(thetas),max(thetas)]

        # Use read parameters if available (params holds force constant and theta_0)
        if angles in modes_from_FF:
            print("Using parameters for this angle from the supplied FF-file(s)...")
            if "dft" in qc_types:
                DFT_params = np.array(FF_dict["angles"][angles]["DFT"][1:])
                DFT_params[1] = DFT_params[1]*np.pi/180.0
            if "mp2" in qc_types:
                MP2_params = np.array(FF_dict["angles"][angles]["MP2"][1:])
                MP2_params[1] = MP2_params[1]*np.pi/180.0

        # Calculate parameters (params holds force constant and theta_0)    
        # NOTE: the force constant is in units of kcal/mol/rad^2 but the equilibrium angle is in units of degrees (lammps harmonic angle convention). The eq_angle is 
        #       converted at the end when it is saved to FF_dict
        else:
            if 'dft' in qc_types:
                DFT_params, DFT_errors = curve_fit(harmonic, thetas, DFT_E,p0=[100.0,Geo_dict[str(int(round(len(thetas)/2.0)))]["angle"]],maxfev=10000)
            if 'mp2' in qc_types:
                MP2_params, MP2_errors = curve_fit(harmonic, thetas, MP2_E,p0=[100.0,Geo_dict[str(int(round(len(thetas)/2.0)))]["angle"]],maxfev=10000)

        ######################################
        ############# Write Info #############
        ######################################

        # Return to parent/"angles" folder
        os.chdir(working_dir)

        # Open file and write header
        with open(save_folder+'/'+"_".join(angles)+'_info.txt', 'w') as f:
            if 'dft' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol rad^-2\n{:20s} {:< 12.6f} degrees\n".format("k (DFT) =",DFT_params[0],"Theta_0 (DFT) =",DFT_params[1]*180/np.pi))
            if 'mp2' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol rad^-2\n{:20s} {:< 12.6f} degrees\n".format("k (MP2) =",MP2_params[0],"Theta_0 (MP2) =",MP2_params[1]*180/np.pi))
            if 'dft' in qc_types and 'mp2' in qc_types:
                f.write("{:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s}\n".format("Angle (degrees)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)","MP2_Energy","E-E0 (MP2;kcal/mol)","E-E0 (MP2_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["angle"]*180/np.pi,Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["angle"],DFT_params[0],DFT_params[1]),\
                            Geo_dict[str(i)]["MP2"],Geo_dict[str(i)]["MP2"],harmonic(Geo_dict[str(i)]["angle"],MP2_params[0],MP2_params[1])))
            elif 'dft' in qc_types:
                f.write("{:<25s} {:<25s} {:<25s} {:<25s}\n".format("Angle (degrees)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["angle"]*180/np.pi,Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["angle"],DFT_params[0],DFT_params[1])))
                    

        #### Generate 2D representation of the molecule KNOWN ISSUE WITH PLOTTING FRAGMENT BASED BONDS AND ANGLES, 2D REPRESENTATION IS TAKEN FROM THE MASTER FRAGMENT
        print("Generating 2D representation of the fit fragment...")
        adj_mat = Table_generator(Geo_dict["0"]["elements"],Geo_dict["0"]["geo"])
        geo_2D = kekule(Geo_dict["0"]["elements"],angle_atomtypes[angles_count],Geo_dict["0"]["geo"],adj_mat)

        ######################################
        ######### Generate Fit Plots #########
        ######################################
        print("Rendering fit plots...")

        if 'dft' in qc_types:

            fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
            draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
            if draw_scale > 1.0: draw_scale = 1.0

            # Add bonds
            linewidth = 2.0*draw_scale*5.0
            if linewidth > 2.0:
                linewidth = 2.0
            if linewidth < 1.0:
                linewidth = 1.0
            for count_i,i in enumerate(adj_mat):
                for count_j,j in enumerate(i):
                    if count_j > count_i:
                        if j == 1:
                            ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

            # Add Highlight
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[0]][0],geo_2D[angle_atoms[1]][0]],[geo_2D[angle_atoms[0]][1],geo_2D[angle_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[1]][0],geo_2D[angle_atoms[2]][0]],[geo_2D[angle_atoms[1]][1],geo_2D[angle_atoms[2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                            
            # Add atom labels
            fontsize = 20*draw_scale*5.0
            if fontsize > 24:
                fontsize = 24
            if fontsize < 9:
                fontsize = 9            
            for count_i,i in enumerate(geo_2D):
                if count_i in angle_atoms:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                else:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
            ax1.axis('image')
            ax1.axis('off')

            # Add data
            step = (max(thetas)-min(thetas))/100.0*180.0/np.pi
            xdata = np.arange(min(thetas)*180.0/np.pi,max(thetas)*180.0/np.pi+step,step)
            ydata = harmonic(xdata*np.pi/180.0,DFT_params[0],DFT_params[1])
            ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
            ax2.scatter(np.array(thetas)*180.0/np.pi,DFT_E,s=120,color=(0.0,0.1,0.8))
            ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
            ax2.set_xlabel(r'$\mathrm{\theta \, (degrees)}$',fontsize=20,labelpad=10)
            ax2.set_xlim([min(thetas)*180.0/np.pi, max(thetas)*180.0/np.pi])
            ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
            ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
            y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            ax2.yaxis.set_major_formatter(y_formatter)
            ax2.xaxis.set_major_formatter(x_formatter)
            [j.set_linewidth(3) for j in ax2.spines.values()]
            plt.tight_layout()

            Name = save_folder+'/'+"_".join(angles)+'-DFT_fit.png'
            plt.savefig(Name, bbox_inches=0,dpi=300)
            plt.savefig(save_folder+'/'+"_".join(angles)+'-DFT_fit.pdf', bbox_inches=0,dpi=300)
            paths_to_angle_plots_DFT+=[Name]
            plt.close(fig)

        if 'mp2' in qc_types:

            fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
            draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
            draw_scale = 1.0
            # Add bonds
            for count_i,i in enumerate(adj_mat):
                for count_j,j in enumerate(i):
                    if count_j > count_i:
                        if j == 1:
                            ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=2.0*draw_scale))
            # Add Highlight
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[0]][0],geo_2D[angle_atoms[1]][0]],[geo_2D[angle_atoms[0]][1],geo_2D[angle_atoms[1]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[1]][0],geo_2D[angle_atoms[2]][0]],[geo_2D[angle_atoms[1]][1],geo_2D[angle_atoms[2]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))
                            
            # Add atom labels
            for count_i,i in enumerate(geo_2D):
                if count_i in angle_atoms:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                else:
                    ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
            ax1.axis('image')
            ax1.axis('off')

            # Add data
            step = (max(thetas)-min(thetas))/100.0*180.0/np.pi
#            xdata = np.arange(min(thetas)*180.0/np.pi,max(thetas)*180.0/np.pi+1,0.1)
            xdata = np.arange(min(thetas)*180.0/np.pi,max(thetas)*180.0/np.pi+step,step)
#            xdata = np.arange(min(thetas)*180.0/np.pi,max(thetas)*180.0/np.pi+1,0.1)
            ydata = harmonic(xdata*np.pi/180.0,MP2_params[0],MP2_params[1])
            ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
            ax2.scatter(np.array(thetas)*180.0/np.pi,MP2_E,s=120,color=(0.0,0.1,0.8))
            ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
            ax2.set_xlabel(r'$\mathrm{\theta \, (degrees)}$',fontsize=20,labelpad=10)
            ax2.set_xlim([min(thetas)*180.0/np.pi, max(thetas)*180.0/np.pi])
            ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
            ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
            y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
            ax2.yaxis.set_major_formatter(y_formatter)
            ax2.xaxis.set_major_formatter(x_formatter)
            [j.set_linewidth(3) for j in ax2.spines.values()]
            plt.tight_layout()

            Name = save_folder+'/'+"_".join(angles)+'-MP2_fit.png'
            plt.savefig(Name, bbox_inches=0,dpi=300)
            plt.savefig(save_folder+'/'+"_".join(angles)+'-MP2_fit.pdf', bbox_inches=0,dpi=300)
            paths_to_angle_plots_MP2+=[Name]
            plt.close(fig)
    
        #############################################################
        ####### Update FF_dict with DFT and MP2 FF-parameters #######
        #############################################################
        # NOTE: the force constant is in units of kcal/mol/rad^2 but the equilibrium angle is in units of degrees (lammps harmonic angle convention). 
        FF_dict["angles"][angles]={}
        if 'dft' in qc_types:
            FF_dict["angles"][angles]["DFT"] = ["harmonic"] + [ DFT_params[0], DFT_params[1]*180.0/np.pi ]
        if 'mp2' in qc_types:
            FF_dict["angles"][angles]["MP2"] = ["harmonic"] + [ MP2_params[0], MP2_params[1]*180.0/np.pi ]

    # Return to parent folder after processing all angles
    os.chdir(working_dir)

    return paths_to_angle_plots_DFT, paths_to_angle_plots_MP2

# This is a messy wrapper for the angle fitting procedure. In brief, the 
# function looks up all the angle data, calculates harmonic fits, and
# generates some plots and data text files. All parameters are saved to FF_dict
def parse_bonds_sc(base_name,bond_types,bond_folders,bond_folders_frag,bond_atomtypes,modes_from_FF,eq_charges,save_folder,qc_types=[],find_restarts=0,second_d_thresh=4000.0,charge_origin=None,save_plots=True,gens=2):

    # In python all global variables must be declared in each scope
    global FF_dict,harm_range

    # Save parent directory variable
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_bond_plots_DFT=[]
    paths_to_bond_plots_MP2=[]

    ######################################
    ###### Process angle directories #####
    ######################################
    for bonds_count,bonds in enumerate(bond_types):

        print("\n{}".format("*"*167))
        print("* {:^163s} *".format("Processing bond {}".format(" ".join(bonds))))
        print("{}".format("*"*167))

        # Change into current bond directory
        os.chdir(bond_folders[bonds_count])

        # Parse angle data
        Geo_dict,bond_atoms,completion_flag = get_bond_angle_data(working_dir+'/'+save_folder+'/'+base_name+'_'+"_".join(bonds)+'_geos.xyz')
        
        # Avoid the parse is the data is incomplete
        if completion_flag == 0:
            os.chdir(working_dir)
            print("******************************************************************************************")
            print("  Warning: Scan jobs did not run to completion for bond {}...Skipping...".format("_".join(bonds)))
            print("******************************************************************************************")
            continue

        # Generate adjacency mat
        adj_mat = Table_generator(Geo_dict["0"]["elements"],Geo_dict["0"]["geo"])        
        #print(adj_mat)

        # ID atomtypes
        atom_types = id_types(Geo_dict["0"]["elements"],adj_mat,gens,geo=Geo_dict["0"]["geo"])
        #print(atom_types)
        
        # Find the mode types in this fragment
        Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(adj_mat,atom_types,return_all=1)

        #print(Angles)

        

        # Determine bonding matrix for the compound        
        frag_info = bond_folders_frag[bonds_count]
        qc_charge_key = bond_folders[bonds_count].split('/')[0]+'/geoopt'
        for k in eq_charges:
            if(fnmatch.fnmatch(k,'*'+frag_info+'*')):
               qc_charge_key = k 
        
        q_tmp   = int(round(sum(eq_charges[qc_charge_key])))
        bond_mat = find_lewis(atom_types, adj_mat, q_tot=q_tmp, b_mat_only=True,verbose=False)

        # Check if the reordering of angle_atoms is required
        if tuple(bond_atoms) not in Bonds:
            bond_atoms = bond_atoms[::-1]
        if tuple(bond_atoms) not in Bonds:            
            print("ERROR in parse_bonds: the fit bond was not discovered in the quantum chemistry geometry. Exiting...")
            print("Bonds: {}".format(Bonds))
            print("bond_atoms: {}".format(bond_atoms))
            quit()        

        # Fit_type
        fit_type = tuple([ atom_types[i] for i in bond_atoms ])

        # Parse modes in each configuration
        # NOTE: the active mode is removed from the list of bonds, but reincluded as the first entry in the iteration list (guarrantees first position for downstream processing)
        modes = []
        Bonds.remove(tuple(bond_atoms))
        for i in natural_sort(list(Geo_dict.keys())):
            tmp = []
            for count_j,j in enumerate([tuple(bond_atoms)] + Bonds + Angles + Dihedrals):

                # Only keep modes that involve the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                if   len(j) == 4 and len(set(bond_atoms).intersection(set(j[1:3]))) < 2: continue
                elif len(j) == 3 and len(set(bond_atoms).intersection(set(j))) < 2: continue
                elif len(j) == 2 and len(set(bond_atoms).intersection(set(j))) < 1: continue

                # Parse bond length
                if len(j) == 2:
                    length = norm(Geo_dict[i]['geo'][j[0]]-Geo_dict[i]['geo'][j[1]])
                    tmp += [((atom_types[j[0]],atom_types[j[1]]),length)]

                # Parse angle degree:
                if len(j) == 3:
                    atom_1 = Geo_dict[i]["geo"][j[0]]
                    atom_2 = Geo_dict[i]["geo"][j[1]]
                    atom_3 = Geo_dict[i]["geo"][j[2]]
                    theta =  np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))
                    tmp += [((atom_types[j[0]],atom_types[j[1]],atom_types[j[2]]),theta)]

            # Dihedrals are handled separately because the indexing of Dihedrals needs to be matched with Dihedral_types
            for count_j,j in enumerate(Dihedrals):                

                # Only keep modes that involve one of the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                if   len(j) == 4 and len(set(bond_atoms).intersection(set(j[1:3]))) == 0: continue
                elif len(set(bond_atoms).intersection(set(j))) == 0: continue

                # Parse harmonic dihedral degree
                if 2 in [ int(k[j[1],j[2]]) for k in bond_mat ]:
                    v1 = (Geo_dict[i]["geo"][j[1]] - Geo_dict[i]["geo"][j[0]])
                    v2 = (Geo_dict[i]["geo"][j[2]] - Geo_dict[i]["geo"][j[1]])
                    v3 = (Geo_dict[i]["geo"][j[3]] - Geo_dict[i]["geo"][j[2]])
                    angle = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )
                    if angle >  np.pi*3.0/2.0: angle = angle - 2.0*np.pi                        
                    if angle < -np.pi/2.0: angle = angle + 2.0*np.pi
                    tmp += [((atom_types[j[0]],atom_types[j[1]],atom_types[j[2]],atom_types[j[3]]),angle)]

            # Append the mode information for the current configuration
            modes += [tmp]

        # Append the scanned mode angles to Geo_dict (this is a holdover from the way this function was originally written, with Geo_dict containing the scanned angles
        for count_i,i in enumerate(modes):
            Geo_dict[str(count_i)]["length"] = i[0][1]

        ######################################
        ###### Calculate Fit Parameters ######
        ######################################


        # Collect list of energies and angles
        lengths = []
        DFT_E   = []
        MP2_E   = []        
        for i in range(len(list(Geo_dict.keys()))):
            lengths += [modes[i][0][1]]
            if 'dft' in qc_types:
                DFT_E += [Geo_dict[str(i)]["DFT"]]
            if 'mp2' in qc_types:
                MP2_E += [Geo_dict[str(i)]["MP2"]]

        # Update the charges from the force-field if all of them are available.
        frag_info = bond_folders_frag[bonds_count]
        qc_charge_key = bond_folders[bonds_count].split('/')[0]+'/geoopt'
        for k in eq_charges:
            if(fnmatch.fnmatch(k,'*'+frag_info+'*')):
               qc_charge_key = k 
        fit_charges   = deepcopy(eq_charges[qc_charge_key])
        if charge_origin is not None:
            if True not in [ charge_origin[i] == "eq_config" for i in atom_types ]:
                for i in range(len(fit_charges)): fit_charges[i] = FF_dict["charges"][atom_types[i]]
            else:
                print("WARNING: Missing partial charges for some or all of the atomtypes in this fragment, defaulting to equilibrium partial charges.") 
        elif False not in [ atom_types[j] in list(FF_dict["charges"].keys()) for j in range(len(fit_charges)) ]:
            for i in range(len(fit_charges)): fit_charges[i] = FF_dict["charges"][atom_types[i]]

        # Calculate electrostatic and vdw energy corrections in each configuration
        fit_LJ        = initialize_VDW(atom_types,[],verbose=False)  
        EC_corr = []
        LJ_corr = []

        # Reassemble modes (since a bond was removed) and calculate the electrostatic (EC) and lennard-jones (LJ) contributions to each configuration
        Bonds,Angles,Dihedrals,One_fives = Find_modes(adj_mat,atom_types,return_all=0)
        for i in range(len(list(Geo_dict.keys()))):            
            EC_corr += [E_coul(Geo_dict[str(i)]["geo"],adj_mat,atom_types,fit_charges,[Bonds,Angles,Dihedrals],(0.0,0.0,0.0),"AA")]
            LJ_corr += [E_LJ(Geo_dict[str(i)]["geo"],adj_mat,atom_types,fit_charges,[Bonds,Angles,Dihedrals],(0.0,0.0,0.0),"AA",LJ_dict=fit_LJ)]

        # Sort by lengths
        if 'dft' in qc_types and 'mp2' in qc_types:
            together = sorted(zip(lengths,DFT_E,MP2_E,EC_corr,LJ_corr,modes))
            lengths  = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            MP2_E    = [ i[2] for i in together ]            
            EC_corr  = [ i[3] for i in together ]            
            LJ_corr  = [ i[4] for i in together ]            
            modes    = [ i[5] for i in together ]
        elif 'dft' in qc_types:
            together = sorted(zip(lengths,DFT_E,EC_corr,LJ_corr,modes))
            lengths  = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            EC_corr  = [ i[2] for i in together ]
            LJ_corr  = [ i[3] for i in together ]
            modes    = [ i[4] for i in together ]
        elif 'mp2' in qc_types:
            together = sorted(zip(lengths,MP2_E,EC_corr,LJ_corr,modes))
            lengths  = [ i[0] for i in together ]
            MP2_E    = [ i[1] for i in together ]
            EC_corr  = [ i[2] for i in together ]
            LJ_corr  = [ i[3] for i in together ]
            modes    = [ i[4] for i in together ]

        # Create a duplicate of these lists in case a reversion is required after extensive pruning and job regeneration
        lengths_bkp  = deepcopy(lengths)
        DFT_E_bkp    = deepcopy(DFT_E)
        MP2_E_bkp    = deepcopy(MP2_E)
        EC_corr_bkp  = deepcopy(EC_corr)
        LJ_corr_bkp  = deepcopy(LJ_corr)
        modes_bkp    = deepcopy(modes)        

        # If the minimum is near the edge of the scanned range then automatically perform a reoptimization
        min_ind = DFT_E.index(min(DFT_E)) 
        resub_flag = 0
        if min_ind in [0,1,len(DFT_E)-1,len(DFT_E)-2]:
            print("Minimum lies near the edge of the scanned range...")
            resub_flag = 1
        
        # Remove edge effects (i.e. edges with the improper derivative)
        # The value of correction_flag is used to toggle a diagnostic print statement
        delta = np.abs(lengths[1]-lengths[0])
        start_ind = 0
        end_ind = len(lengths) - 1
        correction_flag = 0
        if resub_flag == 0:
            for count_i,i in enumerate(lengths[:-1]):
                if DFT_E[count_i+1] - DFT_E[count_i] < 0: 
                    start_ind = count_i
                    break
            if start_ind != 0: correction_flag = 1

            for i in range(len(lengths))[::-1]:
                if DFT_E[i]-DFT_E[i-1] > 0: 
                    end_ind = i
                    break
            if end_ind != len(lengths)-1: correction_flag = 1
        
        # Assembled edited lists and print diagnostic if trimming occured
        if correction_flag == 1:
            print("Data points with improper derivatives were removed to increase the quality of fit...")
        if 'dft' in qc_types: DFT_E = DFT_E[start_ind:end_ind+1]
        if 'mp2' in qc_types: MP2_E = MP2_E[start_ind:end_ind+1]
        lengths = lengths[start_ind:end_ind+1]
        modes  = modes[start_ind:end_ind+1]
        EC_corr = EC_corr[start_ind:end_ind+1]
        LJ_corr = LJ_corr[start_ind:end_ind+1]

        # Remove derivative discontinuities
        # NOTE: the threshold of 4000.0 kcal^2/rad is empirical, it was found to consistently remove most derivative 
        #       discontinuities without removing any continuous data.
        second_d_thresh = float(second_d_thresh)

        # Remove datapoints until derivative discontinuities have been removed
        correction_flag = 0
        clean_flag = 0
        while clean_flag == 0 and len(DFT_E) > 2:

            # Loop over the energies and calculate the second derivatives. 
            # NOTE: the if/else construction is necessary to deal with the endpoints
            clean_flag = 1
            del_list = []            
            for i in range(len(DFT_E))[1:-1]:                

                second_d = ( ( DFT_E[i+1] - DFT_E[i] ) / ( lengths[i+1] - lengths[i] ) - ( DFT_E[i] - DFT_E[i-1] ) / ( lengths[i] - lengths[i-1] ) ) / ( lengths[i] - lengths[i-1] ) 

                if second_d < 0 or np.abs(second_d) > second_d_thresh:
                    correction_flag = 1
                    clean_flag = 0
                    del_list += [i]
                    break
                
            DFT_E   = [ i for count_i,i in enumerate(DFT_E) if count_i not in del_list ]
            lengths = [ i for count_i,i in enumerate(lengths) if count_i not in del_list ]
            modes   = [ i for count_i,i in enumerate(modes) if count_i not in del_list ]
            EC_corr = [ i for count_i,i in enumerate(EC_corr) if count_i not in del_list ]
            LJ_corr = [ i for count_i,i in enumerate(LJ_corr) if count_i not in del_list ]

        if correction_flag == 1:
            print("Data points with positive second derivatives or discontinuities were removed to increase the quality of the fit...")

        # If more than a third of the data has been discarded then the script generates a reoptimization job using the lowest
        # energy conformation from the scan.
        if (len(lengths) <= len(list(Geo_dict.keys()))*4.0/5.0 or resub_flag == 1) and find_restarts == 1:

            min_folder = [ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ].index(min([ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ]))
            template_file = str(min_folder) + '/' + bond_folders[bonds_count].split('/')[0] + '_' + '-'.join(bond_folders[bonds_count].split('/')[2].split('_')) + '_' + str(min_folder) + '.in'
            if os.path.isdir("REDO"):
                print("Input files for resubmission have already been generated for this mode...")
            else:
                print("Generating the input files for a reoptimization based on the lowest energy configuration...")
                gen_reoptimization(template=template_file,geo=Geo_dict[str(min_folder)]['geo'],elements=Geo_dict[str(min_folder)]["elements"],delta=delta,N_steps=len(list(Geo_dict.keys())))
            #os.chdir(working_dir)
            #continue

            # Revert to the original modes so that the user can assess the data and/or use the best fit as is. 
            lengths = lengths_bkp 
            DFT_E   = DFT_E_bkp
            MP2_E   = MP2_E_bkp   
            EC_corr = EC_corr_bkp 
            LJ_corr = LJ_corr_bkp 
            modes   = modes_bkp   

        # If the program is only searching for restarts then the fit is avoided
        if find_restarts == 1:
            os.chdir(working_dir)
            continue

        # If too much of the curve has been removed, proceed with the original data and use the best fit as is
        if len(lengths) < 3:
            print("Too much data was removed by the smoothing algorithm, reverting to raw data...")
            lengths = lengths_bkp 
            DFT_E   = DFT_E_bkp
            MP2_E   = MP2_E_bkp   
            EC_corr = EC_corr_bkp 
            LJ_corr = LJ_corr_bkp 
            modes   = modes_bkp   

        # Subtract the EC and LJ corrections from the DFT fit 
        min_ind_0 = next( count_i for count_i,i in enumerate(DFT_E) if i == min(DFT_E) )
        EC_corr = np.array(EC_corr) 
        LJ_corr = np.array(LJ_corr) 
        DFT_E   = np.array(DFT_E) - EC_corr - LJ_corr

        # Normalize energies by the minimum value
        if 'dft' in qc_types: DFT_E = np.array(DFT_E)-min(DFT_E)
        if 'mp2' in qc_types: MP2_E = np.array(MP2_E)-min(MP2_E)

        # Print diagnostic
        print("constrained atoms: {}".format(" ".join([ str(i) for i in bond_atoms])))
        print("length range: {: <6.4f}-{: <6.4f}".format(min(np.array(lengths)),max(np.array(lengths))))
        print("Calculating harmonic fits...")

        # Use read parameters if available (params holds force constant and theta_0)
        read_flag = False
        fallback_flag = 0
        if bonds in modes_from_FF:
            print("Using parameters for this bond from the supplied FF-file(s)...")
            if "dft" in qc_types:
                read_flag = True
                DFT_params = np.array(FF_dict["bonds"][bonds]["DFT"][1:])
            if "mp2" in qc_types:
                MP2_params = np.array(FF_dict["bonds"][bonds]["MP2"][1:])

        # Calculate parameters (params holds force constant and theta_0)    
        # NOTE: the force constant is in units of kcal/mol/ang^2.
        # NOTE: the initial guess for the mode being fit is set to the direct fitting parameters
        else:

            initial_guess = []
            min_ind = next( count_i for count_i,i in enumerate(DFT_E) if i == min(DFT_E) )
            harm_range=[lengths[min_ind]-0.2,lengths[min_ind]+0.2]
            if len(modes[min_ind][0][0]) == 2:
                lstsq_guess,errs = curve_fit(harmonic_norm, lengths, DFT_E,p0=[100.0,lengths[min_ind],0.0],maxfev=1000000)
            else:
                lstsq_guess,errs = curve_fit(harmonic_norm, thetas, DFT_E,p0=[50.0,thetas[min_ind],0.0],maxfev=1000000)                

            # Check if the equilibrium displacement falls outside the range that was actually scanned.
            # If so, then the program falls back on a direct lstsq fit to the QC data
            fallback_flag = 0
            if lstsq_guess[1] < (min(lengths) - 0.1) or lstsq_guess[1] > (max(lengths) + 0.1):
                DFT_E = DFT_E + EC_corr + LJ_corr
                DFT_E = np.array(DFT_E)-min(DFT_E)
                print("Optimal equilibrium length lies outside the scanned range. Falling back on direct least-squares fitting...")
                DFT_params = curve_fit(harmonic_norm, lengths, DFT_E,p0=[100.0,lengths[min_ind],0.0],maxfev=1000000)[0][:2]
                print("xhi2 (kcal^2/mol^2): {}".format(np.mean((harmonic(lengths,DFT_params[0],DFT_params[1])-DFT_E)**(2.0))))
                fallback_flag = 1

            # Else, proceed with self-consistent fit of the mode 
            else:

                # Apply minimum condition to the fit
                if lstsq_guess[0] < 50.0: lstsq_guess[0] = 50.0
                if lstsq_guess[0] > 5000.0: lstsq_guess[0] = 5000.0

                # Collect the unique mode types in the current fit
                mode_types = []
                mode_eq    = []
                for i in modes[min_ind]:
                    mode_types += [i[0]]
                    mode_eq += [i[1]]

                # Generate initial guess parameters and fit bounds
                # NOTE: the normalization value is seeded with 0.0 and is fit without bounds
                # NOTE: the force constant for the scanned mode is bound from above by the lstsq fit value
                # NOTE: the approximate force constants for all unscanned bonds are derived from UFF
                # NOTE: the force constants for all unscanned angles are approximated as 100 kcal/mol as a reasonable initial guess
                # NOTE: for harmonic dihedrals the equilibrium position isn't used so it is set to a dummy value of 0.0
                bond_dict = UFF_bonds(Bond_types)
                initial_guess = [lstsq_guess[2]]
                bounds = [(None,None)]            
                min_scale = 0.25
                for count_i,i in enumerate(mode_types):

                    # Assign bounds
                    if count_i == 0:
                        lower_k = max([lstsq_guess[0]*min_scale,50.0])                    
                        if len(i) == 2:                            
                            bounds += [(lower_k,lstsq_guess[0]),(mode_eq[count_i]-0.2,mode_eq[count_i]+0.2)]
                        elif len(i) == 3:
                            bounds += [(lower_k,lstsq_guess[0]),(mode_eq[count_i]-10.0*np.pi/180.0,mode_eq[count_i]+10.0*np.pi/180.0)]
                        elif len(i) == 4:
                            bounds += [(lower_k,lstsq_guess[0]),(0,None)]
                    else:
                        bounds += [(0,None),(0,None)]

                    # Assign initial mode parameter guesses
                    if count_i == 0:
                        initial_guess += list(lstsq_guess[:2])
                    elif len(i) == 2:
                        if i in list(FF_dict["bonds"].keys()):
                            initial_guess += [FF_dict["bonds"][i]["DFT"][1],mode_eq[count_i]]
                        elif i[::-1] in list(FF_dict["bonds"].keys()):
                            initial_guess += [FF_dict["bonds"][i[::-1]]["DFT"][1],mode_eq[count_i]]                        
                        else:
                            initial_guess += [600.0,mode_eq[count_i]]
                    elif len(i) == 3:
                        if i in list(FF_dict["angles"].keys()):
                            initial_guess += [FF_dict["angles"][i]["DFT"][1],mode_eq[count_i]]                  
                        elif i[::-1] in list(FF_dict["angles"].keys()):
                            initial_guess += [FF_dict["angles"][i[::-1]]["DFT"][1],mode_eq[count_i]]
                        else:
                            initial_guess += [60.0,mode_eq[count_i]]
                    elif len(i) == 4:
                        if i in list(FF_dict["dihedrals_harmonic"].keys()):
                            initial_guess += [FF_dict["dihedrals_harmonic"][i]["DFT"][1],0.0]
                        elif i[::-1] in list(FF_dict["dihedrals_harmonic"].keys()):
                            initial_guess += [FF_dict["dihedrals_harmonic"][i[::-1]]["DFT"][1],0.0]
                        else:
                            initial_guess += [5.0,0.0]

                # Perform a fit of only the scanned mode while keeping all others constant (mutually adjusts the force constant and equilibrium position)
                fit_func = lambda x: fit_bonds_angles_sc_single(x,*initial_guess[3:],modes=modes,fit_modes=mode_types,E_0=DFT_E)            
                params = minimize(fit_func,initial_guess[:3],method='L-BFGS-B',bounds=[(None,None),bounds[1],bounds[2]],options={'gtol':0.0,'ftol':1.0E-20,'maxiter':1000000,'maxfun':1000000,'maxls':20}).x
                initial_guess[:3] = params

                # Perform a 2-D optimization to find the best force constant for this mode while holding the other parameters constant
                for j in range(1):
                    xhi2_min = None
                    for i in np.arange(initial_guess[1]*min_scale,initial_guess[1]+1.0):
                        initial_guess[1] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    xhi2_min = None
                    for i in np.linspace(best_guess-1.0,best_guess+1.0,10):
                        initial_guess[1] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    initial_guess[1] = best_guess
                    xhi2_min = None
                    for i in np.arange(initial_guess[2]-0.01,initial_guess[2]+0.01,0.001):
                        initial_guess[2] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    xhi2_min = None
                    for i in np.linspace(best_guess-0.001,best_guess+0.001,300):
                        initial_guess[2] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    initial_guess[2] = best_guess

                # Generate fits with varying initial guess for the force constant of the mode being fit
                xhi2_min = None
                mode_guesses = [initial_guess[1]]
                for z in mode_guesses:

                    # Perform fit
                    initial_guess[1] = z
                    if 'dft' in qc_types:

                        # Perform the fit for the current initial guess
                        fit_func = lambda x: fit_bonds_angles_sc_2(*x,modes=modes,fit_modes=mode_types,E_0=DFT_E)
                        params = minimize(fit_func,initial_guess,method='L-BFGS-B',bounds=bounds,options={'gtol':0.0,'ftol':1.0E-20,'maxiter':1000000,'maxfun':1000000,'maxls':20}).x
                        xhi2_current = fit_bonds_angles_sc_2(*params,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)

                        # Update the fit parameters if a lower energy was obtained
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            DFT_params = deepcopy(params[1:3])                        
                            params[1] = 0.0
                            E_corr = fit_bonds_angles_sc_2(*params,modes=modes,fit_modes=mode_types,penalty=False,E_0=DFT_E)
                            xhi2_min = xhi2_current

                    if 'mp2' in qc_types:
                        fit_func = lambda x: fit_bonds_angles_sc(*x,modes=modes,E_0=MP2_E)
                        params = minimize(fit_func,initial_guess,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':1E-20}).x

                # Fall back on the least-squares guess if the fit was bad
                if xhi2_min > 1E-7:
                    DFT_E = DFT_E + EC_corr + LJ_corr
                    DFT_E = np.array(DFT_E)-min(DFT_E)
                    print("Self-consistent fit failed to converge. This is probably a crowded mode with a lot of electrostatic contributions to the scan potential. Falling back on direct least-squares fitting...")
                    DFT_params = curve_fit(harmonic_norm, lengths, DFT_E,p0=[100.0,lengths[min_ind],0.0],maxfev=1000000)[0][:2]
                    print("xhi2 (kcal^2/mol^2): {}".format(np.mean((harmonic(lengths,DFT_params[0],DFT_params[1])-DFT_E)**(2.0))))
                    fallback_flag = 1

                # Print diagnostic
                else:
                    print("xhi2 (kcal^2/mol^2): {}".format(xhi2_min))

            ####### END OF TYPE BASED FIT COMMANDS #######

        #####################################
        ############# Write Info ############
        #####################################

        # Return to parent/"bonds" folder
        os.chdir(working_dir)

        # Open file and write header
        with open(save_folder+'/'+"_".join(bonds)+'_info.txt', 'w') as f:
            if 'dft' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol Ang^-2\n{:20s} {:< 12.6f} Ang\n".format("k (DFT) =",DFT_params[0],"r_0 (DFT) =",DFT_params[1]))
            if 'mp2' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol Ang^-2\n{:20s} {:< 12.6f} Ang\n".format("k (MP2) =",MP2_params[0],"r_0 (MP2) =",MP2_params[1]))
            if 'dft' in qc_types and 'mp2' in qc_types:
                f.write("{:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s}\n".format("r (Angstroms)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)","MP2_Energy","E-E0 (MP2;kcal/mol)","E-E0 (MP2_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["length"],Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["length"],DFT_params[0],DFT_params[1]),\
                            Geo_dict[str(i)]["MP2"],Geo_dict[str(i)]["MP2"],harmonic(Geo_dict[str(i)]["length"],MP2_params[0],MP2_params[1])))
            elif 'dft' in qc_types: 
                f.write("{:<25s} {:<25s} {:<25s} {:<25s}\n".format("r (Angstroms)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["length"],Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["length"],DFT_params[0],DFT_params[1])))


        # Save plots
        if save_plots is True:

            #### Generate 2D representation of the molecule KNOWN ISSUE WITH PLOTTING FRAGMENT BASED BONDS AND ANGLES, 2D REPRESENTATION IS TAKEN FROM THE MASTER FRAGMENT
            print("Generating 2D representation of the fit fragment...")
            adj_mat = Table_generator(Geo_dict["0"]["elements"],Geo_dict["0"]["geo"])
            geo_2D = kekule(Geo_dict["0"]["elements"],bond_atomtypes[bonds_count],Geo_dict["0"]["geo"],adj_mat)

            ######################################
            ######### Generate Fit Plots #########
            ######################################
            print("Rendering fit plots...")

            if 'dft' in qc_types:

                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                if draw_scale > 1.0: draw_scale = 1.0

                # Add bonds
                linewidth = 2.0*draw_scale*5.0
                if linewidth > 2.0:
                    linewidth = 2.0
                if linewidth < 1.0:
                    linewidth = 1.0
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[bond_atoms[0]][0],geo_2D[bond_atoms[1]][0]],[geo_2D[bond_atoms[0]][1],geo_2D[bond_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                # Add atom labels
                fontsize = 20*draw_scale*5.0
                if fontsize > 24:
                    fontsize = 24
                if fontsize < 9:
                    fontsize = 9            
                for count_i,i in enumerate(geo_2D):
                    if count_i in bond_atoms:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Add data
                xdata = np.array(lengths)
                if fallback_flag == 1 or read_flag is True:
                    ydata = harmonic(xdata,DFT_params[0],DFT_params[1]) 
                else:
                    ydata = harmonic(xdata,DFT_params[0],DFT_params[1]) + E_corr + EC_corr + LJ_corr
                min_corr = min(ydata)
                ax2.plot(xdata,ydata-min_corr,'--',linewidth=3,color=(0.0,0.1,0.8))
                min_corr = min(DFT_E+EC_corr+LJ_corr)
                if fallback_flag == 1 or read_flag is True:
                    ax2.scatter(np.array(lengths),DFT_E,s=120,color=(0.0,0.1,0.8))
                else:
                    ax2.scatter(np.array(lengths),DFT_E+EC_corr+LJ_corr-min_corr,s=120,color=(0.0,0.1,0.8))
                if fallback_flag == 1:
                    ax2.text(0.9, 0.8,'*', ha='center', fontsize=40, va='center', transform=ax2.transAxes)
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{Length (\AA)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(lengths), max(lengths)])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()
                Name = save_folder+'/'+"_".join(bonds)+'-DFT_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(save_folder+'/'+"_".join(bonds)+'-DFT_fit.pdf', bbox_inches=0,dpi=300)
                paths_to_bond_plots_DFT+=[Name]
                plt.close(fig)

            if 'mp2' in qc_types:

                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                draw_scale = 1.0
                # Add bonds
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=2.0*draw_scale))
                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[bond_atoms[0]][0],geo_2D[bond_atoms[1]][0]],[geo_2D[bond_atoms[0]][1],geo_2D[bond_atoms[1]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))

                # Add atom labels
                for count_i,i in enumerate(geo_2D):
                    if count_i in bond_atoms:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Plot Data
                step = (max(lengths)-min(lengths))/100.0
                xdata = np.arange(min(lengths),max(lengths)+step,step)
                ydata = harmonic(xdata,MP2_params[0],MP2_params[1])
                ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
                ax2.scatter(np.array(lengths),MP2_E,s=120,color=(0.0,0.1,0.8))
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{Length \, (\AA)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(lengths), max(lengths)])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = save_folder+'/'+"_".join(bonds)+'-MP2_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(save_folder+'/'+"_".join(bonds)+'-MP2_fit.pdf', bbox_inches=0,dpi=300)
                paths_to_bond_plots_MP2+=[Name]
                plt.close(fig)
        
        ##############################################################
        ######## Update FF_dict with DFT and MP2 FF-parameters #######
        ##############################################################
        FF_dict["bonds"][bonds]={}
        if 'dft' in qc_types:
            FF_dict["bonds"][bonds]["DFT"] = ["harmonic"] + DFT_params.tolist()
        if 'mp2' in qc_types:
            FF_dict["bonds"][bonds]["MP2"] = ["harmonic"] + MP2_params.tolist()

    # Return to parent folder after processing all bonds
    os.chdir(working_dir)

    return paths_to_bond_plots_DFT, paths_to_bond_plots_MP2

# This is a messy wrapper for the angle fitting procedure. In brief, the 
# function looks up all the angle data, calculates harmonic fits, and
# generates some plots and data text files. All parameters are saved to FF_dict
def parse_angles_sc(base_name,angle_types,angle_folders,angle_folders_frag,angle_atomtypes,modes_from_FF,eq_charges,save_folder,qc_types=[],find_restarts=0,second_d_thresh=500.0,charge_origin=None,save_plots=True,gens=2):

    # In python all global variables must be declared in each scope
    global FF_dict,harm_range

    # Save parent directory
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_angle_plots_DFT=[]
    paths_to_angle_plots_MP2=[]

    ######################################
    ###### Process angle directories #####
    ######################################
    for angles_count,angles in enumerate(angle_types):

        print("\n{}".format("*"*167))
        print("* {:^163s} *".format("Processing angle {}".format(" ".join(angles))))
        print("{}".format("*"*167))

        # Change into current bond directory
        os.chdir(angle_folders[angles_count])

        # Parse angle data
        #### DEBUG: angle atoms are the same with the old driver

        Geo_dict,angle_atoms,completion_flag = get_bond_angle_data(working_dir+'/'+save_folder+'/'+base_name+'_'+"_".join(angles)+'_geos.xyz')
        
        # Avoid the parse is the data is incomplete
        if completion_flag == 0:
            os.chdir(working_dir)
            print("******************************************************************************************")
            print("  Warning: Scan jobs did not run to completion for angle {}...Skipping...".format("_".join(angles)))
            print("******************************************************************************************")
            continue

        # Generate adjacency mat
        adj_mat = Table_generator(Geo_dict["0"]["elements"],Geo_dict["0"]["geo"])        

        # ID atomtypes
        #### DEBUG: atom_types are correct
        #### DEBUG: adj_mat is different
        #### Angels are not correct (cuz adj_mat)
        #### Cause the error
        atom_types = id_types(Geo_dict["0"]["elements"],adj_mat,gens,geo=Geo_dict["0"]["geo"])
        #print(atom_types)
        #print(adj_mat)
        # Find the mode types in this fragment
        Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(adj_mat,atom_types,return_all=1)

        # Determine bonding matrix for the compound        
        frag_info = angle_folders_frag[angles_count]
        qc_charge_key = angle_folders[angles_count].split('/')[0]+'/geoopt'
        for k in eq_charges:
            if(fnmatch.fnmatch(k,'*'+frag_info+'*')):
               qc_charge_key = k 
        q_tmp   = int(round(sum(eq_charges[qc_charge_key])))
        bond_mat = find_lewis(atom_types, adj_mat, q_tot=q_tmp, b_mat_only=True,verbose=False)

        # Check if the reordering of angle_atoms is required
        if tuple(angle_atoms) not in Angles:
            angle_atoms = angle_atoms[::-1]
        if tuple(angle_atoms) not in Angles:
            print("ERROR in parse_angles: the fit angle was not discovered in the quantum chemistry geometry. Exiting...")
            quit()        

        # Parse modes in each configuration
        # NOTE: the active mode is removed from the list of angles, but reincluded as the first entry in the iteration list (guarrantees first position for downstream processing)
        modes = []
        Angles.remove(tuple(angle_atoms))
        for i in natural_sort(list(Geo_dict.keys())):
            tmp = []
            for count_j,j in enumerate([tuple(angle_atoms)] + Bonds + Angles):

                # Only keep modes that involve the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                if   len(j) == 4 and len(set(angle_atoms).intersection(set(j[1:3]))) < 2: continue
                elif len(j) == 3 and len(set(angle_atoms).intersection(set(j))) < 2: continue
                elif len(j) == 2 and len(set([angle_atoms[1]]).intersection(set(j))) < 1: continue

                # # Only keep modes that involve one of the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                # if   len(j) == 4 and len(set(angle_atoms).intersection(set(j[1:3]))) == 0: continue
                # elif len(set(angle_atoms).intersection(set(j))) == 0: continue

                # Parse bond length
                if len(j) == 2:
                    length = norm(Geo_dict[i]['geo'][j[0]]-Geo_dict[i]['geo'][j[1]])
                    tmp += [((atom_types[j[0]],atom_types[j[1]]),length)]

                # Parse angle degree:
                if len(j) == 3:
                    atom_1 = Geo_dict[i]["geo"][j[0]]
                    atom_2 = Geo_dict[i]["geo"][j[1]]
                    atom_3 = Geo_dict[i]["geo"][j[2]]
                    theta =  np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))
                    tmp += [((atom_types[j[0]],atom_types[j[1]],atom_types[j[2]]),theta)]

            # Dihedrals are handled separately because the indexing of Dihedrals needs to be matched with Dihedral_types
            for count_j,j in enumerate(Dihedrals):                

                # Only keep modes that involve one of the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                if   len(j) == 4 and len(set(angle_atoms).intersection(set(j[1:3]))) == 0: continue
                elif len(set(angle_atoms).intersection(set(j))) == 0: continue

                # Parse harmonic dihedral degree
                if 2 in [ int(k[j[1],j[2]]) for k in bond_mat ]:
                    v1 = (Geo_dict[i]["geo"][j[1]] - Geo_dict[i]["geo"][j[0]])
                    v2 = (Geo_dict[i]["geo"][j[2]] - Geo_dict[i]["geo"][j[1]])
                    v3 = (Geo_dict[i]["geo"][j[3]] - Geo_dict[i]["geo"][j[2]])
                    angle = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )
                    if angle >  np.pi*3.0/2.0: angle = angle - 2.0*np.pi                        
                    if angle < -np.pi/2.0: angle = angle + 2.0*np.pi
                    tmp += [((atom_types[j[0]],atom_types[j[1]],atom_types[j[2]],atom_types[j[3]]),angle)]

            # Append the mode information for the current configuration
            modes += [tmp]

        # Append the scanned mode angles to Geo_dict (this is a holdover from the way this function was originally written, with Geo_dict containing the scanned angles
        for count_i,i in enumerate(modes):
            Geo_dict[str(count_i)]["angle"] = i[0][1]

        ######################################
        ###### Calculate Fit Parameters ######
        ######################################

        # Collect list of energies and angles
        DFT_E  = []
        MP2_E  = []        
        thetas = []
        for i in range(len(list(Geo_dict.keys()))):
            thetas += [modes[i][0][1]]
            if 'dft' in qc_types:
                DFT_E += [Geo_dict[str(i)]["DFT"]]
            if 'mp2' in qc_types:
                MP2_E += [Geo_dict[str(i)]["MP2"]]

        # Update the charges from the force-field if all of them are available.
        frag_info = angle_folders_frag[angles_count]
        qc_charge_key = angle_folders[angles_count].split('/')[0]+'/geoopt'
        #print(frag_info)
        for k in eq_charges:
            if(fnmatch.fnmatch(k,'*'+frag_info+'*')):
               qc_charge_key = k 
        #print(qc_charge_key)
        fit_charges   = deepcopy(eq_charges[qc_charge_key])
        if charge_origin is not None:
            if True not in [ charge_origin[i] == "eq_config" for i in atom_types ]:
                for i in range(len(fit_charges)): fit_charges[i] = FF_dict["charges"][atom_types[i]]
            else:
                print("WARNING: Missing partial charges for some or all of the atomtypes in this fragment, defaulting to equilibrium partial charges.") 
        elif False not in [ atom_types[j] in list(FF_dict["charges"].keys()) for j in range(len(fit_charges)) ]:
            for i in range(len(fit_charges)): fit_charges[i] = FF_dict["charges"][atom_types[i]]

        # Calculate electrostatic and vdw energy corrections in each configuration
        fit_LJ        = initialize_VDW(atom_types,[],verbose=False)  
        EC_corr = []
        LJ_corr = []

        # Reassemble modes (since an angle was removed) and calculate the electrostatic (EC) and lennard-jones (LJ) contributions to each configuration
        Bonds,Angles,Dihedrals,One_fives = Find_modes(adj_mat,atom_types,return_all=0)
        for i in range(len(list(Geo_dict.keys()))):            
            EC_corr += [E_coul(Geo_dict[str(i)]["geo"],adj_mat,atom_types,fit_charges,[Bonds,Angles,Dihedrals],(0.0,0.0,0.0),"AA")]
            LJ_corr += [E_LJ(Geo_dict[str(i)]["geo"],adj_mat,atom_types,fit_charges,[Bonds,Angles,Dihedrals],(0.0,0.0,0.0),"AA",LJ_dict=fit_LJ)]

        # Check for linear angle
        # (i) a change in gradient and (ii) vicinity to 180.0 are used to check for linearity                
        linearity_flag = 0
        if thetas[1] - thetas[0] >= 0.0:            
            if False in [ (i-thetas[count_i]) > 0.0 for count_i,i in enumerate(thetas[1:]) ] and (180.0 - np.mean(thetas)*180.0/np.pi) < 5.0:
                linearity_flag = 1
                turn = [ (i-thetas[count_i]) > 0.0 for count_i,i in enumerate(thetas[1:]) ].index(False)+1
                thetas[turn:] = [ i+2.0*(np.pi-i) for i in thetas[turn:] ]
        elif thetas[1] - thetas[0] < 0:
            if False in [ (i-thetas[count_i]) < 0.0 for count_i,i in enumerate(thetas[1:]) ] and (180.0 - np.mean(thetas)*180.0/np.pi) < 5.0:
                linearity_flag = 1
                turn = [ (i-thetas[count_i]) < 0.0 for count_i,i in enumerate(thetas[1:]) ].index(False)+1
                thetas[turn:] = [ i+2.0*(np.pi-i) for i in thetas[turn:] ]
        if linearity_flag == 1: print("Angles were unwrapped due to the linearity of the angle.")

        # Sort by angles
        if 'dft' in qc_types and 'mp2' in qc_types:
            together = sorted(zip(thetas,DFT_E,MP2_E,EC_corr,LJ_corr,modes))
            thetas   = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            MP2_E    = [ i[2] for i in together ]            
            EC_corr  = [ i[3] for i in together ]            
            LJ_corr  = [ i[4] for i in together ]            
            modes    = [ i[5] for i in together ]
        elif 'dft' in qc_types:
            together = sorted(zip(thetas,DFT_E,EC_corr,LJ_corr,modes))
            thetas   = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            EC_corr  = [ i[2] for i in together ]
            LJ_corr  = [ i[3] for i in together ]
            modes    = [ i[4] for i in together ]
        elif 'mp2' in qc_types:
            together = sorted(zip(thetas,MP2_E,EC_corr,LJ_corr,modes))
            thetas   = [ i[0] for i in together ]
            MP2_E    = [ i[1] for i in together ]
            EC_corr  = [ i[2] for i in together ]
            LJ_corr  = [ i[3] for i in together ]
            modes    = [ i[4] for i in together ]

        # Create a duplicate of these lists in case a reversion is required after extensive pruning and job regeneration
        thetas_bkp   = deepcopy(thetas)
        DFT_E_bkp    = deepcopy(DFT_E)
        MP2_E_bkp    = deepcopy(MP2_E)
        EC_corr_bkp  = deepcopy(EC_corr)
        LJ_corr_bkp  = deepcopy(LJ_corr)
        modes_bkp    = deepcopy(modes)        

        # If the minimum is near the edge of the scanned range then automatically perform a reoptimization
        min_ind = DFT_E.index(min(DFT_E)) 
        resub_flag = 0
        if min_ind in [0,1,len(DFT_E)-1,len(DFT_E)-2]:
            print("Minimum lies near the edge of the scanned range...")
            resub_flag = 1
        
        # Remove edge effects (i.e. edges with the improper derivative)
        # The value of correction_flag is used to toggle a diagnostic print statement
        delta = np.abs(thetas[1]-thetas[0])*180.0/np.pi
        start_ind = 0
        end_ind = len(thetas) - 1
        correction_flag = 0
        if resub_flag == 0:
            for count_i,i in enumerate(thetas[:-1]):
                if DFT_E[count_i+1] - DFT_E[count_i] < 0: 
                    start_ind = count_i
                    break
            if start_ind != 0: correction_flag = 1

            for i in range(len(thetas))[::-1]:
                if DFT_E[i]-DFT_E[i-1] > 0: 
                    end_ind = i
                    break
            if end_ind != len(thetas)-1: correction_flag = 1
        
        # Assembled edited lists and print diagnostic if trimming occured
        if correction_flag == 1:
            print("Data points with improper derivatives were removed to increase the quality of fit...")
        if 'dft' in qc_types: DFT_E = DFT_E[start_ind:end_ind+1]
        if 'mp2' in qc_types: MP2_E = MP2_E[start_ind:end_ind+1]
        thetas = thetas[start_ind:end_ind+1]
        modes  = modes[start_ind:end_ind+1]
        EC_corr = EC_corr[start_ind:end_ind+1]
        LJ_corr = LJ_corr[start_ind:end_ind+1]

        # Remove derivative discontinuities
        # NOTE: the threshold of 500.0 kcal^2/rad is empirical, it was found to consistently remove most derivative 
        #       discontinuities without removing any continuous data.
        second_d_thresh = float(second_d_thresh)

        # Remove datapoints until derivative discontinuities have been removed
        correction_flag = 0
        clean_flag = 0
        while clean_flag == 0 and len(DFT_E) > 2:

            # Loop over the energies and calculate the second derivatives. 
            # NOTE: the if/else construction is necessary to deal with the endpoints
            clean_flag = 1
            del_list = []            
            for i in range(len(DFT_E))[1:-1]:                

                second_d = ( ( DFT_E[i+1] - DFT_E[i] ) / ( thetas[i+1] - thetas[i] ) - ( DFT_E[i] - DFT_E[i-1] ) / ( thetas[i] - thetas[i-1] ) ) / (thetas[i] - thetas[i-1] ) 

                if second_d < 0 or np.abs(second_d) > second_d_thresh:
                    correction_flag = 1
                    clean_flag = 0
                    del_list += [i]
                    break
                
            DFT_E = [ i for count_i,i in enumerate(DFT_E) if count_i not in del_list ]
            thetas = [ i for count_i,i in enumerate(thetas) if count_i not in del_list ]
            modes  = [ i for count_i,i in enumerate(modes) if count_i not in del_list ]
            EC_corr = [ i for count_i,i in enumerate(EC_corr) if count_i not in del_list ]
            LJ_corr = [ i for count_i,i in enumerate(LJ_corr) if count_i not in del_list ]

        if correction_flag == 1:
            print("Data points with positive second derivatives or discontinuities were removed to increase the quality of the fit...")

        # If more than a third of the data has been discarded then the script generates a reoptimization job using the lowest
        # energy conformation from the scan.
        if ( len(thetas) <= len(list(Geo_dict.keys()))*4.0/5.0 or resub_flag == 1 ) and find_restarts == 1:
            min_folder = [ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ].index(min([ Geo_dict[str(i)]["DFT"] for i in range(len(list(Geo_dict.keys()))) ]))
            template_file = str(min_folder) + '/' + angle_folders[angles_count].split('/')[0] + '_' + '-'.join(angle_folders[angles_count].split('/')[2].split('_')) + '_' + str(min_folder) + '.in'
            if os.path.isdir("REDO"):
                print("Input files for resubmission have already been generated for this mode...")
            else:
                print("Generating the input files for a reoptimization based on the lowest energy configuration...")
                gen_reoptimization(template=template_file,geo=Geo_dict[str(min_folder)]['geo'],elements=Geo_dict[str(min_folder)]["elements"],delta=delta,N_steps=len(list(Geo_dict.keys())))
            #os.chdir(working_dir)
            #continue

            # Revert to the original modes so that the user can assess the data and/or use the best fit as is. 
            thetas  = thetas_bkp 
            DFT_E   = DFT_E_bkp
            MP2_E   = MP2_E_bkp   
            EC_corr = EC_corr_bkp 
            LJ_corr = LJ_corr_bkp 
            modes   = modes_bkp   

        # If the program is only searching for restarts then the fit is avoided
        if find_restarts == 1:
            os.chdir(working_dir)
            continue

        # If too much of the curve has been removed, proceed with the original data and use the best fit as is
        if len(thetas) < 3:
            print("Too much data was removed by the smoothing algorithm, reverting to raw data...")
            thetas  = thetas_bkp 
            DFT_E   = DFT_E_bkp
            MP2_E   = MP2_E_bkp   
            EC_corr = EC_corr_bkp 
            LJ_corr = LJ_corr_bkp 
            modes   = modes_bkp   

        # Subtract the EC and LJ corrections from the DFT fit 
        EC_corr = np.array(EC_corr) 
        LJ_corr = np.array(LJ_corr) 
        DFT_E   = np.array(DFT_E) - EC_corr - LJ_corr

        # Normalize energies by the minimum value
        if 'dft' in qc_types: DFT_E = np.array(DFT_E)-min(DFT_E)
        if 'mp2' in qc_types: MP2_E = np.array(MP2_E)-min(MP2_E)

        # Print diagnostic
        print("constrained atoms: {}".format(" ".join([ str(i) for i in angle_atoms])))
        print("theta range: {: <6.4f}-{: <6.4f}".format(min(np.array(thetas))*180.0/np.pi,max(np.array(thetas))*180.0/np.pi))
        print("Calculating harmonic fits...")

        # Use read parameters if available (params holds force constant and theta_0)
        read_flag = False
        fallback_flag = 0
        if angles in modes_from_FF:
            print("Using parameters for this angle from the supplied FF-file(s)...")
            if "dft" in qc_types:
                read_flag = True
                DFT_params = np.array(FF_dict["angles"][angles]["DFT"][1:])
                DFT_params[1] = DFT_params[1]*np.pi/180.0
            if "mp2" in qc_types:
                MP2_params = np.array(FF_dict["angles"][angles]["MP2"][1:])
                MP2_params[1] = MP2_params[1]*np.pi/180.0

        # Calculate parameters (params holds force constant and theta_0)    
        # NOTE: the force constant is in units of kcal/mol/rad^2 but the equilibrium angle is in units of degrees 
        #       (lammps harmonic angle convention). The eq_angle is converted at the end when it is saved to FF_dict
        # NOTE: the initial guess for the mode being fit is set to the direct fitting parameters
        else:

            # Get initial guess
            min_ind = next( count_i for count_i,i in enumerate(DFT_E) if i == min(DFT_E) )
            harm_range=[thetas[min_ind]-5.0*np.pi/180.0,thetas[min_ind]+5.0*np.pi/180.0]
            if len(modes[min_ind][0][0]) == 2:
                lstsq_guess,errs = curve_fit(harmonic_norm, lengths, DFT_E,p0=[100.0,lengths[min_ind],0.0],maxfev=1000000)
            else:
                lstsq_guess,errs = curve_fit(harmonic_norm, thetas, DFT_E,p0=[50.0,thetas[min_ind],0.0],maxfev=1000000)                

            # Check if the equilibrium displacement falls outside the range that was actually scanned.
            # If so, then the program falls back on a direct lstsq fit to the QC data
            fallback_flag = 0
            if lstsq_guess[1] < (min(thetas)-5.0*np.pi/180.0) or lstsq_guess[1] > (max(thetas)+5.0*np.pi/180.0):
                DFT_E = DFT_E + EC_corr + LJ_corr
                DFT_E = np.array(DFT_E)-min(DFT_E)
                print("Optimal equilibrium angle lies outside the scanned range. Falling back on direct least-squares fitting...")
                DFT_params = curve_fit(harmonic_norm, thetas, DFT_E,p0=[50.0,thetas[min_ind],0.0],maxfev=1000000)[0][:2]
                print("xhi2 (kcal^2/mol^2): {}".format(np.mean((harmonic(thetas,DFT_params[0],DFT_params[1])-DFT_E)**(2.0))))
                fallback_flag = 1

            # Else, proceed with self-consistent fit of the mode 
            else:

                # Apply minimum condition to the fit
                if lstsq_guess[0] < 10.0: lstsq_guess[0] = 10.0
                if lstsq_guess[0] > 500: lstsq_guess[0] = 500.0

                # Collect the unique mode types in the current fit
                mode_types = []
                mode_eq    = []
                for i in modes[min_ind]:
                    mode_types += [i[0]]
                    mode_eq += [i[1]]

                # Generate initial guess parameters and fit bounds
                # NOTE: the normalization value is seeded with 0.0 and is fit without bounds
                # NOTE: the force constant for the scanned mode is bound from above by the lstsq fit value
                # NOTE: the approximate force constants for all unscanned bonds are derived from UFF
                # NOTE: the force constants for all unscanned angles are approximated as 100 kcal/mol as a reasonable initial guess
                # NOTE: for harmonic dihedrals the equilibrium position isn't used so it is set to a dummy value of 0.0
                bond_dict = UFF_bonds(Bond_types)
                initial_guess = [lstsq_guess[2]]
                bounds = [(None,None)]
                min_scale = 0.25
                for count_i,i in enumerate(mode_types):

                    # Assign bounds
                    if count_i == 0:
                        lower_k = max([lstsq_guess[0]*min_scale,10.0])                    
                        if len(i) == 2:
                            bounds += [(lower_k,lstsq_guess[0]),(mode_eq[count_i]-0.2,mode_eq[count_i]+0.2)]
                        elif len(i) == 3:
                            bounds += [(lower_k,lstsq_guess[0]),(mode_eq[count_i]-10.0*np.pi/180.0,mode_eq[count_i]+10.0*np.pi/180.0)]
                        elif len(i) == 4:
                            bounds += [(lower_k,lstsq_guess[0]),(0,None)]
                    else:
                        bounds += [(0,None),(0,None)]

                    # # Assign bounds
                    # if count_i == 0:
                    #     if len(i) == 2:
                    #         bounds += [(50.0,lstsq_guess[0]),(mode_eq[count_i]-0.1,mode_eq[count_i]+0.1)]
                    #     elif len(i) == 3:
                    #         bounds += [(10.0,lstsq_guess[0]),(mode_eq[count_i]-5.0,mode_eq[count_i]+5.0)]
                    #     elif len(i) == 4:
                    #         bounds += [(1.0,lstsq_guess[0]),(0,None)]
                    # else:
                    #     bounds += [(0,None),(0,None)]

                    # Assign initial mode parameter guesses
                    if count_i == 0:
                        initial_guess += list(lstsq_guess[:2])
                    elif len(i) == 2:
                        if i in list(FF_dict["bonds"].keys()):
                            initial_guess += [FF_dict["bonds"][i]["DFT"][1],mode_eq[count_i]]
                        elif i[::-1] in list(FF_dict["bonds"].keys()):
                            initial_guess += [FF_dict["bonds"][i[::-1]]["DFT"][1],mode_eq[count_i]]                        
                        else:
                            initial_guess += [600.0,mode_eq[count_i]]
                    elif len(i) == 3:
                        if i in list(FF_dict["angles"].keys()):
                            initial_guess += [FF_dict["angles"][i]["DFT"][1],mode_eq[count_i]]                  
                        elif i[::-1] in list(FF_dict["angles"].keys()):
                            initial_guess += [FF_dict["angles"][i[::-1]]["DFT"][1],mode_eq[count_i]]
                        else:
                            initial_guess += [60.0,mode_eq[count_i]]
                    elif len(i) == 4:
                        if i in list(FF_dict["dihedrals_harmonic"].keys()):
                            initial_guess += [FF_dict["dihedrals_harmonic"][i]["DFT"][1],0.0]
                        elif i[::-1] in list(FF_dict["dihedrals_harmonic"].keys()):
                            initial_guess += [FF_dict["dihedrals_harmonic"][i[::-1]]["DFT"][1],0.0]
                        else:
                            initial_guess += [5.0,0.0]

                    # # Assign initial mode parameter guesses
                    # if count_i == 0:
                    #     initial_guess += list(lstsq_guess[:2])
                    # elif len(i) == 2:
                    #     if i in FF_dict["bonds"].keys():
                    #         print "READING {} FROM PREVIOUS RUN!".format(i)
                    #         initial_guess += FF_dict["bonds"][i]["DFT"][1:]                        
                    #     elif i[::-1] in FF_dict["bonds"].keys():
                    #         print "READING {} FROM PREVIOUS RUN!".format(i)
                    #         initial_guess += FF_dict["bonds"][i[::-1]]["DFT"][1:]                        
                    #     else:
                    #         initial_guess += [600.0,mode_eq[count_i]]
                    # elif len(i) == 3:
                    #     if i in FF_dict["angles"].keys():
                    #         print "READING {} FROM PREVIOUS RUN!".format(i)
                    #         initial_guess += FF_dict["angles"][i]["DFT"][1:]                  
                    #     elif i[::-1] in FF_dict["angles"].keys():
                    #         print "READING {} FROM PREVIOUS RUN!".format(i)
                    #         initial_guess += FF_dict["angles"][i[::-1]]["DFT"][1:]
                    #     else:
                    #         initial_guess += [60.0,mode_eq[count_i]]
                    # elif len(i) == 4:
                    #     if i in FF_dict["dihedrals_harmonic"].keys():
                    #         initial_guess += [FF_dict["dihedrals_harmonic"][i]["DFT"][1],0.0]
                    #     elif i[::-1] in FF_dict["dihedrals_harmonic"].keys():
                    #         initial_guess += [FF_dict["dihedrals_harmonic"][i[::-1]]["DFT"][1],0.0]
                    #     else:
                    #         initial_guess += [5.0,0.0]

                # Perform a fit of only the scanned mode while keeping all others constant 
                fit_func = lambda x: fit_bonds_angles_sc_single(x,*initial_guess[3:],modes=modes,fit_modes=mode_types,E_0=DFT_E)
                params = minimize(fit_func,initial_guess[:3],method='L-BFGS-B',bounds=[(None,None),bounds[1],bounds[2]],options={'gtol':0.0,'ftol':1.0E-20,'maxiter':1000000,'maxfun':1000000,'maxls':20}).x
                initial_guess[:3] = params

                # Perform a 1-D optimization to find the best force constant for this mode while holding the other parameters constant
                for j in range(1):
                    xhi2_min = None
                    for i in np.arange(initial_guess[1]*min_scale,initial_guess[1]+1.0):
                        initial_guess[1] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    xhi2_min = None
                    for i in np.linspace(best_guess-1.0,best_guess+1.0,300):
                        initial_guess[1] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    initial_guess[1] = best_guess
                    xhi2_min = None
                    for i in np.arange(initial_guess[2]-1.0*np.pi/180.0,initial_guess[2]+1.0*np.pi/180.0,0.01*np.pi/180.0):
                        initial_guess[2] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    xhi2_min = None
                    for i in np.linspace(best_guess-0.01*np.pi/180.0,best_guess+0.01*np.pi/180.0,300):
                        initial_guess[2] = i
                        xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            best_guess = i
                            xhi2_min = xhi2_current
                    initial_guess[2] = best_guess

                # Generate fits with varying initial guess for the force constant of the mode being fit
                xhi2_min = None
                mode_guesses = [initial_guess[1]]
                for z in mode_guesses:

                    # Perform fit
                    initial_guess[1] = z
                    if 'dft' in qc_types:

                        # Perform the fit for the current initial guess
                        fit_func = lambda x: fit_bonds_angles_sc_2(*x,modes=modes,fit_modes=mode_types,E_0=DFT_E)
                        params = minimize(fit_func,initial_guess,method='L-BFGS-B',bounds=bounds,options={'gtol':0.0,'ftol':1.0E-20,'maxiter':1000000,'maxfun':1000000,'maxls':20}).x
                        xhi2_current = fit_bonds_angles_sc_2(*params,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)

                        # Update the fit parameters if a lower energy was obtained
                        if xhi2_min is None or xhi2_current < xhi2_min:
                            DFT_params = deepcopy(params[1:3])                        
                            params[1] = 0.0
                            E_corr = fit_bonds_angles_sc_2(*params,modes=modes,fit_modes=mode_types,penalty=False,E_0=DFT_E)
                            xhi2_min = xhi2_current

                    if 'mp2' in qc_types:
                        fit_func = lambda x: fit_bonds_angles_sc(*x,modes=modes,E_0=MP2_E)
                        params = minimize(fit_func,initial_guess,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':1E-20}).x

                # Fall back on the least-squares guess if the fit was bad
                if xhi2_min > 1E-7:
                    DFT_E = DFT_E + EC_corr + LJ_corr
                    DFT_E = np.array(DFT_E)-min(DFT_E)
                    print("Self-consistent fit failed to converge. This is probably a crowded mode with a lot of electrostatic contributions to the scan potential. Falling back on direct least-squares fitting...")
                    DFT_params = curve_fit(harmonic_norm, thetas, DFT_E,p0=[50.0,thetas[min_ind],0.0],maxfev=1000000)[0][:2]
                    print("xhi2 (kcal^2/mol^2): {}".format(np.mean((harmonic(thetas,DFT_params[0],DFT_params[1])-DFT_E)**(2.0))))
                    fallback_flag = 1

                # Print diagnostic
                else:
                    print("xhi2 (kcal^2/mol^2): {}".format(xhi2_min))

                ####### END OF TYPE BASED FIT COMMANDS #######


        ######################################
        ############# Write Info #############
        ######################################

        # Return to parent/"angles" folder
        os.chdir(working_dir)

        # Open file and write header
        with open(save_folder+'/'+"_".join(angles)+'_info.txt', 'w') as f:
            if 'dft' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol rad^-2\n{:20s} {:< 12.6f} degrees\n".format("k (DFT) =",DFT_params[0],"Theta_0 (DFT) =",DFT_params[1]*180/np.pi))
            if 'mp2' in qc_types:
                f.write("{:20s} {:< 12.6f} kcal/mol rad^-2\n{:20s} {:< 12.6f} degrees\n".format("k (MP2) =",MP2_params[0],"Theta_0 (MP2) =",MP2_params[1]*180/np.pi))
            if 'dft' in qc_types and 'mp2' in qc_types:
                f.write("{:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s} {:<25s}\n".format("Angle (degrees)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)","MP2_Energy","E-E0 (MP2;kcal/mol)","E-E0 (MP2_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["angle"]*180/np.pi,Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["angle"],DFT_params[0],DFT_params[1]),\
                            Geo_dict[str(i)]["MP2"],Geo_dict[str(i)]["MP2"],harmonic(Geo_dict[str(i)]["angle"],MP2_params[0],MP2_params[1])))
            elif 'dft' in qc_types:
                f.write("{:<25s} {:<25s} {:<25s} {:<25s}\n".format("Angle (degrees)","DFT_Energy (kcal/mol)",\
                        "E-E0 (DFT;kcal/mol)","E-E0 (DFT_FIT;kcal/mol)"))
                for i in range(len(list(Geo_dict.keys()))):
                    f.write("{:< 25.6f} {:< 25.6f} {:< 25.6f} {:< 25.6f}\n".format(Geo_dict[str(i)]["angle"]*180/np.pi,Geo_dict[str(i)]["DFT"],\
                            Geo_dict[str(i)]["DFT"],harmonic(Geo_dict[str(i)]["angle"],DFT_params[0],DFT_params[1])))
                    

        # Save plots
        if save_plots is True:

            #### Generate 2D representation of the molecule KNOWN ISSUE WITH PLOTTING FRAGMENT BASED BONDS AND ANGLES, 2D REPRESENTATION IS TAKEN FROM THE MASTER FRAGMENT
            print("Generating 2D representation of the fit fragment...")
            adj_mat = Table_generator(Geo_dict["0"]["elements"],Geo_dict["0"]["geo"])
            geo_2D = kekule(Geo_dict["0"]["elements"],angle_atomtypes[angles_count],Geo_dict["0"]["geo"],adj_mat)

            ######################################
            ######### Generate Fit Plots #########
            ######################################
            print("Rendering fit plots...")

            if 'dft' in qc_types:

                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                if draw_scale > 1.0: draw_scale = 1.0

                # Add bonds
                linewidth = 2.0*draw_scale*5.0
                if linewidth > 2.0:
                    linewidth = 2.0
                if linewidth < 1.0:
                    linewidth = 1.0
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[0]][0],geo_2D[angle_atoms[1]][0]],[geo_2D[angle_atoms[0]][1],geo_2D[angle_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[1]][0],geo_2D[angle_atoms[2]][0]],[geo_2D[angle_atoms[1]][1],geo_2D[angle_atoms[2]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                # Add atom labels
                fontsize = 20*draw_scale*5.0
                if fontsize > 24:
                    fontsize = 24
                if fontsize < 9:
                    fontsize = 9            
                for count_i,i in enumerate(geo_2D):
                    if count_i in angle_atoms:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Add data
                xdata = np.array(thetas)*180.0/np.pi
                harm_range=[0,2.0*np.pi]
                if fallback_flag == 1 or read_flag is True:
                    ydata = harmonic(xdata*np.pi/180.0,DFT_params[0],DFT_params[1]) 
                else:
                    ydata = harmonic(xdata*np.pi/180.0,DFT_params[0],DFT_params[1]) + E_corr + EC_corr + LJ_corr
                min_corr = min(ydata)
                ax2.plot(xdata,ydata-min_corr,'--',linewidth=3,color=(0.0,0.1,0.8))
                if fallback_flag == 1 or read_flag is True:
                    ax2.scatter(np.array(thetas)*180.0/np.pi,DFT_E,s=120,color=(0.0,0.1,0.8))
                else:
                    min_corr = min(DFT_E+EC_corr+LJ_corr)
                    ax2.scatter(np.array(thetas)*180.0/np.pi,DFT_E+EC_corr+LJ_corr-min_corr,s=120,color=(0.0,0.1,0.8))
                if fallback_flag == 1:
                    ax2.text(0.9, 0.8,'*', ha='center', fontsize=40, va='center', transform=ax2.transAxes)
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{\theta \, (degrees)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(thetas)*180.0/np.pi, max(thetas)*180.0/np.pi])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = save_folder+'/'+"_".join(angles)+'-DFT_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(save_folder+'/'+"_".join(angles)+'-DFT_fit.pdf', bbox_inches=0,dpi=300)
                paths_to_angle_plots_DFT+=[Name]
                plt.close(fig)

            if 'mp2' in qc_types:

                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                draw_scale = 1.0
                # Add bonds
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=2.0*draw_scale))
                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[0]][0],geo_2D[angle_atoms[1]][0]],[geo_2D[angle_atoms[0]][1],geo_2D[angle_atoms[1]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[1]][0],geo_2D[angle_atoms[2]][0]],[geo_2D[angle_atoms[1]][1],geo_2D[angle_atoms[2]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))

                # Add atom labels
                for count_i,i in enumerate(geo_2D):
                    if count_i in angle_atoms:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1], Geo_dict["0"]["elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Add data
                xdata = np.array(thetas)*180.0/np.pi
                ydata = harmonic(xdata*np.pi/180.0,MP2_params[0],MP2_params[1]) + E_corr
                ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
                ax2.scatter(np.array(thetas)*180.0/np.pi,MP2_E,s=120,color=(0.0,0.1,0.8))
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{\theta \, (degrees)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(thetas)*180.0/np.pi, max(thetas)*180.0/np.pi])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = save_folder+'/'+"_".join(angles)+'-MP2_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(save_folder+'/'+"_".join(angles)+'-MP2_fit.pdf', bbox_inches=0,dpi=300)
                paths_to_angle_plots_MP2+=[Name]
                plt.close(fig)
        
        #############################################################
        ####### Update FF_dict with DFT and MP2 FF-parameters #######
        #############################################################
        # NOTE: the force constant is in units of kcal/mol/rad^2 but the equilibrium angle is in units of degrees (lammps harmonic angle convention). 
        FF_dict["angles"][angles]={}
        if 'dft' in qc_types:
            FF_dict["angles"][angles]["DFT"] = ["harmonic"] + [ DFT_params[0], DFT_params[1]*180.0/np.pi ]
        if 'mp2' in qc_types:
            FF_dict["angles"][angles]["MP2"] = ["harmonic"] + [ MP2_params[0], MP2_params[1]*180.0/np.pi ]

    # Return to parent folder after processing all angles
    os.chdir(working_dir)

    return paths_to_angle_plots_DFT, paths_to_angle_plots_MP2

# This is a messy wrapper for the dihedral fitting procedure. In brief, the 
# function looks up all the dihedral data, calculates harmonic fits, and
# generates some plots and data text files. All parameters are saved to FF_dict
def parse_harmonic_dihedrals_sc(base_name,list_of_dihedral_types,dihedral_folders,dihedral_folders_frag,dihedral_atomtypes,charge_dict,modes_from_FF,eq_charges,save_folder,fit_type='AA',delta_xhi2_thresh=0.000001,min_cycles=10,\
                                max_cycles=100000,qc_types=[],corr_list=[],find_restarts=0,second_d_thresh=500.0,charge_origin=None,save_plots=True,gens=2):
    
    # In python all global variables must be declared in each scope
    global FF_dict,QC_type,Dihedral_Harmonic_Data,dihedral_scan,harm_range

    # Initialize dictionary for mapping atomic number to element
    atom_to_element = { 1:'H' ,  2:'He',\
                        3:'Li',  4:'Be',  5:'B' ,  6:'C' ,  7:'N' ,  8:'O' ,  9:'F' , 10:'Ne',\
                       11:'Na', 12:'Mg', 13:'Al', 14:'Si', 15:'P' , 16:'S' , 17:'Cl', 18:'Ar',\
                       19:'K' , 20:'Ca', 21:'Sc', 22:'Ti', 23:'V' , 24:'Cr', 25:'Mn', 26:'Fe', 27:'Co', 28:'Ni', 29:'Cu', 30:'Zn', 31:'Ga', 32:'Ge', 33:'As', 34:'Se', 35:'Br', 36:'Kr',\
                       37:'Rb', 38:'Sr', 39:'Y' , 40:'Zr', 41:'Nb', 42:'Mo', 43:'Tc', 44:'Ru', 45:'Rh', 46:'Pd', 47:'Ag', 48:'Cd', 49:'In', 50:'Sn', 51:'Sb', 52:'Te', 53:'I' , 54:'Xe',\
                       55:'Cs', 56:'Ba', 57:'La', 72:'Hf', 73:'Ta', 74:'W' , 75:'Re', 76:'Os', 77:'Ir', 78:'Pt', 79:'Au', 80:'Hg', 81:'Tl', 82:'Pb', 83:'Bi', 84:'Po', 85:'At', 86:'Rn'}    

    # Save parent directory
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_dihedral_fitpot_plots_DFT=[]
    paths_to_dihedral_fitpot_plots_MP2=[]
    paths_to_dihedral_totpot_plots_DFT=[]
    paths_to_dihedral_totpot_plots_MP2=[]

    # List of Processed Dihedrals
    processed_dihedrals = []

    # Grab the dihedral data (returns a dictionary holding dihedral_type,angle tuples for each configuration and the DFT and MP2 energies
    Dihedral_Harmonic_Data = get_harmonic_dihedral_data(list_of_dihedral_types,dihedral_folders,fit_type,dihedral_atomtypes,charge_dict,qc_types=qc_types,corr_list=[])
    keys = list(Dihedral_Harmonic_Data.keys())

    # Return to working directory and save the list of catenated geometries
    for i in list(Dihedral_Harmonic_Data.keys()):

        # Check if the MD subfolder needs to be created
        if fit_type == 'UA': savename = base_name+"/dihedrals_harmonic/"+"_".join([ k+"-UA" for k in i ])+"_geos.xyz"
        else: savename = base_name+"/dihedrals_harmonic/"+"_".join([ k for k in i ])+"_geos.xyz"
        f = open(savename, 'w')
    
        # Loop over all discovered geometries and save their geometries to Geo_dict and
        # catenate to a single file for viewing.
        for j in natural_sort(list(Dihedral_Harmonic_Data[i].keys())):
            f.write("{:d}\n\n".format(len(Dihedral_Harmonic_Data[i][j]["Geo"])))
            for count_k,k in enumerate(Dihedral_Harmonic_Data[i][j]["Geo"]):
                f.write(" {:<10s} {:< 15.8f} {:< 15.8f} {:< 15.8f}\n".format(Dihedral_Harmonic_Data[i][j]["Elements"][count_k],k[0],k[1],k[2]))
        f.close()

    # dihedral types are added to already_fit_list as they are fit. Fitting is skipped for any dihedrals in this list, in order
    # to avoid redundant fits when the same dihedral happens to be in two scans
    already_fit_list = modes_from_FF

    # Iterate over each scan's data
    for dihedral_count,dihedrals in enumerate(list_of_dihedral_types):

        # Print banner
        if fit_type == "AA": print("\n{}\n* {:^163s} *\n{}".format("*"*167,"Parsing AA-harmonic_dihedral Types from Scan {}".format(dihedrals),"*"*167))
        if fit_type == "UA": print("\n{}\n* {:^163s} *\n{}".format("*"*167,"Parsing UA-harmonic_dihedral Types from Scan {}".format(dihedrals),"*"*167))

        # Change into current bond directory
        os.chdir(dihedral_folders[dihedral_count])

        # Generate adjacency mat
        adj_mat = Table_generator(Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"],Dihedral_Harmonic_Data[dihedrals]["0"]["Geo"])        

        # ID atomtypes
        atom_types = id_types(Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"],adj_mat,gens,geo=Dihedral_Harmonic_Data[dihedrals]["0"]["Geo"])
        
        # Find the mode types in this fragment
        Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(adj_mat,atom_types,return_all=1)

        # Parse dihedral atom indices
        dihedral_atoms = Dihedral_Harmonic_Data[dihedrals]["0"]["Dihedral_Atoms"][0][0]

        # Determine bonding matrix for the compound        
        frag_info = dihedral_folders_frag[dihedral_count]
        qc_charge_key = dihedral_folders[dihedral_count].split('/')[0]+'/geoopt'
        for k in eq_charges:
            if(fnmatch.fnmatch(k,'*'+frag_info+'*')):
               qc_charge_key = k 
        q_tmp   = int(round(sum(eq_charges[qc_charge_key])))
        bond_mat = find_lewis(atom_types, adj_mat, q_tot=q_tmp, b_mat_only=True,verbose=False)

        # Check if the reordering of angle_atoms is required
        if tuple(dihedral_atoms) not in Dihedrals:
            dihedral_atoms = dihedral_atoms[::-1]
        if tuple(dihedral_atoms) not in Dihedrals:
            print("ERROR in parse_harmonic_dihedrals: the fit dihedral was not discovered in the quantum chemistry geometry. Exiting...")
            quit()        

        # Parse modes in each configuration
        # NOTE: the active mode is removed from the list of angles, but reincluded as the first entry in the iteration list (guarrantees first position for downstream processing)
        modes = []
        Dihedrals.remove(tuple(dihedral_atoms))
        for i in natural_sort(list(Dihedral_Harmonic_Data[dihedrals].keys())):
            tmp = []

            # Dihedrals are handled separately because the indexing of Dihedrals needs to be matched with Dihedral_types
            for count_j,j in enumerate([tuple(dihedral_atoms)]+Dihedrals):                

                # Only keep modes that involve the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                if   len(j) == 4 and len(set(dihedral_atoms[1:3]).intersection(set(j[1:3]))) < 2: continue
                elif len(j) == 3 and len(set(dihedral_atoms[1:3]).intersection(set(j))) < 2: continue
                elif len(j) == 2 and len(set(dihedral_atoms[1:3]).intersection(set(j))) < 1: continue

                # Parse harmonic dihedral degree (NEW criteria based on if a double bond exists in any of the resonance structures)
                # ( NEW CHECK USES BOND_MAT) If j is a harmonic dihedral then redundancy is assessed on the basis of all four mode 
                # indices matching (this is necessary beacuse of the way that harmonic dihedrals are scanned)
#                elif 2 in [ int(k[dihedral_atoms[1],dihedral_atoms[2]]) for k in bond_mat ]:
                if 2 in [ int(k[j[1],j[2]]) for k in bond_mat ]:
                    v1 = (Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[1]] - Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[0]])
                    v2 = (Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[2]] - Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[1]])
                    v3 = (Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[3]] - Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[2]])
                    angle = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )
                    if angle >  np.pi*3.0/2.0: angle = angle - 2.0*np.pi                        
                    if angle < -np.pi/2.0: angle = angle + 2.0*np.pi
                    tmp += [((atom_types[j[0]],atom_types[j[1]],atom_types[j[2]],atom_types[j[3]]),angle)]
                else:
                    print("ERROR in parse_harmonic_dihedral_sc: Harmonic dihedral {} could not be read. Exiting...".format([ atom_types[_] for _ in j ]))
                
            for count_j,j in enumerate(Bonds + Angles):

                # Only keep modes that involve one of the scanned bond atoms (for dihedrals one of the scanned bond atoms must participate)
                if   len(j) == 4 and len(set(dihedral_atoms[1:3]).intersection(set(j[1:3]))) == 0: continue
                elif len(set(dihedral_atoms[1:3]).intersection(set(j))) == 0: continue

                # Parse bond length
                if len(j) == 2:
                    length = norm(Dihedral_Harmonic_Data[dihedrals][i]['Geo'][j[0]]-Dihedral_Harmonic_Data[dihedrals][i]['Geo'][j[1]])
                    tmp += [((atom_types[j[0]],atom_types[j[1]]),length)]

                # Parse angle degree:
                if len(j) == 3:
                    atom_1 = Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[0]]
                    atom_2 = Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[1]]
                    atom_3 = Dihedral_Harmonic_Data[dihedrals][i]["Geo"][j[2]]
                    theta =  np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))                    
                    tmp += [((atom_types[j[0]],atom_types[j[1]],atom_types[j[2]]),theta)]

            # Append the mode information for the current configuration
            modes += [tmp]

        # Append the scanned mode angles to Geo_dict (this is a holdover from the way this function was originally written, with Geo_dict containing the scanned angles
        for count_i,i in enumerate(modes):
            Dihedral_Harmonic_Data[dihedrals][str(count_i)]["angle"] = i[0][1]

        ######################################
        ###### Calculate Fit Parameters ######
        ######################################

        # Collect list of energies and angles
        DFT_E  = []
        MP2_E  = []        
        thetas = []
        for i in range(len(list(Dihedral_Harmonic_Data[dihedrals].keys()))):
            thetas += [modes[i][0][1]]
            if 'dft' in qc_types:
                DFT_E += [Dihedral_Harmonic_Data[dihedrals][str(i)]["DFT"]]
            if 'mp2' in qc_types:
                MP2_E += [Dihedral_Harmonic_Data[dihedrals][str(i)]["MP2"]]

        # print diagnostic
        #print "\n".join([ "{:< 12.6f} {:< 12.6}".format(_ * 180.0/np.pi, DFT_E[count]) for count,_ in enumerate(thetas) ])
        
        # Update the charges from the force-field if all of them are available.
        frag_info = dihedral_folders_frag[dihedral_count]
        qc_charge_key = dihedral_folders[dihedral_count].split('/')[0]+'/geoopt'
        for k in eq_charges:
            if(fnmatch.fnmatch(k,'*'+frag_info+'*')):
               qc_charge_key = k 
        fit_charges   = deepcopy(eq_charges[qc_charge_key])
        if charge_origin is not None:
            if True not in [ charge_origin[i] == "eq_config" for i in atom_types ]:
                for i in range(len(fit_charges)): fit_charges[i] = FF_dict["charges"][atom_types[i]]
            else:
                print("WARNING: Missing partial charges for some or all of the atomtypes in this fragment, defaulting to equilibrium partial charges.") 
        elif False not in [ atom_types[j] in list(FF_dict["charges"].keys()) for j in range(len(fit_charges)) ]:
            for i in range(len(fit_charges)): fit_charges[i] = FF_dict["charges"][atom_types[i]]

        # Calculate electrostatic and vdw energy corrections in each configuration
        fit_LJ        = initialize_VDW(atom_types,[],verbose=False)  
        EC_corr = []
        LJ_corr = []

        # Reassemble modes (since an angle was removed) and calculate the electrostatic (EC) and lennard-jones (LJ) contributions to each configuration
        Bonds,Angles,Dihedrals,One_fives = Find_modes(adj_mat,atom_types,return_all=0)
        for i in range(len(list(Dihedral_Harmonic_Data[dihedrals].keys()))):            
            EC_corr += [E_coul(Dihedral_Harmonic_Data[dihedrals][str(i)]["Geo"],adj_mat,atom_types,fit_charges,[Bonds,Angles,Dihedrals],(0.0,0.0,0.0),"AA")]
            LJ_corr += [E_LJ(Dihedral_Harmonic_Data[dihedrals][str(i)]["Geo"],adj_mat,atom_types,fit_charges,[Bonds,Angles,Dihedrals],(0.0,0.0,0.0),"AA",LJ_dict=fit_LJ)]

        # Sort by angles
        if 'dft' in qc_types and 'mp2' in qc_types:
            together = sorted(zip(thetas,DFT_E,MP2_E,EC_corr,LJ_corr,modes))
            thetas   = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            MP2_E    = [ i[2] for i in together ]            
            EC_corr  = [ i[3] for i in together ]            
            LJ_corr  = [ i[4] for i in together ]            
            modes    = [ i[5] for i in together ]
        elif 'dft' in qc_types:
            together = sorted(zip(thetas,DFT_E,EC_corr,LJ_corr,modes))
            thetas   = [ i[0] for i in together ]
            DFT_E    = [ i[1] for i in together ]
            EC_corr  = [ i[2] for i in together ]
            LJ_corr  = [ i[3] for i in together ]
            modes    = [ i[4] for i in together ]
        elif 'mp2' in qc_types:
            together = sorted(zip(thetas,MP2_E,EC_corr,LJ_corr,modes))
            thetas   = [ i[0] for i in together ]
            MP2_E    = [ i[1] for i in together ]
            EC_corr  = [ i[2] for i in together ]
            LJ_corr  = [ i[3] for i in together ]
            modes    = [ i[4] for i in together ]

        # Create a duplicate of these lists in case a reversion is required after extensive pruning and job regeneration
        thetas_bkp   = deepcopy(thetas)
        DFT_E_bkp    = deepcopy(DFT_E)
        MP2_E_bkp    = deepcopy(MP2_E)
        EC_corr_bkp  = deepcopy(EC_corr)
        LJ_corr_bkp  = deepcopy(LJ_corr)
        modes_bkp    = deepcopy(modes)        

        # If the minimum is near the edge of the scanned range then automatically perform a reoptimization
        min_ind = DFT_E.index(min(DFT_E)) 
        resub_flag = 0
        if min_ind in [0,1,len(DFT_E)-1,len(DFT_E)-2]:
            print("Minimum lies near the edge of the scanned range...")
            resub_flag = 1
        
        # Remove edge effects (i.e. edges with the improper derivative)
        # The value of correction_flag is used to toggle a diagnostic print statement
        delta = np.abs(thetas[1]-thetas[0])*180.0/np.pi
        start_ind = 0
        end_ind = len(thetas) - 1
        correction_flag = 0
        if resub_flag == 0:
            for count_i,i in enumerate(thetas[:-1]):
                if DFT_E[count_i+1] - DFT_E[count_i] < 0: 
                    start_ind = count_i
                    break
            if start_ind != 0: correction_flag = 1

            for i in range(len(thetas))[::-1]:
                if DFT_E[i]-DFT_E[i-1] > 0: 
                    end_ind = i
                    break
            if end_ind != len(thetas)-1: correction_flag = 1
        
        # Assembled edited lists and print diagnostic if trimming occured
        if correction_flag == 1:
            print("Data points with improper derivatives were removed to increase the quality of fit...")
        if 'dft' in qc_types: DFT_E = DFT_E[start_ind:end_ind+1]
        if 'mp2' in qc_types: MP2_E = MP2_E[start_ind:end_ind+1]
        thetas = thetas[start_ind:end_ind+1]
        modes  = modes[start_ind:end_ind+1]
        EC_corr = EC_corr[start_ind:end_ind+1]
        LJ_corr = LJ_corr[start_ind:end_ind+1]

        # Remove derivative discontinuities
        # NOTE: the threshold of 500.0 kcal^2/rad is empirical, it was found to consistently remove most derivative 
        #       discontinuities without removing any continuous data.
        second_d_thresh = float(second_d_thresh)

        # Remove datapoints until derivative discontinuities have been removed
        correction_flag = 0
        clean_flag = 0
        while clean_flag == 0 and len(DFT_E) > 2:

            # Loop over the energies and calculate the second derivatives. 
            # NOTE: the if/else construction is necessary to deal with the endpoints
            clean_flag = 1
            del_list = []            
            for i in range(len(DFT_E))[1:-1]:                

                second_d = ( ( DFT_E[i+1] - DFT_E[i] ) / ( thetas[i+1] - thetas[i] ) - ( DFT_E[i] - DFT_E[i-1] ) / ( thetas[i] - thetas[i-1] ) ) / (thetas[i] - thetas[i-1] ) 
                if second_d < 0.0 or np.abs(second_d) > second_d_thresh:
                    correction_flag = 1
                    clean_flag = 0
                    del_list += [i]
                    break
                
            DFT_E = [ i for count_i,i in enumerate(DFT_E) if count_i not in del_list ]
            thetas = [ i for count_i,i in enumerate(thetas) if count_i not in del_list ]
            modes  = [ i for count_i,i in enumerate(modes) if count_i not in del_list ]
            EC_corr = [ i for count_i,i in enumerate(EC_corr) if count_i not in del_list ]
            LJ_corr = [ i for count_i,i in enumerate(LJ_corr) if count_i not in del_list ]

        if correction_flag == 1:
            print("Data points with positive second derivatives or discontinuities were removed to increase the quality of the fit...")

        # If more than a third of the data has been discarded then the script generates a reoptimization job using the lowest
        # energy conformation from the scan.
#        if ( len(thetas) <= len(Dihedral_Harmonic_Data[dihedrals].keys())*4.0/5.0 or resub_flag == 1 ) and find_restarts == 1:
        if ( len(thetas) <= len(list(Dihedral_Harmonic_Data[dihedrals].keys()))*1.0/2.0 or resub_flag == 1 ) and find_restarts == 1:
            min_folder = [ Dihedral_Harmonic_Data[dihedrals][str(i)]["DFT"] for i in range(len(list(Dihedral_Harmonic_Data[dihedrals].keys()))) ].index(min([ Dihedral_Harmonic_Data[dihedrals][str(j)]["DFT"] \
                           for j in range(len(list(Dihedral_Harmonic_Data[dihedrals].keys()))) ]))
            template_file = str(min_folder) + '/' + dihedral_folders[dihedral_count].split('/')[0] + '_' + '-'.join(dihedral_folders[dihedral_count].split('/')[2].split('_')) + '_' + str(min_folder) + '.in'
            if os.path.isdir("REDO"):
                print("Input files for resubmission have already been generated for this mode...")
            else:
                print("Generating the input files for a reoptimization based on the lowest energy ({}) configuration...".format(min_folder))
                delta = 0.1 # THIS IS A KLUDGE, FOR SOME REASON AFTER MULTIPLE RESTARTS THE DELTA COLLAPSES (I THINK DUE TO THE IMPERFECT RESTRAINT OF THE DIHEDRALS IN THE QC)
                gen_reoptimization(template=template_file,geo=Dihedral_Harmonic_Data[dihedrals][str(min_folder)]['Geo'],\
                                   elements=Dihedral_Harmonic_Data[dihedrals][str(min_folder)]["Elements"],delta=delta,N_steps=len(list(Dihedral_Harmonic_Data[dihedrals].keys()))) # NEW 
            #os.chdir(working_dir)
            #continue
        
            # Revert to the original modes so that the user can assess the data and/or use the best fit as is. 
            thetas  = thetas_bkp 
            DFT_E   = DFT_E_bkp
            MP2_E   = MP2_E_bkp   
            EC_corr = EC_corr_bkp 
            LJ_corr = LJ_corr_bkp 
            modes   = modes_bkp   

        # If the program is only searching for restarts then the fit is avoided
        if find_restarts == 1:
            os.chdir(working_dir)
            continue

        # Subtract the EC and LJ corrections from the DFT fit 
        EC_corr = np.array(EC_corr) 
        LJ_corr = np.array(LJ_corr) 
        DFT_E   = np.array(DFT_E) - EC_corr - LJ_corr

        # Normalize energies by the minimum value
        if 'dft' in qc_types: DFT_E = np.array(DFT_E)-min(DFT_E)
        if 'mp2' in qc_types: MP2_E = np.array(MP2_E)-min(MP2_E)

        # Print diagnostic
        print("constrained atoms: {}".format(" ".join([ str(i) for i in dihedral_atoms])))
        print("theta range: {: <6.4f}-{: <6.4f}".format(min(np.array(thetas))*180.0/np.pi,max(np.array(thetas))*180.0/np.pi))
        print("Calculating harmonic fits...")

        # Use read parameters if available (params holds force constant and theta_0)
        read_flag = False
        if dihedrals in modes_from_FF:
            print("Using parameters for this dihedral from the supplied FF-file(s)...")
            if "dft" in qc_types:
                read_flag = True
                DFT_params = np.array(FF_dict["dihedrals_harmonic"][dihedrals]["DFT"][1:])
                DFT_params[1] = DFT_params[1]*np.pi/180.0

                # Initialize variables necessary to make the plots
                thetas = np.array(thetas)
                offset  = np.pi - thetas[min_ind] 
                fallback_flag = 0
                
            if "mp2" in qc_types:
                MP2_params = np.array(FF_dict["dihedrals_harmonic"][dihedrals]["MP2"][1:])
                MP2_params[1] = MP2_params[1]*np.pi/180.0

        # Calculate parameters (params holds force constant and theta_0)    
        # NOTE: the force constant is in units of kcal/mol/rad^2 but the equilibrium angle is in units of degrees 
        #       (lammps harmonic angle convention). The eq_angle is converted at the end when it is saved to FF_dict
        # NOTE: the initial guess for the mode being fit is set to the direct fitting parameters
        else:

#             ####### START OF UNCONSTRAINED BASED FIT COMMANDS #######
            
#             # Get initial guess
#             harm_range=[0.0,180.0]
#             # harm_range=[min(thetas),max(thetas)]
#             min_ind = next( count_i for count_i,i in enumerate(DFT_E) if i == min(DFT_E) )
#             if len(modes[min_ind][0][0]) == 2:
#                 lstsq_guess,errs = curve_fit(harmonic_norm, lengths, DFT_E,p0=[100.0,lengths[min_ind],0.0],maxfev=10000)
#             else:
#                 lstsq_guess,errs = curve_fit(harmonic_norm, thetas, DFT_E,p0=[50.0,thetas[min_ind],0.0],maxfev=10000)

#             print "lstsq_guess: {}".format(lstsq_guess)
#             # Generate initial guess parameters and fit bounds
#             # NOTE: the normalization value is seeded with 0.0 and is fit without bounds
#             # NOTE: the force constant for the scanned mode is bound from above by the lstsq fit value
#             # NOTE: the approximate force constants for all unscanned bonds are derived from UFF
#             # NOTE: the force constants for all unscanned angles are approximated as 100 kcal/mol as a reasonable initial guess
#             bond_dict = UFF_bonds(Bond_types)
#             initial_guess = [lstsq_guess[2]]
#             bounds = [(None,None)]
#             for count_i,i in enumerate(modes[min_ind]):
# #                bounds += [(10.0,lstsq_guess[0]),(0,None)]
#                 bounds += [(0,None),(0,None)]
#                 if count_i == 0:
#                     initial_guess += list(lstsq_guess[:2])
#                 elif len(i[0]) == 2:
#                     initial_guess += [bond_dict[i[0]],i[1]]
#                 elif len(i[0]) == 3:
#                     initial_guess += [100.0,i[1]]

#             # Perform a 1-D optimization to find the best force constant for this mode while holding the other parameters constant
#             xhi2_min = None
#             for i in np.arange(10.0,lstsq_guess[0]+100.0):
#                 initial_guess[1] = i
#                 xhi2_current = fit_bonds_angles_sc(*initial_guess,modes=modes,penalty=True,E_0=DFT_E)
#                 if xhi2_min is None or xhi2_current < xhi2_min:
#                     best_guess = i
#                     xhi2_min = xhi2_current
#             xhi2_min = None
#             for i in np.linspace(best_guess-1.0,best_guess+1.0,300):
#                 initial_guess[1] = i
#                 xhi2_current = fit_bonds_angles_sc(*initial_guess,modes=modes,penalty=True,E_0=DFT_E)
#                 if xhi2_min is None or xhi2_current < xhi2_min:
#                     best_guess = i
#                     xhi2_min = xhi2_current

#             # Generate fits with varying initial guess for the force constant of the mode being fit
#             print "initial guess err: {}".format(xhi2_min)
#             xhi2_min = None
# #            mode_guesses = list(np.linspace(10.0,lstsq_guess[0],5))
# #            mode_guesses = mode_guesses + [mode_guesses[-1]+mode_guesses[1]]
#             mode_guesses = [best_guess]
#             print "guesses: {}".format(mode_guesses)
#             for z in mode_guesses:

#                 # Perform fit
#                 initial_guess[1] = z
#                 if 'dft' in qc_types:

#                     # Perform the fit for the current initial guess
#                     print "initial params: {} ({})".format(initial_guess[:3],fit_bonds_angles_sc(*initial_guess,modes=modes,penalty=True,E_0=DFT_E))
#                     fit_func = lambda x: fit_bonds_angles_sc(*x,modes=modes,E_0=DFT_E)
#                     params = minimize(fit_func,initial_guess,method='L-BFGS-B',bounds=bounds,options={'gtol':0.0,'ftol':0.0,'maxiter':1000000,'maxfun':1000000,'maxls':100}).x
#                     xhi2_current = fit_bonds_angles_sc(*params,modes=modes,penalty=True,E_0=DFT_E)
#                     print "params: {} ({})".format(params[:3],xhi2_current)
#                     # Update the fit parameters if a lower energy was obtained
#                     if xhi2_min is None or xhi2_current < xhi2_min:
#                         print "{}".format(params[:3])
#                         DFT_params = deepcopy(params[1:3])                        
#                         params[1] = 0.0
#                         E_corr = fit_bonds_angles_sc(*params,modes=modes,penalty=False,E_0=DFT_E)
#                         xhi2_min = xhi2_current

#                 if 'mp2' in qc_types:
#                     fit_func = lambda x: fit_bonds_angles_sc(*x,modes=modes,E_0=MP2_E)
#                     params = minimize(fit_func,initial_guess,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':1E-20}).x

#             # Print diagnostic
#             print "xhi2 (kcal^2/mol^2): {}".format(xhi2_min)

#             ####### END OF UNCONSTRAINED BASED FIT COMMANDS #######

            ####### START OF TYPE BASED FIT COMMANDS #######

            # Center at 180 degrees before fitting
            # Get initial guess
            min_ind = next( count_i for count_i,i in enumerate(DFT_E) if i == min(DFT_E) )
            thetas = np.array(thetas)
            offset  = np.pi - thetas[min_ind] # NEW XXX
            thetas  = thetas + offset  # NEW XXX
            if len(modes[min_ind][0][0]) == 2:
                lstsq_guess,errs = curve_fit(harmonic_norm, lengths, DFT_E,p0=[100.0,lengths[min_ind],0.0],maxfev=1000000)
            elif len(modes[min_ind][0][0]) == 3:
                lstsq_guess,errs = curve_fit(harmonic_norm, thetas, DFT_E,p0=[50.0,thetas[min_ind],0.0],maxfev=1000000)                
            elif len(modes[min_ind][0][0]) == 4:
                fit_func = lambda x,y,z: dihedral_harmonic_norm(x,y,d=-1.0,n=2.0,norm=z)
                lstsq_guess,errs = curve_fit(fit_func, thetas, DFT_E,p0=[10.0,0.0],maxfev=10000)                

            # Apply minimum condition to the fit
            if lstsq_guess[0] < 1.0: lstsq_guess[0] = 1.0
            if lstsq_guess[0] > 30: lstsq_guess[0] = 30.0

            # Collect the unique mode types in the current fit
            mode_types = []
            mode_eq    = []
            for i in modes[min_ind]:
                mode_types += [i[0]]
                mode_eq += [i[1]]

            # Generate initial guess parameters and fit bounds
            # NOTE: the normalization value is seeded with 0.0 and is fit without bounds
            # NOTE: the force constant for the scanned mode is bound from above by the lstsq fit value
            # NOTE: the force constants for all unscanned bonds approximated as 100 kcal/mol as a reasonable initial guess
            # NOTE: the force constants for all unscanned angles are approximated as 100 kcal/mol as a reasonable initial guess
            # NOTE: the force constants for all harmonic dihedrals are approximated as 0.0 kcal/mol as a reasonable initial guess
            # NOTE: for harmonic dihedrals the equilibrium position isn't used so it is set to a dummy value of 0.0
            bond_dict = UFF_bonds(Bond_types)
            initial_guess = [lstsq_guess[1]]
            bounds = [(None,None)]
            min_scale = 0.25
            for count_i,i in enumerate(mode_types):

                # Assign bounds
                if count_i == 0:
                    lower_k = max([lstsq_guess[0]*min_scale,1.0])                    
                    if len(i) == 2:
                        bounds += [(lower_k,lstsq_guess[0]),(mode_eq[count_i]-0.1,mode_eq[count_i]+0.1)]
                    elif len(i) == 3:
                        bounds += [(lower_k,lstsq_guess[0]),(mode_eq[count_i]-5.0*np.pi/180.0,mode_eq[count_i]+5.0*np.pi/180.0)]
                    elif len(i) == 4:
                        bounds += [(lower_k,lstsq_guess[0]),(0,None)]
                else:
                    bounds += [(0,None),(0,None)]

                # Assign initial mode parameter guesses
                if count_i == 0:
                    initial_guess += [lstsq_guess[0],0.0]
                elif len(i) == 2:
                    if i in list(FF_dict["bonds"].keys()):
                        initial_guess += [FF_dict["bonds"][i]["DFT"][1],mode_eq[count_i]]
                    elif i[::-1] in list(FF_dict["bonds"].keys()):
                        initial_guess += [FF_dict["bonds"][i[::-1]]["DFT"][1],mode_eq[count_i]]                        
                    else:
                        initial_guess += [600.0,mode_eq[count_i]]
                elif len(i) == 3:
                    if i in list(FF_dict["angles"].keys()):
                        initial_guess += [FF_dict["angles"][i]["DFT"][1],mode_eq[count_i]]                  
                    elif i[::-1] in list(FF_dict["angles"].keys()):
                        initial_guess += [FF_dict["angles"][i[::-1]]["DFT"][1],mode_eq[count_i]]
                    else:
                        initial_guess += [60.0,mode_eq[count_i]]
                elif len(i) == 4:
                    if i in list(FF_dict["dihedrals_harmonic"].keys()):
                        initial_guess += [FF_dict["dihedrals_harmonic"][i]["DFT"][1],0.0]
                    elif i[::-1] in list(FF_dict["dihedrals_harmonic"].keys()):
                        initial_guess += [FF_dict["dihedrals_harmonic"][i[::-1]]["DFT"][1],0.0]
                    else:
                        initial_guess += [5.0,0.0]

            # Perform a 1-D optimization to find the best force constant for this mode while holding the other parameters constant
            xhi2_min = None
            for i in np.arange(lstsq_guess[0]*min_scale,lstsq_guess[0]+1.0):
                initial_guess[1] = i
                xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                if xhi2_min is None or xhi2_current < xhi2_min:
                    best_guess = i
                    xhi2_min = xhi2_current
            xhi2_min = None
            for i in np.linspace(best_guess-1.0,best_guess+1.0,300):
                initial_guess[1] = i
                xhi2_current = fit_bonds_angles_sc_2(*initial_guess,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)
                if xhi2_min is None or xhi2_current < xhi2_min:
                    best_guess = i
                    xhi2_min = xhi2_current

            # Generate fits with varying initial guess for the force constant of the mode being fit
            xhi2_min = None
            mode_guesses = [best_guess]
            for z in mode_guesses:

                # Perform fit
                initial_guess[1] = z
                if 'dft' in qc_types:

                    # Perform the fit for the current initial guess
                    fit_func = lambda x: fit_bonds_angles_sc_2(*x,modes=modes,fit_modes=mode_types,E_0=DFT_E)
                    params = minimize(fit_func,initial_guess,method='L-BFGS-B',bounds=bounds,options={'gtol':0.0,'ftol':1.0E-20,'maxiter':1000000,'maxfun':1000000,'maxls':20}).x
                    xhi2_current = fit_bonds_angles_sc_2(*params,modes=modes,fit_modes=mode_types,penalty=True,E_0=DFT_E)

                    # Update the fit parameters if a lower energy was obtained
                    if xhi2_min is None or xhi2_current < xhi2_min:
                        DFT_params = deepcopy(params[1:3])                        
                        params[1] = 0.0
                        E_corr = fit_bonds_angles_sc_2(*params,modes=modes,fit_modes=mode_types,penalty=False,E_0=DFT_E)
                        xhi2_min = xhi2_current

                if 'mp2' in qc_types:
                    fit_func = lambda x: fit_bonds_angles_sc(*x,modes=modes,E_0=MP2_E)
                    params = minimize(fit_func,initial_guess,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':1E-20}).x

            # Fall back on the least-squares guess if the fit was bad
            fallback_flag = 0
            if xhi2_min > 1E-4:

                # Add corrections to QC energies
                print("Self-consistent fit failed to converge. This is probably a crowded mode with a lot of electrostatic contributions to the scan potential. Falling back on direct least-squares fitting...")
                DFT_E = DFT_E + EC_corr + LJ_corr
                DFT_E = np.array(DFT_E)-min(DFT_E)

                # Adjust thetas to align with 180
                min_ind = next( count_i for count_i,i in enumerate(DFT_E) if i == min(DFT_E) )
                offset  = np.pi - thetas[min_ind] # NEW XXX
                thetas  = thetas + offset  # NEW XXX
                fit_func = lambda x,y,z: dihedral_harmonic_norm(x,y,d=-1.0,n=2.0,norm=z)
                DFT_params = curve_fit(fit_func, thetas, DFT_E,p0=[0.0,0.0],maxfev=10000)[0][0:2]

                # Apply minimum condition to the fit
                if DFT_params[0] < 1.0: DFT_params[0] = 1.0
                if DFT_params[0] > 30: DFT_params[0] = 30.0

                print("xhi2 (kcal^2/mol^2): {}".format(np.mean((fit_func(thetas,DFT_params[0],DFT_params[1])-DFT_E)**(2.0))))
                fallback_flag = 1
                
            # Print diagnostic
            else:
                print("xhi2 (kcal^2/mol^2): {}".format(xhi2_min))

            ####### END OF TYPE BASED FIT COMMANDS #######

        # Return to parent/dihedrals_harmonic"" folder
        os.chdir(working_dir)

        # Save plots
        if save_plots is True:
        
            #### Generate 2D representation of the molecule
            print("Generating 2D representation of the fit fragment...")
            atomtypes = Dihedral_Harmonic_Data[dihedrals]["0"]["Atomtypes"]
            adj_mat = Table_generator(Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"],Dihedral_Harmonic_Data[dihedrals]["0"]["Geo"])
            geo_2D = kekule(Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"],atomtypes,Dihedral_Harmonic_Data[dihedrals]["0"]["Geo"],adj_mat)

            ######################################
            ######### Generate Fit Plots #########
            ######################################
            print("Rendering fit plots...")

            if 'dft' in qc_types:

                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                if draw_scale > 1.0: draw_scale = 1.0

                # Add bonds
                linewidth = 2.0*draw_scale*5.0
                if linewidth > 2.0:
                    linewidth = 2.0
                if linewidth < 1.0:
                    linewidth = 1.0
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[0]][0],geo_2D[dihedral_atoms[1]][0]],[geo_2D[dihedral_atoms[0]][1],geo_2D[dihedral_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[1]][0],geo_2D[dihedral_atoms[2]][0]],[geo_2D[dihedral_atoms[1]][1],geo_2D[dihedral_atoms[2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[2]][0],geo_2D[dihedral_atoms[3]][0]],[geo_2D[dihedral_atoms[2]][1],geo_2D[dihedral_atoms[3]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                # Add atom labels
                fontsize = 20*draw_scale*5.0
                if fontsize > 24:
                    fontsize = 24
                if fontsize < 9:
                    fontsize = 9            
                for count_i,i in enumerate(geo_2D):
                    if count_i in dihedral_atoms:
                        ax1.text(i[0], i[1],Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"][count_i], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1],Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"][count_i], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Add data
                xdata = np.array(thetas)*180.0/np.pi
                if fallback_flag == 1 or read_flag is True:
                    ydata = dihedral_harmonic_norm(xdata*np.pi/180.0,DFT_params[0],d=-1.0,n=2.0,norm=0.0)
                else:
                    ydata = dihedral_harmonic_norm(xdata*np.pi/180.0,DFT_params[0],d=-1.0,n=2.0,norm=0.0) + E_corr + EC_corr + LJ_corr
                thetas = thetas-offset # NEW
                xdata = np.array(thetas)*180.0/np.pi # NEW
                min_corr = min(ydata)
                ax2.plot(xdata,ydata-min_corr,'--',linewidth=3,color=(0.0,0.1,0.8))
                min_corr = min(DFT_E+EC_corr+LJ_corr)
                if fallback_flag == 1 or read_flag is True:
                    ax2.scatter(np.array(thetas)*180.0/np.pi,DFT_E,s=120,color=(0.0,0.1,0.8))
                else:
                    ax2.scatter(np.array(thetas)*180.0/np.pi,DFT_E+EC_corr+LJ_corr-min_corr,s=120,color=(0.0,0.1,0.8))
                if fallback_flag == 1:
                    ax2.text(0.9, 0.8,'*', ha='center', fontsize=40, va='center', transform=ax2.transAxes)

                # ydata = dihedral_harmonic_norm(xdata*np.pi/180.0,DFT_params[0],d=-1.0,n=2.0,norm=0.0) + E_corr + EC_corr + LJ_corr
                # min_corr = min(ydata)
                # ax2.plot(xdata,ydata-min_corr,'--',linewidth=3,color=(0.0,0.1,0.8))
                # min_corr = min(DFT_E+EC_corr+LJ_corr)
                # ax2.scatter(np.array(thetas)*180.0/np.pi,DFT_E+EC_corr+LJ_corr-min_corr,s=120,color=(0.0,0.1,0.8))
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{\theta \, (degrees)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(thetas)*180.0/np.pi, max(thetas)*180.0/np.pi])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = save_folder+'/'+"_".join(dihedrals)+'-DFT_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(save_folder+'/'+"_".join(dihedrals)+'-DFT_fit.pdf', bbox_inches=0,dpi=300)
                paths_to_dihedral_fitpot_plots_DFT+=[Name]
                paths_to_dihedral_totpot_plots_DFT+=[Name]
                plt.close(fig)

            if 'mp2' in qc_types:

                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                draw_scale = 1.0
                # Add bonds
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=2.0*draw_scale))
                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[0]][0],geo_2D[angle_atoms[1]][0]],[geo_2D[angle_atoms[0]][1],geo_2D[angle_atoms[1]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[angle_atoms[1]][0],geo_2D[angle_atoms[2]][0]],[geo_2D[angle_atoms[1]][1],geo_2D[angle_atoms[2]][1]],color=(0,0.1,0.8),linewidth=2.0*draw_scale))

                # Add atom labels
                for count_i,i in enumerate(geo_2D):
                    if count_i in angle_atoms:
                        ax1.text(i[0], i[1],Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1],Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=16*draw_scale, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Add data
                xdata = np.array(thetas)*180.0/np.pi
                ydata = harmonic(xdata*np.pi/180.0,MP2_params[0],MP2_params[1]) + E_corr
                ax2.plot(xdata,ydata,'--',linewidth=3,color=(0.0,0.1,0.8))
                ax2.scatter(np.array(thetas)*180.0/np.pi,MP2_E,s=120,color=(0.0,0.1,0.8))
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{\theta \, (degrees)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(thetas)*180.0/np.pi, max(thetas)*180.0/np.pi])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = save_folder+'/'+"_".join(angles)+'-MP2_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(save_folder+'/'+"_".join(angles)+'-MP2_fit.pdf', bbox_inches=0,dpi=300)
                paths_to_angle_plots_MP2+=[Name]
                plt.close(fig)
        
        #############################################################
        ####### Update FF_dict with DFT and MP2 FF-parameters #######
        #############################################################
        # NOTE: the force constant is in units of kcal/mol
        FF_dict["dihedrals_harmonic"][dihedrals]={}
        FF_dict["dihedrals_harmonic"][tuple([j+'-UA' for j in dihedrals])]={}
        if 'dft' in qc_types:
            FF_dict["dihedrals_harmonic"][dihedrals]["DFT"] = ["harmonic"] + [ DFT_params[0], -1,2 ]
            # Save the UA types as duplicate entries (in the future other scripts should instead expect UA and AA harm dihedrals to be equivalent
            print("adding duplicated UA entry: {} -> {}".format(dihedrals, tuple([j+'-UA' for j in dihedrals])))
            FF_dict["dihedrals_harmonic"][tuple([j+'-UA' for j in dihedrals])]["DFT"] = ["harmonic"] + [ DFT_params[0], -1,2 ]
        if 'mp2' in qc_types:
            FF_dict["dihedrals_harmonic"][dihedrals]["MP2"] = ["harmonic"] + [ MP2_params[0], -1,2 ]
    
    # Return to parent folder after processing all angles
    os.chdir(working_dir)

    return paths_to_dihedral_fitpot_plots_DFT, paths_to_dihedral_fitpot_plots_MP2, paths_to_dihedral_totpot_plots_DFT, paths_to_dihedral_totpot_plots_MP2

# # ORIGINAL Fit function used for self-consistently fitting bonds and angles with respect to other modes that are shifting
# def fit_bonds_angles_sc(k_0  = None,r_0  = None,k_1  = None,r_1  = None,k_2  = None,r_2  = None,k_3  = None,r_3  = None, k_4  = None,r_4  = None,\
#                         k_5  = None,r_5  = None,k_6  = None,r_6  = None,k_7  = None,r_7  = None,k_8  = None,r_8  = None, k_9  = None,r_9  = None,\
#                         k_10 = None,r_10 = None,k_11 = None,r_11 = None,k_12 = None,r_12 = None,k_13 = None,r_13 = None, k_14 = None,r_14 = None,\
#                         k_15 = None,r_15 = None,k_16 = None,r_16 = None,k_17 = None,r_17 = None,k_18 = None,r_18 = None, k_19 = None,r_19 = None,\
#                         k_20 = None,r_20 = None,k_21 = None,r_21 = None,k_22 = None,r_22 = None,k_23 = None,r_23 = None, k_24 = None,r_24 = None,\
#                         k_25 = None,r_25 = None,k_26 = None,r_26 = None,k_27 = None,r_27 = None,k_28 = None,r_28 = None, k_29 = None,r_29 = None,\
#                         k_30 = None,r_30 = None,k_31 = None,r_31 = None,k_32 = None,r_32 = None,k_33 = None,r_33 = None, k_34 = None,r_34 = None,\
#                         k_35 = None,r_35 = None,k_36 = None,r_36 = None,k_37 = None,r_37 = None,k_38 = None,r_38 = None, k_39 = None,r_39 = None,\
#                         k_40 = None,r_40 = None,k_41 = None,r_41 = None,k_42 = None,r_42 = None,k_43 = None,r_43 = None, k_44 = None,r_44 = None,\
#                         k_45 = None,r_45 = None,k_46 = None,r_46 = None,k_47 = None,r_47 = None,k_48 = None,r_48 = None, k_49 = None,r_49 = None,\
#                         k_50 = None,r_50 = None,k_51 = None,r_51 = None,k_52 = None,r_52 = None,k_53 = None,r_53 = None, k_54 = None,r_54 = None,\
#                         k_55 = None,r_55 = None,k_56 = None,r_56 = None,k_57 = None,r_57 = None,k_58 = None,r_58 = None, k_59 = None,r_59 = None,\
#                         k_60 = None,r_60 = None,k_61 = None,r_61 = None,k_62 = None,r_62 = None,k_63 = None,r_63 = None, k_64 = None,r_64 = None,\
#                         k_65 = None,r_65 = None,k_66 = None,r_66 = None,k_67 = None,r_67 = None,k_68 = None,r_68 = None, k_69 = None,r_69 = None,\
#                         k_70 = None,r_70 = None,k_71 = None,r_71 = None,k_72 = None,r_72 = None,k_73 = None,r_73 = None, k_74 = None,r_74 = None,\
#                         k_75 = None,r_75 = None,k_76 = None,r_76 = None,k_77 = None,r_77 = None,k_78 = None,r_78 = None, k_79 = None,r_79 = None,\
#                         k_80 = None,r_80 = None,k_81 = None,r_81 = None,k_82 = None,r_82 = None,k_83 = None,r_83 = None, k_84 = None,r_84 = None,\
#                         k_85 = None,r_85 = None,k_86 = None,r_86 = None,k_87 = None,r_87 = None,k_88 = None,r_88 = None, k_89 = None,r_89 = None,\
#                         k_90 = None,r_90 = None,k_91 = None,r_91 = None,k_92 = None,r_92 = None,k_93 = None,r_93 = None, k_94 = None,r_94 = None,\
#                         K_95 = None,r_95 = None,k_96 = None,r_96 = None,k_97 = None,r_97 = None,k_98 = None,r_98 = None, k_99 = None,r_99 = None,\
#                         modes=None, fit_type=None,E_0=None,w_pot=1.0,w_hyper=0.0,b_hyper=0.1,w_exp=0.0,w_harm=0.0,T=298.0,kb=0.0019872041,penalty=True):
                      
#     # Initialize local variable dictionary (used for determining what has been defined)
#     local_vars = locals()

#     # Find defined k_* and r_* variables. 
#     k_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'k_' in i and "__" not in i and local_vars[i] is not None ]) ])
#     r_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'r_' in i and "__" not in i and local_vars[i] is not None ]) ])

#     # Consistency checks
#     # if unique_fit_dihedral_types is None or angles is None or V_range is None or dihedral_types is None or QC_dihedral is None or FF_dict is None:
#     #     print "ERROR in OPLS_OPLS_Fit_MIN: {} must be supplied to the function, it/they "\
#     #     .format(", ".join([ i for i in ["unique_fit_dihedral_types","angles","V_range","dihedral_types","QC_dihedral","FF_dict"] if locals()[i] is None ]))
#     #     print "                            are defined as optional arguments for convenient anonymization. Exiting..."
#     #     quit()
#     # if len(angles[0]) != len(dihedral_types):
#     #     print "ERROR in OPLS_OPLS_Fit_MIN: dihedral_types must have the same length as each tuple in angles. Exiting..."
#     #     quit()
#     if len(modes) != len(E_0):
#         print "ERROR in fit_bonds_angles_sc: E_0 and modes must have the same length. Exiting..."
#         quit()
#     if len(k_vals) != len(r_vals) or len(k_vals) != len(modes[0]):
#         print "ERROR in fit_bonds_angles_sc: the function expects the len(modes) and number of non-zero k*_* and r*_* variables to be equal. Exiting..."
#         quit()

#     # Initialize a list of bond and angle harmonic parameters indexed to the modes
#     # The if/else controls whether the coefficients are assigned to the input fit parameters or
#     # taken from FF_dict (dihedrals not being fit on this scan)
#     params = []
#     for count_i,i in enumerate(modes[0]):
#         params += [ (k_vals[count_i],r_vals[count_i]) ]

#     # Initialize E_tot for each configuration
#     E_tot = np.zeros(len(modes))

#     # Cumulatively add up the bond and angle energies (fit type should be defined as a global variable outside the function)
#     for count_i,i in enumerate(modes):

#         # modes holds a list of tuples with bonds and angles of each coincident type being fit
#         for count_j,j in enumerate(i):
#             E_tot[count_i] += params[count_j][0]*(j[1]-params[count_j][1])**(2.0)

#     # Return the penalty function
#     if penalty is True:
#         #      potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes                                        hyperbolic constraint on parameter magnitudes
#         return w_pot*np.mean((E_tot - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
#     else:
#         return E_tot


# (NEW: contains normalization constant) Fit function used for self-consistently fitting bonds and angles with respect to other modes that are shifting
def fit_bonds_angles_sc(norm = 0.0, k_0  = None,r_0  = None,k_1  = None,r_1  = None,k_2  = None,r_2  = None,k_3  = None,r_3  = None, k_4  = None,r_4  = None,\
                                    k_5  = None,r_5  = None,k_6  = None,r_6  = None,k_7  = None,r_7  = None,k_8  = None,r_8  = None, k_9  = None,r_9  = None,\
                                    k_10 = None,r_10 = None,k_11 = None,r_11 = None,k_12 = None,r_12 = None,k_13 = None,r_13 = None, k_14 = None,r_14 = None,\
                                    k_15 = None,r_15 = None,k_16 = None,r_16 = None,k_17 = None,r_17 = None,k_18 = None,r_18 = None, k_19 = None,r_19 = None,\
                                    k_20 = None,r_20 = None,k_21 = None,r_21 = None,k_22 = None,r_22 = None,k_23 = None,r_23 = None, k_24 = None,r_24 = None,\
                                    k_25 = None,r_25 = None,k_26 = None,r_26 = None,k_27 = None,r_27 = None,k_28 = None,r_28 = None, k_29 = None,r_29 = None,\
                                    k_30 = None,r_30 = None,k_31 = None,r_31 = None,k_32 = None,r_32 = None,k_33 = None,r_33 = None, k_34 = None,r_34 = None,\
                                    k_35 = None,r_35 = None,k_36 = None,r_36 = None,k_37 = None,r_37 = None,k_38 = None,r_38 = None, k_39 = None,r_39 = None,\
                                    k_40 = None,r_40 = None,k_41 = None,r_41 = None,k_42 = None,r_42 = None,k_43 = None,r_43 = None, k_44 = None,r_44 = None,\
                                    k_45 = None,r_45 = None,k_46 = None,r_46 = None,k_47 = None,r_47 = None,k_48 = None,r_48 = None, k_49 = None,r_49 = None,\
                                    k_50 = None,r_50 = None,k_51 = None,r_51 = None,k_52 = None,r_52 = None,k_53 = None,r_53 = None, k_54 = None,r_54 = None,\
                                    k_55 = None,r_55 = None,k_56 = None,r_56 = None,k_57 = None,r_57 = None,k_58 = None,r_58 = None, k_59 = None,r_59 = None,\
                                    k_60 = None,r_60 = None,k_61 = None,r_61 = None,k_62 = None,r_62 = None,k_63 = None,r_63 = None, k_64 = None,r_64 = None,\
                                    k_65 = None,r_65 = None,k_66 = None,r_66 = None,k_67 = None,r_67 = None,k_68 = None,r_68 = None, k_69 = None,r_69 = None,\
                                    k_70 = None,r_70 = None,k_71 = None,r_71 = None,k_72 = None,r_72 = None,k_73 = None,r_73 = None, k_74 = None,r_74 = None,\
                                    k_75 = None,r_75 = None,k_76 = None,r_76 = None,k_77 = None,r_77 = None,k_78 = None,r_78 = None, k_79 = None,r_79 = None,\
                                    k_80 = None,r_80 = None,k_81 = None,r_81 = None,k_82 = None,r_82 = None,k_83 = None,r_83 = None, k_84 = None,r_84 = None,\
                                    k_85 = None,r_85 = None,k_86 = None,r_86 = None,k_87 = None,r_87 = None,k_88 = None,r_88 = None, k_89 = None,r_89 = None,\
                                    k_90 = None,r_90 = None,k_91 = None,r_91 = None,k_92 = None,r_92 = None,k_93 = None,r_93 = None, k_94 = None,r_94 = None,\
                                    K_95 = None,r_95 = None,k_96 = None,r_96 = None,k_97 = None,r_97 = None,k_98 = None,r_98 = None, k_99 = None,r_99 = None,\
                                    modes=None, fit_type=None,E_0=None,w_pot=1.0,w_hyper=0.0,b_hyper=0.1,w_exp=0.0,w_harm=0.0,T=298.0,kb=0.0019872041,penalty=True):
                      
    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined k_* and r_* variables. 
    k_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'k_' in i and "__" not in i and local_vars[i] is not None ]) ])
    r_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'r_' in i and "__" not in i and local_vars[i] is not None ]) ])

    # Consistency checks
    # if unique_fit_dihedral_types is None or angles is None or V_range is None or dihedral_types is None or QC_dihedral is None or FF_dict is None:
    #     print "ERROR in OPLS_OPLS_Fit_MIN: {} must be supplied to the function, it/they "\
    #     .format(", ".join([ i for i in ["unique_fit_dihedral_types","angles","V_range","dihedral_types","QC_dihedral","FF_dict"] if locals()[i] is None ]))
    #     print "                            are defined as optional arguments for convenient anonymization. Exiting..."
    #     quit()
    # if len(angles[0]) != len(dihedral_types):
    #     print "ERROR in OPLS_OPLS_Fit_MIN: dihedral_types must have the same length as each tuple in angles. Exiting..."
    #     quit()
    if len(modes) != len(E_0):
        print("ERROR in fit_bonds_angles_sc: E_0 and modes must have the same length. Exiting...")
        quit()
    if len(k_vals) != len(r_vals) or len(k_vals) != len(modes[0]):
        print("ERROR in fit_bonds_angles_sc: the function expects the len(modes) and number of non-zero k*_* and r*_* variables to be equal. Exiting...")
        quit()

    # Initialize a list of bond and angle harmonic parameters indexed to the modes
    # The if/else controls whether the coefficients are assigned to the input fit parameters or
    # taken from FF_dict (dihedrals not being fit on this scan)
    params = []
    for count_i,i in enumerate(modes[0]):
        params += [ (k_vals[count_i],r_vals[count_i]) ]

    # Initialize E_tot for each configuration
    E_tot = np.ones(len(modes))*norm

    # Cumulatively add up the bond and angle energies (fit type should be defined as a global variable outside the function)
    for count_i,i in enumerate(modes):

        # modes holds a list of tuples with bonds and angles of each coincident type being fit
        for count_j,j in enumerate(i):
            E_tot[count_i] += params[count_j][0]*(j[1]-params[count_j][1])**(2.0)
        
    # Return the penalty function
    if penalty is True:
        #      potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes                                        hyperbolic constraint on parameter magnitudes
        return w_pot*np.mean((E_tot - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
    else:
        return E_tot


# # Fit function used for self-consistently fitting bonds and angles with respect to other modes that are shifting
# def fit_bonds_angles_sc_2(norm = 0.0, k_0  = None,r_0  = None,k_1  = None,r_1  = None,k_2  = None,r_2  = None,k_3  = None,r_3  = None, k_4  = None,r_4  = None,\
#                                       k_5  = None, r_5  = None,k_6  = None,r_6  = None,k_7  = None,r_7  = None,k_8  = None,r_8  = None, k_9  = None,r_9  = None,\
#                                       k_10 = None, r_10 = None,k_11 = None,r_11 = None,k_12 = None,r_12 = None,k_13 = None,r_13 = None, k_14 = None,r_14 = None,\
#                                       k_15 = None, r_15 = None,k_16 = None,r_16 = None,k_17 = None,r_17 = None,k_18 = None,r_18 = None, k_19 = None,r_19 = None,\
#                                       k_20 = None, r_20 = None,k_21 = None,r_21 = None,k_22 = None,r_22 = None,k_23 = None,r_23 = None, k_24 = None,r_24 = None,\
#                                       k_25 = None, r_25 = None,k_26 = None,r_26 = None,k_27 = None,r_27 = None,k_28 = None,r_28 = None, k_29 = None,r_29 = None,\
#                                       k_30 = None, r_30 = None,k_31 = None,r_31 = None,k_32 = None,r_32 = None,k_33 = None,r_33 = None, k_34 = None,r_34 = None,\
#                                       k_35 = None, r_35 = None,k_36 = None,r_36 = None,k_37 = None,r_37 = None,k_38 = None,r_38 = None, k_39 = None,r_39 = None,\
#                                       k_40 = None, r_40 = None,k_41 = None,r_41 = None,k_42 = None,r_42 = None,k_43 = None,r_43 = None, k_44 = None,r_44 = None,\
#                                       k_45 = None, r_45 = None,k_46 = None,r_46 = None,k_47 = None,r_47 = None,k_48 = None,r_48 = None, k_49 = None,r_49 = None,\
#                                       k_50 = None, r_50 = None,k_51 = None,r_51 = None,k_52 = None,r_52 = None,k_53 = None,r_53 = None, k_54 = None,r_54 = None,\
#                                       k_55 = None, r_55 = None,k_56 = None,r_56 = None,k_57 = None,r_57 = None,k_58 = None,r_58 = None, k_59 = None,r_59 = None,\
#                                       k_60 = None, r_60 = None,k_61 = None,r_61 = None,k_62 = None,r_62 = None,k_63 = None,r_63 = None, k_64 = None,r_64 = None,\
#                                       k_65 = None, r_65 = None,k_66 = None,r_66 = None,k_67 = None,r_67 = None,k_68 = None,r_68 = None, k_69 = None,r_69 = None,\
#                                       k_70 = None, r_70 = None,k_71 = None,r_71 = None,k_72 = None,r_72 = None,k_73 = None,r_73 = None, k_74 = None,r_74 = None,\
#                                       k_75 = None, r_75 = None,k_76 = None,r_76 = None,k_77 = None,r_77 = None,k_78 = None,r_78 = None, k_79 = None,r_79 = None,\
#                                       k_80 = None, r_80 = None,k_81 = None,r_81 = None,k_82 = None,r_82 = None,k_83 = None,r_83 = None, k_84 = None,r_84 = None,\
#                                       k_85 = None, r_85 = None,k_86 = None,r_86 = None,k_87 = None,r_87 = None,k_88 = None,r_88 = None, k_89 = None,r_89 = None,\
#                                       k_90 = None, r_90 = None,k_91 = None,r_91 = None,k_92 = None,r_92 = None,k_93 = None,r_93 = None, k_94 = None,r_94 = None,\
#                                       K_95 = None, r_95 = None,k_96 = None,r_96 = None,k_97 = None,r_97 = None,k_98 = None,r_98 = None, k_99 = None,r_99 = None,\
#                                       modes=None, fit_type=None,E_0=None,w_pot=1.0,w_hyper=0.0,b_hyper=0.1,w_exp=0.0,w_harm=0.0,fit_modes=None,penalty=True,T=298.0,kb=0.0019872041):
                      
#     # Initialize local variable dictionary (used for determining what has been defined)
#     local_vars = locals()

#     # Find defined k_* and r_* variables. 
#     k_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'k_' in i and "__" not in i and local_vars[i] is not None ]) ])
#     r_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'r_' in i and "__" not in i and local_vars[i] is not None ]) ])

#     # Consistency checks
#     # if unique_fit_dihedral_types is None or angles is None or V_range is None or dihedral_types is None or QC_dihedral is None or FF_dict is None:
#     #     print "ERROR in OPLS_OPLS_Fit_MIN: {} must be supplied to the function, it/they "\
#     #     .format(", ".join([ i for i in ["unique_fit_dihedral_types","angles","V_range","dihedral_types","QC_dihedral","FF_dict"] if locals()[i] is None ]))
#     #     print "                            are defined as optional arguments for convenient anonymization. Exiting..."
#     #     quit()
#     # if len(angles[0]) != len(dihedral_types):
#     #     print "ERROR in OPLS_OPLS_Fit_MIN: dihedral_types must have the same length as each tuple in angles. Exiting..."
#     #     quit()
#     if len(modes) != len(E_0):
#         print "ERROR in fit_bonds_angles_sc: E_0 and modes must have the same length. Exiting..."
#         quit()
#     if len(k_vals) != len(r_vals) or len(k_vals) != len(fit_modes):
#         print "ERROR in fit_bonds_angles_sc: the function expects the len(modes) and number of non-zero k*_* and r*_* variables to be equal. Exiting..."
#         quit()

#     # Initialize a list of bond and angle harmonic parameters indexed to the modes
#     # The if/else controls whether the coefficients are assigned to the input fit parameters or
#     # taken from FF_dict (dihedrals not being fit on this scan)
#     params = []
#     for count_i,i in enumerate(modes[0]):
#         params += [ (k_vals[fit_modes.index(i[0])],r_vals[fit_modes.index(i[0])]) ]

#     # Initialize E_tot for each configuration
#     E_tot = np.ones(len(modes))*norm

#     # Cumulatively add up the bond and angle energies (fit type should be defined as a global variable outside the function)
#     for count_i,i in enumerate(modes):

#         # modes holds a list of tuples with bonds and angles of each coincident type being fit
#         for count_j,j in enumerate(i):
#             E_tot[count_i] += params[count_j][0]*(j[1]-params[count_j][1])**(2.0)

#     # Return the penalty function
#     #          potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes                                        hyperbolic constraint on parameter magnitudes
#     if penalty is True:
#         return w_pot*np.mean((E_tot - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
#     else:
#         return E_tot


# # Fit function used for self-consistently fitting bonds and angles with respect to other modes that are shifting
# # Self-consistent, both force constants and equilibrium displacements are set equal for like mode types
# def fit_bonds_angles_sc_2(norm = 0.0, k_0  = None,r_0  = None,k_1  = None,r_1  = None,k_2  = None,r_2  = None,k_3  = None,r_3  = None, k_4  = None,r_4  = None,\
#                                       k_5  = None, r_5  = None,k_6  = None,r_6  = None,k_7  = None,r_7  = None,k_8  = None,r_8  = None, k_9  = None,r_9  = None,\
#                                       k_10 = None, r_10 = None,k_11 = None,r_11 = None,k_12 = None,r_12 = None,k_13 = None,r_13 = None, k_14 = None,r_14 = None,\
#                                       k_15 = None, r_15 = None,k_16 = None,r_16 = None,k_17 = None,r_17 = None,k_18 = None,r_18 = None, k_19 = None,r_19 = None,\
#                                       k_20 = None, r_20 = None,k_21 = None,r_21 = None,k_22 = None,r_22 = None,k_23 = None,r_23 = None, k_24 = None,r_24 = None,\
#                                       k_25 = None, r_25 = None,k_26 = None,r_26 = None,k_27 = None,r_27 = None,k_28 = None,r_28 = None, k_29 = None,r_29 = None,\
#                                       k_30 = None, r_30 = None,k_31 = None,r_31 = None,k_32 = None,r_32 = None,k_33 = None,r_33 = None, k_34 = None,r_34 = None,\
#                                       k_35 = None, r_35 = None,k_36 = None,r_36 = None,k_37 = None,r_37 = None,k_38 = None,r_38 = None, k_39 = None,r_39 = None,\
#                                       k_40 = None, r_40 = None,k_41 = None,r_41 = None,k_42 = None,r_42 = None,k_43 = None,r_43 = None, k_44 = None,r_44 = None,\
#                                       k_45 = None, r_45 = None,k_46 = None,r_46 = None,k_47 = None,r_47 = None,k_48 = None,r_48 = None, k_49 = None,r_49 = None,\
#                                       k_50 = None, r_50 = None,k_51 = None,r_51 = None,k_52 = None,r_52 = None,k_53 = None,r_53 = None, k_54 = None,r_54 = None,\
#                                       k_55 = None, r_55 = None,k_56 = None,r_56 = None,k_57 = None,r_57 = None,k_58 = None,r_58 = None, k_59 = None,r_59 = None,\
#                                       k_60 = None, r_60 = None,k_61 = None,r_61 = None,k_62 = None,r_62 = None,k_63 = None,r_63 = None, k_64 = None,r_64 = None,\
#                                       k_65 = None, r_65 = None,k_66 = None,r_66 = None,k_67 = None,r_67 = None,k_68 = None,r_68 = None, k_69 = None,r_69 = None,\
#                                       k_70 = None, r_70 = None,k_71 = None,r_71 = None,k_72 = None,r_72 = None,k_73 = None,r_73 = None, k_74 = None,r_74 = None,\
#                                       k_75 = None, r_75 = None,k_76 = None,r_76 = None,k_77 = None,r_77 = None,k_78 = None,r_78 = None, k_79 = None,r_79 = None,\
#                                       k_80 = None, r_80 = None,k_81 = None,r_81 = None,k_82 = None,r_82 = None,k_83 = None,r_83 = None, k_84 = None,r_84 = None,\
#                                       k_85 = None, r_85 = None,k_86 = None,r_86 = None,k_87 = None,r_87 = None,k_88 = None,r_88 = None, k_89 = None,r_89 = None,\
#                                       k_90 = None, r_90 = None,k_91 = None,r_91 = None,k_92 = None,r_92 = None,k_93 = None,r_93 = None, k_94 = None,r_94 = None,\
#                                       K_95 = None, r_95 = None,k_96 = None,r_96 = None,k_97 = None,r_97 = None,k_98 = None,r_98 = None, k_99 = None,r_99 = None,\
#                                       modes=None, fit_type=None,E_0=None,w_pot=1.0,w_hyper=0.0,b_hyper=0.1,w_exp=0.0,w_harm=0.0,fit_modes=None,penalty=True,T=298.0,kb=0.0019872041):
                      
#     # Initialize local variable dictionary (used for determining what has been defined)
#     local_vars = locals()

#     # Find defined k_* and r_* variables. 
#     k_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'k_' in i and "__" not in i and local_vars[i] is not None ]) ])
#     r_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'r_' in i and "__" not in i and local_vars[i] is not None ]) ])

#     # Consistency checks
#     # if unique_fit_dihedral_types is None or angles is None or V_range is None or dihedral_types is None or QC_dihedral is None or FF_dict is None:
#     #     print "ERROR in OPLS_OPLS_Fit_MIN: {} must be supplied to the function, it/they "\
#     #     .format(", ".join([ i for i in ["unique_fit_dihedral_types","angles","V_range","dihedral_types","QC_dihedral","FF_dict"] if locals()[i] is None ]))
#     #     print "                            are defined as optional arguments for convenient anonymization. Exiting..."
#     #     quit()
#     # if len(angles[0]) != len(dihedral_types):
#     #     print "ERROR in OPLS_OPLS_Fit_MIN: dihedral_types must have the same length as each tuple in angles. Exiting..."
#     #     quit()
#     if len(modes) != len(E_0):
#         print "ERROR in fit_bonds_angles_sc: E_0 and modes must have the same length. Exiting..."
#         quit()
#     if len(k_vals) != len(r_vals) or len(k_vals) != len(fit_modes):
#         print "ERROR in fit_bonds_angles_sc: the function expects the len(modes) and number of non-zero k*_* and r*_* variables to be equal. Exiting..."
#         quit()

#     # Initialize a list of bond and angle harmonic parameters indexed to the modes
#     # The if/else controls whether the coefficients are assigned to the input fit parameters or
#     # taken from FF_dict (dihedrals not being fit on this scan)
#     params = []
#     for count_i,i in enumerate(modes[0]):
        
#         # harmonic dihedrals are handled specially
#         if len(i[0]) == 4:
#             params += [ (k_vals[fit_modes.index(i[0])],-1.0,2.0) ]
#         else:
#             params += [ (k_vals[fit_modes.index(i[0])],r_vals[fit_modes.index(i[0])]) ]

#     # Initialize E_tot for each configuration
#     E_tot = np.ones(len(modes))*norm

#     # Cumulatively add up the bond and angle energies (fit type should be defined as a global variable outside the function)
#     for count_i,i in enumerate(modes):

#         # modes holds a list of tuples with bonds and angles of each coincident type being fit
#         for count_j,j in enumerate(i):

#             # harmonic dihedrals are handled specially with a k(1-np.cos(2*x)) function  
#             if len(j[0]) == 4:
#                 E_tot[count_i] += params[count_j][0]*(1.0+params[count_j][1]*np.cos(params[count_j][2]*j[1]))
#             # bonds and angles are fit to harmonics
#             else:                
#                 E_tot[count_i] += params[count_j][0]*(j[1]-params[count_j][1])**(2.0)

#     # Return the penalty function
#     #          potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes                                        hyperbolic constraint on parameter magnitudes
#     if penalty is True:
#         return w_pot*np.mean((E_tot - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
#     else:
#         return E_tot

# Fit function used for self-consistently fitting bonds and angles with respect to other modes that are shifting
# Self-consistent, force constants are equal for like mode types and equilibrium displacements are unqiue for each instance. 
def fit_bonds_angles_sc_2(norm = 0.0, k_0  = None,r_0  = None,k_1  = None,r_1  = None,k_2  = None,r_2  = None,k_3  = None,r_3  = None, k_4  = None,r_4  = None,\
                                      k_5  = None, r_5  = None,k_6  = None,r_6  = None,k_7  = None,r_7  = None,k_8  = None,r_8  = None, k_9  = None,r_9  = None,\
                                      k_10 = None, r_10 = None,k_11 = None,r_11 = None,k_12 = None,r_12 = None,k_13 = None,r_13 = None, k_14 = None,r_14 = None,\
                                      k_15 = None, r_15 = None,k_16 = None,r_16 = None,k_17 = None,r_17 = None,k_18 = None,r_18 = None, k_19 = None,r_19 = None,\
                                      k_20 = None, r_20 = None,k_21 = None,r_21 = None,k_22 = None,r_22 = None,k_23 = None,r_23 = None, k_24 = None,r_24 = None,\
                                      k_25 = None, r_25 = None,k_26 = None,r_26 = None,k_27 = None,r_27 = None,k_28 = None,r_28 = None, k_29 = None,r_29 = None,\
                                      k_30 = None, r_30 = None,k_31 = None,r_31 = None,k_32 = None,r_32 = None,k_33 = None,r_33 = None, k_34 = None,r_34 = None,\
                                      k_35 = None, r_35 = None,k_36 = None,r_36 = None,k_37 = None,r_37 = None,k_38 = None,r_38 = None, k_39 = None,r_39 = None,\
                                      k_40 = None, r_40 = None,k_41 = None,r_41 = None,k_42 = None,r_42 = None,k_43 = None,r_43 = None, k_44 = None,r_44 = None,\
                                      k_45 = None, r_45 = None,k_46 = None,r_46 = None,k_47 = None,r_47 = None,k_48 = None,r_48 = None, k_49 = None,r_49 = None,\
                                      k_50 = None, r_50 = None,k_51 = None,r_51 = None,k_52 = None,r_52 = None,k_53 = None,r_53 = None, k_54 = None,r_54 = None,\
                                      k_55 = None, r_55 = None,k_56 = None,r_56 = None,k_57 = None,r_57 = None,k_58 = None,r_58 = None, k_59 = None,r_59 = None,\
                                      k_60 = None, r_60 = None,k_61 = None,r_61 = None,k_62 = None,r_62 = None,k_63 = None,r_63 = None, k_64 = None,r_64 = None,\
                                      k_65 = None, r_65 = None,k_66 = None,r_66 = None,k_67 = None,r_67 = None,k_68 = None,r_68 = None, k_69 = None,r_69 = None,\
                                      k_70 = None, r_70 = None,k_71 = None,r_71 = None,k_72 = None,r_72 = None,k_73 = None,r_73 = None, k_74 = None,r_74 = None,\
                                      k_75 = None, r_75 = None,k_76 = None,r_76 = None,k_77 = None,r_77 = None,k_78 = None,r_78 = None, k_79 = None,r_79 = None,\
                                      k_80 = None, r_80 = None,k_81 = None,r_81 = None,k_82 = None,r_82 = None,k_83 = None,r_83 = None, k_84 = None,r_84 = None,\
                                      k_85 = None, r_85 = None,k_86 = None,r_86 = None,k_87 = None,r_87 = None,k_88 = None,r_88 = None, k_89 = None,r_89 = None,\
                                      k_90 = None, r_90 = None,k_91 = None,r_91 = None,k_92 = None,r_92 = None,k_93 = None,r_93 = None, k_94 = None,r_94 = None,\
                                      K_95 = None, r_95 = None,k_96 = None,r_96 = None,k_97 = None,r_97 = None,k_98 = None,r_98 = None, k_99 = None,r_99 = None,\
                                      modes=None, fit_type=None,E_0=None,w_pot=1.0,w_hyper=0.0,b_hyper=0.1,w_exp=0.0,w_harm=0.0,fit_modes=None,penalty=True,T=298.0,kb=0.0019872041):
                      
    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined k_* and r_* variables. 
    k_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'k_' in i and "__" not in i and local_vars[i] is not None ]) ])
    r_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'r_' in i and "__" not in i and local_vars[i] is not None ]) ])

    # Consistency checks
    # if unique_fit_dihedral_types is None or angles is None or V_range is None or dihedral_types is None or QC_dihedral is None or FF_dict is None:
    #     print "ERROR in OPLS_OPLS_Fit_MIN: {} must be supplied to the function, it/they "\
    #     .format(", ".join([ i for i in ["unique_fit_dihedral_types","angles","V_range","dihedral_types","QC_dihedral","FF_dict"] if locals()[i] is None ]))
    #     print "                            are defined as optional arguments for convenient anonymization. Exiting..."
    #     quit()
    # if len(angles[0]) != len(dihedral_types):
    #     print "ERROR in OPLS_OPLS_Fit_MIN: dihedral_types must have the same length as each tuple in angles. Exiting..."
    #     quit()
    if len(modes) != len(E_0):
        print("ERROR in fit_bonds_angles_sc: E_0 and modes must have the same length. Exiting...")
        quit()
    if len(k_vals) != len(r_vals) or len(k_vals) != len(fit_modes):
        print("ERROR in fit_bonds_angles_sc: the function expects the len(modes) and number of non-zero k*_* and r*_* variables to be equal. Exiting...")
        quit()

    # Initialize a list of bond and angle harmonic parameters indexed to the modes
    # The if/else controls whether the coefficients are assigned to the input fit parameters or
    # taken from FF_dict (dihedrals not being fit on this scan)
    params = []
    for count_i,i in enumerate(modes[0]):
        
        # harmonic dihedrals are handled specially
        if len(i[0]) == 4:
            params += [ (k_vals[fit_modes.index(i[0])],-1.0,2.0) ]

        # Use the same force constant for all like modes but use unique equilibrium positions
        else:
            params += [ (k_vals[fit_modes.index(i[0])],r_vals[count_i]) ]

    # Initialize E_tot for each configuration
    E_tot = np.ones(len(modes))*norm

    # Cumulatively add up the bond and angle energies (fit type should be defined as a global variable outside the function)
    for count_i,i in enumerate(modes):

        # modes holds a list of tuples with bonds and angles of each coincident type being fit
        for count_j,j in enumerate(i):

            # harmonic dihedrals are handled specially with a k(1-np.cos(2*x)) function  
            if len(j[0]) == 4:
                E_tot[count_i] += params[count_j][0]*(1.0+params[count_j][1]*np.cos(params[count_j][2]*j[1]))
            # bonds and angles are fit to harmonics
            else:                
                E_tot[count_i] += params[count_j][0]*(j[1]-params[count_j][1])**(2.0)

    # Return the penalty function
    #          potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes                                        hyperbolic constraint on parameter magnitudes
    if penalty is True:
        return w_pot*np.mean((E_tot - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
    else:
        return E_tot



# Fit function used for self-consistently fitting bonds and angles with respect to other modes that are shifting
def fit_bonds_angles_sc_single(params = None,                          k_1  = None,r_1  = None,k_2  = None,r_2  = None,k_3  = None,r_3  = None, k_4  = None,r_4  = None,\
                                              k_5  = None, r_5  = None,k_6  = None,r_6  = None,k_7  = None,r_7  = None,k_8  = None,r_8  = None, k_9  = None,r_9  = None,\
                                              k_10 = None, r_10 = None,k_11 = None,r_11 = None,k_12 = None,r_12 = None,k_13 = None,r_13 = None, k_14 = None,r_14 = None,\
                                              k_15 = None, r_15 = None,k_16 = None,r_16 = None,k_17 = None,r_17 = None,k_18 = None,r_18 = None, k_19 = None,r_19 = None,\
                                              k_20 = None, r_20 = None,k_21 = None,r_21 = None,k_22 = None,r_22 = None,k_23 = None,r_23 = None, k_24 = None,r_24 = None,\
                                              k_25 = None, r_25 = None,k_26 = None,r_26 = None,k_27 = None,r_27 = None,k_28 = None,r_28 = None, k_29 = None,r_29 = None,\
                                              k_30 = None, r_30 = None,k_31 = None,r_31 = None,k_32 = None,r_32 = None,k_33 = None,r_33 = None, k_34 = None,r_34 = None,\
                                              k_35 = None, r_35 = None,k_36 = None,r_36 = None,k_37 = None,r_37 = None,k_38 = None,r_38 = None, k_39 = None,r_39 = None,\
                                              k_40 = None, r_40 = None,k_41 = None,r_41 = None,k_42 = None,r_42 = None,k_43 = None,r_43 = None, k_44 = None,r_44 = None,\
                                              k_45 = None, r_45 = None,k_46 = None,r_46 = None,k_47 = None,r_47 = None,k_48 = None,r_48 = None, k_49 = None,r_49 = None,\
                                              k_50 = None, r_50 = None,k_51 = None,r_51 = None,k_52 = None,r_52 = None,k_53 = None,r_53 = None, k_54 = None,r_54 = None,\
                                              k_55 = None, r_55 = None,k_56 = None,r_56 = None,k_57 = None,r_57 = None,k_58 = None,r_58 = None, k_59 = None,r_59 = None,\
                                              k_60 = None, r_60 = None,k_61 = None,r_61 = None,k_62 = None,r_62 = None,k_63 = None,r_63 = None, k_64 = None,r_64 = None,\
                                              k_65 = None, r_65 = None,k_66 = None,r_66 = None,k_67 = None,r_67 = None,k_68 = None,r_68 = None, k_69 = None,r_69 = None,\
                                              k_70 = None, r_70 = None,k_71 = None,r_71 = None,k_72 = None,r_72 = None,k_73 = None,r_73 = None, k_74 = None,r_74 = None,\
                                              k_75 = None, r_75 = None,k_76 = None,r_76 = None,k_77 = None,r_77 = None,k_78 = None,r_78 = None, k_79 = None,r_79 = None,\
                                              k_80 = None, r_80 = None,k_81 = None,r_81 = None,k_82 = None,r_82 = None,k_83 = None,r_83 = None, k_84 = None,r_84 = None,\
                                              k_85 = None, r_85 = None,k_86 = None,r_86 = None,k_87 = None,r_87 = None,k_88 = None,r_88 = None, k_89 = None,r_89 = None,\
                                              k_90 = None, r_90 = None,k_91 = None,r_91 = None,k_92 = None,r_92 = None,k_93 = None,r_93 = None, k_94 = None,r_94 = None,\
                                              K_95 = None, r_95 = None,k_96 = None,r_96 = None,k_97 = None,r_97 = None,k_98 = None,r_98 = None, k_99 = None,r_99 = None,\
                                              modes=None, fit_type=None,E_0=None,w_pot=1.0,w_hyper=0.0,b_hyper=0.1,w_exp=0.0,w_harm=0.0,fit_modes=None,penalty=True,T=298.0,kb=0.0019872041):
                   
    # Assign the norm, force constant, and equilibrium displacement for the mode being fit
    norm = params[0]
    k_0  = params[1]
    r_0  = params[2]

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined k_* and r_* variables. 
    k_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'k_' in i and "__" not in i and local_vars[i] is not None ]) ])
    r_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'r_' in i and "__" not in i and local_vars[i] is not None ]) ])

    # Consistency checks
    if len(modes) != len(E_0):
        print("ERROR in fit_bonds_angles_sc: E_0 and modes must have the same length. Exiting...")
        quit()
    if len(k_vals) != len(r_vals) or len(k_vals) != len(fit_modes):
        print("ERROR in fit_bonds_angles_sc: the function expects the len(modes) and number of non-zero k*_* and r*_* variables to be equal. Exiting...")
        quit()

    # Initialize a list of bond and angle harmonic parameters indexed to the modes
    # The if/else controls whether the coefficients are assigned to the input fit parameters or
    # taken from FF_dict (dihedrals not being fit on this scan)
    params = []
    for count_i,i in enumerate(modes[0]):
        
        # harmonic dihedrals are handled specially
        if len(i[0]) == 4:
            params += [ (k_vals[fit_modes.index(i[0])],-1.0,2.0) ]

        # Use the same force constant for all like modes but use unique equilibrium positions
        else:
            params += [ (k_vals[fit_modes.index(i[0])],r_vals[count_i]) ]

    # Initialize E_tot for each configuration
    E_tot = np.ones(len(modes))*norm

    # Cumulatively add up the bond and angle energies (fit type should be defined as a global variable outside the function)
    for count_i,i in enumerate(modes):

        # modes holds a list of tuples with bonds and angles of each coincident type being fit
        for count_j,j in enumerate(i):

            # harmonic dihedrals are handled specially with a k(1-np.cos(2*x)) function  
            if len(j[0]) == 4:
                E_tot[count_i] += params[count_j][0]*(1.0+params[count_j][1]*np.cos(params[count_j][2]*j[1]))
            # bonds and angles are fit to harmonics
            else:                
                E_tot[count_i] += params[count_j][0]*(j[1]-params[count_j][1])**(2.0)

    # Return the penalty function
    #          potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes                                        hyperbolic constraint on parameter magnitudes
    if penalty is True:
        return w_pot*np.mean((E_tot - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
    else:
        return E_tot

# Description: write function for generating a geometry optimization for an unconverged scan
#
# Inputs:     template: string holding the filename of the input file for the job that failed to converge.
#             geo:      array holding the cartesian coordinates of the initial guess for the new geometry
#             elements: list, indexed to geo, which holds the elements in the system
#             delta:    the step size for displacements in the scan
#            N_steps:  the number of steps in the scan
#
# Returns:    None
def gen_reoptimization(template,geo,elements,delta,N_steps):

    # Create folder to hold the reoptimized job
    save_folder='REDO'
    current_path='./'
    while os.path.isdir(current_path+'REDO'):
        save_folder  += '/REDO'
        current_path += 'REDO/'       

#        save_folder+='REDO/'+save_folder        
#        current_path += '/REDO'       
    os.mkdir(save_folder)
    os.mkdir(save_folder+'/geo_opt')
    
    # Parse information from the original job
    orca_dict = orca_in_parse(template)

    # Write the new input file
    with open(save_folder+'/geo_opt/geo_opt.in','w') as f:        
        f.write("! "+orca_dict["0"]["header_commands"])
        f.write("\n%base \"geo_opt\"\n\n")
        f.write("%geom MaxIter 1000\nend\n\n")
        f.write("* xyz {:d} {:d}\n".format(int(orca_dict["0"]["charge"]),int(orca_dict["0"]["multiplicity"])))
        for count_i,i in enumerate(elements):
            f.write("  {:<20s} {:<20.8f} {:<20.8f} {:<20.8f}\n".format(i,geo[count_i,0],geo[count_i,1],geo[count_i,2]))
        f.write("*\n")                

    # Write the constraint file. This is used by the restart_scans.py program to set the constraints in the reoptimized job
    with open(save_folder+'/constraints.txt','w') as f:
        constraints = [ i.split()[1:-1] for i in orca_dict["0"]["constraints"].split('\n') if len(i.split()[1:-1]) > 0 ]

        # Save only the atom indices that are part of each constraint
        for count_i,i in enumerate(constraints):
            if i[0] == "B":            
                constraints[count_i] = [ str(int(j)) for j in i[1:3] ]
            elif i[0] == "A":            
                constraints[count_i] = [ str(int(j)) for j in i[1:4] ]
            elif i[0] == "D":            
                constraints[count_i] = [ str(int(j)) for j in i[1:5] ]

        # Write the constraint file. 
        for count_i,i in enumerate(constraints):
            if count_i == 0:
                f.write("{} {} {}\n".format(" ".join(i),delta,N_steps))
            else:
                f.write("{}\n".format(" ".join(i)))

    return

# Function for reading in run data, ensuring parameters for all bond/angle modes are present (required for the algorithm)
# and performing a self-consistent least-sq fit of the bonds and angles. 
# The function generates the usual plots and data text files. All parameters are saved to FF_dict
def parse_all_bonds_and_angles(base_name,adj_mat,atomtypes,elements,modes_from_FF,save_folder,qc_types=[]):

    # In python all global variables must be declared in each scope
    global FF_dict,qc_type,harm_range

    # Print banner
    print("\n{}".format("*"*167))
    print("* {:^163s} *".format("Self-Consistently Fitting Bonds and Angles Using an Iterative Algorithm"))
    print("{}".format("*"*167))

    # Save parent directory
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_bond_plots_DFT=[]
    paths_to_bond_plots_MP2=[]
    paths_to_angle_plots_DFT=[]
    paths_to_angle_plots_MP2=[]

    # Initialize a list of all Bonds and Angles (Dihedrals and One_fives are not used, they are just included because the function outputs them
    Bonds,Angles,Dihedrals,One_fives = Find_modes(adj_mat,atomtypes,return_all=0)
    all_modes = [ (atomtypes[i[0]],atomtypes[i[1]]) for i in Bonds ] 
    all_modes += [ (atomtypes[i[0]],atomtypes[i[1]],atomtypes[i[2]]) for i in Angles ]
    all_modes = sorted(set(all_modes))
    all_modes_ind = Bonds + Angles

    #################
    # Process modes #
    #################

    # Find bond and angle directories
    os.chdir("angles")
    ba_dirs = [ "angles/"+names for names in os.listdir('.') if os.path.isdir(os.path.join('.',names)) ]        
    os.chdir(working_dir)
    os.chdir("bonds")
    ba_dirs += [ "bonds/"+names for names in os.listdir('.') if os.path.isdir(os.path.join('.',names)) ]        
    os.chdir(working_dir)    

    # Save the list of modes with explicit parametrization data
    param_modes = [ tuple(i.split('/')[-1].split("_")) for i in ba_dirs ]

    # Check that a complete set of parameters exist 
    complete = [0]*len(all_modes)
    for count_i,i in enumerate(all_modes):
        if i in param_modes or i in modes_from_FF:
            complete[count_i] = 1

    # Print diagnostic and exit if incomplete mode data is present for the parse. 
    if 0 in complete:
        print("\nERROR: Missing parametrization data for the following bond(s)/angles(s):\n")
        for count_i,i in enumerate(complete):
            if i == 0:
                print("\t{}".format(all_modes[count_i]))
        print("\nThese types must be explicitly parametrized in this run or supplied in an external database via the -FF command. Exiting...")
        quit()

    #####################################
    ###### Process mode directories #####
    #####################################

    # Initialize dictionary to hold geometries and the lengths/angles of modes in the structure
    Scans = {}
    Geometries = []
    Modes      = []

    # Collect data from the parametrization trajectories
    for modes in ba_dirs:

        # Change into current bond directory
        os.chdir(modes)
        mode_name = modes.split('/')[-1]
        mode_type = tuple(mode_name.split("_"))

        # Initialize sub dictionary within Scans
        Scans[mode_type] = {"geos":[], "modes":[], "active":[]}
        if 'dft' in qc_types: Scans[mode_type]["DFT_Energies"] = []
        if 'mp2' in qc_types: Scans[mode_type]["MP2_Energies"] = []

        # Parse scan data
        Geo_dict,Scans[mode_type]["con"],completion_flag = get_bond_angle_data(working_dir+'/'+save_folder+'/'+base_name+'_'+mode_name+'_geos.xyz')

        # Add geometry and modes to the dictionaries
        for i in natural_sort(list(Geo_dict.keys())):
            Scans[mode_type]["geos"] += [Geo_dict[i]['geo']]
            if 'dft' in qc_types: Scans[mode_type]["DFT_Energies"] += [Geo_dict[i]['DFT']]
            if 'mp2' in qc_types: Scans[mode_type]["MP2_Energies"] += [Geo_dict[i]['MP2']]

        # Parse modes in each configuration
        for i in natural_sort(list(Geo_dict.keys())):
            tmp = []
            for j in all_modes_ind:

                # Parse active mode
                if list(j) == Scans[mode_type]["con"] or list(j)[::-1] == Scans[mode_type]["con"]:

                    # Parse bond length:
                    if len(j) == 2:
                        length = norm(Geo_dict[i]['geo'][j[0]]-Geo_dict[i]['geo'][j[1]])
                        Scans[mode_type]["active"] += [length]

                    # Parse angle degree:
                    if len(j) == 3:
                        atom_1 = Geo_dict[i]["geo"][j[0]]
                        atom_2 = Geo_dict[i]["geo"][j[1]]
                        atom_3 = Geo_dict[i]["geo"][j[2]]
                        theta =  np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))
                        Scans[mode_type]["active"] += [theta]

                # Parse all other modes
                else:
                    # Parse bond length:
                    if len(j) == 2:
                        length = norm(Geo_dict[i]['geo'][j[0]]-Geo_dict[i]['geo'][j[1]])
                        tmp += [((atomtypes[j[0]],atomtypes[j[1]]),length)]

                    # Parse angle degree:
                    if len(j) == 3:
                        atom_1 = Geo_dict[i]["geo"][j[0]]
                        atom_2 = Geo_dict[i]["geo"][j[1]]
                        atom_3 = Geo_dict[i]["geo"][j[2]]
                        theta =  np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))
                        tmp += [((atomtypes[j[0]],atomtypes[j[1]],atomtypes[j[2]]),theta)]
            Scans[mode_type]["modes"] += [tmp]

        # Avoid the parse is the data is incomplete
        if completion_flag == 0:
            os.chdir(working_dir)
            print("******************************************************************************************")
            print("  Warning: Scan jobs did not run to completion for angle {}...Skipping...".format(angles))
            print("******************************************************************************************")
            continue
        os.chdir(working_dir)    

    # Initialize a list of the modes being fit 
    fit_modes = [ i for i in all_modes if i not in modes_from_FF ]
    if len(fit_modes) == 0:
        return fit_modes

    #####################
    # Fit the DFT modes #
    #####################

    if "dft" in qc_types:

        # Initialize entries in the FF_dict 
        for i in fit_modes:
            min_ind = Scans[i]["DFT_Energies"].index(min(Scans[i]["DFT_Energies"])) 
            if len(i) == 2: 
                FF_dict['bonds'][i] = {}
                FF_dict['bonds'][i]['DFT'] = ('harmonic',100.0,Scans[i]["active"][min_ind])
            if len(i) == 3: 
                FF_dict['angles'][i] = {}
                FF_dict['angles'][i]['DFT'] = ('harmonic',100.0,180.0/np.pi*Scans[i]["active"][min_ind])

        # Normalize DFT energies to the minimum value
        for i in fit_modes:
            Scans[i]["DFT_Energies"] = np.array(Scans[i]["DFT_Energies"]) - min(Scans[i]["DFT_Energies"])

        # Fit loop
        min_cycles = 10
        max_cycles = 1000
        delta_xhi2_thresh = 1E-6
        cycle_count =-1
        xhi2_previous = 0.0
        C = 1.0
        while (1):

            # Increment the cycle count and print diagnostic
            cycle_count += 1
            print("\ncycle {}:\n".format(cycle_count))

            # Each cycle consists of looping over all pair-types being fit
            # To eliminate ordering bias, the list of types is shuffled every cycle.
            random.shuffle(fit_modes)
            qc_type = 'DFT'

            for i in fit_modes:
                print("\tFitting {}...".format(i))

                # Calculate and normalize fit potential
                Scans[i]["DFT_Fit"] = Scans[i]["DFT_Energies"] - bond_angle_energy(Scans[i]["modes"])
                Scans[i]["DFT_Fit"] = Scans[i]["DFT_Fit"] - min(Scans[i]["DFT_Fit"])
                
                # Define harmonic fit range
                harm_range = [min(Scans[i]["active"]),max(Scans[i]["active"])]
                harm_range = [harm_range[0]-(harm_range[1]-harm_range[0])*0.1,harm_range[1]+(harm_range[1]-harm_range[0])*0.1]

                # Perform fit
                if len(i) == 2:
                    params, errors = curve_fit(harmonic, Scans[i]["active"], Scans[i]["DFT_Fit"],p0=[FF_dict['bonds'][i]['DFT'][1],FF_dict['bonds'][i]['DFT'][2]])
                    # debug diagnostics
                    #print "{:12s}  {:12s} {:12s} {:12s} {:12s} {:12s}".format("eq_val","fit_pot","FF_pot_old","FF_pot_new","DFT_pot","hand")
                    #for count_j,j in enumerate(Scans[i]["DFT_Fit"]):
                    #    print "{:<12.6f}: {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f}".format(Scans[i]["active"][count_j],j,harmonic(Scans[i]["active"][count_j],FF_dict['bonds'][i]['DFT'][1],FF_dict['bonds'][i]['DFT'][2]),\
                    #                                                               harmonic(Scans[i]["active"][count_j],params[0],params[1]),Scans[i]["DFT_Energies"][count_j],params[0]*(Scans[i]["active"][count_j]-params[1])**(2.0))
                    #print "{} {} --> {} {}".format(FF_dict['bonds'][i]['DFT'][1],FF_dict['bonds'][i]['DFT'][2],params[0],params[1])
                    #print "harm_range: {}".format(harm_range)
                elif len(i) == 3:
                    params, errors = curve_fit(harmonic, Scans[i]["active"], Scans[i]["DFT_Fit"],p0=[FF_dict['angles'][i]['DFT'][1],np.pi/180.0*FF_dict['angles'][i]['DFT'][2]])
                    # debug diagnostics
                    #print "{:12s}  {:12s} {:12s} {:12s} {:12s} {:12s}".format("eq_val","fit_pot","FF_pot_old","FF_pot_new","DFT_pot","hand")
                    #for count_j,j in enumerate(Scans[i]["DFT_Fit"]):
                    #    print "{:<12.6f}: {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f}".format(Scans[i]["active"][count_j],j,harmonic(Scans[i]["active"][count_j],FF_dict['angles'][i]['DFT'][1],np.pi/180.0*FF_dict['angles'][i]['DFT'][2]),\
                    #                                                           harmonic(Scans[i]["active"][count_j],params[0],params[1]),Scans[i]["DFT_Energies"][count_j],params[0]*(Scans[i]["active"][count_j]-params[1])**(2.0))
                    #print "{} {} --> {} {}".format(FF_dict['angles'][i]['DFT'][1],FF_dict['angles'][i]['DFT'][2]*np.pi/180.0,params[0],params[1])
                    #print "harm_range: {}".format(harm_range)

                # Update parameters 
                if len(i) == 2: FF_dict['bonds'][i]['DFT'] = ('harmonic',params[0],params[1])
                if len(i) == 3: FF_dict['angles'][i]['DFT'] = ('harmonic',params[0],180.0/np.pi*params[1])

            # Calculate Xhi2 (after all new parameters are in place)
            xhi2_current = 0.0
            for i in fit_modes:
                harm_range = [min(Scans[i]["active"]),max(Scans[i]["active"])]
                harm_range = [harm_range[0]-(harm_range[1]-harm_range[0])*0.1,harm_range[1]+(harm_range[1]-harm_range[0])*0.1]
                Scans[i]["DFT_Fit"] = Scans[i]["DFT_Energies"] - bond_angle_energy(Scans[i]["modes"])
                Scans[i]["DFT_Fit"] = Scans[i]["DFT_Fit"] - min(Scans[i]["DFT_Fit"])
                if len(i) == 2: xhi2_current += np.mean((Scans[i]["DFT_Fit"] - harmonic(Scans[i]["active"],FF_dict['bonds'][i]['DFT'][1],FF_dict['bonds'][i]['DFT'][2]))**(2.0))
                if len(i) == 3: xhi2_current += np.mean((Scans[i]["DFT_Fit"] - harmonic(Scans[i]["active"],FF_dict['angles'][i]['DFT'][1],np.pi/180.0*FF_dict['angles'][i]['DFT'][2]))**(2.0))

            # Calculate xhi2 and delta_xhi2 for this cycle (NOTE: first cycle has no delta_xhi2)
            xhi2_current = xhi2_current / float(len(fit_modes))
            if cycle_count == 0:
                delta_xhi2 = 0.0
            else:
                delta_xhi2   = xhi2_previous - xhi2_current

            # Print diagnostic
            print("\n\t{:<20s} {:<12.8e} kcal^2/mol^2".format("np.mean(xhi2):",xhi2_current))
            print("\t{:<20s} {:<12.8e} kcal^2/mol^2".format("np.mean(delta_xhi2):",delta_xhi2))

            # Check break conditions
            if cycle_count >= min_cycles: 
                if cycle_count == max_cycles:
                    print("\nMaximum number of cycles ({}) reached, breaking iterative fit loop...\n".format(max_cycles))
                    break
                elif np.abs(delta_xhi2) < delta_xhi2_thresh:
                    print("\nDelta Xhi2 convergence achieved (np.abs({}) kcal^2/mol^2 < {} kcal^2/mol^2), breaking iterative fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
                    break

            # Update xhi2_previous
            xhi2_previous = xhi2_current

        os.chdir(working_dir)    

    #####################
    # Fit the MP2 modes #
    #####################

    if 'mp2' in qc_types:

        # Initialize entries in the FF_dict 
        for i in fit_modes:
            min_ind = Scans[i]["MP2_Energies"].index(min(Scans[i]["MP2_Energies"])) 
            if len(i) == 2: 
                FF_dict['bonds'][i] = {}
                FF_dict['bonds'][i]['MP2'] = ('harmonic',100.0,Scans[i]["active"][min_ind])
            if len(i) == 3: 
                FF_dict['angles'][i] = {}
                FF_dict['angles'][i]['MP2'] = ('harmonic',100.0,180.0/np.pi*Scans[i]["active"][min_ind])

        # Normalize MP2 energies to the minimum value
        for i in fit_modes:
            Scans[i]["MP2_Energies"] = np.array(Scans[i]["MP2_Energies"]) - min(Scans[i]["MP2_Energies"])

        # Fit loop
        min_cycles = 10
        max_cycles = 1000
        delta_xhi2_thresh = 1E-6
        cycle_count =-1
        xhi2_previous = 0.0
        C = 1.0
        while (1):

            # Increment the cycle count and print diagnostic
            cycle_count += 1
            print("\ncycle {}:\n".format(cycle_count))

            # Each cycle consists of looping over all pair-types being fit
            # To eliminate ordering bias, the list of types is shuffled every cycle.
            random.shuffle(fit_modes)
            qc_type = 'MP2'

            for i in fit_modes:
                print("\tFitting {}...".format(i))

                # Calculate and normalize fit potential
                Scans[i]["MP2_Fit"] = Scans[i]["MP2_Energies"] - bond_angle_energy(Scans[i]["modes"])
                Scans[i]["MP2_Fit"] = Scans[i]["MP2_Fit"] - min(Scans[i]["MP2_Fit"])
                
                # Define harmonic fit range
                harm_range = [min(Scans[i]["active"]),max(Scans[i]["active"])]
                harm_range = [harm_range[0]-(harm_range[1]-harm_range[0])*0.1,harm_range[1]+(harm_range[1]-harm_range[0])*0.1]

                # Perform fit
                if len(i) == 2:
                    params, errors = curve_fit(harmonic, Scans[i]["active"], Scans[i]["MP2_Fit"],p0=[FF_dict['bonds'][i]['MP2'][1],FF_dict['bonds'][i]['MP2'][2]])
                    # debug diagnostics
                    #print "{:12s}  {:12s} {:12s} {:12s} {:12s} {:12s}".format("eq_val","fit_pot","FF_pot_old","FF_pot_new","MP2_pot","hand")
                    #for count_j,j in enumerate(Scans[i]["MP2_Fit"]):
                    #    print "{:<12.6f}: {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f}".format(Scans[i]["active"][count_j],j,harmonic(Scans[i]["active"][count_j],FF_dict['bonds'][i]['MP2'][1],FF_dict['bonds'][i]['MP2'][2]),\
                    #                                                               harmonic(Scans[i]["active"][count_j],params[0],params[1]),Scans[i]["MP2_Energies"][count_j],params[0]*(Scans[i]["active"][count_j]-params[1])**(2.0))
                    #print "{} {} --> {} {}".format(FF_dict['bonds'][i]['MP2'][1],FF_dict['bonds'][i]['MP2'][2],params[0],params[1])
                    #print "harm_range: {}".format(harm_range)
                elif len(i) == 3:
                    params, errors = curve_fit(harmonic, Scans[i]["active"], Scans[i]["MP2_Fit"],p0=[FF_dict['angles'][i]['MP2'][1],np.pi/180.0*FF_dict['angles'][i]['MP2'][2]])
                    # debug diagnostics
                    #print "{:12s}  {:12s} {:12s} {:12s} {:12s} {:12s}".format("eq_val","fit_pot","FF_pot_old","FF_pot_new","MP2_pot","hand")
                    #for count_j,j in enumerate(Scans[i]["MP2_Fit"]):
                    #    print "{:<12.6f}: {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f} {:<12.6f}".format(Scans[i]["active"][count_j],j,harmonic(Scans[i]["active"][count_j],FF_dict['angles'][i]['MP2'][1],np.pi/180.0*FF_dict['angles'][i]['MP2'][2]),\
                    #                                                           harmonic(Scans[i]["active"][count_j],params[0],params[1]),Scans[i]["MP2_Energies"][count_j],params[0]*(Scans[i]["active"][count_j]-params[1])**(2.0))
                    #print "{} {} --> {} {}".format(FF_dict['angles'][i]['MP2'][1],FF_dict['angles'][i]['MP2'][2]*np.pi/180.0,params[0],params[1])
                    #print "harm_range: {}".format(harm_range)

                # Update parameters and add to xhi2 value
                if len(i) == 2: FF_dict['bonds'][i]['MP2'] = ('harmonic',params[0],params[1])
                if len(i) == 3: FF_dict['angles'][i]['MP2'] = ('harmonic',params[0],180.0/np.pi*params[1])

            # Calculate Xhi2 (after all new parameters are in place)
            xhi2_current = 0.0
            for i in fit_modes:
                harm_range = [min(Scans[i]["active"]),max(Scans[i]["active"])]
                harm_range = [harm_range[0]-(harm_range[1]-harm_range[0])*0.1,harm_range[1]+(harm_range[1]-harm_range[0])*0.1]
                Scans[i]["MP2_Fit"] = Scans[i]["MP2_Energies"] - bond_angle_energy(Scans[i]["modes"])
                Scans[i]["MP2_Fit"] = Scans[i]["MP2_Fit"] - min(Scans[i]["MP2_Fit"])
                if len(i) == 2: xhi2_current += np.mean((Scans[i]["MP2_Fit"] - harmonic(Scans[i]["active"],FF_dict['bonds'][i]['MP2'][1],FF_dict['bonds'][i]['MP2'][2]))**(2.0))
                if len(i) == 3: xhi2_current += np.mean((Scans[i]["MP2_Fit"] - harmonic(Scans[i]["active"],FF_dict['angles'][i]['MP2'][1],np.pi/180.0*FF_dict['angles'][i]['MP2'][2]))**(2.0))

            # Calculate xhi2 and delta_xhi2 for this cycle (NOTE: first cycle has no delta_xhi2)
            xhi2_current = xhi2_current / float(len(fit_modes))
            if cycle_count == 0:
                delta_xhi2 = 0.0
            else:
                delta_xhi2   = xhi2_previous - xhi2_current

            # Print diagnostic
            print("\n\t{:<20s} {:<12.8e} kcal^2/mol^2".format("np.mean(xhi2):",xhi2_current))
            print("\t{:<20s} {:<12.8e} kcal^2/mol^2".format("np.mean(delta_xhi2):",delta_xhi2))

            # Check break conditions
            if cycle_count >= min_cycles: 
                if cycle_count == max_cycles:
                    print("\nMaximum number of cycles ({}) reached, breaking iterative fit loop...\n".format(max_cycles))
                    break
                elif np.abs(delta_xhi2) < delta_xhi2_thresh:
                    print("\nDelta Xhi2 convergence achieved (np.abs({}) kcal^2/mol^2 < {} kcal^2/mol^2), breaking iterative fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
                    break

            # Update xhi2_previous
            xhi2_previous = xhi2_current

        os.chdir(working_dir)    

    return fit_modes

# Used for self-consistent fitting of bonds and angles
# k is the harmonic force-constant
# x0 is the harmonic equilibrium displacement
# C is the normalization constant for the energies
def bond_angle_energy(Modes):

    global FF_dict,qc_type
    Energies = np.zeros(len(Modes))

    # Iterate over the mode lists for each configuration
    for count_i,i in enumerate(Modes):

        # Iterate over the individual modes in this configurations
        for j in i:
            if len(j[0]) == 2: Energies[count_i] += FF_dict['bonds'][j[0]][qc_type][1]*(j[1]-FF_dict['bonds'][j[0]][qc_type][2])**(2.0)
            elif len(j[0]) == 3: Energies[count_i] += FF_dict['angles'][j[0]][qc_type][1]*(j[1]-np.pi/180.0*FF_dict['angles'][j[0]][qc_type][2])**(2.0)
    return Energies
            
# This is a messy wrapper for the dihedral fitting procedure. In brief, the 
# function looks up all the dihedral data, calculates opls fits, and
# generates some plots and data text files. All parameters are saved to FF_dict
def parse_harmonic_dihedrals(base_name,list_of_dihedral_types,dihedral_folders,dihedral_atomtypes,charge_dict,modes_from_ff,save_folder,fit_type='AA',delta_xhi2_thresh=0.000001,min_cycles=10,\
                            max_cycles=100000,qc_types=[],corr_list=[]):

    # In python all global variables must be declared in each scope
    global FF_dict,QC_type,Dihedral_Harmonic_Data,dihedral_scan

    # Initialize dictionary for mapping atomic number to element
    atom_to_element = { 1:'H' ,  2:'He',\
                        3:'Li',  4:'Be',  5:'B' ,  6:'C' ,  7:'N' ,  8:'O' ,  9:'F' , 10:'Ne',\
                       11:'Na', 12:'Mg', 13:'Al', 14:'Si', 15:'P' , 16:'S' , 17:'Cl', 18:'Ar',\
                       19:'K' , 20:'Ca', 21:'Sc', 22:'Ti', 23:'V' , 24:'Cr', 25:'Mn', 26:'Fe', 27:'Co', 28:'Ni', 29:'Cu', 30:'Zn', 31:'Ga', 32:'Ge', 33:'As', 34:'Se', 35:'Br', 36:'Kr',\
                       37:'Rb', 38:'Sr', 39:'Y' , 40:'Zr', 41:'Nb', 42:'Mo', 43:'Tc', 44:'Ru', 45:'Rh', 46:'Pd', 47:'Ag', 48:'Cd', 49:'In', 50:'Sn', 51:'Sb', 52:'Te', 53:'I' , 54:'Xe',\
                       55:'Cs', 56:'Ba', 57:'La', 72:'Hf', 73:'Ta', 74:'W' , 75:'Re', 76:'Os', 77:'Ir', 78:'Pt', 79:'Au', 80:'Hg', 81:'Tl', 82:'Pb', 83:'Bi', 84:'Po', 85:'At', 86:'Rn'}    

    # Save parent directory
    working_dir = os.getcwd()

    # Initialize path list for DFT and MP2 data
    paths_to_dihedral_fitpot_plots_DFT=[]
    paths_to_dihedral_fitpot_plots_MP2=[]
    paths_to_dihedral_totpot_plots_DFT=[]
    paths_to_dihedral_totpot_plots_MP2=[]

    # List of Processed Dihedrals
    processed_dihedrals = []

    # Grab the dihedral data (returns a dictionary holding dihedral_type,angle tuples for each configuration and the DFT and MP2 energies
    Dihedral_Harmonic_Data = get_harmonic_dihedral_data(list_of_dihedral_types,dihedral_folders,fit_type,dihedral_atomtypes,charge_dict,qc_types=qc_types,corr_list=[])
    keys = list(Dihedral_Harmonic_Data.keys())

    # Return to working directory and save the list of catenated geometries
    for i in list(Dihedral_Harmonic_Data.keys()):

        # Check if the MD subfolder needs to be created
        if fit_type == 'UA': savename = base_name+"/dihedrals_harmonic/"+"_".join([ k+"-UA" for k in i ])+"_geos.xyz"
        else: savename = base_name+"/dihedrals_harmonic/"+"_".join([ k for k in i ])+"_geos.xyz"
        f = open(savename, 'w')
    
        # Loop over all discovered geometries and save their geometries to Geo_dict and
        # catenate to a single file for viewing.
        for j in natural_sort(list(Dihedral_Harmonic_Data[i].keys())):
            f.write("{:d}\n\n".format(len(Dihedral_Harmonic_Data[i][j]["Geo"])))
            for count_k,k in enumerate(Dihedral_Harmonic_Data[i][j]["Geo"]):
                f.write(" {:<10s} {:< 15.8f} {:< 15.8f} {:< 15.8f}\n".format(Dihedral_Harmonic_Data[i][j]["Elements"][count_k],k[0],k[1],k[2]))
        f.close()

    # dihedral types are added to already_fit_list as they are fit. Fitting is skipped for any dihedrals in this list, in order
    # to avoid redundant fits when the same dihedral happens to be in two scans
    already_fit_list = modes_from_ff

    # Iterate over each scan's data
    for dihedral_count,dihedrals in enumerate(list_of_dihedral_types):

        # Set dihedral_scan to the current folder (used in iterative_fit_all)
        dihedral_scan = dihedrals

        # Check if any dihedral types are fitting in this scan (for UA some hydrogen-only containing scans are skipped)
        dihedral_types = sorted(set([ j[0] for i in list(Dihedral_Harmonic_Data[dihedrals].keys()) for j in Dihedral_Harmonic_Data[dihedrals][i]["Dihedrals"] ]))        
        dihedral_types = [ i for i in dihedral_types if i not in already_fit_list ]
        already_fit_list += dihedral_types
        if len(dihedral_types) == 0: continue

        # Print banner
        if fit_type == "AA": print("\n{}\n* {:^163s} *\n{}".format("*"*167,"Parsing AA-harmonic_dihedral Types from Scan {}".format(dihedrals),"*"*167))
        if fit_type == "UA": print("\n{}\n* {:^163s} *\n{}".format("*"*167,"Parsing UA-harmonic_dihedral Types from Scan {}".format(dihedrals),"*"*167))

        # Print diagnostic
        print("\nHarmonic dihedral types being fit:\n")
        for i in dihedral_types:
            print("\t{}".format(i))

        # Collect list of energies and angles
        DFT_E = []
        MP2_E = []
        x_vals = []
        for i in list(Dihedral_Harmonic_Data[dihedrals].keys()):
            x_vals += [i] 
            if 'dft' in qc_types:
                DFT_E += [Dihedral_Harmonic_Data[dihedrals][i]["DFT_Fit"]]
            if 'mp2' in qc_types:
                MP2_E += [Dihedral_Harmonic_Data[dihedrals][i]["MP2_Fit"]]

        # Calculate parameters (params holds force constant and theta_0). The algorithm first performs a 
        # best fit using only the first dihedral type. Then the fitted potential energy is subtracted off the 
        # real potential, and the remainder is fit to the first coincident dihedral type, and so on. The algorithm
        # is flexible enough to handle fitting multiple instances of the same dihedral type. 
        if 'dft' in qc_types:
            print("\nIteratively fitting harmonic potentials to all DFT-D dihedral configurations in the scan to acheive self-consistency...")
            QC_type = 'DFT'
            fit_vals,DFT_RMSD = iterative_harmonic_fit_all(dihedrals,x_vals,DFT_E,dihedral_types,weight=1.0,fit='harmonic',delta_xhi2_thresh=delta_xhi2_thresh,min_cycles=min_cycles,max_cycles=max_cycles)
        if 'mp2' in qc_types:
            print("\nIteratively fitting harmonic potentials to all MP2 dihedral configurations in the scan to acheive self-consistency...")        
            QC_type = 'MP2'
            fit_vals,MP2_RMSD = iterative_harmonic_fit_all(dihedrals,x_vals,MP2_E,dihedral_types,weight=1.0,fit='harmonic',delta_xhi2_thresh=delta_xhi2_thresh,min_cycles=min_cycles,max_cycles=max_cycles)

        if 'dft' in qc_types:
            print("\nAverage Absolute Error in DFT fit: {} kcal/mol".format(DFT_RMSD))
        if 'mp2' in qc_types:
            print("Average Absolute Error in MP2 fit: {} kcal/mol".format(MP2_RMSD))
        
        # Calculate DFT Fit Potential
        if 'dft' in qc_types:
            for i in list(Dihedral_Harmonic_Data[dihedrals].keys()):            
                Dihedral_Harmonic_Data[dihedrals][i]["E_DFT_Pot"] = dihedral_harmonic(Dihedral_Harmonic_Data[dihedrals][i]["Dihedrals"],QC='DFT')

        # Calculate MP2 Fit Potential
        if 'mp2' in qc_types:
            for i in list(Dihedral_Harmonic_Data[dihedrals].keys()):
                Dihedral_Harmonic_Data[dihedrals][i]["E_MP2_Pot"] = dihedral_harmonic(Dihedral_Harmonic_Data[dihedrals][i]["Dihedrals"],QC='MP2')

        # ##############################
        # # Write catenated geometries #
        # ##############################

        # # Open file and initialize dictionary to hold catenated geometries 
        # f = open(working_dir+'/'+save_folder+'/'+base_name+'_'+dihedrals+'_geos.xyz', 'w')

        # # Loop over all discovered geometries and save their geometries to Geo_dict and
        # # catenate to a single file for viewing.
        # for g in range(len(Dihedral_Harmonic_Data[dihedrals].keys())):
        #     f.write("{:d}\n\n".format(len(Dihedral_Harmonic_Data[dihedrals][str(g)]["Geo"])))
        #     for count_i,i in enumerate(Dihedral_Harmonic_Data[dihedrals][str(g)]["Geo"]):
        #         f.write(" {:<10s} {:< 15.8f} {:< 15.8f} {:< 15.8f}\n".format(Dihedral_Harmonic_Data[dihedrals][str(g)]["Elements"][count_i],i[0],i[1],i[2]))
        # f.close()

        #### Generate 2D representation of the molecule
        print("Generating 2D representation of the fit fragment...")
        atomtypes = Dihedral_Harmonic_Data[dihedrals]["0"]["Atomtypes"]
        adj_mat = Table_generator(Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"],Dihedral_Harmonic_Data[dihedrals]["0"]["Geo"])
        geo_2D = kekule(Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"],atomtypes,Dihedral_Harmonic_Data[dihedrals]["0"]["Geo"],adj_mat)
        
        #######################################
        ######### Generate Fit Plots ##########
        #######################################
        
        # Initialize a list of all unique dihedral types in this scan
        dihedral_types = sorted(set([ i[0] for i in Dihedral_Harmonic_Data[dihedrals]["0"]["Dihedrals"] ]))
        
        # Generate fit plots for each dihedral_type that was fit in this scan
        print("\nSaving figures of harmonic dihedral fits...")
        for count_i,i in enumerate(dihedral_types):

            # Parse list of atom ids from one of the scan configurations
            idx = [ count_j for count_j,j in enumerate(Dihedral_Harmonic_Data[dihedrals]["0"]["Dihedrals"]) if i == j[0] ]
            dihedral_atoms = Dihedral_Harmonic_Data[dihedrals]["0"]["Dihedral_Atoms"][idx[0]][0]

            # Generate Label
            label = atom_to_element[int(i[0].split('[')[1].split(']')[0])]+'-'+\
                    atom_to_element[int(i[1].split('[')[1].split(']')[0])]+'-'+\
                    atom_to_element[int(i[2].split('[')[1].split(']')[0])]+'-'+\
                    atom_to_element[int(i[3].split('[')[1].split(']')[0])]

            # Generate savename
            savename = i[0]+'_'+i[1]+'_'+i[2]+'_'+i[3]

            # Collect fit data for this type
            angles    = []
            E_DFT_Fit = []
            E_MP2_Fit = []
            E_DFT_Pot = []
            E_MP2_Pot = []
            E_DFT_Fit_Tot = []
            E_MP2_Fit_Tot = []
            E_DFT_Pot_Tot = []
            E_MP2_Pot_Tot = []
            for j in range(len(list(Dihedral_Harmonic_Data[dihedrals].keys()))):
                E_Coul         = Dihedral_Harmonic_Data[dihedrals][str(j)]["E_coul"]
                if 'dft' in qc_types:
                    E_bond_DFT     = Dihedral_Harmonic_Data[dihedrals][str(j)]["E_bond_DFT"]
                    E_angle_DFT    = Dihedral_Harmonic_Data[dihedrals][str(j)]["E_angle_DFT"]
                    E_DFT_Fit     += [Dihedral_Harmonic_Data[dihedrals][str(j)]["DFT_Fit"]]
                    E_DFT_Pot     += [Dihedral_Harmonic_Data[dihedrals][str(j)]["E_DFT_Pot"]]
                    E_DFT_Fit_Tot += [Dihedral_Harmonic_Data[dihedrals][str(j)]["DFT_Fit"]+E_Coul+E_bond_DFT+E_angle_DFT]
                    E_DFT_Pot_Tot += [Dihedral_Harmonic_Data[dihedrals][str(j)]["E_DFT_Pot"]+E_Coul+E_bond_DFT+E_angle_DFT]
                if 'mp2' in qc_types:
                    E_bond_MP2     = Dihedral_Harmonic_Data[dihedrals][str(j)]["E_bond_MP2"]
                    E_angle_MP2    = Dihedral_Harmonic_Data[dihedrals][str(j)]["E_angle_MP2"]
                    E_MP2_Fit     += [Dihedral_Harmonic_Data[dihedrals][str(j)]["MP2_Fit"]]
                    E_MP2_Pot     += [Dihedral_Harmonic_Data[dihedrals][str(j)]["E_MP2_Pot"]]
                    E_MP2_Fit_Tot += [Dihedral_Harmonic_Data[dihedrals][str(j)]["MP2_Fit"]+E_Coul+E_bond_MP2+E_angle_MP2]
                    E_MP2_Pot_Tot += [Dihedral_Harmonic_Data[dihedrals][str(j)]["E_MP2_Pot"]+E_Coul+E_bond_MP2+E_angle_MP2]
                for k in Dihedral_Harmonic_Data[dihedrals][str(j)]["Dihedrals"]:
                    if k[0] == i: 
                        if k[1] <-np.pi/2.0:
                            angles += [k[1]+2.0*np.pi]
                        else:
                            angles += [k[1]]
                        break

            # Sort the data based on the angles
            ind = sorted(list(range(len(angles))), key=lambda k: angles[k])
            if 'dft' in qc_types:
                E_DFT_Fit     = np.array([ E_DFT_Fit[j] for j in ind ])
                E_DFT_Pot     = np.array([ E_DFT_Pot[j] for j in ind ])
                E_DFT_Fit_Tot = np.array([ E_DFT_Fit_Tot[j] for j in ind ])
                E_DFT_Pot_Tot = np.array([ E_DFT_Pot_Tot[j] for j in ind ])
            if 'mp2' in qc_types:
                E_MP2_Fit     = np.array([ E_MP2_Fit[j] for j in ind ])
                E_MP2_Pot     = np.array([ E_MP2_Pot[j] for j in ind ])
                E_MP2_Fit_Tot = np.array([ E_MP2_Fit_Tot[j] for j in ind ])
                E_MP2_Pot_Tot = np.array([ E_MP2_Pot_Tot[j] for j in ind ])
            angles        = np.array([ angles[j] for j in ind ])

            # Normalize *_Tot arrays by the minimum value
            if 'dft' in qc_types:
                E_DFT_Fit_Tot =  E_DFT_Fit_Tot - min(E_DFT_Fit_Tot)
                E_DFT_Pot_Tot =  E_DFT_Pot_Tot - min(E_DFT_Pot_Tot)
            if 'mp2' in qc_types:
                E_MP2_Fit_Tot =  E_MP2_Fit_Tot - min(E_MP2_Fit_Tot)
                E_MP2_Pot_Tot =  E_MP2_Pot_Tot - min(E_MP2_Pot_Tot)
            
            # Plot DFT comparisons
            if 'dft' in qc_types:

                # Initialize figure for the total fit potential plus molecule
                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                if draw_scale > 1.0: draw_scale = 1.0

                # Draw molecule in ax1 slot
                # Add bonds
                linewidth = 2.0*draw_scale*5.0
                if linewidth > 2.0:
                    linewidth = 2.0
                if linewidth < 1.0:
                    linewidth = 1.0
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[0]][0],geo_2D[dihedral_atoms[1]][0]],[geo_2D[dihedral_atoms[0]][1],geo_2D[dihedral_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[1]][0],geo_2D[dihedral_atoms[2]][0]],[geo_2D[dihedral_atoms[1]][1],geo_2D[dihedral_atoms[2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[2]][0],geo_2D[dihedral_atoms[3]][0]],[geo_2D[dihedral_atoms[2]][1],geo_2D[dihedral_atoms[3]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                # Add atom labels
                fontsize = 20*draw_scale*5.0
                if fontsize > 24:
                    fontsize = 24
                if fontsize < 9:
                    fontsize = 9            
                for count_i,i in enumerate(geo_2D):
                    if count_i in dihedral_atoms:
                        ax1.text(i[0], i[1], Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"][count_i], style='normal', color=(0.0,0.1,0.8),\
                                 fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1], Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"][count_i], style='normal', color=(0.0,0.0,0.0),\
                                 fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Plot the exact potential determined from DFT calculations
                ax2.scatter(angles*180.0/np.pi,E_DFT_Fit,s=120,color=(0.05,0.05,0.05,1))
                ax2.scatter(angles*180.0/np.pi,E_DFT_Pot,s=80,color=(0.0,0.1,0.8,1.0))
                ax2.plot(angles*180.0/np.pi,E_DFT_Pot,'--',linewidth=3,color=(0.0,0.1,0.8,1.0))
                maxval = max(list(E_DFT_Fit)+list(E_DFT_Pot))
                om = float(np.floor(log10(maxval)))

                # Format the plot            
                ax2.set_title('{} Harmonic Dihedral {}'.format(base_name,label),fontsize=12,y=1.08)
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(angles)*180.0/np.pi, max(angles)*180.0/np.pi])
                ax2.set_ylim([-10.0**(om-1.0), maxval+10.0**(om-1)])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = working_dir+'/'+save_folder+'/'+savename+'-DFT_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                paths_to_dihedral_fitpot_plots_DFT+=[Name]
                plt.close(fig)

                # Initialize figure for the total fit potential plus molecule
                fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                if draw_scale > 1.0: draw_scale = 1.0

                # Draw molecule in ax1 slot
                # Add bonds
                linewidth = 2.0*draw_scale*5.0
                if linewidth > 2.0:
                    linewidth = 2.0
                if linewidth < 1.0:
                    linewidth = 1.0
                for count_i,i in enumerate(adj_mat):
                    for count_j,j in enumerate(i):
                        if count_j > count_i:
                            if j == 1:
                                ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_i][0],geo_2D[count_j][0]],[geo_2D[count_i][1],geo_2D[count_j][1]],color=(0,0,0),linewidth=linewidth))

                # Add Highlight
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[0]][0],geo_2D[dihedral_atoms[1]][0]],[geo_2D[dihedral_atoms[0]][1],geo_2D[dihedral_atoms[1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[1]][0],geo_2D[dihedral_atoms[2]][0]],[geo_2D[dihedral_atoms[1]][1],geo_2D[dihedral_atoms[2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                ax1.add_line(matplotlib.lines.Line2D([geo_2D[dihedral_atoms[2]][0],geo_2D[dihedral_atoms[3]][0]],[geo_2D[dihedral_atoms[2]][1],geo_2D[dihedral_atoms[3]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                # Add atom labels
                fontsize = 20*draw_scale*5.0
                if fontsize > 24:
                    fontsize = 24
                if fontsize < 9:
                    fontsize = 9            
                for count_i,i in enumerate(geo_2D):
                    if count_i in dihedral_atoms:
                        ax1.text(i[0], i[1], Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"][count_i], style='normal', color=(0.0,0.1,0.8),\
                                 fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                    else:
                        ax1.text(i[0], i[1], Dihedral_Harmonic_Data[dihedrals]["0"]["Elements"][count_i], style='normal', color=(0.0,0.0,0.0),\
                                 fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                ax1.axis('image')
                ax1.axis('off')

                # Add data
                ax2.scatter(np.array(angles)*180.0/np.pi,E_DFT_Fit_Tot,s=120,color=(0.05,0.05,0.05,1.0))
                ax2.plot(angles*180.0/np.pi,E_DFT_Pot_Tot,'--',linewidth=3,color=(0.0,0.1,0.8,1.0))
                ax2.scatter(np.array(angles)*180.0/np.pi,E_DFT_Pot_Tot,s=120,color=(0.0,0.1,0.8,1.0))
                ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax2.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
                ax2.set_xlim([min(angles)*180.0/np.pi, max(angles)*180.0/np.pi])
                ax2.set_ylim([-10.0**(om-1.0), maxval+10.0**(om-1)])
                ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                y_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                x_formatter = matplotlib.ticker.ScalarFormatter(useOffset=False)
                ax2.yaxis.set_major_formatter(y_formatter)
                ax2.xaxis.set_major_formatter(x_formatter)
                [j.set_linewidth(3) for j in ax2.spines.values()]
                plt.tight_layout()

                Name = working_dir+'/'+save_folder+'/'+savename+'-DFT_tot.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                plt.savefig(working_dir+'/'+save_folder+'/'+savename+'-DFT_tot.pdf', bbox_inches=0,dpi=300)
                paths_to_dihedral_totpot_plots_DFT+=[Name]
                plt.close(fig)

            # Plot MP2 comparison with the fit potential
            if 'mp2' in qc_types:
                fig = plt.figure(figsize=(6,5))
                ax = plt.subplot(111)

                # Plot the exact potential determined from MP2 calculations
                ax.scatter(angles*180.0/np.pi,E_MP2_Fit,s=120,color=(0.05,0.05,0.05,1))
                ax.scatter(angles*180.0/np.pi,E_MP2_Pot,s=80,color=(0.0,0.1,0.8,1.0))
                ax.plot(angles*180.0/np.pi,E_MP2_Pot,'--',linewidth=3,color=(0.0,0.1,0.8,1.0))

                # Format the plot
                ax.set_title('{} Harmonic Dihedral {}'.format(base_name,label),fontsize=12,y=1.08)
                ax.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
                ax.set_xlim([round(min(angles)*180.0/np.pi), round(max(angles)*180.0/np.pi)])
                ax.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                [j.set_linewidth(3) for j in ax.spines.values()]
                plt.tight_layout()

                Name = working_dir+'/'+save_folder+'/'+savename+'-MP2_fit.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                paths_to_dihedral_fitpot_plots_MP2+=[Name]
                plt.close(fig)

                # Plot MP2 comparison with the total potential
                fig = plt.figure(figsize=(6,5))
                ax = plt.subplot(111)

                # Plot the exact potential determined from MP2 calculations
                ax.scatter(angles*180.0/np.pi,E_MP2_Fit_Tot,s=120,color=(0.05,0.05,0.05,1))
                ax.scatter(angles*180.0/np.pi,E_MP2_Pot_Tot,s=80,color=(0.0,0.1,0.8,1.0))
                ax.plot(angles*180.0/np.pi,E_MP2_Pot_Tot,'--',linewidth=3,color=(0.0,0.1,0.8,1.0))

                # Format the plot
                ax.set_title('{} Harmonic Dihedral {}'.format(base_name,label),fontsize=12,y=1.08)
                ax.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                ax.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
                ax.set_xlim([round(min(angles)*180.0/np.pi), round(max(angles)*180.0/np.pi)])
                ax.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                ax.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                [j.set_linewidth(3) for j in ax.spines.values()]
                plt.tight_layout()

                Name = working_dir+'/'+save_folder+'/'+savename+'-MP2_tot.png'
                plt.savefig(Name, bbox_inches=0,dpi=300)
                paths_to_dihedral_totpot_plots_MP2+=[Name]
                plt.close(fig)

        # Update the force-field entries
        for i in dihedral_types:
            if 'dft' in qc_types:
                if FF_dict["dihedrals_harmonic"][i]['DFT'][0] == "quadratic":
                    FF_dict["dihedrals_harmonic"][i]['DFT'] = (FF_dict["dihedrals_harmonic"][i]['DFT'][0],FF_dict["dihedrals_harmonic"][i]['DFT'][1],FF_dict["dihedrals_harmonic"][i]['DFT'][2]*180.0/np.pi)
                if FF_dict["dihedrals_harmonic"][i]['DFT'][0] == "harmonic":
                    FF_dict["dihedrals_harmonic"][i]['DFT'] = (FF_dict["dihedrals_harmonic"][i]['DFT'][0],FF_dict["dihedrals_harmonic"][i]['DFT'][1],int(FF_dict["dihedrals_harmonic"][i]['DFT'][2]))
            if 'mp2' in qc_types:
                if FF_dict["dihedrals_harmonic"][i]['MP2'][0] == "quadratic":
                    FF_dict["dihedrals_harmonic"][i]['MP2'] = (FF_dict["dihedrals_harmonic"][i]['MP2'][0],FF_dict["dihedrals_harmonic"][i]['MP2'][1],FF_dict["dihedrals_harmonic"][i]['MP2'][2]*180.0/np.pi)
                if FF_dict["dihedrals_harmonic"][i]['MP2'][0] == "harmonic":
                    FF_dict["dihedrals_harmonic"][i]['MP2'] = (FF_dict["dihedrals_harmonic"][i]['MP2'][0],FF_dict["dihedrals_harmonic"][i]['MP2'][1],int(FF_dict["dihedrals_harmonic"][i]['MP2'][2]))
    
    # Return to parent folder after processing all angles
    os.chdir(working_dir)

    return paths_to_dihedral_fitpot_plots_DFT, paths_to_dihedral_fitpot_plots_MP2, paths_to_dihedral_totpot_plots_DFT, paths_to_dihedral_totpot_plots_MP2

def OPLS_new(dihedrals,QC='DFT'):
            
    global FF_dict
    E_tot = 0.0
    for i in dihedrals:
        V1,V2,V3,V4 = FF_dict["dihedrals"][i[0][0]+'_'+i[0][1]+'_'+i[0][2]+'_'+i[0][3]][QC] 
        E_tot += 0.5*V1*(1.0+np.cos(1.0*i[1])) + 0.5*V2*(1.0-np.cos(2.0*i[1])) + 0.5*V3*(1.0+np.cos(3*i[1])) + 0.5*V4*(1.0-np.cos(4*i[1]))

    return E_tot

def OPLS_EFF(params,dihedrals):
            
    E_tot = 0.0
    for i in dihedrals:
        if (i[0][0],i[0][1],i[0][2],i[0][3]) in list(params.keys()):
            V0,V1,V2,V3,V4 = params[(i[0][0],i[0][1],i[0][2],i[0][3])]
            E_tot += V0+0.5*V1*(1.0+np.cos(1.0*i[1])) + 0.5*V2*(1.0-np.cos(2.0*i[1])) + 0.5*V3*(1.0+np.cos(3*i[1])) + 0.5*V4*(1.0-np.cos(4*i[1]))

    return E_tot

def OPLS_new_norm(dihedrals,QC='DFT'):
            
    global FF_dict
    E_tot = 0.0
    for i in dihedrals:
        if len(FF_dict["dihedrals"][i[0]][QC]) == 6:
            V0,V1,V2,V3,V4 = FF_dict["dihedrals"][i[0]][QC][1:]
        else:
            V0,V1,V2,V3,V4 = [0.0] + FF_dict["dihedrals"][i[0]][QC][1:]
        E_tot += V0 + 0.5*V1*(1.0+np.cos(1.0*i[1])) + 0.5*V2*(1.0-np.cos(2.0*i[1])) + 0.5*V3*(1.0+np.cos(3*i[1])) + 0.5*V4*(1.0-np.cos(4*i[1]))

    return E_tot

def dihedral_harmonic(dihedrals,QC='DFT'):
            
    global FF_dict
    E_tot = 0.0
    for i in dihedrals:
        d_type,K,d = FF_dict["dihedrals_harmonic"][i[0]][QC] 
        if d_type == "harmonic":
            E_tot += K*(1.0+d*np.cos(i[1]))
        elif d_type == "quadratic":
            E_tot += K*(i[1]-d)**(2.0)
    return E_tot

# This is a messy wrapper for FF_dict write commands.
def write_params(base_name,save_folder,charges_AA,masses_AA,atomtypes_AA,elements_AA,charges_UA,masses_UA,atomtypes_UA,elements_UA,qc_types=[]):

    # In python all global variables must be declared in each scope
    global FF_dict,VDW_dict

    # Create a temporary copy of FF_dict
    FF_dict_tmp = deepcopy(FF_dict)

    # First clean up VDW_dict by removing reverse types with the wrong priority ordering
    # NOTE: parameters are saved using type_a > type_b ordering
    # NOTE: Hydrogen containing pairs are removed when the fit_type is 'UA'
    pop_list = []
    keys = list(FF_dict_tmp["vdw"].keys())
    for i in keys:
        if i[0] > i[1]:
            if (i[1],i[0]) in list(FF_dict_tmp["vdw"].keys()): FF_dict_tmp["vdw"].pop((i[1],i[0]))

    ##############################
    #  Write DFT FF information  #
    ##############################
    if 'dft' in qc_types:

        # The formatting of this database file is compatible 
        # with the input for the polygen.py program.
        f = open(save_folder+'/'+base_name+'-DFT.db', 'w')

        # Write atom type definitions
        f.write("\n# Atom type definitions\n#\n{:10s} {:40s} {:40s} {:19s}\n".format("#","Atom_type","Label","Mass"))        
        for i in atomtypes_AA:
            f.write("{:10s} {:40s} {:39s} {:< 20.6f}\n"\
            .format("atom",str(i),elements_AA[i]+'-'+str(i),masses_AA[i]))
        for i in atomtypes_UA:
            f.write("{:10s} {:40s} {:39s} {:< 20.6f}\n"\
            .format("atom",str(i),elements_UA[i]+'-'+str(i),masses_UA[i]))

        # Write VDW definitions
        f.write("\n# VDW definitions\n#\n{:10s} {:60s} {:60s} {:11s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","Potential","eps (kcal/mol)","sigma (ang)"))        
        # AA-types
        for i in list(FF_dict_tmp["vdw"].keys()):
            if "-UA" in i[0]: continue
            f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0]),str(i[1]),'lj'," ".join([ "{:< 20.6f}".format(j) for j in FF_dict_tmp["vdw"][i][1:]])))        
        # UA-types
        for i in list(FF_dict_tmp["vdw"].keys()):
            if "-UA" not in i[0]: continue
            f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0]),str(i[1]),'lj'," ".join([ "{:< 20.6f}".format(j) for j in FF_dict_tmp["vdw"][i][1:]])))        

        # Write bond definitions
        f.write("\n# Bond type definitions\n#\n{:10s} {:40s} {:40s} {:21s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","style","k (kcal/mol ang^-2)","R0 (ang)"))        
        for i in list(FF_dict_tmp["bonds"].keys()):
            f.write("{:10s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f}\n".format("bond",i[0],i[1],FF_dict_tmp["bonds"][i]["DFT"][0],FF_dict_tmp["bonds"][i]["DFT"][1],FF_dict_tmp["bonds"][i]["DFT"][2]))

        # Write angle definitions
        f.write("\n# Angle type definitions\n#\n{:10s} {:40s} {:40s} {:40s} {:21s} {:20s} {:20s}\n"\
                .format("#","Atom_type","Atom_type","Atom_type","style","k (kcal/mol rad^-2)","Theta_0 (degrees)"))        
        for i in list(FF_dict_tmp["angles"].keys()):
            f.write("{:10s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f}\n"\
                    .format("angle",i[0],i[1],i[2],FF_dict_tmp["angles"][i]["DFT"][0],FF_dict_tmp["angles"][i]["DFT"][1],FF_dict_tmp["angles"][i]["DFT"][2]))

        # Write dihedral definitions
        f.write("\n# Dihedral/Torsional type definitions\n#\n{:10s} {:40s} {:40s} {:40s} {:40s} {:21s} {:20s}\n".format("#","Atom_type","Atom_type","Atom_type","Atom_type","style","Params"))
        for i in list(FF_dict_tmp["dihedrals"].keys()):

            # Skip any dummy types that were added
            if FF_dict_tmp["dihedrals"][i]["DFT"][-1] == "dummy":
                continue
            
            f.write("{:10s} {:40s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f} {:< 20.6f} {:< 20.6f}\n"\
             .format("torsion",i[0],i[1],i[2],i[3],FF_dict_tmp["dihedrals"][i]["DFT"][0],FF_dict_tmp["dihedrals"][i]["DFT"][1],FF_dict_tmp["dihedrals"][i]["DFT"][2],FF_dict_tmp["dihedrals"][i]["DFT"][3],\
                     FF_dict_tmp["dihedrals"][i]["DFT"][4]))
        for i in list(FF_dict_tmp["dihedrals_harmonic"].keys()):
            if FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][0] == "harmonic":
                f.write("{:10s} {:40s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20d} {:< 20d}\n"\
                 .format("torsion",i[0],i[1],i[2],i[3],FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][0],FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][1],int(FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][2]),int(FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][3])))
            elif FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][0] == "quadratic":
                f.write("{:10s} {:40s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f}\n"\
                 .format("torsion",i[0],i[1],i[2],i[3],FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][0],FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][1],float(FF_dict_tmp["dihedrals_harmonic"][i]["DFT"][2]),1))

        # Write charge definitions
        f.write("\n# Charge definitions\n#\n{:10s} {:41s} {:20s}\n".format("#","Atom_type","Charge"))
        for i in atomtypes_AA:
            f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_AA[i]))
        for i in atomtypes_UA:
            f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_UA[i]))

    ##############################
    #  Write MP2 FF information  #
    ##############################

    if 'mp2' in qc_types:

        # The formatting of this database file is compatible 
        # with the input for the polygen.py program.
        f = open(save_folder+'/'+base_name+'-MP2.db', 'w')

        # Write atom type definitions
        f.write("\n# Atom type definitions\n#\n{:10s} {:40s} {:40s} {:19s}\n".format("#","Atom_type","Label","Mass"))        
        for i in atomtypes_AA:                                                                                
            f.write("{:10s} {:40s} {:39s} {:< 20.6f}\n"\
            .format("atom",str(i),elements_AA[i]+'-'+str(i),masses_AA[i]))
        for i in atomtypes_UA:
            f.write("{:10s} {:40s} {:39s} {:< 20.6f}\n"\
            .format("atom",str(i),elements_UA[i]+'-'+str(i),masses_UA[i]))

        # Write VDW definitions
        f.write("\n# VDW definitions\n#\n{:10s} {:60s} {:60s} {:11s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","Potential","eps (kcal/mol)","sigma (ang)"))        
        # AA-types
        for i in list(FF_dict_tmp["vdw"].keys()):
            if "-UA" in i[0]: continue
            f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0]),str(i[1]),'lj'," ".join([ "{:< 20.6f}".format(j) for j in FF_dict_tmp["vdw"][i][1:]])))        
        # UA-types
        for i in list(FF_dict_tmp["vdw"].keys()):
            if "-UA" not in i[0]: continue
            f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0]),str(i[1]),'lj'," ".join([ "{:< 20.6f}".format(j) for j in FF_dict_tmp["vdw"][i][1:]])))        

        # Write bond definitions
        f.write("\n# Bond type definitions\n#\n{:10s} {:40s} {:40s} {:21s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","style","k (kcal/mol ang^-2)","R0 (ang)"))        
        for i in list(FF_dict_tmp["bonds"].keys()):
            f.write("{:10s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f}\n".format("bond",i[0],i[1],FF_dict_tmp["bonds"][i]["MP2"][0],FF_dict_tmp["bonds"][i]["MP2"][1],FF_dict_tmp["bonds"][i]["MP2"][2]))

        # Write angle definitions
        f.write("\n# Angle type definitions\n#\n{:10s} {:40s} {:40s} {:40s} {:21s} {:20s} {:20s}\n"\
                .format("#","Atom_type","Atom_type","Atom_type","style","k (kcal/mol rad^-2)","Theta_0 (degrees)"))        
        for i in list(FF_dict_tmp["angles"].keys()):
            f.write("{:10s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f}\n"\
                    .format("angle",i[0],i[1],i[2],FF_dict_tmp["angles"][i]["MP2"][0],FF_dict_tmp["angles"][i]["MP2"][1],FF_dict_tmp["angles"][i]["MP2"][2]))

        # Write dihedral definitions
        f.write("\n# Dihedral/Torsional type definitions\n#\n{:10s} {:40s} {:40s} {:40s} {:40s} {:21s} {:20s}\n".format("#","Atom_type","Atom_type","Atom_type","Atom_type","style","Params"))
        for i in list(FF_dict_tmp["dihedrals"].keys()):
            f.write("{:10s} {:40s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20.6f} {:< 20.6f} {:< 20.6f}\n"\
             .format("torsion",i[0],i[1],i[2],i[3],FF_dict_tmp["dihedrals"][i]["MP2"][0],FF_dict_tmp["dihedrals"][i]["MP2"][1],FF_dict_tmp["dihedrals"][i]["MP2"][2],FF_dict_tmp["dihedrals"][i]["MP2"][3],\
                     FF_dict_tmp["dihedrals"][i]["MP2"][4]))
        for i in list(FF_dict_tmp["dihedrals_harmonic"].keys()):
            f.write("{:10s} {:40s} {:40s} {:40s} {:40s} {:20s} {:< 20.6f} {:< 20d} {:< 20d}\n"\
             .format("torsion",i[0],i[1],i[2],i[3],"harmonic",FF_dict_tmp["dihedrals_harmonic"][i]["MP2"][0],int(FF_dict_tmp["dihedrals_harmonic"][i]["MP2"][1]),1))

        # Write charge definitions
        f.write("\n# Charge definitions\n#\n{:10s} {:41s} {:20s}\n".format("#","Atom_type","Charge"))
        for i in atomtypes_AA:
            f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_AA[i]))
        for i in atomtypes_UA:
            f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_UA[i]))

        f.close()

    return

# This is a messy wrapper for the commands to write the powerpoint summary
def write_pptx(base_name,save_folder,paths_to_bond_plots_DFT,paths_to_bond_plots_MP2,paths_to_angle_plots_DFT,paths_to_angle_plots_MP2,\
               paths_to_dihedral_fitpot_plots_DFT, paths_to_dihedral_fitpot_plots_MP2, paths_to_dihedral_totpot_plots_DFT, paths_to_dihedral_totpot_plots_MP2,\
               paths_to_dihedral_fitpot_plots_DFT_UA, paths_to_dihedral_fitpot_plots_MP2_UA, paths_to_dihedral_totpot_plots_DFT_UA, paths_to_dihedral_totpot_plots_MP2_UA,\
               paths_to_h_dihedral_fitpot_plots_DFT_AA, paths_to_h_dihedral_fitpot_plots_MP2_AA, paths_to_h_dihedral_totpot_plots_DFT_AA, paths_to_h_dihedral_totpot_plots_MP2_AA,\
               paths_to_h_dihedral_fitpot_plots_DFT_UA, paths_to_h_dihedral_fitpot_plots_MP2_UA, paths_to_h_dihedral_totpot_plots_DFT_UA, paths_to_h_dihedral_totpot_plots_MP2_UA,qc_types=[]):

    # Load presentation for DFT Results
    if 'dft' in qc_types:

        prs = Presentation()   

        # Iterate over the bond plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_bond_plots_DFT)/9)+1)):
            if len(paths_to_bond_plots_DFT) == 0: break
            slide = add_slide(prs,"Bond Fits")
            try:
                Current_Plots = paths_to_bond_plots_DFT[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_bond_plots_DFT[j*9:]
            except:
                print("ERROR: There is a problem reading the list of bond fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                


        # Iterate over the low energy approach plots in batches of 9 and
        # place them on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_angle_plots_DFT)/9)+1)):
            if len(paths_to_angle_plots_DFT) == 0: break
            slide = add_slide(prs,"Angle Fits")
            try:
                Current_Plots = paths_to_angle_plots_DFT[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_angle_plots_DFT[j*9:]
            except:
                print("ERROR: There is a problem reading the list of angle fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                


        # Iterate over the dihedral-AA fit plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_fitpot_plots_DFT)/9)+1)):            
            if len(paths_to_dihedral_fitpot_plots_DFT) == 0: break
            slide = add_slide(prs,"Dihedral-AA Fit Potential")
            try:
                Current_Plots = paths_to_dihedral_fitpot_plots_DFT[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_fitpot_plots_DFT[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the dihedral-UA fit plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_fitpot_plots_DFT_UA)/9)+1)):            
            if len(paths_to_dihedral_fitpot_plots_DFT_UA) == 0: break
            slide = add_slide(prs,"Dihedral-UA Fit Potential")
            try:
                Current_Plots = paths_to_dihedral_fitpot_plots_DFT_UA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_fitpot_plots_DFT_UA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the dihedral-AA tot plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_totpot_plots_DFT)/9)+1)):            
            if len(paths_to_dihedral_totpot_plots_DFT) == 0: break
            slide = add_slide(prs,"Dihedral-AA Total Potential")
            try:
                Current_Plots = paths_to_dihedral_totpot_plots_DFT[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_totpot_plots_DFT[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the dihedral-UA tot plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_totpot_plots_DFT_UA)/9)+1)):            
            if len(paths_to_dihedral_totpot_plots_DFT_UA) == 0: break
            slide = add_slide(prs,"Dihedral-UA Total Potential")
            try:
                Current_Plots = paths_to_dihedral_totpot_plots_DFT_UA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_totpot_plots_DFT_UA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # # Iterate over the average plots in batches of 9 and place them
        # # on the slides in 3x3 arrays
        # for j in range(int(int(len(paths_to_h_dihedral_fitpot_plots_DFT_AA)/9)+1)):            
        #     if len(paths_to_h_dihedral_fitpot_plots_DFT_AA) == 0: break
        #     slide = add_slide(prs,"Harmonic Dihedral-AA Fit Potential")
        #     try:
        #         Current_Plots = paths_to_h_dihedral_fitpot_plots_DFT_AA[j*9:(j+1)*9]

        #     except IndexError:
        #         Current_Plots = paths_to_h_dihedral_fitpot_plots_DFT_AA[j*9:]
        #     except:
        #         print "ERROR: There is a problem reading the list of harmonic dihedral fit *.png filenames. Exiting..."
        #         quit()
        #     add_9(slide,Current_Plots)                

        # # Iterate over the average plots in batches of 9 and place them
        # # on the slides in 3x3 arrays
        # for j in range(int(int(len(paths_to_h_dihedral_fitpot_plots_DFT_UA)/9)+1)):            
        #     if len(paths_to_h_dihedral_fitpot_plots_DFT_UA) == 0: break
        #     slide = add_slide(prs,"Harmonic Dihedral-UA Fit Potential")
        #     try:
        #         Current_Plots = paths_to_h_dihedral_fitpot_plots_DFT_UA[j*9:(j+1)*9]

        #     except IndexError:
        #         Current_Plots = paths_to_h_dihedral_fitpot_plots_DFT_UA[j*9:]
        #     except:
        #         print "ERROR: There is a problem reading the list of harmonic dihedral fit *.png filenames. Exiting..."
        #         quit()
        #     add_9(slide,Current_Plots)                

        # Iterate over the harmonic-dihedral-AA tot plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_h_dihedral_totpot_plots_DFT_AA)/9)+1)):            
            if len(paths_to_h_dihedral_totpot_plots_DFT_AA) == 0: break
            slide = add_slide(prs,"Harmonic Dihedral-AA Total Potential")
            try:
                Current_Plots = paths_to_h_dihedral_totpot_plots_DFT_AA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_h_dihedral_totpot_plots_DFT_AA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the harmonic-dihedral-UA tot plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_h_dihedral_totpot_plots_DFT_UA)/9)+1)):            
            if len(paths_to_h_dihedral_totpot_plots_DFT_UA) == 0: break
            slide = add_slide(prs,"Harmonic Dihedral-UA Total Potential")
            try:
                Current_Plots = paths_to_h_dihedral_totpot_plots_DFT_UA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_h_dihedral_totpot_plots_DFT_UA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Save presentation
        prs.save(save_folder+'/{}-DFT-fit_summary.pptx'.format(base_name))

    # Load presentation for MP2 Results
    if 'mp2' in qc_types:

        prs = Presentation()   

        # Iterate over the bond plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_bond_plots_MP2)/9)+1)):
            if len(paths_to_bond_plots_MP2) == 0: break
            slide = add_slide(prs,"Bond Fits")
            try:
                Current_Plots = paths_to_bond_plots_MP2[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_bond_plots_MP2[j*9:]
            except:
                print("ERROR: There is a problem reading the list of bond fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the low energy approach plots in batches of 9 and
        # place them on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_angle_plots_MP2)/9)+1)):
            if len(paths_to_angle_plots_MP2) == 0: break
            slide = add_slide(prs,"Angle Fits")
            try:
                Current_Plots = paths_to_angle_plots_MP2[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_angle_plots_MP2[j*9:]
            except:
                print("ERROR: There is a problem reading the list of angle fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                


        # Iterate over the average plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_fitpot_plots_MP2)/9)+1)):            
            if len(paths_to_dihedral_fitpot_plots_MP2) == 0: break
            slide = add_slide(prs,"Dihedral-AA Fit Potential")
            try:
                Current_Plots = paths_to_dihedral_fitpot_plots_MP2[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_fitpot_plots_MP2[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the average plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_fitpot_plots_MP2_UA)/9)+1)):            
            if len(paths_to_dihedral_fitpot_plots_MP2_UA) == 0: break
            slide = add_slide(prs,"Dihedral-UA Fit Potential")
            try:
                Current_Plots = paths_to_dihedral_fitpot_plots_MP2_UA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_fitpot_plots_MP2_UA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the average plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_totpot_plots_MP2)/9)+1)):            
            if len(paths_to_dihedral_totpot_plots_MP2) == 0: break
            slide = add_slide(prs,"Dihedral-AA Total Potential")
            try:
                Current_Plots = paths_to_dihedral_totpot_plots_MP2[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_totpot_plots_MP2[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the average plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_dihedral_totpot_plots_MP2_UA)/9)+1)):            
            if len(paths_to_dihedral_totpot_plots_MP2_UA) == 0: break
            slide = add_slide(prs,"Dihedral-UA Total Potential")
            try:
                Current_Plots = paths_to_dihedral_totpot_plots_MP2_UA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_dihedral_totpot_plots_MP2_UA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the average plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_h_dihedral_fitpot_plots_MP2_AA)/9)+1)):            
            if len(paths_to_h_dihedral_fitpot_plots_MP2_AA) == 0: break
            slide = add_slide(prs,"Harmonic Dihedral-AA Fit Potential")
            try:
                Current_Plots = paths_to_h_dihedral_fitpot_plots_MP2_AA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_h_dihedral_fitpot_plots_MP2_AA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of harmonic dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Iterate over the average plots in batches of 9 and place them
        # on the slides in 3x3 arrays
        for j in range(int(int(len(paths_to_h_dihedral_fitpot_plots_MP2_UA)/9)+1)):            
            if len(paths_to_h_dihedral_fitpot_plots_MP2_UA) == 0: break
            slide = add_slide(prs,"Harmonic Dihedral-UA Fit Potential")
            try:
                Current_Plots = paths_to_h_dihedral_fitpot_plots_MP2_UA[j*9:(j+1)*9]

            except IndexError:
                Current_Plots = paths_to_h_dihedral_fitpot_plots_MP2_UA[j*9:]
            except:
                print("ERROR: There is a problem reading the list of harmonic dihedral fit *.png filenames. Exiting...")
                quit()
            add_9(slide,Current_Plots)                

        # Save presentation
        prs.save(save_folder+'/{}-MP2-fit_summary.pptx'.format(base_name))

    return

# Description: This function drives the fit of the opls torsional parameters. It accepts
#              various arguments that modify the convergence criteria and optimization
#              algorithm. 
# 
# Inputs:      x_vals            : a list holding the keys in Data for configurations being fit
#              y_vals            : a list holding the interaction energies for each configuration being fit 
#              Pairs             : a list of the unique pair-types being fit.
#              weight            : a scalar that controls the weight of the update at each step.
#              delta_xhi2_thresh : a scalar that controls the xhi^2 convergence threshold for the fit. 
#              min_cycles        : an integer that controls the minimum number of fit cycles to perform.
#              max_cycles        : an integer that controls the maximum number of fit cycles to perform. 
# 
# Returns:     fit_vals          : a list of arrays, each holding the FF based interaction energies in each cycle of the fit.
#
def iterative_fit_all(dihedral_scan,x_vals,y_vals,Dihedrals,weight,fit,delta_xhi2_thresh=0.00001,min_cycles=10,max_cycles=10000,QC_type="DFT"):

    # In python you need to expand the scope of global variables at every level
    global FF_dict,fit_type,Dihedral_Data,dihedral_fit_type,V_range,QC_dihedral
    
    QC_dihedral = QC_type
    V_range = max(y_vals)-min(y_vals)
    V_range *= 10.0
#    print "V_range for this fit: {}".format(V_range)


    # Initialize FF_dict with dihedral parameters 
    for i in Dihedrals:

        # Initialization commands for new dihedrals
        if i not in list(FF_dict["dihedrals"].keys()):
            FF_dict["dihedrals"][i]={}
            FF_dict["dihedrals"][i][QC_type] = ["opls",0.0,0.0,0.0,0.0,0.0]
        elif QC_type not in list(FF_dict["dihedrals"][i].keys()):
            FF_dict["dihedrals"][i][QC_type] = ["opls",0.0,0.0,0.0,0.0,0.0]

        # Command to add the V0 term to the dihedral parameters (needed for the MD-based refits)
        if len(FF_dict["dihedrals"][i][QC_type]) == 5:
            FF_dict["dihedrals"][i][QC_type] = ["opls",0.0,0.0,0.0,0.0,0.0]

    # Intialize cycle count and array of fit values
    cycle_count = 0
    dihedral_fit_type = 0 
    fit_vals = [ E_OPLS_Fit_all(x_vals,0.0,0.0,0.0,0.0,0.0) ] * (max_cycles+1)
    xhi2_previous = np.mean((y_vals-fit_vals[cycle_count])**2) 

    # Fit loop
    while (1):
        
        # Increment the cycle count and print diagnostic
        cycle_count += 1

        # Each cycle consists of looping over all pair-types being fit
        # To eliminate ordering bias, the list of types is shuffled every cycle.
        random.shuffle(Dihedrals)
        for i in Dihedrals:

            # Set the pair-type that is being fit
            dihedral_fit_type = i                            

            # Fit the potential
            if fit.lower() == "opls":
                params,err = curve_fit(E_OPLS_Fit_all,x_vals,y_vals,p0=FF_dict["dihedrals"][i][QC_type][1:],maxfev=10000) # ORIGINAL

            # Update FF_dict with a weighted update (when weight is 1, the updated params are simply the params returned by curve_fit)   
            # Note: FF_dict carries around the potential type as the first element, so the parameters start at index 1.
            params = [ j*(1.0-weight)+params[count_j]*weight for count_j,j in enumerate(FF_dict["dihedrals"][i][QC_type][1:]) ]              
            FF_dict["dihedrals"][i][QC_type] = [FF_dict["dihedrals"][i][QC_type][0]] + params 

        # Update the fit data array
        dihedral_fit_type = 0   # dummy value so that all pair-potentials are taken from VDW_dict 
        fit_vals[cycle_count] = E_OPLS_Fit_all(x_vals,0.0,0.0,0.0,0.0,0.0)

        # Calculate mean squared deviation
        xhi2 = np.mean((y_vals-fit_vals[cycle_count])**2)
        delta_xhi2 = xhi2_previous - xhi2
        xhi2_previous = xhi2

        # Check break conditions
        if cycle_count >= min_cycles:
            if cycle_count == max_cycles:
#                print "Maximum number of cycles reached ({}), breaking fit loop...".format(max_cycles)
                break
            elif delta_xhi2 < delta_xhi2_thresh:
#                print "Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2) after {} cycles, breaking fit loop...".format(delta_xhi2,delta_xhi2_thresh,cycle_count)
                break

    return fit_vals[:cycle_count+1],xhi2**(0.5)

# fit_type is a globally defined tuple
def E_OPLS_Fit_all(ind,V0_fit,V1_fit,V2_fit,V3_fit,V4_fit):

    # In python you need to expand the scope of global variables at every level
    global FF_dict,dihedral_fit_type,Dihedral_Data,QC_dihedral,dihedral_scan,V_range

    # Initialize E_OPLS for the current configuration
    E_OPLS = np.zeros(len(ind))

    # Cumulatively add up the OPLS energy (fit type should be defined as a global variable outside the function)
    # NOTE: Dihedral_Data[dihedral_scan][i]["Dihedrals"] only includes the coincident dihedrals being scanned in this fit 
    for count_i,i in enumerate(ind):
        for count_j,j in enumerate(Dihedral_Data[dihedral_scan][i]["Dihedrals"]):

            # Order the pair to conform to FF_dict key definitions
            # The type being fitted is a variable, others are pulled from FF_dict
            # Note: for UA fits, Hydrogen containing pairs are omitted from the FF_dict
            #       so that they do not participate in the fit.
            dihedral_type = j[0]

            if dihedral_type == dihedral_fit_type:
                V0,V1,V2,V3,V4 = (V0_fit,V1_fit,V2_fit,V3_fit,V4_fit)

            # If the data is being read from the FF_dict check for the V0 term (Note: FF_dict carries
            # around the potential type as the first argument, so the parameters start at index 1
            elif dihedral_type in list(FF_dict["dihedrals"].keys()):
                if len(FF_dict["dihedrals"][dihedral_type][QC_dihedral]) == 5:
                    V0,V1,V2,V3,V4 = [0.0] + FF_dict["dihedrals"][dihedral_type][QC_dihedral][1:]
                elif len(FF_dict["dihedrals"][dihedral_type][QC_dihedral]) == 6:
                    V0,V1,V2,V3,V4 = FF_dict["dihedrals"][dihedral_type][QC_dihedral][1:]
            else:
                continue

            # Place some sanity checks on the fit range
            if True in [ np.abs(V1) > V_range, np.abs(V2) > V_range, np.abs(V3) > V_range, np.abs(V4) > V_range ]: E_OPLS[count_i] = 1E10; continue
#            if np.abs(V1) > 50.0 or np.abs(V2) > 50.0 or np.abs(V3) > 50.0 or np.abs(V4) > 50.0: E_OPLS[count_i] = 1E10; continue            
            E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(j[1])) + 0.5*V2*(1-np.cos(2.0*j[1])) + 0.5*V3*(1+np.cos(3.0*j[1])) + 0.5*V4*(1-np.cos(4.0*j[1]))
#            E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(j[1])) + 0.5*V2*(1-np.cos(2.0*j[1])) + 0.5*V3*(1+np.cos(3.0*j[1]))

    return E_OPLS

# Fit function used in the IBI dihedral parse function
# angle is in units of degrees, V_range is globally defined
# outside the function and used to avoid overfitting the potential
def IBI_OPLS_Fit(angles,V0,V1,V2,V3,V4):

    global V_range

    # Initialize E_OPLS for the current configuration
    E_OPLS = np.zeros(len(angles))

    # Cumulatively add up the OPLS energy (fit type should be defined as a global variable outside the function)
    # NOTE: Dihedral_Data[dihedral_scan][i]["Dihedrals"] only includes the coincident dihedrals being scanned in this fit 
    for count_i,i in enumerate(angles):

        # Place some sanity checks on the fit range
        if True in [ np.abs(V1) > V_range, np.abs(V2) > V_range, np.abs(V3) > V_range, np.abs(V4) > V_range ]: E_OPLS[count_i] = 1E10; continue
        E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(i*np.pi/180.0)) + 0.5*V2*(1-np.cos(2.0*i*np.pi/180.0)) + 0.5*V3*(1+np.cos(3.0*i*np.pi/180.0)) + 0.5*V4*(1-np.cos(4.0*i*np.pi/180.0))

    return E_OPLS

# Function for calculating the total OPLS potential for use in the parse_dihedrals_qcmatch function
def OPLS_Calc(angles,dihedral_types,FF,QC_dihedral):

    # Initialize E_OPLS for the current configuration
    E_OPLS = np.zeros(len(angles))

    # Cumulatively add up the OPLS energy (fit type should be defined as a global variable outside the function)
    # NOTE: Dihedral_Data[dihedral_scan][i]["Dihedrals"] only includes the coincident dihedrals being scanned in this fit 
    for count_i,i in enumerate(angles):

        # In OPLS_OPLS_Fit, angles holds a list of tuples with the angle of each coincident type being fit
        for count_j,j in enumerate(i):

            # If the data is being read from the FF_dict check for the V0 term (Note: FF_dict carries
            # around the potential type as the first argument, so the parameters start at index 1
            if dihedral_types[count_j] in list(FF["dihedrals"].keys()):
                if len(FF["dihedrals"][dihedral_types[count_j]][QC_dihedral]) == 5:
                    V0,V1,V2,V3,V4 = [0.0] + FF["dihedrals"][dihedral_types[count_j]][QC_dihedral][1:]
                elif len(FF["dihedrals"][dihedral_types[count_j]][QC_dihedral]) == 6:
                    V0,V1,V2,V3,V4 = FF["dihedrals"][dihedral_types[count_j]][QC_dihedral][1:]
            else:
                print("MISSING DIHEDRAL: {}".format(dihedral_types[count_j]))
                continue

            # Place some sanity checks on the fit range
            if True in [ np.abs(V1) > V_range, np.abs(V2) > V_range, np.abs(V3) > V_range, np.abs(V4) > V_range ]: E_OPLS[count_i] = 1E10; continue
            E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(j*np.pi/180.0)) + 0.5*V2*(1-np.cos(2.0*j*np.pi/180.0)) + 0.5*V3*(1+np.cos(3.0*j*np.pi/180.0)) + 0.5*V4*(1-np.cos(4.0*j*np.pi/180.0))

    return E_OPLS

# # Fit function used in the IBI dihedral parse function
# # angle is in units of degrees, V_range is globally defined
# # outside the function and used to avoid overfitting the potential
# def OPLS_OPLS_Fit(angles,V0_fit,V1_fit,V2_fit,V3_fit,V4_fit):

#     global V_range,dihedral_types,dihedral_fit_type,QC_dihedral

#     # Initialize E_OPLS for the current configuration
#     E_OPLS = np.zeros(len(angles))

#     # Cumulatively add up the OPLS energy (fit type should be defined as a global variable outside the function)
#     # NOTE: Dihedral_Data[dihedral_scan][i]["Dihedrals"] only includes the coincident dihedrals being scanned in this fit 
#     for count_i,i in enumerate(angles):

#         # In OPLS_OPLS_Fit, angles holds a list of tuples with the angle of each coincident type being fit
#         for count_j,j in enumerate(i):

#             if dihedral_types[count_j] == dihedral_fit_type:
#                 V0,V1,V2,V3,V4 = (V0_fit,V1_fit,V2_fit,V3_fit,V4_fit)

#             # If the data is being read from the FF_dict check for the V0 term (Note: FF_dict carries
#             # around the potential type as the first argument, so the parameters start at index 1
#             elif dihedral_types[count_j] in FF_dict["dihedrals"].keys():
#                 if len(FF_dict["dihedrals"][dihedral_types[count_j]][QC_dihedral]) == 5:
#                     V0,V1,V2,V3,V4 = [0.0] + FF_dict["dihedrals"][dihedral_types[count_j]][QC_dihedral][1:]
#                 elif len(FF_dict["dihedrals"][dihedral_types[count_j]][QC_dihedral]) == 6:
#                     V0,V1,V2,V3,V4 = FF_dict["dihedrals"][dihedral_types[count_j]][QC_dihedral][1:]
#             else:
#                 continue

#             # Place some sanity checks on the fit range
#             if True in [ np.abs(V1) > V_range, np.abs(V2) > V_range, np.abs(V3) > V_range, np.abs(V4) > V_range ]: E_OPLS[count_i] = 1E10; continue
#             E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(j*np.pi/180.0)) + 0.5*V2*(1-np.cos(2.0*j*np.pi/180.0)) + 0.5*V3*(1+np.cos(3.0*j*np.pi/180.0)) + 0.5*V4*(1-np.cos(4.0*j*np.pi/180.0))

#     return E_OPLS

# Fit function used in the IBI dihedral parse function
# angle is in units of degrees, V_range is globally defined
# outside the function and used to avoid overfitting the potential
def OPLS_OPLS_Fit_MIN(V0_0  = None,V1_0  = None,V2_0  = None,V3_0  = None,V4_0  = None,\
                      V0_1  = None,V1_1  = None,V2_1  = None,V3_1  = None,V4_1  = None,\
                      V0_2  = None,V1_2  = None,V2_2  = None,V3_2  = None,V4_2  = None,\
                      V0_3  = None,V1_3  = None,V2_3  = None,V3_3  = None,V4_3  = None,\
                      V0_4  = None,V1_4  = None,V2_4  = None,V3_4  = None,V4_4  = None,\
                      V0_5  = None,V1_5  = None,V2_5  = None,V3_5  = None,V4_5  = None,\
                      V0_6  = None,V1_6  = None,V2_6  = None,V3_6  = None,V4_6  = None,\
                      V0_7  = None,V1_7  = None,V2_7  = None,V3_7  = None,V4_7  = None,\
                      V0_8  = None,V1_8  = None,V2_8  = None,V3_8  = None,V4_8  = None,\
                      V0_9  = None,V1_9  = None,V2_9  = None,V3_9  = None,V4_9  = None,\
                      V0_10 = None,V1_10 = None,V2_10 = None,V3_10 = None,V4_10 = None,\
                      V0_11 = None,V1_11 = None,V2_11 = None,V3_11 = None,V4_11 = None,\
                      V0_12 = None,V1_12 = None,V2_12 = None,V3_12 = None,V4_12 = None,\
                      V0_13 = None,V1_13 = None,V2_13 = None,V3_13 = None,V4_13 = None,\
                      V0_14 = None,V1_14 = None,V2_14 = None,V3_14 = None,V4_14 = None,\
                      V0_15 = None,V1_15 = None,V2_15 = None,V3_15 = None,V4_15 = None,\
                      V0_16 = None,V1_16 = None,V2_16 = None,V3_16 = None,V4_16 = None,\
                      unique_fit_dihedral_types=None,angles=None,E_0=None,V_range=None,dihedral_types=None,QC_dihedral=None,FF_dict=None,w_pot=1.0,w_hyper=0.01,b_hyper=0.1,w_exp=0.0,w_harm=0.0,T=298.0,kb=0.0019872041):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined V*_* variables. 
    V0_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'V0_' in i and "__" not in i and local_vars[i] is not None ]) ])
    V1_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'V1_' in i and "__" not in i and local_vars[i] is not None ]) ])
    V2_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'V2_' in i and "__" not in i and local_vars[i] is not None ]) ])
    V3_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'V3_' in i and "__" not in i and local_vars[i] is not None ]) ])
    V4_vals = np.array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'V4_' in i and "__" not in i and local_vars[i] is not None ]) ])

    # Consistency checks
    if unique_fit_dihedral_types is None or angles is None or V_range is None or dihedral_types is None or QC_dihedral is None or FF_dict is None:
        print("ERROR in OPLS_OPLS_Fit_MIN: {} must be supplied to the function, it/they "\
        .format(", ".join([ i for i in ["unique_fit_dihedral_types","angles","V_range","dihedral_types","QC_dihedral","FF_dict"] if locals()[i] is None ])))
        print("                            are defined as optional arguments for convenient anonymization. Exiting...")
        quit()
    if len(angles) != len(E_0):
        print("ERROR in OPLS_OPLS_Fit_MIN: E_0, angles, and dihedral_types must have the same length. Exiting...")
        quit()
    if len(angles[0]) != len(dihedral_types):
        print("ERROR in OPLS_OPLS_Fit_MIN: dihedral_types must have the same length as each tuple in angles. Exiting...")
        quit()
    if len(V0_vals) != len(unique_fit_dihedral_types):
        print("ERROR in OPLS_OPLS_Fit_MIN: the function expects the len(unique_fit_dihedral_types) and number of non-zero V*_* variables to be equal. Exiting...")
        quit()

    # Initialize a list of fourier coefficients that are properly indexed to dihedral types
    # The if/else controls whether the coefficients are assigned to the input fit parameters or
    # taken from FF_dict (dihedrals not being fit on this scan)
    params = []
    for count_i,i in enumerate(dihedral_types):
        if i in unique_fit_dihedral_types:
            idx = unique_fit_dihedral_types.index(i)
            params += [ (V0_0,V1_vals[idx],V2_vals[idx],V3_vals[idx],V4_vals[idx]) ]
        elif i in list(FF_dict["dihedrals"].keys()):
            if len(FF_dict["dihedrals"][i][QC_dihedral]) == 5:
                params += [tuple([0.0] + FF_dict["dihedrals"][i][QC_dihedral][1:])]
            elif len(FF_dict["dihedrals"][i][QC_dihedral]) == 6:
                params += [tuple(FF_dict["dihedrals"][i][QC_dihedral][1:])]
        else:
            params += [(0.0,0.0,0.0,0.0,0.0)]

    # Initialize E_OPLS for the current configuration
    E_OPLS = np.zeros(len(angles))

    # Cumulatively add up the OPLS energy (fit type should be defined as a global variable outside the function)
    for count_i,i in enumerate(angles):

        # In OPLS_OPLS_Fit, angles holds a list of tuples with the angle of each coincident type being fit
        for count_j,j in enumerate(i):

            # Unpack params
            V0,V1,V2,V3,V4 = params[count_j]

            # Calculate the dihedral energy for this configuration
            E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(j*np.pi/180.0)) + 0.5*V2*(1-np.cos(2.0*j*np.pi/180.0)) + 0.5*V3*(1+np.cos(3.0*j*np.pi/180.0)) + 0.5*V4*(1-np.cos(4.0*j*np.pi/180.0))

    # Return the penalty function
    #          potential fit                       boltzman weighting                  harmonic constraint on parameter magnitudes (L2 regularization)                       hyperbolic constraint on parameter magnitudes
    return w_pot*np.mean((E_OPLS - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_harm * np.mean(np.array([ j for i in params for j in i ])**(2.0)) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )
#    return w_pot*np.mean((E_OPLS - E_0)**(2.0)*(1-w_exp+w_exp*np.exp(-E_0/(kb*T)))) + w_hyper * np.mean( (np.array([ j for i in params for j in i ])**(2.0) + b_hyper**(2.0) )**(0.5) - b_hyper )

# Description: This function drives the fit of the harmonic dihedrals parameters. It accepts
#              various arguments that modify the convergence criteria and optimization
#              algorithm. 
# 
# Inputs:      x_vals            : a list holding the keys in Data for configurations being fit
#              y_vals            : a list holding the interaction energies for each configuration being fit 
#              Pairs             : a list of the unique pair-types being fit.
#              weight            : a scalar that controls the weight of the update at each step.
#              delta_xhi2_thresh : a scalar that controls the xhi^2 convergence threshold for the fit. 
#              min_cycles        : an integer that controls the minimum number of fit cycles to perform.
#              max_cycles        : an integer that controls the maximum number of fit cycles to perform. 
# 
# Returns:     fit_vals          : a list of arrays, each holding the FF based interaction energies in each cycle of the fit.
#
def iterative_harmonic_fit_all(dihedral_scan,x_vals,y_vals,Dihedrals,weight,fit,delta_xhi2_thresh=0.00001,min_cycles=10,max_cycles=10000):

    # In python you need to expand the scope of global variables at every level
    global FF_dict,fit_type,Dihedral_Harmonic_Data,QC_type,dihedral_fit_type,min_max_dict
        
    # Initialize FF_dict with dihedral parameters (K,d values are stored)
    for i in Dihedrals:
        if i not in list(FF_dict["dihedrals_harmonic"].keys()):
            FF_dict["dihedrals_harmonic"][i]={}

            # If the first angle is closer to 0, then d is set to -1
            # NOTE: I've found that if the guess is "too good" for a multi dimensional lstsq fit, then the optimizer won't change that value. 
            #       For this reason, using the "first" angle value from the dictionary actually outperforms the better avg_val initial guess for the optimum angle. 
            first_val = [ k[1] for j in list(Dihedral_Harmonic_Data[dihedral_scan].keys()) for k in Dihedral_Harmonic_Data[dihedral_scan][j]["Dihedrals"] if k[0] == i ][0]
            avg_val = np.mean([ k[1] for j in list(Dihedral_Harmonic_Data[dihedral_scan].keys()) for k in Dihedral_Harmonic_Data[dihedral_scan][j]["Dihedrals"] if k[0] == i ])
#            first_val = avg_val

            if first_val <= np.pi/2.0 and first_val > -np.pi/2.0:
                if np.abs(first_val*180.0/np.pi) < 2.0:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("harmonic",0.0,-1.0)
                else:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("quadratic",30.0,first_val)

            # If the first angle is closer to 180, then d is set to -1
            else:
                if np.abs((np.pi-np.abs(first_val))*180.0/np.pi) < 2.0:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("harmonic",0.0,1.0)
                else:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("quadratic",30.0,first_val)

        elif QC_type not in list(FF_dict["dihedrals_harmonic"][i].keys()):

            # If the first angle is closer to 0 then d is set to -1
            first_val = [ k[1] for j in list(Dihedral_Harmonic_Data[dihedral_scan].keys()) for k in Dihedral_Harmonic_Data[dihedral_scan][j]["Dihedrals"] if k[0] == i ][0]
            avg_val = np.mean([ k[1] for j in list(Dihedral_Harmonic_Data[dihedral_scan].keys()) for k in Dihedral_Harmonic_Data[dihedral_scan][j]["Dihedrals"] if k[0] == i ])

            if first_val <= np.pi/2.0 and first_val > -np.pi/2.0:
                if np.abs(first_val*180.0/np.pi) < 2.0:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("harmonic",0.0,-1.0)
                else:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("quadratic",30.0,first_val)

            # If the first angle is closer to 180, then d is set to -1
            else:
                if np.abs((np.pi-np.abs(first_val))*180.0/np.pi) < 2.0:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("harmonic",0.0,1.0)
                else:
                    FF_dict["dihedrals_harmonic"][i][QC_type] = ("quadratic",30.0,first_val)

    # Initialize min_max_dictionary for limits on the quadratic dihedral fits    
    min_max_dict = {}
    for j in [ i for i in list(FF_dict["dihedrals_harmonic"].keys()) if FF_dict["dihedrals_harmonic"][i][QC_type][0] == "quadratic" ]:
        min = None
        max = None
        for k in list(Dihedral_Harmonic_Data[dihedral_scan].keys()):
            for m in Dihedral_Harmonic_Data[dihedral_scan][k]["Dihedrals"]:
                if m[0] == j:
                    if min == None or m[1] < min:
                        min = m[1]
                    if max == None or m[1] > max:
                        max = m[1]
        min_max_dict[j] = (min,max)

    # Intialize cycle count and array of fit values
    cycle_count = 0
    dihedral_fit_type = 0 
    fit_vals = [ E_harmonic_dihedral_fit_all(x_vals,0.0) ] * (max_cycles+1)
    xhi2_previous = np.mean((y_vals-fit_vals[cycle_count])**2) 
    xhi2_running = None
    
    # Save the old FF dictionary
    FF_running = deepcopy(FF_dict["dihedrals_harmonic"])
    if FF_dict["dihedrals_harmonic"][Dihedrals[0]][QC_type][0] == "quadratic":
        initial_guesses = np.linspace(min_max_dict[Dihedrals[0]][0],min_max_dict[Dihedrals[0]][1],10)
        quadratic_flag = 1
    else:
        initial_guesses = ['dummy']
        quadratic_flag = 0

    # Outer loop over initial guesses for the quadratic dihedral angle
    for m in initial_guesses:
        if quadratic_flag == 1:
            FF_dict["dihedrals_harmonic"][Dihedrals[0]][QC_type] = ( "quadratic", 1.0, m )

        # Enter fit loop
        while (1):

            # Increment the cycle count and print diagnostic
            cycle_count += 1

            # Each cycle consists of looping over all pair-types being fit
            # To eliminate ordering bias, the list of types is shuffled every cycle.
            random.shuffle(Dihedrals)
            for i in Dihedrals:

                # Set the pair-type that is being fit
                dihedral_fit_type = i                            

                # Fit the potential (two slightly different fit functinos for harmonic vs quadratic 
                if FF_dict["dihedrals_harmonic"][i][QC_type][0] == "harmonic":
                    params,err = curve_fit(E_harmonic_dihedral_fit_all,x_vals,y_vals,p0=FF_dict["dihedrals_harmonic"][i][QC_type][1],maxfev=10000)

                    # Update FF_dict with a weighted update (when weight is 1, the updated params are simply the params returned by curve_fit)
                    params = FF_dict["dihedrals_harmonic"][i][QC_type][1]*(1.0-weight)+params*weight
                    FF_dict["dihedrals_harmonic"][i][QC_type] = (FF_dict["dihedrals_harmonic"][i][QC_type][0],params[0],FF_dict["dihedrals_harmonic"][i][QC_type][2])

                elif FF_dict["dihedrals_harmonic"][i][QC_type][0] == "quadratic":
                    params,err = curve_fit(E_harmonic_dihedral_fit_all,x_vals,y_vals,p0=FF_dict["dihedrals_harmonic"][i][QC_type][1:3],maxfev=10000)

                    # Update FF_dict with a weighted update (when weight is 1, the updated params are simply the params returned by curve_fit)
                    params = np.array(FF_dict["dihedrals_harmonic"][i][QC_type][1:3])*(1.0-weight)+params*weight
                    FF_dict["dihedrals_harmonic"][i][QC_type] = (FF_dict["dihedrals_harmonic"][i][QC_type][0],params[0],params[1])

            # Update the fit data array
            dihedral_fit_type = 0   # dummy value so that all torsional-potentials are taken from FF_dict 
            fit_vals[cycle_count] = E_harmonic_dihedral_fit_all(x_vals,0.0) 

            # Calculate mean squared deviation
            xhi2 = np.mean((y_vals-fit_vals[cycle_count])**2)
            delta_xhi2 = xhi2_previous - xhi2
            xhi2_previous = xhi2

            # Check break conditions
            if cycle_count >= min_cycles:
                if cycle_count == max_cycles:
#                    print "Maximum number of cycles reached ({}), breaking fit loop...".format(max_cycles)
                    break
                elif delta_xhi2 < delta_xhi2_thresh:
#                    print "Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2) after {} cycles, breaking fit loop...".format(delta_xhi2,delta_xhi2_thresh,cycle_count)
                    break
        if xhi2_running == None or xhi2 < xhi2_running:
            FF_running = deepcopy(FF_dict["dihedrals_harmonic"])
            xhi2_running = deepcopy(xhi2)

    FF_dict["dihedrals_harmonic"] = deepcopy(FF_running)

    return fit_vals[:cycle_count+1],xhi2**(0.5)

# fit_type is a globally defined tuple
def E_harmonic_dihedral_fit_all(ind,K_fit,theta0=None):

    # In python you need to expand the scope of global variables at every level
    global FF_dict,dihedral_fit_type,Dihedral_Harmonic_Data,QC_type,dihedral_scan,min_max_dict

    # Initialize E_OPLS for the current configuration
    E_harmonic = np.zeros(len(ind))

    # Set quadratic flag if theta0 != []
    if theta0 != None: 
        quad_flag = 1
    else: 
        quad_flag = 0

    # Cumulatively add up the OPLS energy (fit type should be defined as a global variable outside the function)
    for count_i,i in enumerate(ind):
        for count_j,j in enumerate(Dihedral_Harmonic_Data[dihedral_scan][i]["Dihedrals"]):

            # Order the pair to conform to FF_dict key definitions
            # The type being fitted is a variable, others are pulled from FF_dict
            # Note: for UA fits, Hydrogen containing pairs are omitted from the FF_dict
            #       so that they do not participate in the fit.
            dihedral_type = j[0]

            if dihedral_type == dihedral_fit_type:                
                if quad_flag == 1:
                    d_type,K,K2 = (FF_dict["dihedrals_harmonic"][dihedral_type][QC_type][0],K_fit,theta0)
                else:
                    d_type,K,K2 = (FF_dict["dihedrals_harmonic"][dihedral_type][QC_type][0],K_fit,FF_dict["dihedrals_harmonic"][dihedral_type][QC_type][2])
            elif dihedral_type in list(FF_dict["dihedrals_harmonic"].keys()):
                d_type,K,K2 = FF_dict["dihedrals_harmonic"][dihedral_type][QC_type]
            else:                
                continue

            # Calculate the harmonic energy
            if d_type == "harmonic":
                E_harmonic[count_i] += K*(1.0+K2*np.cos(j[1]))
            elif d_type == "quadratic":
                E_harmonic[count_i] += K*(j[1]-K2)**(2.0)
                if K2 < min_max_dict[dihedral_type][0] and K2 > min_max_dict[dihedral_type][1]:
                    return 1.0E10
                else:
                    E_harmonic[count_i] += K*(j[1]-K2)**(2.0)

    return E_harmonic

# fit_type is a globally defined tuple
def E_quadratic_dihedral_fit_all(ind,K_fit,theta_0):

    # In python you need to expand the scope of global variables at every level
    global FF_dict,dihedral_fit_type,Dihedral_Harmonic_Data,QC_type,dihedral_scan

    # Initialize E_quadratic for the current configuration
    E_quadratic = np.zeros(len(ind))

    # Cumulatively add up the dihedral energy
    for count_i,i in enumerate(ind):
        for count_j,j in enumerate(Dihedral_Harmonic_Data[dihedral_scan][i]["Dihedrals"]):

            # Order the pair to conform to FF_dict key definitions
            # The type being fitted is a variable, others are pulled from FF_dict
            # Note: for UA fits, Hydrogen containing pairs are omitted from the FF_dict
            #       so that they do not participate in the fit.
            dihedral_type = j[0]

            if dihedral_type == dihedral_fit_type:                
                d_type,K,K2 = (FF_dict["dihedrals_harmonic"][dihedral_type][QC_type][0],K_fit,theta_0)
            elif dihedral_type in list(FF_dict["dihedrals_harmonic"].keys()):
                d_type,K,K2 = FF_dict["dihedrals_harmonic"][dihedral_type][QC_type]
            else:                
                continue

            # Calculate the harmonic energy
            if d_type == "harmonic":
                E_harmonic[count_i] += K*(1.0+K2*np.cos(j[1]))
            elif d_type == "quadratic":
                E_harmonic[count_i] += K*(j[1]-K2)**(2.0)

    return E_harmonic


# Wrapper script for the molecule.db writing commands
def write_molecule_db(folder,master_files,charges):

    # Write examples of each molecule in the sim to a molecule.db file. This is used when
    # adding the intramolecular parameters to the master database. Each parameter in the master is linked
    # to the molecule(s) used to generate it. In the case of intramolecular params, the globally optimized
    # geometry is saved to the db.
    with open(folder+'/molecules.db','w') as f:
        for count_i,i in enumerate(master_files):
            f.write("\nmol {:6d} start\n".format(count_i+1))

            # Grab the elements, geometry, and atomtypes
            elements,geo = xyz_parse(i)    
            atomtypes = get_scan_atomtypes(i)   

            # Save an xyz for viewing the configuration in VMD
            f.write('{}\n\n'.format(len(geo)))
            for count_j,j in enumerate(geo):
                f.write('  {:<10s} {:< 20.6f} {:< 20.6f} {:< 20.6f} {:< 12.6f} {:<60s}\n'.\
                format(elements[count_j],j[0],j[1],j[2],charges[atomtypes[count_j]],atomtypes[count_j]))
            f.write("mol {:6d} end\n".format(count_i+1))
    return

# Parses the output files for each scan and calculates potential corrections, dihedral angles, and the other prerequisites for the
# fitting calculations. All information is stored in the dihedral_data dictionary.
def get_dihedral_data(list_of_dihedral_types,dihedral_dirs,fit_type='AA',atomtypes_lists=[],charge_dict=[],qc_types=[],corr_list=[],scope='scan',one_four_scale_coul=0.0,one_four_scale_vdw=0.0):

    global FF_dict

    # Initialize dictionary
    dihedral_data = {}

    # Loop over all dihedral scans
    current_dir = os.getcwd()

    for count_d,dihedrals in enumerate(list_of_dihedral_types):

        # Initialize subdictionary to hold the data for this scan
        dihedral_data[dihedrals] = {}
        
        # Change into current dihedral directory
        os.chdir(dihedral_dirs[count_d])

        # Check for a scan run
        dirs = [ dn for dp, dn, filenames in os.walk('.') if True in [ "scan" in i for i in dn ] ]
        if len(dirs) > 0:
            dirs = dirs[0]
        
        # Parse bond data for continuous scan
        if True in [ 'scan' in i for i in  dirs]:

            # Since there may be multiple scans the follow lists are used to ensure that the various configurations properly register with one another
            scan_indices  = []
            scan_angles   = []
            scan_energies = []
            scan_data     = {dihedrals:{}}
            
            # Loop over the discovered scans
            for count_s,s in enumerate(dirs):

                # Change into the scan directory
                os.chdir(s)

                # Find the output file for this scan
                output_file = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ][0]

                # Grab a copy of the atomtypes (with the link- labels removed for assigning charges)
#                atomtypes_N = [ next( j for j in i.split('link-') if j != '' ) for i in atomtypes_lists[count_d] ] # with the link- labels removed
                atomtypes_N = atomtypes_lists[count_d]
#                charges_N   = [ charge_dict[next( j for j in i.split('link-') if j != '' )] for i in atomtypes_lists[count_d] ]  # OLD WAY, the charges are calculated here for the coul correction, but it is no longer used by the fit.
                charges_N   = [ 0.0 for i in atomtypes_lists[count_d] ]  # Avoids a conflict with the xyz_fit option

                # Find the number of jobs in this run
                N_jobs          = 1
                with open(output_file,'r') as f:
                    for lines in f:
                        if "JOBS TO BE PROCESSED THIS RUN" in lines: 
                            N_jobs = int(lines.split()[3])
                            break

                # If there are multiple jobs, then that means a preliminary relaxation and/or mode scan is performed prior to the run data
                # so the first job is skipped
                if N_jobs == 1:
                    parse_job = 1
                else: 
                    parse_job = 0

                # Find the optimized geometries and single point energies
                parse_flag      = -1
                job_count       = 0
                geo_flag        = 0
                completion_flag = 1
                atom_count      = 0
                scan_data[dihedrals] = {}
                Geo_dict = {}        
                with open(output_file,'r') as f:
                    for lines in f:

                        # Only the last job is parsed, the parse_job flag and job_count flags control appropriate parsing
                        if parse_job == 1:

                            if "RELAXED SURFACE SCAN STEP" in lines: 
                                parse_flag += 1
                                scan_data[dihedrals][str(parse_flag)] = {}
                                scan_data[dihedrals][str(parse_flag)]["Geo"] = np.zeros([len(atomtypes_N),3])
                                scan_data[dihedrals][str(parse_flag)]["Elements"] = ["X"]*len(atomtypes_N)
                                scan_data[dihedrals][str(parse_flag)]["Completion"] = 0
                            if "CARTESIAN COORDINATES (ANGSTROEM)" in lines: geo_flag = 1; continue
                            if geo_flag == 1:
                                fields = lines.split()
                                if len(fields) == 1: continue
                                scan_data[dihedrals][str(parse_flag)]["Geo"][atom_count,:] = np.array([float(fields[1]),float(fields[2]),float(fields[3])])
                                scan_data[dihedrals][str(parse_flag)]["Elements"][atom_count] = fields[0]
                                atom_count += 1
                                if atom_count == len(atomtypes_N): 
                                    atom_count = 0 
                                    geo_flag = 0

                            if "FINAL SINGLE POINT ENERGY" in lines: scan_data[dihedrals][str(parse_flag)]["DFT"] = float(lines.split()[4])*Htokcalmol
                            if "MP2 TOTAL ENERGY" in lines: scan_data[dihedrals][str(parse_flag)]["MP2"] = float(lines.split()[3])*Htokcalmol

                            if "****ORCA TERMINATED NORMALLY****" in lines: 
                                for i in list(scan_data[dihedrals].keys()): scan_data[dihedrals][i]["Completion"] = 1

                        # Only parse the last job
                        elif "JOB NUMBER" in lines: 
                            job_count += 1
                            if job_count == 2:
                                parse_job = 1

                # If no data was found (i.e. parse_job == 0) then initialize the incomplete flag
                if parse_job == 0:
                    scan_data[dihedrals][str(0)] = {}
                    scan_data[dihedrals][str(0)]["Completion"] = 0

                # Skip the rest of the parse if not all configurations completed.
                skip_flag = 0
                for i in list(scan_data[dihedrals].keys()):
                    if scan_data[dihedrals][i]["Completion"] == 0:
                        skip_flag = 1
                if skip_flag == 1: 
                    os.chdir('..')
                    continue

                # Initialize min_ind to keep track of which geometry should be used for the normalization
                min_key = "0"
                delta_180 = 1000.0

                # Loop over the dihedral data for the scan and collect the modes being parametrized and corrections to the potential energy 
                for i in natural_sort(list(scan_data[dihedrals].keys())):

                    # Calculate hybridization (needed for identifying atom types)
                    Adj_mat = Table_generator(scan_data[dihedrals][i]["Elements"],scan_data[dihedrals][i]["Geo"])
                    Hybridizations = Hybridization_finder(scan_data[dihedrals][i]["Elements"],Adj_mat)

                    # Add adj_mat, charge, and atomtype information to the scan_data dictionary (used by other functions for generating the MD corrections, when requested)                
                    scan_data[dihedrals][i]["Adj_mat"]   = Adj_mat
                    scan_data[dihedrals][i]["Atomtypes"] = atomtypes_N
                    scan_data[dihedrals][i]["Charges"]   = charges_N
                    scan_data[dihedrals][i]["folder"]    = dihedral_dirs[count_d].split('/')[0]

                    # Find the dihedral types and instances being scanned. Find_scanned_dihedrals also returns tuples with the atom indices and angles, indexed to the 
                    # scan_data[dihedrals][str(count_g)]["Dihedrals"] list. This functionality was added to average over indentical coincident scans that share the 1,2,3 atoms with the main dihedral. 
                    if scope == "scan":
                        scan_data[dihedrals][i]["Dihedrals"],scan_data[dihedrals][i]["Dihedral_Atoms"] =\
                            Find_scanned_dihedrals(Adj_mat,scan_data[dihedrals][i]["Geo"],atomtypes_N,output_file,fit_type=fit_type)
                    elif scope == "all":
                        scan_data[dihedrals][i]["Dihedrals"] = Find_all_dihedrals(Adj_mat,scan_data[dihedrals][i]["Geo"],atomtypes_N,output_file,fit_type=fit_type)[0]
                        scan_data[dihedrals][i]["Scanned_Dihedrals"],scan_data[dihedrals][i]["Dihedral_Atoms"] =\
                            Find_scanned_dihedrals(Adj_mat,scan_data[dihedrals][i]["Geo"],atomtypes_N,output_file,fit_type=fit_type)

                    # Find the dihedrals that aren't scanned. These are sometimes weakly constrained during the optimizations to avoid conformational reorganization.
                    if len(scan_data[dihedrals][i]["Dihedral_Atoms"]) > 0:
                        center_atoms = [scan_data[dihedrals][i]["Dihedral_Atoms"][0][0][1],scan_data[dihedrals][i]["Dihedral_Atoms"][0][0][2]]           # Find the indices for the 2-3 atoms in the scanned dihedral
                        tmp_d,tmp_atoms = Find_all_dihedrals(Adj_mat,scan_data[dihedrals][i]["Geo"],atomtypes_N,output_file,fit_type=fit_type)           # Find all dihedrals
                        keep_ind = set([ count_j for count_j,j in enumerate(tmp_atoms) if j[0][1] not in center_atoms or j[0][2] not in center_atoms ])  # keep all non-scanned dihedrals
                        scan_data[dihedrals][i]["Unscanned_Dihedrals"] = [ j for count_j,j in enumerate(tmp_d) if count_j in keep_ind ]                       # save the unscanned dihedral types (type_tuple,angle) to the dictionary
                        scan_data[dihedrals][i]["Unscanned_Dihedral_Atoms"] = [ j for count_j,j in enumerate(tmp_atoms) if count_j in keep_ind ]              # save the unscanned dihedral atoms (index_tuple,angle) to the dictionary

                    # Check if the current configuration will be used in the master dihedral_data 
                    # If this is the first scan, then all configurations are kept and the various scan lists are populated. 
                    if count_s == 0:

                        # populate the dihedral_data dictionary
                        if i not in list(dihedral_data[dihedrals].keys()):
                            dihedral_data[dihedrals][i] = {}
                        for j in list(scan_data[dihedrals][i].keys()):
                            dihedral_data[dihedrals][i][j] = scan_data[dihedrals][i][j]

                        # Populate the match lists
                        # NOTE: sometimes there are no relevant dihedrals being parsed (e.g. UA-fit type where only UA-Hydrogens are parse of the dihedral scanned dihedral)
                        if len(scan_data[dihedrals][i]["Dihedrals"]) > 0:
                            scan_indices  += [i]
                            angle = scan_data[dihedrals][i]["Dihedrals"][0][1]
                            if angle < 0.0: angle += 2.0*np.pi
                            if angle > 2.0*np.pi: angle -= 2.0*np.pi                                
                            scan_angles   += [angle]
                            scan_energies += [scan_data[dihedrals][i]["DFT"]]

                    # For multiple scans, information is only kept for configurations with lower energies than the existing entry in dihedral_data/scan_energies
                    # Configurations are registered with one another via the angle of the constrained dihedral. 
                    else:

                        # Skip if there are no relevent dihedrals
                        # NOTE: sometimes there are no relevant dihedrals being parsed (e.g. UA-fit type where only UA-Hydrogens are parse of the dihedral scanned dihedral)
                        if len(scan_data[dihedrals][i]["Dihedrals"]) == 0:
                            continue

                        # find matching index
                        angle = scan_data[dihedrals][i]["Dihedrals"][0][1]
                        if angle < 0.0: angle += 2.0*np.pi
                        if angle > 2.0*np.pi: angle -= 2.0*np.pi                                
                        idx = np.argmin(np.abs(np.array(scan_angles)-angle))
                        if np.abs(angle-scan_angles[idx]) > 0.01:

                            print("\nWARNING in get_dihedral_data: Check that the scan constraints are consistent for scan {}.".format(dihedrals)) 
                            print("                              The error in registration between scans is > 0.01 radians (np.abs(angle-scan_angles[idx]) = {} radians)\n".format(np.abs(angle-scan_angles[idx])))

                        # Check if this configuration should be used
                        if scan_data[dihedrals][i]["DFT"] < scan_energies[idx]:
                            match = scan_indices[idx]
                            scan_energies[idx] = scan_data[dihedrals][i]["DFT"]
                            
                            # Update dihedral_data dictionary entry
                            for j in list(scan_data[dihedrals][i].keys()):
                                dihedral_data[dihedrals][match][j] = scan_data[dihedrals][i][j]

                        # Else, skip dihedral_data update
                        else:                            
                            continue

                    # If this is the first scan then match = i, else match is set by the above else statement
                    if count_s == 0:
                        match = i

                    # Parse bonds, angles, and dihedrals (needed for calculating force-field energy correction)
                    bond_tuples,angle_tuples,dihedral_tuples,one_five_tuples = Find_modes(Adj_mat,atomtypes_N)                                

                    # Calculate force-field corrections
                    if "dft" in qc_types:
                        if "bond" in corr_list:
                            dihedral_data[dihedrals][match]["E_bond_DFT"]  = E_bond(dihedral_data[dihedrals][match]["Geo"],atomtypes_N,bond_tuples,QC='DFT')
                        else:
                            dihedral_data[dihedrals][match]["E_bond_DFT"]  = 0.0
                        if "angle" in corr_list:
                            dihedral_data[dihedrals][match]["E_angle_DFT"] = E_angle(dihedral_data[dihedrals][match]["Geo"],atomtypes_N,angle_tuples,QC='DFT')
                        else:
                            dihedral_data[dihedrals][match]["E_angle_DFT"] = 0.0

                    if "mp2" in qc_types:
                        if "bond" in corr_list:
                            dihedral_data[dihedrals][match]["E_bond_MP2"]  = E_bond(dihedral_data[dihedrals][match]["Geo"],atomtypes_N,bond_tuples,QC='MP2')
                        else:
                            dihedral_data[dihedrals][match]["E_bond_MP2"]  = 0.0
                        if "angle" in corr_list:
                            dihedral_data[dihedrals][match]["E_angle_MP2"] = E_angle(dihedral_data[dihedrals][match]["Geo"],atomtypes_N,angle_tuples,QC='MP2')
                        else:
                            dihedral_data[dihedrals][match]["E_angle_MP2"] = 0.0

                    if "coul" in corr_list:
                        dihedral_data[dihedrals][match]["E_coul"] = E_coul(dihedral_data[dihedrals][match]["Geo"],Adj_mat,atomtypes_N,charges_N,[bond_tuples,angle_tuples,dihedral_tuples],\
                                                                      (0.0,0.0,one_four_scale_coul),fit_type)
                    else:
                        dihedral_data[dihedrals][match]["E_coul"] = 0.0

                    if "vdw" in corr_list:
                        dihedral_data[dihedrals][match]["E_VDW"] = E_LJ(dihedral_data[dihedrals][match]["Geo"],Adj_mat,atomtypes_N,charges_N,[bond_tuples,angle_tuples,dihedral_tuples],\
                                                                   (0.0,0.0,one_four_scale_vdw),fit_type)
                    else:
                        dihedral_data[dihedrals][match]["E_VDW"] = 0.0

                    # The energies of each scan are normalized to the energy of the closest dihedral to 180 #
                    for count_j,j in enumerate(dihedral_data[dihedrals][match]["Dihedrals"]):
                        if 180.0-np.abs(j[1]*180.0/np.pi)<delta_180:
                            delta_180 = 180.0-np.abs(j[1]*180.0/np.pi)
                            min_key   = match

                    # Save the fit energies (QC energy less the corrections)
                    if 'dft' in qc_types:
                        dihedral_data[dihedrals][match]["DFT_Fit"] = dihedral_data[dihedrals][match]["DFT"] -\
                                                                       dihedral_data[dihedrals][match]["E_bond_DFT"] -\
                                                                       dihedral_data[dihedrals][match]["E_angle_DFT"] -\
                                                                       dihedral_data[dihedrals][match]["E_coul"] -\
                                                                       dihedral_data[dihedrals][match]["E_VDW"]
                    if 'mp2' in qc_types:
                        dihedral_data[dihedrals][match]["MP2_Fit"] = dihedral_data[dihedrals][match]["MP2"] -\
                                                                       dihedral_data[dihedrals][match]["E_bond_MP2"] -\
                                                                       dihedral_data[dihedrals][match]["E_angle_MP2"] -\
                                                                       dihedral_data[dihedrals][match]["E_coul"] -\
                                                                       dihedral_data[dihedrals][match]["E_VDW"]
                os.chdir('..')

        # Parse parallel scan information (the scan is split across several jobs and stored in separate folders)
        else:

            # Find the optimized geometry files and starting geometries, and sort by number.
            # NOTE: starting_geometries have the atom_type mapping for the configuration which is necessary for fragment based parsing
            all_geometries = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.xyz") ]
            ind = [ count_i for count_i,i in enumerate(all_geometries) if "geo_opt.xyz" in i ]
            geometries = [ i for count_i,i in enumerate(all_geometries) if count_i in ind ]
            starting_geometries = [ i for count_i,i in enumerate(all_geometries) if count_i not in ind ]
            sort_nicely(geometries)        
            sort_nicely(starting_geometries)        

            # Find output files and sort by number (matched to geometries)
            outputs = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ]
            sort_nicely(outputs)

            # Initialize min_ind to keep track of which geometry should be used for the normalization
            min_key = "0"
            delta_180 = 1000.0

            # Loop over all discovered geometries and save their geometries to Geo_dict and
            # catenate to a single file for viewing.
            skip_flag = 0 
            for count_g,g in enumerate(geometries):

                # Initialize sub-sub-dictionary for the current geometry's info
                dihedral_data[dihedrals][str(count_g)] = {}

                # Check for job completion based on the discovery of the necessary energies
                dihedral_data[dihedrals][str(count_g)]["Completion"] = 0
                with open(outputs[count_g],'r') as out:                    
                    for lines in out:
                        if "****ORCA TERMINATED NORMALLY****" in lines: dihedral_data[dihedrals][str(count_g)]["Completion"] = 1
                        
                # Skip the rest of the parse if not all configurations completed.
                if dihedral_data[dihedrals][str(count_g)]["Completion"] == 0 : skip_flag = 1; continue  

                # Parse elements, geometry, and adj_mat (for the first geo)
                Elements, Geo = xyz_parse(g)
                if count_g == 0: Adj_mat = Table_generator(Elements,Geo)

                # Grab a copy of the atomtypes (with the link- labels removed)
                atomtypes_N = [ next( j for j in i.split('link-') if j != '' ) for i in atomtypes_lists[count_d] ]
#                charges_N   = [ charge_dict[i] for i in atomtypes_N ]  # OLD WAY, these are used to calculate the coulomb correction, but this is no longer used by the dihedral fit function. 
                charges_N   = [ 0.0 for i in atomtypes_lists[count_d] ]  # Avoids a conflict with the xyz_fit option

                # Add adj_mat, charge, and atomtype information to the dihedral_data dictionary (used by other functions for generating the MD corrections, when requested)                
                dihedral_data[dihedrals][str(count_g)]["Adj_mat"]   = Adj_mat
                dihedral_data[dihedrals][str(count_g)]["Atomtypes"] = atomtypes_N
                dihedral_data[dihedrals][str(count_g)]["Charges"]   = charges_N
                dihedral_data[dihedrals][str(count_g)]["folder"]    = dihedral_dirs[count_d].split('/')[0]

                # Also save the Elements and Geo to the dictionary (used later in fit functions for saving the catenated geometries)
                dihedral_data[dihedrals][str(count_g)]["Elements"] = Elements
                dihedral_data[dihedrals][str(count_g)]["Geo"]      = Geo

                # Calculate hybridization (needed for identifying atom types)
                Hybridizations = Hybridization_finder(Elements,Adj_mat)

                # Id atom types in the geometry # UNCOMMENT THIS LINE IF YOU WANT THE PROGRAM TO AUTOMATICALLY ID TYPES, NOT SURE WHY YOU'D WANT TO BUT APPARENTLY I THOUGHT IT WAS A GOOD IDEA ASP
    #            Atom_types = id_types(Elements,Adj_mat,2,Hybridizations,Geo)
    #            Atom_types = atomtypes

                # Find the dihedral types and instances being scanned. Find_scanned_dihedrals also returns tuples with the atom indices and angles, indexed to the 
                # dihedral_data[dihedrals][str(count_g)]["Dihedrals"] list. This functionality was added to average over indentical coincident scans that share the 1,2,3 atoms with the main dihedral. 
                if scope == "scan":
                    dihedral_data[dihedrals][str(count_g)]["Dihedrals"],dihedral_data[dihedrals][str(count_g)]["Dihedral_Atoms"] =\
                        Find_scanned_dihedrals(Adj_mat,Geo,atomtypes_N,outputs[count_g],fit_type=fit_type)
                elif scope == "all":
                    dihedral_data[dihedrals][str(count_g)]["Dihedrals"] = Find_all_dihedrals(Adj_mat,Geo,atomtypes_N,outputs[count_g],fit_type=fit_type)[0]
                    dihedral_data[dihedrals][str(count_g)]["Scanned_Dihedrals"],dihedral_data[dihedrals][str(count_g)]["Dihedral_Atoms"] =\
                         Find_scanned_dihedrals(Adj_mat,Geo,atomtypes_N,outputs[count_g],fit_type=fit_type)

                # Find the dihedrals that aren't scanned. These are sometimes weakly constrained during the optimizations to avoid conformational reorganization.
                if len(scan_data[dihedrals][i]["Dihedral_Atoms"]) > 0:
                    center_atoms = [dihedral_data[dihedrals][i]["Dihedral_Atoms"][0][0][1],dihedral_data[dihedrals][i]["Dihedral_Atoms"][0][0][2]]            # Find the indices for the 2-3 atoms in the scanned dihedral
                    tmp_d,tmp_atoms = Find_all_dihedrals(Adj_mat,scan_data[dihedrals][i]["Geo"],atomtypes_N,output_file,fit_type=fit_type)                    # Find all dihedrals
                    keep_ind = set([ count_j for count_j,j in enumerate(tmp_atoms) if j[0][1] not in center_atoms or j[0][2] not in center_atoms ])           # keep all non-scanned dihedrals
                    dihedral_data[dihedrals][i]["Unscanned_Dihedrals"] = [ j for count_j,j in enumerate(tmp_d) if count_j in keep_ind ]                       # save the unscanned dihedral types (type_tuple,angle) to the dictionary
                    dihedral_data[dihedrals][i]["Unscanned_Dihedral_Atoms"] = [ j for count_j,j in enumerate(tmp_atoms) if count_j in keep_ind ]              # save the unscanned dihedral atoms (index_tuple,angle) to the dictionary

                # Parse bonds, angles, and dihedrals (needed for calculating force-field energy correction)
                bond_tuples,angle_tuples,dihedral_tuples,one_five_tuples = Find_modes(Adj_mat,atomtypes_N)

                # Calculate force-field corrections
                if "dft" in qc_types:
                    if "bond" in corr_list:
                        dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"]  = E_bond(Geo,atomtypes_N,bond_tuples,QC='DFT')
                    else:
                        dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"]  = 0.0
                    if "angle" in corr_list:
                        dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"] = E_angle(Geo,atomtypes_N,angle_tuples,QC='DFT')
                    else:
                        dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"] = 0.0

                if "mp2" in qc_types:
                    if "bond" in corr_list:
                        dihedral_data[dihedrals][str(count_g)]["E_bond_MP2"]  = E_bond(Geo,atomtypes_N,bond_tuples,QC='MP2')
                    else:
                        dihedral_data[dihedrals][str(count_g)]["E_bond_MP2"]  = 0.0
                    if "angle" in corr_list:
                        dihedral_data[dihedrals][str(count_g)]["E_angle_MP2"] = E_angle(Geo,atomtypes_N,angle_tuples,QC='MP2')
                    else:
                        dihedral_data[dihedrals][str(count_g)]["E_angle_MP2"] = 0.0

                if "coul" in corr_list:
                    dihedral_data[dihedrals][str(count_g)]["E_coul"] = E_coul(Geo,Adj_mat,atomtypes_N,charges_N,[bond_tuples,angle_tuples,dihedral_tuples],(0.0,0.0,one_four_scale_coul),fit_type)
                else:
                    dihedral_data[dihedrals][str(count_g)]["E_coul"] = 0.0

                if "vdw" in corr_list:
                    dihedral_data[dihedrals][str(count_g)]["E_VDW"] = E_LJ(Geo,Adj_mat,atomtypes_N,charges_N,[bond_tuples,angle_tuples,dihedral_tuples],(0.0,0.0,one_four_scale_vdw),fit_type)
                else:
                    dihedral_data[dihedrals][str(count_g)]["E_VDW"] = 0.0

                # The energies of each scan are normalized to the energy of the closest dihedral to 180 #
                for count_i,i in enumerate(dihedral_data[dihedrals][str(count_g)]["Dihedrals"]):
                    if 180.0-np.abs(i[1]*180.0/np.pi)<delta_180:
                        delta_180 = 180.0-np.abs(i[1]*180.0/np.pi)
                        min_key   = str(count_g)

                # Collect energies 
                # NOTE: energies are converted to kcal/mol by default and assume hartree output
                with open(outputs[count_g],'r') as out:
                    
                    for lines in out:
                        fields = lines.split()

                        if len(fields) == 5 and fields[0].lower() == "final" and fields[1].lower() == "single" and fields[2].lower() == "point" and fields[3].lower() == "energy" and 'dft' in qc_types:
                            dihedral_data[dihedrals][str(count_g)]["DFT"] = float(fields[4])*Htokcalmol
                            dihedral_data[dihedrals][str(count_g)]["DFT_Fit"] = float(fields[4])*Htokcalmol - \
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"] -\
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"] -\
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_coul"] -\
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_VDW"]

                        if len(fields) == 5 and fields[0].lower() == "mp2" and fields[1].lower() == "total" and fields[2].lower() == "energy:" and 'mp2' in qc_types:
                            dihedral_data[dihedrals][str(count_g)]["MP2"] = float(fields[3])*Htokcalmol 
                            dihedral_data[dihedrals][str(count_g)]["MP2_Fit"] = float(fields[3])*Htokcalmol - \
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_bond_MP2"] -\
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_angle_MP2"] -\
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_coul"] -\
                                                                                   dihedral_data[dihedrals][str(count_g)]["E_VDW"]
                            break

            # Check if all coincident dihedrals are identical. Averaging is performed or indentical configurations
            # This helps average over asymmetries in the configuration. More elaborate symmetries will be accounted for
            # in the future. 
            # if len(dihedral_atoms) > 1:
            #     symmetry = symmetry_check(dihedral_atoms,dihedral_data[dihedrals]["0"]["Dihedrals"])
            #     if symmetry == 1:
            #         dihedral_data[dihedrals] = average_symmetric(dihedral_data[dihedrals],qc_types,corr_list)
        
        # If not all of the jobs for this scan completed, skip the normalization
        if skip_flag == 1:
            os.chdir('..')
            continue
        
        # Normalize all of the scan energies by the most "trans" dihedral energy
        for keys in list(dihedral_data[dihedrals].keys()):
            if keys == min_key: continue
            if "dft" in qc_types:
                dihedral_data[dihedrals][keys]["DFT"] = dihedral_data[dihedrals][keys]["DFT"] - dihedral_data[dihedrals][min_key]["DFT"]
                dihedral_data[dihedrals][keys]["DFT_Fit"] = dihedral_data[dihedrals][keys]["DFT_Fit"] - dihedral_data[dihedrals][min_key]["DFT_Fit"]
            if "mp2" in qc_types:
                dihedral_data[dihedrals][keys]["MP2"] = dihedral_data[dihedrals][keys]["MP2"] - dihedral_data[dihedrals][min_key]["MP2"]
                dihedral_data[dihedrals][keys]["MP2_Fit"] = dihedral_data[dihedrals][keys]["MP2_Fit"] - dihedral_data[dihedrals][min_key]["MP2_Fit"]
                    
        # I encountered bizarre behavior when trying to simply copy the normalization values to separate objects and normalize all data at once
        # I had to resort to normalizing all non_min values then normalize the minimum configuration last
        if "dft" in qc_types:
            dihedral_data[dihedrals][min_key]["DFT"] = dihedral_data[dihedrals][min_key]["DFT"] - dihedral_data[dihedrals][min_key]["DFT"]
            dihedral_data[dihedrals][min_key]["DFT_Fit"] = dihedral_data[dihedrals][min_key]["DFT_Fit"] - dihedral_data[dihedrals][min_key]["DFT_Fit"]
        if "mp2" in qc_types:
            dihedral_data[dihedrals][min_key]["MP2"] = dihedral_data[dihedrals][min_key]["MP2"] - dihedral_data[dihedrals][min_key]["MP2"]
            dihedral_data[dihedrals][min_key]["MP2_Fit"] = dihedral_data[dihedrals][min_key]["MP2_Fit"] - dihedral_data[dihedrals][min_key]["MP2_Fit"]
        
        os.chdir('..')
        
        
        # Align geometries across the scan so that the 1-2 and 2-3 bonds are aligned and the geometry is centered about the 2-3 bond. 
        # Alignment is based on the first configuation. 
        for key_count,keys in enumerate(natural_sort(list(dihedral_data[dihedrals].keys()))):

            # NOTE: sometimes there are no relevant dihedrals being parsed (e.g. UA-fit type where only UA-Hydrogens are parse of the dihedral scanned dihedral)
            if len(dihedral_data[dihedrals][keys]["Dihedrals"]) == 0:
                break
            
            # Store alignment data for the first key
            if key_count == 0:
                atom_1,atom_2,atom_3,atom_4 = dihedral_data[dihedrals][keys]["Dihedral_Atoms"][0][0]
                v1 = dihedral_data[dihedrals][keys]["Geo"][atom_1] - dihedral_data[dihedrals][keys]["Geo"][atom_2]
                v2 = dihedral_data[dihedrals][keys]["Geo"][atom_2] - dihedral_data[dihedrals][keys]["Geo"][atom_3]
            

            # Perform 1-2 alignment (if statement avoids alignment if it is already aligned. Sometimes there are roundoff errors when the argument of acos is near 1.0)
            c_v1    = dihedral_data[dihedrals][keys]["Geo"][atom_1] - dihedral_data[dihedrals][keys]["Geo"][atom_2]
            if np.dot(c_v1,v1) / ( norm(c_v1)*norm(v1) ) < 1.0:
                angle   = acos( np.dot(c_v1,v1) / ( norm(c_v1)*norm(v1) ) )
                rot_vec = np.cross(v1,c_v1)
                if norm(rot_vec) > 0.0:

                    # try positive rotation
                    tmp = axis_rot(dihedral_data[dihedrals][keys]["Geo"][atom_1],rot_vec,dihedral_data[dihedrals][keys]["Geo"][atom_2],angle,mode="radian") - dihedral_data[dihedrals][keys]["Geo"][atom_2]
                    diff = acos( np.dot(tmp,v1) / ( norm(tmp)*norm(v1) ) )
                    if np.abs(diff) < 0.01:
                        for count_i,i in enumerate(dihedral_data[dihedrals][keys]["Geo"]):
                            dihedral_data[dihedrals][keys]["Geo"][count_i] = axis_rot(i,rot_vec,dihedral_data[dihedrals][keys]["Geo"][atom_2],angle,mode="radian")
                    # else, perform negative rotation
                    else:
                        for count_i,i in enumerate(dihedral_data[dihedrals][keys]["Geo"]):
                            dihedral_data[dihedrals][keys]["Geo"][count_i] = axis_rot(i,rot_vec,dihedral_data[dihedrals][keys]["Geo"][atom_2],-angle,mode="radian")

            # Perform 2-3 alignment
            c_v2     = dihedral_data[dihedrals][keys]["Geo"][atom_2] - dihedral_data[dihedrals][keys]["Geo"][atom_3]
            rot_vec  = dihedral_data[dihedrals][keys]["Geo"][atom_1] - dihedral_data[dihedrals][keys]["Geo"][atom_2]
            min_diff = 100.0
            steps    = [0.1,0.01,0.001,0.0001]

            # Find rotation angle. The nested loops are simple bracketing algorithm
            for count_i,i in enumerate(steps):
                if count_i == 0:
                    a=0.0
                    b=2.0*np.pi
                else:
                    a=angle-steps[count_i-1]
                    b=angle+steps[count_i-1]
                for j in np.arange(a,b,i):
                    tmp = dihedral_data[dihedrals][keys]["Geo"][atom_2] - axis_rot(dihedral_data[dihedrals][keys]["Geo"][atom_3],rot_vec,dihedral_data[dihedrals][keys]["Geo"][atom_2],j,mode="radian")

                    # if statement avoids domain error in the acos call. Sometimes there are roundoff errors when the argument of acos is near 1.0.
                    if np.dot(tmp,v2) / ( norm(tmp)*norm(v2) ) < 1.0:
                        diff = np.abs(acos( np.dot(tmp,v2)/( norm(tmp)*norm(v2) )))
                    else:
                        diff = 100.0

                    # Assign alignment angle
                    if diff < min_diff:
                        min_diff = diff
                        angle = j

            # Align about 2-3 based on the best angle
            for count_i,i in enumerate(dihedral_data[dihedrals][keys]["Geo"]):
                dihedral_data[dihedrals][keys]["Geo"][count_i] = axis_rot(i,rot_vec,dihedral_data[dihedrals][keys]["Geo"][atom_2],angle,mode="radian")

            # Center the geometry about the 2-3 centroid
            centroid = ( dihedral_data[dihedrals][keys]["Geo"][atom_2] + dihedral_data[dihedrals][keys]["Geo"][atom_3] ) / 2.0
            for count_i,i in enumerate(dihedral_data[dihedrals][keys]["Geo"]):
                dihedral_data[dihedrals][keys]["Geo"][count_i] = dihedral_data[dihedrals][keys]["Geo"][count_i] - centroid 

        # For UA types, the dihedrals list might be empty
        if len(dihedral_data[dihedrals]["0"]["Dihedrals"]) == 0:
            clean_flag = 1
        else:
            clean_flag = 0

            # Determine threshold based on the mode of the second derivatives in the set
            # Collect the keys, angles, and energies, calculate the second derivatives and delete data points at the location of derivative discontinuities
            keys = natural_sort(list(dihedral_data[dihedrals].keys()))        
            energies = [ dihedral_data[dihedrals][e]["DFT"] for e in keys ]
            angles   = [ dihedral_data[dihedrals][e]["Dihedrals"][0][1]*180.0/np.pi for e in keys ]
            angles,energies,keys = list(zip(*sorted(zip(angles,energies,keys))))
            d_list   = []

            # Loop over the energies and calculate the second derivatives. 
            # NOTE: the if/else construction is necessary to deal with the endpoints
            for i in range(len(energies)):                

                # Start case
                if i == 0:    
                    # Calculate deltas with care for the pi -> -pi jump
                    delta_1 = angles[i+1] - angles[i]
                    delta_2 = angles[i] - angles[-1]
                    if delta_1 >  180.0: delta_1 = delta_1-360.0
                    if delta_1 < -180.0: delta_1 = delta_1+360.0
                    if delta_2 >  180.0: delta_2 = delta_2-360.0
                    if delta_2 < -180.0: delta_2 = delta_2+360.0

                    #                            forwards-derivative @ i                            forwards-derivative @ i-1       delta angle (i - i-1)                                 
                    d_list += [np.abs(( ( energies[i+1] - energies[i] ) / ( delta_1 ) - ( energies[i] - energies[-1]  ) / ( delta_2 ) ) / (delta_2 ))]

                # End case
                elif i == len(energies)-1:
                    # Calculate deltas with care for the pi -> -pi jump
                    delta_1 = angles[0] - angles[i]
                    delta_2 = angles[i] - angles[i-1]
                    if delta_1 >  180.0: delta_1 = delta_1-360.0
                    if delta_1 < -180.0: delta_1 = delta_1+360.0
                    if delta_2 >  180.0: delta_2 = delta_2-360.0
                    if delta_2 < -180.0: delta_2 = delta_2+360.0

                    #                            forwards-derivative @ i                            forwards-derivative @ i-1       delta angle (i - i-1)                                 
                    d_list += [np.abs(( ( energies[0] - energies[i]   ) / ( delta_1 ) - ( energies[i] - energies[i-1] ) / ( delta_2 ) ) / ( delta_2 ))]

                # Standard case
                else:
                    # Calculate deltas with care for the pi -> -pi jump
                    delta_1 = angles[i+1] - angles[i]
                    delta_2 = angles[i] - angles[i-1]
                    if delta_1 >  180.0: delta_1 = delta_1-360.0
                    if delta_1 < -180.0: delta_1 = delta_1+360.0
                    if delta_2 >  180.0: delta_2 = delta_2-360.0
                    if delta_2 < -180.0: delta_2 = delta_2+360.0

                    #                            forwards-derivative @ i                            forwards-derivative @ i-1       delta angle (i - i-1)                                 
                    d_list += [np.abs(( ( energies[i+1] - energies[i] ) / ( delta_1 ) - ( energies[i] - energies[i-1] ) / ( delta_2 ) ) / (delta_2 ))]

            # print "\nsecond derivatives for {}".format(dihedrals)
            # for count_z,z in enumerate(d_list):
            #     print "{: 20.4f} {: 20.4f}".format(angles[count_z],z)
            # print "mode: {}".format(mode_estimator(d_list))
            second_d_thresh = mode_estimator(d_list)*5.0
            # print "thresh: {}".format(second_d_thresh)

        # Remove datapoints until derivative discontinuities have been removed    
        while clean_flag == 0:

            # Collect the keys, angles, and energies, calculate the second derivatives and delete data points at the location of derivative discontinuities
            keys = natural_sort(list(dihedral_data[dihedrals].keys()))        
            energies = [ dihedral_data[dihedrals][e]["DFT"] for e in keys ]
            angles   = [ dihedral_data[dihedrals][e]["Dihedrals"][0][1]*180.0/np.pi for e in keys ]
            angles,energies,keys = list(zip(*sorted(zip(angles,energies,keys))))
            clean_flag = 1

            # Loop over the energies and calculate the second derivatives. 
            # The largest derivative discontinuitiy is removed at each cycle
            # NOTE: the if/else construction is necessary to deal with the endpoints
            largest = -1.0
            del_ind = 0
            for i in range(len(energies)):                

                # Start case
                if i == 0:    

                    # Calculate deltas with care for the pi -> -pi jump
                    delta_1 = angles[i+1] - angles[i]
                    delta_2 = angles[i] - angles[-1]
                    if delta_1 >  180.0: delta_1 = delta_1-360.0
                    if delta_1 < -180.0: delta_1 = delta_1+360.0
                    if delta_2 >  180.0: delta_2 = delta_2-360.0
                    if delta_2 < -180.0: delta_2 = delta_2+360.0

                    #                            forwards-derivative @ i                           forwards-derivative @ i-1       delta angle (i - i-1)                                 
                    second_d = np.abs(( ( energies[i+1] - energies[i] ) / ( delta_1 ) - ( energies[i] - energies[-1]  ) / ( delta_2 ) ) / (delta_2 ))

                # End case
                elif i == len(energies)-1:

                    # Calculate deltas with care for the pi -> -pi jump
                    delta_1 = angles[0] - angles[i]
                    delta_2 = angles[i] - angles[i-1]
                    if delta_1 >  180.0: delta_1 = delta_1-360.0
                    if delta_1 < -180.0: delta_1 = delta_1+360.0
                    if delta_2 >  180.0: delta_2 = delta_2-360.0
                    if delta_2 < -180.0: delta_2 = delta_2+360.0

                    #                            forwards-derivative @ i                           forwards-derivative @ i-1       delta angle (i - i-1)                                 
                    second_d = np.abs(( ( energies[0] - energies[i]   ) / ( delta_1 ) - ( energies[i] - energies[i-1] ) / ( delta_2 ) ) / ( delta_2 ))

                # Standard case
                else:

                    # Calculate deltas with care for the pi -> -pi jump
                    delta_1 = angles[i+1] - angles[i]
                    delta_2 = angles[i] - angles[i-1]
                    if delta_1 >  180.0: delta_1 = delta_1-360.0
                    if delta_1 < -180.0: delta_1 = delta_1+360.0
                    if delta_2 >  180.0: delta_2 = delta_2-360.0
                    if delta_2 < -180.0: delta_2 = delta_2+360.0

                    #                            forwards-derivative @ i                           forwards-derivative @ i-1       delta angle (i - i-1)                                 
                    second_d = np.abs(( ( energies[i+1] - energies[i] ) / ( delta_1 ) - ( energies[i] - energies[i-1] ) / ( delta_2 ) ) / (delta_2 ))

                # Update the largest second derivative value
                if np.abs(second_d) > largest:
                    largest = second_d
                    if i == 0: del_ind = len(energies)-1
                    else: del_ind = i - 1

            # Remove the configuration with the largest qualifying discontinuity 
            if np.abs(largest) > second_d_thresh:
                # print "removing {} ({}) {} > {}".format(del_ind,angles[del_ind],np.abs(largest),second_d_thresh)
                clean_flag = 0
                del dihedral_data[dihedrals][keys[del_ind]]

        # Rekey the dictionary in case elements were removed
        keys = natural_sort(list(dihedral_data[dihedrals].keys()))        
        for i in range(len(keys)):
            if str(i) != keys[i]:
                dihedral_data[dihedrals][str(i)] = dihedral_data[dihedrals][keys[i]]
                del dihedral_data[dihedrals][keys[i]]
        os.chdir(current_dir)

        # Print final second derivatives, energies, and angles
        # clean_flag = 0
        # second_d_thresh = 0.5
        # print "\nclean second derivatives of {}:".format(dihedrals)
        # keys = natural_sort(dihedral_data[dihedrals].keys())        
        # energies = [ dihedral_data[dihedrals][e]["DFT"] for e in keys ]
        # angles   = [ dihedral_data[dihedrals][e]["Dihedrals"][0][1]*180.0/np.pi for e in keys ]
        # clean_flag = 1
        # for i in range(len(energies)):                
        #     if i == 0:                    
        #         second_d = ( ( energies[i] - energies[-1]  ) / ( angles[i] - angles[-1]  ) - ( energies[i+1] - energies[i] ) / ( angles[i+1] - angles[i] ) ) / (angles[i] - angles[-1]  )
        #     elif i == len(energies)-1:
        #         second_d = ( ( energies[i] - energies[i-1] ) / ( angles[i] - angles[i-1] ) - ( energies[0] - energies[i]   ) / ( angles[0] - angles[i]   ) ) / (angles[i] - angles[i-1] )  
        #     else:
        #         second_d = ( ( energies[i] - energies[i-1] ) / ( angles[i] - angles[i-1] ) - ( energies[i+1] - energies[i] ) / ( angles[i+1] - angles[i] ) ) / (angles[i] - angles[i-1] )
        #     print "{:< 20.6f} {:< 20.6f} {:< 20.6f}".format(second_d,angles[i],energies[i])

    return dihedral_data

# Estimates the mode of a distribution using guassian kernel estimation 
# In this method each datapoint is convolved with a gaussian of width h 
# and the mode is determined by the maximum of the summed guassians.
def mode_estimator(y,h=None,step=None):
    x = deepcopy(y)
    x.sort()

    # If unset, set h to the average separation between datapoints
    if h is None:
        h = np.mean((x-np.roll(x,1))[1:])

    # If unset, set step to a tenth of h
    if step is None:
        step = float(h)/10.0
    search = np.arange(x[0],x[-1],step)
    return search[np.argmax(gaussian_convolution(x,np.ones(len(x)),search,h))]

# function expects vector of X values corresponding to the gaussian centers
# a vector of Y values, corresponding to the area of the gaussian
# a vector of XX values, corresponding to the desired range of gaussian
# evaluation (e.g. np.linspace(xmin,xmax,1000))
# and the gaussian sigma value
def gaussian_convolution(X,Y,XX,Sigma):
    Gaussian_Sum=np.zeros(len(XX));
    for i in range(len(X)):
        for j in range(len(XX)):
            Gaussian_Sum[j]=Y[i]*(Sigma*(2.0*np.pi)**(1./2.))**(-1)*np.exp(-(XX[j]-X[i])**(2)/(2*Sigma)**2)+Gaussian_Sum[j];

    return Gaussian_Sum


# Average over like configurations of identical dihedrals
def average_symmetric(d_data,qc_types,corr_list):

    geo_keys = list(d_data.keys())
    print("{}".format("*"*100))
    print("initial_interpolation data:\n{}".format(d_data))
    # Initialize arrays
    angles_0 = np.zeros(len(geo_keys))

    if "dft" in qc_types:
        DFT_0         = np.zeros(len(geo_keys))
        DFT_Fit_0     = np.zeros(len(geo_keys))
        E_bond_DFT_0  = np.zeros(len(geo_keys))
        E_angle_DFT_0 = np.zeros(len(geo_keys))
    if "mp2" in qc_types:
        MP2_0         = np.zeros(len(geo_keys)) 
        MP2_Fit_0     = np.zeros(len(geo_keys))
        E_bond_MP2_0  = np.zeros(len(geo_keys))
        E_angle_MP2_0 = np.zeros(len(geo_keys))
    E_coul_0 = np.zeros(len(geo_keys))

    # Collect the energies for the main dihedral    
    for count_g,g in enumerate(geo_keys):

        angles_0[count_g] = d_data[g]["Dihedrals"][0][1]
        E_coul_0[count_g] = d_data[g]["E_coul"]

        if "dft" in qc_types:
            DFT_0[count_g]         = d_data[g]["DFT"]
            DFT_Fit_0[count_g]     = d_data[g]["DFT_Fit"]
            E_bond_DFT_0[count_g]  = d_data[g]["E_bond_DFT"]
            E_angle_DFT_0[count_g] = d_data[g]["E_angle_DFT"]

        if "mp2" in qc_types:
            MP2_0[count_g]         = d_data[g]["MP2"]
            MP2_Fit_0[count_g]     = d_data[g]["MP2_Fit"]
            E_bond_MP2_0[count_g]  = d_data[g]["E_bond_MP2"]
            E_angle_MP2_0[count_g] = d_data[g]["E_angle_MP2"]

    # Average over symmetrically identical datapoints
    for i in range(1,len(d_data[geo_keys[0]]["Dihedrals"])):

        # Initialize lists
        angles_N = np.zeros(len(geo_keys))

        if "dft" in qc_types:
            DFT_N         = np.zeros(len(geo_keys))
            DFT_Fit_N     = np.zeros(len(geo_keys))
            E_bond_DFT_N  = np.zeros(len(geo_keys))
            E_angle_DFT_N = np.zeros(len(geo_keys))

        if "mp2" in qc_types:
            MP2_N         = np.zeros(len(geo_keys))
            MP2_Fit_N     = np.zeros(len(geo_keys))
            E_bond_MP2_N  = np.zeros(len(geo_keys))
            E_angle_MP2_N = np.zeros(len(geo_keys))
        E_coul_N = np.zeros(len(geo_keys))

        # Collect energies for the coincident dihedral
        for count_g,g in enumerate(geo_keys):

            # collect angles for the current dihedral
            angles_N[count_g] = d_data[g]["Dihedrals"][i][1]
            if angles_N[count_g] < min(angles_0): angles_N[count_g] = angles_N[count_g]+2.0*np.pi
            elif angles_N[count_g] > max(angles_0): angles_N[count_g] = angles_N[count_g]-2.0*np.pi

            # collect energies for the current dihedral
            if "dft" in qc_types:
                DFT_N[count_g]         = d_data[g]["DFT"]
                DFT_Fit_N[count_g]     = d_data[g]["DFT_Fit"]
                E_bond_DFT_N[count_g]  = d_data[g]["E_bond_DFT"]
                E_angle_DFT_N[count_g] = d_data[g]["E_angle_DFT"]
            if "mp2" in qc_types:
                MP2_N[count_g]         = d_data[g]["MP2"]
                MP2_Fit_N[count_g]     = d_data[g]["MP2_Fit"]
                E_bond_MP2_N[count_g]  = d_data[g]["E_bond_MP2"]
                E_angle_MP2_N[count_g] = d_data[g]["E_angle_MP2"]

            E_coul_N[count_g] = d_data[g]["E_coul"]
        
        # Interpolate the energies on the basis of angles and apply average
        min_0 = [ count_i for count_i,i in enumerate(angles_0) if i == min(angles_0) ][0]
        min_N = [ count_i for count_i,i in enumerate(angles_N) if i == min(angles_N) ][0]
        max_0 = [ count_i for count_i,i in enumerate(angles_0) if i == max(angles_0) ][0]
        max_N = [ count_i for count_i,i in enumerate(angles_N) if i == max(angles_N) ][0]
        angles_N[min_N] = angles_0[min_0]
        angles_N[max_N] = angles_0[max_0]

#         print ""
#         for count_z,z, in enumerate(sorted(angles_0)):
#             print "{} {}".format(z,sorted(angles_N)[count_z])

        E_coul_0 = E_coul_0 + interp_fun(angles_N,E_coul_N,angles_0)
        print(interp_fun(sorted(angles_N),E_coul_N,sorted(angles_0)))
        if "dft" in qc_types:
            DFT_0         = DFT_0         + interp_fun(angles_N,DFT_N,angles_0)
            DFT_Fit_0     = DFT_Fit_0     + interp_fun(angles_N,DFT_Fit_N,angles_0)
            E_bond_DFT_0  = E_bond_DFT_0  + interp_fun(angles_N,E_bond_DFT_N,angles_0)
            E_angle_DFT_0 = E_angle_DFT_0 + interp_fun(angles_N,E_angle_DFT_N,angles_0)
        if "mp2" in qc_types:
            MP2_0         = MP2_0         + interp_fun(angles_N,MP2_N,angles_0)
            MP2_Fit_0     = MP2_Fit_0     + interp_fun(angles_N,MP2_Fit_N,angles_0)
            E_bond_MP2_0  = E_bond_MP2_0  + interp_fun(angles_N,E_bond_MP2_N,angles_0)
            E_angle_MP2_0 = E_angle_MP2_0 + interp_fun(angles_N,E_angle_MP2_N,angles_0)
    
    # Assign new averaged values to the original d_data dictionary
    divisor = float(len(d_data[geo_keys[0]]["Dihedrals"]))
    for count_g,g in enumerate(geo_keys):
        if "dft" in qc_types:
            d_data[g]["DFT"] = DFT_0[count_g]/divisor
            d_data[g]["DFT_Fit"] = DFT_Fit_0[count_g]/divisor
            d_data[g]["E_bond_DFT"] = E_bond_DFT_0[count_g]/divisor
            d_data[g]["E_angle_DFT"] = E_angle_DFT_0[count_g]/divisor
        if "mp2" in qc_types:
            d_data[g]["MP2"] = DFT_0[count_g]/divisor
            d_data[g]["MP2_Fit"] = DFT_Fit_0[count_g]/divisor
            d_data[g]["E_bond_MP2"] = E_bond_DFT_0[count_g]/divisor
            d_data[g]["E_angle_MP2"] = E_angle_DFT_0[count_g]/divisor
    print("{}".format("*"*100))
    print("final_interpolation data:\n{}".format(d_data))
    return d_data

# Wrapper for the array interpolation function. The heart of this is the interpolation object 
# generated by the scipy interpolate.interp2d function.
def interp_fun(x,y,x_new):
    together = sorted([ (i,y[count_i]) for count_i,i in enumerate(x) ])
    x = [ i[0] for i in together ]
    y = [ i[1] for i in together ]
    f = interpolate.interp1d(x,y,kind='linear',bounds_error=False)
    y_new = f(x_new)
    return y_new
        
            
# Check for symmetry amongst coincident atom types
def symmetry_check(dihedral_atoms,dihedral_types):

    a,b,c,d = dihedral_atoms[0][0]

    # Check if all dihedrals share 1 and 4 atom types
    a_set = [dihedral_types[0][0][0]]
    d_set = [dihedral_types[0][0][3]]
    for count_i,i in enumerate(dihedral_atoms):
        if i[0][1] == b:
            a_set += [dihedral_types[count_i][0][0]]
            d_set += [dihedral_types[count_i][0][3]]
        else:
            a_set += [dihedral_types[count_i][0][3]]
            d_set += [dihedral_types[count_i][0][0]]

    # If there is only one type of 1 atomtype and 4 atomtype then the dihedral is symmetric. Return symmetric flag.
    if len(set(a_set)) == 1 and len(set(d_set)) == 1:
        return 1

    # Return non-symmetric flag.
    else:
        return 0

# Function for scraping harmonic parametrization data from output files
def get_harmonic_dihedral_data(list_of_dihedral_types,dihedral_dirs,fit_type='AA',atomtype_lists=[],charge_dict=[],qc_types=[],corr_list=[],scope='scan',one_four_scale_coul=0.0,one_four_scale_vdw=0.0):

    global FF_dict

    # Initialize dictionary
    dihedral_data = {}

    # Loop over all dihedral scans
    current_dir = os.getcwd()
    for count_d,dihedrals in enumerate(list_of_dihedral_types):

        # Initialize subdictionary to hold the data for this scan
        dihedral_data[dihedrals] = {}

        # Change into current dihedral directory
        os.chdir(dihedral_dirs[count_d])

        # Find output files and sort by number (matched to geometries)
        outputs = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ]
        sort_nicely(outputs)

        # Find the optimized geometry files and starting geometries, and sort by number.
        # NOTE: starting_geometries have the atom_type mapping for the configuration which is necessary for fragment based parsing
        all_geometries = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.xyz") ]
        ind = [ count_i for count_i,i in enumerate(all_geometries) if "geo_opt.xyz" in i ]
        geometries = [ i for count_i,i in enumerate(all_geometries) if count_i in ind ]
        starting_geometries = [ i for count_i,i in enumerate(all_geometries) if count_i not in ind ]
        sort_nicely(geometries)        
        sort_nicely(starting_geometries)        

        # Find output files and sort by number (matched to geometries)
        outputs = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ]
        sort_nicely(outputs)

        # Initialize min_ind to keep track of which geometry should be used for the normalization
        min_ind_MP2 = -1
        min_ind_DFT = -1

        # Loop over all discovered geometries and save their geometries to Geo_dict and
        # catenate to a single file for viewing.
        for count_g,g in enumerate(geometries):
            
            # Initialize sub-sub-dictionary for the current geometry's info
            dihedral_data[dihedrals][str(count_g)] = {}

            # Check for job completion based on the discovery of the necessary energies
            dihedral_data[dihedrals][str(count_g)]["Completion"] = 0
            with open(outputs[count_g],'r') as out:                    
                for lines in out:
                    if "****ORCA TERMINATED NORMALLY****" in lines: dihedral_data[dihedrals][str(count_g)]["Completion"] = 1

            # Skip the rest of the parse if not all configurations completed.
            if dihedral_data[dihedrals][str(count_g)]["Completion"] == 0 : skip_flag = 1; os.chdir(current_dir); continue  

            # Parse elements, geometry, and adj_mat (for the first geo)
            Elements, Geo = xyz_parse(g)
            if count_g == 0: Adj_mat = Table_generator(Elements,Geo)

            # Grab a copy of the atomtypes (with the link- labels removed)
            atomtypes_N = [ next( j for j in i.split('link-') if j != '' ) for i in atomtype_lists[count_d] ]
            charges_N   = [ charge_dict[i] for i in atomtypes_N ]

            # Add adj_mat, charge, and atomtype information to the dihedral_data dictionary (used by other functions for generating the MD corrections, when requested)                
            dihedral_data[dihedrals][str(count_g)]["Adj_mat"]   = Adj_mat
            dihedral_data[dihedrals][str(count_g)]["Atomtypes"] = atomtypes_N
            dihedral_data[dihedrals][str(count_g)]["Charges"]   = charges_N

            # Also save the Elements and Geo to the dictionary (used later in fit functions for saving the catenated geometries)
            dihedral_data[dihedrals][str(count_g)]["Elements"] = Elements
            dihedral_data[dihedrals][str(count_g)]["Geo"]      = Geo
            
            # Calculate hybridization (needed for identifying atom types)
            Hybridizations = Hybridization_finder(Elements,Adj_mat)

            # Id atom types in the geometry # UNCOMMENT THIS LINE IF YOU WANT THE PROGRAM TO AUTOMATICALLY ID TYPES, NOT SURE WHY YOU'D WANT TO BUT APPARENTLY I THOUGHT IT WAS A GOOD IDEA ASP
#            Atom_types = id_types(Elements,Adj_mat,2,Hybridizations,Geo)
#            Atom_types = atomtypes
            
            # Find the dihedral types and instances being scanned. Find_scanned_dihedrals also returns tuples with the atom indices and angles, indexed to the 
            # dihedral_data[dihedrals][str(count_g)]["Dihedrals"] list. This functionality was added to average over indentical coincident scans that share the 1,2,3 atoms with the main dihedral. 
            if scope == "scan":
                dihedral_data[dihedrals][str(count_g)]["Dihedrals"],dihedral_data[dihedrals][str(count_g)]["Dihedral_Atoms"] =\
                    Find_scanned_dihedrals(Adj_mat,Geo,atomtypes_N,outputs[count_g],fit_type=fit_type)
            elif scope == "all":
                dihedral_data[dihedrals][str(count_g)]["Dihedrals"] = Find_all_dihedrals(Adj_mat,Geo,atomtypes_N,outputs[count_g],fit_type=fit_type)

            # Parse bonds, angles, and dihedrals (needed for calculating force-field energy correction)
            bond_tuples,angle_tuples,dihedral_tuples,one_five_tuples = Find_modes(Adj_mat,atomtypes_N)
            
            # Calculate force-field corrections
            if 'dft' in qc_types:
                if "bond" in corr_list:
                    dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"]  = E_bond(Geo,atomtypes_N,bond_tuples,QC='DFT')
                else:
                    dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"]  = 0.0
                if "angle" in corr_list:
                    dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"] = E_angle(Geo,atomtypes_N,angle_tuples,QC='DFT')
                else:
                    dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"] = 0.0

            if 'mp2' in qc_types:
                if "bond" in corr_list:
                    dihedral_data[dihedrals][str(count_g)]["E_bond_MP2"]  = E_bond(Geo,atomtypes_N,bond_tuples,QC='MP2')
                else:
                    dihedral_data[dihedrals][str(count_g)]["E_bond_MP2"]  = 0.0
                if "angle" in corr_list:
                    dihedral_data[dihedrals][str(count_g)]["E_angle_MP2"] = E_angle(Geo,atomtypes_N,angle_tuples,QC='MP2')
                else:
                    dihedral_data[dihedrals][str(count_g)]["E_angle_MP2"] = 0.0

            if "coul" in corr_list:
                dihedral_data[dihedrals][str(count_g)]["E_coul"] = E_coul(Geo,Adj_mat,atomtypes_N,charges_N,[bond_tuples,angle_tuples,dihedral_tuples],(0.0,0.0,one_four_scale_coul),fit_type)
            else:
                dihedral_data[dihedrals][str(count_g)]["E_coul"] = 0.0

            if "vdw" in corr_list:
                dihedral_data[dihedrals][str(count_g)]["E_VDW"] = E_LJ(Geo,Adj_mat,atomtypes_N,charges_N,[bond_tuples,angle_tuples,dihedral_tuples],\
                                                                       (0.0,0.0,one_four_scale_vdw),fit_type)
            else:
                dihedral_data[dihedrals][str(count_g)]["E_VDW"] = 0.0

            # Collect energies
            # NOTE: energies are converted to kcal/mol by default and assume hartree output
            with open(outputs[count_g],'r') as out:
                for lines in out:
                    fields = lines.split()

                    if len(fields) == 5 and fields[0].lower() == "final" and fields[1].lower() == "single" and fields[2].lower() == "point" and fields[3].lower() == "energy" and 'dft' in qc_types:
                        dihedral_data[dihedrals][str(count_g)]["DFT"] = float(fields[4])*Htokcalmol
                        dihedral_data[dihedrals][str(count_g)]["DFT_Fit"] = float(fields[4])*Htokcalmol - \
                                                                               dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"] -\
                                                                               dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"] -\
                                                                               dihedral_data[dihedrals][str(count_g)]["E_coul"] -\
                                                                               dihedral_data[dihedrals][str(count_g)]["E_VDW"]
                              
                    if len(fields) == 5 and fields[0].lower() == "mp2" and fields[1].lower() == "total" and fields[2].lower() == "energy:" and 'mp2' in qc_types:
                        dihedral_data[dihedrals][str(count_g)]["MP2"] = float(fields[3])*Htokcalmol 
                        dihedral_data[dihedrals][str(count_g)]["MP2_Fit"] = float(fields[3])*Htokcalmol - \
                                                                               dihedral_data[dihedrals][str(count_g)]["E_bond_MP2"] -\
                                                                               dihedral_data[dihedrals][str(count_g)]["E_angle_MP2"] -\
                                                                               dihedral_data[dihedrals][str(count_g)]["E_coul"] -\
                                                                               dihedral_data[dihedrals][str(count_g)]["E_VDW"]

                        # break after all energies are in hand
                        break
        
            # The energies of each scan are normalized to the minimum fit energy
            if 'dft' in qc_types and (min_ind_DFT == -1 or dihedral_data[dihedrals][str(count_g)]["DFT_Fit"] < dihedral_data[dihedrals][min_ind_DFT]["DFT_Fit"]):
                min_ind_DFT = str(count_g)
            if 'mp2' in qc_types and (min_ind_MP2 == -1 or dihedral_data[dihedrals][str(count_g)]["MP2_Fit"] < dihedral_data[dihedrals][min_ind_MP2]["MP2_Fit"]):
                min_ind_MP2 = str(count_g)

        # Normalize all of the scan energies by the most "trans" dihedral energy
        if 'dft' in qc_types:
            for count_g,g in enumerate(geometries):
                if str(count_g) == min_ind_DFT: continue
                dihedral_data[dihedrals][str(count_g)]["DFT"] = dihedral_data[dihedrals][str(count_g)]["DFT"] - dihedral_data[dihedrals][str(min_ind_DFT)]["DFT"]
                dihedral_data[dihedrals][str(count_g)]["DFT_Fit"] = dihedral_data[dihedrals][str(count_g)]["DFT_Fit"] - dihedral_data[dihedrals][str(min_ind_DFT)]["DFT_Fit"]

        # Normalize all of the scan energies by the most "trans" dihedral energy
        if 'mp2' in qc_types:
            for count_g,g in enumerate(geometries):
                if str(count_g) == min_ind_MP2: continue
                dihedral_data[dihedrals][str(count_g)]["MP2"] = dihedral_data[dihedrals][str(count_g)]["MP2"] - dihedral_data[dihedrals][str(min_ind_MP2)]["MP2"]
                dihedral_data[dihedrals][str(count_g)]["MP2_Fit"] = dihedral_data[dihedrals][str(count_g)]["MP2_Fit"] - dihedral_data[dihedrals][str(min_ind_MP2)]["MP2_Fit"]

        # I encountered bizarre behavior when trying to simply copy the normalization values to separate objects and normalize all data at once
        # I had to resort to normalizing all non_min values then normalize the minimum configuration last
        if 'dft' in qc_types:
            dihedral_data[dihedrals][str(min_ind_DFT)]["DFT"] = dihedral_data[dihedrals][str(min_ind_DFT)]["DFT"] - dihedral_data[dihedrals][str(min_ind_DFT)]["DFT"]
            dihedral_data[dihedrals][str(min_ind_DFT)]["DFT_Fit"] = dihedral_data[dihedrals][str(min_ind_DFT)]["DFT_Fit"] - dihedral_data[dihedrals][str(min_ind_DFT)]["DFT_Fit"]
        if 'mp2' in qc_types:
            dihedral_data[dihedrals][str(min_ind_MP2)]["MP2"] = dihedral_data[dihedrals][str(min_ind_MP2)]["MP2"] - dihedral_data[dihedrals][str(min_ind_MP2)]["MP2"]
            dihedral_data[dihedrals][str(min_ind_MP2)]["MP2_Fit"] = dihedral_data[dihedrals][str(min_ind_MP2)]["MP2_Fit"] - dihedral_data[dihedrals][str(min_ind_MP2)]["MP2_Fit"]
        
        os.chdir(current_dir)
    
#     # Print all energies    
#     for dihedrals in dihedral_dirs:
#         print "\n{}:\n".format(dihedrals)
#         for count_g in range(count_g+1):
#             print "\tdihedral_data[{}][{}]['DFT']: {}".format(dihedrals,str(count_g),dihedral_data[dihedrals][str(count_g)]["DFT"])
#             print "\tdihedral_data[{}][{}]['DFT_Fit']: {}".format(dihedrals,str(count_g),dihedral_data[dihedrals][str(count_g)]["DFT_Fit"])
#             print "\tdihedral_data[{}][{}]['E_bond_DFT']: {}".format(dihedrals,str(count_g),dihedral_data[dihedrals][str(count_g)]["E_bond_DFT"])
#             print "\tdihedral_data[{}][{}]['E_angle_DFT']: {}".format(dihedrals,str(count_g),dihedral_data[dihedrals][str(count_g)]["E_angle_DFT"])
#             print "\tdihedral_data[{}][{}]['E_coul']: {}\n".format(dihedrals,str(count_g),dihedral_data[dihedrals][str(count_g)]["E_coul"])

    return dihedral_data

# This function finds the dihedrals (and their angles) that are fit during each scan
def Find_scanned_dihedrals(adj_mat,Geo,atomtypes,filename,fit_type='AA'):

    # First find the atoms that are constrained (several constraints may be present, but the first corresponds to the dihedral being scanned)
    with open(filename,'r') as out:
        flag = 0
        for lines in out:
            fields = lines.split()
            if (len(fields) == 4 and fields[2].lower() == "%geom" and fields[3].lower() == "constraints") or (len(fields) == 3 and fields[0] == "|" and fields[2].lower()=="constraints"):
                flag = 1
                continue

            # Break after parsing the first constraint
            if flag == 1:
                fields = lines.split("{")[1].split()
                dihedral_atoms = [ (int(fields[1]), int(fields[2]), int(fields[3]), int(fields[4])) ]
                dihedral_types = [ (atomtypes[int(fields[1])],atomtypes[int(fields[2])],atomtypes[int(fields[3])],atomtypes[int(fields[4])]) ]                
                break

#     # For harmonic fits (i.e. dihedrals where the second and third atoms belong to a ring) coincident dihedrals
#     # are not fit simultaneously.
#     if dihedral_types[0][1][0] == "R" and dihedral_types[0][2][0] == "R":
#         left_atoms = []
#         right_atoms = []

    # Iterate through Adj_mat finding coincident dihedrals
#    else:
    left_atoms = [ count_a for count_a,a in enumerate(adj_mat[dihedral_atoms[0][1],:]) if a == 1 and count_a != dihedral_atoms[0][2] ]    
    right_atoms = [ count_a for count_a,a in enumerate(adj_mat[dihedral_atoms[0][2],:]) if a == 1 and count_a != dihedral_atoms[0][1] ]
    ml_atom = dihedral_atoms[0][1]
    mr_atom = dihedral_atoms[0][2]

    # Find all dihedrals (and types) being sampled about this rotation. Leave out dihedrals
    # involving atoms not listed in keep_ind (usually hydrogens if united atom FF params are desired.)
    for l in left_atoms:
        for r in right_atoms:

            # avoid repetitions (probably a redunant precaution)
            if (l,ml_atom,mr_atom,r) in dihedral_atoms or (r,mr_atom,ml_atom,l) in dihedral_atoms:
                continue
            else:

                # dihedral types are written so that the lesser *atom_type* between 1 and 4 is first.
                # In the event that 1 and 4 are of the same type, then the lesser of 2 and 3 goes first
                if atomtypes[l] == atomtypes[r]:
                    if atomtypes[ml_atom] <= atomtypes[mr_atom]:
                        dihedral_types += [(atomtypes[l],atomtypes[ml_atom],atomtypes[mr_atom],atomtypes[r])]
                        dihedral_atoms += [(l,ml_atom,mr_atom,r)]        
                    else:
                        dihedral_types += [(atomtypes[r],atomtypes[mr_atom],atomtypes[ml_atom],atomtypes[l])]
                        dihedral_atoms += [(r,mr_atom,ml_atom,l)]        
                elif atomtypes[l] < atomtypes[r]:
                    dihedral_types += [(atomtypes[l],atomtypes[ml_atom],atomtypes[mr_atom],atomtypes[r])]
                    dihedral_atoms += [(l,ml_atom,mr_atom,r)]        
                else:
                    dihedral_types += [(atomtypes[r],atomtypes[mr_atom],atomtypes[ml_atom],atomtypes[l])]            
                    dihedral_atoms += [(r,mr_atom,ml_atom,l)]                            

    # Add the dihedrals and types, and canonicalize all types (Above loops should be removed at some point and everything should utilize these canonicalization functions)
    for count_i,i in enumerate(dihedral_atoms):
        dihedral_types[count_i],dihedral_atoms[count_i] = canon_dihedral(dihedral_types[count_i],dihedral_atoms[count_i])

    # Remove hydrogen containing dihedrals if type == 'UA'
    if fit_type == "UA":

        H_ind = [ count_i for count_i,i in enumerate(atomtypes) if int(i.split('[')[1].split(']')[0]) == 1 ]

        # Remove hydrogens from the H_ind list that aren't attached to sp3 carbon
        del_list = []
        for count_i,i in enumerate(H_ind):
            if len([ a for count_a,a in enumerate(adj_mat[i]) if a == 1 and int(atomtypes[count_a].split('[')[1].split(']')[0]) == 6 and int(sum(adj_mat[count_a])) == 4 ]) == 0:
                del_list += [count_i]
        
        # Save hydrogen types that are attached to carbon to a list for comparisons
        H_types = [ atomtypes[i] for count_i,i in enumerate(H_ind) if count_i not in del_list ]            

        # Remove dihedrals involving UA-hydrogens
        del_list       = [ count_i for count_i,i in enumerate(dihedral_types) if i[0] in H_types or i[3] in H_types ]
        dihedral_atoms = [ i for count_i,i in enumerate(dihedral_atoms) if count_i not in del_list ]
        dihedral_types = [ (i[0]+'-UA',i[1]+'-UA',i[2]+'-UA',i[3]+'-UA') for count_i,i in enumerate(dihedral_types) if count_i not in del_list ]

    # Calculate all dihedral angles
    for count_x,x in enumerate(dihedral_atoms):
       atom_1 = Geo[x[0]]
       atom_2 = Geo[x[1]]
       atom_3 = Geo[x[2]]
       atom_4 = Geo[x[3]]
       v1 = atom_2-atom_1
       v2 = atom_3-atom_2
       v3 = atom_4-atom_3
       angle = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )
       dihedral_types[count_x] = (dihedral_types[count_x],angle)
       dihedral_atoms[count_x] = (x,angle)

    return dihedral_types,dihedral_atoms

def Find_all_dihedrals(adj_mat,Geo,atomtypes,filename,fit_type='AA'):

    dihedrals = []
    
    # Each row in the adjacency matrix is used to seed a 4-bond deep search
    for count_i,i in enumerate(adj_mat):

        # current holds the dihedral seedlings
        current = [ [count_i,count_j] for count_j,j in enumerate(i) if j == 1 ]

        # new holds the updated list of dihedrals at the end of each bond search
        new = []
        iteration = 0 

        # recursively add bonds to new until no new bonds are found
        while new != current:
            if iteration == 0: new = deepcopy(current)
            current = deepcopy(new)

            # Iterate over the seedlings and add connections/spawn new seedlings
            for count_j,j in enumerate(current):
                if len(j) == 4: continue                                                                            # avoid adding new elements to full dihedrals
                connections = [ count_k for count_k,k in enumerate(adj_mat[j[-1]]) if k == 1 and count_k not in j ] # add new connections not already in the dihedral
                for count_k,k in enumerate(connections):
                    if count_k == 0: new[count_j] += [k]                                                            # add the first connection to the existing seedling
                    else: new += [ j + [k] ]                                                                        # add the rest as new elements
            iteration += 1             

        # Only add full dihedrals and new dihedrals to the dihedrals list
        tmp = [ j for j in new if len(j) == 4 ]
        dihedrals += [ j for j in tmp if ( j not in dihedrals and [j[3],j[2],j[1],j[0]] not in dihedrals ) ]

    # Add the dihedrals and types, and canonicalize all types. 
    dihedral_atoms = [ tuple(i) for i in dihedrals ]
    dihedral_types = [ (atomtypes[i[0]],atomtypes[i[1]],atomtypes[i[2]],atomtypes[i[3]]) for i in dihedral_atoms ]
    for count_i,i in enumerate(dihedral_atoms):
        dihedral_types[count_i],dihedral_atoms[count_i] = canon_dihedral(dihedral_types[count_i],dihedral_atoms[count_i])

    # Remove hydrogen containing dihedrals if type == 'UA'
    if fit_type == "UA":

        H_ind = [ count_i for count_i,i in enumerate(atomtypes) if int(i.split('[')[1].split(']')[0]) == 1 ]

        # Remove hydrogens that aren't attached to carbon sp3 carbon
        del_list = []
        for count_i,i in enumerate(H_ind):
            if len([ a for count_a,a in enumerate(adj_mat[i]) if a == 1 and int(atomtypes[count_a].split('[')[1].split(']')[0]) == 6 and int(sum(adj_mat[count_a])) == 4 ]) == 0:
                del_list += [count_i]

        # Save hydrogen types that are attached to carbon to a list for comparisons
        H_types = [ atomtypes[i] for count_i,i in enumerate(H_ind) if count_i not in del_list ]            

        # Remove dihedrals involving UA-hydrogens
        del_list       = [ count_i for count_i,i in enumerate(dihedral_types) if i[0] in H_types or i[3] in H_types ]
        dihedral_atoms = [ i for count_i,i in enumerate(dihedral_atoms) if count_i not in del_list ]
        dihedral_types = [ (i[0]+'-UA',i[1]+'-UA',i[2]+'-UA',i[3]+'-UA') for count_i,i in enumerate(dihedral_types) if count_i not in del_list ]

    # Calculate all dihedral angles
    for count_x,x in enumerate(dihedral_atoms):
       atom_1 = Geo[x[0]]
       atom_2 = Geo[x[1]]
       atom_3 = Geo[x[2]]
       atom_4 = Geo[x[3]]
       v1 = atom_2-atom_1
       v2 = atom_3-atom_2
       v3 = atom_4-atom_3
       angle = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )
       dihedral_types[count_x] = (dihedral_types[count_x],angle)
       dihedral_atoms[count_x] = (x,angle)

    return dihedral_types,dihedral_atoms


def E_bond(Geo,atomtypes,bond_tuples,QC='DFT'):

    E_bonds_tot = 0.0            # Initialize the total coulombic energy
    seps = cdist(Geo,Geo)        # Initialize the separation matrix

    # Iterate over all bonds
    for i in bond_tuples:

        # For fragment based parametrizations, any modes involving the added hydrogens (conventionally [0] types) are omitted.
        if "[0]" in [ atomtypes[i[0]], atomtypes[i[1]] ] or "[1]" in [ atomtypes[i[0]], atomtypes[i[1]] ]:
            continue

        # Check for key ordering and add the energy to the total
        if (atomtypes[i[0]],atomtypes[i[1]]) in list(FF_dict["bonds"].keys()): key = (atomtypes[i[0]],atomtypes[i[1]])
        else: key = (atomtypes[i[1]],atomtypes[i[0]])


        # If key isn't in the FF_dict then print an error
        if key not in list(FF_dict["bonds"].keys()):
            print("ERROR: Force-field information for bond {} was not found. Cannot calculate bond correction to the QC potential. Exiting...".format(key))
            quit()

        # Add energy to the total (note: functional type is the first element in the FF_dict["bonds"], for harmonic, k is the second element, eq_bond length is the third
        E_bonds_tot += FF_dict["bonds"][key][QC][1]*(seps[i[0],i[1]]-FF_dict["bonds"][key][QC][2])**2

    return E_bonds_tot

def E_angle(Geo,atomtypes,angle_tuples,QC='DFT'):

    E_angles_tot = 0.0           # Initialize the total coulombic energy

    # Interate over all angles
    for i in angle_tuples:

        # For fragment based parametrizations, any modes involving the added hydrogens (conventionally [0] types) are omitted.
        if "[0]" in [ atomtypes[i[0]], atomtypes[i[1]],atomtypes[i[2]] ] or "[1]" in [ atomtypes[i[0]], atomtypes[i[1]], atomtypes[i[2]] ]:
            continue

        # Check for key ordering and add the energy to the total
        if (atomtypes[i[0]],atomtypes[i[1]],atomtypes[i[2]]) in list(FF_dict["angles"].keys()): key = (atomtypes[i[0]],atomtypes[i[1]],atomtypes[i[2]])
        else: key = (atomtypes[i[2]],atomtypes[i[1]],atomtypes[i[0]])

        # If key isn't in the FF_dict then print an error
        if key not in list(FF_dict["angles"].keys()):
            print("ERROR: Force-field information for angle {} was not found. Cannot calculate angle correction to QC potential. Exiting...".format(key))
            quit()
        
        # Calculate angle
        atom_1 = Geo[i[0]]
        atom_2 = Geo[i[1]]
        atom_3 = Geo[i[2]]
        theta = np.arccos(np.dot(atom_1-atom_2,atom_3-atom_2)/(norm(atom_1-atom_2)*norm(atom_3-atom_2)))

        # Add energy to the total (note: functional type is the first element in the FF_dict["angle"], for harmonic, k is the second element, eq_angle is the third
        E_angles_tot += FF_dict["angles"][key][QC][1]*(theta-FF_dict["angles"][key][QC][2]*np.pi/180.0)**2

    return E_angles_tot

def E_dihedral(Geo,atomtypes,dihedral_tuples,QC='DFT'):

    E_dihedral_tot = 0.0  # Initialize the total coulombic energy

    # Interate over all dihedrals
    for i in dihedral_tuples:

        # For fragment based parametrizations, any modes involving the added hydrogens (conventionally [0] types) are omitted.
        if "[0]" in [ atomtypes[i[0]], atomtypes[i[1]], atomtypes[i[2]], atomtypes[i[3]] ] or "[1]" in [ atomtypes[i[0]], atomtypes[i[1]], atomtypes[i[2]], atomtypes[i[3]] ]:
            continue

        # Check for key ordering and add the energy to the total
        if (atomtypes[i[0]],atomtypes[i[1]],atomtypes[i[2]],atomtypes[i[3]]) in (list(FF_dict["dihedrals"].keys())+list(FF_dict["dihedrals_harmonic"].keys())): key = (atomtypes[i[0]],atomtypes[i[1]],atomtypes[i[2]],atomtypes[i[3]])
        else: key = (atomtypes[i[3]],atomtypes[i[2]],atomtypes[i[1]],atomtypes[i[0]])

        # Calculate all dihedral angles
        atom_1 = Geo[i[0]]
        atom_2 = Geo[i[1]]
        atom_3 = Geo[i[2]]
        atom_4 = Geo[i[3]]
        v1 = atom_2-atom_1
        v2 = atom_3-atom_2
        v3 = atom_4-atom_3
        theta = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )

        # Use OPLS expression is if is a non-harmonic dihedral:
        if key in list(FF_dict["dihedrals"].keys()):

            # Add energy to the total (note: functional type is the first element in the FF_dict["angle"], for harmonic, k is the second element, eq_angle is the third
            E_dihedral_tot += 0.5*FF_dict["dihedrals"][key][QC][1]*(1+np.cos(theta)) + 0.5*FF_dict["dihedrals"][key][QC][2]*(1-np.cos(2.0*theta)) +\
                               0.5*FF_dict["dihedrals"][key][QC][3]*(1+np.cos(3.0*theta)) + 0.5*FF_dict["dihedrals"][key][QC][4]*(1-np.cos(4.0*theta))

        elif key in FF_dict["dihedrals_harmonic"] and FF_dict["dihedrals_harmonic"][key][QC][0] == "harmonic":

            # Add energy to the total (note: functional type is the first element in the FF_dict["angle"], for quadratic, k is the second element, eq_angle is the third
            E_dihedral_tot += FF_dict["dihedrals_harmonic"][key][QC][1]*(1+FF_dict["dihedrals_harmonic"][key][QC][2]*np.cos(theta))

        elif key in FF_dict["dihedrals_harmonic"] and FF_dict["dihedrals_harmonic"][key][QC][0] == "quadratic":

            # Add energy to the total (note: functional type is the first element in the FF_dict["angle"], for quadratic, k is the second element, eq_angle is the third
            E_dihedral_tot += 0.0  # THESE AREN'T PROPERLY IMPLEMENTED YET
#            E_dihedral_tot += FF_dict["dihedrals_harmonic"][key][QC][1] * ( theta * 180.0/np.pi - FF_dict["dihedrals_harmonic"][key][QC][2] )**(2.0)
            
        else:
            print("ERROR: Force-field information for dihedral {} was not found. Cannot calculate dihedral correction to the QC potential. Exiting...".format(key))
            quit()
        


    return E_dihedral_tot

# This is the newer version with support for scaled 1-2 1-3 and 1-4 interactions
# The function expects connections to be a list of lists with 1-2 pairs as the first list, 1-3 pairs as the second, and 1-4 pairs as the third
# scale_factors is a tuple or list with the 1-2, 1-3, and 1-4 scaling factors as the first, second, and third elements, respectively. 
def E_coul(Geo,Adj_mat,atomtypes,charges,connections,scale_factors=(0.0,0.0,0.0),fit_type='AA'):

    # For UA types the UA-H charges are added to the corresponding UA-carbons
    if fit_type == "UA":
        
        # Add hydrogen charges to bonded atoms
        for count_i,i in enumerate(Adj_mat):
            
            # Add UA-hydrogens to carbon atoms
            if int(atomtypes[count_i].split('[')[1].split(']')[0]) == 6 and sum(i) == 4:
                H_ind = [ count_j for count_j,j in enumerate(i) if j == 1 and int(atomtypes[count_j].split('[')[1].split(']')[0]) == 1 ]            
                charges[count_i] += sum([ charges[j] for j in H_ind ])

                # Set hydrogen charges to zero
                for j in H_ind: 
                    charges[j] = 0.0            
    
    # Calculate the electrostatic energy (neglecting 1-2,1-3,1-4 (possibly also 1-5) interactions)
    E_coul_tot = 0.0                                           # Initialize the total coulombic energy
    seps = cdist(Geo,Geo)                                      # Initialize the separation matrix
    for count_i,i in enumerate(seps):                          # Iterate over all rows in the separation matrix

        # Initialize the special bonds lists
        one_two_list = []
        one_three_list = []
        one_four_list = []

        # Find 1-2 interactions
        for j in connections[0]:
            if count_i in j: one_two_list += [j[0],j[1]]

        # Find 1-3 interactions
        for j in connections[1]:
            if count_i == j[0] or count_i == j[2]: one_three_list += [j[0],j[2]]

        # Find 1-4 interactions
        for j in connections[2]:
            if count_i == j[0] or count_i == j[3]: one_four_list += [j[0],j[-1]]

        # Loop over pairs (i.e. elements in each seps row)
        for count_j,j in enumerate(i):
            if count_j > count_i:

                # Check for 1-2 condition
                if count_j in one_two_list:
                    E_coul_tot += charges[count_i]*charges[count_j]/j*Ecoul_Const * scale_factors[0]
                # Check for 1-3 condition
                elif count_j in one_three_list:
                    E_coul_tot += charges[count_i]*charges[count_j]/j*Ecoul_Const * scale_factors[1]
                # Check for 1-4 condition
                elif count_j in one_four_list:
                    E_coul_tot += charges[count_i]*charges[count_j]/j*Ecoul_Const * scale_factors[2]
                # Default to unscaled interaction
                else:
                    E_coul_tot += charges[count_i]*charges[count_j]/j*Ecoul_Const

    return E_coul_tot

# Calculates the LJ energy of a configuration neglecting 1-2 1-3 and 1-4 interactions
def E_LJ(Geo,Adj_mat,atomtypes_0,charges,connections,scale_factors=(0.0,0.0,0.0),fit_type='AA',LJ_dict=None):
    
    global FF_dict

    # Use the global LJ parameters if no specific LJ dictionary is supplied
    if LJ_dict is None:
        LJ_dict = FF_dict["vdw"]

    # A new atomtypes list is used here to avoid altering the supplied list (python sometimes allows functions to alter unreturned objects)
    if fit_type == "UA":

        # Append UA labels and initialize UA-H list
        atomtypes = [ i+'-UA' for i in atomtypes_0 ]
        H_ind = []

        # Determine UA-H based on attachments to carbon
        for count_i,i in enumerate(Adj_mat):

            # Add UA-hydrogens to carbon atoms
            if int(atomtypes[count_i].split('[')[1].split(']')[0]) == 6 and sum(i) == 4:
                H_ind += [ count_j for count_j,j in enumerate(i) if j == 1 and int(atomtypes[count_j].split('[')[1].split(']')[0]) == 1 ]

    # Else, initialize atomtypes list and an empty UA-H list
    else:
        atomtypes =  np.copy(atomtypes_0)
        H_ind = []

    # Calculate the LJ energy (neglecting 1-2,1-3,1-4 interactions)
    E_LJ_tot = 0.0                                             # Initialize the total LJ
    seps = cdist(Geo,Geo)                                      # Initialize the separation matrix
    for count_i,i in enumerate(seps):                          # Iterate over all rows in the separation matrix

        # For UA parametrizations H_ind holds the indices of UA-H atoms
        if count_i in H_ind: continue                          

        # Initialize the special bonds lists
        one_two_list = []
        one_three_list = []
        one_four_list = []

        # Find 1-2 interactions
        for j in connections[0]:
            if count_i in j: one_two_list += [j[0],j[1]]

        # Find 1-3 interactions
        for j in connections[1]:
            if count_i == j[0] or count_i == j[2]: one_three_list += [j[0],j[2]]

        # Find 1-4 interactions
        for j in connections[2]:
            if count_i == j[0] or count_i == j[3]: one_four_list += [j[0],j[-1]]

        # Perform summation while avoiding UA-H types (if fit_type=="UA") and 1-2 1-3 and 1-4 interactions (via avoid list)
        for count_j,j in enumerate(i):
            if count_j > count_i and count_j not in H_ind:
                pair_type = (atomtypes[count_i],atomtypes[count_j])
                if "[0]" in pair_type or "[1]" in pair_type: continue
                elif "[0]-UA" in pair_type or "[1]-UA" in pair_type: continue

                # Check for 1-2 condition
                elif count_j in one_two_list:
                    E_LJ_tot += 4.0*LJ_dict[pair_type][1]*( (LJ_dict[pair_type][2]/j)**(12.0) - (LJ_dict[pair_type][2]/j)**(6.0) ) * scale_factors[0]

                # Check for 1-3 condition
                elif count_j in one_three_list:
                    E_LJ_tot += 4.0*LJ_dict[pair_type][1]*( (LJ_dict[pair_type][2]/j)**(12.0) - (LJ_dict[pair_type][2]/j)**(6.0) ) * scale_factors[1]

                # Check for 1-4 condition
                elif count_j in one_four_list:
                    E_LJ_tot += 4.0*LJ_dict[pair_type][1]*( (LJ_dict[pair_type][2]/j)**(12.0) - (LJ_dict[pair_type][2]/j)**(6.0) ) * scale_factors[2]

                # Default to unscaled interaction
                else:
                    E_LJ_tot += 4.0*LJ_dict[pair_type][1]*( (LJ_dict[pair_type][2]/j)**(12.0) - (LJ_dict[pair_type][2]/j)**(6.0) )
    return E_LJ_tot

def get_dihedral_data_scan_min(dihedral_dirs,fit_type='AA',charges=[],one_four_scale_coul=0.0):

    # Initialize dictionary
    dihedral_data = {}

    delta_180 = [1E20]*len(dihedral_dirs)
    min_ind   = [(0,0)]*len(dihedral_dirs)

    # Loop over all dihedral scans
    for count_d,dihedrals in enumerate(dihedral_dirs):

        # Change into current dihedral directory
        os.chdir(dihedrals)

        # Find the optimized geometry files, and sort by number.
        geometries = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*geo_opt.xyz") ]
        sort_nicely(geometries)        
    
        # Find output files and sort by number (matched to geometries)
        outputs = [ os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,"*.out") ]
        sort_nicely(outputs)

        # Loop over all discovered geometries and save their geometries to Geo_dict and
        # catenate to a single file for viewing.
        for count_g,g in enumerate(geometries):
            dihedral_data[dihedrals+'-'+str(count_g)] = {}
            Elements, Geo = xyz_parse(g)
            if count_g == 0: Adj_mat = Table_generator(Elements,Geo)
            Hybridizations = Hybridization_finder(Elements,Adj_mat)
            Atom_types = id_types(Elements,Adj_mat,2,Hybridizations,Geo)
#            dihedral_data[dihedrals+'-'+str(count_g)]["Dihedrals"] = Find_dihedrals(Adj_mat,Geo,Atom_types,fit_type)
            dihedral_data[dihedrals+'-'+str(count_g)]["Dihedrals"],bond_tuples,angle_tuples,dihedral_tuples, = Find_dihedrals(Adj_mat,Geo,Atom_types,fit_type)
            dihedral_data[dihedrals+'-'+str(count_g)]["E_Coul"] = E_coul(Geo,charges,[bond_tuples,angle_tuples,dihedral_tuples],(0.0,0.0,one_four_scale_coul))

            # The energies of each scan are normalized to the energy of the closest dihedral to 180 #
            for count_i,i in enumerate(dihedral_data[dihedrals+'-'+str(count_g)]["Dihedrals"]):
                if 180.0-np.abs(i[1]*180.0/np.pi)<delta_180[count_d]:
                    delta_180[count_d] = 180.0-np.abs(i[1])
                    min_ind[count_d]   = (count_g,count_i)

            # Collect energies
            # NOTE: energies are converted to kcal/mol by default and assume hartree output
            with open(outputs[count_g],'r') as out:
                for lines in out:
                    fields = lines.split()
                    if len(fields) == 5 and fields[0].lower() == "final" and fields[1].lower() == "single" and fields[2].lower() == "point" and fields[3].lower() == "energy":
#                        dihedral_data[dihedrals+'-'+str(count_g)]["DFT"] = float(fields[4])*Htokcalmol
                        dihedral_data[dihedrals+'-'+str(count_g)]["DFT"] = float(fields[4])*Htokcalmol - dihedral_data[dihedrals+'-'+str(count_g)]["E_Coul"]                        
                    if len(fields) == 5 and fields[0].lower() == "mp2" and fields[1].lower() == "total" and fields[2].lower() == "energy:":
#                        dihedral_data[dihedrals+'-'+str(count_g)]["MP2"] = float(fields[3])*Htokcalmol
                        dihedral_data[dihedrals+'-'+str(count_g)]["MP2"] = float(fields[3])*Htokcalmol - dihedral_data[dihedrals+'-'+str(count_g)]["E_Coul"]

        os.chdir('..')

    # Normalize all energies by the minimum configuration in each scan
    for count_d,dihedrals in enumerate(dihedral_dirs):
        DFT_correction = dihedral_data[dihedrals+'-'+str(min_ind[count_d][0])]["DFT"]
        MP2_correction = dihedral_data[dihedrals+'-'+str(min_ind[count_d][0])]["MP2"]
        for j in range(count_g+1):
            dihedral_data[dihedrals+'-'+str(j)]["DFT"] -= DFT_correction
            dihedral_data[dihedrals+'-'+str(j)]["MP2"] -= MP2_correction

    # Print all energies
    for dihedrals in dihedral_dirs:
        print("\n{}:\n".format(dihedrals))
        for count_g in range(count_g+1):
            print("\tdihedral_data[{}]['DFT']: {}".format(dihedrals+'-'+str(count_g),dihedral_data[dihedrals+'-'+str(count_g)]["DFT"]))
            print("\tdihedral_data[{}]['E_Coul']: {}\n".format(dihedrals+'-'+str(count_g),dihedral_data[dihedrals+'-'+str(count_g)]["E_Coul"]))

    return dihedral_data


# Description: Simple wrapper function for grabbing the coordinates and
#              elements from an xyz file
#
# Inputs      input: string holding the filename of the xyz
# Returns     Elements: list of element types (list of strings)
#             Geometry: Nx3 array holding the cartesian coordinates of the
#                       geometry (atoms are indexed to the elements in Elements)
#
def xyz_parse(input):

    # Open file and read contents into variable
    with open(input,'r') as f:
            content=f.readlines()

    # Find number of atoms and initialize various matrices
    Atom_Number = int(content[0].split()[0])
    Elements = ['']*Atom_Number
    Geometry = np.zeros([Atom_Number,3])

    # Iterate over the remainder of contents and read the
    # geometry and elements into variable. Note that the
    # first two lines are considered a header
    count=0
    for lines in content[2:]:
        fields=lines.split()

        # Skip empty lines
        if len(fields) == 0:
            continue

        # Write geometry containing lines to variable
        if len(fields) > 3:
            Elements[count]=fields[0]
            Geometry[count,:]=np.array([float(fields[1]),float(fields[2]),float(fields[3])])
            count = count + 1

    return Elements,Geometry

# A wrapper for the commands to parse the dihedrals from the adjacency matrix and geometry.
#           Atom_types isn't necessary here, this section of code just hasn't been cleaned up.
# Returns:  list of (dihedral_type,angle) tuples. 
def Find_modes(Adj_mat,Atom_types,return_all=0):

    # Initialize lists of each instance and type of FF object.
    # instances are stored as tuples of the atoms involved 
    # (e.g., bonds between atoms 1 and 13 and 17 and 5 would be stored as [(1,13),(17,5)] 
    # Similarly, types are stored as tuples of atom types.
    Atom_types = [ next( j for j in i.split('link-') if j != '' ) for i in Atom_types ]    #  split('-link') call is necessary for handling fragment atoms
    Bonds = []
    Bond_types = []
    Angles = []
    Angle_types = []
    Dihedrals = []
    Dihedral_types = []
    One_fives = []
    VDW_types = []

    # Find bonds #
    for count_i,i in enumerate(Adj_mat):        
        Tmp_Bonds = [ (count_i,count_j) for count_j,j in enumerate(i) if j == 1 and count_j > count_i ]

        # Store bond tuple so that lowest atom *type* between the first and the second atom is placed first
        # and avoid redundant placements
        for j in Tmp_Bonds:
            if Atom_types[j[1]] < Atom_types[j[0]] and (j[1],j[0]) not in Bonds and (j[0],j[1]) not in Bonds:
                Bonds = Bonds + [ (j[1],j[0]) ]
                Bond_types = Bond_types + [ (Atom_types[j[1]],Atom_types[j[0]]) ]
            elif (j[0],j[1]) not in Bonds and (j[1],j[0]) not in Bonds:
                Bonds = Bonds + [ (j[0],j[1]) ]
                Bond_types = Bond_types + [ (Atom_types[j[0]],Atom_types[j[1]]) ]

    # Remove -UA tag from Bond_types (united-atom has no meaning for bonds)
    Bond_types = [ (i[0].split('-UA')[0],i[1].split('-UA')[0]) for i in Bond_types ]

    # Find angles #
    for i in Bonds:        

        # Find angles based on connections to first index of Bonds
        Tmp_Angles = [ (count_j,i[0],i[1]) for count_j,j in enumerate(Adj_mat[i[0]]) if j == 1 and count_j != i[1] ]

        # Store angle tuple so that lowest atom *type* between the first and the third is placed first
        # and avoid redundant placements
        for j in Tmp_Angles:
            if Atom_types[j[2]] < Atom_types[j[0]] and (j[2],j[1],j[0]) not in Angles and (j[0],j[1],j[2]) not in Angles:
                Angles = Angles + [(j[2],j[1],j[0])]
                Angle_types = Angle_types + [ (Atom_types[j[2]],Atom_types[j[1]],Atom_types[j[0]]) ]
            elif (j[0],j[1],j[2]) not in Angles and (j[2],j[1],j[0]) not in Angles:
                Angles = Angles + [(j[0],j[1],j[2])]
                Angle_types = Angle_types + [ (Atom_types[j[0]],Atom_types[j[1]],Atom_types[j[2]]) ]

        # Find angles based on connections to second index of Bonds
        Tmp_Angles = [ (i[0],i[1],count_j) for count_j,j in enumerate(Adj_mat[i[1]]) if j == 1 and count_j != i[0] ]

        # Store angle tuple so that lowest atom *type* between the first and the third is placed first
        # and avoid redundant placements
        for j in Tmp_Angles:
            if Atom_types[j[2]] < Atom_types[j[0]] and (j[2],j[1],j[0]) not in Angles and (j[0],j[1],j[2]) not in Angles:
                Angles = Angles + [(j[2],j[1],j[0])]
                Angle_types = Angle_types + [ (Atom_types[j[2]],Atom_types[j[1]],Atom_types[j[0]]) ]
            elif (j[0],j[1],j[2]) not in Angles and (j[2],j[1],j[0]) not in Angles:
                Angles = Angles + [(j[0],j[1],j[2])]
                Angle_types = Angle_types + [ (Atom_types[j[0]],Atom_types[j[1]],Atom_types[j[2]]) ]

    # Remove -UA tag from Angle_types (united-atom has no meaning for angles)
    Angle_types = [ (i[0].split('-UA')[0],i[1].split('-UA')[0],i[2].split('-UA')[0]) for i in Angle_types ]
        
    # Find dihedrals #
    for i in Angles:
        
        # Find atoms attached to first atom of each angle
        Tmp_Dihedrals = [ (count_j,i[0],i[1],i[2]) for count_j,j in enumerate(Adj_mat[i[0]]) if j == 1 and count_j not in [i[1],i[2]] ]
        
        # Store dihedral tuple so that the lowest atom *type* between the first and fourth is placed first
        # and avoid redundant placements        
        for j in Tmp_Dihedrals:

            # If the first and fourth atoms are equal, then sorting is based on the second and third
            if Atom_types[j[3]] == Atom_types[j[0]] and (j[3],j[2],j[1],j[0]) not in Dihedrals and (j[0],j[1],j[2],j[3]) not in Dihedrals:
                if Atom_types[j[2]] < Atom_types[j[1]]:
                    Dihedrals = Dihedrals + [(j[3],j[2],j[1],j[0])]
                    Dihedral_types = Dihedral_types + [ (Atom_types[j[3]],Atom_types[j[2]],Atom_types[j[1]],Atom_types[j[0]]) ]
                else:
                    Dihedrals = Dihedrals + [(j[0],j[1],j[2],j[3])]
                    Dihedral_types = Dihedral_types + [ (Atom_types[j[0]],Atom_types[j[1]],Atom_types[j[2]],Atom_types[j[3]]) ]

            elif Atom_types[j[3]] < Atom_types[j[0]] and (j[3],j[2],j[1],j[0]) not in Dihedrals and (j[0],j[1],j[2],j[3]) not in Dihedrals:
                Dihedrals = Dihedrals + [(j[3],j[2],j[1],j[0])]
                Dihedral_types = Dihedral_types + [ (Atom_types[j[3]],Atom_types[j[2]],Atom_types[j[1]],Atom_types[j[0]]) ]
            elif (j[0],j[1],j[2],j[3]) not in Dihedrals and (j[3],j[2],j[1],j[0]) not in Dihedrals:
                Dihedrals = Dihedrals + [(j[0],j[1],j[2],j[3])]
                Dihedral_types = Dihedral_types + [ (Atom_types[j[0]],Atom_types[j[1]],Atom_types[j[2]],Atom_types[j[3]]) ]

        # Find atoms attached to the third atom of each angle
        Tmp_Dihedrals = [ (i[0],i[1],i[2],count_j) for count_j,j in enumerate(Adj_mat[i[2]]) if j == 1 and count_j not in [i[0],i[1]] ]
        
        # Store dihedral tuple so that the lowest atom *type* between the first and fourth is placed first
        # and avoid redundant placements        
        for j in Tmp_Dihedrals:

            # If the first and fourth atoms are equal, then sorting is based on the second and third
            if Atom_types[j[3]] == Atom_types[j[0]] and (j[3],j[2],j[1],j[0]) not in Dihedrals and (j[0],j[1],j[2],j[3]) not in Dihedrals:
                if Atom_types[j[2]] < Atom_types[j[1]]:
                    Dihedrals = Dihedrals + [(j[3],j[2],j[1],j[0])]
                    Dihedral_types = Dihedral_types + [ (Atom_types[j[3]],Atom_types[j[2]],Atom_types[j[1]],Atom_types[j[0]]) ]
                else:
                    Dihedrals = Dihedrals + [(j[0],j[1],j[2],j[3])]
                    Dihedral_types = Dihedral_types + [ (Atom_types[j[0]],Atom_types[j[1]],Atom_types[j[2]],Atom_types[j[3]]) ]

            elif Atom_types[j[3]] < Atom_types[j[0]] and (j[3],j[2],j[1],j[0]) not in Dihedrals and (j[0],j[1],j[2],j[3]) not in Dihedrals:
                Dihedrals = Dihedrals + [(j[3],j[2],j[1],j[0])]
                Dihedral_types = Dihedral_types+ [ (Atom_types[j[3]],Atom_types[j[2]],Atom_types[j[1]],Atom_types[j[0]]) ]
            elif (j[0],j[1],j[2],j[3]) not in Dihedrals and (j[3],j[2],j[1],j[0]) not in Dihedrals:
                Dihedrals = Dihedrals + [(j[0],j[1],j[2],j[3])]
                Dihedral_types = Dihedral_types + [ (Atom_types[j[0]],Atom_types[j[1]],Atom_types[j[2]],Atom_types[j[3]]) ]

    # Find 1-5s
    # NOTE: no effort is made to sort based on types because these are only used for coul and lj corrections
    for i in Dihedrals:
        
        # Find atoms attached to first atom of each dihedral
        One_fives += [ (count_j,i[0],i[1],i[2],i[3]) for count_j,j in enumerate(Adj_mat[i[0]]) if j == 1 and count_j not in [i[1],i[2],i[3]] ]

        # Find atoms attached to the fourth atom of each dihedral
        One_fives += [ (i[0],i[1],i[2],i[3],count_j) for count_j,j in enumerate(Adj_mat[i[3]]) if j == 1 and count_j not in [i[0],i[1],i[2]] ]

    One_five_types = [ (Atom_types[i[0]],Atom_types[i[1]],Atom_types[i[2]],Atom_types[i[3]],Atom_types[i[4]]) for i in One_fives ]

    if return_all == 1: return Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types
    else: return Bonds,Angles,Dihedrals,One_fives

# Checks the first output file discovered during the subdirectory walk for the types of calculations performed (dft/mp2)
# returns a string to main that is used during the parsing of various modes.
def check_qc():
    
    # Grab the first output file
    output_file = next( os.path.join(dp, f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,'*.out') )

    # Parse the Orca input file to determine the types of calculations present
    parse_flag = 0
    dft_flag = 0
    mp2_flag = 0
    with open(output_file,'r') as f:
        for lc,lines in enumerate(f):
            if "INPUT FILE" in lines:
                parse_flag = 1
            if parse_flag == 1 and "DFT" in lines:
                dft_flag = 1
            if parse_flag == 1 and "MP2" in lines:
                mp2_flag = 1
            if "END OF INPUT" in lines:
                break

    # Based on what was found fill in the qc_types list
    qc_types = []
    if dft_flag == 1:
        qc_types += ['dft']
    if mp2_flag == 1:
        qc_types += ['mp2']

    # Return the list of calculation types
    return  qc_types

# Reads in the charges from the supplied databases and replaces the corresponding charges in charges_AA and charges_UA
def FF_charges(FF,charges_AA,charges_UA,origin_dict=None,add_more=[]):

    # Read in the charges from the *.db file(s)
    FF_charges = {}
    for i in FF:
        with open(i,'r') as f:
            for lines in f:
                fields = lines.split()
                if len(fields) > 0 and fields[0] == "charge":
                    FF_charges[fields[1]] = float(fields[2])
                    if origin_dict is not None:
                        origin_dict[fields[1]] = i

    # Replace elements in charges_AA and charges_UA
    for i in list(charges_AA.keys()):
        if i in list(FF_charges.keys()):
            charges_AA[i] = FF_charges[i]
    for i in list(charges_UA.keys()):
        if i in list(FF_charges.keys()):
            charges_UA[i] = FF_charges[i]

    # Add more option
    for i in list(FF_charges.keys()):
        if i in add_more:
            charges_AA[i] = FF_charges[i]
        if "-UA" in i and i.replace('-UA',"") in add_more:
            charges_UA[i] = FF_charges[i]

    if origin_dict is not None:
        return charges_AA,charges_UA,origin_dict
    else:
        return charges_AA,charges_UA

# Reads in the masses from the supplied databases and replaces the corresponding masses in masses_AA and masses_UA
def FF_masses(FF,masses_AA,masses_UA,add_more=[]):

    # Read in the charges from the *.db file(s)
    FF_masses = {}
    for i in FF:
        with open(i,'r') as f:
            for lines in f:
                fields = lines.split()
                if len(fields) > 0 and fields[0] == "atom":
                    FF_masses[fields[1]] = float(fields[3])

    # Replace elements in masses_AA and masses_UA
    for i in list(masses_AA.keys()):
        if i in list(FF_masses.keys()):
            masses_AA[i] = FF_masses[i]
    for i in list(masses_UA.keys()):
        if i in list(FF_masses.keys()):
            masses_UA[i] = FF_masses[i]

    # Add more option, add masses from FF_masses that are in the add_more
    # set. the -UA versions of add_more atoms are also added (second if)
    for i in list(FF_masses.keys()):
        if i in add_more:
            masses_AA[i] = FF_masses[i]
        if "-UA" in i and i.replace('-UA',"") in add_more:
            masses_UA[i] = FF_masses[i]

    return masses_AA,masses_UA

# Description: Initialize VDW_dict based on UFF parameters for the initial guess of the fit.
def initialize_VDW(atomtypes_AA,atomtypes_UA,FF_files=[],verbose=True,mixing_rule=None):

    # Initialize UFF parameters (tuple corresponds to eps,sigma pairs for each element)
    # Taken from UFF (Rappe et al. JACS 1992)
    # Note: LJ parameters in the table are specificed in the eps,r_min form rather than eps,sigma
    #       the conversion between r_min and sigma is sigma = r_min/2^(1/6)
    # Note: Units for sigma = angstroms and eps = kcal/mol 
    UFF_dict = { 1:(0.044,2.5711337005530193),  2:(0.056,2.1043027722474816),  3:(0.025,2.183592758161972),  4:(0.085,2.4455169812952313),\
                 5:(0.180,3.6375394661670053),  6:(0.105,3.4308509635584463),  7:(0.069,3.260689308393642),  8:(0.060,3.1181455134911875),\
                 9:(0.050,2.996983287824101),  10:(0.042,2.88918454292912),   11:(0.030,2.657550876212632), 12:(0.111,2.6914050275019648),\
                13:(0.505,4.008153332913386),  14:(0.402,3.82640999441276),   15:(0.305,3.694556984127987), 16:(0.274,3.594776327696269),\
                17:(0.227,3.5163772404999194), 18:(0.185,3.4459962417668324), 19:(0.035,3.396105913550973), 20:(0.238,3.0281647429590133),\
                21:(0.019,2.935511276272418),  22:(0.017,2.828603430095577),  23:(0.016,2.800985569833227), 24:(0.015,2.6931868249382456),\
                25:(0.013,2.6379511044135446), 26:(0.013,2.5942970672246677), 27:(0.014,2.558661118499054), 28:(0.015,2.5248069672097215),\
                29:(0.005,3.113691019900486),  30:(0.124,2.4615531582217574), 31:(0.415,3.904809081609107), 32:(0.379,3.813046513640652),\
                33:(0.309,3.7685015777336357), 34:(0.291,3.746229109780127),  35:(0.251,3.731974730289881), 36:(0.220,3.689211591819145),\
                37:(0.040,3.6651573264293558), 38:(0.235,3.2437622327489755), 39:(0.072,2.980056212179435), 40:(0.069,2.78316759547042),\
                41:(0.059,2.819694442914174),  42:(0.056,2.7190228877643157), 43:(0.048,2.670914356984738), 44:(0.056,2.6397329018498255),\
                45:(0.053,2.6094423454330538), 46:(0.048,2.5827153838888437), 47:(0.036,2.804549164705788), 48:(0.228,2.537279549263686),\
                49:(0.599,3.976080979060334),  50:(0.567,3.9128271700723705), 51:(0.449,3.937772334180300), 52:(0.398,3.982317270087316),\
                53:(0.339,4.009044231631527),  54:(0.332,3.923517954690054),  55:(0.045,4.024189509839913), 56:(0.364,3.2989979532736764),\
                72:(0.072,2.798312873678806),  73:(0.081,2.8241489365048755), 74:(0.067,2.734168165972701), 75:(0.066,2.631714813386562),\
                76:(0.037,2.7796040005978586), 77:(0.073,2.5301523595185635), 78:(0.080,2.453535069758495), 79:(0.039,2.9337294788361374),\
                80:(0.385,2.4098810325696176), 81:(0.680,3.872736727756055),  82:(0.663,3.828191791849038), 83:(0.518,3.893227398273283),\
                84:(0.325,4.195242063722858),  85:(0.284,4.231768911166611),  86:(0.248,4.245132391938716) }

    # Initialize *.db based VDW_dictionary (by default, any LJ parameters found in the supplied database file are used instead of the UFF parameters)
    VDW_FF = {}
    for i in FF_files:
        with open(i,'r') as f:
            for lines in f:
                fields = lines.split()
# OLD                if len(fields) > 0 and fields[0] == "vdw" and fields[3] == "lj":
                if len(fields) > 0 and fields[0] == "vdw" and fields[3] == "lj" and ( fields[1] in atomtypes_AA or fields[1] in atomtypes_UA ) and (fields[2] in atomtypes_AA or fields[2] in atomtypes_UA ) :
                    VDW_FF[(fields[1],fields[2])] = [fields[3],float(fields[4]),float(fields[5])]
                    VDW_FF[(fields[2],fields[1])] = VDW_FF[(fields[1],fields[2])]

    # Apply mixing rules to the VDW params read from file
    if mixing_rule is not None:
        self_types = [ i[0] for i in list(VDW_FF.keys()) if i[0] == i[1] ]
        for i in self_types:
            for j in self_types:
                if (i,j) not in list(VDW_FF.keys()):
                    if mixing_rule == "lb":
                        eps    = (VDW_FF[(i,i)][1]*VDW_FF[(j,j)][1])**(0.5) 
                        sigma  = (VDW_FF[(i,i)][2]+VDW_FF[(j,j)][2])/2.0 
                        VDW_FF[(i,j)] = ["lj",eps,sigma]
                        VDW_FF[(j,i)] = VDW_FF[(i,j)]
                    elif mixing_rule == "wh":
                        sigma  = ((VDW_FF[(i,i)][2]**(6.0)+VDW_FF[(j,j)][2]**(6.0))/2.0)**(1.0/6.0)
                        eps    = (VDW_FF[(i,i)][1]*VDW_FF[(i,i)][2]**(6.0) * VDW_FF[(j,j)][1]*VDW_FF[(j,j)][2]**(6.0) )**(0.5) / sigma**(6.0)
                        VDW_FF[(i,j)] = ["lj",eps,sigma]
                        VDW_FF[(j,i)] = VDW_FF[(i,j)]                        

    # Initialize AA-parameters for the VDW_dict (master dictionary used for fits) based VDW_FF or UFF parameters with Lorentz-Berthelot mixing rules
    VDW_dict = {}
    for count_i,i in enumerate(atomtypes_AA):
        for count_j,j in enumerate(atomtypes_AA):
            if count_i < count_j:
                continue

            type_1 = int(i.split('[')[1].split(']')[0])
            type_2 = int(j.split('[')[1].split(']')[0])
            eps    = (UFF_dict[type_1][0]*UFF_dict[type_2][0])**(0.5) 
            sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0 
            if (i,j) in VDW_FF: VDW_dict[(i,j)] = ["lj",VDW_FF[(i,j)][1],VDW_FF[(i,j)][2]]
            else: VDW_dict[(i,j)] = ['lj',eps,sigma]
            if (j,i) in VDW_FF: VDW_dict[(j,i)] = ["lj",VDW_FF[(j,i)][1],VDW_FF[(j,i)][2]]
            else: VDW_dict[(j,i)] = ['lj',eps,sigma]

    # Initialize UA-parameters for the VDW_dict (master dictionary used for fits) based VDW_FF or UFF parameters with Lorentz-Berthelot mixing rules
    for count_i,i in enumerate(atomtypes_UA):
        for count_j,j in enumerate(atomtypes_UA):
            if count_i < count_j:
                continue

            type_1 = int(i.split('[')[1].split(']')[0])
            type_2 = int(j.split('[')[1].split(']')[0])
            eps    = (UFF_dict[type_1][0]*UFF_dict[type_2][0])**(0.5)
            sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0
            if (i,j) in VDW_FF: VDW_dict[(i,j)] = ["lj",VDW_FF[(i,j)][1],VDW_FF[(i,j)][2]]
            else: VDW_dict[(i,j)] = ['lj',eps,sigma]
            if (j,i) in VDW_FF: VDW_dict[(j,i)] = ["lj",VDW_FF[(j,i)][1],VDW_FF[(j,i)][2]]
            else: VDW_dict[(j,i)] = ['lj',eps,sigma]

    # Print summary
    if verbose is True:
        print("\n{}".format("*"*167))
        print("* {:^163s} *".format("Initializing/Reading VDW parameters for the fit (those with * were read from the FF file(s))"))
        print("*{}*".format("-"*165))
        print("* {:<50s} {:<50s} {:<20s}  {:<18s} {:<18s}   *".format("Type","Type","VDW_type","eps (kcal/mol)","sigma (angstroms)"))
        print("{}".format("*"*167))

        # Print AA parameters
        printed_list = []
        for j in list(VDW_dict.keys()):            
            if j[0] == j[1]:
                if j in printed_list or "-UA" in j[0]: continue
                if j in VDW_FF:
                    print("  {:<50s} {:<50s} {:<20s} {:< 18.4f} {:< 18.4f} *".format(j[0],j[1],VDW_dict[j][0],VDW_dict[j][1],VDW_dict[j][2]))
                else:
                    print("  {:<50s} {:<50s} {:<20s} {:< 18.4f} {:< 18.4f}".format(j[0],j[1],VDW_dict[j][0],VDW_dict[j][1],VDW_dict[j][2]))
                printed_list += [j,(j[1],j[0])]

        # Print UA parameters
        for j in list(VDW_dict.keys()):
            if j[0] == j[1]:
                if j in printed_list: continue
                if j in VDW_FF:
                    print("  {:<50s} {:<50s} {:<20s} {:< 18.4f} {:< 18.4f} *".format(j[0],j[1],VDW_dict[j][0],VDW_dict[j][1],VDW_dict[j][2]))
                else:
                    print("  {:<50s} {:<50s} {:<20s} {:< 18.4f} {:< 18.4f}".format(j[0],j[1],VDW_dict[j][0],VDW_dict[j][1],VDW_dict[j][2]))
                printed_list += [j,(j[1],j[0])]
        print("")        

    return VDW_dict

# Description: Initialize a dictionary containing the UFF bond force-constants for the supplied taffi bond types.
# 
# Inputs       bond_types: A list of tuples containing taffi bond types
#
# Returns      bond_dict:  A dictionary with UFF bond force-constant keyed to the taffi bond types
def UFF_bonds(bond_types):

    if getattr(UFF_bonds, "UFF_Z_dict", None) is None:

        # Initialize UFF parameters (tuple corresponds to eps,sigma pairs for each element)
        # Taken from UFF (Rappe et al. JACS 1992)
        # Note: LJ parameters in the table are specificed in the eps,r_min form rather than eps,sigma
        #       the conversion between r_min and sigma is sigma = r_min/2^(1/6)
        # Note: Units for sigma = angstroms and eps = kcal/mol 
        UFF_bonds.UFF_Z_dict = { 1:0.712,  2:0.098,  3:1.026,  4:1.565,\
                                 5:1.755,  6:1.912,  7:2.544,  8:2.300,\
                                 9:1.735,  10:0.194, 11:1.081, 12:1.787,\
                                13:1.792,  14:2.323, 15:2.863, 16:2.703,\
                                17:2.348,  18:0.300, 19:1.165, 20:2.141,\
                                21:2.592,  22:2.659, 23:2.679, 24:2.463,\
                                25:2.430,  26:2.430, 27:2.430, 28:2.430,\
                                29:1.756,  30:1.308, 31:1.821, 32:2.789,\
                                33:2.864,  34:2.764, 35:2.519, 36:0.452,\
                                37:1.592,  38:2.449, 39:3.257, 40:3.667,\
                                41:3.618,  42:3.400, 43:3.400, 44:3.400,\
                                45:3.608,  46:3.210, 47:1.956, 48:1.650,\
                                49:2.070,  50:2.961, 51:2.704, 52:2.882,\
                                53:2.650,  54:0.556, 55:1.573, 56:2.727,\
                                72:3.921,  73:4.075, 74:3.700, 75:3.700,\
                                76:3.700,  77:3.731, 78:3.382, 79:2.625,\
                                80:1.750,  81:2.068, 82:2.846, 83:2.470,\
                                84:2.330,  85:2.240, 86:0.583 }

    if getattr(UFF_bonds, "UFF_radii", None) is None:

        # Initialize UFF bond radii (Rappe et al. JACS 1992)
        # NOTE: Units of angstroms 
        # NOTE: These radii neglect the bond-order and electronegativity corrections in the original paper. Where several values exist for the same atom, the largest was used. 
        UFF_bonds.UFF_radii = {  1:0.354,  2:0.849,\
                                 3:1.336,  4:1.074,                                                                                                      5:0.838,  6:0.757,  7:0.700,  8:0.658,  9:0.668, 10:0.920,\
                                11:1.539, 12:1.421,                                                                                                     13:1.244, 14:1.117, 15:1.117, 16:1.064, 17:1.044, 18:1.032,\
                                19:1.953, 20:1.761, 21:1.513, 22:1.412, 23:1.402, 24:1.345, 25:1.382, 26:1.335, 27:1.241, 28:1.164, 29:1.302, 30:1.193, 31:1.260, 32:1.197, 33:1.211, 34:1.190, 35:1.192, 36:1.147,\
                                37:2.260, 38:2.052, 39:1.698, 40:1.564, 41:1.473, 42:1.484, 43:1.322, 44:1.478, 45:1.332, 46:1.338, 47:1.386, 48:1.403, 49:1.459, 50:1.398, 51:1.407, 52:1.386, 53:1.382, 54:1.267,\
                                55:2.570, 56:2.277, 57:1.943, 72:1.611, 73:1.511, 74:1.526, 75:1.372, 76:1.372, 77:1.371, 78:1.364, 79:1.262, 80:1.340, 81:1.518, 82:1.459, 83:1.512, 84:1.500, 85:1.545, 86:1.420}

    # Initialize *.db based VDW_dictionary (by default, any LJ parameters found in the supplied database file are used instead of the UFF parameters)
    Bond_dict = {}
    for i in bond_types:
        type_1 = int(i[0].split('[')[1].split(']')[0])
        type_2 = int(i[1].split('[')[1].split(']')[0])
        Bond_dict[i] = 664.12 * UFF_bonds.UFF_Z_dict[type_1] * UFF_bonds.UFF_Z_dict[type_2] / ( UFF_bonds.UFF_radii[type_1] + UFF_bonds.UFF_radii[type_2] )**(3.0)
        # print "bond {}: {}".format(i,Bond_dict[i])

    return Bond_dict

def natural_sort(l): 
    convert = lambda text: int(text) if text.isdigit() else text.lower() 
    alphanum_key = lambda key: [ convert(c) for c in re.split('([0-9]+)', key) ] 
    return sorted(l, key = alphanum_key)

# Description: A wrapper function for the commands to parse the masses and charges
#              from the optimized geometry located in the first bond being parametrized.
#
# Inputs       elements:   A list of elements (from the *master.xyz file parsed prior)
#              bond_steps: The number of bond configurations used for the parametrization
#                          This is used here to id the optimized (unperturbed) geometry.
#
# Returns      charges:    A list of the CHELPG partial charges for each atom type (indexed to elements).
#              masses:     A list of the masses of each atomtypes (indexed to elements)
def parse_spectator_charges(geo_folders,UA_opt=False):

    #########################################################
    #  Process Charges from the Geometry Optimization Runs  #
    #########################################################

    # Check whether the information is present to perform a TAFFI fit
    # if not, print an error to the user
    taffi_flag = 1
    for count_i,i in enumerate(geo_folders):
        charge_vpot = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"*.vpot") ]
        charge_xyz = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"geoopt.xyz") ]
        if len(charge_vpot) == 0 or len(charge_xyz) == 0:
            taffy_flag = 0
            break

    # Print diagnostic
    if taffi_flag == 1:
        print("Fitting the partial charges for the spectator atoms using the TAFFI two-step procedure...")
    else:
        print("ERROR in parse_spectator_charges: The requisite information for performing a taffi fit was not discovered. Check job completion...")
        exit

    # Collect taffi charges for spectator atoms.
    charge_list = []
    atomtypes = []
    taffy_flag = 0
    charges = {}
    for count_i,i in enumerate(geo_folders):

        # Find the output file of the globally optimized geometry (i.e., the iteration 
        # corresponding to args.bond_steps is the unperturbed globally optimized geometry)
        # and check for *.vpot and *.xyz files
        output = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"*.out") ]
        if len(output) == 0:
            print("ERROR in parse_spectator_charges: No output for the globally optimized geometry was found. Check that the run completed.")
            quit()
        charge_vpot = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"*.vpot") ]
        charge_xyz = [ i+'/'+names for names in os.listdir(i) if fnmatch.fnmatch(names,"geoopt.xyz") ]

        # Parse the charges using the TAFFI procedure

        # parse the total charge from the output file
        qtot = None
        with open(output[0],'r') as f:
            for lines in f:
                fields = lines.split()
                if len(fields) == 3 and fields[0] == "Total" and fields[1] == "charge:":
                    qtot = round(float(fields[2]))

        # safety in case something went wrong during the parse
        if qtot is None:
            print("ERROR in parse_spectator_charges: could not parse the total charge on the molecule from {}. Exiting...".format(output[0]))
            quit()

        # parse the partial charges using the taffi two-step procedure
        charge_dict,errs = fit_charges(charge_vpot[0],xyz_file=charge_xyz[0],out_file=output[0],qtot=qtot,gens=2,w_pot=1.0,w_qtot=1.0,w_hyper=0.0,w_dipole=0.1,UA_opt=UA_opt,symmetrize=False,two_step=True)

        # get the atomtypes
        e,g = xyz_parse(charge_xyz[0]) 
        a   = Table_generator(e,g,File=charge_xyz[0])
        atom_types = id_types(e,a,gens=2,geo=g)      
        if UA_opt is True:
            charges[i] = np.array([ charge_dict[j+'-UA'][0] if return_UA_H(j) == 0 else 0.0 for j in atom_types ])
        else:
            charges[i] = np.array([ charge_dict[j][0] for j in atom_types ])
        
    return charges

# Fit is performed OPLS style, but with each fit dihedral constrained to the values observed in the QC scan.
# At each configuration all other modes are allowed to relax, and the dihedral potentials are fit to the residual
# of the QC_potential minus the FF_potential without the fit modes. 
def parse_dihedrals_qcmatch(Folder,geo_folders,dihedral_folders_frag,list_of_dihedral_types,dihedral_folders,dihedral_atomtypes,charge_dict,eq_charges,qc,modes_from_FF,lammps_exe,UA_opt=0,weight=1.0,delta_xhi2_thresh=1.0E-5,\
                            min_cycles=10,max_cycles=1000,one_four_scale_coul=0.0,one_four_scale_vdw=0.0,link_mode="sub",save_intermediates=1,ramp_coul=0,charge_origin=None,gens=2):
                            
    # In python all global variables must be declared in each scope
    global FF_dict,Dihedral_Data,dihedral_scan,V_range,dihedral_types,dihedral_fit_type,QC_dihedral
    kb_kcalmol = 0.0019872041
    conv_thresh = 1.0  # This controls the convergence threshold for breaking out of the self-consistency fit loop. This is a constraint on the percentage change in the dihedral barrier height for each dihedral being fit. 
    inner_cycles = 1   # This controls the number of self-consistency steps when fitting each individual scan. Presently it is unused (i.e. one fit is used). 
    QC_dihedral = qc

    # opt_steps controls the degree increments for the opls dihedral scan
    opt_steps = 360
    if opt_steps > 360:
        print("ERROR: opt_steps must be <= 360, otherwise it breaks the parsing algorithm which uses a round call to determine which configurations to print")
        quit()

    # Initialize dictionary for mapping atomic number to element
    atom_to_element = { 1:'H' ,  2:'He',\
                        3:'Li',  4:'Be',  5:'B' ,  6:'C' ,  7:'N' ,  8:'O' ,  9:'F' , 10:'Ne',\
                       11:'Na', 12:'Mg', 13:'Al', 14:'Si', 15:'P' , 16:'S' , 17:'Cl', 18:'Ar',\
                       19:'K' , 20:'Ca', 21:'Sc', 22:'Ti', 23:'V' , 24:'Cr', 25:'Mn', 26:'Fe', 27:'Co', 28:'Ni', 29:'Cu', 30:'Zn', 31:'Ga', 32:'Ge', 33:'As', 34:'Se', 35:'Br', 36:'Kr',\
                       37:'Rb', 38:'Sr', 39:'Y' , 40:'Zr', 41:'Nb', 42:'Mo', 43:'Tc', 44:'Ru', 45:'Rh', 46:'Pd', 47:'Ag', 48:'Cd', 49:'In', 50:'Sn', 51:'Sb', 52:'Te', 53:'I' , 54:'Xe',\
                       55:'Cs', 56:'Ba', 57:'La', 72:'Hf', 73:'Ta', 74:'W' , 75:'Re', 76:'Os', 77:'Ir', 78:'Pt', 79:'Au', 80:'Hg', 81:'Tl', 82:'Pb', 83:'Bi', 84:'Po', 85:'At', 86:'Rn'}    

    # Save working directory
    current_dir = os.getcwd()
    
    # Initialize path list for IBI/QC comparison figures
    paths_to_dihedral_fitpot_plots = []
    paths_to_dihedral_fitpot_dict  = {}   
    paths_to_dihedral_totpot_plots = []
    paths_to_dihedral_totpot_dict  = {}   

    # Grab the dihedral data (returns a dictionary holding dihedral_type,angle tuples for each configuration and the DFT and MP2 energies
    if UA_opt == 0:
        Dihedral_Data = get_dihedral_data(list_of_dihedral_types,dihedral_folders,"AA",dihedral_atomtypes,charge_dict,[qc.lower()],corr_list=[],one_four_scale_coul=one_four_scale_coul,one_four_scale_vdw=one_four_scale_vdw)
    else:
        Dihedral_Data = get_dihedral_data(list_of_dihedral_types,dihedral_folders,"UA",dihedral_atomtypes,charge_dict,[qc.lower()],corr_list=[],one_four_scale_coul=one_four_scale_coul,one_four_scale_vdw=one_four_scale_vdw)

    # Return to working directory and save the list of catenated geometries
    os.chdir(current_dir)
    for i in list(Dihedral_Data.keys()):

        # Check if the MD subfolder needs to be created
        if UA_opt == 1: savename = Folder+"/dihedrals/"+"_".join([ k+"-UA" for k in i ])+"_geos.xyz"
        else: savename = Folder+"/dihedrals/"+"_".join([ k for k in i ])+"_geos.xyz"
        f = open(savename, 'w')
    
        # Loop over all discovered geometries and save their geometries to Geo_dict and
        # catenate to a single file for viewing.
        for j in natural_sort(list(Dihedral_Data[i].keys())):
            f.write("{:d}\n\n".format(len(Dihedral_Data[i][j]["Geo"])))
            for count_k,k in enumerate(Dihedral_Data[i][j]["Geo"]):
                f.write(" {:<10s} {:< 15.8f} {:< 15.8f} {:< 15.8f}\n".format(Dihedral_Data[i][j]["Elements"][count_k],k[0],k[1],k[2]))
        f.close()

    # Initialize/Zero out the initial guess for the dihedrals being fit        
    # NOTE: get_dihedral_data() populates the UA-types with the appropriate labels so there is no need to append the -UA string(s)
    # NOTE: only dihedrals that weren't read from the user-supplied FF file are fit (the second line accomplishes this).
    new_dihedrals = list(set([ k[0] for j in list(Dihedral_Data.keys()) for k in Dihedral_Data[j]["0"]["Dihedrals"] ]))
    new_dihedrals = [ i for i in new_dihedrals if i+tuple(["opls"]) not in modes_from_FF ]    
    for j in new_dihedrals:
        if j not in list(FF_dict["dihedrals"].keys()):
            FF_dict["dihedrals"][j] = {}
        FF_dict["dihedrals"][j][qc] = ["opls",0.0,0.0,0.0,0.0]
    
    # return if no dihedrals are being fit
    if new_dihedrals == []:
        max_cycles = 0
        return [],[]
    
    # Print banner
    if UA_opt == 0:
        print("\n{}\n* {:^163s} *\n{}".format("*"*167,"Beginning Self-Consistent OPLS-style Fit of All AA-{} Dihedral Scans".format(qc),"*"*167))
    else:
        print("\n{}\n* {:^163s} *\n{}".format("*"*167,"Beginning Self-Consistent OPLS-style Fit of All UA-{} Dihedral Scans".format(qc),"*"*167))

    # Initialize break flags for each scan and a list of the scans being parsed (i.e. the keys from the Dihedral_Data dictionary)
    break_flags = { i:0 for i in list(Dihedral_Data.keys()) }
    keys = list(Dihedral_Data.keys())
    other_inchi = []
    for i in eq_charges:
      i = i.split('/')
      if(len(i) !=2 ) and i[-4] not in other_inchi:
         other_inchi.append(i[-4])
    #print(other_inchi)
    Direct_keys = []
    for count_i,i in enumerate(keys):
      if(dihedral_folders_frag[count_i] in other_inchi):
         Direct_keys += ['../../{}/Intra/dihedrals/geoopt'.format(dihedral_folders_frag[count_i])]
      else:
         Direct_keys += ['dihedrals/geoopt']
         

    # Initialize old_fit_pot dictionary which holds the most recent composite fit potential for each scan (used for determining when self-consistency has been reached)
    # and the percent_updates dictionary which holds the fractional change in composite dihedral range for each scan
    old_fit_pot = { i:[0.0] for i in list(Dihedral_Data.keys()) }
    percent_updates = { i:0.0 for i in list(Dihedral_Data.keys()) }
    fit_err_dict = { i:0.0 for i in list(Dihedral_Data.keys()) }

    # List used for temporary dihedrals that need to be removed before exiting the function
    remove_at_end = []

    ########################
    # Begin FIT Outer Loop #
    ########################

    # Self-consistency is acheived by looping over all dihedral fits and refitting until mutual convergence is acheived.
    for c1 in range(max_cycles):

        # Check for break condition and print diagnostic
        if False not in [ break_flags[i]==1 for i in list(break_flags.keys()) ]:
            print("\n{:^167s}\n{:^167s}\n{:^167}\n".format("*"*87,"* {:^83s} *".format("Fit complete (all dihedrals converged within {:1.2f}%)".format(conv_thresh)),"*"*87))
            break

        # Reinitialize break conditions and print diagnostic    
        else:
            print("\n{:^167s}\n{:^167s}\n{:^167}\n".format("*"*87,"* {:^83s} *".format("Self-consistency cycle {}".format(c1)),"*"*87))
            break_flags = { i:0 for i in list(Dihedral_Data.keys()) }

        # Shuffle the keys to randomize the order in which the dihedrals are fit and initialize a list of fit dihedrals to avoid the same
        # dihedral being fit in multiple scans
        # random.shuffle(keys)
        already_fit_list = []
        CURRENT_FF_dict = deepcopy(FF_dict)

        ########################
        # Begin FIT Inner Loop #
        ########################

        # Loop over the dihedral scans and perform fits to all coincident dihedrals in each scan
        for count_i,i in enumerate(keys):
            
            #############################################################################################################################
            # Generate prerequisite data for the scan (NOTE: the program assumes atom ordering is retained across the scan. This is     #
            # guarranteed by the generation scripts, but can lead it errors if other programs don't also abide by the same convention.  #
            #############################################################################################################################

            # Skip if UA_opt == 1 and there are any UA-hydrogens are in this dihedral
            if UA_opt == 1 and 1 in [ return_UA_H(k) for k in i ]: 
                break_flags[i] = 1
                continue

            # Check if all runs for this scan succesfully completed
            incomplete = []
            for k in list(Dihedral_Data[i].keys()):
                if Dihedral_Data[i][k]["Completion"] == 0: incomplete += [ int(i) ]
            if len(incomplete) > 0:
                break_flags[i] = 1
                continue

            # Print diagnostic
            print("\n\tPerforming OPLS fit on scan {}".format(i))
                
            # Initialize the save_folder name and check if the subfolder needs to be created
            if UA_opt == 1: save_folder = Folder+'/dihedrals/'+"_".join([ k+"-UA" for k in i ])+"-MD"
            else: save_folder = Folder+'/dihedrals/'+"_".join(i)+"-MD"
            if os.path.isdir(save_folder) == False:
                os.makedirs(save_folder)

            # Change directory into the working subfolder
            os.chdir(save_folder)

            # Update the charges from the force-field if all of them are available.
            Charges = deepcopy(eq_charges[Direct_keys[count_i]])
            if UA_opt == 0:
                tmp_types = [ next( m for m in k.split('link-') if m != '') for k in Dihedral_Data[i]["0"]["Atomtypes"] ]
            elif UA_opt == 1:
                tmp_types = [ next( m for m in k.split('link-') if m != '')+'-UA' for k in Dihedral_Data[i]["0"]["Atomtypes"] ]
            if charge_origin is not None:
                
                # Update the charges with the FF charges is they are all available
                if True not in [ charge_origin[k] == "eq_config" for k in tmp_types ]:
                    #print("Charges")
                    #print(len(Charges))
                    for k in range(len(Charges)): Charges[k] = FF_dict["charges"][tmp_types[k]]

                # Else, proceed with eq_charges
                else:
                    print("\tWARNING: Missing partial charges for some or all of the atomtypes in this fragment, defaulting to equilibrium partial charges.") 

            # If the provenance of the charges is unknown, then proceed by using the FF charges if they are all available
            elif False not in [ tmp_types[k] in list(FF_dict["charges"].keys()) for k in range(len(Charges)) ]:
                for k in range(len(Charges)): Charges[k] = FF_dict["charges"][tmp_types[k]]

            # Generate the charges for the scan (taffi charges for the scanned dihedrals and eq_charges for the rest)
            scanned_atoms = set([ m for k in Dihedral_Data[i]["0"]["Dihedral_Atoms"] for m in k[0] ])

            # Recalculate the atom types if this is a ring dihedral being fit to a linear analog
            # These are used later for reassigning the bond and angle types
            if "R" in i[1] and "R" in i[2]:
                atomtypes_lin = id_types([ atom_to_element[int(k.split('[')[1].split(']')[0])] for k in Dihedral_Data[i]["0"]["Atomtypes"] ],Dihedral_Data[i]["0"]["Adj_mat"],gens=gens)

            # Flag link atoms for removal.(only matters for fragment based parametrizations)
            # NOTE: the most up to date convention is to use link-* for spectator atoms, but [0] and [1] are included for legacy reasons.
            # NOTE: when "del" is specified the link atoms are deleted, otherwise the link atoms are substituted for their appropriate atomtype. 
            if link_mode == 'del':
                keep_ind = [ count_k for count_k,k in enumerate(Dihedral_Data[i]["0"]["Atomtypes"]) if ( k != "[0]" and k != "[1]" and "link-" not in k ) ]
            elif link_mode == 'sub':
                keep_ind = [ count_k for count_k,k in enumerate(Dihedral_Data[i]["0"]["Atomtypes"]) if ( k != "[0]" and k != "[1]" ) ]

            # If UA_opt == 1 then remove UA-hydrogens. return_UA_H determines if an atomtype is a UA-H based on its label.
            if UA_opt == 1:
                keep_ind = [ k for k in keep_ind if return_UA_H(Dihedral_Data[i]["0"]["Atomtypes"][k]) == 0 ]

            # Initialize requisite lists/arrays, minus any spectator hydrogens
            if UA_opt == 0: Atomtypes = [ Dihedral_Data[i]["0"]["Atomtypes"][k] for k in keep_ind ]
            if UA_opt == 1: Atomtypes = [ Dihedral_Data[i]["0"]["Atomtypes"][k]+"-UA" for k in keep_ind ]            
            Atomtypes = [ next( m for m in k.split('link-') if m != '') for k in Atomtypes ]
            Adj_mat   = Dihedral_Data[i]["0"]["Adj_mat"][keep_ind,:]
            Adj_mat   = Adj_mat[:,keep_ind]
            Elements  = [ atom_to_element[int(k.split('[')[1].split(']')[0])] for k in Atomtypes ]
            Charges = [ Charges[k] for k in keep_ind ]
            Masses    = [ FF_dict["masses"][next( m for m in k.split('link-') if m != '' )] for k in Atomtypes ]
            Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(Adj_mat,Atomtypes,return_all=1)

            # If this is a flexible ring dihedral then the bonds and angles needs to be reassigned as the corresponding linear types
            # This is because the flexible ring dihedrals are fit with linear analaogs of the cycles they belong to. 
            if "R" in i[1] and "R" in i[2]:

                # Trim atomtypes_lin to conform with keep_ind 
                atomtypes_lin = [ atomtypes_lin[k] for k in keep_ind ]
                
                # Reassign bonds using atomtypes_lin atomtypes
                for count_j,j in enumerate(Bonds):
                    Bond_types[count_j],Bonds[count_j] = canon_bond(tuple([ atomtypes_lin[k] for k in j ]),ind=Bonds[count_j])

                # Reassign angles using atomtypes_lin atomtypes
                for count_j,j in enumerate(Angles):
                    Angle_types[count_j],Angles[count_j] = canon_angle(tuple([ atomtypes_lin[k] for k in j ]),ind=Angles[count_j])

            # Determine bonding matrix for the compound
            q_tmp = int(round(sum(eq_charges[Direct_keys[count_i]])))
            bond_mat = find_lewis([ next( m for m in k.split('link-') if m != '') for k in Dihedral_Data[i]["0"]["Atomtypes"] ],\
                                  Dihedral_Data[i]["0"]["Adj_mat"], q_tot=q_tmp, b_mat_only=True,verbose=False)
            bond_mat = [ k[keep_ind,:][:,keep_ind] for k in bond_mat ]### IMPT
 
            # Add the opls/harmonic type to the Dihedral_types tuples
            # NOTE: due to the hydrogenation procedure of fragments, sometimes non-scanned dihedrals end up needing to be
            #       modeled as harmonic or flexible. The following if/else statements check the need to swapping the type 
            #       of each dihedral as needed. 
            for j in range(len(Dihedrals)):
                if 2 in [ k[Dihedrals[j][1],Dihedrals[j][2]] for k in bond_mat ]:
                    if Dihedral_types[j] in list(FF_dict["dihedrals_harmonic"].keys()):                        
                        dihedral_type = "harmonic"
                    elif Dihedral_types[j] not in list(FF_dict["dihedrals"].keys()):
                        print("ERROR in parse_dihedrals_qcmatch: Dihedral {} ({}) is being processed as a harmonic type, but it is not in either the harmonic or flexible dihedral dictionaries.".format(Dihedral_types[j],Dihedrals[j]))
                        quit()
                    else:
                        print("\t\tWARNING: Dihedral {} ({}) is being processed as a harmonic type, but it is being modeled as a flexible dihedral.".format(Dihedral_types[j],Dihedrals[j]))
                        # print "this geometry is geometry: {}".format(i)
                        dihedral_type = "opls"
                else:
                    if Dihedral_types[j] in list(FF_dict["dihedrals"].keys()):                        
                        dihedral_type = "opls"
                    elif Dihedral_types[j] not in list(FF_dict["dihedrals_harmonic"].keys()):
                        print("ERROR in parse_dihedrals_qcmatch: Dihedral {} ({}) is being processed as a flexible type, but it is not in either the flexible or harmonic dihedral dictionaries.".format(Dihedral_types[j],Dihedrals[j]))
                        quit()
                    else:
                        print("\t\tWARNING: Dihedral {} ({}) is being processed as a flexible type, but it is being modeled as a harmonic dihedral.".format(Dihedral_types[j],Dihedrals[j]))
                        # print "this geometry is geometry: {}".format(i)
                        dihedral_type = "harmonic"

                # Add dihedral type to the types information
                Dihedral_types[j] = tuple(list(Dihedral_types[j])+[dihedral_type])

            # Check that all modes neccessary for the MD run are in the FF_dict
            for j in set(Bond_types):
                if j not in list(FF_dict["bonds"].keys()) or qc not in list(FF_dict["bonds"][j].keys()): print("\nERROR: Bond {} is missing from the force-field. Exiting...".format(j)); quit()
            for j in set(Angle_types):
                if j not in list(FF_dict["angles"].keys()) or qc not in list(FF_dict["angles"][j].keys()): print("\nERROR: Angle {} is missing from the force-field. Exiting...".format(j)); quit()
            for j in set(Dihedral_types):

                # Handles the non-UA case
                if "UA" not in j[0]:

                    # opls-like check
                    if j[4] == "opls":
                        if j[:4] not in list(FF_dict["dihedrals"].keys()) or qc not in list(FF_dict["dihedrals"][j[:4]].keys()):
                            print("\nERROR: Dihedral {} is missing from the force-field. Exiting...".format(j))
                            quit()
                    # harmonic-like check
                    elif j[4] == "harmonic":
                        if ( j[:4] not in list(FF_dict["dihedrals_harmonic"].keys()) or qc not in list(FF_dict["dihedrals_harmonic"][j[:4]].keys()) ):
                            print("\nERROR: Dihedral {} is missing from the force-field. Exiting...".format(j))
                            quit()
                                                                                
                # Handles the UA case
                elif "UA" in j[0]:

                    # opls-like check
                    if j[4] == "opls":
                        if ( j[:4] not in list(FF_dict["dihedrals"].keys()) or qc not in list(FF_dict["dihedrals"][j[:4]].keys()) ):
                             print("\nWARNING: Dihedral {} is missing from the force-field. This is occasionally possible for -UA types and shouldn't affect the fit.".format(j))
                             FF_dict["dihedrals"][j[:4]] = {}
                             FF_dict["dihedrals"][j[:4]]["DFT"] = ["opls", 0.0, 0.0, 0.0, 0.0]
                             FF_dict["dihedrals"][j[:4]]["MP2"] = ["opls", 0.0, 0.0, 0.0, 0.0]
                             CURRENT_FF_dict = deepcopy(FF_dict)
                             remove_at_end += [j]

                    # harmonic-like check
                    elif j[4] == "harmonic":
                        if ( j[:4] not in list(FF_dict["dihedrals_harmonic"].keys()) or qc not in list(FF_dict["dihedrals_harmonic"][j[:4]].keys()) ):
                             print("\nWARNING: Dihedral {} is missing from the force-field. This is occasionally possible for -UA types and shouldn't affect the fit.".format(j))
                             FF_dict["dihedrals_harmonic"][j[:4]] = {}
                             FF_dict["dihedrals_harmonic"][j[:4]]["DFT"] = ["harmonic", 0.0, -1, 2]
                             FF_dict["dihedrals_harmonic"][j[:4]]["MP2"] = ["harmonic", 0.0, -1, 2]
                             CURRENT_FF_dict = deepcopy(FF_dict)
                             remove_at_end += [j]
                                             
            # Update the constrained atoms indices 
            # NOTE: since some atoms may have been removed, the "Dihedral_Atoms" entry is no longer accurate. The
            # use of the tmp array here is a simple way of getting the appropriate indexing.
            fit_dihedrals = []
            fit_dihedral_types = []
            con_map = {}
            for j in (Dihedral_Data[i]["0"]["Dihedral_Atoms"]+Dihedral_Data[i]["0"]["Unscanned_Dihedral_Atoms"]):
                tmp = [0] * len(Dihedral_Data[i]["0"]["Atomtypes"])                                 # All zeros initialization of tmp of equal length to the original atomtypes list
                for count_k,k in enumerate(j[0]): tmp[k] = count_k+1                                # label the 1,2,3, and 4 atoms of the dihedral
                tmp = [ tmp[k] for k in keep_ind ]                                                  # only keep the non-spectator atoms (as in the lists above)
                
                # Be cautious to only add "fit" dihedrals to fit_dihedrals
                if j in Dihedral_Data[i]["0"]["Dihedral_Atoms"]:                    
                    fit_dihedrals += [tuple([ tmp.index(k) for k in range(1,5) ])]                  # identify the indices based on where the 1,2,3, and 4 labels are in tmp, add the current dihedral to the list
                    fit_dihedral_types += [tuple([ Atomtypes[tmp.index(k)] for k in range(1,5) ])]  # identify the types based on where the 1,2,3, and 4 labels are in tmp, add the current dihedral to the list
                    con_map[j[0]] = fit_dihedrals[-1]                                               # add the mapping of the current dihedral to the constrained dictionary.

                # Add the rest of the dihedrals provided they haven't been trimmed during
                # the removal of link atoms
                elif all([ k in tmp for k in range(1,5) ]):
                    con_map[j[0]] = tuple([ tmp.index(k) for k in range(1,5) ])                     # add the mapping of the current Link-atom dihedral to the constrained dictionary.                    

            # Remove any dihedral types that were fitted in another scan
            del_ind = [ count_j for count_j,j in enumerate(fit_dihedral_types) if j in already_fit_list or j+tuple(["opls"]) in modes_from_FF ]
            fit_dihedrals = [ j for count_j,j in enumerate(fit_dihedrals) if count_j not in del_ind ]
            fit_dihedral_types = [ j for count_j,j in enumerate(fit_dihedral_types) if count_j not in del_ind ]

            # Collect the angle constraints for each dihedral instance being fit and the QC energies for each configuration
            fit_constraints = {}
            E_QC = np.zeros(len(list(Dihedral_Data[i].keys())))
            for count_j,j in enumerate(natural_sort(list(Dihedral_Data[i].keys()))):
                E_QC[count_j] = Dihedral_Data[i][j][qc]
                fit_constraints[count_j] = []

                # Add scanned atom constraints
                for count_k,k in enumerate(Dihedral_Data[i][j]["Dihedral_Atoms"]):
                    fit_constraints[count_j] += [(con_map[k[0]],k[1]*180.0/np.pi)]

            # Adjust the angles by +/- 360 if necessary to perform the shortest interpolation between the two configurations
            for j in sorted(fit_constraints.keys())[1:]:
                for count_k,k in enumerate(fit_constraints[j]):
                    angle = k[1]
                    diff = np.abs(fit_constraints[j-1][count_k][1] - angle)
                    if np.abs(fit_constraints[j-1][count_k][1] - (angle+360.0)) < diff:
                        diff = np.abs(fit_constraints[j-1][count_k][1] - (angle+360.0))
                        fit_constraints[j][count_k] = (fit_constraints[j][count_k][0],angle+360.0)
                    if np.abs(fit_constraints[j-1][count_k][1] - (angle-360.0)) < diff:
                        fit_constraints[j][count_k] = (fit_constraints[j][count_k][0],angle-360.0)

            # Collect the angle constraints for secondary dihedrals that aren't being fit
            other_constraints = {}
            for count_j,j in enumerate(natural_sort(list(Dihedral_Data[i].keys()))):
                other_constraints[count_j] = []

                # Add link atom contraints
                if link_mode == 'sub' or link_mode == 'del':
                    for count_k,k in enumerate(Dihedral_Data[i][j]["Unscanned_Dihedral_Atoms"]):
                        if k[0] in con_map:
                            other_constraints[count_j] += [(con_map[k[0]],k[1]*180.0/np.pi)]

            # Adjust the angles by +/- 360 if necessary to perform the shortest interpolation between the two configurations
            for j in sorted(other_constraints.keys())[1:]:
                for count_k,k in enumerate(other_constraints[j]):
                    angle = k[1]
                    diff = np.abs(other_constraints[j-1][count_k][1] - angle)
                    if np.abs(other_constraints[j-1][count_k][1] - (angle+360.0)) < diff:
                        diff = np.abs(other_constraints[j-1][count_k][1] - (angle+360.0))
                        other_constraints[j][count_k] = (other_constraints[j][count_k][0],angle+360.0)
                    if np.abs(other_constraints[j-1][count_k][1] - (angle-360.0)) < diff:
                        other_constraints[j][count_k] = (other_constraints[j][count_k][0],angle-360.0)

            # Save the unique fit types and update already_fit_list (used to avoid fitting the same dihedral as part of multiple scans)
            unique_fit_dihedral_types = sorted(set(fit_dihedral_types))
            already_fit_list += unique_fit_dihedral_types
            
            ####################################################################################################################
            # Run dynamics with the fit potentials set to zero and calculate boltmann inverted potentials effective potentials #
            ####################################################################################################################
            for c in range(inner_cycles):

                # Initialize the start time and print message and gather dihedral_atoms for all instances of the type being fit.
                start_time = Time.time()

                # Gather the different dihedral_styles
                Dihedral_styles = [ FF_dict["dihedrals"][k][qc][0] for k in list(FF_dict["dihedrals"].keys()) ]
                Dihedral_styles += [ FF_dict["dihedrals_harmonic"][k][qc][0] for k in list(FF_dict["dihedrals_harmonic"].keys()) ]
                Dihedral_styles = set(Dihedral_styles)

                # Perform the optimized scans, while varying the intial angle constraint (z).
                E_con = np.zeros(len(list(fit_constraints.keys())))
                angle_tuples = [()]*len(list(fit_constraints.keys()))
                geos = [np.zeros([len(keep_ind),3])]*len(list(fit_constraints.keys()))

                # Assign charge scale factor (when enabled, this linearly ramps the coulombic interactions from 0 to 1 over cycles c1=0..min_cycles)
                # NOTE: the sqrt factor comes from the fact that we are scaling the charges and so the coulomb interaction scales as the square of the charge scaling factor.
                if ramp_coul == 1 and (c1+1) < min_cycles:
                    coul_scale = ((1.0-float(min_cycles-c1-1)/float(min_cycles-1)))**(0.5)
                else:
                    coul_scale = 1.0

                # For all cycles besides the first, the dihedrals are fit to the residual difference between the FF potential, less the modes being fit, and the QC potential.
                # NOTE: this option is intentionally disabled for now. During testing this idea was mooted, but in general better stability was acheived by fitting to the residual from the start.
                if c1 > -1:

                    # Print message 
                    #print "\n\t\tGenerating relaxed FF-based configurations via lammps..."                

                    # Iterate over starting configurations and perform geometry optimizations
                    for z in range(len(list(fit_constraints.keys()))):

                        # Grab the QC optimized geometry for this scan
                        Geometry  = np.array([ Dihedral_Data[i][str(z)]["Geo"][k] for k in keep_ind ])

                        # Initialize Sim_Box
                        Sim_Box = np.array([min(Geometry[:,0])-100.0,max(Geometry[:,0])+100.0,min(Geometry[:,1])-100.0,max(Geometry[:,1])+100.0,min(Geometry[:,2])-100.0,max(Geometry[:,2])+100.0])

                        # Write the lammps datafile (the function returns a dictionary, Atom_type_dict, holding the mapping between
                        # atom_types and the lammps id numbers; this mapping is needed for setting fixes)                    
                        Atom_type_dict = Write_OPLSFIT_MD_data("optscan",Atomtypes,Sim_Box,Elements,Geometry,Bonds,Bond_types,FF_dict["bonds"],Angles,Angle_types,FF_dict["angles"],\
                                                               Dihedrals,Dihedral_types,FF_dict["dihedrals"],FF_dict["dihedrals_harmonic"],One_fives,Charges,FF_dict["vdw"],FF_dict["masses"],fit_dihedrals,\
                                                               qc,one_five_opt=0,coul_scale=coul_scale)

                        # Write the input file for the forward constrained scan starting from configuration z
                        Write_QCMATCH_input("optscan",100,0.0,"lj/cut/coul/cut 100.0 100.0",Dihedral_styles,fit_constraints,other_constraints,mode='single',start_ind=z,k_sec=10.0)

                        # Run the LAMMPS minimization
                        lammps_call = "{} -in {}.in.init -screen {}.screen -log {}.log".format(lammps_exe,"optscan","optscan","optscan")
                        subprocess.call(lammps_call.split())

                    # Grab the minimized geometry from the forward scan, the angles of all dihedral instances being fit, and the FF energy
                    for j in range(len(list(fit_constraints.keys()))):
                        current_ind = j
                        if current_ind >= len(list(fit_constraints.keys())): current_ind = current_ind - len(list(fit_constraints.keys()))
                        new_geo,current_E_FF,angles = parse_OPLSFIT_traj("{}.lammpstrj".format(j),Atomtypes,[ _*coul_scale for _ in Charges ],Adj_mat,fit_dihedrals,qc,UA_opt,\
                                                                         one_five_opt=0,include_con=0,one_four_scale_coul=one_four_scale_coul,one_four_scale_vdw=one_four_scale_vdw)

                        # Remove run data 
                        os.remove("{}.lammpstrj".format(j))

                        # assign energies/angles/geoemetries
                        E_con[current_ind]         = current_E_FF
                        angle_tuples[current_ind]  = tuple([ angles[k] for k in fit_dihedrals ])
                        geos[current_ind]          = new_geo                

                    # Align the geometries
                    align_geos(geos,fit_constraints[list(fit_constraints.keys())[0]][0][0])                

                    # Write the optimized scan 
                    with open("{}_optscan.xyz".format("_".join(i)),'w') as f:
                        for count_j,j in enumerate(geos):
                            f.write("{}\nE_FF: {:< 12.6f}\n".format(len(j),E_con[count_j]))
                            for count_k,k in enumerate(j):
                                f.write("{:<12s} {:< 12.6f} {:< 12.6f} {:< 12.6f}\n".format(Elements[count_k],k[0],k[1],k[2]))

                    # Update the residual potential for this fit
                    res_potential = E_QC - E_con
                    res_potential = res_potential - np.mean(res_potential)
                    
                # For the first cycle, the dihedrals are fit directly to the raw quantum chemistry configurations/energies. This generates a reasonable initial guess entering into the relaxed fits.
                else:

                    # Print message 
                    print("\n\t\tPerforming OPLS fit directly to the QC potential and configurations...")

                    # Use the raw QC potential for the fit
                    res_potential = E_QC

                    # Iterate over starting configurations and populate qc dihedrals
                    for z in range(len(list(fit_constraints.keys()))):

                        # I don't think this is needed.
                        current_ind = z

                        # Grab the QC optimized geometry for this scan
                        geos[current_ind] = np.array([ Dihedral_Data[i][str(z)]["Geo"][k] for k in keep_ind ])
                        geo = geos[current_ind]
                        
                        # Calculate the constrained dihedral angles (this section of code is taken from the OPLS_parse function)
                        angles = { i:[] for i in fit_dihedrals }
                        for k in fit_dihedrals:

                            # Calculate all dihedral angles
                            atom_1 = geo[k[0]]
                            atom_2 = geo[k[1]]
                            atom_3 = geo[k[2]]
                            atom_4 = geo[k[3]]
                            v1 = atom_2-atom_1
                            v2 = atom_3-atom_2
                            v3 = atom_4-atom_3
                            angles[k] = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )*180.0/np.pi        
                        angle_tuples[current_ind] = tuple([ angles[k] for k in fit_dihedrals ])

                ################################################################
                # Fit OPLS potential(s) to the residual potential(s), save the #
                # fit(s) to FF_dict and plot a comparison                      #
                ################################################################
                
                # Set the v_range, qc_style, angles and dihedral_types for the fit
                # NOTE: under the current implementation, the fit is never "Unrestricted"
                if c1 >= max_cycles:
                    V_range = 1E10  # put no constraints on the V_range. Might be ammended in the future.
                    print("\n\t\t{:<15s} {}".format("V_range:","Unrestricted"))

                # Use the QC range as the fit range with a practical floor and ceiling
                else:
#                    V_range = 2.0*(max(E_QC) - min(E_QC))                     #ORIGINAL
                    V_range = 2.0*(max(res_potential) - min(res_potential))                     
                    if V_range < 1.0:
                        V_range = 1.0
                    elif V_range > 100.0:
                        V_range = 100.0
                    print("\n\t\t{:<15s} {:<4.2f}".format("V_range:",V_range))

                # Initialize variables for the fit 
                w_pot = 1.0
                w_exp = 0.0
                w_hyper = 0.0
                b_hyper = 0.1
#                if max(E_QC) - min(E_QC) > 5.0:
                if max(E_QC) - min(E_QC) > 0.0:
#                    print "\t\tThe range of the fit potential is greater than 5 kcal/mol, applying harmonic constraint during the fit."
                    w_harm = (max(res_potential) - min(res_potential))*0.001 # SCALED RELATIVE TO THE FIX BEING PERFORMED
                else:
                    w_harm = 0.0
                dihedral_types = fit_dihedral_types
                unique_fit_dihedral_types = list(set(fit_dihedral_types))
                last_err = 0.0
                print("\t\t{:<15s} {:<4.2f}".format("coul_scale:",coul_scale))
                print("\t\t{:<15s} {:<4.2f}".format("L2 scale:",w_harm))
                print("\t\t{:<15s} {:<4.2f}".format("update_weight:",weight))


                # Initialize guesses and parameter bounds
                # the initial_guess is simple a long list of parameters with each 5 indexed to a particular dihedral
                # the first of the five is a V0 term that is presently unused, but might be useful in some future context
                # as a normalization constant
                initial_guess = []
                bounds = []
                for count_f,f in enumerate(unique_fit_dihedral_types):
                    if len(FF_dict["dihedrals"][f][qc]) == 6:
                        initial_guess += FF_dict["dihedrals"][f][qc][1:]
                    else:
                        initial_guess += ([0.0]+FF_dict["dihedrals"][f][qc][1:])                    
#                    initial_guess += [0.0,0.0,0.0,0.0,0.0]
                    bounds += [(None,None),(-V_range,V_range),(-V_range,V_range),(-V_range,V_range),(-V_range,V_range)]

                # Check for 1-2-3 or 2-3-4 linearity.
                lin_break_flag = 0 
                for g in geos:
                    for f in fit_dihedrals:
                        if np.arccos(np.dot(g[f[0]]-g[f[1]],g[f[2]]-g[f[1]])/(norm(g[f[0]]-g[f[1]])*norm(g[f[2]]-g[f[1]]))) * 180.0/np.pi > 160.0 or \
                            np.arccos(np.dot(g[f[3]]-g[f[2]],g[f[1]]-g[f[2]])/(norm(g[f[3]]-g[f[2]])*norm(g[f[1]]-g[f[2]]))) * 180.0/np.pi > 160.0:
                            lin_break_flag = 1
                        if lin_break_flag == 1:
                            fit_err = 0.0
                            break
                    if lin_break_flag == 1:
                        break
                
                # Avoid fitting linear dihedrals
                if lin_break_flag == 0:

                    # All coincident dihedrals are self-consistently fit at once, the registration between the input arguments in 
                    # OPLS_OPLS_Fit_MIN and the individual dihedral types is acheived using the unique_fit_dihedral_tOAypes list
                    fit_func = lambda x: OPLS_OPLS_Fit_MIN(*x,unique_fit_dihedral_types=unique_fit_dihedral_types,angles=angle_tuples,E_0=res_potential,V_range=V_range,\
                                                              dihedral_types=dihedral_types,QC_dihedral=QC_dihedral,FF_dict=FF_dict,w_pot=w_pot,w_hyper=w_hyper,b_hyper=b_hyper,w_exp=w_exp,w_harm=w_harm)                        

                    params = minimize(fit_func,initial_guess,bounds=bounds,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':1E-20}).x
                    print(params)

                    # Perform a weighted update of the force-field dictionary 
                    for count_f,f in enumerate(unique_fit_dihedral_types):
                        CURRENT_FF_dict["dihedrals"][f][qc] = [CURRENT_FF_dict["dihedrals"][f][qc][0]] + [ params[0] ] + list(params[(count_f*5+1):(count_f+1)*5])
                        CURRENT_FF_dict["dihedrals"][f][qc] = CURRENT_FF_dict["dihedrals"][f][qc][0:2] + [ weight*CURRENT_FF_dict["dihedrals"][f][qc][j+1]+(1.0-weight)*FF_dict["dihedrals"][f][qc][j] for j in range(1,5) ]

                    # Calculate the average error in the fit and check the satisfaction of break condition(s)
                    fit_err = np.mean((OPLS_Calc(angle_tuples,dihedral_types,CURRENT_FF_dict,QC_dihedral)-res_potential)**(2.0))
                    fit_err_dict[i] = fit_err

                    # Remove the V0 term from CURRENT_FF_dict
                    for f in unique_fit_dihedral_types:
                        CURRENT_FF_dict["dihedrals"][f][qc] = [CURRENT_FF_dict["dihedrals"][f][qc][0]] + CURRENT_FF_dict["dihedrals"][f][qc][2:]

                # Else, print message about what has occured
                else:
                    print("\t\t* Linearity in the 1-2-3/2-3-4 angle(s) was discovered. Leaving dihedral undefined *")
                    
                # Recalculate the fit potential using the weighted update force-field parameters and update the percent change (difference in range over average range of the dihedral)  
                fit_potential = OPLS_Calc(angle_tuples,dihedral_types,CURRENT_FF_dict,QC_dihedral)
                new_range = max(fit_potential) - min(fit_potential)
                old_range = max(old_fit_pot[i]) - min(old_fit_pot[i])
                old_fit_pot[i]     = fit_potential
                if new_range + old_range > 0.0:
                    percent_updates[i] = np.abs( new_range - old_range ) / ( 0.5 * (new_range + old_range) ) * 100.0
                else:
                    percent_updates[i] = 0.0

                # Print diagnostic on updated parameters
                print("\n\t\tInitial params for this scan:")
                for f in unique_fit_dihedral_types:
                    print("\t\t{} {}".format(" ".join([ "{:<40s}".format(j) for j in f ])+": ","{:<10s} ".format(FF_dict["dihedrals"][f][qc][0])+" ".join(["{:< 6.4f}".format(float(j)) for j in FF_dict["dihedrals"][f][qc][1:]])))
                print("\n\t\tUpdated params for this scan:")
                for f in unique_fit_dihedral_types:
                    print("\t\t{} {}".format(" ".join([ "{:<40s}".format(j) for j in f ])+": ","{:<10s} ".format(CURRENT_FF_dict["dihedrals"][f][qc][0])+" ".join(["{:< 6.4f}".format(float(j)) for j in CURRENT_FF_dict["dihedrals"][f][qc][1:]])))

                # Test calc thermal potentials at each step
#                t0 = time.time()
#                thermal_potentials,thermal_angles = gen_thermal_distribution(np.array([Dihedral_Data[i]["0"]["Geo"][k] for k in keep_ind ]),Adj_mat,Atomtypes,Charges,fit_dihedrals,fit_dihedral_types,\
#                                                                             lammps_exe,t_sample=1E6,freq=100,T=298,bin_width=20.0,qc=qc)
 #               print "\t\tIt took {} seconds to generate the thermal potentials...".format(time.time()-t0)

                # Generate 2D representation of the dihedral
#                print "\n\t\tGenerating 2D representation of the fit fragment..."                    
#                geo_2D = kekule(Elements,Atomtypes,geos[0],Adj_mat)
                geo_2D = kekule(Elements,id_types(Elements,Table_generator(Elements,geos[0]),gens=gens),geos[0],Table_generator(Elements,geos[0]))

                # Generate comparisdons of the QC, FF, and fit potentials 
                # NOTE: all potentials were generated relative to -180,180 interval of the scanned dihedral
                #       thus the values need to be shifted when plotting the other angles
                if save_intermediates == 1:
                    for d in unique_fit_dihedral_types:

                        # Generate Label
                        label = atom_to_element[int(d[0].split('[')[1].split(']')[0])]+'-'+\
                                atom_to_element[int(d[1].split('[')[1].split(']')[0])]+'-'+\
                                atom_to_element[int(d[2].split('[')[1].split(']')[0])]+'-'+\
                                atom_to_element[int(d[3].split('[')[1].split(']')[0])]

                        # Generate savename
                        savename = d[0]+'_'+d[1]+'_'+d[2]+'_'+d[3]

                        # Shift angle values and sort
                        dihedral_ind   = [ count_e for count_e,e in enumerate(fit_dihedral_types) if e == d ][0]
                        current_angles = [ angle_tuples[j][dihedral_ind] for j in range(len(angle_tuples)) ]
                        current_QC     = deepcopy(E_QC)
                        current_fit    = deepcopy(fit_potential)
                        current_res    = deepcopy(res_potential)
                        current_E_con  = deepcopy(E_con)
                        current_angles,current_QC,current_fit,current_res,current_E_con = list(zip(*sorted(zip(current_angles,current_QC,current_fit,current_res,current_E_con))))
                        current_angles = np.array(current_angles)
                        current_QC     = np.array(current_QC)
                        current_fit    = np.array(current_fit)
                        current_res    = np.array(current_res)
                        current_E_con  = np.array(current_E_con)

                        # Initialize Figure with two sets of axes (top is for chemical structure bottom is for the fit)
                        fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                        draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                        if draw_scale > 1.0: draw_scale = 1.0

                        # Add bonds
                        linewidth = 2.0*draw_scale*5.0
                        if linewidth > 2.0:
                            linewidth = 2.0
                        if linewidth < 1.0:
                            linewidth = 1.0
                        for count_j,j in enumerate(Adj_mat):
                            for count_k,k in enumerate(j):
                                if count_k > count_j:
                                    if k == 1:
                                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_j][0],geo_2D[count_k][0]],[geo_2D[count_j][1],geo_2D[count_k][1]],color=(0,0,0),linewidth=linewidth))

                        # Add Highlights
                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][0]][0],geo_2D[fit_dihedrals[dihedral_ind][1]][0]],[geo_2D[fit_dihedrals[dihedral_ind][0]][1],geo_2D[fit_dihedrals[dihedral_ind][1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][1]][0],geo_2D[fit_dihedrals[dihedral_ind][2]][0]],[geo_2D[fit_dihedrals[dihedral_ind][1]][1],geo_2D[fit_dihedrals[dihedral_ind][2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][2]][0],geo_2D[fit_dihedrals[dihedral_ind][3]][0]],[geo_2D[fit_dihedrals[dihedral_ind][2]][1],geo_2D[fit_dihedrals[dihedral_ind][3]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                        # Add atom labels
                        fontsize = 20*draw_scale*5.0
                        if fontsize > 24:
                            fontsize = 24
                        if fontsize < 9:
                            fontsize = 9            
                        for count_j,j in enumerate(geo_2D):
                            if count_j in fit_dihedrals[dihedral_ind]:
                                ax1.text(j[0], j[1], Elements[count_j], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                            else:
                                ax1.text(j[0], j[1], Elements[count_j], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                        ax1.axis('image')
                        ax1.axis('off')

                        # Plot the fit potential comparisons and total potential comparisons
                        ax2.scatter(current_angles,current_QC-current_QC[0],s=80,color=(0.05,0.05,0.05,1.0))
                        ax2.plot(current_angles,current_QC-current_QC[0],'--',linewidth=3,color=(0.05,0.05,0.05,1.0))
                        ax2.scatter(current_angles,(current_fit-current_fit[0])+(current_E_con-current_E_con[0]),s=80,color=(0.0,0.1,0.8,1.0))
                        ax2.plot(current_angles,(current_fit-current_fit[0])+(current_E_con-current_E_con[0]),'--',linewidth=3,color=(0.0,0.1,0.8,1.0))
                        ax2.scatter(current_angles,current_res-current_res[0],s=80,color=(0.8,0.1,0.0,1.0))
                        ax2.plot(current_angles,current_res-current_res[0],'--',linewidth=3,color=(0.8,0.1,0.0,1.0))
                        ax2.scatter(current_angles,current_fit-current_fit[0],s=80,color=(0.0,0.8,0.1,1.0))
                        ax2.plot(current_angles,current_fit-current_fit[0],'--',linewidth=3,color=(0.0,0.8,0.1,1.0))

                        # Thermal potentials
    #                    ax2.scatter(thermal_angles,thermal_potentials[d]-thermal_potentials[d][0],s=80,color=(0.7,0.1,1.0,1.0))
    #                    ax2.plot(thermal_angles,thermal_potentials[d]-thermal_potentials[d][0],'--',linewidth=3,color=(0.7,0.1,1.0,1.0))

                        # Format the plot            
    #                    ax2.set_title('Dihedral {}'.format(label),fontsize=12,y=1.08)
                        ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                        ax2.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
                        ax2.set_xlim([-180, 180])
                        ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                        ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                        [k.set_linewidth(3) for k in ax2.spines.values()]
                        plt.tight_layout()

                        # Create save name, and if this is the final cycle save the name to the figure list variable 
                        # for inclusion in the results summary.
                        Name = '{}_Fit_{}.png'.format(savename,c1)
                        paths_to_dihedral_fitpot_dict[d] = [save_folder+'/'+Name]
                        plt.savefig(Name, bbox_inches=0,dpi=300)
                        plt.close(fig)

                        # Initialize Figure with two sets of axes (top is for chemical structure bottom is for the fit)
                        fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
                        draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
                        if draw_scale > 1.0: draw_scale = 1.0

                        # Add bonds
                        linewidth = 2.0*draw_scale*5.0
                        if linewidth > 2.0:
                            linewidth = 2.0
                        if linewidth < 1.0:
                            linewidth = 1.0
                        for count_j,j in enumerate(Adj_mat):
                            for count_k,k in enumerate(j):
                                if count_k > count_j:
                                    if k == 1:
                                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_j][0],geo_2D[count_k][0]],[geo_2D[count_j][1],geo_2D[count_k][1]],color=(0,0,0),linewidth=linewidth))

                        # Add Highlights
                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][0]][0],geo_2D[fit_dihedrals[dihedral_ind][1]][0]],[geo_2D[fit_dihedrals[dihedral_ind][0]][1],geo_2D[fit_dihedrals[dihedral_ind][1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][1]][0],geo_2D[fit_dihedrals[dihedral_ind][2]][0]],[geo_2D[fit_dihedrals[dihedral_ind][1]][1],geo_2D[fit_dihedrals[dihedral_ind][2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                        ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][2]][0],geo_2D[fit_dihedrals[dihedral_ind][3]][0]],[geo_2D[fit_dihedrals[dihedral_ind][2]][1],geo_2D[fit_dihedrals[dihedral_ind][3]][1]],color=(0,0.1,0.8),linewidth=linewidth))

                        # Add atom labels
                        fontsize = 20*draw_scale*5.0
                        if fontsize > 24:
                            fontsize = 24
                        if fontsize < 9:
                            fontsize = 9            
                        for count_j,j in enumerate(geo_2D):
                            if count_j in fit_dihedrals[dihedral_ind]:
                                ax1.text(j[0], j[1], Elements[count_j], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                            else:
                                ax1.text(j[0], j[1], Elements[count_j], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                        ax1.axis('image')
                        ax1.axis('off')

                        # Plot the fit potential comparisons and total potential comparisons
                        ax2.scatter(current_angles,current_QC-current_QC[0],s=80,color=(0.05,0.05,0.05,1.0))
                        ax2.plot(current_angles,current_QC-current_QC[0],'--',linewidth=3,color=(0.05,0.05,0.05,1.0))
                        ax2.scatter(current_angles,(current_fit-current_fit[0])+(current_E_con-current_E_con[0]),s=80,color=(0.0,0.1,0.8,1.0))
                        ax2.plot(current_angles,(current_fit-current_fit[0])+(current_E_con-current_E_con[0]),'--',linewidth=3,color=(0.0,0.1,0.8,1.0))

                        # Format the plot            
                        # ax2.set_title('Dihedral {}'.format(label),fontsize=12,y=1.08)
                        ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
                        ax2.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
                        ax2.set_xlim([-180, 180])
                        ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
                        ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
                        [k.set_linewidth(3) for k in ax2.spines.values()]
                        plt.tight_layout()

                        # Create save name, and if this is the final cycle save the name to the figure list variable 
                        # for inclusion in the results summary.
                        Name = '{}_Tot_{}.png'.format(savename,c1)
                        plt.savefig(Name, bbox_inches=0,dpi=300)
                        plt.close(fig)
                
                # Print the statistics for this cycle
                print("\n\t\t{:50s} {}".format("Xhi2 error in scan fit (kcal^2/mol^2):",fit_err))                                
                print("\t\t{:50s} {}".format("Time for running this cycle (seconds):",Time.time()-start_time))
                if c1 > 0 and lin_break_flag == 0:
                    print("\t\t{:50s} {}".format("Mean percent change in potential range:","{:1.6f} %".format(percent_updates[i])))
                                             
                # Update break condition for this scan
                if c1 >= min_cycles and percent_updates[i] < conv_thresh:
                    break_flags[i] = 1
        
            # Return to working directory
            os.chdir(current_dir)
        
        # Remove V0 terms from CURRENT_FF_dict
        for i in  list(CURRENT_FF_dict["dihedrals"].keys()):
            if len(CURRENT_FF_dict["dihedrals"][i][qc]) == 6:
                CURRENT_FF_dict["dihedrals"][i][qc] = [CURRENT_FF_dict["dihedrals"][i][qc][0]] + CURRENT_FF_dict["dihedrals"][i][qc][2:]
        
        # Update the FF_dict after each full cycle of scans
        FF_dict["dihedrals"] = CURRENT_FF_dict["dihedrals"]
    
    # If the last cycle was reached and the dihedrals are unconverged, print a message
    if max_cycles > 0 and c1 == max_cycles-1 and False in [ percent_updates[i] < 1.0 for i in keys ]:
        print("\n{:^167s}\n{:^167s}\n{:^167}\n".format("*"*87,"* {:^83s} *".format("Fit complete (max cycles reached)"),"*"*87))
    # Print completion message (only gets triggered if the job completed on the last cycle)
    elif max_cycles > 0 and c1 == max_cycles-1 and False not in [ break_flags[i]==1 for i in list(break_flags.keys()) ]:
        print("\n{:^167s}\n{:^167s}\n{:^167}\n".format("*"*87,"* {:^83s} *".format("Fit complete (all dihedrals converged within {:1.2f}%)".format(conv_thresh)),"*"*87))

    ################################################################
    # Run scans with the final set of parameters. The code between #
    # the "for i" loop start and the "for z" loop start is copied  #
    # directly from the self-consistency loop above. The rest is   #
    # unique to the final scan.                                    #
    ################################################################

    # Loop over the dihedral scans and perform fits to all coincident dihedrals in each scan
    already_fit_list = []
    for count_i,i in enumerate(keys):

        #############################################################################################################################
        # Generate prerequisite data for the scan (NOTE: the program assumes atom ordering is retained across the scan. This is     #
        # guarranteed by the generation scripts, but can lead it errors if other programs don't also abide by the same convention.  #
        #############################################################################################################################

        # Skip if UA_opt == 1 and there are any UA-hydrogens are in this dihedral
        if UA_opt == 1 and 1 in [ return_UA_H(k) for k in i ]: continue

        # Check if all runs for this scan succesfully completed
        incomplete = []
        for k in list(Dihedral_Data[i].keys()):
            if Dihedral_Data[i][k]["Completion"] == 0: incomplete += [ int(i) ]
        if len(incomplete) > 0:
            continue

        # Print diagnostic
        print("\n\t\tRunning MD to sample the final optimized potential(s) for scan {}...".format(i))

        # Initialize the save_folder name and check if the subfolder needs to be created
        if UA_opt == 1: save_folder = Folder+'/dihedrals/'+"_".join([ k+"-UA" for k in i ])+"-MD"
        else: save_folder = Folder+'/dihedrals/'+"_".join(i)+"-MD"
        if os.path.isdir(save_folder) == False:
            os.makedirs(save_folder)

        # Change directory into the working subfolder
        os.chdir(save_folder)

        # Update the charges from the force-field if all of them are available.
        Charges = deepcopy(eq_charges[Direct_keys[count_i]])
        if UA_opt == 0:
            tmp_types = [ next( m for m in k.split('link-') if m != '') for k in Dihedral_Data[i]["0"]["Atomtypes"] ]
        elif UA_opt == 1:
            tmp_types = [ next( m for m in k.split('link-') if m != '')+'-UA' for k in Dihedral_Data[i]["0"]["Atomtypes"] ]
        if charge_origin is not None:

            # Update the charges with the FF charges is they are all available
            if True not in [ charge_origin[k] == "eq_config" for k in tmp_types ]:
                for k in range(len(Charges)): Charges[k] = FF_dict["charges"][tmp_types[k]]

            # Else, proceed with eq_charges
            else:
                print("WARNING: Missing partial charges for some or all of the atomtypes in this fragment, defaulting to equilibrium partial charges.") 

        # If the provenance of the charges is unknown, then proceed by using the FF charges if they are all available
        elif False not in [ tmp_types[k] in list(FF_dict["charges"].keys()) for k in range(len(Charges)) ]:
            for k in range(len(Charges)): Charges[k] = FF_dict["charges"][tmp_types[k]]

        # Generate the charges for the scan (taffi charges for the scanned dihedrals and eq_charges for the rest)
        scanned_atoms = set([ m for k in Dihedral_Data[i]["0"]["Dihedral_Atoms"] for m in k[0] ])

        # Recalculate the atom types if this is a ring dihedral being fit to a linear analog
        # These are used later for reassigning the bond and angle types
        if "R" in i[1] and "R" in i[2]:
            atomtypes_lin = id_types([ atom_to_element[int(k.split('[')[1].split(']')[0])] for k in Dihedral_Data[i]["0"]["Atomtypes"] ],Dihedral_Data[i]["0"]["Adj_mat"],gens=gens)

        # Flag link atoms for removal.(only matters for fragment based parametrizations)
        # NOTE: the most up to date convention is to use link-* for spectator atoms, but [0] and [1] are included for legacy reasons.
        # NOTE: when "del" is specified the link atoms are deleted, otherwise the link atoms are substituted for their appropriate atomtype. 
        if link_mode == 'del':
            keep_ind = [ count_k for count_k,k in enumerate(Dihedral_Data[i]["0"]["Atomtypes"]) if ( k != "[0]" and k != "[1]" and "link-" not in k ) ]
        elif link_mode == 'sub':
            keep_ind = [ count_k for count_k,k in enumerate(Dihedral_Data[i]["0"]["Atomtypes"]) if ( k != "[0]" and k != "[1]" ) ]

        # If UA_opt == 1 then remove UA-hydrogens. return_UA_H determines if an atomtype is a UA-H based on its label.
        if UA_opt == 1:
            keep_ind = [ k for k in keep_ind if return_UA_H(Dihedral_Data[i]["0"]["Atomtypes"][k]) == 0 ]

        # Initialize requisite lists/arrays, minus any spectator hydrogens
        if UA_opt == 0: Atomtypes = [ Dihedral_Data[i]["0"]["Atomtypes"][k] for k in keep_ind ]
        if UA_opt == 1: Atomtypes = [ Dihedral_Data[i]["0"]["Atomtypes"][k]+"-UA" for k in keep_ind ]
        Atomtypes = [ next( m for m in k.split('link-') if m != '') for k in Atomtypes ]
        Adj_mat   = Dihedral_Data[i]["0"]["Adj_mat"][keep_ind,:]
        Adj_mat   = Adj_mat[:,keep_ind]
        Elements  = [ atom_to_element[int(k.split('[')[1].split(']')[0])] for k in Atomtypes ]
        Charges = [ Charges[k] for k in keep_ind ]
        Masses    = [ FF_dict["masses"][next( m for m in k.split('link-') if m != '' )] for k in Atomtypes ]
        Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(Adj_mat,Atomtypes,return_all=1)

        # If this is a flexible ring dihedral then the bonds and angles needs to be reassigned as the corresponding linear types
        # This is because the flexible ring dihedrals are fit with linear analaogs of the cycles they belong to. 
        if "R" in i[1] and "R" in i[2]:

            # Trim atomtypes_lin to conform with keep_ind             
            atomtypes_lin = [ atomtypes_lin[k] for k in keep_ind ]

            # Reassign bonds using atomtypes_lin atomtypes
            for count_j,j in enumerate(Bonds):
                Bond_types[count_j],Bonds[count_j] = canon_bond(tuple([ atomtypes_lin[k] for k in j ]),ind=Bonds[count_j])

            # Reassign angles using atomtypes_lin atomtypes
            for count_j,j in enumerate(Angles):
                Angle_types[count_j],Angles[count_j] = canon_angle(tuple([ atomtypes_lin[k] for k in j ]),ind=Angles[count_j])

        # Determine bonding matrix for the compound
        q_tmp = int(round(sum(eq_charges[Direct_keys[count_i]])))
        bond_mat = find_lewis([ next( m for m in k.split('link-') if m != '') for k in Dihedral_Data[i]["0"]["Atomtypes"] ],\
                              Dihedral_Data[i]["0"]["Adj_mat"], q_tot=q_tmp, b_mat_only=True,verbose=False)
        bond_mat = [ k[keep_ind,:][:,keep_ind] for k in bond_mat ]### IMPT

        # Add the opls/harmonic type to the Dihedral_types tuples
        # NOTE: due to the hydrogenation procedure of fragments, sometimes non-scanned dihedrals end up needing to be
        #       modeled as harmonic or flexible. The following if/else statements check the need to swapping the type 
        #       of each dihedral as needed. 
        for j in range(len(Dihedrals)):
            if 2 in [ k[Dihedrals[j][1],Dihedrals[j][2]] for k in bond_mat ]:
                if Dihedral_types[j] in list(FF_dict["dihedrals_harmonic"].keys()):                        
                    dihedral_type = "harmonic"
                elif Dihedral_types[j] not in list(FF_dict["dihedrals"].keys()):
                    print("ERROR in parse_dihedrals_qcmatch: Dihedral {} ({}) is being processed as a harmonic type, but it is not in either the harmonic or flexible dihedral dictionaries.".format(Dihedral_types[j],Dihedral[j]))
                    quit()
                else:
                    print("\t\tWARNING: Dihedral {} ({}) is being processed as a harmonic type, but it is being modeled as a flexible dihedral.".format(Dihedral_types[j],Dihedrals[j]))
                    dihedral_type = "opls"
            else:
                if Dihedral_types[j] in list(FF_dict["dihedrals"].keys()):                        
                    dihedral_type = "opls"
                elif Dihedral_types[j] not in list(FF_dict["dihedrals_harmonic"].keys()):
                    print("ERROR in parse_dihedrals_qcmatch: Dihedral {} ({}) is being processed as a flexible type, but it is not in either the flexible or harmonic dihedral dictionaries.".format(Dihedral_types[j],Dihedral[j]))
                    quit()
                else:
                    print("\t\tWARNING: Dihedral {} ({}) is being processed as a flexible type, but it is being modeled as a harmonic dihedral.".format(Dihedral_types[j],Dihedrals[j]))
                    dihedral_type = "harmonic"

            # Add dihedral type to the types information
            Dihedral_types[j] = tuple(list(Dihedral_types[j])+[dihedral_type])

        # Update the constrained atoms indices 
        # NOTE: since some atoms may have been removed, the "Dihedral_Atoms" entry is no longer accurate. The
        # use of the tmp array here is a simple way of getting the appropriate indexing.
        fit_dihedrals = []
        fit_dihedral_types = []
        con_map = {}
        for j in (Dihedral_Data[i]["0"]["Dihedral_Atoms"]+Dihedral_Data[i]["0"]["Unscanned_Dihedral_Atoms"]):
            tmp = [0] * len(Dihedral_Data[i]["0"]["Atomtypes"])                                 # All zeros initialization of tmp of equal length to the original atomtypes list
            for count_k,k in enumerate(j[0]): tmp[k] = count_k+1                                # label the 1,2,3, and 4 atoms of the dihedral
            tmp = [ tmp[k] for k in keep_ind ]                                                  # only keep the non-spectator atoms (as in the lists above)

            # Be cautious to only add "fit" dihedrals to fit_dihedrals
            if j in Dihedral_Data[i]["0"]["Dihedral_Atoms"]:                    
                fit_dihedrals += [tuple([ tmp.index(k) for k in range(1,5) ])]                  # identify the indices based on where the 1,2,3, and 4 labels are in tmp, add the current dihedral to the list
                fit_dihedral_types += [tuple([ Atomtypes[tmp.index(k)] for k in range(1,5) ])]  # identify the types based on where the 1,2,3, and 4 labels are in tmp, add the current dihedral to the list
                con_map[j[0]] = fit_dihedrals[-1]                                               # add the mapping of the current dihedral to the constrained dictionary.

            # Add the rest of the dihedrals provided they haven't been trimmed during
            # the removal of link atoms
            elif all([ k in tmp for k in range(1,5) ]):
                con_map[j[0]] = tuple([ tmp.index(k) for k in range(1,5) ])                     # add the mapping of the current Link-atom dihedral to the constrained dictionary.                    

        # Remove any dihedral types that were fitted in another scan
        del_ind = [ count_j for count_j,j in enumerate(fit_dihedral_types) if j in already_fit_list or j+tuple("opls") in modes_from_FF ]
        fit_dihedrals = [ j for count_j,j in enumerate(fit_dihedrals) if count_j not in del_ind ]
        fit_dihedral_types = [ j for count_j,j in enumerate(fit_dihedral_types) if count_j not in del_ind ]

        # Collect the angle constraints for each dihedral instance being fit and the QC energies for each configuration
        fit_constraints = {}
        E_QC = np.zeros(len(list(Dihedral_Data[i].keys())))
        for count_j,j in enumerate(natural_sort(list(Dihedral_Data[i].keys()))):
            E_QC[count_j] = Dihedral_Data[i][j][qc]
            fit_constraints[count_j] = []

            # Add scanned atom constraints
            for count_k,k in enumerate(Dihedral_Data[i][j]["Dihedral_Atoms"]):
                fit_constraints[count_j] += [(con_map[k[0]],k[1]*180.0/np.pi)]

        # Adjust the angles by +/- 360 if necessary to perform the shortest interpolation between the two configurations
        for j in sorted(fit_constraints.keys())[1:]:
            for count_k,k in enumerate(fit_constraints[j]):
                angle = k[1]
                diff = np.abs(fit_constraints[j-1][count_k][1] - angle)
                if np.abs(fit_constraints[j-1][count_k][1] - (angle+360.0)) < diff:
                    diff = np.abs(fit_constraints[j-1][count_k][1] - (angle+360.0))
                    fit_constraints[j][count_k] = (fit_constraints[j][count_k][0],angle+360.0)
                if np.abs(fit_constraints[j-1][count_k][1] - (angle-360.0)) < diff:
                    fit_constraints[j][count_k] = (fit_constraints[j][count_k][0],angle-360.0)

        # Collect the angle constraints for secondary dihedrals that aren't being fit
        other_constraints = {}
        for count_j,j in enumerate(natural_sort(list(Dihedral_Data[i].keys()))):
            other_constraints[count_j] = []

            # Add link atom contraints
            if link_mode == 'sub' or link_mode == 'del':
                for count_k,k in enumerate(Dihedral_Data[i][j]["Unscanned_Dihedral_Atoms"]):
                    if k[0] in con_map:
                        other_constraints[count_j] += [(con_map[k[0]],k[1]*180.0/np.pi)]

        # Adjust the angles by +/- 360 if necessary to perform the shortest interpolation between the two configurations
        for j in sorted(other_constraints.keys())[1:]:
            for count_k,k in enumerate(other_constraints[j]):
                angle = k[1]
                diff = np.abs(other_constraints[j-1][count_k][1] - angle)
                if np.abs(other_constraints[j-1][count_k][1] - (angle+360.0)) < diff:
                    diff = np.abs(other_constraints[j-1][count_k][1] - (angle+360.0))
                    other_constraints[j][count_k] = (other_constraints[j][count_k][0],angle+360.0)
                if np.abs(other_constraints[j-1][count_k][1] - (angle-360.0)) < diff:
                    other_constraints[j][count_k] = (other_constraints[j][count_k][0],angle-360.0)

        # Save the unique fit types and update already_fit_list (used to avoid fitting the same dihedral as part of multiple scans)
        unique_fit_dihedral_types = sorted(set(fit_dihedral_types))
        already_fit_list += unique_fit_dihedral_types

        # Print diagnostic
        print("\n\t\tFinal params for this scan:")
        for j in unique_fit_dihedral_types:
            print("\t\t{} {}".format(" ".join([ "{:<40s}".format(k) for k in j ])+": ","{:<10s} ".format(FF_dict["dihedrals"][j][qc][0])+" ".join(["{:< 6.4f}".format(float(k)) for k in FF_dict["dihedrals"][j][qc][1:]])))

        # Gather the different dihedral_styles
        Dihedral_styles = [ FF_dict["dihedrals"][k][qc][0] for k in list(FF_dict["dihedrals"].keys()) ]
        Dihedral_styles += [ FF_dict["dihedrals_harmonic"][k][qc][0] for k in list(FF_dict["dihedrals_harmonic"].keys()) ]
        Dihedral_styles = set(Dihedral_styles)

        # Perform the optimized scans, while varying the intial angle constraint (z).
        E_con = np.zeros(len(list(fit_constraints.keys())))
        angle_tuples = [()]*len(list(fit_constraints.keys()))
        geos = [np.zeros([len(keep_ind),3])]*len(list(fit_constraints.keys()))

        # Iterate over starting configurations and perform geometry optimizations
        for z in range(len(list(fit_constraints.keys()))):

            # Grab the QC optimized geometry for this scan
            Geometry  = np.array([ Dihedral_Data[i][str(z)]["Geo"][k] for k in keep_ind ])

            # Initialize Sim_Box
            Sim_Box = np.array([min(Geometry[:,0])-100.0,max(Geometry[:,0])+100.0,min(Geometry[:,1])-100.0,max(Geometry[:,1])+100.0,min(Geometry[:,2])-100.0,max(Geometry[:,2])+100.0])

            # Write the lammps datafile (the function returns a dictionary, Atom_type_dict, holding the mapping between
            # atom_types and the lammps id numbers; this mapping is needed for setting fixes)                    
            Atom_type_dict = Write_OPLSFIT_MD_data("finalscan",Atomtypes,Sim_Box,Elements,Geometry,Bonds,Bond_types,FF_dict["bonds"],Angles,Angle_types,FF_dict["angles"],\
                                                   Dihedrals,Dihedral_types,FF_dict["dihedrals"],FF_dict["dihedrals_harmonic"],One_fives,Charges,FF_dict["vdw"],FF_dict["masses"],fit_dihedrals,qc,one_five_opt=0)

            # Write the input file for the forward constrained scan starting from configuration z
            Write_QCMATCH_input("finalscan",100,0.0,"lj/cut/coul/cut 100.0 100.0",Dihedral_styles,fit_constraints,other_constraints,mode='single',start_ind=z,k_sec=10.0)

            # Run the LAMMPS minimization
            lammps_call = "{} -in {}.in.init -screen {}.screen -log {}.log".format(lammps_exe,"finalscan","finalscan","finalscan")
            subprocess.call(lammps_call.split())

        # Grab the minimized geometry from each optimization, the angles of all dihedral instances being fit, and the FF energy
        # NOTE: This loop differs from the self-consistency version by the include_con=1 option in the parse function. 
        for j in range(len(list(fit_constraints.keys()))):
            current_ind = j
            if current_ind >= len(list(fit_constraints.keys())): current_ind = current_ind - len(list(fit_constraints.keys()))
            new_geo,current_E_FF,angles = parse_OPLSFIT_traj("{}.lammpstrj".format(j),Atomtypes,Charges,Adj_mat,fit_dihedrals,qc,UA_opt,\
                                                             one_five_opt=0,include_con=1,one_four_scale_coul=one_four_scale_coul,one_four_scale_vdw=one_four_scale_vdw)  

            # Remove run data 
            os.remove("{}.lammpstrj".format(j))

            # assign energies/angles/geoemetries
            E_con[current_ind]         = current_E_FF
            angle_tuples[current_ind]  = tuple([ angles[k] for k in fit_dihedrals ])
            geos[current_ind]          = new_geo                

        # Align the geometries
        align_geos(geos,fit_constraints[list(fit_constraints.keys())[0]][0][0])                

        # Calculate fit potential (this was done solely as a consistency check with the above for diagnostic purposes. The include_con=1 
        # option in the parse function leads to the fit potential terms automatically being included in the E_con array). 
        dihedral_fit_type = 'dummy'
        dihedral_types = fit_dihedral_types
        V_range = 1E10
        fit_potential = OPLS_Calc(angle_tuples,dihedral_types,FF_dict,QC_dihedral)

        # Write the optimized scan 
        with open("{}_finalscan.xyz".format("_".join(i)),'w') as f:
            for count_j,j in enumerate(geos):
                f.write("{}\nE_FF: {:< 12.6f}\n".format(len(j),E_con[count_j]))
                for count_k,k in enumerate(j):
                    f.write("{:<12s} {:< 12.6f} {:< 12.6f} {:< 12.6f}\n".format(Elements[count_k],k[0],k[1],k[2]))

        # Generate 2D representation of the dihedral
        print("\n\t\tGenerating 2D representation of the fit fragment...")                    
        geo_2D = kekule(Elements,id_types(Elements,Table_generator(Elements,geos[0]),gens=gens),geos[0],Table_generator(Elements,geos[0]))
#        geo_2D = kekule(Elements,Atomtypes,geos[0],Adj_mat)

        # If the agreement with the scan potential was poor then generate the thermal potential
        for _ in list(fit_err_dict.keys()): fit_err_dict[_] = 0.0
        if fit_err_dict[i] >= 1.0:
            print("\t\tThe minimum energy FF scan had poor agreement with the QC scan. Generating a thermal potential at 500K to assess the behavior of the fit...")
            t0 = Time.time()
            thermal_potentials,thermal_angles = gen_thermal_distribution(np.array([Dihedral_Data[i]["0"]["Geo"][k] for k in keep_ind ]),Adj_mat,Atomtypes,Charges,fit_dihedrals,fit_dihedral_types,\
                                                                             lammps_exe,t_sample=1E6,freq=50,T=500,bin_width=10.0,qc=qc)
 
            # for count_m,m in enumerate(Dihedral_Data[i].keys()[0]):
            #     if count_m == 0:
            #         thermal_potentials,thermal_angles = gen_thermal_distribution(np.array([Dihedral_Data[i][m]["Geo"][k] for k in keep_ind ]),Adj_mat,Atomtypes,Charges,fit_dihedrals,fit_dihedral_types,\
            #                                                                      lammps_exe,t_sample=1E5,freq=1000,T=500,bin_width=10.0,qc=qc)
            #     else:
            #         tmp_potentials,tmp_angles = gen_thermal_distribution(np.array([Dihedral_Data[i][m]["Geo"][k] for k in keep_ind ]),Adj_mat,Atomtypes,Charges,fit_dihedrals,fit_dihedral_types,\
            #                                                                     lammps_exe,t_sample=1E5,freq=1000,T=500,bin_width=10.0,qc=qc)

            #         for n in thermal_potentials.keys():
            #             thermal_potentials[n] += tmp_potentials[n]
            # for n in thermal_potentials.keys():
            #     thermal_potentials[n] = thermal_potentials[n] / float(len(Dihedral_Data[i].keys()))

            print("\t\tIt took {} seconds to generate the thermal potential...".format(Time.time()-t0))

        # Generate comparisdons of the QC, FF, and fit potentials 
        # NOTE: all potentials were generated relative to -180,180 interval of the scanned dihedral
        #       thus the values need to be shifted when plotting the other angles
        print("\t\tSaving fit plots...")
        for d in unique_fit_dihedral_types:

            # Generate Label
            label = atom_to_element[int(d[0].split('[')[1].split(']')[0])]+'-'+\
                    atom_to_element[int(d[1].split('[')[1].split(']')[0])]+'-'+\
                    atom_to_element[int(d[2].split('[')[1].split(']')[0])]+'-'+\
                    atom_to_element[int(d[3].split('[')[1].split(']')[0])]

            # Generate savename
            savename = d[0]+'_'+d[1]+'_'+d[2]+'_'+d[3]

            # Shift angle values and sort
            dihedral_ind   = [ count_e for count_e,e in enumerate(fit_dihedral_types) if e == d ][0]
            current_angles = [ angle_tuples[j][dihedral_ind] for j in range(len(angle_tuples)) ]
            current_QC     = deepcopy(E_QC)
            current_E_con  = deepcopy(E_con)
            current_fit    = deepcopy(fit_potential)
            current_angles,current_QC,current_fit,current_E_con = list(zip(*sorted(zip(current_angles,current_QC,current_fit,current_E_con))))
            current_angles = np.array(current_angles)
            current_QC     = np.array(current_QC)
            current_E_con  = np.array(current_E_con)
            current_fit    = np.array(current_fit)

            # Initialize Figure with two sets of axes (top is for chemical structure bottom is for the fit)
            fig, (ax1, ax2) = plt.subplots(2,1,gridspec_kw = {'height_ratios':[1, 3]},figsize=(6,5))
            draw_scale = 1.0/(3.0*(max(geo_2D[:,1])-min(geo_2D[:,1])))   # The final drawing ends up in the 1:3 box, so the scaling of the lines/labels should match the eventual scaling of the drawing.
            if draw_scale > 1.0: draw_scale = 1.0

            # Add bonds
            linewidth = 2.0*draw_scale*5.0
            if linewidth > 2.0:
                linewidth = 2.0
            if linewidth < 1.0:
                linewidth = 1.0
            for count_j,j in enumerate(Adj_mat):
                for count_k,k in enumerate(j):
                    if count_k > count_j:
                        if k == 1:
                            ax1.add_line(matplotlib.lines.Line2D([geo_2D[count_j][0],geo_2D[count_k][0]],[geo_2D[count_j][1],geo_2D[count_k][1]],color=(0,0,0),linewidth=linewidth))

            # Add Highlights
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][0]][0],geo_2D[fit_dihedrals[dihedral_ind][1]][0]],[geo_2D[fit_dihedrals[dihedral_ind][0]][1],geo_2D[fit_dihedrals[dihedral_ind][1]][1]],color=(0,0.1,0.8),linewidth=linewidth))
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][1]][0],geo_2D[fit_dihedrals[dihedral_ind][2]][0]],[geo_2D[fit_dihedrals[dihedral_ind][1]][1],geo_2D[fit_dihedrals[dihedral_ind][2]][1]],color=(0,0.1,0.8),linewidth=linewidth))
            ax1.add_line(matplotlib.lines.Line2D([geo_2D[fit_dihedrals[dihedral_ind][2]][0],geo_2D[fit_dihedrals[dihedral_ind][3]][0]],[geo_2D[fit_dihedrals[dihedral_ind][2]][1],geo_2D[fit_dihedrals[dihedral_ind][3]][1]],color=(0,0.1,0.8),linewidth=linewidth))
                            
            # Add atom labels
            fontsize = 20*draw_scale*5.0
            if fontsize > 24:
                fontsize = 24
            if fontsize < 9:
                fontsize = 9            
            for count_j,j in enumerate(geo_2D):
                if count_j in fit_dihedrals[dihedral_ind]:
                    ax1.text(j[0], j[1], Elements[count_j], style='normal', color=(0.0,0.1,0.8), fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
                else:
                    ax1.text(j[0], j[1], Elements[count_j], style='normal', color=(0.0,0.0,0.0),fontweight='bold', fontsize=fontsize, ha='center', va='center',bbox={'facecolor':'white', 'edgecolor':'None','pad':1.0*draw_scale})
            ax1.axis('image')
            ax1.axis('off')
        
            # Plot the fit potential comparisons and total potential comparisons
            ax2.scatter(current_angles,current_QC-current_QC[0],s=80,color=(0.05,0.05,0.05,1.0))
            ax2.plot(current_angles,current_QC-current_QC[0],'--',linewidth=3,color=(0.05,0.05,0.05,1.0))
            ax2.scatter(current_angles,current_E_con-current_E_con[0],s=80,color=(0.0,0.1,0.8,1.0))
            ax2.plot(current_angles,current_E_con-current_E_con[0],'--',linewidth=3,color=(0.0,0.1,0.8,1.0))

            # Thermal potentials
            if fit_err_dict[i] >= 1.0:
                ax2.text(0.9, 0.8,'*', ha='center', fontsize=40, va='center', transform=ax2.transAxes)
                ax2.scatter(thermal_angles,thermal_potentials[d]-thermal_potentials[d][0],s=80,color=(0.7,0.1,1.0,1.0))
                ax2.plot(thermal_angles,thermal_potentials[d]-thermal_potentials[d][0],'--',linewidth=3,color=(0.7,0.1,1.0,1.0))
                
            # These print commands are here solely for diagnostic purposes.
            # ax2.scatter(current_angles,(current_fit-current_fit[0])+(current_E_con-current_E_con[0]),s=80,color=(0.0,0.8,0.1,1.0))
            # ax2.plot(current_angles,(current_fit-current_fit[0])+(current_E_con-current_E_con[0]),'--',linewidth=3,color=(0.0,0.8,0.1,1.0))

            # Format the plot            
#            ax2.set_title('Dihedral {}'.format(label),fontsize=12,y=1.08)
            ax2.set_ylabel(r'$\mathrm{Energy \, (kcal \cdot mol^{-1})}$',fontsize=20,labelpad=10)
            ax2.set_xlabel(r'$\mathrm{\phi \, (degrees)}$',fontsize=20,labelpad=10)
            ax2.set_xlim([-180, 180])
            ax2.tick_params(axis='both', which='major',labelsize=16,pad=10,direction='out',width=3,length=6)
            ax2.tick_params(axis='both', which='minor',labelsize=16,pad=10,direction='out',width=2,length=4)
            [k.set_linewidth(3) for k in ax2.spines.values()]
            plt.tight_layout()

            # Create save name, and if this is the final cycle save the name to the figure list variable 
            # for inclusion in the results summary.
            Name = '{}_Final.png'.format(savename)
            paths_to_dihedral_totpot_dict[d] = [save_folder+'/'+Name]
            plt.savefig(Name, bbox_inches=0,dpi=300)
            plt.savefig('{}_Final.pdf'.format(savename), bbox_inches=0,dpi=300)
            plt.close(fig)

        # Return to working directory
        os.chdir(current_dir)
        continue

    # Assemble paths/filenames for the final plots
    for d in list(paths_to_dihedral_totpot_dict.keys()):
        if d in list(paths_to_dihedral_fitpot_dict.keys()):
            paths_to_dihedral_fitpot_plots += paths_to_dihedral_fitpot_dict[d]
        paths_to_dihedral_totpot_plots += paths_to_dihedral_totpot_dict[d]

    # Remove temporary types 
    for i in remove_at_end:
        del FF_dict["dihedrals"][i]

    return paths_to_dihedral_fitpot_plots,paths_to_dihedral_totpot_plots

# This function reads in an atomtype label, generates its adjacency matrix, 
# and determines if it is a UA-hydrogen based on the first atomtype and its connections
# NOTE: this function only works with gens >= 2, for gens = 1 all hydrogens would
#       have to be treated as either UA or non-UA to be consistent with that definition.
def return_UA_H(type):

    # Remove -link label if it is present
    type = next( i for i in type.split('link-') if i != '' )    

    # Initialize lists/array 
    brackets = [ i for i in type if i == "[" or i == "]" ]
    atoms    = type.replace('['," ").replace(']'," ").split()
    adj_mat = np.zeros([len([ i for i in brackets if i == "[" ]),len([ i for i in brackets if i == "[" ])])

    # The basis of adj_mat corresponds to the atoms in "atoms", equivalently the "[" objects in brackets
    for count_i,i in enumerate(adj_mat):        

        # Indices are counted based on the number of left_brackets encountered (i.e. left_count's values)
        # Connections occur when bracket_count is equal to 1 and the current object in brackets is a left-bracket.
        # The algorithm works by incrementing bracket_count by +/- 1 for every bracket encountered.
        bracket_count = -1
        left_count = -1

        # Loop over the brackets list, bracket_count only gets incremented for objects in brackets satisfying left_count >= count_i 
        for count_j,j in enumerate(brackets):
            if j == "[": left_count += 1
            if left_count < count_i: continue
            if j == "[": bracket_count += 1; 
            if j == "]": bracket_count -= 1
            if left_count == count_i: continue
            if bracket_count == 1 and j == "[":
                adj_mat[count_i,left_count] = 1
                adj_mat[left_count,count_i] = 1

    # If the current atom is a hydrogen (i.e. atomtype "1") and it is connected to at least on carbon (i.e. atomtype "6") then it is a UA-type
    if atoms[0] == "1" and len([ i for count_i,i in enumerate(adj_mat[0]) if i == 1 and atoms[count_i] == "6" and sum(adj_mat[count_i]) == 4 ]) > 0:
        return 1
    # Else, it is not a UA-type
    else:
        return 0

# Adds parameters from the FF file(s) to the FF_dict.
def parse_FF_params(FF_files,qc_types,FF_dict=None,charges_AA=None,charges_UA=None):

    # Initialize lists/dictionaries if they were not supplied
    if FF_dict is None:
        FF_dict = { "masses":{}, "charges":{}, "vdw":{}, "bonds":{}, "angles":{}, "dihedrals":{}, "dihedrals_harmonic":{} }
    if charges_AA is None:
        charges_AA = {}
    if charges_UA is None:
        charges_UA = {}

    # Iterate over all supplied FF_files
    modes_from_FF = []
    for i in FF_files:
        with open(i,'r') as f:

            # Iterate over the lines of force-field file i
            for lines in f:
                fields = lines.split()

                # Skip empty lines
                if len(fields) == 0: continue

                # The remainder of the if/elif/else statements control the parsing of FF details
                # and the updating of the FF_dict dictionary and modes_from_FF list.
                if fields[0].lower() == "atom":   FF_dict["masses"][fields[1]] = float(fields[3])
                if fields[0].lower() == "charge": FF_dict["charges"][fields[1]] = float(fields[2])
                if fields[0].lower() == "bond":   
                    modes_from_FF += [(fields[1],fields[2])]
                    FF_dict["bonds"][(fields[1],fields[2])] = {}
                    if "dft" in qc_types:
                        FF_dict["bonds"][(fields[1],fields[2])]["DFT"] = [fields[3],float(fields[4]),float(fields[5])]
                    if "mp2" in qc_types:
                        FF_dict["bonds"][(fields[1],fields[2])]["MP2"] = [fields[3],float(fields[4]),float(fields[5])]
                if fields[0].lower() == "angle":
                    modes_from_FF += [(fields[1],fields[2],fields[3])]
                    FF_dict["angles"][(fields[1],fields[2],fields[3])] = {}
                    if "dft" in qc_types:
                        FF_dict["angles"][(fields[1],fields[2],fields[3])]["DFT"] = [fields[4],float(fields[5]),float(fields[6])]
                    if "mp2" in qc_types:
                        FF_dict["angles"][(fields[1],fields[2],fields[3])]["MP2"] = [fields[4],float(fields[5]),float(fields[6])]
                if fields[0].lower() in ["dihedral","torsion"]: 
                    modes_from_FF += [(fields[1],fields[2],fields[3],fields[4],fields[5])]
                    if fields[5] == "opls":       
                        FF_dict["dihedrals"][(fields[1],fields[2],fields[3],fields[4])] = {}                    
                        if "dft" in qc_types:
                            FF_dict["dihedrals"][(fields[1],fields[2],fields[3],fields[4])]["DFT"] = [fields[5]] + [ float(i) for i in fields[6:10] ]
                        if "mp2" in qc_types:
                            FF_dict["dihedrals"][(fields[1],fields[2],fields[3],fields[4])]["MP2"] = [fields[5]] + [ float(i) for i in fields[6:10] ]
                    elif fields[5] == "harmonic":
                        FF_dict["dihedrals_harmonic"][(fields[1],fields[2],fields[3],fields[4])] = {}                    
                        if "dft" in qc_types:
                            FF_dict["dihedrals_harmonic"][(fields[1],fields[2],fields[3],fields[4])]["DFT"] = [fields[5]] + [ float(fields[6]),int(float(fields[7])),int(float(fields[8])) ] 
                        if "mp2" in qc_types:
                            FF_dict["dihedrals_harmonic"][(fields[1],fields[2],fields[3],fields[4])]["MP2"] = [fields[5]] + [ float(fields[6]),int(float(fields[7])),int(float(fields[8])) ] 
                    elif fields[5] == "quadratic":
                        FF_dict["dihedrals_harmonic"][(fields[1],fields[2],fields[3],fields[4])] = {}                    
                        if "dft" in qc_types:
                            FF_dict["dihedrals_harmonic"][(fields[1],fields[2],fields[3],fields[4])]["DFT"] = [fields[5]] + [ float(fields[6]),float(fields[7]) ]
                        if "mp2" in qc_types:
                            FF_dict["dihedrals_harmonic"][(fields[1],fields[2],fields[3],fields[4])]["MP2"] = [fields[5]] + [ float(fields[6]),float(fields[7]) ]

                # Charges are updated in two places since some parts of the code rely on 
                # the charges_* lists in addition to the FF_dict.
                if fields[0].lower() == "charge":
                    if "-UA" in fields[1]:
                        charges_UA[fields[1]] = float(fields[2])
                    else:
                        charges_AA[fields[1]] = float(fields[2])                        
                if fields[0].lower() == "vdw":    
                    FF_dict["vdw"][(fields[1],fields[2])] = [fields[3],float(fields[4]),float(fields[5])]
                    FF_dict["vdw"][(fields[2],fields[1])] = [fields[3],float(fields[4]),float(fields[5])]

    return FF_dict,modes_from_FF,charges_AA,charges_UA

# Wrapper function for writing the lammps data file which is read by the .in.init file during run initialization. This function is 
# nearly identical to the function of the same name within the gen_md_for_vdw.py script, except the qc option has been added to be compatible
# with the FF_dict structure and the Molecule variable has been removed because only one molecule is present in each simulation.
def Write_MD_data(Filename,Atom_types,Sim_Box,Elements,Geometry,Bonds,Bond_types,Bond_params,Angles,Angle_types,Angle_params,Dihedrals,Dihedral_types,Dihedral_params,One_fives,\
                      Charges,VDW_params,Masses,Scanned_Dihedrals,qc,one_five_opt=0,tables=[]):
    
    # Write an xyz for easy viewing
    with open(Filename+'.xyz','w') as f:
        f.write('{}\n\n'.format(len(Geometry)))
        for count_i,i in enumerate(Geometry):
            f.write('{:20s} {:< 20.6f} {:< 20.6f} {:< 20.6f}\n'.format(Elements[count_i],i[0],i[1],i[2]))

    # Create type dictionaries (needed to convert each atom,bond,angle, and dihedral type to consecutive numbers as per LAMMPS convention)
    # Note: LAMMPS orders atomtypes, bonds, angles, dihedrals, etc as integer types. Each of the following dictionaries holds the mapping between
    #       the true type (held in the various types lists) and the lammps type_id, obtained by enumerated iteration over the respective set(types).
    Atom_type_dict = {}
    for count_i,i in enumerate(sorted(set(Atom_types))):
        for j in Atom_types:
            if i == j:
                Atom_type_dict[i]=count_i+1
            if i in list(Atom_type_dict.keys()):
                break
    Bond_type_dict = {}
    for count_i,i in enumerate(sorted(set(Bond_types))):
        for j in Bond_types:
            if i == j:
                Bond_type_dict[i]=count_i+1
            if i in list(Bond_type_dict.keys()):
                break
    
    # Add 1-5 bonds as the last type (used for turning off 1-5 non-bonded interactions
    if one_five_opt == 1:
        Bond_type_dict["one_fives"] = count_i+2   

    Angle_type_dict = {}
    for count_i,i in enumerate(sorted(set(Angle_types))):
        for j in Angle_types:
            if i == j:
                Angle_type_dict[i]=count_i+1
            if i in list(Angle_type_dict.keys()):
                break
    Dihedral_type_dict = {}
    for count_i,i in enumerate(sorted(set(Dihedral_types))):
        for j in Dihedral_types:
            if i == j:
                Dihedral_type_dict[i]=count_i+1
            if i in list(Dihedral_type_dict.keys()):
                break

    # Write the data file
    with open(Filename+'.data','w') as f:
        
        # Write system properties
        f.write("LAMMPS data file via vdw_self_gen.py, on {}\n\n".format(datetime.datetime.now()))

        f.write("{} atoms\n".format(len(Elements)))
        f.write("{} atom types\n".format(len(set(Atom_types))))
        if len(Bonds) > 0:
            if one_five_opt == 1:
                f.write("{} bonds\n".format(len(Bonds)+len(One_fives)))
                f.write("{} bond types\n".format(len(set(Bond_types))+1))
            else:
                f.write("{} bonds\n".format(len(Bonds)))
                f.write("{} bond types\n".format(len(set(Bond_types))))
        if len(Angles) > 0:
            f.write("{} angles\n".format(len(Angles)))
            f.write("{} angle types\n".format(len(set(Angle_types))))
        if len(Dihedrals) > 0:
            f.write("{} dihedrals\n".format(len(Dihedrals)))
            f.write("{} dihedral types\n\n".format(len(set(Dihedral_types))))

        # Write box dimensions
        f.write("{:< 20.16f} {:< 20.16f} xlo xhi\n".format(Sim_Box[0],Sim_Box[1]))
        f.write("{:< 20.16f} {:< 20.16f} ylo yhi\n".format(Sim_Box[2],Sim_Box[3]))
        f.write("{:< 20.16f} {:< 20.16f} zlo zhi\n\n".format(Sim_Box[4],Sim_Box[5]))

        # Write Masses
        f.write("Masses\n\n")
        for count_i,i in enumerate(sorted(set(Atom_types))):
            for j in set(Atom_types):
                if Atom_type_dict[j] == count_i+1:
                    f.write("{} {:< 8.6f}\n".format(count_i+1,Masses[str(j)])) # count_i+1 bc of LAMMPS 1-indexing
        f.write("\n")

        # Write Bond Coeffs
        f.write("Bond Coeffs\n\n")
        for count_i,i in enumerate(set(Bond_types)):
            for j in set(Bond_types):
                if Bond_type_dict[j] == count_i+1:
                    f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+1,Bond_params[j][qc][1],Bond_params[j][qc][2])) # count_i+1 bc of LAMMPS 1-indexing

        # Add one_five ghost bonds
        if one_five_opt == 1:
            f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+2,0.0,2.0)) # dummy bond for 1-5 interactions
        f.write("\n")

        # Write Angle Coeffs
        f.write("Angle Coeffs\n\n")
        for count_i,i in enumerate(set(Angle_types)):
            for j in set(Angle_types):
                if Angle_type_dict[j] == count_i+1:
                    f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+1,Angle_params[j][qc][1],Angle_params[j][qc][2])) # count_i+1 bc of LAMMPS 1-indexing
        f.write("\n")

        # Write Atoms
        f.write("Atoms\n\n")
        for count_i,i in enumerate(Atom_types):
            f.write("{:<8d} {:< 4d} {:< 4d} {:< 20.16f} {:< 20.16f} {:< 20.16f} {:< 20.16f} {:d} {:d} {:d}\n"\
            .format(count_i+1,0,Atom_type_dict[i],Charges[count_i],Geometry[count_i,0],Geometry[count_i,1],Geometry[count_i,2],0,0,0))

        # Write Bonds
        bond_count = 1
        if len(Bonds) > 0:
            f.write("\nBonds\n\n")
            for count_i,i in enumerate(Bonds):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d}\n".format(bond_count,Bond_type_dict[Bond_types[count_i]],i[0]+1,i[1]+1))
                bond_count += 1

        # Write 1-5 ghost bonds
        if one_five_opt == 1: 
            for i in One_fives:
                f.write("{:<8d} {:<8d} {:<8d} {:<8d}\n".format(bond_count,Bond_type_dict["one_fives"],i[0]+1,i[-1]+1))
                bond_count += 1

        # Write Angles
        if len(Angles) > 0:
            f.write("\nAngles\n\n")
            for count_i,i in enumerate(Angles):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Angle_type_dict[Angle_types[count_i]],i[0]+1,i[1]+1,i[2]+1))

        # Write Dihedrals
        if len(Dihedrals) > 0: 
            f.write("\nDihedrals\n\n")
            for count_i,i in enumerate(Dihedrals):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Dihedral_type_dict[Dihedral_types[count_i]],i[0]+1,i[1]+1,i[2]+1,i[3]+1))

    # Write the settings file
    with open(Filename+'.in.settings','w') as f:

        # Write non-bonded interactions (the complicated form of this loop is owed
        # to desire to form a nicely sorted list in terms of the lammps atom types
        # Note: Atom_type_dict was initialize according to sorted(set(Atom_types) 
        #       so iterating over this list (twice) orders the pairs, too.
        f.write("     {}\n".format("# Non-bonded interactions (pair-wise)"))
        for count_i,i in enumerate(sorted(set(Atom_types))):     
            for count_j,j in enumerate(sorted(set(Atom_types))): 

                # Skip duplicates
                if count_j < count_i:
                    continue

                # Conform to LAMMPS i <= j formatting
                if Atom_type_dict[i] <= Atom_type_dict[j]:
                    f.write("     {:20s} {:<10d} {:<10d} ".format("pair_coeff",Atom_type_dict[i],Atom_type_dict[j]))
                else:
                    f.write("     {:20s} {:<10d} {:<10d} ".format("pair_coeff",Atom_type_dict[j],Atom_type_dict[i]))

                # Determine key (ordered by initialize_VDW function such that i > j)
                if i > j:
                    key = (i,j)
                else:
                    key = (j,i)

                # Write the parameters (note: the internal type (e.g., lj) is converted into the corresponding lammps type (e.g., lj/cut/coul/long))
                for k in VDW_params[key]:
                    if k == "lj": k = "lj/cut/coul/cut"
                    if type(k) is str:
                        f.write("{:20s} ".format(k))
                    if type(k) is float:
                        f.write("{:< 20.6f} ".format(k))
                f.write("\n")                        

        # Write stretching interactions
        # Note: Bond_type_dict was initialized by looping over sorted(set(Bond_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Stretching interactions"))
        for i in sorted(set(Bond_types)):
            f.write("     {:20s} {:<10d} ".format("bond_coeff",Bond_type_dict[i]))
            for j in Bond_params[i][qc][1:]:
                if type(j) is str:
                    f.write("{:20s} ".format(j))
                if type(j) is float:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is np.double:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is int:
                    f.write("{:< 20d} ".format(j))
            f.write("\n")

        # Write 1-5 ghost bond parameter
        if one_five_opt == 1:
            f.write("     {:20s} {:<10d} {:< 20.6f} {:< 20.6f}\n".format("bond_coeff",Bond_type_dict["one_fives"],0.0,2.0))

        # Write bending interactions
        # Note: Angle_type_dict was initialized by looping over sorted(set(Angle_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Bending interactions"))
        for i in sorted(set(Angle_types)):
            f.write("     {:20s} {:<10d} ".format("angle_coeff",Angle_type_dict[i]))
            for j in Angle_params[i][qc][1:]:
                if type(j) is str:
                    f.write("{:20s} ".format(j))
                if type(j) is float:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is np.double:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is int:
                    f.write("{:< 20d} ".format(j))
            f.write("\n")

        # Write dihedral interactions
        # Note: Dihedral_type_dict was initialized by looping over sorted(set(Dihedral_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Dihedral interactions"))
        for i in sorted(set(Dihedral_types)):

            # The potential fo the scanned dihedrals are set to zero
            if i in tables:
                f.write("     {:20s} {:<10d} {:20s} {:20s} {:20s}".format("dihedral_coeff",Dihedral_type_dict[i],"table",tables[i][0],tables[i][1]))
            elif i in [ j[0] for j in Scanned_Dihedrals ]:
                f.write("     {:20s} {:<10d} opls 0.0 0.0 0.0 0.0".format("dihedral_coeff",Dihedral_type_dict[i]))
            else:
                f.write("     {:20s} {:<10d} ".format("dihedral_coeff",Dihedral_type_dict[i]))
                for j in Dihedral_params[i][qc]:
                    if type(j) is str:
                        f.write("{:20s} ".format(j))
                    if type(j) is float:
                        f.write("{:< 20.6f} ".format(j))
                    if type(j) is np.double:
                        f.write("{:< 20.6f} ".format(j))
                    if type(j) is int:
                        f.write("{:< 20d} ".format(j))
            f.write("\n")

    return Atom_type_dict

# Wrapper function for writing the lammps data file which is read by the .in.init file during run initialization. This function is 
# nearly identical to the function of the same name within the gen_md_for_vdw.py script, except the qc option has been added to be compatible
# with the FF_dict structure and the Molecule variable has been removed because only one molecule is present in each simulation.
def Write_OPLSFIT_MD_data(Filename,Atom_types,Sim_Box,Elements,Geometry,Bonds,Bond_types,Bond_params,Angles,Angle_types,Angle_params,Dihedrals,Dihedral_types,Dihedral_params,Dihedral_harmonic_params,One_fives,\
                          Charges,VDW_params,Masses,Scanned_Dihedrals,qc,one_five_opt=0,coul_scale=1.0,tables=[]):
    
    # Write an xyz for easy viewing
    with open(Filename+'.xyz','w') as f:
        f.write('{}\n\n'.format(len(Geometry)))
        for count_i,i in enumerate(Geometry):
            f.write('{:20s} {:< 20.6f} {:< 20.6f} {:< 20.6f}\n'.format(Elements[count_i],i[0],i[1],i[2]))

    # Create type dictionaries (needed to convert each atom,bond,angle, and dihedral type to consecutive numbers as per LAMMPS convention)
    # Note: LAMMPS orders atomtypes, bonds, angles, dihedrals, etc as integer types. Each of the following dictionaries holds the mapping between
    #       the true type (held in the various types lists) and the lammps type_id, obtained by enumerated iteration over the respective set(types).
    Atom_type_dict = {}
    for count_i,i in enumerate(sorted(set(Atom_types))):
        for j in Atom_types:
            if i == j:
                Atom_type_dict[i]=count_i+1
            if i in list(Atom_type_dict.keys()):
                break
    Bond_type_dict = {}
    for count_i,i in enumerate(sorted(set(Bond_types))):
        for j in Bond_types:
            if i == j:
                Bond_type_dict[i]=count_i+1
            if i in list(Bond_type_dict.keys()):
                break
    
    # Add 1-5 bonds as the last type (used for turning off 1-5 non-bonded interactions
    if one_five_opt == 1:
        Bond_type_dict["one_fives"] = count_i+2   

    Angle_type_dict = {}
    for count_i,i in enumerate(sorted(set(Angle_types))):
        for j in Angle_types:
            if i == j:
                Angle_type_dict[i]=count_i+1
            if i in list(Angle_type_dict.keys()):
                break
    Dihedral_type_dict = {}
    for count_i,i in enumerate(sorted(set(Dihedral_types))):
        for j in Dihedral_types:
            if i == j:
                Dihedral_type_dict[i]=count_i+1
            if i in list(Dihedral_type_dict.keys()):
                break

    # Add the special fit type as last LAMMPS dihedral type
    if len(Scanned_Dihedrals) > 0:
        Dihedral_type_dict["fit"]=len(list(Dihedral_type_dict.keys()))+1

    # Write the data file
    with open(Filename+'.data','w') as f:
        
        # Write system properties
        f.write("LAMMPS data file via vdw_self_gen.py, on {}\n\n".format(datetime.datetime.now()))

        f.write("{} atoms\n".format(len(Elements)))
        f.write("{} atom types\n".format(len(set(Atom_types))))
        if len(Bonds) > 0:
            if one_five_opt == 1:
                f.write("{} bonds\n".format(len(Bonds)+len(One_fives)))
                f.write("{} bond types\n".format(len(set(Bond_types))+1))
            else:
                f.write("{} bonds\n".format(len(Bonds)))
                f.write("{} bond types\n".format(len(set(Bond_types))))
        if len(Angles) > 0:
            f.write("{} angles\n".format(len(Angles)))
            f.write("{} angle types\n".format(len(set(Angle_types))))
        if len(Dihedrals) > 0:
            f.write("{} dihedrals\n".format(len(Dihedrals)))
            f.write("{} dihedral types\n\n".format(len(list(Dihedral_type_dict.keys()))))

        # Write box dimensions
        f.write("{:< 20.16f} {:< 20.16f} xlo xhi\n".format(Sim_Box[0],Sim_Box[1]))
        f.write("{:< 20.16f} {:< 20.16f} ylo yhi\n".format(Sim_Box[2],Sim_Box[3]))
        f.write("{:< 20.16f} {:< 20.16f} zlo zhi\n\n".format(Sim_Box[4],Sim_Box[5]))

        # Write Masses
        f.write("Masses\n\n")
        for count_i,i in enumerate(sorted(set(Atom_types))):
            for j in set(Atom_types):
                if Atom_type_dict[j] == count_i+1:
                    f.write("{} {:< 8.6f}\n".format(count_i+1,Masses[str(j)])) # count_i+1 bc of LAMMPS 1-indexing
        f.write("\n")

        # Write Bond Coeffs
        f.write("Bond Coeffs\n\n")
        for count_i,i in enumerate(set(Bond_types)):
            for j in set(Bond_types):
                if Bond_type_dict[j] == count_i+1:
                    f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+1,Bond_params[j][qc][1],Bond_params[j][qc][2])) # count_i+1 bc of LAMMPS 1-indexing

        # Add one_five ghost bonds
        if one_five_opt == 1:
            f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+2,0.0,2.0)) # dummy bond for 1-5 interactions
        f.write("\n")

        # Write Angle Coeffs
        f.write("Angle Coeffs\n\n")
        for count_i,i in enumerate(set(Angle_types)):
            for j in set(Angle_types):
                if Angle_type_dict[j] == count_i+1:
                    f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+1,Angle_params[j][qc][1],Angle_params[j][qc][2])) # count_i+1 bc of LAMMPS 1-indexing
        f.write("\n")

        # Write Atoms
        f.write("Atoms\n\n")
        for count_i,i in enumerate(Atom_types):
            f.write("{:<8d} {:< 4d} {:< 4d} {:< 20.16f} {:< 20.16f} {:< 20.16f} {:< 20.16f} {:d} {:d} {:d}\n"\
            .format(count_i+1,0,Atom_type_dict[i],Charges[count_i]*coul_scale,Geometry[count_i,0],Geometry[count_i,1],Geometry[count_i,2],0,0,0))

        # Write Bonds
        bond_count = 1
        if len(Bonds) > 0:
            f.write("\nBonds\n\n")
            for count_i,i in enumerate(Bonds):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d}\n".format(bond_count,Bond_type_dict[Bond_types[count_i]],i[0]+1,i[1]+1))
                bond_count += 1

        # Write 1-5 ghost bonds
        if one_five_opt == 1: 
            for i in One_fives:
                f.write("{:<8d} {:<8d} {:<8d} {:<8d}\n".format(bond_count,Bond_type_dict["one_fives"],i[0]+1,i[-1]+1))
                bond_count += 1

        # Write Angles
        if len(Angles) > 0:
            f.write("\nAngles\n\n")
            for count_i,i in enumerate(Angles):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Angle_type_dict[Angle_types[count_i]],i[0]+1,i[1]+1,i[2]+1))

        # Write Dihedrals
        if len(Dihedrals) > 0: 
            f.write("\nDihedrals\n\n")
            for count_i,i in enumerate(Dihedrals):
                if i in Scanned_Dihedrals:
                    f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Dihedral_type_dict["fit"],i[0]+1,i[1]+1,i[2]+1,i[3]+1))
                else:
                    f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Dihedral_type_dict[Dihedral_types[count_i]],i[0]+1,i[1]+1,i[2]+1,i[3]+1))

    # Write the settings file
    with open(Filename+'.in.settings','w') as f:

        # Write non-bonded interactions (the complicated form of this loop is owed
        # to desire to form a nicely sorted list in terms of the lammps atom types
        # Note: Atom_type_dict was initialize according to sorted(set(Atom_types) 
        #       so iterating over this list (twice) orders the pairs, too.
        f.write("     {}\n".format("# Non-bonded interactions (pair-wise)"))
        for count_i,i in enumerate(sorted(set(Atom_types))):     
            for count_j,j in enumerate(sorted(set(Atom_types))): 

                # Skip duplicates
                if count_j < count_i:
                    continue

                # Conform to LAMMPS i <= j formatting
                if Atom_type_dict[i] <= Atom_type_dict[j]:
                    f.write("     {:20s} {:<10d} {:<10d} ".format("pair_coeff",Atom_type_dict[i],Atom_type_dict[j]))
                else:
                    f.write("     {:20s} {:<10d} {:<10d} ".format("pair_coeff",Atom_type_dict[j],Atom_type_dict[i]))

                # Determine key (ordered by initialize_VDW function such that i > j)
                if i > j:
                    key = (i,j)
                else:
                    key = (j,i)

                # Write the parameters (note: the internal type (e.g., lj) is converted into the corresponding lammps type (e.g., lj/cut/coul/long))
                for k in VDW_params[key]:
                    if k == "lj": k = "lj/cut/coul/cut"
                    if type(k) is str:
                        f.write("{:20s} ".format(k))
                    if type(k) is float:
                        f.write("{:< 20.6f} ".format(k))
                f.write("\n")                        

        # Write stretching interactions
        # Note: Bond_type_dict was initialized by looping over sorted(set(Bond_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Stretching interactions"))
        for i in sorted(set(Bond_types)):
            f.write("     {:20s} {:<10d} ".format("bond_coeff",Bond_type_dict[i]))
            for j in Bond_params[i][qc][1:]:
                if type(j) is str:
                    f.write("{:20s} ".format(j))
                if type(j) is float:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is np.double:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is int:
                    f.write("{:< 20d} ".format(j))
            f.write("\n")

        # Write 1-5 ghost bond parameter
        if one_five_opt == 1:
            f.write("     {:20s} {:<10d} {:< 20.6f} {:< 20.6f}\n".format("bond_coeff",Bond_type_dict["one_fives"],0.0,2.0))

        # Write bending interactions
        # Note: Angle_type_dict was initialized by looping over sorted(set(Angle_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Bending interactions"))
        for i in sorted(set(Angle_types)):
            f.write("     {:20s} {:<10d} ".format("angle_coeff",Angle_type_dict[i]))
            for j in Angle_params[i][qc][1:]:
                if type(j) is str:
                    f.write("{:20s} ".format(j))
                if type(j) is float:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is np.double:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is int:
                    f.write("{:< 20d} ".format(j))
            f.write("\n")

        # Write dihedral interactions
        # Note: Dihedral_type_dict was initialized by looping over sorted(set(Dihedral_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Dihedral interactions"))
        for i in sorted(set(Dihedral_types)):
            
                # The potential for the scanned dihedrals are set to zero
                if i in tables:
                    f.write("     {:20s} {:<10d} {:20s} {:20s} {:20s}".format("dihedral_coeff",Dihedral_type_dict[i],"table",tables[i][0],tables[i][1]))
                elif i in Scanned_Dihedrals:
                    f.write("     {:20s} {:<10d} opls 0.0 0.0 0.0 0.0".format("dihedral_coeff",Dihedral_type_dict["fit"]))
                else:

                    # If it is an opls type then it needs to reference the Dihedral_params dictionary
                    if i[4] == "opls":
                        f.write("     {:20s} {:<10d} ".format("dihedral_coeff",Dihedral_type_dict[i]))

                        # This check is performed in case a V0 normalization constant is included in the dihedral definition
                        if len(Dihedral_params[i[:4]][qc]) == 6:
                            tmp = Dihedral_params[i[:4]][qc][0] + Dihedral_params[i[:4]][qc][2:]
                        else:
                            tmp = Dihedral_params[i[:4]][qc]
                        for j in tmp:
                            if type(j) is str:
                                f.write("{:20s} ".format(j))
                            if type(j) is float:
                                f.write("{:< 20.6f} ".format(j))
                            if type(j) is np.double:
                                f.write("{:< 20.6f} ".format(j))
                            if type(j) is int:
                                f.write("{:< 20d} ".format(j))

                    # If it is a harmomnic type then it needs to reference the Dihedral_harmonic_params dictionary
                    elif i[4] == "harmonic":
                        f.write("     {:20s} {:<10d} ".format("dihedral_coeff",Dihedral_type_dict[i]))
                        for j in Dihedral_harmonic_params[i[:4]][qc]:
                            if type(j) is str:
                                f.write("{:20s} ".format(j))
                            if type(j) is float:
                                f.write("{:< 20.6f} ".format(j))
                            if type(j) is np.double:
                                f.write("{:< 20.6f} ".format(j))
                            if type(j) is int:
                                f.write("{:< 20d} ".format(j))
                    else:
                        print("ERROR: Dihedral parameters for {} could not be found. Exiting...".format(i))
                        quit()
                f.write("\n")
            
        # Write the fit type last
        if len(Scanned_Dihedrals) > 0:
            f.write("     {:20s} {:<10d} opls 0.0 0.0 0.0 0.0\n".format("dihedral_coeff",Dihedral_type_dict["fit"]))

    return Atom_type_dict


# Description: A wrapper for the write commands for generating the lammps input file
def Write_OPLSFIT_input(Filename,frequency,Onefourscale,time,Pair_styles,Dihedral_styles,con,Temp,opt_steps,init_angle):

    with open(Filename+'.in.init','w') as f:
        f.write("# lammps input file for polymer simulation with dilute ions\n\n"+\
                "# VARIABLES\n"+\
                "variable        data_name       index   {}\n".format(Filename+'.data')+\
                "variable        settings_name   index   {}\n".format(Filename+'.in.settings')+\
                "variable        coords_freq     index   {}\n".format(frequency)+\
                "variable        vseed           index   {: <6d}\n".format(int(random.random()*100000))+\
                
                "#===========================================================\n"+\
                "# GENERAL PROCEDURES\n"+\
                "#===========================================================\n"+\
                "units		real	# g/mol, angstroms, fs, kcal/mol, K, atm, charge*angstrom\n"+\
                "dimension	3	# 3 dimensional simulation\n"+\
                "newton		off	# use Newton's 3rd law\n"+\
                "boundary	f f f	# periodic boundary conditions \n"+\
                "atom_style	full    # molecular + charge\n\n"+\

                "#===========================================================\n"+\
                "# FORCE FIELD DEFINITION\n"+\
                "#===========================================================\n"+\
                'special_bonds  lj/coul  0.0 0.0 0.0     # NO 1-4 LJ/Coul interactions\n'+\
                'pair_style     hybrid {} # outer_LJ outer_Coul (cutoff values, see LAMMPS Doc)\n'.format(Pair_styles)+\
                'bond_style     harmonic             # parameters needed: k_bond, r0\n'+\
                'angle_style    harmonic             # parameters needed: k_theta, theta0\n')
        if len(Dihedral_styles) >= 1:
            f.write('dihedral_style hybrid {}            # parameters needed: V1, V2, V3, V4\n'.format(' '.join(Dihedral_styles)))                    
        f.write('pair_modify    shift yes mix arithmetic table 0       # using Lorenz-Berthelot mixing rules\n\n'+\

                "#===========================================================\n"+\
                "# SETUP SIMULATIONS\n"+\
                "#===========================================================\n\n"+\
                "# READ IN COEFFICIENTS/COORDINATES/TOPOLOGY\n"+\
                'read_data ${data_name}\n'+\
                'include ${settings_name}\n\n'+\

                "# SET RUN PARAMETERS\n"+\
                "timestep	1.0		# fs\n"+\
                "run_style	verlet 		# Velocity-Verlet integrator\n"+\
                "neigh_modify every 1 delay 0 check no # More relaxed rebuild criteria can be used\n\n"+\
        
                "# Define Theromstyle variables for print\n"+\
                "thermo_style    custom step temp vol density etotal pe ebond eangle edihed ecoul elong evdwl\n"+\
                "variable        my_pe   equal   pe\n"+\
                "variable        my_ebon equal   ebond\n"+\
                "variable        my_eang equal   eangle\n"+\
                "variable        my_edih equal   edihed\n"+\
                "variable        my_evdw equal   evdwl\n"+\
                "variable        my_ecol equal   ecoul\n\n")

        f.write("#===========================================================\n"+\
                "# RUN MINIMIZATION\n"+\
                "#===========================================================\n\n")
        
        R_seed = int(random.random()*1000000.0)        

        # Constraints are gradually introduced during an annealing process
        # NOTE: LAMMPS atoms are 1-indexed, so the constraints get incremented by 1
        f.write("# minimize molecule energy with restraints\n"+\
                "dump           minimization all atom ${coords_freq}"+" {}.lammpstrj\n".format(Filename)+\
                "dump_modify    minimization scale no\n"+\
                "dump_modify    minimization sort  id\n"+\
                "velocity all create {} {} mom yes rot yes dist gaussian\n".format(Temp,R_seed)+\
                "fix NVE all nve\n"+\
                "fix TFIX all langevin {} 0.0 100 {} zero yes\n".format(Temp,R_seed)+\
                "fix CON all restrain dihedral {} {} {} {} 0.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,init_angle-180.0)+\
                "fix_modify CON energy yes\n"+\
                "run {}\n".format(time)+\
                "fix TFIX all langevin 0.0 0.0 100 {} zero yes\n".format(R_seed)+\
                "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,init_angle-180.0)+\
                "fix_modify CON energy yes\n"+\
                "run 1000\n"+\
                "undump minimization\n")

        for count_j,j in enumerate(np.arange(init_angle,init_angle+360,360.0/opt_steps)):
            f.write("\n# PERFORM A PRELIMINARY POSITIVE ROTATION TO REMOVE INITIAL GUESS DEPENDENCE\n"+\
                    "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,j-180.0)+\
                    "fix_modify CON energy yes\n"+\
                    "minimize 1e-6 1e-9 1000 100000\n"+\
                    "unfix CON\n")

        counter = 0
        for count_j,j in enumerate(np.arange(init_angle,init_angle+360,360.0/opt_steps)):
            if int(round(j)) % 10 == 0:
                f.write("\n# APPLY CONTRAINT AND PERFORM POSITIVE SCAN\n"+\
                        "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,j-180.0)+\
                        "fix_modify CON energy yes\n"+\
                        "dump               minimization all atom ${coords_freq}"+" {}.plus.lammpstrj\n".format(counter)+\
                        "dump_modify        minimization scale no\n"+\
                        "dump_modify        minimization sort  id\n"+\
                        "minimize 1e-6 1e-9 1000 100000\n"+\
                        "undump minimization\n"+\
                        "unfix CON\n")
                counter += 1
            else:
                f.write("\n# APPLY CONTRAINT AND PERFORM POSITIVE SCAN\n"+\
                        "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,j-180.0)+\
                        "fix_modify CON energy yes\n"+\
                        "minimize 1e-6 1e-9 1000 100000\n"+\
                        "unfix CON\n")

        for count_j,j in enumerate(np.arange(init_angle,init_angle+360,360.0/opt_steps)[::-1]):
            f.write("\n# PERFORM A PRELIMINARY ROTATION TO REMOVE INITIAL GUESS DEPENDENCE\n"+\
                    "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,j-180.0)+\
                    "fix_modify CON energy yes\n"+\
                    "minimize 1e-6 1e-9 1000 100000\n"+\
                    "unfix CON\n")

        counter -= 1
        for count_j,j in enumerate(np.arange(init_angle,init_angle+360,360.0/opt_steps)[::-1]):
            if int(round(j)) % 10 == 0:
                f.write("\n# APPLY CONTRAINT AND PERFORM NEGATIVE SCAN\n"+\
                        "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,j-180.0)+\
                        "fix_modify CON energy yes\n"+\
                        "dump               minimization all atom ${coords_freq}"+" {}.minus.lammpstrj\n".format(counter)+\
                        "dump_modify        minimization scale no\n"+\
                        "dump_modify        minimization sort  id\n"+\
                        "minimize 1e-6 1e-9 1000 100000\n"+\
                        "undump minimization\n"+\
                        "unfix CON\n")
                counter -= 1
            else:
                f.write("\n# APPLY CONTRAINT AND PERFORM NEGATIVE SCAN\n"+\
                        "fix CON all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n".format(con[0]+1,con[1]+1,con[2]+1,con[3]+1,j-180.0)+\
                        "fix_modify CON energy yes\n"+\
                        "minimize 1e-6 1e-9 1000 100000\n"+\
                        "unfix CON\n")

    f.close()
    return

# Description: A wrapper for the write commands for generating the lammps input file based on constrained minimizations
def Write_QCMATCH_input(Filename,frequency,Onefourscale,Pair_styles,Dihedral_styles,cons,other_constraints,mode='forward',start_ind=0,k_sec=10.0):

    with open(Filename+'.in.init','w') as f:
        f.write("# lammps input file for polymer simulation with dilute ions\n\n"+\
                "# VARIABLES\n"+\
                "variable        data_name       index   {}\n".format(Filename+'.data')+\
                "variable        settings_name   index   {}\n".format(Filename+'.in.settings')+\
                "variable        coords_freq     index   {}\n".format(frequency)+\
                "variable        vseed           index   {: <6d}\n".format(int(random.random()*100000))+\
                
                "#===========================================================\n"+\
                "# GENERAL PROCEDURES\n"+\
                "#===========================================================\n"+\
                "units		real	# g/mol, angstroms, fs, kcal/mol, K, atm, charge*angstrom\n"+\
                "dimension	3	# 3 dimensional simulation\n"+\
                "newton		off	# use Newton's 3rd law\n"+\
                "boundary	f f f	# fixed boundary conditions \n"+\
                "atom_style	full    # molecular + charge\n\n"+\

                "#===========================================================\n"+\
                "# FORCE FIELD DEFINITION\n"+\
                "#===========================================================\n"+\
                'special_bonds  lj/coul  0.0 0.0 0.0     # NO 1-4 LJ/Coul interactions\n'+\
                'pair_style     hybrid {} # outer_LJ outer_Coul (cutoff values, see LAMMPS Doc)\n'.format(Pair_styles)+\
                'pair_modify    shift yes\n'+\
                'bond_style     harmonic             # parameters needed: k_bond, r0\n'+\
                'angle_style    harmonic             # parameters needed: k_theta, theta0\n')
        if len(Dihedral_styles) >= 1:
            f.write('dihedral_style hybrid {}            # parameters depend on the kinds of dihedrals that are present\n'.format(' '.join(Dihedral_styles)))                    
        f.write('pair_modify    shift yes mix arithmetic table 0       # using Lorenz-Berthelot mixing rules\n\n'+\

                "#===========================================================\n"+\
                "# SETUP SIMULATIONS\n"+\
                "#===========================================================\n\n"+\
                "# READ IN COEFFICIENTS/COORDINATES/TOPOLOGY\n"+\
                'read_data ${data_name}\n'+\
                'include ${settings_name}\n\n'+\

                "# SET RUN PARAMETERS\n"+\
                "timestep	1.0		# fs\n"+\
                "run_style	verlet 		# Velocity-Verlet integrator\n"+\
                "neigh_modify every 1 delay 0 check no # More relaxed rebuild criteria can be used\n\n"+\
        
                "# Define Theromstyle variables for print\n"+\
                "thermo_style    custom step temp vol density etotal pe ebond eangle edihed ecoul elong evdwl\n"+\
                "variable        my_pe   equal   pe\n"+\
                "variable        my_ebon equal   ebond\n"+\
                "variable        my_eang equal   eangle\n"+\
                "variable        my_edih equal   edihed\n"+\
                "variable        my_evdw equal   evdwl\n"+\
                "variable        my_ecol equal   ecoul\n\n")

        f.write("#===========================================================\n"+\
                "# RUN MINIMIZATION\n"+\
                "#===========================================================\n\n")

        f.write('# Set minimization style (hftn seems to be the most robust)\n')
        f.write('min_style hftn\n\n')

        # Initialize the list of constraint indices. This allows the function to start at any arbitrary configuration
        max_evals = 1000
        con_ind = list(range(len(list(cons.keys()))))
        con_ind = con_ind[start_ind:] + con_ind[:start_ind]

        # Run minimization on the initial structure
        # NOTE: LAMMPS atoms are 1-indexed, so the constraints get incremented by 1
        counter = 0
        f.write("# minimize molecule energy with restraints\n"+\
                "dump           minimization all atom ${coords_freq}"+" {}.lammpstrj\n".format(start_ind)+\
                "dump_modify    minimization scale no\n"+\
                "dump_modify    minimization sort  id\n")
        for count_j,j in enumerate(cons[con_ind[0]]):
            f.write("fix CON_{} all restrain dihedral {} {} {} {} 1000000.0 1000000.0 {}\n".format(count_j,j[0][0]+1,j[0][1]+1,j[0][2]+1,j[0][3]+1,j[1]+180.0)+\
                    "fix_modify CON_{} energy yes\n".format(count_j))
        for count_j,j in enumerate(other_constraints[con_ind[0]]):
            # f.write("fix SEC_CON_{} all restrain dihedral {} {} {} {} 10.0 10.0 {}\n".format(count_j,j[0][0]+1,j[0][1]+1,j[0][2]+1,j[0][3]+1,j[1]+180.0)+\
            #         "fix_modify SEC_CON_{} energy yes\n".format(count_j))

            f.write("fix SEC_CON_{} all restrain dihedral {} {} {} {} {:<6.1f} {:<6.1f} {}\n".format(count_j,j[0][0]+1,j[0][1]+1,j[0][2]+1,j[0][3]+1,k_sec,k_sec,j[1]+180.0)+\
                    "fix_modify SEC_CON_{} energy yes\n".format(count_j))

        f.write("minimize 0.0 0.0 {} 100000\n".format(max_evals)+\
                "undump minimization\n")

        # Break if the scan is not requested
        if mode not in ["forward","reverse"]:
            return
        print("I RAN!!!!!!!")
        ############################################################################################################################
        # NOTE: when the above return command is removed, the rest of these loops are used to perform a full scan of the dihedral. #
        #       These are being kept for now until the current protocol is thoroughly vetted.                                      #
        ############################################################################################################################

        # Initialize the number of interpolation steps so that interpolation occurs in steps of 1 degree. A minimum of 2 steps is allowed (i.e. the configuration values themselves)
        num_steps = int(round(360.0/float(len(list(cons.keys()))))) + 1
        if num_steps < 2: num_steps = 2

        # Loop over the configurational constraints in forward direction
        # NOTE: +[0] is added to the j loop to repeat the first configuration at the end of the scan.
        if mode == "forward":
            counter += 1        
            for j in con_ind[1:]+[con_ind[0]]:
                prev_ind = j - 1
                current_ind = j

                # Special condition for the first configuration (second if statement pertains for the reverse scan
                if prev_ind == -1: prev_ind = len(list(cons.keys()))-1
                if prev_ind == len(list(cons.keys())): prev_ind = 0
                               
                # Run interpolation between the constraints. num_steps controls the spacing (initialized above to ensure a 1 degree step)
                f.write("\n# Performing {}->{} interpolation\n".format(prev_ind,current_ind))
                for k in np.linspace(0,1,num_steps)[1:-1]:
                    for count_n,n in enumerate(cons[j]):

                        # A quirk of the forward scan is that the on the *last* step the previous step and current step are separated by 360+dihedral_step which breaks the interpolation algorithm
                        # For this special case, the current constraint gets shifted by 360.
                        if j == 0:
                            angle = (1.0-k)*cons[prev_ind][count_n][1] + k*(cons[current_ind][count_n][1]+360.0)
                        else:
                            angle = (1.0-k)*cons[prev_ind][count_n][1] + k*cons[current_ind][count_n][1]
                        f.write("fix CON_{} all restrain dihedral {} {} {} {} 1000000.0 1000000.0 {}\n".format(count_n,n[0][0]+1,n[0][1]+1,n[0][2]+1,n[0][3]+1,angle+180.0)+\
                                "fix_modify CON_{} energy yes\n".format(count_n))
                    f.write("minimize 0.0 0.0 {} 100000\n".format(max_evals))
                    for count_n,n in enumerate(cons[current_ind]):
                        f.write("unfix CON_{}\n".format(count_n))

                # Run minimization at the current configuration (dump geometry)
                f.write("\n# Performing optimized evaluation at step {}\n".format(j))
                for count_k,k in enumerate(cons[current_ind]):
                    f.write("fix CON_{} all restrain dihedral {} {} {} {} 1000000.0 1000000.0 {}\n".format(count_k,k[0][0]+1,k[0][1]+1,k[0][2]+1,k[0][3]+1,k[1]+180.0)+\
                            "fix_modify CON_{} energy yes\n".format(count_k))                            
                f.write("dump               minimization all atom ${coords_freq}"+" {}.lammpstrj\n".format(j)+\
                        "dump_modify        minimization scale no\n"+\
                        "dump_modify        minimization sort  id\n")
                f.write("minimize 0.0 0.0 {} 100000\n".format(max_evals)+\
                        "undump minimization\n")
                for count_k,k in enumerate(cons[current_ind]):
                    f.write("unfix CON_{}\n".format(count_k))

                # Increment counter
                counter += 1

        # Loop over the configurational constraints in reverse direction
        # NOTE: +[0] is added to the j loop to repeat the first configuration at the end of the scan.
        if mode == "reverse":
            counter = len(list(cons.keys()))-1
            for j in con_ind[1:][::-1]+[con_ind[0]]:
                prev_ind = j + 1
                current_ind = j


                # Special condition for the last configuration (first if statement pertains for the forward scan)
                if prev_ind == -1: prev_ind = len(list(cons.keys()))-1
                if prev_ind == len(list(cons.keys())): prev_ind = 0

                # Run interpolation between the constraints. num_steps controls the spacing (initialized above to ensure a 1 degree step)
                f.write("\n# Performing {}->{} interpolation\n".format(prev_ind,current_ind))
                for k in np.linspace(0,1,num_steps)[1:-1]:
                    for count_n,n in enumerate(cons[j]):

                        # A quirk of the reverse scan is that on the *first* step the previous step and current step are separated by 360+dihedral_step which breaks the interpolation algorithm
                        # For this special case, the previous constraint gets shifted by 360.
                        if j == len(list(cons.keys()))-1:
                            angle = (1.0-k)*(cons[prev_ind][count_n][1]+360.0) + k*cons[current_ind][count_n][1]
                        else:
                            angle = (1.0-k)*cons[prev_ind][count_n][1] + k*cons[current_ind][count_n][1]
                        f.write("fix CON_{} all restrain dihedral {} {} {} {} 1000000.0 1000000.0 {}\n".format(count_n,n[0][0]+1,n[0][1]+1,n[0][2]+1,n[0][3]+1,angle+180.0)+\
                                "fix_modify CON_{} energy yes\n".format(count_n))
                    f.write("minimize 0.0 0.0 {} 100000\n".format(max_evals))
                    for count_n,n in enumerate(cons[current_ind]):
                        f.write("unfix CON_{}\n".format(count_n))

                # Run minimization at the current configuration (dump geometry)
                f.write("\n# Performing optimized evaluation at step {}\n".format(j))
                for count_k,k in enumerate(cons[current_ind]):
                    f.write("fix CON_{} all restrain dihedral {} {} {} {} 1000000.0 1000000.0 {}\n".format(count_k,k[0][0]+1,k[0][1]+1,k[0][2]+1,k[0][3]+1,k[1]+180.0)+\
                            "fix_modify CON_{} energy yes\n".format(count_k))                            
                f.write("dump               minimization all atom ${coords_freq}"+" {}.lammpstrj\n".format(j)+\
                        "dump_modify        minimization scale no\n"+\
                        "dump_modify        minimization sort  id\n")
                f.write("minimize 0.0 0.0 {} 100000\n".format(max_evals)+\
                        "undump minimization\n")
                for count_k,k in enumerate(cons[current_ind]):
                    f.write("unfix CON_{}\n".format(count_k))

                # Increment counter
                counter -= 1

    f.close()    
    return

# Description: A wrapper for the write commands for generating the lammps input file
def Write_MD_input_2(Filename,frequency,Onefourscale,Pair_styles,Dihedral_styles,Dihedral_Data,ind):

    with open(Filename+'.in.init','w') as f:
        f.write("# lammps input file for polymer simulation with dilute ions\n\n"+\
                "# VARIABLES\n"+\
                "variable        data_name       index   {}\n".format(Filename+'.data')+\
                "variable        settings_name   index   {}\n".format(Filename+'.in.settings')+\
                "variable        coords_freq     index   {}\n".format(frequency)+\
                "variable        vseed           index   {: <6d}\n".format(int(random.random()*100000))+\
                
                "#===========================================================\n"+\
                "# GENERAL PROCEDURES\n"+\
                "#===========================================================\n"+\
                "units		real	# g/mol, angstroms, fs, kcal/mol, K, atm, charge*angstrom\n"+\
                "dimension	3	# 3 dimensional simulation\n"+\
                "newton		off	# use Newton's 3rd law\n"+\
                "boundary	s s s	# periodic boundary conditions \n"+\
                "atom_style	full    # molecular + charge\n\n"+\

                "#===========================================================\n"+\
                "# FORCE FIELD DEFINITION\n"+\
                "#===========================================================\n"+\
                'special_bonds  lj   0.0 0.0 0.0     # NO     1-4 LJ interactions\n'+\
                'special_bonds  coul 0.0 0.0 {}      # REDUCE 1-4 electrostatics\n'.format(Onefourscale)+\
                'pair_style     hybrid {} # outer_LJ outer_Coul (cutoff values, see LAMMPS Doc)\n'.format(Pair_styles)+\
                'bond_style     harmonic             # parameters needed: k_bond, r0\n'+\
                'angle_style    harmonic             # parameters needed: k_theta, theta0\n')
        if len(Dihedral_styles) >= 1:
            f.write('dihedral_style hybrid {}            # parameters needed: V1, V2, V3, V4\n'.format(' '.join(Dihedral_styles)))            
        f.write('pair_modify    shift yes mix arithmetic       # using Lorenz-Berthelot mixing rules\n\n'+\

                "#===========================================================\n"+\
                "# SETUP SIMULATIONS\n"+\
                "#===========================================================\n\n"+\
                "# READ IN COEFFICIENTS/COORDINATES/TOPOLOGY\n"+\
                'read_data ${data_name}\n'+\
                'include ${settings_name}\n\n'+\

                "# SET RUN PARAMETERS\n"+\
                "timestep	1.0		# fs\n"+\
                "run_style	verlet 		# Velocity-Verlet integrator\n"+\
                "neigh_modify every 1 delay 0 check no # More relaxed rebuild criteria can be used\n\n")

        f.write("#===========================================================\n"+\
                "# RUN MINIMIZATIONS\n"+\
                "#===========================================================\n\n")

        R_seed = int(random.random()*1000000.0)        
        for i in range(36):

            # NOTE: LAMMPS atoms are 1-indexed, so the constraints get incremented by 1
            f.write("# minimize molecule energy with restraints\n"+\
                    "dump               minimization all atom ${coords_freq}"+" {}.lammpstrj\n".format(i)+\
                    "dump_modify	minimization scale no\n"+\
                    "dump_modify	minimization sort  id\n"+\
                    "velocity all create 100.0 {} mom yes rot yes dist gaussian\n".format(R_seed)+\
                    "fix NVE all nve\n"+\
                    "fix TFIX all langevin 1000.0 0.0 100 {} zero yes\n".format(R_seed))
            for count_j,j in enumerate(Dihedral_Data[ind][str(i)]["Dihedral_Atoms"]):
                f.write("fix r_{} all restrain dihedral {} {} {} {} 0.0 500.0 {}\n".format(count_j,j[0][0]+1,j[0][1]+1,j[0][2]+1,j[0][3]+1,j[1]*180.0/np.pi-180.0))
            f.write("run 1000\n"
                    "fix TFIX all langevin 0.0 0.0 100 {} zero yes\n".format(R_seed))
            for count_j,j in enumerate(Dihedral_Data[ind][str(i)]["Dihedral_Atoms"]):
                f.write("fix r_{} all restrain dihedral {} {} {} {} 500.0 500.0 {}\n".format(count_j,j[0][0]+1,j[0][1]+1,j[0][2]+1,j[0][3]+1,j[1]*180.0/np.pi-180.0))
            f.write("run 1000\n"+\
                    "# sanity check for convergence\n"+\
                    "minimize 1e-6 1e-9 1000 100000\n"+\
                    "# undump/unfix\n"+\
                    "undump minimization\n")
            for count_j,j in enumerate(Dihedral_Data[ind][str(i)]["Dihedral_Atoms"]):
                f.write("unfix r_{}\n".format(count_j))
            f.write("\n")

#             f.write("# Apply constraint to the dihedral\n"+\
#                     "fix holdem all restrain dihedral {} {} {} {} 1000.0 1000.0 {}\n\n".format(constraint[0]+1,constraint[1]+1,constraint[2]+1,constraint[3]+1,angle+10.0*i)) # LAMMPS atoms are 1-indexed

#             f.write("# CREATE COORDINATE DUMPS FOR MINIMIZATION\n"+\
#                     "dump           minimization all atom 10 {}.lammpstrj\n".format(i)+\
#                     "dump_modify	minimization scale no\n"+\
#                     "dump_modify	minimization sort  id\n"+\
#                     "minimize 0.0 1.0e-8 10000 10000000\n"+\
#                     "undump minimization\n"+\
#                     "unfix holdem\n\n")

    f.close()

    return

# Grab the optimized geometry from an MD minimization and calculate the FF potential for the configuration
def parse_OPLSFIT_traj(trj_file,atomtypes,charges,adj_mat,con_dihedrals,qc,UA_opt,one_five_opt=1,include_con=0,one_four_scale_coul=0.0,one_four_scale_vdw=0.0):
    
    # Find the number of steps in the minimization
    frames = 0
    with open(trj_file,'r') as f: 
        for lines in f: 
            if "ITEM: TIMESTEP" in lines: frames += 1

    # Grab optimized geometry
    count = 0
    frame_flag = 0
    parse_flag = 0
    atom_count = 0
    new_geo = np.zeros([len(atomtypes),3])
    with open(trj_file,'r') as f:
        for lines in f:

            if len(lines) == 0: continue 
            if "ITEM: TIMESTEP" in lines: 
                count += 1
                if count == frames: 
                    frame_flag = 1
                    continue
            if frame_flag == 1 and "ITEM: ATOMS" in lines:
                parse_flag = 1
                continue
            if parse_flag == 1:
                new_geo[atom_count,:] = np.array([float(lines.split()[2]),float(lines.split()[3]),float(lines.split()[4])])
                atom_count += 1

    # Calculate the total FF energy (less the constrained mode(s))
    Bonds,Angles,Dihedrals,One_fives = Find_modes(adj_mat,atomtypes,return_all=0)

    # Sum up the separate contributions to the total energy
    E_Tot = 0.0
    E_Tot += E_bond(new_geo,[ i.split("-UA")[0] for i in atomtypes ],Bonds,QC=qc)
    E_Tot += E_angle(new_geo,[ i.split("-UA")[0] for i in atomtypes ],Angles,QC=qc)
    if include_con == 0:
        E_Tot += E_dihedral(new_geo,atomtypes,[ i for i in Dihedrals if i not in con_dihedrals and (i[3],i[2],i[1],i[0]) not in con_dihedrals ],QC=qc)  # List compression removes the constrained modes from the dihedral energy calculation
    elif include_con == 1:
        E_Tot += E_dihedral(new_geo,atomtypes,Dihedrals,QC=qc)  # List compression removes the constrained modes from the dihedral energy calculation
    if UA_opt == 1:
        if one_five_opt == 1:
            E_Tot += E_coul(new_geo,adj_mat,atomtypes,charges,[Bonds,Angles,Dihedrals+One_fives],(0.0,0.0,one_four_scale_coul),"UA")
            E_Tot += E_LJ(new_geo,adj_mat,[ i.split("-UA")[0] for i in atomtypes ],charges,[Bonds,Angles,Dihedrals+One_fives],(0.0,0.0,one_four_scale_vdw),"UA")
        elif one_five_opt == 0:
            E_Tot += E_coul(new_geo,adj_mat,atomtypes,charges,[Bonds,Angles,Dihedrals],(0.0,0.0,one_four_scale_coul),"UA")
            E_Tot += E_LJ(new_geo,adj_mat,[ i.split("-UA")[0] for i in atomtypes ],charges,[Bonds,Angles,Dihedrals],(0.0,0.0,one_four_scale_vdw),"UA")
    else:
        if one_five_opt == 1:
            E_Tot += E_coul(new_geo,adj_mat,atomtypes,charges,[Bonds,Angles,Dihedrals+One_fives],(0.0,0.0,one_four_scale_coul),"AA")
            E_Tot += E_LJ(new_geo,adj_mat,[ i.split("-UA")[0] for i in atomtypes ],charges,[Bonds,Angles,Dihedrals+One_fives],(0.0,0.0,one_four_scale_vdw),"AA")
        elif one_five_opt == 0:
            E_Tot += E_coul(new_geo,adj_mat,atomtypes,charges,[Bonds,Angles,Dihedrals],(0.0,0.0,one_four_scale_coul),"AA")
            E_Tot += E_LJ(new_geo,adj_mat,[ i.split("-UA")[0] for i in atomtypes ],charges,[Bonds,Angles,Dihedrals],(0.0,0.0,one_four_scale_vdw),"AA")

    # Calculate the constrained dihedral angles
    angles = { i:[] for i in con_dihedrals }
    for i in con_dihedrals:

        # Calculate all dihedral angles
        atom_1 = new_geo[i[0]]
        atom_2 = new_geo[i[1]]
        atom_3 = new_geo[i[2]]
        atom_4 = new_geo[i[3]]
        v1 = atom_2-atom_1
        v2 = atom_3-atom_2
        v3 = atom_4-atom_3
        angles[i] = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )*180.0/np.pi
        
    return new_geo,E_Tot,angles

# Grab the optimized geometry from an MD minimization and calculate the FF potential for the configuration
def parse_torsample_traj(trj_file,dihedral_atoms,dihedral_types,Temp,bin_width=10.0):

    # initialize histogram and maximum atom index
    # NOTE: histogram is initialized with ones to avoid division by zero for any bins that 
    #       aren't sampled. The error is negligible for any reasonable sample size.
    hist = { i:np.zeros(int(360/bin_width)) for i in set(dihedral_types) }
    potential = {}
    max_atom = max([ j for i in dihedral_atoms for j in i ])
    for i in dihedral_types:
        potential[i] = 1

    # parse dihedrals
    frame_count = 0
    frame_flag  = 0
    parse_flag  = 0
    atom_count  = 0
    Geo = np.zeros([max_atom+1,3])
    with open(trj_file,'r') as f:
        for lines in f:

            # Skip empty lines
            if len(lines) == 0: continue 

            # Find frame start
            if "ITEM: TIMESTEP" in lines: 
                frame_count += 1
                frame_flag = 1
                continue

            # Find geo start
            if frame_flag == 1 and "ITEM: ATOMS" in lines:
                frame_flag = 0
                parse_flag = 1
                continue

            # Parse geometry (once all atoms necessary to calculate all dihedrals are in hand, the inner loop is
            # triggered to histogram the dihedral values)
            if parse_flag == 1:
                Geo[atom_count,:] = np.array([float(lines.split()[2]),float(lines.split()[3]),float(lines.split()[4])])

                # If all of the necessary atoms have been parsed for this geometry then calculate the 
                # various dihedral angles and update the histogram(s)
                if atom_count == max_atom:
                    for count_i,i in enumerate(dihedral_atoms):

                        # Calculate all dihedral angles
                        atom_1 = Geo[i[0]]
                        atom_2 = Geo[i[1]]
                        atom_3 = Geo[i[2]]
                        atom_4 = Geo[i[3]]
                        v1 = atom_2-atom_1
                        v2 = atom_3-atom_2
                        v3 = atom_4-atom_3
                        theta = np.arctan2( np.dot(v1,np.cross(v2,v3))*(np.dot(v2,v2))**(0.5) , np.dot(np.cross(v1,v2),np.cross(v2,v3)) )*180.0/np.pi

                        # NO SYMMETRIZATION
                        # Calculate ind based on middle-binning (the if statement is required to fold the 175-180 interval
                        # into the (-)180-(-)175 bin
                        ind = int( ( theta + 180.0 + (float(bin_width)/2.0) ) / float(bin_width) )
                        if ind >= float(360/bin_width): ind = 0
#                        hist[dihedral_types[count_i]][int((theta+180.0)/10.0)] += 1
                        hist[dihedral_types[count_i]][ind] += 1

                    # Reset flag(s)
                    parse_flag = 0
                    atom_count = 0

                # Increment atom_count
                else:
                    atom_count += 1

    # remove zeros from any of the histograms to avoid division by zero
    # then convert to kcal/mol
    kb_kcalmol = 0.0019872041
    for i in dihedral_types:
        hist[i][np.where(hist[i] == 0)] = 0.0000001
        max_prob = max(hist[i])
        potential[i] = -kb_kcalmol*Temp*log(hist[i]/max_prob)
    
    # NOTE: it might be convenient to duplicate the 0 entry and append it to the end of the potentials
    #       and add the 180 degree entry to the angle return
    return potential,np.arange(-180,180,bin_width)

def lstsq_fit(Dihedral_Data,dihedral_types,ind,qc):


    # Collect the dihedral types being fit
#    dihedral_types = sorted(set([ j[0] for i in Dihedral_Data[ind].keys() for j in Dihedral_Data[ind][i]["Dihedrals"] ]))

    # Perform the least-squares fit
#    print "\nPerforming least-squares fit to fit potential (solving A.x=b)..."
    params = {}
    Keys = list(Dihedral_Data[ind].keys())
    b = np.array([ Dihedral_Data[ind][i][qc+"_Fit"] for i in Keys ])
    A = np.zeros([len(b),len(dihedral_types)*5]) 
#    print "Buildng A matrix..."
    for count_i,i in enumerate(Keys):
        for count_j,j in enumerate(dihedral_types):
            for k in [ m for m in Dihedral_Data[ind][i]["Dihedrals"] if m[0] == j ]:

                # generate coefficients for the fit expression  E_OPLS[count_i] += V0 + 0.5*V1*(1+np.cos(j[1])) + 0.5*V2*(1-np.cos(2.0*j[1])) + 0.5*V3*(1+np.cos(3.0*j[1])) + 0.5*V4*(1-np.cos(4.0*j[1]))
                A[count_i,(count_j*5)+0] += 1.0
                A[count_i,(count_j*5)+1] += 0.5*(1.0+np.cos(1.0*k[1]))  
                A[count_i,(count_j*5)+2] += 0.5*(1.0-np.cos(2.0*k[1]))  
                A[count_i,(count_j*5)+3] += 0.5*(1.0+np.cos(3.0*k[1]))  
                A[count_i,(count_j*5)+4] += 0.5*(1.0-np.cos(4.0*k[1]))  

    # Calculate inverse
#    print "Calculating inverse(A)..."
    A_inv = linalg.pinv(A)

    # Calculate fit coefficients
#    print "Calculating fit coefficients (x)..."
    x = np.dot(A_inv,b)

    # Initialize the params dictionary with dihedral parameters (NOTE: for all fit, there is a constant factor added to each dihedral
    for i in dihedral_types:
        if i not in list(params.keys()):
            params[i] = {}
    for count_i,i in enumerate(dihedral_types):
        params[i] = tuple([x[(count_i*5)+0],x[(count_i*5)+1],x[(count_i*5)+2],x[(count_i*5)+3],x[(count_i*5)+4]])

    # Calculate Fit Potential and xhi2
    xhi2 = 0.0
    for count_i,i in enumerate(Dihedral_Data[ind].keys()):
        Dihedral_Data[ind][i][qc+"_Fit_lstsq"] = OPLS_EFF(params,Dihedral_Data[ind][i]["Dihedrals"])
        Dihedral_Data[ind][i]["E_FF_lstsq"] = Dihedral_Data[ind][i][qc] - Dihedral_Data[ind][i][qc+"_Fit_lstsq"]
        xhi2 += (Dihedral_Data[ind][i][qc+"_Fit_lstsq"] - Dihedral_Data[ind][i]["DFT_Fit"])**(2.0)
#    print "Average Absolute Error in least-squares fit: {} kcal/mol".format(xhi2/len(Dihedral_Data[ind].keys()))

    return params

# Write a dihedral table file for tabulated MD simulations
def write_dihedral_table(name,fit_potentials,dihedrals):
    tables = {}
    with open(name+'.dihedral.table','w') as f:
        for count_i,i in enumerate(dihedrals):
            active_ind = len(list(fit_potentials[i].keys()))-1
            tables[i] = (name+'.dihedral.table',"DIH_TABLE_{}".format(count_i))
            f.write("DIH_TABLE_{}\n".format(count_i))
            f.write("N {} NOF DEGREES CHECKU {}\n\n".format(len(fit_potentials[i][active_ind]),name+'.DIH_TABLE_{}.table'.format(count_i)))
            for count_j,j in enumerate(fit_potentials[i][active_ind]):
                f.write("{:< 8d} {:< 12.5f} {:< 12.5f}\n".format(count_j+1,-180.0+(10.0*count_j),j)) 
            f.write("\n")
    return tables


# Return a dictionary of interpolation objects for the QC potentials of each scan
# with respect to the angle of the first instance of each type being fit. 
def interp_QC(Dihedral_Data,qc):

    QC_interp_dict = {}
    for i in Dihedral_Data:
        QC_interp_dict[i] = {}
        for d in set([ z[0] for z in Dihedral_Data[i]["0"]["Dihedrals"] ]):

            # Grab the index of the first instance of this dihedral type
            dihedral_ind = [ count_k for count_k,k in enumerate(Dihedral_Data[i]["0"]["Dihedrals"]) if k[0] == d ][0]

            # Grab angles for the qc potential plot
            # NOTE: the angles are wrt to the first instance of the dihedral type, d,
            #       being fit. 
            angles = []
            E_QC   = []
            for k in list(Dihedral_Data[i].keys()):
                E_QC += [Dihedral_Data[i][k][qc]]
                angles += [Dihedral_Data[i][k]["Dihedrals"][dihedral_ind][1]]
                
            # Sort wrt to the angle of the current dihedral type
            together = sorted([ (k,E_QC[count_k]) for count_k,k in enumerate(angles) ])
            angles   = np.array([ k[0] for k in together ])*180.0/np.pi
            E_QC     = np.array([ k[1] for k in together ])

            # Create an interpolation of the QC curve
            if np.abs(np.abs(angles[0])-180.0) < np.abs(np.abs(angles[-1])-180.0):
                QC_interp = interpolate.interp1d(angles,E_QC,kind='linear',fill_value=E_QC[0],bounds_error=False)
            else:
                QC_interp = interpolate.interp1d(angles,E_QC,kind='linear',fill_value=E_QC[-1],bounds_error=False)

            E_QC_interp = QC_interp(np.arange(-180,180,10))

            # Normalize the data by the first element (arbitrary choice)
            E_QC = (E_QC_interp-E_QC_interp[0])

            # Create an interpolation of the QC curve
            if np.abs(np.abs(angles[0])-180.0) < np.abs(np.abs(angles[-1])-180.0):
                QC_interp_dict[i][d] = interpolate.interp1d(np.arange(-180,180,10),E_QC,kind='linear',fill_value=E_QC[0],bounds_error=False)
            else:
                QC_interp_dict[i][d] = interpolate.interp1d(np.arange(-180,180,10),E_QC,kind='linear',fill_value=E_QC[-1],bounds_error=False)

    return QC_interp_dict


# Description:
# Rotate Point by an angle, theta, about the vector with an orientation of v1 passing through v2. 
# Performs counter-clockwise rotations (i.e., if the direction vector were pointing
# at the spectator, the rotations would appear counter-clockwise)
# For example, a 90 degree rotation of a 0,0,1 about the canonical 
# y-axis results in 1,0,0.
#
# Point: 1x3 array, coordinates to be rotated
# v1: 1x3 array, point the rotation passes through
# v2: 1x3 array, rotation direction vector
# theta: scalar, magnitude of the rotation (defined by default in degrees)
def axis_rot(Point,v1,v2,theta,mode='angle'):

    # Temporary variable for performing the transformation
    rotated=np.array([Point[0],Point[1],Point[2]])

    # If mode is set to 'angle' then theta needs to be converted to radians to be compatible with the
    # definition of the rotation vectors
    if mode == 'angle':
        theta = theta*np.pi/180.0

    # Rotation carried out using formulae defined here (11/22/13) http://inside.mines.edu/fs_home/gmurray/ArbitraryAxisRotation/)
    # Adapted for the assumption that v1 is the direction vector and v2 is a point that v1 passes through
    a = v2[0]
    b = v2[1]
    c = v2[2]
    u = v1[0]
    v = v1[1]
    w = v1[2]
    L = u**2 + v**2 + w**2

    # Rotate Point
    x=rotated[0]
    y=rotated[1]
    z=rotated[2]

    # x-transformation
    rotated[0] = ( a * ( v**2 + w**2 ) - u*(b*v + c*w - u*x - v*y - w*z) )\
             * ( 1.0 - np.cos(theta) ) + L*x*np.cos(theta) + L**(0.5)*( -c*v + b*w - w*y + v*z )*np.sin(theta)

    # y-transformation
    rotated[1] = ( b * ( u**2 + w**2 ) - v*(a*u + c*w - u*x - v*y - w*z) )\
             * ( 1.0 - np.cos(theta) ) + L*y*np.cos(theta) + L**(0.5)*(  c*u - a*w + w*x - u*z )*np.sin(theta)

    # z-transformation
    rotated[2] = ( c * ( u**2 + v**2 ) - w*(a*u + b*v - u*x - v*y - w*z) )\
             * ( 1.0 - np.cos(theta) ) + L*z*np.cos(theta) + L**(0.5)*( -b*u + a*v - v*x + u*y )*np.sin(theta)

    rotated = rotated/L
    return rotated

# Align a set of geometries (list of arrays) about the 1-2 and 2-3 bonds of a supplied dihedral (tupel of indices)
def align_geos(geos,dihedral):

    # Align geometries across the scan so that the 1-2 and 2-3 bonds are aligned and the geometry is centered about the 2-3 bond. 
    # Alignment is based on the first configuation. 
    for geo_count,geo in enumerate(geos):

        # Store alignment data for the first key
        if geo_count == 0:
            atom_1,atom_2,atom_3,atom_4 = dihedral
            v1 = geo[atom_1] - geo[atom_2]
            v2 = geo[atom_2] - geo[atom_3]

        # Perform 1-2 alignment (if statement avoids alignment if it is already aligned. Sometimes there are roundoff errors when the argument of acos is near 1.0)
        c_v1    = geo[atom_1] - geo[atom_2]
        if np.dot(c_v1,v1) / ( norm(c_v1)*norm(v1) ) < 1.0:
            angle   = acos( np.dot(c_v1,v1) / ( norm(c_v1)*norm(v1) ) )
            rot_vec = np.cross(v1,c_v1)
            if norm(rot_vec) > 0.0:

                # try positive rotation
                tmp = axis_rot(geo[atom_1],rot_vec,geo[atom_2],angle,mode="radian") - geo[atom_2]
                diff = acos( np.dot(tmp,v1) / ( norm(tmp)*norm(v1) ) )
                if np.abs(diff) < 0.01:
                    for count_i,i in enumerate(geo):
                        geo[count_i] = axis_rot(i,rot_vec,geo[atom_2],angle,mode="radian")
                # else, perform negative rotation
                else:
                    for count_i,i in enumerate(geo):
                        geo[count_i] = axis_rot(i,rot_vec,geo[atom_2],-angle,mode="radian")

        # Perform 2-3 alignment
        c_v2     = geo[atom_2] - geo[atom_3]
        rot_vec  = geo[atom_1] - geo[atom_2]
        min_diff = 100.0
        steps    = [0.1,0.01,0.001,0.0001]

        # Find rotation angle. The nested loops are simple bracketing algorithm
        for count_i,i in enumerate(steps):
            if count_i == 0:
                a=0.0
                b=2.0*np.pi
            else:
                a=angle-steps[count_i-1]
                b=angle+steps[count_i-1]
            for j in np.arange(a,b,i):
                tmp = geo[atom_2] - axis_rot(geo[atom_3],rot_vec,geo[atom_2],j,mode="radian")

                # if statement avoids domain error in the acos call. Sometimes there are roundoff errors when the argument of acos is near 1.0.
                if np.dot(tmp,v2) / ( norm(tmp)*norm(v2) ) < 1.0:
                    diff = np.abs(acos( np.dot(tmp,v2)/( norm(tmp)*norm(v2) )))
                else:
                    diff = 100.0

                # Assign alignment angle
                if diff < min_diff:
                    min_diff = diff
                    angle = j

        # Align about 2-3 based on the best angle
        for count_i,i in enumerate(geo):
            geo[count_i] = axis_rot(i,rot_vec,geo[atom_2],angle,mode="radian")

        # Center the geometry about the 2-3 centroid
        centroid = ( geo[atom_2] + geo[atom_3] ) / 2.0
        for count_i,i in enumerate(geo):
            geo[count_i] = geo[count_i] - centroid 
        geos[geo_count] = geo

    return geos

# Commands for writing the lammps inputs, running the lammps job, and parsing the output to generate the thermal distribution for a dihedral(s)
def gen_thermal_distribution(geo,adj_mat,atomtypes,charges,dihedrals,dihedral_types,lammps_exe,t_sample=10E6,freq=100,T=298,onefourscale=0.0,keep_files=0,bin_width=10.0,qc='DFT'):

    global FF_dict

    # Initialize Sim_Box
    Sim_Box = np.array([min(geo[:,0])-100,max(geo[:,0])+100,min(geo[:,1])-100,max(geo[:,1])+100,min(geo[:,2])-100,max(geo[:,2])+100])

    # Write the lammps datafile (the function returns a dictionary, Atom_type_dict, holding the mapping between
    # atom_types and the lammps id numbers; this mapping is needed for setting fixes)                    
    Bonds,Angles,Dihedrals,One_fives,Bond_types,Angle_types,Dihedral_types,One_five_types = Find_modes(adj_mat,atomtypes,return_all=1)
    Dihedral_styles = set([ FF_dict["dihedrals"][i]["DFT"][0] for i in Dihedral_types ])
    Atom_type_dict = Write_MD_data("thermal_sample",atomtypes,Sim_Box,geo,Bonds,Bond_types,Angles,Angle_types,\
                                   Dihedrals,Dihedral_types,One_fives,charges,[],qc,one_five_opt=0,tables=[])

    # Write the lammps input files
    Write_TORSAMPLE_input("thermal_sample",freq,onefourscale,t_sample,"lj/cut/coul/cut 100.0 100.0",Dihedral_styles,T)

    # Run the LAMMPS minimization
    lammps_call = "{} -in {}.in.init -screen none -log none".format(lammps_exe,"thermal_sample","thermal_sample","thermal_sample")
    subprocess.call(lammps_call.split())

    # Grab the minimized geometry and FF energy
    potentials,angles = parse_torsample_traj("thermal_sample.lammpstrj",dihedrals,dihedral_types,T,bin_width)
    
    # Remove run files if option enabled
    if keep_files == 0:
        os.remove('thermal_sample.lammpstrj')
        os.remove('thermal_sample.in.init')
        os.remove('thermal_sample.data')
        os.remove('thermal_sample.in.settings')

    return potentials,angles

# Wrapper function for writing the lammps data file which is read by the .in.init file during run initialization. This function is 
# nearly identical to the function of the same name within the gen_md_for_vdw.py script, except the qc option has been added to be compatible
# with the FF_dict structure and the Molecule variable has been removed because only one molecule is present in each simulation.
def Write_MD_data(Filename,Atom_types,Sim_Box,Geometry,Bonds,Bond_types,Angles,Angle_types,Dihedrals,Dihedral_types,One_fives,\
                      Charges,Scanned_Dihedrals,qc,one_five_opt=0,tables=[]):

    global FF_dict
    
    # Create type dictionaries (needed to convert each atom,bond,angle, and dihedral type to consecutive numbers as per LAMMPS convention)
    # Note: LAMMPS orders atomtypes, bonds, angles, dihedrals, etc as integer types. Each of the following dictionaries holds the mapping between
    #       the true type (held in the various types lists) and the lammps type_id, obtained by enumerated iteration over the respective set(types).
    Atom_type_dict = {}
    for count_i,i in enumerate(sorted(set(Atom_types))):
        for j in Atom_types:
            if i == j:
                Atom_type_dict[i]=count_i+1
            if i in list(Atom_type_dict.keys()):
                break
    Bond_type_dict = {}
    for count_i,i in enumerate(sorted(set(Bond_types))):
        for j in Bond_types:
            if i == j:
                Bond_type_dict[i]=count_i+1
            if i in list(Bond_type_dict.keys()):
                break
    
    # Add 1-5 bonds as the last type (used for turning off 1-5 non-bonded interactions
    if one_five_opt == 1:
        Bond_type_dict["one_fives"] = count_i+2   

    Angle_type_dict = {}
    for count_i,i in enumerate(sorted(set(Angle_types))):
        for j in Angle_types:
            if i == j:
                Angle_type_dict[i]=count_i+1
            if i in list(Angle_type_dict.keys()):
                break
    Dihedral_type_dict = {}
    for count_i,i in enumerate(sorted(set(Dihedral_types))):
        for j in Dihedral_types:
            if i == j:
                Dihedral_type_dict[i]=count_i+1
            if i in list(Dihedral_type_dict.keys()):
                break

    # Write the data file
    with open(Filename+'.data','w') as f:
        
        # Write system properties
        f.write("LAMMPS data file via vdw_self_gen.py, on {}\n\n".format(datetime.datetime.now()))

        f.write("{} atoms\n".format(len(Atom_types)))
        f.write("{} atom types\n".format(len(set(Atom_types))))
        if len(Bonds) > 0:
            if one_five_opt == 1:
                f.write("{} bonds\n".format(len(Bonds)+len(One_fives)))
                f.write("{} bond types\n".format(len(set(Bond_types))+1))
            else:
                f.write("{} bonds\n".format(len(Bonds)))
                f.write("{} bond types\n".format(len(set(Bond_types))))
        if len(Angles) > 0:
            f.write("{} angles\n".format(len(Angles)))
            f.write("{} angle types\n".format(len(set(Angle_types))))
        if len(Dihedrals) > 0:
            f.write("{} dihedrals\n".format(len(Dihedrals)))
            f.write("{} dihedral types\n\n".format(len(set(Dihedral_types))))

        # Write box dimensions
        f.write("{:< 20.16f} {:< 20.16f} xlo xhi\n".format(Sim_Box[0],Sim_Box[1]))
        f.write("{:< 20.16f} {:< 20.16f} ylo yhi\n".format(Sim_Box[2],Sim_Box[3]))
        f.write("{:< 20.16f} {:< 20.16f} zlo zhi\n\n".format(Sim_Box[4],Sim_Box[5]))

        # Write FF_dict["masses"]
        f.write("Masses\n\n")
        for count_i,i in enumerate(sorted(set(Atom_types))):
            for j in set(Atom_types):
                if Atom_type_dict[j] == count_i+1:
                    f.write("{} {:< 8.6f}\n".format(count_i+1,FF_dict["masses"][str(j)])) # count_i+1 bc of LAMMPS 1-indexing
        f.write("\n")

        # Write Bond Coeffs
        f.write("Bond Coeffs\n\n")
        for count_i,i in enumerate(set(Bond_types)):
            for j in set(Bond_types):
                if Bond_type_dict[j] == count_i+1:
                    f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+1,FF_dict["bonds"][j][qc][1],FF_dict["bonds"][j][qc][2])) # count_i+1 bc of LAMMPS 1-indexing

        # Add one_five ghost bonds
        if one_five_opt == 1:
            f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+2,0.0,2.0)) # dummy bond for 1-5 interactions
        f.write("\n")

        # Write Angle Coeffs
        f.write("Angle Coeffs\n\n")
        for count_i,i in enumerate(set(Angle_types)):
            for j in set(Angle_types):
                if Angle_type_dict[j] == count_i+1:
                    f.write("{} {:< 12.6f} {:< 12.6f}\n".format(count_i+1,FF_dict["angles"][j][qc][1],FF_dict["angles"][j][qc][2])) # count_i+1 bc of LAMMPS 1-indexing
        f.write("\n")

        # Write Atoms
        f.write("Atoms\n\n")
        for count_i,i in enumerate(Atom_types):
            f.write("{:<8d} {:< 4d} {:< 4d} {:< 20.16f} {:< 20.16f} {:< 20.16f} {:< 20.16f} {:d} {:d} {:d}\n"\
            .format(count_i+1,0,Atom_type_dict[i],Charges[count_i],Geometry[count_i,0],Geometry[count_i,1],Geometry[count_i,2],0,0,0))

        # Write Bonds
        bond_count = 1
        if len(Bonds) > 0:
            f.write("\nBonds\n\n")
            for count_i,i in enumerate(Bonds):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d}\n".format(bond_count,Bond_type_dict[Bond_types[count_i]],i[0]+1,i[1]+1))
                bond_count += 1

        # Write 1-5 ghost bonds
        if one_five_opt == 1: 
            for i in One_fives:
                f.write("{:<8d} {:<8d} {:<8d} {:<8d}\n".format(bond_count,Bond_type_dict["one_fives"],i[0]+1,i[-1]+1))
                bond_count += 1

        # Write Angles
        if len(Angles) > 0:
            f.write("\nAngles\n\n")
            for count_i,i in enumerate(Angles):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Angle_type_dict[Angle_types[count_i]],i[0]+1,i[1]+1,i[2]+1))

        # Write Dihedrals
        if len(Dihedrals) > 0: 
            f.write("\nDihedrals\n\n")
            for count_i,i in enumerate(Dihedrals):
                f.write("{:<8d} {:<8d} {:<8d} {:<8d} {:<8d} {:<8d}\n".format(count_i+1,Dihedral_type_dict[Dihedral_types[count_i]],i[0]+1,i[1]+1,i[2]+1,i[3]+1))

    # Write the settings file
    with open(Filename+'.in.settings','w') as f:

        # Write non-bonded interactions (the complicated form of this loop is owed
        # to desire to form a nicely sorted list in terms of the lammps atom types
        # Note: Atom_type_dict was initialize according to sorted(set(Atom_types) 
        #       so iterating over this list (twice) orders the pairs, too.
        f.write("     {}\n".format("# Non-bonded interactions (pair-wise)"))
        for count_i,i in enumerate(sorted(set(Atom_types))):     
            for count_j,j in enumerate(sorted(set(Atom_types))): 

                # Skip duplicates
                if count_j < count_i:
                    continue

                # Conform to LAMMPS i <= j formatting
                if Atom_type_dict[i] <= Atom_type_dict[j]:
                    f.write("     {:20s} {:<10d} {:<10d} ".format("pair_coeff",Atom_type_dict[i],Atom_type_dict[j]))
                else:
                    f.write("     {:20s} {:<10d} {:<10d} ".format("pair_coeff",Atom_type_dict[j],Atom_type_dict[i]))

                # Determine key (ordered by initialize_VDW function such that i > j)
                if i > j:
                    key = (i,j)
                else:
                    key = (j,i)

                # Write the parameters (note: the internal type (e.g., lj) is converted into the corresponding lammps type (e.g., lj/cut/coul/long))
                for k in FF_dict["vdw"][key]:
                    if k == "lj": k = "lj/cut/coul/cut"
                    if type(k) is str:
                        f.write("{:20s} ".format(k))
                    if type(k) is float:
                        f.write("{:< 20.6f} ".format(k))
                f.write("\n")                        

        # Write stretching interactions
        # Note: Bond_type_dict was initialized by looping over sorted(set(Bond_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Stretching interactions"))
        for i in sorted(set(Bond_types)):
            f.write("     {:20s} {:<10d} ".format("bond_coeff",Bond_type_dict[i]))
            for j in FF_dict["bonds"][i][qc][1:]:
                if type(j) is str:
                    f.write("{:20s} ".format(j))
                if type(j) is float:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is np.double:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is int:
                    f.write("{:< 20d} ".format(j))
            f.write("\n")

        # Write 1-5 ghost bond parameter
        if one_five_opt == 1:
            f.write("     {:20s} {:<10d} {:< 20.6f} {:< 20.6f}\n".format("bond_coeff",Bond_type_dict["one_fives"],0.0,2.0))

        # Write bending interactions
        # Note: Angle_type_dict was initialized by looping over sorted(set(Angle_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Bending interactions"))
        for i in sorted(set(Angle_types)):
            f.write("     {:20s} {:<10d} ".format("angle_coeff",Angle_type_dict[i]))
            for j in FF_dict["angles"][i][qc][1:]:
                if type(j) is str:
                    f.write("{:20s} ".format(j))
                if type(j) is float:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is np.double:
                    f.write("{:< 20.6f} ".format(j))
                if type(j) is int:
                    f.write("{:< 20d} ".format(j))
            f.write("\n")

        # Write dihedral interactions
        # Note: Dihedral_type_dict was initialized by looping over sorted(set(Dihedral_types)), so 
        #       iterating over this list resulted in ordered parameters in the in.settings file
        f.write("\n     {}\n".format("# Dihedral interactions"))
        for i in sorted(set(Dihedral_types)):

            # The potential fo the scanned dihedrals are set to zero
            if i in tables:
                f.write("     {:20s} {:<10d} {:20s} {:20s} {:20s}".format("dihedral_coeff",Dihedral_type_dict[i],"table",tables[i][0],tables[i][1]))
            elif i in [ j[0] for j in Scanned_Dihedrals ]:
                f.write("     {:20s} {:<10d} opls 0.0 0.0 0.0 0.0".format("dihedral_coeff",Dihedral_type_dict[i]))
            else:
                f.write("     {:20s} {:<10d} ".format("dihedral_coeff",Dihedral_type_dict[i]))
                for j in FF_dict["dihedrals"][i][qc]:
                    if type(j) is str:
                        f.write("{:20s} ".format(j))
                    if type(j) is float:
                        f.write("{:< 20.6f} ".format(j))
                    if type(j) is np.double:
                        f.write("{:< 20.6f} ".format(j))
                    if type(j) is int:
                        f.write("{:< 20d} ".format(j))
            f.write("\n")

    return Atom_type_dict

# Description: A wrapper for the write commands for generating the lammps input file
def Write_TORSAMPLE_input(Filename,frequency,Onefourscale,runtime,Pair_styles,Dihedral_styles,Temp,tables=[]):

    with open(Filename+'.in.init','w') as f:
        f.write("# lammps input file for polymer simulation with dilute ions\n\n"+\
                "# VARIABLES\n"+\
                "variable        data_name       index   {}\n".format(Filename+'.data')+\
                "variable        settings_name   index   {}\n".format(Filename+'.in.settings')+\
                "variable        coords_freq     index   {}\n".format(frequency)+\
                "variable        vseed           index   {: <6d}\n".format(int(random.random()*100000))+\
                
                "#===========================================================\n"+\
                "# GENERAL PROCEDURES\n"+\
                "#===========================================================\n"+\
                "units		real	# g/mol, angstroms, fs, kcal/mol, K, atm, charge*angstrom\n"+\
                "dimension	3	# 3 dimensional simulation\n"+\
                "newton		off	# use Newton's 3rd law\n"+\
                "boundary	f f f	# periodic boundary conditions \n"+\
                "atom_style	full    # molecular + charge\n\n"+\

                "#===========================================================\n"+\
                "# FORCE FIELD DEFINITION\n"+\
                "#===========================================================\n"+\
                'special_bonds  lj/coul  0.0 0.0 0.0     # NO 1-4 LJ/Coul interactions\n'+\
                'pair_style     hybrid {} # outer_LJ outer_Coul (cutoff values, see LAMMPS Doc)\n'.format(Pair_styles)+\
                'bond_style     harmonic             # parameters needed: k_bond, r0\n'+\
                'angle_style    harmonic             # parameters needed: k_theta, theta0\n')
        if len(Dihedral_styles) >= 1:
            if tables != []:
                f.write('dihedral_style hybrid {} {}         # parameters needed: V1, V2, V3, V4\n'.format(' '.join(Dihedral_styles),"table spline 360"))            
            else:
                f.write('dihedral_style hybrid {}            # parameters needed: V1, V2, V3, V4\n'.format(' '.join(Dihedral_styles)))            
        f.write('pair_modify    shift yes mix arithmetic table 0       # using Lorenz-Berthelot mixing rules\n\n'+\

                "#===========================================================\n"+\
                "# SETUP SIMULATIONS\n"+\
                "#===========================================================\n\n"+\
                "# READ IN COEFFICIENTS/COORDINATES/TOPOLOGY\n"+\
                'read_data ${data_name}\n'+\
                'include ${settings_name}\n\n'+\

                "# SET RUN PARAMETERS\n"+\
                "timestep	1.0		# fs\n"+\
                "run_style	verlet 		# Velocity-Verlet integrator\n"+\
                "neigh_modify every 1 delay 0 check no # More relaxed rebuild criteria can be used\n\n"+\
        
                "# Define Theromstyle variables for print\n"+\
                "thermo_style    custom step temp vol density etotal pe ebond eangle edihed ecoul elong evdwl\n"+\
                "variable        my_pe   equal   pe\n"+\
                "variable        my_ebon equal   ebond\n"+\
                "variable        my_eang equal   eangle\n"+\
                "variable        my_edih equal   edihed\n"+\
                "variable        my_evdw equal   evdwl\n"+\
                "variable        my_ecol equal   ecoul\n\n")

        f.write("#===========================================================\n"+\
                "# RUN DYNAMICS FOR THERMAL SAMPLING\n"+\
                "#===========================================================\n\n")

        R_seed = int(random.random()*1000000.0)        

        # Straight brownian dynamics are run for torsional sampling 
        f.write("# Create dump, fixes, and run the simulation\n"+\
                "dump           {}".format(Filename)+" all atom ${coords_freq}"+" {}.lammpstrj\n".format(Filename)+\
                "dump_modify	{} scale no\n".format(Filename)+\
                "dump_modify	{} sort  id\n".format(Filename)+\
                "velocity all create {} {} mom yes rot yes dist gaussian\n".format(Temp,R_seed)+\
                "fix NVE all nve\n"+\
                "fix TFIX all langevin {} {} 100 {} zero yes\n".format(Temp,Temp,R_seed)+\
                "run {}\n\n".format(int(runtime))+\
                "# undump/unfix\n"+\
                "undump {}\n".format(Filename)+\
                "unfix NVE\n"+\
                "unfix TFIX\n")
    f.close()

    return

# One off function that returns a dictionary of elemnents from the atomtype labels
def update_elements(atomtypes):

    if hasattr(update_elements, 'atomic2ele') is False:
        # Initialize element to atomic number conversion
        atomic2ele = { 1:"H",   2:"He",\
                       3:"Li",  4:"Be",                                                                                                      5: "B",    6:"C",   7: "N",   8: "O",   9: "F",  10:"Ne",\
                      11:"Na", 12:"Mg",                                                                                                     13:"Al",  14:"Si",  15: "P",  16: "S",  17:"Cl",  18:"Ar",\
                      19:"K",  20:"Ca",  21:"Sc",  22:"Ti",  23: "V", 24:"Cr",  25:"Mn",  26:"Fe",  27:"Co",  28:"Ni",  29:"Cu",  30:"Zn",  31:"Ga",  32:"Ge",  33:"As",  34:"Se",  35:"Br",  36:"Kr",\
                      37:"Rb", 38:"Sr",  39: "Y",  40:"Zr",  41:"Nb", 42:"Mo",  43:"Tc",  44:"Ru",  45:"Rh",  46:"Pd",  47:"Ag",  48:"Cd",  49:"In",  50:"Sn",  51:"Sb",  52:"Te",  53: "I",  54:"Xe",\
                      55:"Cs", 56:"Ba",            72:"Hf",  73:"Ta", 74: "W",  75:"Re",  76:"Os",  77:"Ir",  78:"Pt",  79:"Au",  80:"Hg",  81:"Tl",  82:"Pb",  83:"Bi",  84:"Po",  85:"At",  86:"Rn"}

    return { i:atomic2ele[int(i.split('[')[1].split(']')[0])] for i in atomtypes }    

# Create logger to save stdout to logfile
class Logger(object):

    def __init__(self,folder):
        self.terminal = sys.stdout
        self.log = open(folder+"/extract_intramolecular.log", "a")

    def write(self, message):
        self.terminal.write(message)
        self.log.write(message)  

    def flush(self):
        pass

if __name__ == "__main__":
   main(sys.argv[1:])
