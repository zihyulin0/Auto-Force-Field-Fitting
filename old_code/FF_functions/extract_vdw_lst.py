#!/bin/env python                                                                                                                                                              
# Author: Brett Savoie (brettsavoie@gmail.com)

import sys,argparse,os,ast,re,fnmatch,matplotlib,subprocess
matplotlib.use('Agg') # Needed for cluster image generation
from scipy.optimize import leastsq
from scipy.optimize import lsq_linear
from scipy.optimize import minimize
from scipy.optimize import curve_fit
from scipy.optimize import nnls
import scipy.optimize
from scipy.spatial.distance import *
from numpy import *
from numpy.linalg import norm
import numpy.linalg as npl
import numpy as np
from numpy.linalg import matrix_rank
from pylab import *
from copy import deepcopy
import random,time
from matplotlib.offsetbox import AnchoredText
from scipy import stats

# Powerpoint generation dependencies
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches,Pt

def main(argv):

    parser = argparse.ArgumentParser(description='This program collects force field parameters from the output of Orca calculations. '+\
                                                 'This script should be used in conjunction with the gen_jobs_for_vdw.py program, which automatically generates '\
                                                 'the Orca jobs for submission. After the quantum chemistry completes, this script collates the information and '\
                                                 ' generates a database file that can be used in conjunction with polygen to generate the MD job.')

    #optional arguments                                                                                                                                                        
    parser.add_argument('-f', dest='base_name', default='*',
                        help = 'The program operates on all files discovered during a directory walk from the working directory whose name matches this variable.')

    parser.add_argument('-o', dest='output_folder', default='vdw_fits',
                        help = 'All outputs from the fitting are saved to a subfolder of the working directory (determined by the -f argument). '+\
                               'This option determines the folder name. (default: vdw_fits)') 
    
    parser.add_argument('-c', dest='cycles', default='all',
                        help = 'This variable determines which configurations are included in the fit set. The configs folder which holds the sampled configurations can possess cycle labels ( i.e. "cycle_num-config_num" ). The '+\
                               'user can include all configurations by supplying "all" (the default), or supply individual cycle labels as a space-delimited string (i.e. "1 3 5") or use a colon syntax for the start:end of the '+\
                               'cycles to include (i.e. "1:4" would included cycles 1 through 4, inclusive). The last two options can be used in combination. (default: "all")')
    parser.add_argument('-method', dest = 'method_list', default='global',
                        help = 'Specified the method used to fit the VDW potential (options: global, lstsq,Boltzmann,Boltzmann_tot,linear,quadratic,; default: global)')
                        
    parser.add_argument('-fit', dest='fit', default="lj",
                        help = 'Specifies the function used to fit the VDW potential (options: buck, lj, born ; default: lj)')
    
    parser.add_argument('-max_cycles', dest='max_cycles',default=10000,
                        help = 'Specifies the maximum number of iterations to perform before terminating the fit. (default: 1000)')

    parser.add_argument('-min_cycles', dest='min_cycles', default=10,
                        help = 'Specifies the minimum number of iterations to perform before terminating the fit. (default: 100)')

    parser.add_argument('-xhi2_thresh', dest='xhi2_thresh', default=1E-6,
                        help = 'Specifies the termination threshold for differences in xhi2 between cycles. '+\
                               'Only breaks once min_cycles of iterations have been performed (kcal2/mol2; default: 1E-6)')

    parser.add_argument('-weight', dest='weight', default=1.0,
                        help = 'Specifies the weighting factor for updating the VDW parameters at each iteration. For example, '+\
                               '1.0 means the old value is completely placed by the new value, and 0.5 means the updated value '+\
                               'is the average of the new and old. (default: 1.0)')

    parser.add_argument('-N_config', dest='N_config', default=0,
                        help = 'Specifies the number of configurations to use for the fit. By default (0) all discovered configurations '+\
                               'are used. When set to a value a random subset of the discovered configurations is used. This option is '+\
                               'intended to be used for convergence tests with respect to the number of configurations. If N_config is set '+\
                               'to a value larger than the number of discovered configurations, then a warning is printed and the program '+\
                               'defaults to using all configurations. (default: 0, all configurations)') 

    parser.add_argument('-FF_DFT', dest='DFT_db', default=[],
                        help = 'This variable holds the filename of a DFT parameters database. When supplied, the program will avoid parametrizing any '+\
                               'vdw pairs at the DFT level that already exist in the database (default: None)')

    parser.add_argument('-FF_MP2', dest='MP2_db', default=[],
                        help = 'This variable holds the filename of a MP2 parameters database. When supplied, the program will avoid parametrizing any '+\
                               'vdw pairs at the MP2 level that already exist in the database (default: None)')

    parser.add_argument('-E_max', dest='E_max', default=10.0,
                        help = 'This variable holds the maximum energy threshold for using a configuration in the fitting procedure. For example, settings '+\
                               'E_max to 10.0 will only use configurations whose interaction energy is less than 10.0 kcal/mol in the fit. (default: 10.0)')

    parser.add_argument('-E_min', dest='E_min', default=-10000.0,
                        help = 'This variable holds the minimum energy threshold for using a configuration in the fitting procedure. For example, settings '+\
                               'E_max to 0.0 will only use configurations whose interaction energy is greater than 0.0 kcal/mol in the fit. This isn\'t usually '+\
                               'useful since the low energy portion of the interaction potential is most relevant to the fit; it is included for completeness (default: -10000.0)')

    parser.add_argument('-seed', dest='seed', default=444,
                        help = 'The seed for the random number generator. (default: 444)')

    parser.add_argument('-q_a', dest='q_a', default=None,
                        help = 'The molecular charge on molecule a. By default the charge is automatically determined by rounding the total charge to the nearest integer. (default: None)')

    parser.add_argument('-q_b', dest='q_b', default=None,
                        help = 'The molecular charge on molecule b. By default the charge is automatically determined by rounding the total charge to the nearest integer. (default: None)')

    parser.add_argument('-r_min_scale', dest='r_min_scale', default=1.0,
                        help = 'Scales the r_min value for the VDW fit. By default, the minima of fit potetials are retrained to lie within the sampled region of pair separations. '+\
                               'This behavior can be modified by setting this variable to <1 (more relaxed) or >1 (forces larger r_min). (default: 1.0)')

    parser.add_argument('-UA_vol_scale', dest='UA_vol_scale', default=1.0,
                        help = 'Scale factor for UA multiplying the interaction volume of the UA pseudo-atoms. (default: 1.0)')

    parser.add_argument('-mixing_rule', dest='mixing_rule', default="None",
                        help = 'Mixing rule for heteroatomic lennard-jones interactions. Valid options are none (no mixing rules imposed), lb (Lorentz-Berthelot), and wh (Waldman-Hagler). The options are '+\
                               'case insensitive (default: None)')

    parser.add_argument('-L2_sigma', dest='L2_s', default=0.1,
                        help = 'Sets the weight of L2 regularization for sigma values in the fit penalty function. When set to zero, the LJ parameters are fit using least-squares criteria, '+\
                               'non-zero values result in L2 regularization of the sigma values (default: 0.1)')

    parser.add_argument('-L2_eps', dest='L2_e', default=0.01,
                        help = 'Sets the weight of L2 regularization for the epsilon values in the fit penalty function. When set to zero, the LJ parameters are fit using least-squares criteria, '+\
                               'non-zero values result in L2 regularization of the epsilon values (default: 0.01)')                               

    parser.add_argument('-UA_method', dest='UA_method', default='AA_fit',
                        help = 'Determines the method for calculating the UA values. '+\
                               '"radius" uses the volume of the sphere that encompasses all UA atoms in the group, scaled by a factor determined by the -UA_vol_scale variable. '+\
                               '"volume" uses the monte-carlo calculated volume of the UA group, with overlaps naturally subtracted. '+\
                               'Both, radius and volume methods base the UA eps values upon scaled AA-eps values in order to reproduce the intermolecular '+\
                               'LJ interaction energy. "fit" performs an iterative fit of the eps and sigma values to reproduce the DFT data, independent of the AA '+\
                               'parameters. radius and volume are meant for use fitting self interactions (where A and B molecules are identical), the fit method is meant '+\
                               'for fitting heteromolecular interactions (e.g., ion polymer or solute solvent) (default: radius)')

    parser.add_argument('--extract_charges', dest='extract_charges', default=0, const=1, action='store_const',
                        help = 'If this flag is enabled then the charges are parsed and averaged over all QC calculations (default: off)')

    parser.add_argument('--avoid_read', dest='avoid_read', default=0, const=1, action='store_const',
                        help = 'If this flag is enabled then the force-field potentials found in any present "initial_params.db" files are not used. (default: use initial_params.db to initialize LJ guess)')
    parser.add_argument('--remove_outlier', dest='outlier_option', default=False, const=True, action='store_const',
                        help = 'If this flag is enabled then fitting will remove outliers (remove modified z-score>3.5). (default: do not remove outlier)')


    # Declare global variables and save the working directory to variable
    global Htokcalmol,Ecoul_Const,Data,min_pairs
    working_dir = os.getcwd()       # Not sure if this is useful anymore
    Htokcalmol = 627.509            # Converts hartree to kcal/mol
    Ecoul_Const = 332.0702108431797 # Converts elementary_charge^2/ang to kcal/mol 

    # Convert inputs to the proper data types
    args=parser.parse_args(argv)
    args.base_name = str(args.base_name)
    args.output_folder = str(args.output_folder)    
    args.fit = str(args.fit).lower()
    args.mixing_rule = str(args.mixing_rule).lower()
    args.max_cycles = int(args.max_cycles)
    args.min_cycles = int(args.min_cycles)
    args.xhi2_thresh = float(args.xhi2_thresh)
    args.weight = float(args.weight)
    args.N_config = int(args.N_config)
    args.E_max = float(args.E_max)
    args.E_min = float(args.E_min)
    args.UA_vol_scale = float(args.UA_vol_scale)
    args.r_min_scale = float(args.r_min_scale)
    args.L2_s = float(args.L2_s)
    args.L2_e = float(args.L2_e)
    args.method_list = args.method_list.split()
    if args.DFT_db != []:
        args.DFT_db = args.DFT_db.split()
    if args.MP2_db != []:
        args.MP2_db = args.MP2_db.split()    

    # Consistency checks
    if args.fit not in [ 'lj', 'buck', 'born' ]: print("ERROR: an unrecognized fit type was requested by the user. Exiting..."); quit()
    if args.UA_method not in [ 'radius', 'volume', 'AA_fit', 'fit' ]: print("ERROR: an unrecognized -UA_method was requested by the user. Exiting..."); quit()
    if args.mixing_rule not in ['none','lb','wh']: print("ERROR: an unrecognized -mixing_rule was requested by the user. Exiting..."); quit()
    if args.min_cycles > args.max_cycles: print("ERROR: min_cycles must be less than or equal to max_cycles. Exiting..."); quit()
    if args.L2_s < 0.0: print("ERROR: -L2_sigma must be set to a float greater than 0. Exiting..."); quit()
    if args.L2_e < 0.0: print("ERROR: -L2_eps must be set to a float greater than 0. Exiting..."); quit()
    
    # Parse the cycles input. Additional complexity is due to the use of a colon minilanguage that allows the user to supply supply a range cycles.
    args.cycles = args.cycles.split()
    if "all" in args.cycles: args.cycles = ["*"]
    else:
        expanded = [ i for i in args.cycles if ":" in i ]
        args.cycles = [ i for i in args.cycles if i not in expanded ]   
        added = []
        # Expand the colon containing arguments
        for count_i,i in enumerate(expanded): added += [ str(j) for j in range(int(i.split(":")[0]),int(i.split(":")[1])+1) ]
        args.cycles += added
        args.cycles = [ i+'-*' for i in args.cycles ]

    # Use a fixed seed value to ensure reproducibility between runs
    random.seed(int(args.seed))

    # Check that the job directory exists
    if os.path.isdir(args.base_name) == False:
        print("ERROR: Requested directory could not be found. Please check the name. Exiting...")
        quit()
    else:
        os.chdir(args.base_name)
        if os.path.isdir(args.output_folder) == True:
            print("ERROR: Requested output folder ({}/{}) already exists. Exiting to avoid overwriting data...".format(args.base_name,args.output_folder))
            quit()
        else:
            os.makedirs(args.output_folder)
            os.makedirs(args.output_folder+'/figures')
            for i in args.method_list:
               os.makedirs(args.output_folder+'/'+i)
            sys.stdout = Logger(args.output_folder)
            print("PROGRAM CALL: python extract_vdw.py {}\n".format(' '.join([ i for i in argv])))
        if args.DFT_db != []:
            args.DFT_db = [ '../'+i if i[0] != '/' else i for i in args.DFT_db ]
        if args.MP2_db != []:
            args.MP2_db = [ '../'+i if i[0] != '/' else i for i in args.MP2_db ]

    # Find all output files
    Base_Filename = '*out'
    Base_dir = os.getcwd()

    # Use fnmatch to use wildcards and os.walk to perform a directory walk. Names holds the path and filename of each 
    # file to be processed by the program. Name contains the unique thermo filenames that are parsed by the program. 
    # The program expects each trajectory being sampled to have the same number of thermo.avg files, indexed according
    # to the lambda step in the thermodynamic integration simulation. 
    Names = [os.path.join(dp,f) for dp, dn, filenames in os.walk('.') for f in filenames if fnmatch.fnmatch(f,Base_Filename) ]
    Names = keep_complete(Names)
    if len(Names) == 0: print("Sorry, no complete vdw jobs were discovered in the specified directory. Exiting..."); quit()
    Names = [ i for i in Names if True in [ fnmatch.fnmatch(i.split('/')[-1],j) for j in args.cycles ] ] # Only keep files that belong to cycles specifid by the args.cycles argument

    #################
    # PARSE QC Data #
    #################

    # Print header
    print("{}".format("*"*144))
    print("* {:^140s} *".format("Parsing QC Interaction Energies, Charges, and Configurations in {}".format(args.base_name)))
    print("{}".format("*"*144))

    # Returns a dictionary keyed to the filenames, with the information needed for the fits keyed
    # to subdictionaries. See the get_data function for more info on the keys that are populated.
    # NOTE: this is function parses both DFT and MP2 data. The subsequent fit function only makes use
    #       of a subset of the data depending on whether the MP2 or DFT information is being used. 
    Data, job_type = get_data(Names,args.extract_charges,args.DFT_db,q_a=args.q_a,q_b=args.q_b)

    ########################
    # PARSE DFT PARAMETERS #
    ########################

    # If DFT data was discovered in the QC output files then the DFT vdw parameters are parsed
    if "dft" in job_type:

        # Print header
        print("\n{}".format("*"*144))
        print("* {:^140s} *".format("Parsing DFT-Interaction Energy Calculations for Configurations within {}".format(args.base_name)))
        print("{}".format("*"*144))

        # Remove entries based on E_min/E_max criteria
        for i in list(Data.keys()):
            if Data[i]["e_dft"] > args.E_max: Data.pop(i,None)
            elif Data[i]["e_dft"] < args.E_min: Data.pop(i,None)

        # Remove entries if they differ by more than 2 sigma from the mean 
#         mean = average(array([ Data[i]["e_dft_fit_AA"] for i in Data.keys() ]))
#         stdev = (average(array([ (Data[i]["e_dft_fit_AA"]-mean)**2 for i in Data.keys() ])))**(0.5)
#         for i in Data.keys():
#             if abs(Data[i]["e_dft_fit_AA"]-mean) > 2.0*stdev: Data.pop(i,None)

#       # Loop for removing entries based on e_dft value rather than fit value
#         mean = average(array([ Data[i]["e_dft"] for i in Data.keys() ]))
#         stdev = (average(array([ (Data[i]["e_dft"]-mean)**2 for i in Data.keys() ])))**(0.5)
#         for i in Data.keys():
#             if abs(Data[i]["e_dft"]-mean) > 1.0*stdev: Data.pop(i,None)


#        print "mean: {}".format(mean)
#        print "stdev: {}".format(stdev)
        
        # Print diagnostic
        print("\n{:70s} {}".format("Number of Configurations Discovered:",len(Names)))    
        print("{:70s} {}".format("Number of Configurations Satisfying E_min/E_max criteria:",len(list(Data.keys()))))

        # Check that at least 1 configuration has been retained based on the min/max criteria
        if len(list(Data.keys())) == 0:
            print("\nERROR: No configurations were discovered that matched the E_max criteria ({:<2.2f} kcal/mol). Exiting...".format(args.E_max))
            quit()

        # Print a warning if more configurations were requested than were discovered.
        if args.N_config > len(list(Data.keys())):
            print("\n                                 ****************************************************************************")
            print("                                 * WARNING: Less configurations were discovered than requested by N_config. *")
            print("                                 *          Continuing with fit using all configurations.                   *")
            print("                                 ****************************************************************************\n")

        # Use the number of configurations requested by user (by default, args.N_configs=0, all are used)
        if args.N_config != 0 and args.N_config <= len(list(Data.keys())):
            keys = list(Data.keys())
            random.shuffle(keys)
            for i in keys[args.N_config:]:
                Data.pop(i,None)

        print("{:70s} {}".format("Number of Configurations Used for Fit:",len(list(Data.keys()))))

        # Parse the unique molecules types in the simulation
        parse_molecules(args.output_folder,Data)

        # Fit all possible pair-wise interactions
        Fits = Fit_pairs(args.output_folder,Data,args.fit,args.weight,args.xhi2_thresh,args.min_cycles,args.max_cycles,args.DFT_db,args.method_list,\
                         QC_type="DFT",charges_opt=args.extract_charges,UA_vol_scale=args.UA_vol_scale,UA_method=args.UA_method,r_min_scale=args.r_min_scale,\
                         avoid_read_flag=args.avoid_read,mixing_rule=args.mixing_rule,L2_s=args.L2_s,L2_e=args.L2_e,outlier_option=args.outlier_option)

    ########################
    # PARSE MP2 PARAMETERS #
    ########################

    # Returns a dictionary keyed to the filenames, with the information needed for the fits keyed
    # to subdictionaries. See the get_data function for more info on the keys that are populated.
    # NOTE: this is function parses both DFT and MP2 data. The subsequent fit function only makes use
    #       of a subset of the data depending on whether the MP2 or DFT information is being used. 
#    Data, job_type = get_data(Names,args.extract_charges)

    # If MP2 data was discovered in the QC output files then the MP2 vdw parameters are parsed
    if "mp2" in job_type:

        # Print header
        print("{}".format("*"*144))
        print("* {:^140s} *".format("Parsing MP2-Interaction Energy Calculations for Configurations within {}".format(args.base_name)))
        print("{}".format("*"*144))

        # Remove entries based on E_min/E_max criteria
        for i in list(Data.keys()):
            if Data[i]["e_mp2"] > args.E_max: Data.pop(i,None)
            elif Data[i]["e_dft"] < args.E_min: Data.pop(i,None)

        # Print diagnostic
        print("\nNumber of Configurations Discovered: {}".format(len(Names)))    
        print("Number of Configurations Satisfying E_min/E_max criteria: {}".format(len(list(Data.keys()))))

        # Print a warning if more configurations were requested than were discovered.
        if args.N_config > len(list(Data.keys())):
            print("\n                                 ****************************************************************************")
            print("                                 * WARNING: Less configurations were discovered than requested by N_config. *")
            print("                                 *          Continuing with fit using all configurations.                   *")
            print("                                 ****************************************************************************\n")

        # Use the number of configurations requested by user (by default, args.N_configs=0, all are used)
        if args.N_config != 0 and args.N_config <= len(list(Data.keys())):
            keys = list(Data.keys())
            random.shuffle(keys)
            for i in keys[args.N_config:]:
                Data.pop(i,None)
        print("Number of Configurations Used for Fit: {}".format(len(list(Data.keys()))))

        # Parse the unique molecules types in the simulation
        parse_molecules(args.output_folder,Data)

        # Fit all possible pair-wise interactions
        Fits = Fit_pairs(args.output_folder,Data,args.fit,args.weight,args.xhi2_thresh,args.min_cycles,args.max_cycles,args.MP2_db,\
                         QC_type="MP2",charges_opt=args.extract_charges,UA_vol_scale=args.UA_vol_scale,UA_method=args.UA_method,r_min_scale=args.r_min_scale,mixing_rule=args.mixing_rule)

        # Fit all possible pair-wise interactions (old call; the new one hasn't been tested with MP2 yet)
#        Fits = Fit_pairs(args.output_folder,Data,args.fit,args.weight,args.xhi2_thresh,args.min_cycles,args.max_cycles,args.MP2_db,QC_type="MP2")

    print("\n{}".format("*"*144))
    print("* {:^140s} *".format("VDW Parameter Parse Complete!"))
    print("{}".format("*"*144))

    quit()

def find_orientation(Geometry,Elements):
    
    # sum all unique cross-products over the geometry
    o_vec = array([0,0,0])
    for count_i,i in enumerate(Geometry):
        for count_j,j in enumerate(Geometry):
            if count_i < count_j:
                o_vec += cross(i,j)
    o_vec = o_vec/norm(o_vec)

    # find the unique long axis of the molecule by 
    # summing atomic vectorial differences.
    l_vec = array([0,0,0])
    for count_i,i in enumerate(Geometry):
        for count_j,j in enumerate(Geometry):
            if count_i < count_j and global_atomtypes[count_i] in keep_types and global_atomtypes[count_j] in keep_types:
                if norm(l_vec+i-j) > norm(l_vec):
                    l_vec += i-j
                else: 
                    l_vec -= i-j
    l_vec = l_vec/norm(l_vec)

    return o_vec,l_vec
                    
# Description:
# Rotate Point by an angle, theta, about the vector with an orientation of v1 passing through v2. 
# Performs counter-clockwise rotations (i.e., if the direction vector were pointing
# at the spectator, the rotations would appear counter-clockwise)
# For example, a 90 degree rotation of a 0,0,1 about the canonical 
# y-axis results in 1,0,0.
#
# Point: 1x3 array, coordinates to be rotated
# v1: 1x3 array, the direction vector
# v2: 1x3 array, point the  direction vector passes through
# theta: scalar, magnitude of the rotation (defined by default in degrees)
def axis_rot(Point,v1,v2,theta,mode='angle'):

    # Temporary variable for performing the transformation
    rotated=array([Point[0],Point[1],Point[2]])

    # If mode is set to 'angle' then theta needs to be converted to radians to be compatible with the
    # definition of the rotation vectors
    if mode == 'angle':
        theta = theta*pi/180.0

    # Rotation carried out using formulae defined here (11/22/13) http://inside.mines.edu/fs_home/gmurray/ArbitraryAxisRotation/)
    # Adapted for the assumption that v1 is the direction vector and v2 is a point that v1 passes through
    a = v2[0]
    b = v2[1]
    c = v2[2]
    u = v1[0]
    v = v1[1]
    w = v1[2]
    L = u**2 + v**2 + w**2

    # Rotate Point
    x=rotated[0]
    y=rotated[1]
    z=rotated[2]

    # x-transformation
    rotated[0] = ( a * ( v**2 + w**2 ) - u*(b*v + c*w - u*x - v*y - w*z) )\
             * ( 1.0 - cos(theta) ) + L*x*cos(theta) + L**(0.5)*( -c*v + b*w - w*y + v*z )*sin(theta)

    # y-transformation
    rotated[1] = ( b * ( u**2 + w**2 ) - v*(a*u + c*w - u*x - v*y - w*z) )\
             * ( 1.0 - cos(theta) ) + L*y*cos(theta) + L**(0.5)*(  c*u - a*w + w*x - u*z )*sin(theta)

    # z-transformation
    rotated[2] = ( c * ( u**2 + v**2 ) - w*(a*u + b*v - u*x - v*y - w*z) )\
             * ( 1.0 - cos(theta) ) + L*z*cos(theta) + L**(0.5)*( -b*u + a*v - v*x + u*y )*sin(theta)

    rotated = rotated/L
    return rotated

# Generates the adjacency matrix based on UFF bond radii
def Table_generator(Elements,Geometry):

    # Initialize UFF bond radii (Rappe et al. JACS 1992)
    # NOTE: Units of angstroms 
    # NOTE: These radii neglect the bond-order and electronegativity corrections in the original paper. Where several values exist for the same atom, the largest was used. 
    Radii = {  'H':0.354, 'He':0.849,\
              'Li':1.336, 'Be':1.074,                                                                                                                          'B':0.838,  'C':0.757,  'N':0.700,  'O':0.658,  'F':0.668, 'Ne':0.920,\
              'Na':1.539, 'Mg':1.421,                                                                                                                         'Al':1.244, 'Si':1.117,  'P':1.117,  'S':1.064, 'Cl':1.044, 'Ar':1.032,\
               'K':1.953, 'Ca':1.761, 'Sc':1.513, 'Ti':1.412,  'V':1.402, 'Cr':1.345, 'Mn':1.382, 'Fe':1.335, 'Co':1.241, 'Ni':1.164, 'Cu':1.302, 'Zn':1.193, 'Ga':1.260, 'Ge':1.197, 'As':1.211, 'Se':1.190, 'Br':1.192, 'Kr':1.147,\
              'Rb':2.260, 'Sr':2.052,  'Y':1.698, 'Zr':1.564, 'Nb':1.473, 'Mo':1.484, 'Tc':1.322, 'Ru':1.478, 'Rh':1.332, 'Pd':1.338, 'Ag':1.386, 'Cd':1.403, 'In':1.459, 'Sn':1.398, 'Sb':1.407, 'Te':1.386,  'I':1.382, 'Xe':1.267,\
              'Cs':2.570, 'Ba':2.277, 'La':1.943, 'Hf':1.611, 'Ta':1.511,  'W':1.526, 'Re':1.372, 'Os':1.372, 'Ir':1.371, 'Pt':1.364, 'Au':1.262, 'Hg':1.340, 'Tl':1.518, 'Pb':1.459, 'Bi':1.512, 'Po':1.500, 'At':1.545, 'Rn':1.42,\
              'default' : 0.7 }
    
    # Scale factor is used for determining the bonding threshold. 1.2 is a heuristic that give some lattitude in defining bonds since the UFF radii correspond to equilibrium lengths. 
    scale_factor = 1.2

    # Print warning for uncoded elements.
    for i in Elements:
        if i not in list(Radii.keys()):
            print("ERROR: The geometry contains an element ({}) that the Table_generator function doesn't have bonding information for. This needs to be directly added to the Radii".format(i)+\
                  " dictionary before proceeding. Exiting...")
            quit()


    # Generate distance matrix holding atom-atom separations (only save upper right)
    Dist_Mat = triu(cdist(Geometry,Geometry))
    
    # Find plausible connections
    x_ind,y_ind = where( (Dist_Mat > 0.0) & (Dist_Mat < max([ Radii[i]**2.0 for i in list(Radii.keys()) ])) )

    # Initialize Adjacency Matrix
    Adj_mat = zeros([len(Geometry),len(Geometry)])

    # Iterate over plausible connections and determine actual connections
    for count,i in enumerate(x_ind):
        
        # Assign connection if the ij separation is less than the UFF-sigma value times the scaling factor
        if Dist_Mat[i,y_ind[count]] < (Radii[Elements[i]]+Radii[Elements[y_ind[count]]])*scale_factor:            
            Adj_mat[i,y_ind[count]]=1

    # Hermitize Adj_mat
    Adj_mat=Adj_mat + Adj_mat.transpose()

    # Perform some simple checks on bonding to catch errors
    problem_dict = { i:0 for i in list(Radii.keys()) }
    conditions = { "H":1, "C":4, "F":1, "Cl":1, "Br":1, "I":1, "O":2, "N":4, "B":4 }
    for count_i,i in enumerate(Adj_mat):

        if Elements[count_i] in list(conditions.keys()):
            if sum(i) > conditions[Elements[count_i]]:
                problem_dict[Elements[count_i]] += 1

    # Print warning messages for obviously suspicious bonding motifs.
    if sum( [ problem_dict[i] for i in list(problem_dict.keys()) ] ) > 0:
        print("Table Generation Warnings:")
        for i in sorted(problem_dict.keys()):
            if problem_dict[i] > 0:
                if i == "H": print("WARNING: {} hydrogen(s) have more than one bond.".format(problem_dict[i]))
                if i == "C": print("WARNING: {} carbon(s) have more than four bonds.".format(problem_dict[i]))
                if i == "F": print("WARNING: {} fluorine(s) have more than one bond.".format(problem_dict[i]))
                if i == "Cl": print("WARNING: {} chlorine(s) have more than one bond.".format(problem_dict[i]))
                if i == "Br": print("WARNING: {} bromine(s) have more than one bond.".format(problem_dict[i]))
                if i == "I": print("WARNING: {} iodine(s) have more than one bond.".format(problem_dict[i]))
                if i == "O": print("WARNING: {} oxygen(s) have more than two bonds.".format(problem_dict[i]))
                if i == "N": print("WARNING: {} nitrogen(s) have more than four bonds.".format(problem_dict[i]))
                if i == "B": print("WARNING: {} bromine(s) have more than four bonds.".format(problem_dict[i]))
        print("")

    return Adj_mat


def tryint(s):
    try:
        return int(s)
    except:
        return s
     
def alphanum_key(s):
    """ Turn a string into a list of string and number chunks.
        "z23a" -> ["z", 23, "a"]
    """
    return [ tryint(c) for c in re.split('([0-9]+)', s) ]

def sort_nicely(l):
    """ Sort the given list in the way that humans expect.
    """
    l.sort(key=alphanum_key)        

def E_LG(r,epsilon,sigma):
    return 4.0*epsilon*( (sigma/r)**12 - (sigma/r)**6 )

def E_Coul(r,q_1,q_2):
    return q_1*q_2/(r)*Ecoul_Const

##############################
# Slide generation functions #
##############################
# Adds a slide with title and default blue line
# returns the slide handle
def add_slide(pres,title):

    blank_slide = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide)
    shapes = slide.shapes

    # Create title 
    txBox = slide.shapes.add_textbox(Inches(0.25),Inches(0.1),Inches(10.0),Inches(1.0))
    text_frame = txBox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    font = run.font
    font.name = 'Helvetica'
    font.size = Pt(40)
    font.bold = True

    # Create line 
    left = Inches(1.45)
    top = Inches(-2.5)
    width = Inches(7.1)  # lines are instantiated as diagonal lines. The slide width is 10 in. so cos(45)*10 are the necessary width and height.
    height = Inches(7.1) # lines are instantiated as diagonal lines. The slide width is 10 in. so cos(45)*10 are the necessary width and height.
    shape = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left, top, width, height)
    shape.rotation = 45.0
    shape.line.color.rgb = RGBColor(0,50,230)
    shape.line.width = Pt(3.0)

    return slide

# Adds a 2x2 array of pictures to a slide.
# Assumes all images have a 4x3 width to height ratio
def add_4(slide,names,h_pad=0.0,v_pad=0.0):
        
    top = 1.5 # Distance of the 2x2 array from the top of the slide in inches.
    left_start = 1 # Distance of the 2x2 array from the left edge of the slide in inches.
    top_start = 1.25
    counter = 0
    for i in range(2):
        for j in range(2):
            try:
                height = Inches(3)
                left = Inches(left_start+(4+h_pad)*j)
                top = Inches(top_start+(3+v_pad)*i)
                pic = slide.shapes.add_picture(names[counter], left, top, height=height)
                counter += 1
            except:
               break

# Adds a 3x3 array of pictures to a slide.
# Assumes all images have a 4x3 width to height ratio
def add_9(slide,names,h_pad=0.0,v_pad=0.0):
        
    top = 1.5 # Distance of the 3x3 array from the top of the slide in inches.
    left_start = 1 # Distance of the 3x3 array from the left edge of the slide in inches.
    top_start = 1.25
    h_pad = 0.0
    v_pad = 0.0
    counter = 0
    for i in range(3):
        for j in range(3):
            try:
                height = Inches(2)
                left = Inches(left_start+(2.67+h_pad)*j)
                top = Inches(top_start+(2+v_pad)*i)
                pic = slide.shapes.add_picture(names[counter], left, top, height=height)
                counter += 1
            except:
               break

# Scrapes the DFT-D and BSSE-corrected MP2 interaction energies from the discovered orca outputs
# Inputs:      Names      : a list of relative filenames from the working directory
# Returns:     Data       : a dictionary with keys for each filename in "Names", and subkeys holding
#                           each molecule's geometry (geo_a, geo_b), list of types (types_a,types_b)
#                           charge (charges_a, charges_b), and dft-d and BSSE-corrected MP2 interaction
#                           energies (e_dft,e_mp2)
# Dict_info:  "geo_a"     : array holding the geometry of the first molecule
#             "geo_b"     : array holding the geometry of the second molecule
#             "types_a"   : array holding the atom types of the first molecule
#             "types_b"   : array holding the atom types of the second molecule
#             "charges_a" : array holding the atomic partial charges of the first molecule
#             "charges_b" : array holding the atomic partial charges of the second molecule
#             "e_dft"     : array holding the DFT-D interaction energies for each configuration
#             "e_MP2"     : array holding the counterpoise-corrected MP2 interaction energies for each configuration
#             "e_dft_fit" : array holding the DFT-D interaction energies less the electrostatic contribution for each configuration
#             "e_MP2_fit" : array holding the counterpoise-corrected MP2 interaction energies less the electrostatic contribution for each configuration
#             "r_dist"    : array holding the pair-separations for each configurations (rows are indexed to the atoms in geo_a and columns to the atoms in geo_b).
#
def get_data(Names,parse_charges=1,db_files=[],q_a=None,q_b=None):

    # Initialize dictionary for all configurations
    global Data
    Data = {}

    # Initialize dictionary to hold charges being read from the database file
    FF_charges = {}
    if db_files != []:
        FF_charges = grab_db_charges(db_files)

    if parse_charges == 1:

        # Initialize a list of unique atom types
        unique_types = []

        # Iterate over all input files and parse atom types and charges for each molecule
        for count_i,i in enumerate(Names):

            # Initialize dictionary for this configuration
            Data[i] = { 'geo_a': array([]), 'geo_b': array([]), 'types_a': [], 'types_b': [], 'charges_a':array([]), 'charges_b':array([]), 'e_dft':0.0, 'e_mp2':0.0, 
                        'e_dft_fit_AA':0.0, 'e_mp2_fit_AA':0.0, 'e_dft_fit_UA':0.0, 'e_mp2_fit_UA':0.0, 'charges_a_UA':array([]), 'charges_b_UA':array([]) }

            # Grab geometry and datatypes from xyz file
            XYZ_file = '.'.join([ j for j in i.split('.')[:-1] ])+'.xyz'
            Data[i]['elem_a'],Data[i]['elem_b'],Data[i]['geo_a'],Data[i]['geo_b'],Data[i]['types_a'],Data[i]['types_b'] = scrape_xyz(XYZ_file,parse_charges=0)

            # Add atom types to unique list
            for j in Data[i]['types_a']:
                if j not in unique_types: unique_types += [j]
            for j in Data[i]['types_b']:
                if j not in unique_types: unique_types += [j]
                    
            # Get charges from output files
            Data[i]['charges_a'],Data[i]['charges_b'] = scrape_charges(i,Data[i]['types_a'],Data[i]['types_b'])

#             # Overwrite values for molecules already in the FF_charges dictionary
#             if False not in [ j in FF_charges.keys() for j in Data[i]['types_a'] ]:
#                 for count_j,j in enumerate(Data[i]['types_a']):
#                     Data[i]['charges_a'][count_j] = FF_charges[j]
#             if False not in [ j in FF_charges.keys() for j in Data[i]['types_b'] ]:
#                 for count_j,j in enumerate(Data[i]['types_b']):
#                     Data[i]['charges_b'][count_j] = FF_charges[j]

#             # Balance the charges in atoms not being read from the database file to
#             # ensure an integer number of electrons on each molecule
#             residual   = sum(Data[i]['charges_a']) - round(sum(Data[i]['charges_a']))
#             for count_j,j in enumerate(Data[i]['types_a']):
#                 Data[i]['charges_a'][count_j] -= residual/float(len(Data[i]['types_a']))
#             residual   = sum(Data[i]['charges_b']) - round(sum(Data[i]['charges_b']))
#             for count_j,j in enumerate(Data[i]['types_b']):
#                 Data[i]['charges_b'][count_j] -= residual/float(len(Data[i]['types_b']))
    
        # Calculate average of each charge type over all configurations
        charges = zeros(len(unique_types))
        counts  = zeros(len(unique_types))
        for count_i,i in enumerate(unique_types):
            for count_j,j in enumerate(Names):
                for count_k,k in enumerate(Data[j]["types_a"]):
                    if i == k:
                        charges[count_i] += Data[j]["charges_a"][count_k]
                        counts[count_i]  += 1
                for count_k,k in enumerate(Data[j]["types_b"]):
                    if i == k:
                        charges[count_i] += Data[j]["charges_b"][count_k]
                        counts[count_i]  += 1
                        
        # Perform average 
        charges = charges/counts
        
        # Assign charges to each sub-dictionary
        for count_i,i in enumerate(unique_types):
            for count_j,j in enumerate(Names):
                for count_k,k in enumerate(Data[j]["types_a"]):
                    if i == k:
                        Data[j]["charges_a"][count_k] = charges[count_i]
                for count_k,k in enumerate(Data[j]["types_b"]):
                    if i == k:
                        Data[j]["charges_b"][count_k] = charges[count_i]
                
    # Parse configurational energies
    flag_a_AA_charges = 0
    flag_b_AA_charges = 0
    flag_a_UA_charges = 0
    flag_b_UA_charges = 0
    residual_flag_a = 0
    residual_flag_b = 0
    for count_i,i in enumerate(Names):

        # Initialize dictionary for this configuration
        if parse_charges == 0:
            Data[i] = { 'geo_a': array([]), 'geo_b': array([]), 'types_a': [], 'types_b': [], 'e_dft':0.0, 'e_mp2':0.0, 
                        'e_dft_fit_AA':0.0, 'e_mp2_fit_AA':0.0, 'e_dft_fit_UA':0.0, 'e_mp2_fit_UA':0.0, 'charges_a_UA':array([]), 'charges_b_UA':array([]) }

        # Grab geometry and datatypes from xyz file
        XYZ_file = '.'.join([ j for j in i.split('.')[:-1] ])+'.xyz'
        if parse_charges == 1:
            Data[i]['elem_a'],Data[i]['elem_b'],Data[i]['geo_a'],Data[i]['geo_b'],Data[i]['types_a'],Data[i]['types_b'] = scrape_xyz(XYZ_file,parse_charges=0)
        else:
            Data[i]['elem_a'],Data[i]['elem_b'],Data[i]['geo_a'],Data[i]['geo_b'],Data[i]['types_a'],Data[i]['types_b'],Data[i]['charges_a'],Data[i]['charges_b'] = scrape_xyz(XYZ_file,parse_charges=1)

        # Overwrite values for molecules already in the FF_charges dictionary (the charges are only assigned if all of the charges in the molecule are present)
        if False not in [ j in list(FF_charges.keys()) for j in Data[i]['types_a'] ]:
            flag_a_AA_charges = 1
            for count_j,j in enumerate(Data[i]['types_a']):
                Data[i]['charges_a'][count_j] = FF_charges[j]
        if False not in [ j in list(FF_charges.keys()) for j in Data[i]['types_b'] ]:
            flag_b_AA_charges = 1
            for count_j,j in enumerate(Data[i]['types_b']):
                Data[i]['charges_b'][count_j] = FF_charges[j]

        # Overwrite values for UA-molecules already in the FF_charges dictionary (a check is performed to ensure that all of the atom types
        # are in the FF parameters read from file for each molecule. If they are then the charges are read from file.)
        # NOTE: values are first initialized from a copy of the charges* arrays
        Data[i]["charges_a_UA"] = copy(Data[i]["charges_a"])
        Data[i]["charges_b_UA"] = copy(Data[i]["charges_b"])
        if False not in [ j+'-UA' in list(FF_charges.keys()) for j in Data[i]['types_a'] if int(j.split('[')[1].split(']')[0]) != 1 ]:
            flag_a_UA_charges = 1
            for count_j,j in enumerate(Data[i]['types_a']):
                if int(j.split('[')[1].split(']')[0]) == 1: continue
                Data[i]['charges_a_UA'][count_j] = FF_charges[j+'-UA']
        if False not in [ j+'-UA' in list(FF_charges.keys()) for j in Data[i]['types_b'] if int(j.split('[')[1].split(']')[0]) != 1 ]:
            flag_b_UA_charges = 1
            for count_j,j in enumerate(Data[i]['types_b']):
                if int(j.split('[')[1].split(']')[0]) == 1: continue
                Data[i]['charges_b_UA'][count_j] = FF_charges[j+'-UA']

        # Assign the total charge on molecule a
        if q_a is None:
            q_a = round(sum(Data[i]['charges_a']))
        else:
            q_a = float(q_a)

        # Assign the total charge on molecule b
        if q_b is None:
            q_b = round(sum(Data[i]['charges_b']))
        else:
            q_b = float(q_b)

        # Balance the charges to ensure an integer number of electrons on each molecule
        # this can potentially happen if the atom types being read from the database file
        # were derived from a different molecule
        residual   = sum(Data[i]['charges_a']) - q_a
        for count_j,j in enumerate(Data[i]['types_a']):

            # Check to see that the amount of charge being redistributed is small            
            if abs(residual/float(len(Data[i]['types_a']))) > 0.05: 
                residual_flag_a = 1

            # Equally redistribute the residual charge
            Data[i]['charges_a'][count_j] -= residual/float(len(Data[i]['types_a']))

        residual   = sum(Data[i]['charges_b']) - q_b
        for count_j,j in enumerate(Data[i]['types_b']):

            # Check to see that the amount of charge being redistributed is small
            if abs(residual/float(len(Data[i]['types_b']))) > 0.05: 
                residual_flag_b = 1

            # Equally redistribute the residual charge
            Data[i]['charges_b'][count_j] -= residual/float(len(Data[i]['types_b']))

        # Calculate UA charges
        Data[i]['adj_mat_a'] = Table_generator(Data[i]['elem_a'],Data[i]['geo_a'])            
        Data[i]['adj_mat_b'] = Table_generator(Data[i]['elem_b'],Data[i]['geo_b'])            

        # Add hydrogen charges to bonded carbons - for molecule a
        Data[i]["charges_a_UA"] = copy(Data[i]["charges_a"])
        for count_j,j in enumerate(Data[i]['adj_mat_a']):
            if int(Data[i]["types_a"][count_j].split('[')[1].split(']')[0]) != 6: continue
            H_ind = [ count_k for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_a'][count_k].split('[')[1].split(']')[0]) == 1 ]
            Data[i]["charges_a_UA"][count_j] += sum([ Data[i]["charges_a"][k] for k in H_ind ])
            Data[i]["charges_a_UA"][H_ind] = 0.0

        # Add hydrogen charges to bonded carbons - for molecule b
        Data[i]["charges_b_UA"] = copy(Data[i]["charges_b"])
        for count_j,j in enumerate(Data[i]['adj_mat_b']):
            if int(Data[i]["types_b"][count_j].split('[')[1].split(']')[0]) != 6: continue
            H_ind = [ count_k for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_b'][count_k].split('[')[1].split(']')[0]) == 1 ]
            Data[i]["charges_b_UA"][count_j] += sum([ Data[i]["charges_b"][k] for k in H_ind ])
            Data[i]["charges_b_UA"][H_ind] = 0.0
        
#         OLD LOOP (didn't check if all atoms were in the FF dictionary, compared with above)
#         # Overwrite values for any charges that are already in the FF_charges dictionary
#         for count_j,j in enumerate(Data[i]['types_a']):
#             if j+'-UA' in FF_charges.keys(): Data[i]['charges_a_UA'][count_j] = FF_charges[j+'-UA']
#         for count_j,j in enumerate(Data[i]['types_b']):
#             if j+'-UA' in FF_charges.keys(): Data[i]['charges_b_UA'][count_j] = FF_charges[j+'-UA']

        # Determine job type (i.e. 'dft' 'mp2' or combined 'dft' 'mp2')
        if count_i == 0:
            job_type = find_jobtype(i)

        # Parse the orca output files 
        flag = 0
        DFT_energy_count = 0
        MP2_energy_count = 0
        DFT_AB = 0.0
        DFT_A  = 0.0
        DFT_B  = 0.0
        MP2_AB = 0.0
        MP2_A  = 0.0
        MP2_B  = 0.0
        with open(i,'r') as output:
            for lines in output:
                fields = lines.split()

                # The ordering of jobs for 'dft mp2': 1=dimer_DFT, 2=dimer_MP2, 3=A_DFT, 4=A_MP2, 5=B_DFT, 6=B_MP2
                # The ordering of jobs for 'dft':     1=dimer_DFT, 2=A_DFT, 3=B_DFT
                # The ordering of jobs for 'dft mp2': 1=dimer_MP2, 2=A_MP2, 3=B_MP2
                if len(fields) == 5 and fields[0] == "FINAL" and fields[1] == "SINGLE" and fields[2] == "POINT" and fields[3] == "ENERGY":
                    DFT_energy_count += 1
                    # If mp2 jobs are interspersed then the HF jobs will also print "FINAL SINGLE POINT ENERGY"
                    if "dft" in job_type and "mp2" in job_type:
                        if DFT_energy_count == 1:
                            DFT_AB = float(fields[4])
                        elif DFT_energy_count == 3:
                            DFT_A  = float(fields[4])
                        elif DFT_energy_count == 5:
                            DFT_B  = float(fields[4])
                    # If only dft jobs are run then each FINAL SINGLE POINT ENERGY corresponds to one of the DFT jobs
                    elif "dft" in job_type:
                        if DFT_energy_count == 1:
                            DFT_AB = float(fields[4])
                        elif DFT_energy_count == 2:
                            DFT_A  = float(fields[4])
                        elif DFT_energy_count == 3:
                            DFT_B  = float(fields[4])

                if len(fields) == 5 and fields[0] == "MP2" and fields[1] == "TOTAL" and fields[2] == "ENERGY:":
                    MP2_energy_count += 1
                    if MP2_energy_count == 1:
                        MP2_AB = float(fields[3])
                    elif MP2_energy_count == 2:
                        MP2_A  = float(fields[3])
                    elif MP2_energy_count == 3:
                        MP2_B  = float(fields[3])
                        break

        if "dft" in job_type and "mp2" in job_type:
            if DFT_energy_count < 5:
                print("ERROR: Not enough DFT calculations were discovered in file {}".format(i))
                quit()
            elif MP2_energy_count < 3:
                print("ERROR: Not enough MP2 calculations were discovered in file {}".format(i))
                quit()
        elif "dft" in job_type and DFT_energy_count < 3:
            print("ERROR: Not enough DFT calculations were discovered in file {}".format(i))
            quit()
        elif "mp2" in job_type and MP2_energy_count < 3:
            print("ERROR: Not enough MP2 calculations were discovered in file {}".format(i))
            quit()
            
        # Calculate counterpoise-corrected interaction energy
        if "dft" in job_type:
            Data[i]["e_dft"]        = (DFT_AB-DFT_A-DFT_B)*Htokcalmol          # BSSE-corrected DFT interaction energy
        if "mp2" in job_type:
            Data[i]["e_mp2"]        = (MP2_AB-MP2_A-MP2_B)*Htokcalmol          # BSSE-corrected MP2 interaction energy

        # Calculate r_dist matrix, holding the pair-wise separations between atoms on molecule a and b
        Data[i]["r_dist"] = cdist(Data[i]['geo_a'],Data[i]['geo_b'])

        # Calculate pairs lists, holding the pair-wise separations for all instances of each pair-type in this configuration
        Data[i]["pairs"] = {}
        for count_j,j in enumerate(Data[i]["types_a"]):
            for count_k,k in enumerate(Data[i]["types_b"]):
                if (j,k) not in list(Data[i]["pairs"].keys()): Data[i]["pairs"][(j,k)]=[]
                Data[i]["pairs"][(j,k)]+=[Data[i]["r_dist"][count_j,count_k]]

        # Convert lists to arrays
        for j in list(Data[i]["pairs"].keys()):
            Data[i]["pairs"][j] = array(Data[i]["pairs"][j])

        # Calculate arrays for vectorized calculation of pairwise interactions
        # "pair_vector" holds the sum of 1/r elements for each pair type indexed to 
        # the "pair_type_vector" pair_types. That is, if element j in pair_type_vector
        # is type (x,y), then element j in "pair_vector" is sum(1/r) instances for 
        # type (x,y) in configuration i. if/else statements are there to remove redundancies.
        Data[i]["pair_vector-6"] = []
        Data[i]["pair_vector-12"] = []
        Data[i]["pair_type_vector"] = []
        added_types = []
        for j in list(Data[i]["pairs"].keys()):
            if j in added_types:
                continue
            if j[0]==j[1]:
                Data[i]["pair_vector-6"]  += [sum(1.0/Data[i]["pairs"][j]**6.0)]
                Data[i]["pair_vector-12"] += [sum(1.0/Data[i]["pairs"][j]**12.0)]
            elif (j[1],j[0]) in list(Data[i]["pairs"].keys()):
                Data[i]["pair_vector-6"]  += [sum(1.0/Data[i]["pairs"][j]**6.0)  + sum(1.0/Data[i]["pairs"][(j[1],j[0])]**6.0)]
                Data[i]["pair_vector-12"] += [sum(1.0/Data[i]["pairs"][j]**12.0) + sum(1.0/Data[i]["pairs"][(j[1],j[0])]**12.0)]                
            else:
                Data[i]["pair_vector-6"]  += [sum(1.0/Data[i]["pairs"][j]**6.0)]
                Data[i]["pair_vector-12"] += [sum(1.0/Data[i]["pairs"][j]**12.0)]
            if j[0] >= j[1]:
                Data[i]["pair_type_vector"]+=[j]
            else:
                Data[i]["pair_type_vector"]+=[(j[1],j[0])]
            added_types += [j,(j[1],j[0])]
        Data[i]["pair_vector-6"]  = array(Data[i]["pair_vector-6"])
        Data[i]["pair_vector-12"] = array(Data[i]["pair_vector-12"])

    # Print relevant diagnostics
    if residual_flag_a != 0 or residual_flag_b != 0 or flag_a_AA_charges == 1 or flag_b_AA_charges == 1 or flag_a_UA_charges == 1 or flag_b_UA_charges == 1:
        print(" ")
    if residual_flag_a != 0 or residual_flag_b != 0:
        print("{}\n! {:^140} !\n{}\n".format("!"*144,"More than abs(0.05) electrons per atom were required to reach an integer number of electrons","!"*144))
    if flag_a_AA_charges == 1:
        print("AA-charges were read for the a-molecule(s) from file(s): {}".format(', '.join(db_files)))
    if flag_b_AA_charges == 1:
        print("AA-charges were read for the b-molecule(s) from file(s): {}".format(', '.join(db_files)))
    if flag_a_UA_charges == 1:
        print("UA-charges were read for the a-molecule(s) from file(s): {}".format(', '.join(db_files)))
    if flag_b_UA_charges == 1:
        print("UA-charges were read for the b-molecule(s) from file(s): {}".format(', '.join(db_files)))

    ################################################################################
    # Perform charge rescaling to conserve the total electrostatics in the UA_fits #
    ################################################################################

    # Only rescale the UA charges if none of the molecules being fit have their UA_charges read from the FF dictionary.
    # (those being read are expected to be fixed, so charge rescaling needs to be avoided)
    if flag_a_UA_charges == 0 and flag_b_UA_charges == 0:

        # Calculate total intermolecular AA-electrostatic energy
        E_C_Tot_AA = 0.0
        for i in Names:
            E_C_Tot_AA += E_C_Tot(i,'AA')

        # Calculate total intermolecular UA-electrostatic energy (uncorrected)
        E_C_Tot_UA = 0.0
        for i in Names:
            E_C_Tot_UA += E_C_Tot(i,'UA')

        # calculate scale factor (the scale factor is the sqrt of the ratio because coulomb evaluations involve the product of charges)
        scale_factor = (E_C_Tot_AA/E_C_Tot_UA)**(0.5)

        # Print diagnostic and calculate scale factor 
        print("\n{:40s} {:< 12.6f} (kcal/mol)".format("All-atom EC_tot:",E_C_Tot_AA))
        print("{:40s} {:< 12.6f} (kcal/mol)".format("United-atom EC_tot:",E_C_Tot_UA))
        print("{:40s} {:< 12.6f}".format("United-atom EC rescaling factor:",scale_factor))

#         # Scale the UA charges and recalculate the total UA-electrostatic energy
#         E_C_Tot_UA = 0.0
#         for i in Names:
#             Data[i]["charges_a_UA"] *= scale_factor
#             Data[i]["charges_b_UA"] *= scale_factor
#             E_C_Tot_UA += E_C_Tot(i,'UA')

#         # Print diagnostic
#         print "{:40} {:< 12.6f} (kcal/mol)".format("Rescaled United-atom EC_tot:",E_C_Tot_UA)

    # Calculate interaction energies with the coulombic contribution subtracted out
    for i in Names:
        if "dft" in job_type:
            Data[i]["e_dft_fit_AA"] = Data[i]["e_dft"] - E_C_Tot(i,'AA')       # DFT fit potential (subtracted electrostatic component)
            Data[i]["e_dft_fit_UA"] = Data[i]["e_dft"] - E_C_Tot(i,'UA')       # DFT fit potential (subtracted electrostatic component)
        if "mp2" in job_type:
            Data[i]["e_mp2_fit_AA"] = Data[i]["e_mp2"] - E_C_Tot(i,'AA')       # MP2 fit potential (subtracted electrostatic component)
            Data[i]["e_mp2_fit_UA"] = Data[i]["e_mp2"] - E_C_Tot(i,'UA')       # MP2 fit potential (subtracted electrostatic component)

#         # Print data:
#         print "Data[{}]['e_dft']  = {}".format(i,Data[i]["e_dft"])
#         print "Data[{}]['e_mp2']  = {}".format(i,Data[i]["e_mp2"])
#         print "DFT_A: {}".format(DFT_A)
#         print "DFT_B: {}".format(DFT_B)
#         print "DFT_AB: {}\n".format(DFT_AB)

    return Data,job_type

# Description: scrape charges from the supplied db file. Any charges read from the file are used regardless of the
#              charge parsing option.
def grab_db_charges(db_files):
    
    FF_charges = {}
    for i in db_files:
        with open(i,'r') as f:
            for lines in f:
                fields = lines.split()
                if len(fields) > 0 and fields[0] == 'charge':
                    FF_charges[fields[1]] = float(fields[2])

    return FF_charges

# Description: scrape the charges for both molecules from the first CHELPG calculation
#              of the QC calculation. 
def scrape_charges(filename,atomtypes_a,atomtypes_b):

    # Read in the charges
    charges_flag  = 0
    break_flag    = 0
    charges_count = 0
    charges_a = zeros(len(atomtypes_a))
    charges_b = zeros(len(atomtypes_b))
    with open(filename,'r') as f:
        for count,lines in enumerate(f):
            fields = lines.split()

            # Break after all charges have been read
            if break_flag == 1:
                break

            # Find the start of the first CHELPG calculation
            if len(fields) == 2 and fields[0] == "CHELPG" and fields[1] == "Charges":
                charges_flag = 1
                continue

            # Necessary because of a "------" filler line in the Orca CHELPG output
            if charges_flag == 1:
                charges_flag = 2
                continue

            # Start parsing charges in molecule b
            if charges_flag == 3:
                charges_b[charges_count-len(charges_a)] = float(fields[3])
                charges_count += 1

                # Once all charges have been read reset the count and flag.
                if charges_count == len(charges_a)+len(charges_b):
                    break_flag = 1

            # Start parsing charges in molecule a
            if charges_flag == 2:
                charges_a[charges_count] = float(fields[3])
                charges_count += 1
                if charges_count == len(charges_a):
                    charges_flag += 1
    

    # Average over like types of molecule a
    for i in sorted(set(atomtypes_a)):
        ind = []
        charge = 0.0

        # find like types and sum charges
        for count_j,j in enumerate(charges_a):
            if atomtypes_a[count_j] == i:
                ind += [count_j]
                charge += j

        # Assign averaged charges to all like types
        for j in ind:
            charges_a[j] = charge/float(len(ind))        

    # Average over like types of molecule b
    for i in sorted(set(atomtypes_b)):
        ind = []
        charge = 0.0

        # find like types and sum charges
        for count_j,j in enumerate(charges_b):
            if atomtypes_b[count_j] == i:
                ind += [count_j]
                charge += j

        # Assign averaged charges to all like types
        for j in ind:
            charges_b[j] = charge/float(len(ind))        

    return charges_a,charges_b

# Return the coordinates and atom types from an xyz file
# atom types are expected as a fifth column and molecule
# index is expected in the sixth column. These are atypical
# an xyz file, but automatically outputed by the vdw_gen.py 
# program.
def scrape_xyz(name,parse_charges=0):

    atom_count = 0
    with open(name,'r') as f:
        for count_j,j in enumerate(f):
            fields = j.split()

            # Grab the number of atoms in the geometry and initialize the parsing lists
            if count_j == 0:
                elem = ["X"]*int(fields[0])
                geo = zeros([int(fields[0]),3]) 
                types = ["X"]*int(fields[0])
                a_ind = []
                b_ind = []
                
                # Only parse charges if flag is enabled
                if parse_charges == 1:
                    charges =zeros([int(fields[0])])

            # Within the geometry block parse the geometry, charges, mol_ids, and atom types from the enhanced xyz file
            if count_j > 1 and len(fields) >= 7:
                geo[atom_count,:] = array([float(fields[1]),float(fields[2]),float(fields[3])])
                if fields[4] == '0':
                    a_ind += [atom_count]
                else:
                    b_ind += [atom_count]
                elem[atom_count] = fields[0]

                # Only parse charges if flag is enabled
                if parse_charges == 1:
                    charges[atom_count] = float(fields[5])
                types[atom_count] = fields[6]
                atom_count += 1    

    # If flag is enabled return the charges
    if parse_charges == 1:
        return [ i for count_i,i in enumerate(elem) if count_i in a_ind ],\
               [ i for count_i,i in enumerate(elem) if count_i in b_ind ],\
               geo[a_ind],geo[b_ind],\
               [ i for count_i,i in enumerate(types) if count_i in a_ind ],\
               [ i for count_i,i in enumerate(types) if count_i in b_ind ],\
               charges[a_ind],charges[b_ind]

    # Don't return the charges list if the flag isn't enabled
    if parse_charges == 0:
        return [ i for count_i,i in enumerate(elem) if count_i in a_ind ],\
               [ i for count_i,i in enumerate(elem) if count_i in b_ind ],\
               geo[a_ind],geo[b_ind],\
               [ i for count_i,i in enumerate(types) if count_i in a_ind ],\
               [ i for count_i,i in enumerate(types) if count_i in b_ind ]

# This function reads in an atomtype label, generates its adjacency matrix, 
# and determines if it is a UA-hydrogen based on the first atomtype and its connections
def return_UA_H(type):

    # Initialize lists/array 
    brackets = [ i for i in type if i == "[" or i == "]" ]
    atoms    = type.replace('['," ").replace(']'," ").split()
    adj_mat = zeros([len([ i for i in brackets if i == "[" ]),len([ i for i in brackets if i == "[" ])])

    # The basis of adj_mat corresponds to the atoms in "atoms", equivalently the "[" objects in brackets
    for count_i,i in enumerate(adj_mat):        

        # Indices are counted based on the number of left_brackets encountered (i.e. left_count's values)
        # Connections occur when bracket_count is equal to 1 and the current object in brackets is a left-bracket.
        # The algorithm works by incrementing bracket_count by +/- 1 for every bracket encountered.
        bracket_count = -1
        left_count = -1

        # Loop over the brackets list, bracket_count only gets incremented for objects in brackets satisfying left_count >= count_i 
        for count_j,j in enumerate(brackets):
            if j == "[": left_count += 1
            if left_count < count_i: continue
            if j == "[": bracket_count += 1; 
            if j == "]": bracket_count -= 1
            if left_count == count_i: continue
            if bracket_count == 1 and j == "[":
                adj_mat[count_i,left_count] = 1
                adj_mat[left_count,count_i] = 1

    # If the current atom is a hydrogen (i.e. atomtype "1") and it is connected to at least on carbon (i.e. atomtype "6") then it is a UA-type
    if atoms[0] == "1" and len([ i for count_i,i in enumerate(adj_mat[0]) if i == 1 and atoms[count_i] == "6" ]) > 0:
        return 1
    # Else, it is not a UA-type
    else:
        return 0

# This function drives the fitting algorithm for UA and AA pairwise interaction
# Calls subfunctions "iterative_fit" and "plot_convergence" to perform the fitting
# and generate convergence plots.
def Fit_pairs(Folder,Data,fit,weight,delta_xhi2_thresh,min_cycles,max_cycles,FF_db,method_list,QC_type="DFT",charges_opt=1,UA_vol_scale=1.0,UA_method='radius',r_min_scale=1.0,avoid_read_flag=0,mixing_rule="none",L2_s=0.0,L2_e=0.0,outlier_option=False):

    global VDW_dict,fit_type,Pair_min_dict

    # Get unique pairs and number of occurrences
    Fit_Pairs = []                                                   # Initialize list holding unique pair interactions
    Fit_Pairs_hist = {}                                              # Initialize list indexed to pairs to hold the occurence of each pair interaction (there can be multiple per pair of molecules).
    for i in list(Data.keys()):                                            # Iterate over all configurations
        for j in Data[i]["types_a"]:                                 # Iterate over the list of types in molecule A
            for k in Data[i]["types_b"]:                             # Iterate over the list of types in molecule B
                if j >= k:                                           # Higher ranked elements go first in the definition of the type
                    if (j,k) not in Fit_Pairs:                       # If the tuple isn't in the Pairs list then add it and add an element in Pairs_hist
                        Fit_Pairs += [(j,k)]
                        Fit_Pairs_hist[(j,k)] = 1
                    else:                                            # If the pair already exists Pairs then increment the appropriate 
                        Fit_Pairs_hist[(j,k)] += 1                   # element in Fit_Pairs_hist
                else:
                    if (k,j) not in Fit_Pairs:                       # If the tuple isn't in the Pairs list then add it and add an element in Pairs_hist
                        Fit_Pairs += [(k,j)]
                        Fit_Pairs_hist[(k,j)] = 1
                    else:                                            # If the pair already exists Pairs then increment the appropriate 
                        Fit_Pairs_hist[(k,j)] += 1                   # element in Fit_Pairs_hist  
    
    # Generate a list of UA pairs (pairs between atoms that do not contain hydrogen)
    UA_Pairs = []
    for i in Fit_Pairs:
        if return_UA_H(i[0]) != 1 and return_UA_H(i[1]) != 1:
            UA_Pairs += [i]

    # Read in fixed parameters from the fixed_params.db file
    FIXED_Data = {"vdws":{}}
    if os.path.isfile('fixed_params.db'):         
        FIXED_Data = get_FF_data(['fixed_params.db'])

    # If a FF_db file is supplied then the program avoids fitting any of the
    # pairs that are pairs are in the database
    DB_Pairs = []
    FF_Data = {}
    if FF_db != [] or len(list(FIXED_Data['vdws'].keys())) > 0:
        
        # Read in parameters from user-supplied database
        FF_Data = get_FF_data(FF_db,mixing_rule=mixing_rule)

        # Add FIXED_Data parameters while avoiding overwriting those read from the user-supplied database
        for i in [ j for j in list(FIXED_Data["vdws"].keys()) if j not in list(FF_Data["vdws"].keys()) ]:
            FF_Data["vdws"][i] = FIXED_Data["vdws"][i]

        if fit == 'lj':

            # Move AA-pairs that are already in the database to DB_Pairs
            if mixing_rule == "none":
                move_list  = [ count_i for count_i,i in enumerate(Fit_Pairs) if ( (i[0],i[1],'lj') in list(FF_Data["vdws"].keys()) or (i[1],i[0],'lj') in list(FF_Data["vdws"].keys()) ) ]
            elif mixing_rule in ['wh','lb']:
#                move_list  = [ count_i for count_i,i in enumerate(Fit_Pairs) if ( (i[0],i[1],'lj') in FF_Data["vdws"].keys() or (i[1],i[0],'lj') in FF_Data["vdws"].keys() ) and i[0] == i[1] ]
                move_list  = [ count_i for count_i,i in enumerate(Fit_Pairs) if ( (i[0],i[0],'lj') in list(FF_Data["vdws"].keys()) and (i[1],i[1],'lj') in list(FF_Data["vdws"].keys()) ) ]
            else:
                move_list  = []
            DB_Pairs       = [ i for count_i,i in enumerate(Fit_Pairs) if count_i in move_list ]
            Fit_Pairs      = [ i for count_i,i in enumerate(Fit_Pairs) if count_i not in move_list ]

            # Move UA-pairs that are already in the database to DB_Pairs
            if mixing_rule == "none":
                move_list  = [ count_i for count_i,i in enumerate(UA_Pairs) if ( (i[0]+'-UA',i[1]+'-UA','lj') in list(FF_Data["vdws"].keys()) or (i[1]+'-UA',i[0]+'-UA','lj') in list(FF_Data["vdws"].keys()) ) ]
            elif mixing_rule in ['wh','lb']:
#                move_list  = [ count_i for count_i,i in enumerate(UA_Pairs) if ( (i[0]+'-UA',i[1]+'-UA','lj') in FF_Data["vdws"].keys() or (i[1]+'-UA',i[0]+'-UA','lj') in FF_Data["vdws"].keys() ) and i[0] == i[1] ]
                move_list  = [ count_i for count_i,i in enumerate(UA_Pairs) if ( (i[0]+'-UA',i[0]+'-UA','lj') in list(FF_Data["vdws"].keys()) and (i[1]+'-UA',i[1]+'-UA','lj') in list(FF_Data["vdws"].keys()) ) ]
            else:
                move_list  = []
#            DB_Pairs       = [ i for count_i,i in enumerate(Fit_Pairs) if count_i in move_list ]
            UA_Pairs       = [ i for count_i,i in enumerate(UA_Pairs) if count_i not in move_list ]
            
        if fit == 'buck':

            # Move AA-pairs that are already in the database to DB_Pairs
            if mixing_rule == "none":
                move_list  = [ count_i for count_i,i in enumerate(Fit_Pairs) if ( (i[0],i[1],'buck') in list(FF_Data["vdws"].keys()) or (i[1],i[0],'buck') in list(FF_Data["vdws"].keys()) ) ]
            else:
                move_list  = []
            DB_Pairs       = [ i for count_i,i in enumerate(Fit_Pairs) if count_i in move_list ]
            Fit_Pairs      = [ i for count_i,i in enumerate(Fit_Pairs) if count_i not in move_list ]

            # Move UA-pairs that are already in the database to DB_Pairs
            if mixing_rule == "none":                
                move_list  = [ count_i for count_i,i in enumerate(UA_Pairs) if ( (i[0]+'-UA',i[1]+'-UA','buck') in list(FF_Data["vdws"].keys()) or (i[1]+'-UA',i[0]+'-UA','buck') in list(FF_Data["vdws"].keys()) ) ]
            else:
                move_list  = []
            DB_Pairs      += [ (i[0]+'-UA',i[1]+'-UA') for count_i,i in enumerate(UA_Pairs) if count_i in move_list ]
            UA_Pairs       = [ i for count_i,i in enumerate(UA_Pairs) if count_i not in move_list ]

    # Generate a dictionary of minimum separation for pairs (used during fit to ensure that 
    # the minimum in the fit potential does not occur in a region where there is no sampling)
    Pair_min_dict = {}
    Pair_min_dict_UA = {}
    for i in Fit_Pairs: Pair_min_dict[i] = 1000.0; Pair_min_dict[(i[1],i[0])] = 1000.0
    for i in UA_Pairs: Pair_min_dict[i] = 1000.0; Pair_min_dict[(i[1],i[0])] = 1000.0     # Can eventually be removed once separate Pair_min_dict_AA and Pair_min_dict_UA are coded
    for i in list(Data.keys()):
        for count_j,j in enumerate(Data[i]["r_dist"]):
            for count_k,k in enumerate(j):

                # Order the pair to conform to VDW_dict key definitions
                if Data[i]['types_a'][count_j] >= Data[i]['types_b'][count_k]:
                    pair_type = (Data[i]['types_a'][count_j],Data[i]['types_b'][count_k])
                else:
                    pair_type = (Data[i]['types_b'][count_k],Data[i]['types_a'][count_j])

                # Replace minimum value
                if pair_type in list(Pair_min_dict.keys()) and k < Pair_min_dict[pair_type]: Pair_min_dict[pair_type] = k; Pair_min_dict[(pair_type[1],pair_type[0])] = k

    # Scale Pair_min_dict entries by the r_min_scale factor (when set to 1.0 no change is applied)
    for i in list(Pair_min_dict.keys()):
        Pair_min_dict[i] *= r_min_scale
    
    # Print diagnostics
    print("\nUnique AA pairs being fit (occurences,min_sep,min_sigma):\n")
    for count_i,i in enumerate(Fit_Pairs):
        print("\tpair {:80}: {:<10d}  {:<10.4f} {:<10.4f}".format(str(i),Fit_Pairs_hist[i],Pair_min_dict[i],Pair_min_dict[i]*2.0**(-1.0/6.0)))
    print("\nUnique UA pairs being fit (occurences):\n")
    for i in UA_Pairs:
        print("\tpair {:80}: {:<10d}".format(str(i),Fit_Pairs_hist[i]))

    # Generate a separation vs number histogam for all pair types
    print("\n\tGenerating separation histograms...")
#    pair_hist_plot(Folder+'/figures',Fit_Pairs,Data,QC_type=QC_type)

    ###############################
    # Fit DFT all-atom parameters #
    ###############################
    if QC_type == "DFT":
        
        print("\n{}".format("*"*144))
        print("* {:^140s} *".format("Fit DFT All-Atom VDW Parameters (Self-Consistent Pair-Type Fits to All Configurations)"))
        print("{}".format("*"*144))

        # Initialize fit arrays
        x_vals = list(Data.keys())
        y_vals = [ Data[i]["e_dft_fit_AA"] for i in x_vals ]
        E_config_tot = [ Data[i]["e_dft"] for i in x_vals]

        # Initialize VDW_dict based on DFT_FF and UFF parameters    
        VDW_dict = init_VDW(Data,Fit_Pairs,DB_Pairs,Pair_min_dict,fit,FF_db,{},UA=0,avoid_read_flag=avoid_read_flag,mixing_rule=mixing_rule)
        VDW_dict_multi = {}
         

        # Initialize UFF reference dictionary
        VDW_0 = init_VDW(Data,Fit_Pairs,[],Pair_min_dict,fit,[],{},UA=0,avoid_read_flag=1,mixing_rule=mixing_rule,verbose=False)        
        missing_key= [ k for k in VDW_dict if k not in VDW_0]
        for k in missing_key:
            VDW_0[k] =  VDW_dict[k]
        #VDW_dict_multi['UFF'] = VDW_0
        VDW_dict_multi['cycle-19']=dict(VDW_dict)
        fit_vals = {}
        xhi_total = {}
        # Just to get cycle info, this should be modified
        # Cuz lstsq used cycle-19 as initial guess
        fit_vals_tmp,xhi_tmp = fit_method('lstsq',x_vals,y_vals,Fit_Pairs,weight,fit,E_config_tot,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule=mixing_rule,L2_s=L2_s,L2_e=L2_e,VDW_0=VDW_0)
        print("go")
        quit()
        VDW_dict = init_VDW(Data,Fit_Pairs,DB_Pairs,Pair_min_dict,fit,FF_db,{},UA=0,avoid_read_flag=avoid_read_flag,mixing_rule=mixing_rule)
        xhi_total['cycle-19'] = xhi_tmp[0]
        fit_vals['cycle-19']=fit_vals_tmp[0]
        
        for i in method_list:
            fit_vals_tmp,xhi_tmp = fit_method(i,x_vals,y_vals,Fit_Pairs,weight,fit,E_config_tot,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule=mixing_rule,L2_s=L2_s,L2_e=L2_e,VDW_0=VDW_0,outlier_option=outlier_option)
            #fit_vals_tmp,xhi_tmp = global_fit(x_vals,y_vals,Fit_Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule=mixing_rule,L2_s=L2_s,L2_e=L2_e,VDW_0=VDW_0)

            VDW_dict_multi[i] = deepcopy(VDW_dict) # have to use dict to make it local or it's like pointer in c++
            #### global has to be first so that xhi_total would be UFF, if lstsq first then it would actually be cycle-19 
            if(len(fit_vals) == 1):
               fit_vals['UFF'] = fit_vals_tmp[0]
               fit_vals[i] = fit_vals_tmp[1]
               xhi_total['UFF'] = xhi_tmp[0]
               xhi_total[i] = xhi_tmp[1]
            else:
               fit_vals[i] = fit_vals_tmp[1]
               xhi_total[i] = xhi_tmp[1]


            write_params(Folder+'/'+i+'/DFT-AA',fit,type='AA',write_charges=charges_opt) 

            # Generate individual plots of the fit potential
            plot_fit_potentials(Folder+'/'+i,fit,type='DFT-AA',r_min=0.1,r_max=10.0)
            plot_convergence(Folder+'/'+i,Data,[fit_vals[i]],plot_num=5,type='DFT-AA')
            # Print diagnostic indicating which pairs probably have sampling errors based if their sigma values are within
            # 0.05% of their minimum separation threshold (i.e. the element in Pair_min_dict generated based on the configurational distribution)
            Warning_list = [ i for i in Fit_Pairs if VDW_dict[i][1]*2.0**(1.0/6.0) < Pair_min_dict[i]*1.05 ]
            if len(Warning_list) > 0:
                print("{}".format("*"*144))
                print("* {:^140s} *".format("The following pairs have minima in their LJ interaction within 5% of their closest sampled separation."))
                print("* {:^140s} *".format("This is usually indicative of insufficient sampling."))
                print("{}".format("*"*144))
                print("\n\t{:<80s} {:<20s} {:<20s} {:<20s}".format("Pair","sigma (ang)","r_min (ang)","min_sep (ang)"))
                for i in Warning_list:
                    print("\t{:80s} {:<20.6f} {:<20.6f} {:<20.6f}".format(str(i),VDW_dict[i][1],VDW_dict[i][1]*2.0**(1.0/6.0),Pair_min_dict[i]))
                print("")


        # Generate convergence plot
        #method_list = ['cycle-19','UFF']+method_list
        plot_convergence_new(Folder+'/figures',Data,fit_vals,xhi_total,plot_num=5,type='DFT-AA')
        plot_convergence_outlier(Folder+'/figures',Data,fit_vals,xhi_total,plot_num=5,type='DFT-AA')
        #plot_fit_potentials_method(VDW_dict_multi,Folder+'/figures',fit,type='DFT-AA',r_min=0.1,r_max=10.0)
        plot_method_convergence(Folder+'/figures',VDW_dict_multi,xhi_total,VDW_0=VDW_0)
        quit()

    ##################################
    # Fit DFT united-atom parameters #
    ##################################
    if QC_type == "DFT":

        # Initialize fit arrays
        x_vals = list(Data.keys())
        y_vals = [ Data[i]["e_dft_fit_UA"] for i in x_vals ]

        if UA_method != "fit":

            print("{}".format("*"*144))
            print("* {:^140s} *".format("Fit DFT United-Atom VDW Parameters (Conversion from AA-parameters based on configurational fit)"))
            print("{}".format("*"*144))

            # Calculate UA parameters based on a geometric conversion of the AA sigmas and scaled AA-epsilons to reproduce the configurational LJ interaction.                        
#            VDW_dict = calc_eff_UA(Folder+'/figures',VDW_dict,DB_Pairs,Data)            
            VDW_dict = calc_eff_UA(Folder+'/figures',VDW_dict,UA_Pairs,Data)

            # Calculate UA parameters based on a geometric conversion of the AA sigmas and scaled AA-epsilons to reproduce the configurational LJ interaction.                        
#            VDW_dict = calc_UA_parameters(VDW_dict,DB_Pairs,Data,vol_scale=UA_vol_scale,vol_method=UA_method)

            # Calculate the configurational energies and print the average xhi2 difference between the fit and DFT values
            fit_vals = [ E_LJ_configs(x_vals) ]
            print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",mean((y_vals-fit_vals[0])**2)))

            # Calculate VDW energy with united-atom parameters
            E_VDW_after = 0.0
            if fit == 'lj':
                for i in x_vals:
                    E_VDW_after += E_LJ_Tot(i)
            elif fit == 'buck':
                for i in x_vals:
                    E_VDW_after += E_BUCK_Tot(i)

        # When using the "fit" approach, the LJ parameters involving UA-Carbons are fit to the configurational energies directly
        # just like the AA-parameters, with the exception that the AA-parameters are used as an initial guess and all parameters 
        # NOT involving UA-Carbons are held constant.
        if UA_method == "fit":

            print("{}".format("*"*144))
            print("* {:^140s} *".format("Fit DFT United-Atom VDW Parameters (Direct fit to the QC data)"))
            print("{}".format("*"*144))

            # Initialize VDW_dict based on DFT_FF and UFF parameters    
#            VDW_dict = init_VDW(Data,Fit_Pairs,DB_Pairs,Pair_min_dict,fit,FF_db,VDW_dict,UA=0)
            VDW_dict = init_VDW(Data,Fit_Pairs,DB_Pairs,Pair_min_dict,fit,FF_db,VDW_dict,UA=1,avoid_read_flag=avoid_read_flag,mixing_rule=mixing_rule)

            # Remove pairs that don't involve UA-Carbons in the interaction, this has the effect of holding the other parameters
            # constant during the fit.
            UA_Pairs = [ i for i in UA_Pairs if UA_carbon(i[0]) == True or UA_carbon(i[1]) == True ]

            # Perform iterative fit of all parameters involving UA-Carbons
            fit_vals = iterative_fit_UA(x_vals,y_vals,UA_Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles)

            # Fit all epsilon parameters at once.
            # fit_vals = fit_all_UA(x_vals,y_vals,UA_Pairs,fit,delta_xhi2_thresh)

        # Write VDW parameters to database file (note VDW_dict is a global variable so the subfunction has access to the parameters)
        write_params(Folder+'/DFT-UA',fit,type='UA',write_charges=charges_opt)
#        write_params(Folder+'/DFT-UA.db',fit,type='UA',write_charges=1) # Always writes parameters

        # Generate individual plots of the fit potential
        plot_fit_potentials(Folder+'/figures',fit,type='DFT-UA',r_min=0.1,r_max=10.0)

        # Generate convergence plot
        plot_convergence(Folder+'/figures',Data,fit_vals,plot_num=5,type='DFT-UA')

    ###############################
    # Fit MP2 all-atom parameters #
    ###############################
    if QC_type == "MP2":

        print("\n{}".format("*"*144))
        print("* {:^140s} *".format("Fit MP2 All-Atom Parameters (Iterative Pair-Type Fits to All Configurations)"))
        print("{}".format("*"*144))

        # Initialize fit arrays
        x_vals = list(Data.keys())
        y_vals = [ Data[i]["e_mp2_fit_AA"] for i in x_vals ]

        # Initialize VDW_dict based on DFT_FF and UFF parameters    
        VDW_dict = init_VDW(Data,Fit_Pairs,DB_Pairs,Pair_min_dict,fit,FF_Data,{},UA=0,avoid_read_flag=avoid_read_flag,mixing_rule=mixing_rule)

        # Iteratively fit all parameters
        fit_vals = iterative_fit(x_vals,y_vals,Fit_Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles)

        # Write VDW parameters to database file (note VDW_dict is a global variable so the subfunction has access to the parameters)
        write_params(Folder+'/MP2-AA',fit,type='AA',write_charges=charges_opt)

        # Generate individual plots of the fit potential
        plot_fit_potentials(Folder+'/figures',fit,type='MP2-AA',r_min=0.1,r_max=10.0)

        # Generate convergence plot
#        plot_convergence(Folder+'/figures',y_vals,fit_vals,plot_num=5,type='MP2-AA')
        plot_convergence(Folder+'/figures',Data,fit_vals,plot_num=5,type='MP2-AA') 

    ##################################
    # Fit MP2 united-atom parameters #
    ##################################
    if QC_type == "MP2":

        print("{}".format("*"*144))
        print("* {:^140s} *".format("Fit MP2 United-Atom VDW Parameters (Conversion from AA-parameters based on excluded volume)"))
        print("{}".format("*"*144))

        # Initialize fit arrays
        x_vals = list(Data.keys())
        y_vals = [ Data[i]["e_mp2_fit_UA"] for i in x_vals ]

        # Calculate UA parameters based on a geometric conversion of the AA sigmas and scaled AA-epsilons to reproduce the configurational LJ interaction.
        VDW_dict = calc_UA_parameters(VDW_dict,DB_Pairs,Data,vol_scale=UA_vol_scale,vol_method=UA_method)

        # Calculate the configurational energies and print the average xhi2 difference between the fit and DFT values
        fit_vals = [ E_LJ_configs(x_vals) ]
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",mean((y_vals-fit_vals[0])**2)))

        # print "{}".format("*"*144)
        # print "* {:^140s} *".format("Fit MP2 United-Atom Parameters (Iterative Pair-Type Fits to All Configurations)")
        # print "{}".format("*"*144)

        # # Initialize fit arrays
        # x_vals = Data.keys()
        # y_vals = [ Data[i]["e_mp2_fit_UA"] for i in x_vals ]

        # # Initialize UA sigmas based on a geometric conversion of the AA sigmas
        # VDW_dict = set_UA_sigmas(VDW_dict,DB_Pairs,Data,vol_scale=UA_vol_scale,vol_method=UA_method)

        # # Initialize epsilon values in VDW_dict based on DFT_FF and UFF parameters    
        # VDW_dict = init_VDW(Data,Fit_Pairs,DB_Pairs,Pair_min_dict,FF_Data,VDW_dict,UA=1)

        # # Iteratively fit all parameters
        # fit_vals = iterative_fit_UA(x_vals,y_vals,UA_Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles)

        # Write VDW parameters to database file (note VDW_dict is a global variable so the subfunction has access to the parameters)
        write_params(Folder+'/MP2-UA',fit,type='UA',write_charges=charges_opt)

        # Generate individual plots of the fit potential
        plot_fit_potentials(Folder+'/figures',fit,type='MP2-UA',r_min=0.1,r_max=10.0)

        # Generate convergence plot
#        plot_convergence(Folder+'/figures',y_vals,fit_vals,plot_num=5,type='MP2-UA')
        plot_convergence(Folder+'/figures',Data,fit_vals,plot_num=5,type='MP2-UA')

# This function expects data to be a globally defined dictionary. The function
# outputs the coulomb contribution to the interaction energy of a pair of molecules
def E_C_Tot(ind,type):

    # Initialize the energy
    E_C = 0.0

    # Calculate the atoms_A atom_B separations
    r_dist = cdist(Data[str(ind)]['geo_a'],Data[str(ind)]['geo_b'])

    # Cumulatively add up the electrostatic energy
    for count_i,i in enumerate(r_dist):
        for count_j,j in enumerate(i):
            if type == 'AA':
                E_C += ( Data[str(ind)]["charges_a"][count_i]*Data[str(ind)]["charges_b"][count_j] / j )*Ecoul_Const
            elif type == 'UA':
                E_C += ( Data[str(ind)]["charges_a_UA"][count_i]*Data[str(ind)]["charges_b_UA"][count_j] / j )*Ecoul_Const

    return E_C


# This function expects data to be a globally defined dictionary. The function
# outputs the LJ contribution to the interaction energy of a pair of molecules
def E_LJ_Tot(ind):

    global Data,VDW_dict

    # Initialize the energy
    E_LJ = 0
    
    # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
    for i in list(Data[ind]["pairs"].keys()):
        E_LJ += sum( ( 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(12.0)*Data[ind]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(6.0)*Data[ind]["pairs"][i]**(-6.0) ) )

    return E_LJ

# This function expects data to be a globally defined dictionary. The function
# outputs the Buckingham contribution to the interaction energy of a pair of molecules
def E_BUCK_Tot(ind):

    global Data,VDW_dict

    # Initialize the energy
    E_BUCK = 0
    
    # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
    for i in list(Data[ind]["pairs"].keys()):
        E_BUCK += sum( ( VDW_dict[i][0]*exp(-Data[ind]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[ind]["pairs"][i]**(-6.0) ) )

    return E_BUCK

# def E_LJ_Type_test(type,VDW_new):

#     global Data

#     # Initialize the energy
#     E_LJ = 0.0

#     # Iterate over all configurations
#     for d in Data.keys():

#         # Add up all LJ interations involving this type
#         if type in Data[d]["pairs"].keys():
#             E_LJ += sum( ( 4.0*VDW_new[type][0]*VDW_new[type][1]**(12.0)*Data[d]["pairs"][type]**(-12.0) - 4.0*VDW_new[type][0]*VDW_new[type][1]**(6.0)*Data[d]["pairs"][type]**(-6.0) ) )

#         # Add up reverse combinations if applicable
#         if type[0] != type[1] and (type[1],type[0]) in Data[d]["pairs"].keys():
#             E_LJ += sum( ( 4.0*VDW_new[type][0]*VDW_new[type][1]**(12.0)*Data[d]["pairs"][(type[1],type[0])]**(-12.0) - \
#                            4.0*VDW_new[type][0]*VDW_new[type][1]**(6.0)*Data[d]["pairs"][(type[1],type[0])]**(-6.0) ) )
            
#     return E_LJ

# This function expects data to be a globally defined dictionary. The function
# outputs the LJ contribution to the interaction energy of a pair interaction, summed over all configurations
# contributions are only summed for atom types possessing the proper connections (e.g. there may be redundant H types
# but for UA purposes it is important to resolve them by which type of carbon they are attached to. con_0 supplies the atomtype
# type[0] should have a connection with, con_1 supplies the atomtype type[1] should have a connection with.)
def E_LJ_Type_w_con(type,con_0,con_1):

    global Data,VDW_dict

    # Initialize the energy
    E_LJ = 0.0

    # Iterate over all configurations
    for d in list(Data.keys()):

        # If a contraint is present for type 0, find indices of atomtype type[0] that are bonded to the supplied contraint
        if con_0 != []:
            ind_con_0_a  = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == con_0 ]
            ind_type_0_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[0] and True in [ 1 == j for j in Data[d]["adj_mat_a"][count_i][ind_con_0_a] ] ] 

            ind_con_0_b  = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == con_0 ]
            ind_type_0_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[0] and True in [ 1 == j for j in Data[d]["adj_mat_b"][count_i][ind_con_0_b] ] ] 
            
        # Else, return indices of all type[0] atoms
        else:
            ind_type_0_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[0] ]
            ind_type_0_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[0] ]

        # If a constraint is present for type 1, find indices of atomtype type[1] that are bonded to the supplied contraint     
        if con_1 != []:
            ind_con_1_a  = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == con_1 ]
            ind_type_1_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[1] and True in [ 1 == j for j in Data[d]["adj_mat_a"][count_i][ind_con_1_a] ] ] 

            ind_con_1_b  = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == con_1 ]
            ind_type_1_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[1] and True in [ 1 == j for j in Data[d]["adj_mat_b"][count_i][ind_con_1_b] ] ] 

        # Else, return indices of all type[1] atoms
        else:
            ind_type_1_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[1] ]
            ind_type_1_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[1] ]

        # Find pair-wise separations for this configuration
        seps = []
        for i in ind_type_0_a:
            for j in ind_type_1_b:
                seps += [Data[d]["r_dist"][i,j]]
        if type[0] != type[1] or con_0 != con_1:
            for i in ind_type_1_a:
                for j in ind_type_0_b:
                    seps += [Data[d]["r_dist"][i,j]]

        # Add up all LJ interations involving this type
        E_LJ += sum( ( 4.0*VDW_dict[type][0]*VDW_dict[type][1]**(12.0)*array(seps)**(-12.0) - 4.0*VDW_dict[type][0]*VDW_dict[type][1]**(6.0)*array(seps)**(-6.0) ) )

    return E_LJ

# # This function expects data to be a globally defined dictionary. The function
# # outputs the LJ contribution to the interaction energy of a pair interaction, summed over all configurations
# # contributions are only summed for atom types possessing the proper connections (e.g. there may be redundant H types
# # but for UA purposes it is important to resolve them by which type of carbon they are attached to. con_0 supplies the atomtype
# # type[0] should have a connection with, con_1 supplies the atomtype type[1] should have a connection with.)
# def E_LJ_Type_w_con_test(type,con_0,con_1,VDW_new):

#     global Data

#     # Initialize the energy
#     E_LJ = 0.0

#     # Iterate over all configurations
#     for d in Data.keys():

#         # If a contraint is present for type 0, find indices of atomtype type[0] that are bonded to the supplied contraint
#         if con_0 != []:
#             ind_con_0_a  = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == con_0 ]
#             ind_type_0_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[0] and True in [ 1 == j for j in Data[d]["adj_mat_a"][count_i][ind_con_0_a] ] ] 

#             ind_con_0_b  = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == con_0 ]
#             ind_type_0_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[0] and True in [ 1 == j for j in Data[d]["adj_mat_b"][count_i][ind_con_0_b] ] ] 
            
#         # Else, return indices of all type[0] atoms
#         else:
#             ind_type_0_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[0] ]
#             ind_type_0_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[0] ]

#         # If a constraint is present for type 1, find indices of atomtype type[1] that are bonded to the supplied contraint     
#         if con_1 != []:
#             ind_con_1_a  = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == con_1 ]
#             ind_type_1_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[1] and True in [ 1 == j for j in Data[d]["adj_mat_a"][count_i][ind_con_1_a] ] ] 

#             ind_con_1_b  = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == con_1 ]
#             ind_type_1_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[1] and True in [ 1 == j for j in Data[d]["adj_mat_b"][count_i][ind_con_1_b] ] ] 

#         # Else, return indices of all type[1] atoms
#         else:
#             ind_type_1_a = [ count_i for count_i,i in enumerate(Data[d]['types_a']) if i == type[1] ]
#             ind_type_1_b = [ count_i for count_i,i in enumerate(Data[d]['types_b']) if i == type[1] ]

#         # Find pair-wise separations for this configuration
#         seps = []
#         for i in ind_type_0_a:
#             for j in ind_type_1_b:
#                 seps += [Data[d]["r_dist"][i,j]]
#         if type[0] != type[1] or con_0 != con_1:
#             for i in ind_type_1_a:
#                 for j in ind_type_0_b:
#                     seps += [Data[d]["r_dist"][i,j]]

#         # Add up all LJ interations involving this type
#         E_LJ += sum( ( 4.0*VDW_new[type][0]*VDW_new[type][1]**(12.0)*array(seps)**(-12.0) - 4.0*VDW_new[type][0]*VDW_new[type][1]**(6.0)*array(seps)**(-6.0) ) )

#     return E_LJ

# This function expects data to be a globally defined dictionary. The function
# outputs the LJ contribution to the interaction energy of a pair interaction, summed over all configurations
def E_LJ_Type(type):

    global Data,VDW_dict

    # Initialize the energy
    E_LJ = 0.0

    # Iterate over all configurations
    for d in list(Data.keys()):

        # Add up all LJ interations involving this type
        if type in list(Data[d]["pairs"].keys()):
            E_LJ += sum( ( 4.0*VDW_dict[type][0]*VDW_dict[type][1]**(12.0)*Data[d]["pairs"][type]**(-12.0) - 4.0*VDW_dict[type][0]*VDW_dict[type][1]**(6.0)*Data[d]["pairs"][type]**(-6.0) ) )

        # Add up reverse combinations if applicable
        if type[0] != type[1] and (type[1],type[0]) in list(Data[d]["pairs"].keys()):
            E_LJ += sum( ( 4.0*VDW_dict[type][0]*VDW_dict[type][1]**(12.0)*Data[d]["pairs"][(type[1],type[0])]**(-12.0) - \
                           4.0*VDW_dict[type][0]*VDW_dict[type][1]**(6.0)*Data[d]["pairs"][(type[1],type[0])]**(-6.0) ) )
            
    return E_LJ

# Returns the LJ energy for each configuration (as an array) using the parameters currently defined in VDW_dict
# and the supplied parameters for the type being fit (fit_type is globally defined outside of the function). 
# OLD UNVECTORIZED FUNCTION
# def E_LJ_Fit(ind,eps_fit,sigma_fit):

#     # In python you need to expand the scope of global variables at every level
#     global VDW_dict,fit_type,Pair_min_dict

#     # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
#     # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#     if fit_type in Pair_min_dict.keys() and ( eps_fit < 0.001 or eps_fit > 30.0 or sigma_fit > 10.0 or sigma_fit < Pair_min_dict[fit_type]*2.**(-1./6.) ): return ones(len(ind))*1.E20        

#     # Initialize E_LJ for the current configuration
#     E_LJ = zeros(len(ind))

#     # Iterate over the configuration index
#     for count_d,d in enumerate(ind):

#         # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
#         for i in Data[d]["pairs"].keys():
#             if i == fit_type or (i[1],i[0]) == fit_type:
#                 E_LJ[count_d] += sum( ( 4.0*eps_fit*sigma_fit**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*eps_fit*sigma_fit**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )
#             else:
#                 E_LJ[count_d] += sum( ( 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )

#     return E_LJ


# Returns the LJ energy for each configuration (as an array) using the parameters currently defined in VDW_dict
# and the supplied parameters for the type being fit (fit_type is globally defined outside of the function). 
def E_LJ_Fit(ind,eps_fit,sigma_fit):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
    if fit_type in list(Pair_min_dict.keys()) and ( eps_fit < 0.001 or eps_fit > 30.0 or sigma_fit > 10.0 or sigma_fit < Pair_min_dict[fit_type]*2.**(-1./6.) ): return ones(len(ind))*1.E20        

    # Initialize E_LJ for the current configuration 
    E_LJ = zeros(len(ind))

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])
    for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == fit_type or (j[1],j[0]) == fit_type ) ]:        
        eps_array[i]   = eps_fit
        sigma_array[i] = sigma_fit

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        E_LJ[count_d] = sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) )

    return E_LJ
def E_LJ_Fit_custom(ind,VDW_dict_custom):

    # In python you need to expand the scope of global variables at every level
    global fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
    if fit_type in list(Pair_min_dict.keys()) and ( eps_fit < 0.001 or eps_fit > 30.0 or sigma_fit > 10.0 or sigma_fit < Pair_min_dict[fit_type]*2.**(-1./6.) ): return ones(len(ind))*1.E20        

    # Initialize E_LJ for the current configuration 
    E_LJ = zeros(len(ind))

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict_custom[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict_custom[i][1] for i in Data[ind[0]]["pair_type_vector"] ])
    for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == fit_type or (j[1],j[0]) == fit_type ) ]:        
        eps_array[i]   = eps_fit
        sigma_array[i] = sigma_fit

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        E_LJ[count_d] = sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) )

    return E_LJ

# Returns the BUCK energy for each configuration (as an array) using the parameters currently defined in VDW_dict
# and the supplied parameters for the type being fit (fit_type is globally defined outside of the function). 
def E_BUCK_Fit(ind,A,B,C):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
    if fit_type in list(Pair_min_dict.keys()) and ( ( A < 0.001 or B < 0.001 or C < 0.001 ) or True not in [ (A*exp(-i/B) - C*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ):
        return ones(len(ind))*1.E20        

    # Initialize E_BUCK for the current configuration
    E_BUCK = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            if i == fit_type or (i[1],i[0]) == fit_type:
                E_BUCK[count_d] += sum( ( A*exp(-Data[d]["pairs"][i]/B) - C*Data[d]["pairs"][i]**(-6.0) ) )
            else:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )

    return E_BUCK

# Returns the LJ energy for each configuration (as an array) using the parameters currently defined in VDW_dict
def E_LJ_configs(ind):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict

    # Initialize E_LJ for the current configuration
    E_LJ = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            E_LJ[count_d] += sum( ( 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )

    return E_LJ


# Function for iteratively-fitting a single eps value
# fit_type is a globally defined tuple
# OLD UNVECTORIZED FUNCTION
# def E_LJ_Fit_eps(ind,eps_fit):

#     # In python you need to expand the scope of global variables at every level
#     global VDW_dict,fit_type,Pair_min_dict

#     # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
#     # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#     if fit_type in Pair_min_dict.keys() and ( eps_fit < 0.001 or eps_fit > 30.0 ): return ones(len(ind))*1.E20        

#     # Initialize E_LJ for the current configuration
#     E_LJ = zeros(len(ind))

#     # Iterate over the configuration index
#     for count_d,d in enumerate(ind):

#         # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
#         for i in Data[d]["pairs"].keys():
#             if i == fit_type or (i[1],i[0]) == fit_type:
#                 E_LJ[count_d] += sum( ( 4.0*eps_fit*VDW_dict[i][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*eps_fit*VDW_dict[i][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )
#             else:
#                 E_LJ[count_d] += sum( ( 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )

#     return E_LJ


# Function for iteratively-fitting a single eps value
# fit_type is a globally defined tuple
def E_LJ_Fit_eps(ind,eps_fit):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
    if fit_type in list(Pair_min_dict.keys()) and ( eps_fit < 0.001 or eps_fit > 30.0 ): return ones(len(ind))*1.E20        

    # Initialize E_LJ for the current configuration
    E_LJ = zeros(len(ind))

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])
    for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == fit_type or (j[1],j[0]) == fit_type ) ]:        
        eps_array[i] = eps_fit

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        E_LJ[count_d] = sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) )

    return E_LJ


# Function for iteratively-fitting a single sigma value
# fit_type is a globally defined tuple
# OLD UNVECTORIZED FUNCTION
# def E_LJ_Fit_sigma(ind,sigma_fit):

#     # In python you need to expand the scope of global variables at every level
#     global VDW_dict,fit_type,Pair_min_dict

#     # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
#     # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#     if fit_type in Pair_min_dict.keys() and ( sigma_fit > 10.0 or sigma_fit < Pair_min_dict[fit_type]*2.**(-1./6.) ): return ones(len(ind))*1.E20        

#     # Initialize E_LJ for the current configuration
#     E_LJ = zeros(len(ind))

#     # Iterate over the configuration index
#     for count_d,d in enumerate(ind):

#         # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
#         for i in Data[d]["pairs"].keys():
#             if i == fit_type or (i[1],i[0]) == fit_type:
#                 E_LJ[count_d] += sum( ( 4.0*VDW_dict[i][0]*sigma_fit**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*sigma_fit**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )
#             else:
#                 E_LJ[count_d] += sum( ( 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )

#     return E_LJ


# Function for iteratively-fitting a single sigma value
# fit_type is a globally defined tuple
def E_LJ_Fit_sigma(ind,sigma_fit):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
    if fit_type in list(Pair_min_dict.keys()) and ( sigma_fit > 10.0 or sigma_fit < Pair_min_dict[fit_type]*2.**(-1./6.) ): return ones(len(ind))*1.E20        

    # Initialize E_LJ for the current configuration
    E_LJ = zeros(len(ind))

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])
    for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == fit_type or (j[1],j[0]) == fit_type ) ]:        
        sigma_array[i] = sigma_fit

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        E_LJ[count_d] = sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) )

    return E_LJ

# Function for iteratively-fitting a single eps value
# fit_type is a globally defined tuple
def E_BUCK_Fit_ABC(ind,A,B,C):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#    if fit_type in Pair_min_dict.keys() and ( A < 0.001 ): return ones(len(ind))*1.E20        
#    if fit_type in Pair_min_dict.keys() and ( ( A < 0.001 ) or ( True not in [ (A*exp(-i/VDW_dict[fit_type][1]) - VDW_dict[fit_type][2]*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
    if fit_type in list(Pair_min_dict.keys()) and ( ( A < 0 ) or ( B < 0 ) or ( C < 0 ) or ( True not in [ (A*exp(-i/B) - C*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
        return ones(len(ind))*1.E20        


    # Initialize E_BUCK for the current configuration
    E_BUCK = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            if i == fit_type or (i[1],i[0]) == fit_type:
                E_BUCK[count_d] += sum( ( A*exp(-Data[d]["pairs"][i]/B) - C*Data[d]["pairs"][i]**(-6.0) ) )
            else:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )

    return E_BUCK

# Function for iteratively-fitting a single eps value
# fit_type is a globally defined tuple
def E_BUCK_Fit_A(ind,A):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#    if fit_type in Pair_min_dict.keys() and ( A < 0.001 ): return ones(len(ind))*1.E20        
#    if fit_type in Pair_min_dict.keys() and ( ( A < 0.001 ) or ( True not in [ (A*exp(-i/VDW_dict[fit_type][1]) - VDW_dict[fit_type][2]*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
    if fit_type in list(Pair_min_dict.keys()) and ( ( A < 0 ) or ( True not in [ (A*exp(-i/VDW_dict[fit_type][1]) - VDW_dict[fit_type][2]*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
        return ones(len(ind))*1.E20        


    # Initialize E_BUCK for the current configuration
    E_BUCK = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            if i == fit_type or (i[1],i[0]) == fit_type:
                E_BUCK[count_d] += sum( ( A*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )
            else:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )

    return E_BUCK

def E_BUCK_Fit_B(ind,B):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#    if fit_type in Pair_min_dict.keys() and ( B < 0.001 ): return ones(len(ind))*1.E20        
#    if fit_type in Pair_min_dict.keys() and ( ( B < 0.001 ) or ( True not in [ (VDW_dict[fit_type][0]*exp(-i/B) - VDW_dict[fit_type][2]*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
    if fit_type in list(Pair_min_dict.keys()) and ( ( B < 0 ) or ( True not in [ (VDW_dict[fit_type][0]*exp(-i/B) - VDW_dict[fit_type][2]*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
        return ones(len(ind))*1.E20        


    # Initialize E_BUCK for the current configuration
    E_BUCK = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            if i == fit_type or (i[1],i[0]) == fit_type:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/B) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )
            else:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )

    return E_BUCK

def E_BUCK_Fit_C(ind,C):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
#    if fit_type in Pair_min_dict.keys() and ( C < 0.001 ): return ones(len(ind))*1.E20        
#    if fit_type in Pair_min_dict.keys() and ( ( C < 0.001 ) or ( True not in [ (VDW_dict[fit_type][0]*exp(-i/VDW_dict[fit_type][1]) - C*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
    if fit_type in list(Pair_min_dict.keys()) and ( ( C < 0 ) or ( True not in [ (VDW_dict[fit_type][0]*exp(-i/VDW_dict[fit_type][1]) - C*i**(-6.0) ) > 10.0 for i in arange(0.1,5.1,0.01) ] ) ):
        return ones(len(ind))*1.E20        

    # Initialize E_BUCK for the current configuration
    E_BUCK = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            if i == fit_type or (i[1],i[0]) == fit_type:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - C*Data[d]["pairs"][i]**(-6.0) ) )
            else:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )

    return E_BUCK


# Description: Initialize VDW_dict based on UFF parameters for the initial guess of the fit.
# OLD FUNCTION
def init_VDW_old(Data,fit_pairs,read_pairs,Pair_min_dict,fit,FF_db,VDW_dict_prev,UA=0):

    print("\nReading VDW parameters (*) or initializing based on UFF values and Lorentz-Berthelot mixing rules\n")

    # Initialize UFF parameters (tuple corresponds to eps,sigma pairs for each element)
    # Taken from UFF (Rappe et al. JACS 1992)
    # Note: LJ parameters in the table are specificed in the eps,r_min form rather than eps,sigma
    #       the conversion between r_min and sigma is sigma = r_min/2^(1/6)
    # Note: Units for sigma = angstroms and eps = kcal/mol 
    UFF_dict = { 1:(0.044,2.5711337005530193),  2:(0.056,2.1043027722474816),  3:(0.025,2.183592758161972),  4:(0.085,2.4455169812952313),\
                 5:(0.180,3.6375394661670053),  6:(0.105,3.4308509635584463),  7:(0.069,3.260689308393642),  8:(0.060,3.1181455134911875),\
                 9:(0.050,2.996983287824101),  10:(0.042,2.88918454292912),   11:(0.030,2.657550876212632), 12:(0.111,2.6914050275019648),\
                13:(0.505,4.008153332913386),  14:(0.402,3.82640999441276),   15:(0.305,3.694556984127987), 16:(0.274,3.594776327696269),\
                17:(0.227,3.5163772404999194), 18:(0.185,3.4459962417668324), 19:(0.035,3.396105913550973), 20:(0.238,3.0281647429590133),\
                21:(0.019,2.935511276272418),  22:(0.017,2.828603430095577),  23:(0.016,2.800985569833227), 24:(0.015,2.6931868249382456),\
                25:(0.013,2.6379511044135446), 26:(0.013,2.5942970672246677), 27:(0.014,2.558661118499054), 28:(0.015,2.5248069672097215),\
                29:(0.005,3.113691019900486),  30:(0.124,2.4615531582217574), 31:(0.415,3.904809081609107), 32:(0.379,3.813046513640652),\
                33:(0.309,3.7685015777336357), 34:(0.291,3.746229109780127),  35:(0.251,3.731974730289881), 36:(0.220,3.689211591819145),\
                37:(0.040,3.6651573264293558), 38:(0.235,3.2437622327489755), 39:(0.072,2.980056212179435), 40:(0.069,2.78316759547042),\
                41:(0.059,2.819694442914174),  42:(0.056,2.7190228877643157), 43:(0.048,2.670914356984738), 44:(0.056,2.6397329018498255),\
                45:(0.053,2.6094423454330538), 46:(0.048,2.5827153838888437), 47:(0.036,2.804549164705788), 48:(0.228,2.537279549263686),\
                49:(0.599,3.976080979060334),  50:(0.567,3.9128271700723705), 51:(0.449,3.937772334180300), 52:(0.398,3.982317270087316),\
                53:(0.339,4.009044231631527),  54:(0.332,3.923517954690054),  55:(0.045,4.024189509839913), 56:(0.364,3.2989979532736764),\
                72:(0.072,2.798312873678806),  73:(0.081,2.8241489365048755), 74:(0.067,2.734168165972701), 75:(0.066,2.631714813386562),\
                76:(0.037,2.7796040005978586), 77:(0.073,2.5301523595185635), 78:(0.080,2.453535069758495), 79:(0.039,2.9337294788361374),\
                80:(0.385,2.4098810325696176), 81:(0.680,3.872736727756055),  82:(0.663,3.828191791849038), 83:(0.518,3.893227398273283),\
                84:(0.325,4.195242063722858),  85:(0.284,4.231768911166611),  86:(0.248,4.245132391938716) }


    # Collect types that will be converted into UA
    if UA==1:
        UA_types = {}
        UA_H_types = []
        for i in list(Data.keys()):

            # Iterate over the molecule a adj_mat and identify hydrogens attached to carbons. 
            for count_j,j in enumerate(Data[i]["adj_mat_a"]):
                if int(Data[i]["types_a"][count_j].split('[')[1].split(']')[0]) != 6: continue
                UA_H_types += [ Data[i]["types_a"][count_k] for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_a'][count_k].split('[')[1].split(']')[0]) == 1 ]

            # Iterate over the molecule b adj_mat and identify hydrogens attached to carbons.
            for count_j,j in enumerate(Data[i]["adj_mat_b"]):
                if int(Data[i]["types_b"][count_j].split('[')[1].split(']')[0]) != 6: continue
                UA_H_types += [ Data[i]["types_b"][count_k] for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_b'][count_k].split('[')[1].split(']')[0]) == 1 ]

    # Initialize VDW_dict with fit types, first guess based on element types and Lorentz-Berthelot mixing rules
    VDW_dict = {}

    for i in fit_pairs:
#    for i in [ j for j in fit_pairs if j not in read_pairs]:

        # Grab the base atomic number from the atomtype
        type_1 = int(i[0].split('[')[1].split(']')[0])
        type_2 = int(i[1].split('[')[1].split(']')[0])
        
        # for UA fits, the eps of Hydrogen containing pairs are set to zero
        # so that they do not participate in the fitting. 
        if UA==1 and ( i[0] in UA_H_types or i[1] in UA_H_types ):
            eps    = 0.0
            sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0
            VDW_dict[i] = (eps,1.1*sigma)
            VDW_dict[(i[1],i[0])] = (eps,1.1*sigma)

            # Print diagnostic
            print("\tpair {:80s}: {:20s}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]])))

        # Non-hydrogen UA types only get their eps values updated. The set_UA_sigmas
        # function handles the assignment of UA sigma values
        elif UA==1:
            eps    = (UFF_dict[type_1][0]*UFF_dict[type_2][0])**(0.5)
            sigma  = VDW_dict_prev[i][1]
            VDW_dict[i] = (eps,sigma)
            VDW_dict[(i[1],i[0])] = (eps,sigma)

            # Print diagnostic
            print("\tpair {:80s}: {:20s}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]])))

        # For all atom fits, initial guesses for eps and sigma are based on scaled UFF parameters with a 
        # check to make sure that the radii aren't too small for the fit set. 
        else:
            eps    = (UFF_dict[type_1][0]*UFF_dict[type_2][0])**(0.5) 
            sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0
            if sigma < Pair_min_dict[i]*2.**(-1./6.): sigma = Pair_min_dict[i]*2.**(-1./6.)
            VDW_dict[i] = (eps,sigma)
            VDW_dict[(i[1],i[0])] = (eps,sigma)
        
            # Print diagnostic
            print("\tpair {:80s}: {:20s}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]])))

    # Initialize read_pairs with entries from the FF_db.
    for i in read_pairs:

        # Grab the base atomic number from the atomtype
        type_1 = int(i[0].split('[')[1].split(']')[0])
        type_2 = int(i[1].split('[')[1].split(']')[0])
        
        # for UA fits, the eps of Hydrogen containing pairs are set to zero
        # so that they do not participate in the fitting. 
        if UA==1 and ( i[0] in UA_H_types or i[1] in UA_H_types ):
            VDW_dict[i] = (0.0,1.0)
            VDW_dict[(i[1],i[0])] = (0.0,1.0)
            continue
        else:        

            # Add the pair type (along with the reverse combination) using the parameters found in the FF_db
            if (i[0],i[1],'lj') in list(FF_db["vdws"].keys()):
                VDW_dict[i] = (float(FF_db["vdws"][(i[0],i[1],'lj')][3]),float(FF_db["vdws"][(i[0],i[1],'lj')][4]))
                VDW_dict[(i[1],i[0])] = (float(FF_db["vdws"][(i[0],i[1],'lj')][3]),float(FF_db["vdws"][(i[0],i[1],'lj')][4]))            
            else:
                VDW_dict[i] = (float(FF_db["vdws"][(i[1],i[0],'lj')][3]),float(FF_db["vdws"][(i[1],i[0],'lj')][4]))
                VDW_dict[(i[1],i[0])] = (float(FF_db["vdws"][(i[1],i[0],'lj')][3]),float(FF_db["vdws"][(i[1],i[0],'lj')][4]))

            # Print diagnostic
            print("\tpair {:80s}: {:20s} *".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]])))

    # Perform conversions to buckingham style. Conversions follow the approach of Rappe et al. with a "form factor" of 12.0 
    # which reproduces the long-range LJ behavior
    if fit == 'buck':
        for i in list(VDW_dict.keys()):
            A = VDW_dict[i][0]*0.5*exp(12.0)
            B = VDW_dict[i][1]/12.0
            C = VDW_dict[i][0] * 2.0 * VDW_dict[i][1]**(6.0)
            VDW_dict[i] = (A,B,C)

    return VDW_dict

# Description: Initialize VDW_dict based on UFF parameters for the initial guess of the fit.
def init_VDW(Data,fit_pairs,read_pairs,Pair_min_dict,fit,FF_db,VDW_dict_prev,UA=0,avoid_read_flag=0,mixing_rule="lb",verbose=True):

    if verbose is True:
        if avoid_read_flag == 0 and mixing_rule in ["none","lb"]:
            print("\nReading VDW parameters or initializing based on UFF values and Lorentz-Berthelot mixing rules\n")
        elif avoid_read_flag == 0 and mixing_rule in ["wh"]:
            print("\nReading VDW parameters or initializing based on UFF values and Waldman-Hagler mixing rules\n")
        elif avoid_read_flag == 1 and mixing_rule in ["none","lb"]:
            print("\nInitializing VDW parameters based on UFF values and Lorentz-Berthelot mixing rules\n")
        elif avoid_read_flag == 1 and mixing_rule in ["wh"]:
            print("\nInitializing VDW parameters based on UFF values and Waldman-Hagler mixing rules\n")

        print("\t{:85s}  {:22s} {:20s}".format("pair_interaction","params","origin"))

    # Initialize UFF parameters (tuple corresponds to eps,sigma pairs for each element)
    # Taken from UFF (Rappe et al. JACS 1992)
    # Note: LJ parameters in the table are specificed in the eps,r_min form rather than eps,sigma
    #       the conversion between r_min and sigma is sigma = r_min/2^(1/6)
    # Note: Units for sigma = angstroms and eps = kcal/mol 
    UFF_dict = { 1:(0.044,2.5711337005530193),  2:(0.056,2.1043027722474816),  3:(0.025,2.183592758161972),  4:(0.085,2.4455169812952313),\
                 5:(0.180,3.6375394661670053),  6:(0.105,3.4308509635584463),  7:(0.069,3.260689308393642),  8:(0.060,3.1181455134911875),\
                 9:(0.050,2.996983287824101),  10:(0.042,2.88918454292912),   11:(0.030,2.657550876212632), 12:(0.111,2.6914050275019648),\
                13:(0.505,4.008153332913386),  14:(0.402,3.82640999441276),   15:(0.305,3.694556984127987), 16:(0.274,3.594776327696269),\
                17:(0.227,3.5163772404999194), 18:(0.185,3.4459962417668324), 19:(0.035,3.396105913550973), 20:(0.238,3.0281647429590133),\
                21:(0.019,2.935511276272418),  22:(0.017,2.828603430095577),  23:(0.016,2.800985569833227), 24:(0.015,2.6931868249382456),\
                25:(0.013,2.6379511044135446), 26:(0.013,2.5942970672246677), 27:(0.014,2.558661118499054), 28:(0.015,2.5248069672097215),\
                29:(0.005,3.113691019900486),  30:(0.124,2.4615531582217574), 31:(0.415,3.904809081609107), 32:(0.379,3.813046513640652),\
                33:(0.309,3.7685015777336357), 34:(0.291,3.746229109780127),  35:(0.251,3.731974730289881), 36:(0.220,3.689211591819145),\
                37:(0.040,3.6651573264293558), 38:(0.235,3.2437622327489755), 39:(0.072,2.980056212179435), 40:(0.069,2.78316759547042),\
                41:(0.059,2.819694442914174),  42:(0.056,2.7190228877643157), 43:(0.048,2.670914356984738), 44:(0.056,2.6397329018498255),\
                45:(0.053,2.6094423454330538), 46:(0.048,2.5827153838888437), 47:(0.036,2.804549164705788), 48:(0.228,2.537279549263686),\
                49:(0.599,3.976080979060334),  50:(0.567,3.9128271700723705), 51:(0.449,3.937772334180300), 52:(0.398,3.982317270087316),\
                53:(0.339,4.009044231631527),  54:(0.332,3.923517954690054),  55:(0.045,4.024189509839913), 56:(0.364,3.2989979532736764),\
                72:(0.072,2.798312873678806),  73:(0.081,2.8241489365048755), 74:(0.067,2.734168165972701), 75:(0.066,2.631714813386562),\
                76:(0.037,2.7796040005978586), 77:(0.073,2.5301523595185635), 78:(0.080,2.453535069758495), 79:(0.039,2.9337294788361374),\
                80:(0.385,2.4098810325696176), 81:(0.680,3.872736727756055),  82:(0.663,3.828191791849038), 83:(0.518,3.893227398273283),\
                84:(0.325,4.195242063722858),  85:(0.284,4.231768911166611),  86:(0.248,4.245132391938716) }

    # Initialize initial_guess dictionary based on the initial_params.db file (if present). 
    # The default is to use these parameters in place of UFF if the file is present.     
    INIT_dict = {}
    if os.path.isfile('initial_params.db') and avoid_read_flag == 0:
        with open('initial_params.db','r') as f:
            for lines in f:
                fields = lines.split()
                if len(fields) > 0 and fields[0] == "vdw":
                    INIT_dict[(fields[1],fields[2])] = [ float(i) for i in fields[4:] ]
                    INIT_dict[(fields[2],fields[1])] = [ float(i) for i in fields[4:] ]

    # Collect types that will be converted into UA
    if UA==1:
        UA_types = {}
        UA_H_types = []
        for i in list(Data.keys()):

            # Iterate over the molecule a adj_mat and identify hydrogens attached to carbons. 
            for count_j,j in enumerate(Data[i]["adj_mat_a"]):
                if int(Data[i]["types_a"][count_j].split('[')[1].split(']')[0]) != 6: continue
                UA_H_types += [ Data[i]["types_a"][count_k] for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_a'][count_k].split('[')[1].split(']')[0]) == 1 ]

            # Iterate over the molecule b adj_mat and identify hydrogens attached to carbons.
            for count_j,j in enumerate(Data[i]["adj_mat_b"]):
                if int(Data[i]["types_b"][count_j].split('[')[1].split(']')[0]) != 6: continue
                UA_H_types += [ Data[i]["types_b"][count_k] for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_b'][count_k].split('[')[1].split(']')[0]) == 1 ]

    # Initialize VDW_dict with fit types, first guess based on element types and Lorentz-Berthelot mixing rules
    VDW_dict = {}
    for i in fit_pairs:

        # Grab the base atomic number from the atomtype
        type_1 = int(i[0].split('[')[1].split(']')[0])
        type_2 = int(i[1].split('[')[1].split(']')[0])
        
        # for UA fits, the eps of Hydrogen containing pairs are set to zero
        # so that they do not participate in the fitting. 
        if UA==1 and ( i[0] in UA_H_types or i[1] in UA_H_types ):
            eps    = 0.0
            sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0
            VDW_dict[i] = (eps,sigma)
            VDW_dict[(i[1],i[0])] = (eps,sigma)

            # Print diagnostic
            if verbose is True:
                print("\tpair {:80s}: {:20s} {:20s}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"Zeroed"))

        # Non-hydrogen UA types read their values from the AA-fit
        elif UA==1:
            eps    = VDW_dict_prev[i][0]
            sigma  = VDW_dict_prev[i][1]
            VDW_dict[i] = (eps,sigma)
            VDW_dict[(i[1],i[0])] = (eps,sigma)

            # Print diagnostic
            if verbose is True:
                print("\tpair {:80s}: {:20s} {:20s}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"AA-Fit"))

        # For all atom fits, initial guesses for eps and sigma are based on scaled UFF parameters with a 
        # check to make sure that the radii aren't too small for the fit set. 
        else:
            if i in INIT_dict:
                eps    = INIT_dict[i][0]
                sigma  = INIT_dict[i][1]
                if sigma < Pair_min_dict[i]*2.**(-1./6.): sigma = Pair_min_dict[i]*2.**(-1./6.)
                VDW_dict[i] = (eps,sigma)
                VDW_dict[(i[1],i[0])] = (eps,sigma)
                if verbose is True:
                    print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"initial_params.db"))
            else:
                if mixing_rule in ["lb",'none']:
                    eps    = (UFF_dict[type_1][0]*UFF_dict[type_2][0])**(0.5) 
                    sigma  = (UFF_dict[type_1][1]+UFF_dict[type_2][1])/2.0
                elif mixing_rule == "wh":
                    sigma  = ((UFF_dict[type_1][1]**(6.0)+UFF_dict[type_2][1]**(6.0))/2.0)**(1.0/6.0)
                    eps    = (UFF_dict[type_1][0]*UFF_dict[type_1][1]**(6.0) * UFF_dict[type_2][0]*UFF_dict[type_2][1]**(6.0) )**(0.5) / sigma**(6.0)
                else: 
                    print("ERROR in init_VDW: {} is not an implemented mixing rule. Exiting...".format(mixing_rule))
                    quit()
                if sigma < Pair_min_dict[i]*2.**(-1./6.): sigma = Pair_min_dict[i]*2.**(-1./6.)
                VDW_dict[i] = (eps,sigma)
                VDW_dict[(i[1],i[0])] = (eps,sigma)
        
                # Print diagnostic
                if verbose is True:
                    print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"UFF"))

    # For diagnostic purposes, the fixed_params.db and user_supplied FF files are reparsed here in order to 
    # print the proper origin information for the user.
    # Read in fixed parameters from the fixed_params.db file
    FIXED_dict = {"vdws":{}}
    if os.path.isfile('fixed_params.db'):         
        FIXED_dict = get_FF_data(['fixed_params.db'])

    # If a FF_db file is supplied then the program avoids fitting any of the
    # pairs that are pairs are in the database
    READ_dict = {"vdws":{}}
    if FF_db != []:        
        READ_dict = get_FF_data(FF_db,mixing_rule=mixing_rule)

    # Initialize read_pairs with entries from the FF_db.
    for i in read_pairs:

        # Grab the base atomic number from the atomtype
        type_1 = int(i[0].split('[')[1].split(']')[0])
        type_2 = int(i[1].split('[')[1].split(']')[0])
        
        # for UA fits, the eps of Hydrogen containing pairs are set to zero
        # so that they do not participate in the fitting. 
        if UA==1 and ( i[0] in UA_H_types or i[1] in UA_H_types ):
            VDW_dict[i] = (0.0,1.0)
            VDW_dict[(i[1],i[0])] = (0.0,1.0)
            continue
        else:        

            # Mixing rule protocols
            if mixing_rule != "none": 

                # Use mixing rule protocol for read parameters
                if (i[0],i[0],'lj') in list(READ_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(READ_dict["vdws"].keys()):

                    eps_1 = float(READ_dict["vdws"][(i[0],i[0],'lj')][3])
                    sigma_1 = float(READ_dict["vdws"][(i[0],i[0],'lj')][4])
                    eps_2 = float(READ_dict["vdws"][(i[1],i[1],'lj')][3])
                    sigma_2 = float(READ_dict["vdws"][(i[1],i[1],'lj')][4])

                # Use mixing rule protocol for fixed parameters
                elif (i[0],i[0],'lj') in list(FIXED_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(FIXED_dict["vdws"].keys()):

                    eps_1 = float(FIXED_dict["vdws"][(i[0],i[0],'lj')][3])
                    sigma_1 = float(FIXED_dict["vdws"][(i[0],i[0],'lj')][4])
                    eps_2 = float(FIXED_dict["vdws"][(i[1],i[1],'lj')][3])
                    sigma_2 = float(FIXED_dict["vdws"][(i[1],i[1],'lj')][4])

                # Use mixing rule protocol for fixed/read parameters
                elif (i[0],i[0],'lj') in list(FIXED_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(READ_dict["vdws"].keys()):

                    eps_1 = float(FIXED_dict["vdws"][(i[0],i[0],'lj')][3])
                    sigma_1 = float(FIXED_dict["vdws"][(i[0],i[0],'lj')][4])
                    eps_2 = float(READ_dict["vdws"][(i[1],i[1],'lj')][3])
                    sigma_2 = float(READ_dict["vdws"][(i[1],i[1],'lj')][4])

                # Use mixing rule protocol for read/fixed parameters
                elif (i[0],i[0],'lj') in list(READ_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(FIXED_dict["vdws"].keys()):

                    eps_1 = float(READ_dict["vdws"][(i[0],i[0],'lj')][3])
                    sigma_1 = float(READ_dict["vdws"][(i[0],i[0],'lj')][4])
                    eps_2 = float(FIXED_dict["vdws"][(i[1],i[1],'lj')][3])
                    sigma_2 = float(FIXED_dict["vdws"][(i[1],i[1],'lj')][4])

                # Apply mixing rule
                if mixing_rule in ["lb"]:
                    eps    = (eps_1*eps_2)**(0.5)
                    sigma  = (sigma_1+sigma_2)/2.0
                elif mixing_rule == "wh":
                    sigma  =  ( ( sigma_1**(6.0) + sigma_2**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps    =  ( eps_1*sigma_1**(6.0) * eps_2*sigma_2**(6.0) )**(0.5) / sigma**(6.0)    
                else: 
                    print("ERROR in init_VDW: {} is not an implemented mixing rule. Exiting...".format(mixing_rule))
                    quit()
                VDW_dict[i] = (eps,sigma)
                VDW_dict[(i[1],i[0])] = (eps,sigma)

                # Print commands
                if verbose is True:
                    if (i[0],i[0],'lj') in list(READ_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(READ_dict["vdws"].keys()):
                        print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"supplied_file"))
                    elif (i[0],i[0],'lj') in list(FIXED_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(FIXED_dict["vdws"].keys()):
                        print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"fixed_params.db"))
                    elif (i[0],i[0],'lj') in list(FIXED_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(READ_dict["vdws"].keys()):
                        print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"fixed/supplied mixing"))
                    elif (i[0],i[0],'lj') in list(READ_dict["vdws"].keys()) and (i[1],i[1],'lj') in list(FIXED_dict["vdws"].keys()):
                        print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"fixed/supplied mixing"))



            # Add the pair type (along with the reverse combination) using the parameters found in the READ_dict
            elif (i[0],i[1],'lj') in list(READ_dict["vdws"].keys()) or (i[1],i[0],'lj') in list(READ_dict["vdws"].keys()):
                if (i[0],i[1],'lj') in list(READ_dict["vdws"].keys()):
                    VDW_dict[i] = (float(READ_dict["vdws"][(i[0],i[1],'lj')][3]),float(READ_dict["vdws"][(i[0],i[1],'lj')][4]))
                    VDW_dict[(i[1],i[0])] = (float(READ_dict["vdws"][(i[0],i[1],'lj')][3]),float(READ_dict["vdws"][(i[0],i[1],'lj')][4]))            

                else:
                    VDW_dict[i] = (float(READ_dict["vdws"][(i[1],i[0],'lj')][3]),float(READ_dict["vdws"][(i[1],i[0],'lj')][4]))
                    VDW_dict[(i[1],i[0])] = (float(READ_dict["vdws"][(i[1],i[0],'lj')][3]),float(READ_dict["vdws"][(i[1],i[0],'lj')][4]))

                if verbose is True:
                    print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"supplied_file"))

            # Add the pair type (along with the reverse combination) using the parameters found in the FIXED_dict
            elif (i[0],i[1],'lj') in list(FIXED_dict["vdws"].keys()) or (i[1],i[0],'lj') in list(FIXED_dict["vdws"].keys()):

                if (i[0],i[1],'lj') in list(FIXED_dict["vdws"].keys()):
                    VDW_dict[i] = (float(FIXED_dict["vdws"][(i[0],i[1],'lj')][3]),float(FIXED_dict["vdws"][(i[0],i[1],'lj')][4]))
                    VDW_dict[(i[1],i[0])] = (float(FIXED_dict["vdws"][(i[0],i[1],'lj')][3]),float(FIXED_dict["vdws"][(i[0],i[1],'lj')][4]))            

                else:
                    VDW_dict[i] = (float(FIXED_dict["vdws"][(i[1],i[0],'lj')][3]),float(FIXED_dict["vdws"][(i[1],i[0],'lj')][4]))
                    VDW_dict[(i[1],i[0])] = (float(FIXED_dict["vdws"][(i[1],i[0],'lj')][3]),float(FIXED_dict["vdws"][(i[1],i[0],'lj')][4]))

                if verbose is True:
                    print("\tpair {:80s}: {:20s} {:20s}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),"fixed_params.db"))

            else:
                print("ERROR in init_VDW: Information for pair type {} was expected in the the fixed_params.db".format((i[0],i[1],'lj')))
                print("                   and/or the supplied FF files, but it could not be found. Please check that the parameters exist. Exiting...")
                quit()


    # Perform conversions to buckingham style. Conversions follow the approach of Rappe et al. 
    # with a "form factor" of 12.0 which reproduces the long-range LJ behavior
    if fit == 'buck':
        for i in list(VDW_dict.keys()):
            A = VDW_dict[i][0]*0.5*exp(12.0)
            B = VDW_dict[i][1]/12.0
            C = VDW_dict[i][0] * 2.0 * VDW_dict[i][1]**(6.0)
            VDW_dict[i] = (A,B,C)

    return VDW_dict


# Description: This function drives the pair-wise fit of the VDW parameters. It accepts
#              various arguments that modify the convergence criteria and optimization
#              algorithm. 
# 
# Inputs:      x_vals            : a list holding the keys in Data for configurations being fit
#              y_vals            : a list holding the interaction energies for each configuration being fit 
#              Pairs             : a list of the unique pair-types being fit.
#              weight            : a scalar that controls the weight of the update at each step.
#              delta_xhi2_thresh : a scalar that controls the xhi^2 convergence threshold for the fit. 
#              min_cycles        : an integer that controls the minimum number of fit cycles to perform.
#              max_cycles        : an integer that controls the maximum number of fit cycles to perform. 
# 
# Returns:     fit_vals          : a list of arrays, each holding the FF based interaction energies in each cycle of the fit.
#
def lstsq_fit(x_vals,y_vals,Pairs,fit):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # Calculate the Fit xhi2 using UFF parameters (Here E_LJ_Fit is just a function for calculating the fit energies of all configurations)
    fit_type = 0 # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
    fit_vals = [ E_LJ_Fit(x_vals,0.0,0.0) ] 
    xhi2_initial = mean((y_vals-fit_vals[0])**2) 

    b = y_vals                       # The binding energy of each configuration is used as the fit vector
    A = zeros([len(b),len(Pairs)*2]) # For the lj fit there are 2 parameters per pair type

    print("Buildng A matrix...")
    for count_i,i in enumerate(x_vals):
        for count_j,j in enumerate(Pairs):
            if j[0]==j[1]: 
               A[count_i,(count_j*2)+0] += sum(Data[i]["pairs"][j]**(-12.0))
               A[count_i,(count_j*2)+1] -= sum(Data[i]["pairs"][j]**(-6.0))
            else:
               A[count_i,(count_j*2)+0] += (sum(Data[i]["pairs"][j]**(-12.0))+ sum(Data[i]["pairs"][(j[1],j[0])]**(-12.0)))
               A[count_i,(count_j*2)+1] -= (sum(Data[i]["pairs"][j]**(-6.0))+ sum(Data[i]["pairs"][(j[1],j[0])]**(-6.0)))
            #A[count_i,(count_j*2)+0] += sum(Data[i]["pairs"][j]**(-12.0))
            #A[count_i,(count_j*2)+1] -= sum(Data[i]["pairs"][j]**(-6.0))

    # Check for rank deficiency
    rank = matrix_rank(A)
    if rank < len(A[0]): print("A is rank deficient (rank: {}; number of fit parameters: {})".format(rank,len(A[0])))
    else: print("A if full rank (rank: {}; number of fit parameters: {})".format(rank,len(A[0])))

    # Calculate values using lstsq with positive constraints on the fit variables
    #r3 = lsq_linear(A,b,bounds=(0,inf),lsq_solver='exact')['fun']
    #x3 = lsq_linear(A,b,bounds=(0,inf),lsq_solver='exact')['x']
    x, r = nnls(A,b)
    #x2 = np.linalg.lstsq(A, b, rcond=None)[0]
    #x = np.linalg.lstsq(A, b, rcond=None)[0]
    #r2 = np.linalg.lstsq(A, b, rcond=None)[1]
    #r = [i**2 for i in (np.dot(A,x)-b)]
    #print(mean(r))
    #r2 = [i**2 for i in (np.dot(A,x2)-b)]
    #print(mean(r2))
    #r3 = [i**2 for i in (np.dot(A,x3)-b)]
    #print(mean(r))
    
    
    ##### these two results are the same
    x[where(x==0.0)]=0.001
    #r = [i**2 for i in (np.dot(A,x)-b)]
    #print("hey")
    #print(mean(r))

    #for count_i,i  in enumerate(Pairs):
    #    x[count_i*2+0] = 4*VDW_dict[i][0]*VDW_dict[i][1]**12
    #    x[count_i*2+1] = 4*VDW_dict[i][0]*VDW_dict[i][1]**6

    #resid = [i**2 for i in (np.dot(A,x)-b)]
    #print("global")
    #print(mean(resid))



    # Update VDW_dict with a lstsq parameters
    for count_i,i in enumerate(Pairs):
        eps = (x[count_i*2+1]**(2.0))/(4.0*x[count_i*2+0]) # epsilon is equal to B^2/4*A
        sigma = (x[count_i*2+0]/x[count_i*2+1])**(1.0/6.0) # sigma is equal to (A/B)^(1/6)
    #    A = 4*eps*sigma**12
    #    B = 4*eps*sigma**6
    #    print(A-x[count_i*2+0])
    #    print(B-x[count_i*2+1])
    #    if eps < 0.00001: sigma=0.0
    #    if eps<0.00001: eps = 0.0
    #    if eps >30.0: eps=0.0;sigma=0.0
        
        print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<10.4f}".format(j) for j in [eps,sigma]])))
        VDW_dict[i] = (eps,sigma)
        VDW_dict[(i[1],i[0])] = (eps,sigma)

    # Calculate the Fit xhi2 using the lstsq parameters (Here E_LJ_Fit is just a function for calculating the fit energies of all configurations)
    fit_type = 0 # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
    fit_vals += [ E_LJ_Fit(x_vals,0.0,0.0) ] 
    xhi2_lstsq = mean((y_vals-fit_vals[1])**2) 

    xhi_total = [xhi2_initial, xhi2_lstsq]
    
    # Print diagnostic
    print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2_initial:",xhi2_initial))
    print("\t{:<12s} {} (kcal/mol)^2\n".format("xhi2_lstsq:",xhi2_lstsq))
    

    return fit_vals,xhi_total


# Description: This function drives the pair-wise fit of the VDW parameters. It accepts
#              various arguments that modify the convergence criteria and optimization
#              algorithm. 
# 
# Inputs:      x_vals            : a list holding the keys in Data for configurations being fit
#              y_vals            : a list holding the interaction energies for each configuration being fit 
#              Pairs             : a list of the unique pair-types being fit.
#              weight            : a scalar that controls the weight of the update at each step.
#              delta_xhi2_thresh : a scalar that controls the xhi^2 convergence threshold for the fit. 
#              min_cycles        : an integer that controls the minimum number of fit cycles to perform.
#              max_cycles        : an integer that controls the maximum number of fit cycles to perform. 
# 
# Returns:     fit_vals          : a list of arrays, each holding the FF based interaction energies in each cycle of the fit.
#
def global_fit(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh=0.00001,min_cycles=100,max_cycles=1000,mixing_rule='none',L2_s=0.0,L2_e=0.0,VDW_0=None,outlier_option=False):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict,Data

    # BOTH OF THE SELF TERMS NEED TO BE DEFINED FOR MIXING RULES TO BE SENSIBLY APPLIED
    if mixing_rule in ['wh','lb']:
        # Pairs = [ (k,k) for k in set([ j for i in Pairs for j in i ]) ] # OLD
        Pairs = [ i for i in Pairs if i[0] == i[1] ] # NEW
    elif mixing_rule != "none":
        print("ERROR in global_fit: {} is not an implemented mixing rule. Exiting...".format(mixing_rule))
        quit()

    # Intialize cycle count and array of fit values
    fit_type = 0 # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
    cycle_count = 0

    if fit == 'lj':
        fit_vals = [ E_LJ_Fit(x_vals,0.0,0.0) ] * (max_cycles+5)
    elif fit == 'buck':
        fit_vals = [ E_BUCK_Fit(x_vals,0.0,0.0,0.0) ] * (max_cycles+5)
    else:
        print("ERROR in global_fit: {} is not an implemented option for the fit variable. Exiting...".format(fit))
        quit()
    xhi2_previous = mean((y_vals-fit_vals[cycle_count])**2) 

    # If no pairs are being fit (as happens if only UA types are being fit and AA are being read from file) just return the fit_vals
    if len(Pairs) == 0: print(" "); return fit_vals[:cycle_count+1]

    # Check the rank of the fit data by assembling the lstsq fit matrix
    if fit in ['lj','buck'] :
        A = zeros([len(x_vals),len(Pairs)*2]) # For the lj fit there are 2 linearly parameters per pair type (A,B). The rank of the LJ matrix is used in place of the buck 
                                              # A matrix owing to the non-linearity of the buckingham potential
    # Build up A matrix
    for count_i,i in enumerate(x_vals):
        for count_j,j in enumerate(Pairs):

            if fit in ['lj','buck']:
                # Pairs only contains I J or J I combinations whereas Data[i]["pairs"] has both combinations
                # To be consistent the contributions from both I J and J I configurations need to be added to 
                # the elements of A
                tmp_6 = 0.0
                tmp_12 = 0.0

                if (j[0],j[1]) in list(Data[i]["pairs"].keys()):
                    tmp_6  += sum(Data[i]["pairs"][(j[0],j[1])]**(-6.0))
                    tmp_12 += sum(Data[i]["pairs"][(j[0],j[1])]**(-12.0))
                if j[0] != j[1] and (j[1],j[0]) in list(Data[i]["pairs"].keys()):
                    tmp_6  += sum(Data[i]["pairs"][(j[1],j[0])]**(-6.0))
                    tmp_12 += sum(Data[i]["pairs"][(j[1],j[0])]**(-12.0))

                A[count_i,(count_j*2)+0] += tmp_12
                A[count_i,(count_j*2)+1] -= tmp_6

    # Calculate the rank of the fit matrix
    rank = matrix_rank(A)

    # Print diagnostics
    print("\nInitializing L-BFGS-B fit of pair-wise interactions:\n")
    print("\t{:<30s} {}".format("delta_xhi2_thresh:",delta_xhi2_thresh))
    print("\t{:<30s} {}".format("mixing rule:",mixing_rule))
    print("\t{:<30s} {}".format("Fit rank:",rank))
    print("\t{:<30s} {}".format("Fit parameters:",len(A[0])))
    print("\t{:<30s} {}".format("L2 sigma weight:",L2_s))
    print("\t{:<30s} {}".format("L2 epsilon weight:",L2_e))
    print("\t{:<30s} {:<12.6g}".format("Fit Tolerance (kcal/mol):",delta_xhi2_thresh))

    # If Reference VDW dictionary isn't supplied, then the VDW_dict is copied and used for regularization
    if VDW_0 is None:
        VDW_0 = deepcopy(VDW_dict)
    # Initialize anonymous fit function
    if outlier_option:
        Mi,outlier = detect_outlier(y_vals)
        print("Number of outliers removed: {}".format(len([out for out in outlier if out])))
        fit_func = lambda x: global_fit_LJ_outlier(*x,ind=x_vals,outlier=outlier,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                      Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=L2_s,L2_eps=L2_e,VDW_0=VDW_0)
    else:
        
        fit_func = lambda x: global_fit_LJ(*x,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                      Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=L2_s,L2_eps=L2_e,VDW_0=VDW_0)

#    fit_func = lambda x: global_fit_LJ(*x,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_weight=L2_weight)

    # Initialize guesses and parameter bounds
    initial_guess = []
    bounds = []
    for i in Pairs:
#        initial_guess += [ j for j in VDW_dict[i] ] # previous method's result, from initial_params.db in vdw folder if there's one
        initial_guess += [ j for j in VDW_0[i] ] #UFF GUESS
#        initial_guess += [ 0.002,Pair_min_dict[i] ] # MIN GUESS
        bounds += [ (0.001,30.0),(Pair_min_dict[i]*2.**(-1./6.),10.0) ]
    fit_vals[0] = E_LJ_Fit_custom(x_vals,VDW_0)

    # Perform the fit
    t0 = time.time()
    if fit.lower() == "lj":
        results = minimize(fit_func,initial_guess,bounds=bounds,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':delta_xhi2_thresh})
        new_params = results.x

    # Print diagnostics
    if results.success == True:
        print("\t{:<30s} {:<12s}".format("Fit termination condition:","Converged"))
    else:
        print("\t{:<30s} {:<12s}".format("Fit termination condition:","Unconverged"))


    xhi_total = []
    xhi_total.append(global_fit_LJ(*initial_guess,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0))
    xhi_total.append( global_fit_LJ(*new_params,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0))

    print("\t{:<30s} {:<12.6f}".format("Initial xhi^2 (kcal/mol):",global_fit_LJ(*initial_guess,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                                                                 Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0)))
    print("\t{:<30s} {:<12.6f}".format("Final xhi^2 (kcal/mol):",global_fit_LJ(*new_params,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                                                               Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0)))
    print("\t{:<30s} {:<12.6f}".format("Time for completion (s):",time.time()-t0))

    # Update the VDW_dict    
    print("\nParameter update summary:\n")
    counter = 0
    for i in Pairs:      
        #print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<15.4f}".format(j) for j in new_params[counter:counter+len(VDW_dict[i])] ])))
        print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(j) for j in initial_guess[counter:counter+len(VDW_dict[i])]]),", ".join([ "{:<15.4f}".format(j) for j in new_params[counter:counter+len(VDW_dict[i])] ])))
        VDW_dict[i] = tuple([ j for j in new_params[counter:counter+len(VDW_dict[i])] ])
        VDW_dict[(i[1],i[0])] = VDW_dict[i]
        counter += len(VDW_dict[i])

    # Update cross terms if mixing rules were used
    if mixing_rule == "wh":        
        set_of_types = set([ j for i in Pairs for j in i ])
        updated_pairs = [ i for i in list(VDW_dict.keys()) if i[0] > i[1] and ( i[0] in set_of_types or i[1] in set_of_types ) ]
        for i in updated_pairs:
            #old = deepcopy(VDW_dict[i])
            old = deepcopy(VDW_0[i])
            sigma = (( VDW_dict[(i[0],i[0])][1]**(6.0) + VDW_dict[(i[1],i[1])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
            VDW_dict[i] = ( ( VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[0],i[0])][1]**(6.0) * VDW_dict[(i[1],i[1])][0]*VDW_dict[(i[1],i[1])][1]**(6.0) )**(0.5) / sigma**(6.0)  , sigma )
            VDW_dict[(i[1],i[0])] = VDW_dict[i]
            print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(k) for k in old ]),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i] ])))
    if mixing_rule == "lb":
        set_of_types = set([ j for i in Pairs for j in i ])
        updated_pairs = [ i for i in list(VDW_dict.keys()) if i[0] > i[1] and ( i[0] in set_of_types or i[1] in set_of_types ) ]
        for i in updated_pairs:
            #old = deepcopy(VDW_dict[i])
            old = deepcopy(VDW_0[i])
            VDW_dict[i] = ( ( VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[1],i[1])][0] )**(0.5) , (VDW_dict[(i[0],i[0])][1]+VDW_dict[(i[1],i[1])][1]) / 2.0)
            VDW_dict[(i[1],i[0])] = VDW_dict[i]
            print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(k) for k in old ]),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i] ])))
                    
    print(" ")

    # Calculate fit values and xhi2 with final fit parameters
    if fit.lower() == "lj":
        fit_vals[1] = E_LJ_Fit(x_vals,0.0,0.0)
    return fit_vals[:2],xhi_total

def global_fit_Boltz(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh=0.00001,min_cycles=100,max_cycles=1000,mixing_rule='none',L2_s=0.0,L2_e=0.0,VDW_0=None,outlier_option=False,E_config_tot=[],weighting_function='Boltzmann'):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict,Data

    # BOTH OF THE SELF TERMS NEED TO BE DEFINED FOR MIXING RULES TO BE SENSIBLY APPLIED
    if mixing_rule in ['wh','lb']:
        # Pairs = [ (k,k) for k in set([ j for i in Pairs for j in i ]) ] # OLD
        Pairs = [ i for i in Pairs if i[0] == i[1] ] # NEW
    elif mixing_rule != "none":
        print("ERROR in global_fit: {} is not an implemented mixing rule. Exiting...".format(mixing_rule))
        quit()

    # Intialize cycle count and array of fit values
    fit_type = 0 # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
    cycle_count = 0

    if fit == 'lj':
        fit_vals = [ E_LJ_Fit(x_vals,0.0,0.0) ] * (max_cycles+5)
    elif fit == 'buck':
        fit_vals = [ E_BUCK_Fit(x_vals,0.0,0.0,0.0) ] * (max_cycles+5)
    else:
        print("ERROR in global_fit: {} is not an implemented option for the fit variable. Exiting...".format(fit))
        quit()
    xhi2_previous = mean((y_vals-fit_vals[cycle_count])**2) 

    # If no pairs are being fit (as happens if only UA types are being fit and AA are being read from file) just return the fit_vals
    if len(Pairs) == 0: print(" "); return fit_vals[:cycle_count+1]

    # Check the rank of the fit data by assembling the lstsq fit matrix
    if fit in ['lj','buck'] :
        A = zeros([len(x_vals),len(Pairs)*2]) # For the lj fit there are 2 linearly parameters per pair type (A,B). The rank of the LJ matrix is used in place of the buck 
                                              # A matrix owing to the non-linearity of the buckingham potential
    # Build up A matrix
    for count_i,i in enumerate(x_vals):
        for count_j,j in enumerate(Pairs):

            if fit in ['lj','buck']:
                # Pairs only contains I J or J I combinations whereas Data[i]["pairs"] has both combinations
                # To be consistent the contributions from both I J and J I configurations need to be added to 
                # the elements of A
                tmp_6 = 0.0
                tmp_12 = 0.0

                if (j[0],j[1]) in list(Data[i]["pairs"].keys()):
                    tmp_6  += sum(Data[i]["pairs"][(j[0],j[1])]**(-6.0))
                    tmp_12 += sum(Data[i]["pairs"][(j[0],j[1])]**(-12.0))
                if j[0] != j[1] and (j[1],j[0]) in list(Data[i]["pairs"].keys()):
                    tmp_6  += sum(Data[i]["pairs"][(j[1],j[0])]**(-6.0))
                    tmp_12 += sum(Data[i]["pairs"][(j[1],j[0])]**(-12.0))

                A[count_i,(count_j*2)+0] += tmp_12
                A[count_i,(count_j*2)+1] -= tmp_6

    # Calculate the rank of the fit matrix
    rank = matrix_rank(A)

    # Print diagnostics
    print("\nInitializing L-BFGS-B fit of pair-wise interactions:\n")
    print("\t{:<30s} {}".format("delta_xhi2_thresh:",delta_xhi2_thresh))
    print("\t{:<30s} {}".format("mixing rule:",mixing_rule))
    print("\t{:<30s} {}".format("Fit rank:",rank))
    print("\t{:<30s} {}".format("Fit parameters:",len(A[0])))
    print("\t{:<30s} {}".format("L2 sigma weight:",L2_s))
    print("\t{:<30s} {}".format("L2 epsilon weight:",L2_e))
    print("\t{:<30s} {:<12.6g}".format("Fit Tolerance (kcal/mol):",delta_xhi2_thresh))

    # If Reference VDW dictionary isn't supplied, then the VDW_dict is copied and used for regularization
    if VDW_0 is None:
        VDW_0 = deepcopy(VDW_dict)

    # Initialize anonymous fit function
    if outlier_option:
        Mi,outlier = detect_outlier(y_vals)
        print("Number of outliers removed from fitting: {}".format(len([out for out in outlier if out])))
    else: outlier = []
    fit_func = lambda x: global_fit_LJ_Boltz(*x,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=L2_s,L2_eps=L2_e,VDW_0=VDW_0,E_config_tot=E_config_tot,outlier=outlier,weighting_function=weighting_function)
#    fit_func = lambda x: global_fit_LJ(*x,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_weight=L2_weight)

    # Initialize guesses and parameter bounds
    initial_guess = []
    bounds = []
    for i in Pairs:
#        initial_guess += [ j for j in VDW_dict[i] ] # UFF GUESS
        initial_guess += [ j for j in VDW_0[i]]
#        initial_guess += [ 0.002,Pair_min_dict[i] ] # MIN GUESS
        bounds += [ (0.001,30.0),(Pair_min_dict[i]*2.**(-1./6.),10.0) ]
    fit_vals[0] = E_LJ_Fit_custom(x_vals,VDW_0)

    # Perform the fit
    t0 = time.time()
    if fit.lower() == "lj":
        results = minimize(fit_func,initial_guess,bounds=bounds,method='L-BFGS-B',options={'maxiter':1000000,'maxfun':1000000,'ftol':delta_xhi2_thresh})
        new_params = results.x

    # Print diagnostics
    if results.success == True:
        print("\t{:<30s} {:<12s}".format("Fit termination condition:","Converged"))
    else:
        print("\t{:<30s} {:<12s}".format("Fit termination condition:","Unconverged"))


    xhi_total = []
    xhi_total.append(global_fit_LJ(*initial_guess,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0))
    xhi_total.append( global_fit_LJ(*new_params,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0))

    print("\t{:<30s} {:<12.6f}".format("Initial xhi^2 (kcal/mol):",global_fit_LJ(*initial_guess,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                                                                 Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0)))
    print("\t{:<30s} {:<12.6f}".format("Final xhi^2 (kcal/mol):",global_fit_LJ(*new_params,ind=x_vals,E_config=y_vals,VDW_dict=VDW_dict,Data=Data,Fit_Pairs=Pairs,\
                                                                               Pair_min_dict=Pair_min_dict,penalty=0.0,mixing_rule=mixing_rule,L2_sigma=0.0,L2_eps=0.0)))
    print("\t{:<30s} {:<12.6f}".format("Time for completion (s):",time.time()-t0))

    # Update the VDW_dict    
    print("\nParameter update summary:\n")
    counter = 0
    for i in Pairs:      
        #print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<15.4f}".format(j) for j in new_params[counter:counter+len(VDW_dict[i])] ])))
        print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(j) for j in initial_guess[counter:counter+len(VDW_dict[i])]]),", ".join([ "{:<15.4f}".format(j) for j in new_params[counter:counter+len(VDW_dict[i])] ])))
        VDW_dict[i] = tuple([ j for j in new_params[counter:counter+len(VDW_dict[i])] ])
        VDW_dict[(i[1],i[0])] = VDW_dict[i]
        counter += len(VDW_dict[i])

    # Update cross terms if mixing rules were used
    if mixing_rule == "wh":        
        set_of_types = set([ j for i in Pairs for j in i ])
        updated_pairs = [ i for i in list(VDW_dict.keys()) if i[0] > i[1] and ( i[0] in set_of_types or i[1] in set_of_types ) ]
        for i in updated_pairs:
            #old = deepcopy(VDW_dict[i])
            old = deepcopy(VDW_0[i])
            sigma = (( VDW_dict[(i[0],i[0])][1]**(6.0) + VDW_dict[(i[1],i[1])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
            VDW_dict[i] = ( ( VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[0],i[0])][1]**(6.0) * VDW_dict[(i[1],i[1])][0]*VDW_dict[(i[1],i[1])][1]**(6.0) )**(0.5) / sigma**(6.0)  , sigma )
            VDW_dict[(i[1],i[0])] = VDW_dict[i]
            print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(k) for k in old ]),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i] ])))
    if mixing_rule == "lb":
        set_of_types = set([ j for i in Pairs for j in i ])
        updated_pairs = [ i for i in list(VDW_dict.keys()) if i[0] > i[1] and ( i[0] in set_of_types or i[1] in set_of_types ) ]
        for i in updated_pairs:
            #old = deepcopy(VDW_dict[i])
            old = deepcopy(VDW_0[i])
            VDW_dict[i] = ( ( VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[1],i[1])][0] )**(0.5) , (VDW_dict[(i[0],i[0])][1]+VDW_dict[(i[1],i[1])][1]) / 2.0)
            VDW_dict[(i[1],i[0])] = VDW_dict[i]
            print("\tpair {:80}: {} ->     {}".format(str(i),", ".join([ "{:<15.4f}".format(k) for k in old ]),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i] ])))
                    
    print(" ")

    # Calculate fit values and xhi2 with final fit parameters
    if fit.lower() == "lj":
        fit_vals[1] = E_LJ_Fit(x_vals,0.0,0.0)
    return fit_vals[:2],xhi_total


#This function detect outlier based on modified z-score
# Return: corresponding z-score and outlier(true/false) list

def detect_outlier(E_config):
    # modified z-score: Mi=0.6745*(xi-x_m)/MAD 
    # x_m: median; 
    #MAD: mean absolute deviation = median(|xi-x_m|)
    # recommanded threshold: remove Mi>3.5 
    threshold = 3.3
    
    MAD = stats.median_absolute_deviation(E_config)
    x_m = np.median(E_config)
    # Mi holds the corresponding modified z-score
    # outliter is a bunch of true/false indicating whether that point is outlier
    Mi = zeros(len(E_config))
    outlier = zeros(len(E_config))
    for count_i,i in enumerate(E_config):
        Mi[count_i] = abs(0.6745*(i-x_m)/MAD)
        #if Mi[count_i] > 3.5 : outlier[count_i] = True
        if Mi[count_i] > threshold : outlier[count_i] = True
        else: outlier[count_i] = False
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []
    plot_handle, = ax.plot(E_config,Mi,marker='.',markersize=20,color='r',markeredgewidth=0.0,alpha=0.3,linestyle='None')
    ax.set_xlabel('E_config(kcal/mol)')
    ax.set_ylabel('modified z-score')
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]
    # Save the figure
    plot_name = 'z-score'
    savefig(plot_name, dpi=300)
    close(fig)

    return Mi,outlier
    
    
# fit_type is a globally defined tuple
def global_fit_LJ_Boltz(e_0=0, s_0=0, e_1=0,  s_1=0,  e_2=0,  s_2=0,  e_3=0,  s_3=0,  e_4=0,  s_4=0,  e_5=0,  s_5=0,  e_6=0,  s_6=0,  e_7=0,  s_7=0,  e_8=0,  s_8=0,  e_9=0,  s_9=0,  e_10=0, s_10=0,\
                                e_11=0, s_11=0, e_12=0, s_12=0, e_13=0, s_13=0, e_14=0, s_14=0, e_15=0, s_15=0, e_16=0, s_16=0, e_17=0, s_17=0, e_18=0, s_18=0, e_19=0, s_19=0, e_20=0, s_20=0,\
                                e_21=0, s_21=0, e_22=0, s_22=0, e_23=0, s_23=0, e_24=0, s_24=0, e_25=0, s_25=0, e_26=0, s_26=0, e_27=0, s_27=0, e_28=0, s_28=0, e_29=0, s_29=0, e_30=0, s_30=0,\
                                e_31=0, s_31=0, e_32=0, s_32=0, e_33=0, s_33=0, e_34=0, s_34=0, e_35=0, s_35=0, e_36=0, s_36=0, e_37=0, s_37=0, e_38=0, s_38=0, e_39=0, s_39=0, e_40=0, s_40=0,\
                                e_41=0, s_41=0, e_42=0, s_42=0, e_43=0, s_43=0, e_44=0, s_44=0, e_45=0, s_45=0, e_46=0, s_46=0, e_47=0, s_47=0, e_48=0, s_48=0, e_49=0, s_49=0, e_50=0, s_50=0,\
                                e_51=0, s_51=0, e_52=0, s_52=0, e_53=0, s_53=0, e_54=0, s_54=0, e_55=0, s_55=0, e_56=0, s_56=0, e_57=0, s_57=0, e_58=0, s_58=0, e_59=0, s_59=0, e_60=0, s_60=0,\
                                e_61=0, s_61=0, e_62=0, s_62=0, e_63=0, s_63=0, e_64=0, s_64=0, e_65=0, s_65=0, e_66=0, s_66=0, e_67=0, s_67=0, e_68=0, s_68=0, e_69=0, s_69=0, e_70=0, s_70=0,\
                                e_71=0, s_71=0, e_72=0, s_72=0, e_73=0, s_73=0, e_74=0, s_74=0, e_75=0, s_75=0, e_76=0, s_76=0, e_77=0, s_77=0, e_78=0, s_78=0, e_79=0, s_79=0, e_80=0, s_80=0,\
                                e_81=0, s_81=0, e_82=0, s_82=0, e_83=0, s_83=0, e_84=0, s_84=0, e_85=0, s_85=0, e_86=0, s_86=0, e_87=0, s_87=0, e_88=0, s_88=0, e_89=0, s_89=0, e_90=0, s_90=0,\
                                e_91=0, s_91=0, e_92=0, s_92=0, e_93=0, s_93=0, e_94=0, s_94=0, e_95=0, s_95=0, e_96=0, s_96=0, e_97=0, s_97=0, e_98=0, s_98=0, e_99=0, s_99=0, e_100=0,s_100=0,\
                                e_101=0,s_101=0,e_102=0,s_102=0,e_103=0,s_103=0,e_104=0,s_104=0,e_105=0,s_105=0,e_106=0,s_106=0,e_107=0,s_107=0,e_108=0,s_108=0,e_109=0,s_109=0,e_110=0,s_110=0,\
                                e_111=0,s_111=0,e_112=0,s_112=0,e_113=0,s_113=0,e_114=0,s_114=0,e_115=0,s_115=0,e_116=0,s_116=0,e_117=0,s_117=0,e_118=0,s_118=0,e_119=0,s_119=0,e_120=0,s_120=0,\
                                e_121=0,s_121=0,e_122=0,s_122=0,e_123=0,s_123=0,e_124=0,s_124=0,e_125=0,s_125=0,e_126=0,s_126=0,e_127=0,s_127=0,e_128=0,s_128=0,e_129=0,s_129=0,e_130=0,s_130=0,\
                                e_131=0,s_131=0,e_132=0,s_132=0,e_133=0,s_133=0,e_134=0,s_134=0,e_135=0,s_135=0,e_136=0,s_136=0,e_137=0,s_137=0,e_138=0,s_138=0,e_139=0,s_139=0,e_140=0,s_140=0,\
                                e_141=0,s_141=0,e_142=0,s_142=0,e_143=0,s_143=0,e_144=0,s_144=0,e_145=0,s_145=0,e_146=0,s_146=0,e_147=0,s_147=0,e_148=0,s_148=0,e_149=0,s_149=0,e_150=0,s_150=0,\
                                e_151=0,s_151=0,e_152=0,s_152=0,e_153=0,s_153=0,e_154=0,s_154=0,e_155=0,s_155=0,e_156=0,s_156=0,e_157=0,s_157=0,e_158=0,s_158=0,e_159=0,s_159=0,e_160=0,s_160=0,\
                                e_161=0,s_161=0,e_162=0,s_162=0,e_163=0,s_163=0,e_164=0,s_164=0,e_165=0,s_165=0,e_166=0,s_166=0,e_167=0,s_167=0,e_168=0,s_168=0,e_169=0,s_169=0,e_170=0,s_170=0,\
                                e_171=0,s_171=0,e_172=0,s_172=0,e_173=0,s_173=0,e_174=0,s_174=0,e_175=0,s_175=0,e_176=0,s_176=0,e_177=0,s_177=0,e_178=0,s_178=0,e_179=0,s_179=0,e_180=0,s_180=0,\
                                e_181=0,s_181=0,e_182=0,s_182=0,e_183=0,s_183=0,e_184=0,s_184=0,e_185=0,s_185=0,e_186=0,s_186=0,e_187=0,s_187=0,e_188=0,s_188=0,e_189=0,s_189=0,e_190=0,s_190=0,\
                                e_191=0,s_191=0,e_192=0,s_192=0,e_193=0,s_193=0,e_194=0,s_194=0,e_195=0,s_195=0,e_196=0,s_196=0,e_197=0,s_197=0,e_198=0,s_198=0,e_199=0,s_199=0,e_200=0,s_200=0,\
                                ind=[],E_config=[],VDW_dict=[],Data=[],Fit_Pairs=[],Pair_min_dict=[],penalty=100.0,mixing_rule="none",L2_sigma=0.0,L2_eps=0.0,VDW_0=None,E_config_tot=[],outlier=[],weighting_function='Boltzmann'):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined e and s variables. *vars holds the variable names *vals holds the variable values
    e_vals = array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'e_' in i and "__" not in i and local_vars[i] != 0 ]) ])
    s_vals = array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 's_' in i and "__" not in i and local_vars[i] != 0 ]) ])

    # Consistency check
    if len(e_vals) != len(Fit_Pairs):
        print("ERROR in global_fit_LJ: the function expects the len(Fit_Pairs) and number of non-zero e_* variable to be equal. Exiting...")
        quit()

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])

    # Update the elements being fit with the corresponding e_vals and s_vals
    for count_c,c in enumerate(Fit_Pairs):        
        for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == c or (j[1],j[0]) == c ) ]:        
            eps_array[i]   = e_vals[count_c]
            sigma_array[i] = s_vals[count_c]

    # Apply mixing rules
    # If/else conditions handle cases where one or both of the involved types is in the fit set 
    if mixing_rule == "wh":
        fit_types = { i[0]:count_i for count_i,i in enumerate(Fit_Pairs) }
        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] in fit_types or j[1] in fit_types )]:            
            #print "\ntarget: {}".format(i)
            if i[0] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),e_vals[fit_types[i[0]]],s_vals[fit_types[i[0]]])
                if i[1] in fit_types:
                    #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                    sigma_array[count_i] = (( s_vals[fit_types[i[0]]]**(6.0) + s_vals[fit_types[i[1]]]**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]]*s_vals[fit_types[i[0]]]**(6.0) * e_vals[fit_types[i[1]]]*s_vals[fit_types[i[1]]]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
                else:
                    #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])
                    sigma_array[count_i] = (( s_vals[fit_types[i[0]]]**(6.0) + VDW_dict[(i[1],i[1])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]]*s_vals[fit_types[i[0]]]**(6.0) * VDW_dict[(i[1],i[1])][0]*VDW_dict[(i[1],i[1])][1]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
            elif i[1] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])                
                #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])

                sigma_array[count_i] = (( s_vals[fit_types[i[1]]]**(6.0) + VDW_dict[(i[0],i[0])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
                eps_array[count_i]   = (e_vals[fit_types[i[1]]]*s_vals[fit_types[i[1]]]**(6.0) * VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[0],i[0])][1]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
            #print "result: {},{}".format(eps_array[count_i],sigma_array[count_i])
#        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] not in fit_types and j[1] not in fit_types )]:            
            #print "\ntarget (2): {}".format(i)
            #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])    
            #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])            
            #print "result: {},{}".format(VDW_dict[i][0],VDW_dict[i][1])

    elif mixing_rule == "lb":
        fit_types = { i[0]:count_i for count_i,i in enumerate(Fit_Pairs) }

        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] in fit_types or j[1] in fit_types )]:            
            #print "\ntarget: {}".format(i)
            if i[0] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),e_vals[fit_types[i[0]]],s_vals[fit_types[i[0]]])
                if i[1] in fit_types:
                    #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                    sigma_array[count_i] = (s_vals[fit_types[i[0]]] + s_vals[fit_types[i[1]]] ) / 2.0 
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]] * e_vals[fit_types[i[1]]] )**(0.5)
                else:
                    #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])
                    sigma_array[count_i] = (s_vals[fit_types[i[0]]] + VDW_dict[(i[1],i[1])][1] ) / 2.0 
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]] * VDW_dict[(i[1],i[1])][0] )**(0.5)
            elif i[1] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])                
                #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                sigma_array[count_i] = (VDW_dict[(i[0],i[0])][1] + VDW_dict[(i[1],i[1])][1] ) / 2.0 
                eps_array[count_i]   = (VDW_dict[(i[0],i[0])][0] * VDW_dict[(i[1],i[1])][0] )**(0.5)
            #print "result: {},{}".format(eps_array[count_i],sigma_array[count_i])
#        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] not in fit_types and j[1] not in fit_types )]:            
            #print "\ntarget (2): {}".format(i)
            #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])    
            #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])            
            #print "result: {},{}".format(VDW_dict[i][0],VDW_dict[i][1])

    # Initialize L2 regularization sums
    if VDW_0 is not None:

        # only loop over the values in the sigma and eps arrays (they are the only ones that affect the fit!)
        L2_e = mean([ (e_vals[count_i]-VDW_0[i][0])**(2.0) for count_i,i in enumerate(Fit_Pairs) if i[0] == i[1] ]) 
        L2_s = mean([ (s_vals[count_i]-VDW_0[i][1])**(2.0) for count_i,i in enumerate(Fit_Pairs) if i[0] == i[1] ]) 

    else:
        L2_e = mean(eps_array**(2.0))
        L2_s = mean(sigma_array**(2.0))
        
    # Initialize E_LJ for the current configuration
    xhi2 = 0
    if E_config_tot != []:
        E_config_0 =  deepcopy(E_config_tot) 
    else:
        E_config_0 = deepcopy(E_config)
        if outlier != []:
            for count_d,d in enumerate(outlier):
              if d: E_config_0.pop(E_config_0.index(E_config[count_d]))
            
    E_config_min = min(E_config_0)
    E_config_max = max(E_config_0)
    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        if outlier != []:
            if outlier[count_d]: continue

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        if E_config_tot != []:
            if weighting_function == 'Boltzmann':
               xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)*exp(-(E_config_tot[count_d]-E_config_min))
            elif weighting_function == 'linear':
               xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)*(E_config_tot[count_d]-E_config_min)
            elif weighting_function == 'quadratic':
               xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)*(E_config_tot[count_d]-E_config_min)**(2.0)
            else: 
               print("ERROR: {} is not a supported weighting function".format(weighting_function))
               quit()
        else:
            if weighting_function == 'Boltzmann':
               xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)*exp(-(E_config[count_d]-E_config_min))
            elif weighting_function == 'linear':
               xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)*((E_config_max-E_config_min)-(E_config[count_d]-E_config_min))
            elif weighting_function == 'quadratic':
               xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)*((E_config_max-E_config_min)-(E_config[count_d]-E_config_min))**(2.0)
            else:
               print("ERROR: {} is not a supported weighting function".format(weighting_function))
               quit()

    return xhi2/float(len(ind)-1) + L2_sigma*L2_s + L2_eps*L2_e
# fit_type is a globally defined tuple
def global_fit_LJ(e_0=0, s_0=0, e_1=0,  s_1=0,  e_2=0,  s_2=0,  e_3=0,  s_3=0,  e_4=0,  s_4=0,  e_5=0,  s_5=0,  e_6=0,  s_6=0,  e_7=0,  s_7=0,  e_8=0,  s_8=0,  e_9=0,  s_9=0,  e_10=0, s_10=0,\
                                e_11=0, s_11=0, e_12=0, s_12=0, e_13=0, s_13=0, e_14=0, s_14=0, e_15=0, s_15=0, e_16=0, s_16=0, e_17=0, s_17=0, e_18=0, s_18=0, e_19=0, s_19=0, e_20=0, s_20=0,\
                                e_21=0, s_21=0, e_22=0, s_22=0, e_23=0, s_23=0, e_24=0, s_24=0, e_25=0, s_25=0, e_26=0, s_26=0, e_27=0, s_27=0, e_28=0, s_28=0, e_29=0, s_29=0, e_30=0, s_30=0,\
                                e_31=0, s_31=0, e_32=0, s_32=0, e_33=0, s_33=0, e_34=0, s_34=0, e_35=0, s_35=0, e_36=0, s_36=0, e_37=0, s_37=0, e_38=0, s_38=0, e_39=0, s_39=0, e_40=0, s_40=0,\
                                e_41=0, s_41=0, e_42=0, s_42=0, e_43=0, s_43=0, e_44=0, s_44=0, e_45=0, s_45=0, e_46=0, s_46=0, e_47=0, s_47=0, e_48=0, s_48=0, e_49=0, s_49=0, e_50=0, s_50=0,\
                                e_51=0, s_51=0, e_52=0, s_52=0, e_53=0, s_53=0, e_54=0, s_54=0, e_55=0, s_55=0, e_56=0, s_56=0, e_57=0, s_57=0, e_58=0, s_58=0, e_59=0, s_59=0, e_60=0, s_60=0,\
                                e_61=0, s_61=0, e_62=0, s_62=0, e_63=0, s_63=0, e_64=0, s_64=0, e_65=0, s_65=0, e_66=0, s_66=0, e_67=0, s_67=0, e_68=0, s_68=0, e_69=0, s_69=0, e_70=0, s_70=0,\
                                e_71=0, s_71=0, e_72=0, s_72=0, e_73=0, s_73=0, e_74=0, s_74=0, e_75=0, s_75=0, e_76=0, s_76=0, e_77=0, s_77=0, e_78=0, s_78=0, e_79=0, s_79=0, e_80=0, s_80=0,\
                                e_81=0, s_81=0, e_82=0, s_82=0, e_83=0, s_83=0, e_84=0, s_84=0, e_85=0, s_85=0, e_86=0, s_86=0, e_87=0, s_87=0, e_88=0, s_88=0, e_89=0, s_89=0, e_90=0, s_90=0,\
                                e_91=0, s_91=0, e_92=0, s_92=0, e_93=0, s_93=0, e_94=0, s_94=0, e_95=0, s_95=0, e_96=0, s_96=0, e_97=0, s_97=0, e_98=0, s_98=0, e_99=0, s_99=0, e_100=0,s_100=0,\
                                e_101=0,s_101=0,e_102=0,s_102=0,e_103=0,s_103=0,e_104=0,s_104=0,e_105=0,s_105=0,e_106=0,s_106=0,e_107=0,s_107=0,e_108=0,s_108=0,e_109=0,s_109=0,e_110=0,s_110=0,\
                                e_111=0,s_111=0,e_112=0,s_112=0,e_113=0,s_113=0,e_114=0,s_114=0,e_115=0,s_115=0,e_116=0,s_116=0,e_117=0,s_117=0,e_118=0,s_118=0,e_119=0,s_119=0,e_120=0,s_120=0,\
                                e_121=0,s_121=0,e_122=0,s_122=0,e_123=0,s_123=0,e_124=0,s_124=0,e_125=0,s_125=0,e_126=0,s_126=0,e_127=0,s_127=0,e_128=0,s_128=0,e_129=0,s_129=0,e_130=0,s_130=0,\
                                e_131=0,s_131=0,e_132=0,s_132=0,e_133=0,s_133=0,e_134=0,s_134=0,e_135=0,s_135=0,e_136=0,s_136=0,e_137=0,s_137=0,e_138=0,s_138=0,e_139=0,s_139=0,e_140=0,s_140=0,\
                                e_141=0,s_141=0,e_142=0,s_142=0,e_143=0,s_143=0,e_144=0,s_144=0,e_145=0,s_145=0,e_146=0,s_146=0,e_147=0,s_147=0,e_148=0,s_148=0,e_149=0,s_149=0,e_150=0,s_150=0,\
                                e_151=0,s_151=0,e_152=0,s_152=0,e_153=0,s_153=0,e_154=0,s_154=0,e_155=0,s_155=0,e_156=0,s_156=0,e_157=0,s_157=0,e_158=0,s_158=0,e_159=0,s_159=0,e_160=0,s_160=0,\
                                e_161=0,s_161=0,e_162=0,s_162=0,e_163=0,s_163=0,e_164=0,s_164=0,e_165=0,s_165=0,e_166=0,s_166=0,e_167=0,s_167=0,e_168=0,s_168=0,e_169=0,s_169=0,e_170=0,s_170=0,\
                                e_171=0,s_171=0,e_172=0,s_172=0,e_173=0,s_173=0,e_174=0,s_174=0,e_175=0,s_175=0,e_176=0,s_176=0,e_177=0,s_177=0,e_178=0,s_178=0,e_179=0,s_179=0,e_180=0,s_180=0,\
                                e_181=0,s_181=0,e_182=0,s_182=0,e_183=0,s_183=0,e_184=0,s_184=0,e_185=0,s_185=0,e_186=0,s_186=0,e_187=0,s_187=0,e_188=0,s_188=0,e_189=0,s_189=0,e_190=0,s_190=0,\
                                e_191=0,s_191=0,e_192=0,s_192=0,e_193=0,s_193=0,e_194=0,s_194=0,e_195=0,s_195=0,e_196=0,s_196=0,e_197=0,s_197=0,e_198=0,s_198=0,e_199=0,s_199=0,e_200=0,s_200=0,\
                                ind=[],E_config=[],VDW_dict=[],Data=[],Fit_Pairs=[],Pair_min_dict=[],penalty=100.0,mixing_rule="none",L2_sigma=0.0,L2_eps=0.0,VDW_0=None):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined e and s variables. *vars holds the variable names *vals holds the variable values
    e_vals = array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'e_' in i and "__" not in i and local_vars[i] != 0 ]) ])
    s_vals = array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 's_' in i and "__" not in i and local_vars[i] != 0 ]) ])

    # Consistency check
    if len(e_vals) != len(Fit_Pairs):
        print("ERROR in global_fit_LJ: the function expects the len(Fit_Pairs) and number of non-zero e_* variable to be equal. Exiting...")
        quit()

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])

    # Update the elements being fit with the corresponding e_vals and s_vals
    for count_c,c in enumerate(Fit_Pairs):        
        for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == c or (j[1],j[0]) == c ) ]:        
            eps_array[i]   = e_vals[count_c]
            sigma_array[i] = s_vals[count_c]

    # Apply mixing rules
    # If/else conditions handle cases where one or both of the involved types is in the fit set 
    if mixing_rule == "wh":
        fit_types = { i[0]:count_i for count_i,i in enumerate(Fit_Pairs) }
        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] in fit_types or j[1] in fit_types )]:            
            #print "\ntarget: {}".format(i)
            if i[0] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),e_vals[fit_types[i[0]]],s_vals[fit_types[i[0]]])
                if i[1] in fit_types:
                    #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                    sigma_array[count_i] = (( s_vals[fit_types[i[0]]]**(6.0) + s_vals[fit_types[i[1]]]**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]]*s_vals[fit_types[i[0]]]**(6.0) * e_vals[fit_types[i[1]]]*s_vals[fit_types[i[1]]]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
                else:
                    #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])
                    sigma_array[count_i] = (( s_vals[fit_types[i[0]]]**(6.0) + VDW_dict[(i[1],i[1])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]]*s_vals[fit_types[i[0]]]**(6.0) * VDW_dict[(i[1],i[1])][0]*VDW_dict[(i[1],i[1])][1]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
            elif i[1] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])                
                #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])

                sigma_array[count_i] = (( s_vals[fit_types[i[1]]]**(6.0) + VDW_dict[(i[0],i[0])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
                eps_array[count_i]   = (e_vals[fit_types[i[1]]]*s_vals[fit_types[i[1]]]**(6.0) * VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[0],i[0])][1]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
            #print "result: {},{}".format(eps_array[count_i],sigma_array[count_i])
#        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] not in fit_types and j[1] not in fit_types )]:            
            #print "\ntarget (2): {}".format(i)
            #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])    
            #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])            
            #print "result: {},{}".format(VDW_dict[i][0],VDW_dict[i][1])

    elif mixing_rule == "lb":
        fit_types = { i[0]:count_i for count_i,i in enumerate(Fit_Pairs) }

        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] in fit_types or j[1] in fit_types )]:            
            #print "\ntarget: {}".format(i)
            if i[0] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),e_vals[fit_types[i[0]]],s_vals[fit_types[i[0]]])
                if i[1] in fit_types:
                    #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                    sigma_array[count_i] = (s_vals[fit_types[i[0]]] + s_vals[fit_types[i[1]]] ) / 2.0 
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]] * e_vals[fit_types[i[1]]] )**(0.5)
                else:
                    #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])
                    sigma_array[count_i] = (s_vals[fit_types[i[0]]] + VDW_dict[(i[1],i[1])][1] ) / 2.0 
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]] * VDW_dict[(i[1],i[1])][0] )**(0.5)
            elif i[1] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])                
                #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                sigma_array[count_i] = (VDW_dict[(i[0],i[0])][1] + VDW_dict[(i[1],i[1])][1] ) / 2.0 
                eps_array[count_i]   = (VDW_dict[(i[0],i[0])][0] * VDW_dict[(i[1],i[1])][0] )**(0.5)
            #print "result: {},{}".format(eps_array[count_i],sigma_array[count_i])
#        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] not in fit_types and j[1] not in fit_types )]:            
            #print "\ntarget (2): {}".format(i)
            #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])    
            #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])            
            #print "result: {},{}".format(VDW_dict[i][0],VDW_dict[i][1])

    # Initialize L2 regularization sums
    if VDW_0 is not None:

        # only loop over the values in the sigma and eps arrays (they are the only ones that affect the fit!)
        L2_e = mean([ (e_vals[count_i]-VDW_0[i][0])**(2.0) for count_i,i in enumerate(Fit_Pairs) if i[0] == i[1] ]) 
        L2_s = mean([ (s_vals[count_i]-VDW_0[i][1])**(2.0) for count_i,i in enumerate(Fit_Pairs) if i[0] == i[1] ]) 

    else:
        L2_e = mean(eps_array**(2.0))
        L2_s = mean(sigma_array**(2.0))
        
    # Initialize E_LJ for the current configuration
    xhi2 = 0
    min_index = E_config.index(min(E_config))
    # Iterate over the configuration index
    for count_d,d in enumerate(ind):
#        if count_d == min_index: continue
        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)

    return xhi2/float(len(ind)) + L2_sigma*L2_s + L2_eps*L2_e

def global_fit_LJ_outlier(e_0=0, s_0=0, e_1=0,  s_1=0,  e_2=0,  s_2=0,  e_3=0,  s_3=0,  e_4=0,  s_4=0,  e_5=0,  s_5=0,  e_6=0,  s_6=0,  e_7=0,  s_7=0,  e_8=0,  s_8=0,  e_9=0,  s_9=0,  e_10=0, s_10=0,\
                                e_11=0, s_11=0, e_12=0, s_12=0, e_13=0, s_13=0, e_14=0, s_14=0, e_15=0, s_15=0, e_16=0, s_16=0, e_17=0, s_17=0, e_18=0, s_18=0, e_19=0, s_19=0, e_20=0, s_20=0,\
                                e_21=0, s_21=0, e_22=0, s_22=0, e_23=0, s_23=0, e_24=0, s_24=0, e_25=0, s_25=0, e_26=0, s_26=0, e_27=0, s_27=0, e_28=0, s_28=0, e_29=0, s_29=0, e_30=0, s_30=0,\
                                e_31=0, s_31=0, e_32=0, s_32=0, e_33=0, s_33=0, e_34=0, s_34=0, e_35=0, s_35=0, e_36=0, s_36=0, e_37=0, s_37=0, e_38=0, s_38=0, e_39=0, s_39=0, e_40=0, s_40=0,\
                                e_41=0, s_41=0, e_42=0, s_42=0, e_43=0, s_43=0, e_44=0, s_44=0, e_45=0, s_45=0, e_46=0, s_46=0, e_47=0, s_47=0, e_48=0, s_48=0, e_49=0, s_49=0, e_50=0, s_50=0,\
                                e_51=0, s_51=0, e_52=0, s_52=0, e_53=0, s_53=0, e_54=0, s_54=0, e_55=0, s_55=0, e_56=0, s_56=0, e_57=0, s_57=0, e_58=0, s_58=0, e_59=0, s_59=0, e_60=0, s_60=0,\
                                e_61=0, s_61=0, e_62=0, s_62=0, e_63=0, s_63=0, e_64=0, s_64=0, e_65=0, s_65=0, e_66=0, s_66=0, e_67=0, s_67=0, e_68=0, s_68=0, e_69=0, s_69=0, e_70=0, s_70=0,\
                                e_71=0, s_71=0, e_72=0, s_72=0, e_73=0, s_73=0, e_74=0, s_74=0, e_75=0, s_75=0, e_76=0, s_76=0, e_77=0, s_77=0, e_78=0, s_78=0, e_79=0, s_79=0, e_80=0, s_80=0,\
                                e_81=0, s_81=0, e_82=0, s_82=0, e_83=0, s_83=0, e_84=0, s_84=0, e_85=0, s_85=0, e_86=0, s_86=0, e_87=0, s_87=0, e_88=0, s_88=0, e_89=0, s_89=0, e_90=0, s_90=0,\
                                e_91=0, s_91=0, e_92=0, s_92=0, e_93=0, s_93=0, e_94=0, s_94=0, e_95=0, s_95=0, e_96=0, s_96=0, e_97=0, s_97=0, e_98=0, s_98=0, e_99=0, s_99=0, e_100=0,s_100=0,\
                                e_101=0,s_101=0,e_102=0,s_102=0,e_103=0,s_103=0,e_104=0,s_104=0,e_105=0,s_105=0,e_106=0,s_106=0,e_107=0,s_107=0,e_108=0,s_108=0,e_109=0,s_109=0,e_110=0,s_110=0,\
                                e_111=0,s_111=0,e_112=0,s_112=0,e_113=0,s_113=0,e_114=0,s_114=0,e_115=0,s_115=0,e_116=0,s_116=0,e_117=0,s_117=0,e_118=0,s_118=0,e_119=0,s_119=0,e_120=0,s_120=0,\
                                e_121=0,s_121=0,e_122=0,s_122=0,e_123=0,s_123=0,e_124=0,s_124=0,e_125=0,s_125=0,e_126=0,s_126=0,e_127=0,s_127=0,e_128=0,s_128=0,e_129=0,s_129=0,e_130=0,s_130=0,\
                                e_131=0,s_131=0,e_132=0,s_132=0,e_133=0,s_133=0,e_134=0,s_134=0,e_135=0,s_135=0,e_136=0,s_136=0,e_137=0,s_137=0,e_138=0,s_138=0,e_139=0,s_139=0,e_140=0,s_140=0,\
                                e_141=0,s_141=0,e_142=0,s_142=0,e_143=0,s_143=0,e_144=0,s_144=0,e_145=0,s_145=0,e_146=0,s_146=0,e_147=0,s_147=0,e_148=0,s_148=0,e_149=0,s_149=0,e_150=0,s_150=0,\
                                e_151=0,s_151=0,e_152=0,s_152=0,e_153=0,s_153=0,e_154=0,s_154=0,e_155=0,s_155=0,e_156=0,s_156=0,e_157=0,s_157=0,e_158=0,s_158=0,e_159=0,s_159=0,e_160=0,s_160=0,\
                                e_161=0,s_161=0,e_162=0,s_162=0,e_163=0,s_163=0,e_164=0,s_164=0,e_165=0,s_165=0,e_166=0,s_166=0,e_167=0,s_167=0,e_168=0,s_168=0,e_169=0,s_169=0,e_170=0,s_170=0,\
                                e_171=0,s_171=0,e_172=0,s_172=0,e_173=0,s_173=0,e_174=0,s_174=0,e_175=0,s_175=0,e_176=0,s_176=0,e_177=0,s_177=0,e_178=0,s_178=0,e_179=0,s_179=0,e_180=0,s_180=0,\
                                e_181=0,s_181=0,e_182=0,s_182=0,e_183=0,s_183=0,e_184=0,s_184=0,e_185=0,s_185=0,e_186=0,s_186=0,e_187=0,s_187=0,e_188=0,s_188=0,e_189=0,s_189=0,e_190=0,s_190=0,\
                                e_191=0,s_191=0,e_192=0,s_192=0,e_193=0,s_193=0,e_194=0,s_194=0,e_195=0,s_195=0,e_196=0,s_196=0,e_197=0,s_197=0,e_198=0,s_198=0,e_199=0,s_199=0,e_200=0,s_200=0,\
                                ind=[],outlier=[],E_config=[],VDW_dict=[],Data=[],Fit_Pairs=[],Pair_min_dict=[],penalty=100.0,mixing_rule="none",L2_sigma=0.0,L2_eps=0.0,VDW_0=None):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()

    # Find defined e and s variables. *vars holds the variable names *vals holds the variable values
    e_vals = array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 'e_' in i and "__" not in i and local_vars[i] != 0 ]) ])
    s_vals = array([ local_vars[j] for j in natural_sort([ i for i in local_vars if 's_' in i and "__" not in i and local_vars[i] != 0 ]) ])

    # Consistency check
    if len(e_vals) != len(Fit_Pairs):
        print("ERROR in global_fit_LJ: the function expects the len(Fit_Pairs) and number of non-zero e_* variable to be equal. Exiting...")
        quit()

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])

    # Update the elements being fit with the corresponding e_vals and s_vals
    for count_c,c in enumerate(Fit_Pairs):        
        for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == c or (j[1],j[0]) == c ) ]:        
            eps_array[i]   = e_vals[count_c]
            sigma_array[i] = s_vals[count_c]

    # Apply mixing rules
    # If/else conditions handle cases where one or both of the involved types is in the fit set 
    if mixing_rule == "wh":
        fit_types = { i[0]:count_i for count_i,i in enumerate(Fit_Pairs) }
        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] in fit_types or j[1] in fit_types )]:            
            #print "\ntarget: {}".format(i)
            if i[0] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),e_vals[fit_types[i[0]]],s_vals[fit_types[i[0]]])
                if i[1] in fit_types:
                    #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                    sigma_array[count_i] = (( s_vals[fit_types[i[0]]]**(6.0) + s_vals[fit_types[i[1]]]**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]]*s_vals[fit_types[i[0]]]**(6.0) * e_vals[fit_types[i[1]]]*s_vals[fit_types[i[1]]]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
                else:
                    #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])
                    sigma_array[count_i] = (( s_vals[fit_types[i[0]]]**(6.0) + VDW_dict[(i[1],i[1])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]]*s_vals[fit_types[i[0]]]**(6.0) * VDW_dict[(i[1],i[1])][0]*VDW_dict[(i[1],i[1])][1]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
            elif i[1] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])                
                #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])

                sigma_array[count_i] = (( s_vals[fit_types[i[1]]]**(6.0) + VDW_dict[(i[0],i[0])][1]**(6.0) ) / 2.0 )**(1.0/6.0)
                eps_array[count_i]   = (e_vals[fit_types[i[1]]]*s_vals[fit_types[i[1]]]**(6.0) * VDW_dict[(i[0],i[0])][0]*VDW_dict[(i[0],i[0])][1]**(6.0) )**(0.5) / sigma_array[count_i]**(6.0)
            #print "result: {},{}".format(eps_array[count_i],sigma_array[count_i])
#        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] not in fit_types and j[1] not in fit_types )]:            
            #print "\ntarget (2): {}".format(i)
            #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])    
            #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])            
            #print "result: {},{}".format(VDW_dict[i][0],VDW_dict[i][1])

    elif mixing_rule == "lb":
        fit_types = { i[0]:count_i for count_i,i in enumerate(Fit_Pairs) }

        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] in fit_types or j[1] in fit_types )]:            
            #print "\ntarget: {}".format(i)
            if i[0] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),e_vals[fit_types[i[0]]],s_vals[fit_types[i[0]]])
                if i[1] in fit_types:
                    #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                    sigma_array[count_i] = (s_vals[fit_types[i[0]]] + s_vals[fit_types[i[1]]] ) / 2.0 
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]] * e_vals[fit_types[i[1]]] )**(0.5)
                else:
                    #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])
                    sigma_array[count_i] = (s_vals[fit_types[i[0]]] + VDW_dict[(i[1],i[1])][1] ) / 2.0 
                    eps_array[count_i]   = (e_vals[fit_types[i[0]]] * VDW_dict[(i[1],i[1])][0] )**(0.5)
            elif i[1] in fit_types:
                #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])                
                #print "{}: {},{}".format((i[1],i[1]),e_vals[fit_types[i[1]]],s_vals[fit_types[i[1]]])
                sigma_array[count_i] = (VDW_dict[(i[0],i[0])][1] + VDW_dict[(i[1],i[1])][1] ) / 2.0 
                eps_array[count_i]   = (VDW_dict[(i[0],i[0])][0] * VDW_dict[(i[1],i[1])][0] )**(0.5)
            #print "result: {},{}".format(eps_array[count_i],sigma_array[count_i])
#        for count_i,i in [ (count_j,j) for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j[0] not in fit_types and j[1] not in fit_types )]:            
            #print "\ntarget (2): {}".format(i)
            #print "{}: {},{}".format((i[0],i[0]),VDW_dict[(i[0],i[0])][0],VDW_dict[(i[0],i[0])][1])    
            #print "{}: {},{}".format((i[1],i[1]),VDW_dict[(i[1],i[1])][0],VDW_dict[(i[1],i[1])][1])            
            #print "result: {},{}".format(VDW_dict[i][0],VDW_dict[i][1])

    # Initialize L2 regularization sums
    if VDW_0 is not None:

        # only loop over the values in the sigma and eps arrays (they are the only ones that affect the fit!)
        L2_e = mean([ (e_vals[count_i]-VDW_0[i][0])**(2.0) for count_i,i in enumerate(Fit_Pairs) if i[0] == i[1] ]) 
        L2_s = mean([ (s_vals[count_i]-VDW_0[i][1])**(2.0) for count_i,i in enumerate(Fit_Pairs) if i[0] == i[1] ]) 

    else:
        L2_e = mean(eps_array**(2.0))
        L2_s = mean(sigma_array**(2.0))
        
    # Initialize E_LJ for the current configuration
    xhi2 = 0
    # Iterate over the configuration index
    for count_d,d in enumerate(ind):
        if outlier[count_d]: continue
        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        xhi2 += ( sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) ) - E_config[count_d] )**(2.0)

    return xhi2/float(len(ind)) + L2_sigma*L2_s + L2_eps*L2_e

# Description: This function drives the pair-wise fit of the VDW parameters. It accepts
#              various arguments that modify the convergence criteria and optimization
#              algorithm. 
# 
# Inputs:      x_vals            : a list holding the keys in Data for configurations being fit
#              y_vals            : a list holding the interaction energies for each configuration being fit 
#              Pairs             : a list of the unique pair-types being fit.
#              weight            : a scalar that controls the weight of the update at each step.
#              delta_xhi2_thresh : a scalar that controls the xhi^2 convergence threshold for the fit. 
#              min_cycles        : an integer that controls the minimum number of fit cycles to perform.
#              max_cycles        : an integer that controls the maximum number of fit cycles to perform. 
# 
# Returns:     fit_vals          : a list of arrays, each holding the FF based interaction energies in each cycle of the fit.
#
def iterative_fit(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh=0.00001,min_cycles=100,max_cycles=1000):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # Intialize cycle count and array of fit values
    fit_type = 0 # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
    cycle_count = 0

    if fit == 'lj':
        fit_vals = [ E_LJ_Fit(x_vals,0.0,0.0) ] * (max_cycles+5)
    elif fit == 'buck':
        fit_vals = [ E_BUCK_Fit(x_vals,0.0,0.0,0.0) ] * (max_cycles+5)
    xhi2_previous = mean((y_vals-fit_vals[cycle_count])**2) 

    # If no pairs are being fit (as happens if only UA types are being fit and AA are being read from file) just return the fit_vals
    if len(Pairs) == 0: print(" "); return fit_vals[:cycle_count+1]

    # Check the rank of the fit data by assembling the lstsq fit matrix
    if fit in ['lj','buck'] :
        A = zeros([len(x_vals),len(Pairs)*2]) # For the lj fit there are 2 linearly parameters per pair type (A,B). The rank of the LJ matrix is used in place of the buck 
                                              # A matrix owing to the non-linearity of the buckingham potential
    # Build up A matrix
    for count_i,i in enumerate(x_vals):
        for count_j,j in enumerate(Pairs):

            if fit in ['lj','buck']:
                # Pairs only contains I J or J I combinations whereas Data[i]["pairs"] has both combinations
                # To be consistent the contributions from both I J and J I configurations need to be added to 
                # the elements of A
                tmp_6 = 0.0
                tmp_12 = 0.0

                if (j[0],j[1]) in list(Data[i]["pairs"].keys()):
                    tmp_6  += sum(Data[i]["pairs"][(j[0],j[1])]**(-6.0))
                    tmp_12 += sum(Data[i]["pairs"][(j[0],j[1])]**(-12.0))
                if j[0] != j[1] and (j[1],j[0]) in list(Data[i]["pairs"].keys()):
                    tmp_6  += sum(Data[i]["pairs"][(j[1],j[0])]**(-6.0))
                    tmp_12 += sum(Data[i]["pairs"][(j[1],j[0])]**(-12.0))

                A[count_i,(count_j*2)+0] += tmp_12
                A[count_i,(count_j*2)+1] -= tmp_6

    # Calculate the rank of the fit matrix
    rank = matrix_rank(A)

    # Print diagnostics
    print("\nInitializing iterative fit of pair-wise interactions:\n")
    print("\t{:<30s} {}".format("max cycles:",max_cycles))
    print("\t{:<30s} {}".format("min_cycles:",min_cycles))
    print("\t{:<30s} {}".format("update_weight:",weight))
    print("\t{:<30s} {}".format("delta_xhi2_thresh:",delta_xhi2_thresh))
    print("\t{:<30s} {}".format("Fit rank:",rank))
    print("\t{:<30s} {}".format("Fit parameters:",len(A[0])))

    print("\ncycle 0:\n\n\txhi2: {} (kcal/mol)^2\n".format(xhi2_previous))

    # Fit loop
    while (1):
        
        # Increment the cycle count and print diagnostic
        cycle_count += 1
        t0 = time.time()
        print("cycle {}:\n".format(cycle_count))
        
        # Each cycle consists of looping over all pair-types being fit
        # To eliminate ordering bias, the list of types is shuffled every cycle.
        random.shuffle(Pairs)
        
        # All non-hydrogen terms are fit first, starting with self-terms then cross-terms, followed by hydrogen self-terms then hydrogen cross-terms
        Non_H_self  = [ i for i in Pairs if i[0] == i[1] and int(i[0].split('[')[1].split(']')[0]) != 1 and int(i[1].split('[')[1].split(']')[0]) != 1 ]
        Non_H_cross = [ i for i in Pairs if i[0] != i[1] and int(i[0].split('[')[1].split(']')[0]) != 1 and int(i[1].split('[')[1].split(']')[0]) != 1 ]
        H_self      = [ i for i in Pairs if i[0] == i[1] and int(i[0].split('[')[1].split(']')[0]) == 1 and int(i[1].split('[')[1].split(']')[0]) == 1 ]
        H_cross     = [ i for i in Pairs if i[0] != i[1] and (int(i[0].split('[')[1].split(']')[0]) == 1 or int(i[1].split('[')[1].split(']')[0]) == 1 ) ]
        Pairs = Non_H_self + Non_H_cross + H_self + H_cross

        for i in Pairs:

            # Set the pair-type that is being fit
            fit_type = i                            
            old_params = list(VDW_dict[i])

            # Fit the potential (eps is fit first then sigma; this was pursued because fitting both at the same time
            # led to instability when using weighted updates.)
            if fit.lower() == "lj":
                new_eps,err           = curve_fit(E_LJ_Fit_eps,x_vals,y_vals,p0=VDW_dict[i][0],maxfev=10000)
                new_eps               = VDW_dict[i][0]*(1.0-weight) + new_eps[0]*weight
                VDW_dict[i]           = tuple([new_eps,VDW_dict[i][1]])
                VDW_dict[(i[1],i[0])] = tuple([new_eps,VDW_dict[i][1]])
                new_sigma,err         = curve_fit(E_LJ_Fit_sigma,x_vals,y_vals,p0=VDW_dict[i][1],maxfev=10000)
                new_sigma             = VDW_dict[i][1]*(1.0-weight) + new_sigma[0]*weight
                VDW_dict[i]           = tuple([VDW_dict[i][0],new_sigma])
                VDW_dict[(i[1],i[0])] = tuple([VDW_dict[i][0],new_sigma])

            # Fit the potential (A is fit first, then B, then C; this was pursued because fitting all of them at the same time
            # led to instability when using weighted updates.)
            if fit.lower() == "buck":
#                 new_A,err             = curve_fit(E_BUCK_Fit_A,x_vals,y_vals,p0=VDW_dict[i][0],maxfev=10000)
#                 new_A                 = VDW_dict[i][0]*(1.0-weight) + new_A[0]*weight
#                 VDW_dict[i]           = tuple([new_A,VDW_dict[i][1],VDW_dict[i][2]])
#                 VDW_dict[(i[1],i[0])] = tuple([new_A,VDW_dict[i][1],VDW_dict[i][2]])
#                 new_B,err             = curve_fit(E_BUCK_Fit_B,x_vals,y_vals,p0=VDW_dict[i][1],maxfev=10000)
#                 new_B                 = VDW_dict[i][1]*(1.0-weight) + new_B[0]*weight
#                 VDW_dict[i]           = tuple([VDW_dict[i][0],new_B,VDW_dict[i][2]])
#                 VDW_dict[(i[1],i[0])] = tuple([VDW_dict[i][0],new_B,VDW_dict[i][2]])
#                 new_C,err             = curve_fit(E_BUCK_Fit_C,x_vals,y_vals,p0=VDW_dict[i][2],maxfev=10000)
#                 new_C                 = VDW_dict[i][2]*(1.0-weight) + new_C[0]*weight
#                 VDW_dict[i]           = tuple([VDW_dict[i][0],VDW_dict[i][1],new_C])
#                 VDW_dict[(i[1],i[0])] = tuple([VDW_dict[i][0],VDW_dict[i][1],new_C])

                new_ABC,err             = curve_fit(E_BUCK_Fit_ABC,x_vals,y_vals,p0=(VDW_dict[i][0],VDW_dict[i][1],VDW_dict[i][2]),maxfev=10000)
                VDW_dict[i]           = tuple([new_ABC[0],new_ABC[1],new_ABC[2]])
                VDW_dict[(i[1],i[0])] = tuple([new_ABC[0],new_ABC[1],new_ABC[2]])

            print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<15.4f}".format(j) for j in old_params]),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]])))


        # Update the fit data array
        fit_type = 0   # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
        if fit.lower() == 'lj':
            fit_vals[cycle_count] = E_LJ_Fit(x_vals,0.0,0.0)
        elif fit.lower() == 'buck':
            fit_vals[cycle_count] = E_BUCK_Fit(x_vals,0.0,0.0,0.0)

        # Calculate mean squared deviation
        xhi2 = mean((y_vals-fit_vals[cycle_count])**2)
        delta_xhi2 = xhi2_previous - xhi2
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",xhi2))
        print("\t{:<12s} {} (kcal/mol)^2".format("delta_xhi2:",delta_xhi2))
        print("\t{:<12s} {} (seconds)".format("time:",time.time()-t0))
        print("")
        xhi2_previous = copy(xhi2)

        # Check break conditions
        if cycle_count >= min_cycles:
            if cycle_count == max_cycles:
                print("Maximum number of cycles reached, breaking iterative fit loop...\n")
                break
            elif delta_xhi2 < delta_xhi2_thresh:
                print("Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2), breaking iterative fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
                break

    # Avoid the global fit if the number of pairs being fit exceeds the number of configurations (the global fit function complains)
    # Or if there are more than 100 pairs being fit (the maximum number of parameters the function can handle is 100 eps and 100 sigma values respectively)
    if len(Pairs)*2 > len(x_vals) or len(Pairs) > 100:        
        return fit_vals[:cycle_count+1]

    # Perform a final optimization of all parameters at once
    global Current_Fit_Pairs

    # Copy the supplied Fit_Pairs to the local/global variable Current_Fit_Pairs
    # This is performed because the pairs being fit need to be flexibly defined locally
    Current_Fit_Pairs = Pairs

    # Prepare initial guess parameters, initial fit energies, and initial xhi2 value
    initial_guess = [ j for i in Current_Fit_Pairs for j in VDW_dict[i] ]
    if fit == 'lj':
        xhi2_previous = mean((y_vals-E_LJ_Fit_all(x_vals,*initial_guess))**2)
    elif fit == 'buck':
        xhi2_previous = mean((y_vals-E_BUCK_Fit_all(x_vals,*initial_guess))**2)

    # Perform the global fit (all eps and sigma values are fit at once)
    while 1:

        # Increment the cycle count and print diagnostic
        cycle_count += 1
        t0 = time.time()
        print("cycle {} (global fit):\n".format(cycle_count))

        # Perform LJ fit
        if fit == 'lj':
            params,err = curve_fit(E_LJ_Fit_all,x_vals,y_vals,p0=initial_guess,maxfev=10000)

            # Calculate fit values and xhi2 with final fit parameters
            fit_vals[cycle_count] = E_LJ_Fit_all(x_vals,*params)
            xhi2 = mean((y_vals-fit_vals[cycle_count])**2)        
            delta_xhi2 = xhi2_previous - xhi2 

            # Convert fit parameters back into eps,sigma tuples
            initial_guess = params
            params = [ (params[2*count_i],params[2*count_i+1]) for count_i,i in enumerate(Current_Fit_Pairs) ]

            # Print diagnostic and assign fit parameters to the VDW_dict
            for count_i,i in enumerate(Current_Fit_Pairs):
                print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<15.4f}".format(j) for j in params[count_i]])))
                VDW_dict[i] = tuple(params[count_i])  
                VDW_dict[(i[1],i[0])] = tuple(params[count_i])        

        # Perform BUCK fit
        if fit == 'buck':
            params,err = curve_fit(E_BUCK_Fit_all,x_vals,y_vals,p0=initial_guess,maxfev=10000)

            # Calculate fit values and xhi2 with final fit parameters
            fit_vals[cycle_count] = E_BUCK_Fit_all(x_vals,*params)
            xhi2 = mean((y_vals-fit_vals[cycle_count])**2)        
            delta_xhi2 = xhi2_previous - xhi2 

            # Convert fit parameters back into A,B,C tuples
            initial_guess = params
            params = [ (params[3*count_i],params[3*count_i+1],params[3*count_i+2]) for count_i,i in enumerate(Current_Fit_Pairs) ]

            # Print diagnostic and assign fit parameters to the VDW_dict
            for count_i,i in enumerate(Current_Fit_Pairs):
                print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<15.4f}".format(j) for j in params[count_i]])))
                VDW_dict[i] = tuple(params[count_i])  
                VDW_dict[(i[1],i[0])] = tuple(params[count_i])        

        # Print change in xhi2 and update xhi2_previous
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",xhi2))
        print("\t{:<12s} {} (kcal/mol)^2".format("delta_xhi2:",delta_xhi2))
        print("\t{:<12s} {} (seconds)".format("time:",time.time()-t0))
        print("")

        xhi2_previous = copy(xhi2)
        
        # Check break conditions
        if delta_xhi2 < delta_xhi2_thresh:
            print("Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2), breaking global fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
            break    

    return fit_vals[:cycle_count+1]

# Description: This function drives the pair-wise fitting of the UA-VDW parameters. It accepts
#              various arguments that modify the convergence criteria and optimization
#              algorithm. 
# 
# Inputs:      x_vals            : a list holding the keys in Data for configurations being fit
#              y_vals            : a list holding the interaction energies for each configuration being fit 
#              Pairs             : a list of the unique pair-types being fit.
#              weight            : a scalar that controls the weight of the update at each step.
#              delta_xhi2_thresh : a scalar that controls the xhi^2 convergence threshold for the fit. 
#              min_cycles        : an integer that controls the minimum number of fit cycles to perform.
#              max_cycles        : an integer that controls the maximum number of fit cycles to perform. 
# 
# Returns:     fit_vals          : a list of arrays, each holding the FF based interaction energies in each cycle of the fit.
#
def iterative_fit_UA(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh=0.00001,min_cycles=100,max_cycles=1000):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,fit_type,Pair_min_dict

    # Intialize cycle count and array of fit values
    cycle_count = 0
    fit_type = 0 # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
    if fit == 'lj':
        fit_vals = [ E_LJ_Fit(x_vals,0.0,0.0) ] * (max_cycles+5)
    elif fit == 'buck':
        fit_vals = [ E_BUCK_Fit(x_vals,0.0,0.0,0.0) ] * (max_cycles+5)
    xhi2_previous = mean((y_vals-fit_vals[cycle_count])**2) 

    # If no pairs are being fit (as happens if only UA types are being fit and AA are being read from file) just return the fit_vals
    if len(Pairs) == 0: print(" "); return fit_vals[:cycle_count+1]

    # Check the rank of the fit data
    if fit in ['lj','buck']:
        A = zeros([len(x_vals),len(Pairs)*2]) # # For the lj fit there are 2 linearly parameters per pair type (A,B). The rank of the LJ matrix is used in place of the buck matrix.

    for count_i,i in enumerate(x_vals):
        for count_j,j in enumerate(Pairs):

            # Pairs only contains I J or J I combinations whereas Data[i]["pairs"] has both combinations
            # To be consistent the contributions from both I J and J I configurations need to be added to 
            # the elements of A
            tmp_6 = 0.0
            tmp_12 = 0.0

            if (j[0],j[1]) in list(Data[i]["pairs"].keys()):
                tmp_6  += sum(Data[i]["pairs"][(j[0],j[1])]**(-6.0))
                tmp_12 += sum(Data[i]["pairs"][(j[0],j[1])]**(-12.0))
            if j[0] != j[1] and (j[1],j[0]) in list(Data[i]["pairs"].keys()):
                tmp_6  += sum(Data[i]["pairs"][(j[1],j[0])]**(-6.0))
                tmp_12 += sum(Data[i]["pairs"][(j[1],j[0])]**(-12.0))

            A[count_i,(count_j*2)+0] += tmp_12
            A[count_i,(count_j*2)+1] -= tmp_6

    # Define fit_weights (unweighted is ones)
    #    fit_weights = exp(array(y_vals)/(0.0019872041*400))
    #    fit_weights = exp(array([ Data[i]['e_dft'] for i in x_vals])/(0.0019872041*400.0))
    fit_weights = ones(len(x_vals))
    #    print "\nfit_weights:"
    #    for count_z,z in enumerate(fit_weights):
    #        print "{:< 12.6f} {:< 12.6f}".format(y_vals[count_z],z)

    # Calculate the rank of the fit matrix
    rank = matrix_rank(A)

    # Print diagnostics
    print("\nInitializing iterative fit of UA-epsilon values:\n")
    print("\t{:<30s} {}".format("max cycles:",max_cycles))
    print("\t{:<30s} {}".format("min_cycles:",min_cycles))
    print("\t{:<30s} {}".format("update_weight:",weight))
    print("\t{:<30s} {}".format("delta_xhi2_thresh:",delta_xhi2_thresh))
    print("\t{:<30s} {}".format("Fit rank:",rank))
    print("\t{:<30s} {}".format("Fit parameters:",len(A[0])))

    print("\ncycle 0:\n\n\txhi2: {} (kcal/mol)^2\n".format(xhi2_previous))

    # Fit loop
    while (1):
        
        # Increment the cycle count and print diagnostic
        cycle_count += 1
        print("cycle {}:\n".format(cycle_count))
        
        # Each cycle consists of looping over all pair-types being fit
        # To eliminate ordering bias, the list of types is shuffled every cycle.
        random.shuffle(Pairs)
        
        # All non-hydrogen terms are fit first, starting with self-terms then cross-terms, followed by hydrogen self-terms then hydrogen cross-terms
        Non_H_self  = [ i for i in Pairs if i[0] == i[1] and int(i[0].split('[')[1].split(']')[0]) != 1 and int(i[1].split('[')[1].split(']')[0]) != 1 ]
        Non_H_cross = [ i for i in Pairs if i[0] != i[1] and int(i[0].split('[')[1].split(']')[0]) != 1 and int(i[1].split('[')[1].split(']')[0]) != 1 ]
        H_self      = [ i for i in Pairs if i[0] == i[1] and int(i[0].split('[')[1].split(']')[0]) == 1 and int(i[1].split('[')[1].split(']')[0]) == 1 ]
        H_cross     = [ i for i in Pairs if i[0] != i[1] and (int(i[0].split('[')[1].split(']')[0]) == 1 or int(i[1].split('[')[1].split(']')[0]) == 1 ) ]
        Pairs = Non_H_self + Non_H_cross + H_self + H_cross

        for i in Pairs:

            # Set the pair-type that is being fit
            fit_type = i                            
            old_params = list(VDW_dict[i])

            # Fit the potential (For UA fits only eps is fit)

            if fit.lower() == "lj":
#                new_eps,err            = curve_fit(E_LJ_Fit_eps,x_vals,y_vals,p0=VDW_dict[i][0],maxfev=10000)
                new_eps,err            = curve_fit(E_LJ_Fit_eps,x_vals,y_vals,p0=VDW_dict[i][0],maxfev=10000,sigma=fit_weights)
                new_eps                = VDW_dict[i][0]*(1.0-weight) + new_eps[0]*weight
                VDW_dict[i]            = tuple([new_eps,VDW_dict[i][1]])
                VDW_dict[(i[1],i[0])] = tuple([new_eps,VDW_dict[i][1]])
                # COMMENT OUT THE NEXT 4 LINES FOR OLD VERSION #
#                new_sigma,err         = curve_fit(E_LJ_Fit_sigma,x_vals,y_vals,p0=VDW_dict[i][1],maxfev=10000)
                new_sigma,err         = curve_fit(E_LJ_Fit_sigma,x_vals,y_vals,p0=VDW_dict[i][1],maxfev=10000,sigma=fit_weights)
                new_sigma             = VDW_dict[i][1]*(1.0-weight) + new_sigma[0]*weight
                VDW_dict[i]           = tuple([VDW_dict[i][0],new_sigma])
                VDW_dict[(i[1],i[0])] = tuple([VDW_dict[i][0],new_sigma])

            # Fit the potential (A is fit first, then B, then C; this was pursued because fitting all of them at the same time
            # led to instability when using weighted updates.)
            if fit.lower() == "buck":
                new_A,err             = curve_fit(E_BUCK_Fit_A,x_vals,y_vals,p0=VDW_dict[i][0],maxfev=10000)
                new_A                 = VDW_dict[i][0]*(1.0-weight) + new_A[0]*weight
                VDW_dict[i]           = tuple([new_A,VDW_dict[i][1],VDW_dict[i][2]])
                VDW_dict[(i[1],i[0])] = tuple([new_A,VDW_dict[i][1],VDW_dict[i][2]])
                new_B,err             = curve_fit(E_BUCK_Fit_B,x_vals,y_vals,p0=VDW_dict[i][1],maxfev=10000)
                new_B                 = VDW_dict[i][1]*(1.0-weight) + new_B[0]*weight
                VDW_dict[i]           = tuple([VDW_dict[i][0],new_B,VDW_dict[i][2]])
                VDW_dict[(i[1],i[0])] = tuple([VDW_dict[i][0],new_B,VDW_dict[i][2]])
                new_C,err             = curve_fit(E_BUCK_Fit_C,x_vals,y_vals,p0=VDW_dict[i][2],maxfev=10000)
                new_C                 = VDW_dict[i][2]*(1.0-weight) + new_C[0]*weight
                VDW_dict[i]           = tuple([VDW_dict[i][0],VDW_dict[i][1],new_C])
                VDW_dict[(i[1],i[0])] = tuple([VDW_dict[i][0],VDW_dict[i][1],new_C])


            print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<15.4f}".format(j) for j in old_params]),", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]])))                

        # Update the fit data array
        fit_type = 0   # dummy value so that all pair-potentials are taken from VDW_dict (see E_LJ_Fit for more details)
        if fit == 'lj':
            fit_vals[cycle_count] = E_LJ_Fit(x_vals,0.0,0.0)
        elif fit == 'buck':
            fit_vals[cycle_count] = E_BUCK_Fit(x_vals,0.0,0.0,0.0)

        # Calculate mean squared deviation
        xhi2 = mean(((y_vals-fit_vals[cycle_count])/fit_weights)**2)
        delta_xhi2 = xhi2_previous - xhi2
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",xhi2))
        print("\t{:<12s} {} (kcal/mol)^2\n".format("delta_xhi2:",delta_xhi2))
        xhi2_previous = copy(xhi2)

        # Check break conditions
        if cycle_count >= min_cycles:
            if cycle_count == max_cycles:
                print("Maximum number of cycles reached, breaking iterative fit loop...\n")
                break
            elif delta_xhi2 < delta_xhi2_thresh:
                print("Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2), breaking iterative fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
                break

    # Avoid the global fit if the number of pairs being fit exceeds the number of configurations (the global fit function complains)
    # Or if there are more than 100 pairs being fit (the maximum number of parameters the function can handle is 100 eps and 100 sigma values respectively)
    if len(Pairs) > len(x_vals) or len(Pairs) > 100:        
        return fit_vals[:cycle_count+1]

    # Perform a final optimization of all parameters at once    
    global Current_Fit_Pairs

    # Copy the supplied Fit_Pairs to the local/global variable Current_Fit_Pairs
    # This is performed because the pairs being fit need to be flexibly defined locally
    Current_Fit_Pairs = Pairs

    # Prepare initial guess parameters, initial fit energies, and initial xhi2 value
    if fit == 'lj':
        initial_guess = [ VDW_dict[i][0] for i in Current_Fit_Pairs ]
#        xhi2_previous = mean((y_vals-E_LJ_Fit_all_UA(x_vals,*initial_guess))**2)
        xhi2_previous = mean((y_vals-E_LJ_Fit_all(x_vals,*initial_guess))**2)
        print("\nInitializing simultaneous fit of all LJ parameters:\n")
    elif fit == 'buck':
        initial_guess = [ j for i in Current_Fit_Pairs for j in VDW_dict[i] ]
        xhi2_previous = mean((y_vals-E_BUCK_Fit_all(x_vals,*initial_guess))**2)
        print("\nInitializing simultaneous fit of all Buckingham parameters:\n")

    # Print diagnostic
    print("\t{:<30s} {}".format("Initial xhi2:",xhi2_previous))

    # Perform parameter fit until self-consistency is acheived
    while 1:

        # Increment the cycle count and print diagnostic
        cycle_count += 1
        print("cycle {} (global fit):\n".format(cycle_count))

        # Perform LJ fit
        if fit == 'lj':
#            params,err = curve_fit(E_LJ_Fit_all_UA,x_vals,y_vals,p0=initial_guess,maxfev=10000)
#            params,err = curve_fit(E_LJ_Fit_all_UA,x_vals,y_vals,p0=initial_guess,maxfev=10000,sigma=fit_weights)
            params,err = curve_fit(E_LJ_Fit_all,x_vals,y_vals,p0=initial_guess,maxfev=10000,sigma=fit_weights)

            # Calculate fit values and xhi2 with final fit parameters
#            fit_vals[cycle_count] = E_LJ_Fit_all_UA(x_vals,*params)
            fit_vals[cycle_count] = E_LJ_Fit_all(x_vals,*params)
            xhi2 = mean(((y_vals-fit_vals[cycle_count])/fit_weights)**2)        
            delta_xhi2 = xhi2_previous - xhi2 

            # Convert fit parameters back into eps,sigma tuples
            initial_guess = params
            params = [ (params[count_i],VDW_dict[i][1]) for count_i,i in enumerate(Current_Fit_Pairs) ]

            # Print diagnostic and assign fit parameters to the VDW_dict
            for count_i,i in enumerate(Current_Fit_Pairs):
                print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<15.4f}".format(j) for j in params[count_i]])))
                VDW_dict[i] = tuple(params[count_i])  
                VDW_dict[(i[1],i[0])] = tuple(params[count_i])        

        # Perform BUCK fit
        if fit == 'buck':
            params,err = curve_fit(E_BUCK_Fit_all,x_vals,y_vals,p0=initial_guess,maxfev=10000)

            # Calculate fit values and xhi2 with final fit parameters
            fit_vals[cycle_count] = E_BUCK_Fit_all(x_vals,*params)
            xhi2 = mean((y_vals-fit_vals[cycle_count])**2)        
            delta_xhi2 = xhi2_previous - xhi2 

            # Convert fit parameters back into eps,sigma tuples
            initial_guess = params
            params = [ (params[3*count_i],params[3*count_i+1],params[3*count_i+2]) for count_i,i in enumerate(Current_Fit_Pairs) ]

            # Print diagnostic and assign fit parameters to the VDW_dict
            for count_i,i in enumerate(Current_Fit_Pairs):
                print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<15.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<15.4f}".format(j) for j in params[count_i]])))
                VDW_dict[i] = tuple(params[count_i])  
                VDW_dict[(i[1],i[0])] = tuple(params[count_i])        

        # Print change in xhi2 and update xhi2_previous
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",xhi2))
        print("\t{:<12s} {} (kcal/mol)^2\n".format("delta_xhi2:",delta_xhi2))
        xhi2_previous = copy(xhi2)
        
        # Check break conditions
        if delta_xhi2 < delta_xhi2_thresh:
            print("Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2), breaking global fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
            break

    return fit_vals[:cycle_count+1]



# Description: This function is a simple wrapper for the matplotlib plotting commands
#              to generate the scatterplot of the correlation between the QC interaction
#              energies and the FF interaction energies. The function also saves the 
#              convergence data to file.
def plot_convergence(folder,Data,fit_vals,plot_num=5,type='DFT-AA'):

    # Generate fit plot
    color_list = [(0.05,0.35,0.75),(0.05,0.8,0.6),(0.9,0.3,0.05),(0.35,0.7,0.9),(0.9,0.5,0.7),(0.9,0.6,0.05),(0.95,0.9,0.25),(0.05,0.05,0.05)]*10   # Supposedly more CB-friendly
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []

    # Generate subset of plot values ("plot_num" argument determines the number of convergence cycles
    # to explicitly plot, although all are saved to the .txt file; e.g. plot_num plots the first and last
    # cycles plus the 1/4 1/2 and 3/4 converged cycles).
    if len(fit_vals) <= plot_num:
        ind = list(range(0,len(fit_vals)))
    else:
        ind = [ int(round((len(fit_vals)-1)/float(plot_num-1)*(i))) for i in range(plot_num) ]

    #
    # Plot VDW_fit correlation plot
    #
    x_vals = [ i for i in list(Data.keys()) ]

    # Set the y_vals to the total dft energy
    if type == 'DFT-AA':
        y_vals = [ Data[i]["e_dft_fit_AA"] for i in list(Data.keys()) ]    
    if type == 'DFT-UA':
        y_vals = [ Data[i]["e_dft_fit_UA"] for i in list(Data.keys()) ]    
    if type == 'MP2-AA':
        y_vals = [ Data[i]["e_mp2_fit_AA"] for i in list(Data.keys()) ]    
    if type == 'MP2-UA':
        y_vals = [ Data[i]["e_mp2_fit_UA"] for i in list(Data.keys()) ]    

    # Plot the scatter data
    for count_i,i in enumerate(ind):   
        plot_handle, = ax.plot(y_vals,fit_vals[i],marker='.',markersize=20,color=color_list[count_i],markeredgewidth=0.0,alpha=0.3,linestyle='None',label='cycle: {}'.format(i))

    # Find min and max based on the final fit_vals
    y_min_fit = 1000000.0
    y_max_fit =-1000000.0
    for i in fit_vals:
        if min(i) < y_min_fit: y_min_fit = floor(min(i))
        if max(i) > y_min_fit: y_max_fit = ceil(max(i))        
            
    # Set limits based on the largest range
    y_min,y_max = ax.get_ylim()
    x_min,x_max = ax.get_xlim()
    
    # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
    if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
    if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

    if x_min < y_min: y_min = x_min;
    else: x_min = y_min
    if x_max > y_max: y_max = x_max;
    else: x_max = y_max
    ax.set_xlim([x_min,x_max])
    ax.set_ylim([y_min,y_max])

    # Plot the diagonal line
    ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{FF,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    if type == 'DFT-AA':
        ax.set_xlabel("$\mathrm{E_{DFT,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-AA_vdw_convergence.png'
        data_name = 'DFT-AA_vdw_convergence.txt'
    elif type == 'DFT-UA':
        ax.set_xlabel("$\mathrm{E_{DFT,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-UA_vdw_convergence.png'
        data_name = 'DFT-UA_vdw_convergence.txt'
    elif type == 'MP2-AA':
        ax.set_xlabel("$\mathrm{E_{MP2,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-AA_vdw_convergence.png'
        data_name = 'MP2-AA_vdw_convergence.txt'
    elif type == 'MP2-UA':
        ax.set_xlabel("$\mathrm{E_{MP2,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-UA_vdw_convergence.png'
        data_name = 'MP2-UA_vdw_convergence.txt'

    # Generate Legend
    ax.legend(loc='best',frameon=False)
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

    # Format ticks
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
    close(fig)

    # Save the convergence data to file
    with open(folder+'/'+data_name,'w') as f:
        if type in ['DFT-AA','DFT-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_VDW_DFT",' '.join([ "{:<15s}".format("E_INT_VDW_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        elif type in ['MP2-AA','MP2-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_VDW_MP2",' '.join([ "{:<15s}".format("E_INT_VDW_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        for count_i,i in enumerate(y_vals):
            f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(j[count_i]) for j in fit_vals ])))

    #
    # Plot Total_fit correlation plot (Total interaction energy = VDW + Coul contribution
    #
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []

    # Initialize list of indices
    x_vals = [ i for i in list(Data.keys()) ]

    # Set the y_vals to the total dft energy
    if type == 'DFT-AA' or type == 'DFT-UA':
        y_vals = [ Data[i]["e_dft"] for i in list(Data.keys()) ]    

    # Set the y_vals to the total mp2 energy
    if type == 'MP2-AA' or type == 'MP2-UA':
        y_vals = [ Data[i]["e_mp2"] for i in list(Data.keys()) ]    
    
    # Create all-atom/united-atom coulombic array
    if type == 'DFT-AA' or type == 'MP2-AA':
        coul_array = array([ E_C_Tot(i,'AA') for i in x_vals ])
    if type == 'DFT-UA' or type == 'MP2-UA':
        coul_array = array([ E_C_Tot(i,'UA') for i in x_vals ])

    # Add coulombic contribution to the fit values
    fit_vals_tot = deepcopy(fit_vals)
    for i in range(len(fit_vals)):
        fit_vals_tot[i] += coul_array 

    # Plot the scatter data
    for count_i,i in enumerate(ind):   
        plot_handle, = ax.plot(y_vals,fit_vals_tot[i],marker='.',markersize=20,color=color_list[count_i],markeredgewidth=0.0,alpha=0.3,linestyle='None',label='cycle: {}'.format(i))

    # Find min and max based on the final fit_vals
    y_min_fit = 1000000.0
    y_max_fit =-1000000.0
    for i in fit_vals_tot:
        if min(i) < y_min_fit: y_min_fit = floor(min(i))
        if max(i) > y_min_fit: y_max_fit = ceil(max(i))        
            
    # Set limits based on the largest range
    y_min,y_max = ax.get_ylim()
    x_min,x_max = ax.get_xlim()
    
    # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
    if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
    if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

    if x_min < y_min: y_min = x_min;
    else: x_min = y_min
    if x_max > y_max: y_max = x_max;
    else: x_max = y_max
    ax.set_xlim([x_min,x_max])
    ax.set_ylim([y_min,y_max])

    # Plot the diagonal line
    ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{FF,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    if type == 'DFT-AA':
        ax.set_xlabel("$\mathrm{E_{DFT,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-AA_tot_convergence.png'
        data_name = 'DFT-AA_tot_convergence.txt'
    elif type == 'DFT-UA':
        ax.set_xlabel("$\mathrm{E_{DFT,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-UA_tot_convergence.png'
        data_name = 'DFT-UA_tot_convergence.txt'
    elif type == 'MP2-AA':
        ax.set_xlabel("$\mathrm{E_{MP2,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-AA_tot_convergence.png'
        data_name = 'MP2-AA_tot_convergence.txt'
    elif type == 'MP2-UA':
        ax.set_xlabel("$\mathrm{E_{MP2,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-UA_tot_convergence.png'
        data_name = 'MP2-UA_tot_convergence.txt'

    # Generate Legend
    ax.legend(loc='best',frameon=False)
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

    # Format ticks
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
    close(fig)

    # Save the convergence data to file
    with open(folder+'/'+data_name,'w') as f:
        if type in ['DFT-AA','DFT-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_TOT_DFT",' '.join([ "{:<15s}".format("E_INT_TOT_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        elif type in ['MP2-AA','MP2-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_TOT_MP2",' '.join([ "{:<15s}".format("E_INT_TOT_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        for count_i,i in enumerate(y_vals):
            f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(j[count_i]) for j in fit_vals ])))

def plot_convergence_new(folder,Data,fit_vals,xhi_total,plot_num=5,type='DFT-AA'):

    # Generate fit plot
    color_list = [(0.05,0.35,0.75),(0.05,0.8,0.6),(0.9,0.3,0.05),(0.35,0.7,0.9),(0.9,0.5,0.7),(0.9,0.6,0.05),(0.95,0.9,0.25),(0.05,0.05,0.05)]*10   # Supposedly more CB-friendly
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []

    # Generate subset of plot values ("plot_num" argument determines the number of convergence cycles
    # to explicitly plot, although all are saved to the .txt file; e.g. plot_num plots the first and last
    # cycles plus the 1/4 1/2 and 3/4 converged cycles).

    #
    # Plot VDW_fit correlation plot
    #
    x_vals = [ i for i in list(Data.keys()) ]

    # Set the y_vals to the total dft energy
    if type == 'DFT-AA':
        y_vals = [ Data[i]["e_dft_fit_AA"] for i in list(Data.keys()) ]    
    if type == 'DFT-UA':
        y_vals = [ Data[i]["e_dft_fit_UA"] for i in list(Data.keys()) ]    
    if type == 'MP2-AA':
        y_vals = [ Data[i]["e_mp2_fit_AA"] for i in list(Data.keys()) ]    
    if type == 'MP2-UA':
        y_vals = [ Data[i]["e_mp2_fit_UA"] for i in list(Data.keys()) ]    

    # Plot the scatter data
    for count_i,i in enumerate(list(fit_vals.keys())):   
        legend = "$\mathrm{{{{{0}}}, xhi^{{2}}: {{{1:<20.4f}}}  (kcal/mol)^{{2}}}}$".format(i,xhi_total[i])
        plot_handle, = ax.plot(y_vals,fit_vals[i],marker='.',markersize=18,markerfacecolor='none',markeredgewidth=2.0,alpha=0.3,linestyle='None',label=legend,markeredgecolor=color_list[count_i])

    # Find min and max based on the final fit_vals
    y_min_fit = 1000000.0
    y_max_fit =-1000000.0
    for i in fit_vals:
        if min(fit_vals[i]) < y_min_fit: y_min_fit = floor(min(fit_vals[i]))
        if max(fit_vals[i]) > y_min_fit: y_max_fit = ceil(max(fit_vals[i]))        
            
    # Set limits based on the largest range
    y_min,y_max = ax.get_ylim()
    x_min,x_max = ax.get_xlim()
    
    # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
    if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
    if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

    if x_min < y_min: y_min = x_min;
    else: x_min = y_min
    if x_max > y_max: y_max = x_max;
    else: x_max = y_max
    ax.set_xlim([x_min,x_max])
    ax.set_ylim([y_min,y_max])

    # Plot the diagonal line
    ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{FF,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    if type == 'DFT-AA':
        ax.set_xlabel("$\mathrm{E_{DFT,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-AA_vdw_convergence.png'
        data_name = 'DFT-AA_vdw_convergence.txt'
    elif type == 'DFT-UA':
        ax.set_xlabel("$\mathrm{E_{DFT,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-UA_vdw_convergence.png'
        data_name = 'DFT-UA_vdw_convergence.txt'
    elif type == 'MP2-AA':
        ax.set_xlabel("$\mathrm{E_{MP2,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-AA_vdw_convergence.png'
        data_name = 'MP2-AA_vdw_convergence.txt'
    elif type == 'MP2-UA':
        ax.set_xlabel("$\mathrm{E_{MP2,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-UA_vdw_convergence.png'
        data_name = 'MP2-UA_vdw_convergence.txt'

    # Generate Legend
    ax.legend(loc='best',frameon=False)
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

    # Format ticks
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
    close(fig)

    # Save the convergence data to file
    with open(folder+'/'+data_name,'w') as f:
        if type in ['DFT-AA','DFT-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_VDW_DFT",' '.join([ "{:<15s}".format("E_INT_VDW_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        elif type in ['MP2-AA','MP2-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_VDW_MP2",' '.join([ "{:<15s}".format("E_INT_VDW_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        for count_i,i in enumerate(y_vals):
            f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(fit_vals[j][count_i]) for j in fit_vals ])))

    #
    # Plot Total_fit correlation plot (Total interaction energy = VDW + Coul contribution
    #
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []

    # Initialize list of indices
    x_vals = [ i for i in list(Data.keys()) ]

    # Set the y_vals to the total dft energy
    if type == 'DFT-AA' or type == 'DFT-UA':
        y_vals = [ Data[i]["e_dft"] for i in list(Data.keys()) ]    

    # Set the y_vals to the total mp2 energy
    if type == 'MP2-AA' or type == 'MP2-UA':
        y_vals = [ Data[i]["e_mp2"] for i in list(Data.keys()) ]    
    
    # Create all-atom/united-atom coulombic array
    if type == 'DFT-AA' or type == 'MP2-AA':
        coul_array = array([ E_C_Tot(i,'AA') for i in x_vals ])
    if type == 'DFT-UA' or type == 'MP2-UA':
        coul_array = array([ E_C_Tot(i,'UA') for i in x_vals ])


    fit_vals_tot = deepcopy(fit_vals)
    # Add coulombic contribution to the fit values
    for i in fit_vals:
        fit_vals_tot[i] += coul_array 

    # Plot the scatter data
    for count_i,i in enumerate(list(fit_vals.keys())):   
        legend = "{}".format(i)
        plot_handle, = ax.plot(y_vals,fit_vals_tot[i],marker='.',markersize=18,markerfacecolor='none',markeredgewidth=2.0,alpha=0.3,linestyle='None',label=legend,markeredgecolor=color_list[count_i])

    # Find min and max based on the final fit_vals
    y_min_fit = 1000000.0
    y_max_fit =-1000000.0
    for i in fit_vals_tot:
        if min(fit_vals_tot[i]) < y_min_fit: y_min_fit = floor(min(fit_vals_tot[i]))
        if max(fit_vals_tot[i]) > y_min_fit: y_max_fit = ceil(max(fit_vals_tot[i]))        
            
    # Set limits based on the largest range
    y_min,y_max = ax.get_ylim()
    x_min,x_max = ax.get_xlim()
    
    # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
    if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
    if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

    if x_min < y_min: y_min = x_min;
    else: x_min = y_min
    if x_max > y_max: y_max = x_max;
    else: x_max = y_max
    ax.set_xlim([x_min,x_max])
    ax.set_ylim([y_min,y_max])

    # Plot the diagonal line
    ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{FF,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    if type == 'DFT-AA':
        ax.set_xlabel("$\mathrm{E_{DFT,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-AA_tot_convergence.png'
        data_name = 'DFT-AA_tot_convergence.txt'
    elif type == 'DFT-UA':
        ax.set_xlabel("$\mathrm{E_{DFT,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-UA_tot_convergence.png'
        data_name = 'DFT-UA_tot_convergence.txt'
    elif type == 'MP2-AA':
        ax.set_xlabel("$\mathrm{E_{MP2,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-AA_tot_convergence.png'
        data_name = 'MP2-AA_tot_convergence.txt'
    elif type == 'MP2-UA':
        ax.set_xlabel("$\mathrm{E_{MP2,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-UA_tot_convergence.png'
        data_name = 'MP2-UA_tot_convergence.txt'

    # Generate Legend
    ax.legend(loc='best',frameon=False)
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

    # Format ticks
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
    close(fig)

    # Save the convergence data to file
    with open(folder+'/'+data_name,'w') as f:
        if type in ['DFT-AA','DFT-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_TOT_DFT",' '.join([ "{:<15s}".format("E_INT_TOT_FF_"+str(i)) for count_i,i in enumerate(fit_vals) ])))
        elif type in ['MP2-AA','MP2-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_TOT_MP2",' '.join([ "{:<15s}".format("E_INT_TOT_FF_"+str(i)) for count_i,i in enumerate(fit_vals) ])))
        for count_i,i in enumerate(y_vals):
            f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(fit_vals[j][count_i]) for j in fit_vals ])))

def plot_convergence_outlier(folder,Data,fit_vals,xhi_total,plot_num=5,type='DFT-AA'):

    # Generate fit plot
    color_list = [(0.05,0.35,0.75),(0.05,0.8,0.6),(0.9,0.3,0.05),(0.35,0.7,0.9),(0.9,0.5,0.7),(0.9,0.6,0.05),(0.95,0.9,0.25),(0.05,0.05,0.05)]*10   # Supposedly more CB-friendly
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []

    # Generate subset of plot values ("plot_num" argument determines the number of convergence cycles
    # to explicitly plot, although all are saved to the .txt file; e.g. plot_num plots the first and last
    # cycles plus the 1/4 1/2 and 3/4 converged cycles).

    #
    # Plot VDW_fit correlation plot
    #
    x_vals = [ i for i in list(Data.keys()) ]

    # Set the y_vals to the total dft energy
    if type == 'DFT-AA':
        y_vals = [ Data[i]["e_dft_fit_AA"] for i in list(Data.keys()) ]    
    if type == 'DFT-UA':
        y_vals = [ Data[i]["e_dft_fit_UA"] for i in list(Data.keys()) ]    
    if type == 'MP2-AA':
        y_vals = [ Data[i]["e_mp2_fit_AA"] for i in list(Data.keys()) ]    
    if type == 'MP2-UA':
        y_vals = [ Data[i]["e_mp2_fit_UA"] for i in list(Data.keys()) ]    
   
    Mi,outlier = detect_outlier(y_vals)
    print("Number of outliers removed from figure: {}".format(len([i for i in outlier if i])))
    # pop out outlier
    y_vals = [ i for count_i,i in enumerate(y_vals) if not outlier[count_i]]
    pop_index = [ count_i for count_i,i in enumerate(outlier) if i]
    for j in fit_vals:
        fit_vals[j] = np.delete(fit_vals[j],pop_index)

    # Plot the scatter data
    for count_i,i in enumerate(list(fit_vals.keys())):   
        legend = "$\mathrm{{{{{0}}}, xhi^{{2}}: {{{1:<20.4f}}}  (kcal/mol)^{{2}}}}$".format(i,xhi_total[i])
        plot_handle, = ax.plot(y_vals,fit_vals[i],marker='.',markersize=18,markerfacecolor='none',markeredgewidth=2.0,alpha=0.3,linestyle='None',label=legend,markeredgecolor=color_list[count_i])

    # Find min and max based on the final fit_vals
    y_min_fit = 1000000.0
    y_max_fit =-1000000.0
    for i in fit_vals:
        if min(fit_vals[i]) < y_min_fit: y_min_fit = floor(min(fit_vals[i]))
        if max(fit_vals[i]) > y_min_fit: y_max_fit = ceil(max(fit_vals[i]))        
            
    # Set limits based on the largest range
    y_min,y_max = ax.get_ylim()
    x_min,x_max = ax.get_xlim()
    
    # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
    if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
    if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

    if x_min < y_min: y_min = x_min;
    else: x_min = y_min
    if x_max > y_max: y_max = x_max;
    else: x_max = y_max
    ax.set_xlim([x_min,x_max])
    ax.set_ylim([y_min,y_max])

    # Plot the diagonal line
    ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{FF,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    if type == 'DFT-AA':
        ax.set_xlabel("$\mathrm{E_{DFT,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-AA_vdw_convergence_outlier.png'
        data_name = 'DFT-AA_vdw_convergence_outlier.txt'
    elif type == 'DFT-UA':
        ax.set_xlabel("$\mathrm{E_{DFT,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-UA_vdw_convergence.png'
        data_name = 'DFT-UA_vdw_convergence.txt'
    elif type == 'MP2-AA':
        ax.set_xlabel("$\mathrm{E_{MP2,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-AA_vdw_convergence.png'
        data_name = 'MP2-AA_vdw_convergence.txt'
    elif type == 'MP2-UA':
        ax.set_xlabel("$\mathrm{E_{MP2,VDW} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-UA_vdw_convergence.png'
        data_name = 'MP2-UA_vdw_convergence.txt'

    # Generate Legend
    ax.legend(loc='best',frameon=False)
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

    # Format ticks
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
    close(fig)

    # Save the convergence data to file
    with open(folder+'/'+data_name,'w') as f:
        if type in ['DFT-AA','DFT-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_VDW_DFT",' '.join([ "{:<15s}".format("E_INT_VDW_FF_"+str(i)) for count_i,i in enumerate(fit_vals) ])))
        elif type in ['MP2-AA','MP2-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_VDW_MP2",' '.join([ "{:<15s}".format("E_INT_VDW_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        for count_i,i in enumerate(y_vals):
            f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(fit_vals[j][count_i]) for j in fit_vals ])))

    #
    # Plot Total_fit correlation plot (Total interaction energy = VDW + Coul contribution
    #
    fig = plt.figure(figsize=(6,5))
    ax = plt.subplot(111)
    plot_handles = []

    # Initialize list of indices
    x_vals = [ i for i in list(Data.keys()) ]

    # Set the y_vals to the total dft energy
    if type == 'DFT-AA' or type == 'DFT-UA':
        y_vals = [ Data[i]["e_dft"] for i in list(Data.keys()) ]    

    # Set the y_vals to the total mp2 energy
    if type == 'MP2-AA' or type == 'MP2-UA':
        y_vals = [ Data[i]["e_mp2"] for i in list(Data.keys()) ]    
    
    # Create all-atom/united-atom coulombic array
    if type == 'DFT-AA' or type == 'MP2-AA':
        coul_array = array([ E_C_Tot(i,'AA') for i in x_vals ])
    if type == 'DFT-UA' or type == 'MP2-UA':
        coul_array = array([ E_C_Tot(i,'UA') for i in x_vals ])

    coul_array = np.delete(coul_array,[count_i for count_i,i in enumerate(outlier) if i])
    # Add coulombic contribution to the fit values
    fit_vals_tot = deepcopy(fit_vals)
    for i in fit_vals:
        fit_vals_tot[i] += coul_array

    y_vals = [ i for count_i,i in enumerate(y_vals) if not outlier[count_i]]

    # Plot the scatter data
    for count_i,i in enumerate(list(fit_vals.keys())):   
        legend = "{}".format(i)
        plot_handle, = ax.plot(y_vals,fit_vals_tot[i],marker='.',markersize=18,markerfacecolor='none',markeredgewidth=2.0,alpha=0.3,linestyle='None',label=legend,markeredgecolor=color_list[count_i])

    # Find min and max based on the final fit_vals
    y_min_fit = 1000000.0
    y_max_fit =-1000000.0
    for i in fit_vals:
        if min(fit_vals_tot[i]) < y_min_fit: y_min_fit = floor(min(fit_vals_tot[i]))
        if max(fit_vals_tot[i]) > y_min_fit: y_max_fit = ceil(max(fit_vals_tot[i]))        
            
    # Set limits based on the largest range
    y_min,y_max = ax.get_ylim()
    x_min,x_max = ax.get_xlim()
    
    # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
    if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
    if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

    if x_min < y_min: y_min = x_min;
    else: x_min = y_min
    if x_max > y_max: y_max = x_max;
    else: x_max = y_max
    ax.set_xlim([x_min,x_max])
    ax.set_ylim([y_min,y_max])

    # Plot the diagonal line
    ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{FF,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    if type == 'DFT-AA':
        ax.set_xlabel("$\mathrm{E_{DFT,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-AA_tot_convergence_outlier.png'
        data_name = 'DFT-AA_tot_convergence_outlier.txt'
    elif type == 'DFT-UA':
        ax.set_xlabel("$\mathrm{E_{DFT,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'DFT-UA_tot_convergence.png'
        data_name = 'DFT-UA_tot_convergence.txt'
    elif type == 'MP2-AA':
        ax.set_xlabel("$\mathrm{E_{MP2,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-AA_tot_convergence.png'
        data_name = 'MP2-AA_tot_convergence.txt'
    elif type == 'MP2-UA':
        ax.set_xlabel("$\mathrm{E_{MP2,TOT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        plot_name = 'MP2-UA_tot_convergence.png'
        data_name = 'MP2-UA_tot_convergence.txt'

    # Generate Legend
    ax.legend(loc='best',frameon=False)
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

    # Format ticks
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
    close(fig)

    # Save the convergence data to file
    with open(folder+'/'+data_name,'w') as f:
        if type in ['DFT-AA','DFT-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_TOT_DFT",' '.join([ "{:<15s}".format("E_INT_TOT_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        elif type in ['MP2-AA','MP2-UA']:
            f.write(' {:<15s} {}\n'.format("E_INT_TOT_MP2",' '.join([ "{:<15s}".format("E_INT_TOT_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
        for count_i,i in enumerate(y_vals):
            f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(fit_vals[j][count_i]) for j in fit_vals ])))

# # Description: This function is a simple wrapper for the matplotlib plotting commands
# #              to generate the scatterplot of the correlation between the QC interaction
# #              energies and the FF interaction energies. The function also saves the 
# #              convergence data to file.
# def plot_convergence(folder,y_vals,fit_vals,plot_num=5,type='DFT-AA'):

#     # Generate fit plot
#     color_list = [(0.05,0.35,0.75),(0.05,0.8,0.6),(0.9,0.3,0.05),(0.35,0.7,0.9),(0.9,0.5,0.7),(0.9,0.6,0.05),(0.95,0.9,0.25),(0.05,0.05,0.05)]*10   # Supposedly more CB-friendly
#     fig = plt.figure(figsize=(6,5))
#     ax = plt.subplot(111)
#     plot_handles = []

#     # Generate subset of plot values ("plot_num" argument determines the number of convergence cycles
#     # to explicitly plot, although all are saved to the .txt file; e.g. plot_num plots the first and last
#     # cycles plus the 1/4 1/2 and 3/4 converged cycles).
#     if len(fit_vals) <= plot_num:
#         ind = range(0,len(fit_vals))
#     else:
#         ind = [ int(round((len(fit_vals)-1)/float(plot_num-1)*(i))) for i in range(plot_num) ]

#     # Plot the scatter data
#     for count_i,i in enumerate(ind):   
#         plot_handle, = ax.plot(y_vals,fit_vals[i],marker='.',markersize=20,color=color_list[count_i],markeredgewidth=0.0,alpha=0.3,linestyle='None',label='cycle: {}'.format(i))

#     # Find min and max based on the final fit_vals
#     y_min_fit = 1000000.0
#     y_max_fit =-1000000.0
#     for i in fit_vals:
#         if min(i) < y_min_fit: y_min_fit = floor(min(i))
#         if max(i) > y_min_fit: y_max_fit = ceil(max(i))        
            
#     # Set limits based on the largest range
#     y_min,y_max = ax.get_ylim()
#     x_min,x_max = ax.get_xlim()
    
#     # If y_min or y_max are impractically larger than the limits of the final fit values, then the y limits are reassigned
#     if abs(y_min-y_min_fit) > 10: y_min = y_min_fit
#     if abs(y_max-y_max_fit) > 10: y_max = y_max_fit

#     if x_min < y_min: y_min = x_min;
#     else: x_min = y_min
#     if x_max > y_max: y_max = x_max;
#     else: x_max = y_max
#     ax.set_xlim([x_min,x_max])
#     ax.set_ylim([y_min,y_max])

#     # Plot the diagonal line
#     ax.plot(ax.get_xlim(), ax.get_ylim(), ls="-", c="0")

#     # Set Labels and Save Name
#     ax.set_ylabel("$\mathrm{E_{FF} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
#     if type == 'DFT-AA':
#         ax.set_xlabel("$\mathrm{E_{DFT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
#         plot_name = 'DFT-AA_convergence.png'
#         data_name = 'DFT-AA_convergence.txt'
#     elif type == 'DFT-UA':
#         ax.set_xlabel("$\mathrm{E_{DFT} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
#         plot_name = 'DFT-UA_convergence.png'
#         data_name = 'DFT-UA_convergence.txt'
#     elif type == 'MP2-AA':
#         ax.set_xlabel("$\mathrm{E_{MP2} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
#         plot_name = 'MP2-AA_convergence.png'
#         data_name = 'MP2-AA_convergence.txt'
#     elif type == 'MP2-UA':
#         ax.set_xlabel("$\mathrm{E_{MP2} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
#         plot_name = 'MP2-UA_convergence.png'
#         data_name = 'MP2-UA_convergence.txt'

#     # Generate Legend
#     ax.legend(loc='best',frameon=False)
#     handles, labels = ax.get_legend_handles_labels()
#     lgd = ax.legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))

#     # Format ticks
#     ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
#     ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
#     [j.set_linewidth(2) for j in ax.spines.itervalues()]

#     # Save the figure
#     savefig(folder+'/'+plot_name, dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')
#     close(fig)

#     # Save the convergence data to file
#     with open(folder+'/'+data_name,'w') as f:
#         if type in ['DFT-AA','DFT-UA']:
#             f.write(' {:<15s} {}\n'.format("E_INT_DFT",' '.join([ "{:<15s}".format("E_INT_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
#         elif type in ['MP2-AA','MP2-UA']:
#             f.write(' {:<15s} {}\n'.format("E_INT_MP2",' '.join([ "{:<15s}".format("E_INT_FF_"+str(count_i)) for count_i,i in enumerate(fit_vals) ])))
#         for count_i,i in enumerate(y_vals):
#             f.write('{:< 15.8f} {}\n'.format(i,' '.join([ "{:< 15.8f}".format(j[count_i]) for j in fit_vals ])))
    

# Description: A simple wrapper for the plotting commands to generate the pair
#              separation histograms. The function takes arguments for the histogramming
#              specifics, calculates the histogram, plots the data, and saves the data
#              to a text file. 
def pair_hist_plot(folder,Pairs,Data,r_min=0.0,r_max=15.0,bin_width=0.1,QC_type="DFT"):

    # Initialize Pair_Hist List (note: using [array]*x to form an expanded list weirdly
    # yields linked arrays. Not sure about the origin of this behavior but initializing
    # the arrays separately through a list expansion avoids the problem.
    Pair_histograms = [zeros(int(float(r_max-r_min)/bin_width+1)) for i in Pairs ]

    # Iterate over all configurations
    for d in list(Data.keys()):

        # Calculate pair-wise separations
        r_dist = cdist(Data[d]['geo_a'],Data[d]['geo_b'])

        # Histogram all pairs in the current configuration
        for count_i,i in enumerate(r_dist):
            for count_j,j in enumerate(i):

                # Skip long and short pairs
                if j > r_max or j < r_min:
                    continue

                # Order the pair to conform to VDW_dict key definitions
                if Data[d]['types_a'][count_i] >= Data[d]['types_b'][count_j]:
                    pair_type = (Data[d]['types_a'][count_i],Data[d]['types_b'][count_j])
                else:
                    pair_type = (Data[d]['types_b'][count_j],Data[d]['types_a'][count_i])
                    
                # Increment the proper pair histogram, bin index is based on forward binning
                # NOTE: if a FF_db was supplied, only pair_types not in FF_db will be histogrammed.
                if pair_type in Pairs:
                    bin_index = int(floor((j-r_min)/bin_width))     
                    pair_index = Pairs.index(pair_type)
                    Pair_histograms[pair_index][bin_index] += 1
    
    # Generate pair-separation histograms for each pair type
    for count_i,i in enumerate(Pairs):

        # Initialize figure
        fig = plt.figure(figsize=(6,4))

        # Plot pair-separation histogram
        bins = arange(r_min,r_max+bin_width,bin_width)
        ax = plt.subplot(111)
        ax.bar(bins, Pair_histograms[count_i], width=bin_width, linewidth=0.5, color = (0.3,0.4,1.0))
        ax.set_xlim([r_min, r_max])
        ax.set_xlabel(r'$r\mathrm{_{ij} \, (\AA )}$',fontsize=22,fontweight='bold',labelpad=5)
        ax.set_ylabel(r'$\mathrm{Number}$',fontsize=22,fontweight='bold',labelpad=5)
        ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
        ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
        [j.set_linewidth(2) for j in ax.spines.values()]
        plt.tight_layout()
        Name = folder+"/{}_{}-{}_sep_hist.png".format(QC_type,i[0],i[1])
        savefig(Name,bbox_inches=0,dpi=300)
        close(fig)

        # Save pair-separation data as a textfile
        with open(folder+"/{}_{}-{}_sep_hist.txt".format(QC_type,i[0],i[1]),'w') as f:
            f.write("{:<20s} {:<20s}\n".format("r_ij(ang)","Number"))
            for count_j,j in enumerate(bins):
                f.write("{: <20.8f} {:<20d}\n".format(j,int(Pair_histograms[count_i][count_j])))

def plot_fit_potentials(folder,fit='lj',type='DFT-AA',r_min=0.1,r_max=10.0):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict

    # Initialize the rij range
    seps = arange(r_min,r_max,0.01)

    # Plot all pair potentials
    # NOTE: write_params call cleans up VDW_dict.keys to remove redundant pair_types 
    for i in list(VDW_dict.keys()):

        if fit == 'lj':
            y_vals = 4*VDW_dict[i][0]*( (VDW_dict[i][1]/seps)**12 - (VDW_dict[i][1]/seps)**6 )
        elif fit == 'buck':
            y_vals = VDW_dict[i][0]*exp(-seps/VDW_dict[i][1]) - VDW_dict[i][2]*seps**(-6.0) 

        # Initialize figure and generate plot
        fig = plt.figure(figsize=(6,4))
        ax = plt.subplot(111)
        ax.plot(seps,y_vals,color=(0.05,0.35,0.75),linestyle='-',linewidth=2.0)
        
        # Set limits based on the largest range
        y_min = -5.0
        if fit == 'lj':
            while (1):
                if y_min > min(y_vals):
                    y_min -= 5
                else:
                    break
        ax.set_xlim([0,r_max])
        ax.set_ylim([y_min,5])

        # Set Labels and Save Name
        ax.set_ylabel("$\mathrm{E_{LJ} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        ax.set_xlabel("$r\mathrm{_{ij} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')

        # Set save name
        if type == 'DFT-AA':
            plot_name = i[0]+'-'+i[1]+'_DFT-AA.png'
        elif type == 'DFT-UA':
            plot_name = i[0]+'-'+i[1]+'_DFT-UA.png'
        elif type == 'MP2-AA':
            plot_name = i[0]+'-'+i[1]+'_MP2-AA.png'
        elif type == 'MP2-UA':
            plot_name = i[0]+'-'+i[1]+'_MP2-UA.png'

        # Format ticks and plot box
        ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
        ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
        [j.set_linewidth(2) for j in ax.spines.values()]

        # Save the figure
        savefig(folder+'/'+plot_name, dpi=300, bbox_inches='tight')
        close(fig)
        
def plot_fit_potentials_method(VDW_dict_multi,folder,fit='lj',type='DFT-AA',r_min=0.1,r_max=10.0):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict

    # Initialize the rij range
    seps = arange(r_min,r_max,0.01)

    # Plot all pair potentials
    # NOTE: write_params call cleans up VDW_dict.keys to remove redundant pair_types 
    print(VDW_dict.keys()) 
    print(VDW_dict_multi['global'].keys())
    #print(VDW_dict_multi['Boltzmann'].keys())
    for i in list(VDW_dict.keys()):

        y_vals = {}
        # Initialize figure and generate plot
        num_method = len(VDW_dict_multi.keys())
        fig,ax = plt.subplots(1,num_method,figsize=(6*num_method,4))
        if num_method == 1:
            ax = [ax]
        
        for count_j,j in enumerate(list(VDW_dict_multi.keys())):

           if fit == 'lj':
               y_vals[j] = 4*VDW_dict_multi[j][i][0]*( (VDW_dict_multi[j][i][1]/seps)**12 - (VDW_dict_multi[j][i][1]/seps)**6 )
           elif fit == 'buck':
               y_vals[j] = VDW_dict_multi[j][i][0]*exp(-seps/VDW_dict_multi[j][i][1]) - VDW_dict_multi[j][i][2]*seps**(-6.0) 

           ax[count_j].plot(seps,y_vals[j],color=(0.05,0.35,0.75),linestyle='-',linewidth=2.0)
           
           # Set limits based on the largest range
           y_min = -5.0
           if fit == 'lj':
               while (1):
                   if y_min > min(y_vals[j]):
                       y_min -= 5
                   elif y_min <-20:
                       break
                   else:
                       break
           ax[count_j].set_xlim([0,r_max])
           ax[count_j].set_ylim([y_min,5])

           # Set Labels and Save Name
           ax[count_j].set_xlabel("$r\mathrm{_{ij}}$($\mathring{A}$)",fontsize=32,labelpad=10,fontweight='bold')
           ax[count_j].set_title(j,fontsize=32,fontweight='bold')

           # Set arrow for eps and sigma
           Ro = 2**(1.0/6)*VDW_dict_multi[j][i][1]
           if(VDW_dict_multi[j][i][1]>r_max):
               eps_tmp = VDW_dict_multi[j][i][0]
               sig_tmp = VDW_dict_multi[j][i][1]
               at = AnchoredText("$\epsilon$ = {{{0:6.4f}}} (kcal/mol) \n$\sigma$ = {{{1:6.4f}}} ($\mathring{{A}}$)".format(eps_tmp,sig_tmp), prop=dict(size=15), frameon=True,loc='upper right')
               ax[count_j].add_artist(at)
           else:
               ax[count_j].annotate('', xy=(Ro,0 ), xycoords='data', xytext=(Ro, -VDW_dict_multi[j][i][0]), textcoords='data', arrowprops=dict(arrowstyle='-',color='red'))
               ax[count_j].hlines(0, 0, r_max, linestyles='dashed')
               ax[count_j].annotate(r'$\epsilon$ = {:<20.4f}'.format(VDW_dict_multi[j][i][0]),xy=(Ro-1,-VDW_dict_multi[j][i][0]-0.5),xycoords='data')
               ax[count_j].annotate('', xy=(0,0.1), xycoords='data', xytext=(VDW_dict_multi[j][i][1],0.1), textcoords='data', arrowprops=dict(arrowstyle='<->',color='red'))
               ax[count_j].annotate(r'$\sigma$ = {:<20.4f}'.format(VDW_dict_multi[j][i][1]),xy=(VDW_dict_multi[j][i][1]/2-0.1,0.2),xycoords='data')


           # Format ticks and plot box
           ax[count_j].tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
           ax[count_j].tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
           [k.set_linewidth(2) for k in ax[count_j].spines.values()]
        ax[0].set_ylabel("$\mathrm{E_{LJ} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
        # Set save name
        if type == 'DFT-AA':
            plot_name = i[0]+'-'+i[1]+'_DFT-AA.png'
        elif type == 'DFT-UA':
            plot_name = i[0]+'-'+i[1]+'_DFT-UA.png'
        elif type == 'MP2-AA':
            plot_name = i[0]+'-'+i[1]+'_MP2-AA.png'
        elif type == 'MP2-UA':
            plot_name = i[0]+'-'+i[1]+'_MP2-UA.png'

        # Save the figure
        savefig(folder+'/'+plot_name, dpi=300, bbox_inches='tight')
        close(fig)

# Description: A simple wrapper for the file writing commands for the vdw parameters.
#              The function takes arguments for the filename, fit type (e.g. lj), and
#              united-atom (UA) vs all-atom (AA) style.
def write_params(name,fit,type='AA',write_charges=0):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,Data,min_pairs

    # Collect minimum types
    min_types = []
    for i in min_pairs:
        for j in i:
            min_types += [j,j+'-UA']
    min_types = set(min_types)
    
    # First clean up VDW_dict by removing reverse types with the wrong priority ordering
    # NOTE: parameters are saved using type_a > type_b ordering
    # NOTE: Hydrogen containing pairs are removed when the fit_type is 'UA'
    pop_list = []
    keys = list(VDW_dict.keys())
    for i in keys:
        if i[0] > i[1]:
            if (i[1],i[0]) in list(VDW_dict.keys()): VDW_dict.pop((i[1],i[0]))
        if type=="UA" and ( return_UA_H(i[0]) == 1 or return_UA_H(i[1]) == 1 ):
            if i in list(VDW_dict.keys()): VDW_dict.pop(i)
            if (i[1],i[0]) in list(VDW_dict.keys()): VDW_dict.pop((i[1],i[0]))

    # Save params
    with open(name+'.db','w') as f:

        # Write VDW definitions
        if fit == 'lj':
            f.write("\n# VDW definitions\n#\n{:10s} {:60s} {:60s} {:11s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","Potential","eps (kcal/mol)","sigma (ang)"))        
        if fit == 'buck':
            f.write("\n# VDW definitions\n#\n{:10s} {:60s} {:60s} {:11s} {:20s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","Potential","A (kcal/mol)","B (ang)","C (ang kcal/mol)"))        
        if type == 'AA':
            for i in list(VDW_dict.keys()):
                f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0]),str(i[1]),fit," ".join([ "{:< 20.6f}".format(j) for j in VDW_dict[i]])))        
        elif type == 'UA':
            for i in list(VDW_dict.keys()):
                f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0])+'-UA',str(i[1])+'-UA',fit," ".join([ "{:< 20.6f}".format(j) for j in VDW_dict[i]])))        

        # Write charges is the flag is enabled
        if write_charges == 1:

            # Collect AA-charges and UA-charges
            atom_types_AA = []
            atom_types_UA = []
            charges_AA    = []
            charges_UA    = []
            for i in list(Data.keys()):
                for count_j,j in enumerate(Data[i]["types_a"]):
                    if j not in atom_types_AA:
                        atom_types_AA += [j]
                        charges_AA    += [Data[i]["charges_a"][count_j]]
                        if return_UA_H(j) != 1:
                            atom_types_UA += [j+'-UA']
                            charges_UA    += [Data[i]["charges_a_UA"][count_j]]
                for count_j,j in enumerate(Data[i]["types_b"]):
                    if j not in atom_types_AA:
                        atom_types_AA += [j]
                        charges_AA    += [Data[i]["charges_b"][count_j]]
                        if return_UA_H(j) != 1:
                            atom_types_UA += [j+'-UA']
                            charges_UA    += [Data[i]["charges_b_UA"][count_j]]
                    
            # Write charges
            f.write("\n# Charge definitions\n#\n{:10s} {:41s} {:20s}\n".format("#","Atom_type","Charge"))
            if type == 'AA':
                for count_i,i in enumerate(atom_types_AA):
                    f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_AA[count_i]))
            if type == 'UA':
                for count_i,i in enumerate(atom_types_UA):
                    f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_UA[count_i]))

    # Save minimum pair info
    with open(name+'-min.db','w') as f:

        # Write VDW definitions
        if fit == 'lj':
            f.write("\n# VDW definitions\n#\n{:10s} {:60s} {:60s} {:11s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","Potential","eps (kcal/mol)","sigma (ang)"))        
        if fit == 'buck':
            f.write("\n# VDW definitions\n#\n{:10s} {:60s} {:60s} {:11s} {:20s} {:20s} {:20s}\n".format("#","Atom_type","Atom_type","Potential","A (kcal/mol)","B (ang)","C (ang kcal/mol)"))        
        if type == 'AA':
            for i in [ j for j in list(VDW_dict.keys()) if j in min_pairs ]:
                f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0]),str(i[1]),fit," ".join([ "{:< 20.6f}".format(j) for j in VDW_dict[i]])))        
        elif type == 'UA':
            for i in [ j for j in list(VDW_dict.keys()) if j in min_pairs ]:
                f.write("{:10s} {:60s} {:60s} {:10s} {}\n".format("vdw",str(i[0])+'-UA',str(i[1])+'-UA',fit," ".join([ "{:< 20.6f}".format(j) for j in VDW_dict[i]])))        

        # Write charges is the flag is enabled
        if write_charges == 1:

            # Collect AA-charges and UA-charges
            atom_types_AA = []
            atom_types_UA = []
            charges_AA    = []
            charges_UA    = []
            for i in list(Data.keys()):
                for count_j,j in enumerate(Data[i]["types_a"]):
                    if j not in atom_types_AA:
                        atom_types_AA += [j]
                        charges_AA    += [Data[i]["charges_a"][count_j]]
                        if return_UA_H(j) != 1:
                            atom_types_UA += [j+'-UA']
                            charges_UA    += [Data[i]["charges_a_UA"][count_j]]
                for count_j,j in enumerate(Data[i]["types_b"]):
                    if j not in atom_types_AA:
                        atom_types_AA += [j]
                        charges_AA    += [Data[i]["charges_b"][count_j]]
                        if return_UA_H(j) != 1:
                            atom_types_UA += [j+'-UA']
                            charges_UA    += [Data[i]["charges_b_UA"][count_j]]
                    
            # Write charges
            f.write("\n# Charge definitions\n#\n{:10s} {:41s} {:20s}\n".format("#","Atom_type","Charge"))
            if type == 'AA':
                for count_i,i in enumerate(atom_types_AA):
                    if i not in min_types:
                        continue
                    else:
                        f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_AA[count_i]))
            if type == 'UA':
                for count_i,i in enumerate(atom_types_UA):
                    if i not in min_types:
                        continue
                    else:
                        f.write("{:10s} {:40s} {:< 20.6f}\n".format("charge",str(i),charges_UA[count_i]))

# Description: This function searches through the configurations being parsed for unique
#              molecules. All FF_parameters are linked examples of the molecule(s) used
#              for their parametrization. 
# Note:        identifying unique molecules is generally a graph isomorphism problem
#              (P difficult). The algorithm used here assumes that replica molecules have
#              same atomic ordering (which the vdw_gen.py parser should guarrantee) which 
#              makes identifying the molecule types much easier.
#             
# Inputs:      folder: folder for saving the molecule.db file
#              Data:   the dictionary holding the configuration data 
def parse_molecules(folder,Data):

    # Import min_pairs
    global min_pairs

    # All atom types are specified in terms of atomic number. This dictionary is used to convert them to the corresponding element label. 
    atom2element_dict = { 1:'H' ,  2:'He',
                          3:'Li',  4:'Be',                                                                                            5:'B' ,  6:'C' ,  7:'N' ,  8:'O' ,  9:'F' , 10:'Ne',
                         11:'Na', 12:'Mg',                                                                                           13:'Al', 14:'Si', 15:'P' , 16:'S' , 17:'Cl', 18:'Ar',
                         19:'K' , 20:'Ca', 21:'Sc', 22:'Ti', 23:'V' , 24:'Cr', 25:'Mn', 26:'Fe', 27:'Co', 28:'Ni', 29:'Cu', 30:'Zn', 31:'Ga', 32:'Ge', 33:'As', 34:'Se', 35:'Br', 36:'Kr',
                         37:'Rb', 38:'Sr', 39:'Y' , 40:'Zr', 41:'Nb', 42:'Mo', 43:'Tc', 44:'Ru', 45:'Rh', 46:'Pd', 47:'Ag', 48:'Cd', 49:'In', 50:'Sn', 51:'Sb', 52:'Te', 53:'I' , 54:'Xe',
                         55:'Cs', 56:'Ba', 57:'La', 72:'Hf', 73:'Ta', 74:'W' , 75:'Re', 76:'Os', 77:'Ir', 78:'Pt', 79:'Au', 80:'Hg', 81:'Tl', 82:'Pb', 83:'Bi', 84:'Po', 85:'At', 86:'Rn'}

    # Initialize lists to hold unique adj_mats, type_lists, element_lists, and geometries
    mol_adj_mats = []
    mol_types    = []
    mol_elements = []
    mol_geos     = []
    mol_charges  = []

    # Iterate over all configurations
    for i in list(Data.keys()):

        # Find elements for mol_a and mol_b
        elem_a = [ atom2element_dict[int(j.split('[')[1].split(']')[0])] for j in Data[i]['types_a'] ]
        elem_b = [ atom2element_dict[int(j.split('[')[1].split(']')[0])] for j in Data[i]['types_b'] ]

        # Generate adjacency matrix for the geo_a and geo_b
        adj_mat_a = Table_generator(elem_a,Data[i]['geo_a'])            
        adj_mat_b = Table_generator(elem_b,Data[i]['geo_b'])            

        # if adj_mat_a isn't in mol_adj_mats then it is a new molecule
        if True not in [ array_equal(adj_mat_a,j) for j in  mol_adj_mats ]:
            mol_adj_mats += [adj_mat_a]
            mol_types    += [Data[i]['types_a']]
            mol_elements += [elem_a]
            mol_geos     += [Data[i]['geo_a']]
            mol_charges  += [Data[i]['charges_a']]

        # if adj_mat_a is in mol_adj_mats but there isn't a match for type lists then it is a new molecule
        else:
            match = 0
            for count_j,j in enumerate(mol_adj_mats):
                if array_equal(adj_mat_a, j) and Data[i]['types_a'] == mol_types[count_j]: match = 1
            if match == 0:
                mol_adj_mats += [adj_mat_a]
                mol_types    += [Data[i]['types_a']]
                mol_elements += [elem_a]
                mol_geos     += [Data[i]['geo_a']]
                mol_charges  += [Data[i]['charges_a']]

        # if adj_mat_b isn't in mol_adj_mats then it is a new molecule
        if True not in [ array_equal(adj_mat_b,j) for j in  mol_adj_mats ]:
            mol_adj_mats += [adj_mat_b]
            mol_types    += [Data[i]['types_b']]
            mol_elements += [elem_b]
            mol_geos     += [Data[i]['geo_b']]
            mol_charges  += [Data[i]['charges_b']]
        
        # if adj_mat_b is in mol_adj_mats but there isn't a match for type lists then it is a new molecule
        else:
            match = 0
            for count_j,j in enumerate(mol_adj_mats):
                if array_equal(adj_mat_b, j) and Data[i]['types_b'] == mol_types[count_j]: match = 1
            if match == 0:
                mol_adj_mats += [adj_mat_b]
                mol_types    += [Data[i]['types_b']]
                mol_elements += [elem_b]
                mol_geos     += [Data[i]['geo_b']]
                mol_charges  += [Data[i]['charges_b']]

    # Write examples of each molecule in the sim to a molecule.db file. This is used when
    # adding the VDW parameters to the master database. Each parameter in the master is linked
    # to the molecule(s) used to generate it. In the case of VDW params, one example configuration
    # for each molecule type is used. 
    with open(folder+'/molecules.db','w') as f:
        for count_i,i in enumerate(mol_geos):
            f.write("\nmol {:6d} start\n".format(count_i+1))
            # Save an xyz for viewing the configuration in VMD
            f.write('{}\n\n'.format(len(i)))
            for count_j,j in enumerate(i):
                f.write('  {:<10s} {:< 20.6f} {:< 20.6f} {:< 20.6f} {:< 12.6f} {:<60s}\n'.\
                format(mol_elements[count_i][count_j],j[0],j[1],j[2],mol_charges[count_i][count_j],mol_types[count_i][count_j]))
            f.write("mol {:6d} end\n".format(count_i+1))

    # Assemble the list of atomtypes whose molecules are minimal structures
    min_list = []
    for count_i,i in enumerate(mol_types):
        for count_j,j in enumerate(i):
            if minimal_structure(j,mol_geos[count_i],mol_elements[count_i],adj_mat=mol_adj_mats[count_i],atomtypes=mol_types[count_i],gens=2) is True:
                min_list += [j]

    # Assemble all permutations of minimal types                
    min_list = set(min_list)
    min_pairs = []
    for i in min_list:
        for j in min_list:
            min_pairs += [(i,j),(i+'-UA',j+'-UA')]            
    min_pairs = set(min_pairs)

    # Print diagnostic and return
    print("{:70s} {}".format("Number of unique molecules in the fit configurations:",len(mol_types)))
    return

# Description:   Checks is the supplied geometry corresponds to the minimal structure of the molecule
# 
# Inputs:        atomtype:      The taffi atomtype being checked
#                geo:           Geometry of the molecule
#                elements:      elements, indexed to the geometry 
#                adj_mat:       adj_mat, indexed to the geometry (optional)
#                atomtypes:     atomtypes, indexed to the geometry (optional)
#                gens:          number of generations for determining atomtypes (optional, only used if atomtypes are not supplied)
# 
# Outputs:       Boolean:       True if geo is the minimal structure for the atomtype, False if not.
def minimal_structure(atomtype,geo,elements,adj_mat=None,atomtypes=None,gens=2):

    # If required find the atomtypes for the geometry
    if atomtypes is None or adj_mat is None:
        if len(elements) != len(geo):
            print("ERROR in minimal_structure: While trying to automatically assign atomtypes, the elements argument must have dimensions equal to geo. Exiting...")
            quit()

        # Generate the adjacency matrix
        # NOTE: the units are converted back angstroms
        adj_mat = Table_generator(elements,geo)

        # Generate the atomtypes
        atomtypes = id_types(elements,adj_mat,gens,Hybridization_finder(elements,adj_mat),geo)
        
    # Check minimal conditions
    count = 0
    for count_i,i in enumerate(atomtypes):

        # If the current atomtype matches the atomtype being searched for then proceed with minimal geo check
        if i == atomtype:
            count += 1

            # Initialize lists for holding indices in the structure within "gens" bonds of the seed atom (count_i)
            keep_list = [count_i]
            new_list  = [count_i]
            
            # Carry ount a "gens" bond deep search
            for j in range(gens):

                # Id atoms in the next generation
                tmp_new_list = []                
                for k in new_list:
                    tmp_new_list += [ count_m for count_m,m in enumerate(adj_mat[k]) if m == 1 and count_m not in keep_list ]

                # Update lists
                tmp_new_list = list(set(tmp_new_list))
                if len(tmp_new_list) > 0:
                    keep_list += tmp_new_list
                new_list = tmp_new_list
            
            # Check for the minimal condition
            keep_list = set(keep_list)
            if False in [ elements[j] == "H" for j in range(len(elements)) if j not in keep_list ]:
                minimal_flag = False
            else:
                minimal_flag = True
        
    return minimal_flag

# Description: Reads in the FF parameters
# NOTE: db_files
def get_FF_data(db_files,keep_types=[ "atom", "vdw", "bond", "angle", "torsion", "dihedral", "charge" ],mixing_rule=None):

    Data = {"atoms":{},"bonds":{},"angles":{},"dihedrals":{},"vdws":{},"charges":{}}
    
    # Read in Atom data
    for i in db_files:
        with open(i,'r') as f:
            for lines in f:
                fields = lines.split()

                if len(fields) == 0: continue
                if fields[0] == "atom" and "atom" in keep_types:
                    Data["atoms"][fields[1]] = fields[1:]
                if fields[0] == "vdw" and "vdw" in keep_types:
                    Data["vdws"][(fields[1],fields[2],fields[3])] = fields[1:]
                if fields[0] == "bond" and "bond" in keep_types:
                    Data["bonds"][(fields[1],fields[2],fields[3])] = fields[1:]
                if fields[0] == "angle" and "angle" in keep_types:
                    Data["angles"][(fields[1],fields[2],fields[3],fields[4])] = fields[1:]
                if fields[0] == "diehdral" or fields[0] == "torsion" and ( "diehdral" in keep_types or "torsion" in keep_types ):
                    Data["dihedrals"][(fields[1],fields[2],fields[3],fields[4],fields[5])] = fields[1:]
                if fields[0] == "charge" and "charge" in keep_types:
                    Data["charges"][fields[1]] = fields[1:]
  
    # Generate the cross terms via the mixing rule if specified
    if mixing_rule == "lb":

        # AA-Types
        self_types = [ i[0] for i in Data["vdws"] if i[0] == i[1] and i[2] == 'lj' and "-UA" not in i[0] ]
        for count_i,i in enumerate(self_types):
            for count_j,j in enumerate(self_types):
                if count_j > count_i:
                    eps   = ( float(Data["vdws"][(i,i,'lj')][3]) * float(Data["vdws"][(j,j,'lj')][3]) )**(0.5)
                    sigma = ( float(Data["vdws"][(i,i,'lj')][4]) + float(Data["vdws"][(j,j,'lj')][4]) ) / 2.0
                    Data["vdws"][(i,j,'lj')] = [i,j,'lj', eps, sigma]
                    Data["vdws"][(j,i,'lj')] = [j,i,'lj', eps, sigma]

        # UA-Types
        self_types = [ i[0] for i in Data["vdws"] if i[0] == i[1] and i[2] == 'lj' and "-UA" in i[0] ]
        for count_i,i in enumerate(self_types):
            for count_j,j in enumerate(self_types):
                if count_j > count_i:
                    eps   = ( float(Data["vdws"][(i,i,'lj')][3]) * float(Data["vdws"][(j,j,'lj')][3]) )**(0.5)
                    sigma = ( float(Data["vdws"][(i,i,'lj')][4]) + float(Data["vdws"][(j,j,'lj')][4]) ) / 2.0
                    Data["vdws"][(i,j,'lj')] = [i,j,'lj', eps, sigma]
                    Data["vdws"][(j,i,'lj')] = [j,i,'lj', eps, sigma]

    if mixing_rule == "wh":

        # AA-Types
        self_types = [ i[0] for i in Data["vdws"] if i[0] == i[1] and i[2] == 'lj' and "-UA" not in i[0] ]
        for count_i,i in enumerate(self_types):
            for count_j,j in enumerate(self_types):
                if count_j > count_i:
                    sigma  = (( float(Data["vdws"][(i,i,'lj')][4])**(6.0) + float(Data["vdws"][(j,j,'lj')][4])**(6.0))/2.0)**(1.0/6.0)
                    eps    = (float(Data["vdws"][(i,i,'lj')][3])*float(Data["vdws"][(i,i,'lj')][4])**(6.0) * float(Data["vdws"][(j,j,'lj')][3])*float(Data["vdws"][(j,j,'lj')][4])**(6.0) )**(0.5) / sigma**(6.0)
                    Data["vdws"][(i,j,'lj')] = [i,j,'lj', eps, sigma]
                    Data["vdws"][(j,i,'lj')] = [j,i,'lj', eps, sigma]

        # UA-Types
        self_types = [ i[0] for i in Data["vdws"] if i[0] == i[1] and i[2] == 'lj' and "-UA" in i[0] ]
        for count_i,i in enumerate(self_types):
            for count_j,j in enumerate(self_types):
                if count_j > count_i:
                    sigma  = (( float(Data["vdws"][(i,i,'lj')][4])**(6.0) + float(Data["vdws"][(j,j,'lj')][4])**(6.0))/2.0)**(1.0/6.0)
                    eps    = (float(Data["vdws"][(i,i,'lj')][3])*float(Data["vdws"][(i,i,'lj')][4])**(6.0) * float(Data["vdws"][(j,j,'lj')][3])*float(Data["vdws"][(j,j,'lj')][4])**(6.0) )**(0.5) / sigma**(6.0)
                    Data["vdws"][(i,j,'lj')] = [i,j,'lj', eps, sigma]
                    Data["vdws"][(j,i,'lj')] = [j,i,'lj', eps, sigma]
      
    return Data

def find_jobtype(name):

    job_type = []
    with open(name,'r') as f:
        for i in f:
            fields = i.split()
            if len(fields) >= 4 and fields[3] == "B3LYP" and 'dft' not in job_type:
                job_type += ['dft']
            if len(fields) >= 4 and fields[3] == "RI-MP2" and 'mp2' not in job_type:
                job_type += ['mp2']
    return job_type

# Fit all AA parameters at once. This function calls a specialized LJ fit function 
# called E_LJ_Fit_all to simultaneously fit all LJ sigmas and epsilons. The iterative
# approach has some initial guess dependence that this approach avoids. 
def fit_all(x_vals,y_vals,Fit_Pairs,fit,delta_xhi2_thresh):
    
    global VDW_dict,fit_type,Pair_min_dict,Current_Fit_Pairs

    # Copy the supplied Fit_Pairs to the local/global variable Current_Fit_Pairs
    # This is performed because the pairs being fit need to be flexibly defined locally
    Current_Fit_Pairs = Fit_Pairs

    # Prepare initial guess parameters, initial fit energies, and initial xhi2 value
    initial_guess = [ j for i in Current_Fit_Pairs for j in VDW_dict[i] ]
    fit_vals = [ E_LJ_Fit_all(x_vals,*initial_guess) ]
    xhi2_previous = mean((y_vals-fit_vals[0])**2)

    # Print diagnostics
    print("\nInitializing simultaneous fit of all LJ parameters:\n")
    print("\t{:<30s} {}".format("delta_xhi2_thresh:",delta_xhi2_thresh))
    print("\t{:<30s} {}".format("Fit parameters:",len(Fit_Pairs)))
    print("\t{:<30s} {}\n".format("Initial xhi2:",xhi2_previous))

    # Perform parameter fit until self-consistency is acheived
    cycle = 1
    while 1:

        # Perform fit
        params,err = curve_fit(E_LJ_Fit_all,x_vals,y_vals,p0=initial_guess,maxfev=10000)

        # Calculate fit values and xhi2 with final fit parameters
        fit_vals += [E_LJ_Fit_all(x_vals,*params)]
        xhi2 = mean((y_vals-fit_vals[cycle])**2)        
        delta_xhi2 = xhi2_previous - xhi2 

        # Convert fit parameters back into eps,sigma tuples
        initial_guess = params
        params = [ (params[2*count_i],params[2*count_i+1]) for count_i,i in enumerate(Current_Fit_Pairs) ]

        # Print diagnostic and assign fit parameters to the VDW_dict
        print("\nCycle {} parameter shifts:".format(cycle))
        for count_i,i in enumerate(Current_Fit_Pairs):
            print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<10.4f}".format(j) for j in params[count_i]])))
            VDW_dict[i] = tuple(params[count_i])  
            VDW_dict[(i[1],i[0])] = tuple(params[count_i])        

        # Print change in xhi2 and update xhi2_previous
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",xhi2))
        print("\t{:<12s} {} (kcal/mol)^2\n".format("delta_xhi2:",delta_xhi2))
        xhi2_previous = copy(xhi2)
        
        # Check break conditions
        if delta_xhi2 < delta_xhi2_thresh:
            print("Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2), breaking fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
            break
        else:
            cycle += 1

    return fit_vals

# Fit all UA parameters at once. This function calls a specialized LJ fit function 
# called E_LJ_Fit_all_UA to simultaneously fit all LJ sigmas. The iterative
# approach has some initial guess dependence that this approach avoids. The main distinction
# for fitting the UA parameters is that the sigmas are held fixed based on a geometric 
# conversion from the AA sigmas, and only the epsilon values are fit. 
def fit_all_UA(x_vals,y_vals,Fit_Pairs,fit,delta_xhi2_thresh):
    
    global VDW_dict,fit_type,Pair_min_dict,Current_Fit_Pairs

    # Copy the supplied Fit_Pairs to the local/global variable Current_Fit_Pairs
    # This is performed because the pairs being fit need to be flexibly defined locally
    Current_Fit_Pairs = Fit_Pairs

    # Prepare initial guess parameters, initial fit energies, and initial xhi2 value
    initial_guess = [ VDW_dict[i][0] for i in Current_Fit_Pairs ]
    fit_vals = [ E_LJ_Fit_all_UA(x_vals,*initial_guess) ]
    xhi2_previous = mean((y_vals-fit_vals[0])**2)

    # Print diagnostics
    print("\nInitializing simultaneous fit of all UA-epsilon values:\n")
    print("\t{:<30s} {}".format("Initial xhi2:",xhi2_previous))
    print("\t{:<30s} {}\n".format("Fit parameters:",len(Fit_Pairs)))
        
    # Perform parameter fit until self-consistency is acheived
    cycle = 1
    while 1:

        # Perform fit
        params,err = curve_fit(E_LJ_Fit_all_UA,x_vals,y_vals,p0=initial_guess,maxfev=10000)

        # Calculate fit values and xhi2 with final fit parameters
        fit_vals += [E_LJ_Fit_all_UA(x_vals,*params)]
        xhi2 = mean((y_vals-fit_vals[cycle])**2)        
        delta_xhi2 = xhi2_previous - xhi2 

        # Convert fit parameters back into eps,sigma tuples
        initial_guess = params
        params = [ (params[count_i],VDW_dict[i][1]) for count_i,i in enumerate(Current_Fit_Pairs) ]

        # Print diagnostic and assign fit parameters to the VDW_dict
        print("\nCycle {} parameter shifts:\n".format(cycle))
        for count_i,i in enumerate(Current_Fit_Pairs):
            print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<10.4f}".format(j) for j in params[count_i]])))
            VDW_dict[i] = tuple(params[count_i])  
            VDW_dict[(i[1],i[0])] = tuple(params[count_i])        

        # Print change in xhi2 and update xhi2_previous
        print("\n\t{:<12s} {} (kcal/mol)^2".format("xhi2:",xhi2))
        print("\t{:<12s} {} (kcal/mol)^2\n".format("delta_xhi2:",delta_xhi2))
        xhi2_previous = copy(xhi2)
        
        # Check break conditions
        if delta_xhi2 < delta_xhi2_thresh:
            print("Delta Xhi2 convergence achieved ({} kcal^2/mol^2 < {} kcal^2/mol^2), breaking fit loop...\n".format(delta_xhi2,delta_xhi2_thresh))
            break
        else:
            cycle += 1

    return fit_vals

# Simple function for ordinary 1-9 < 10-19 < etc sorting
def natural_sort(l): 
    convert = lambda text: int(text) if text.isdigit() else text.lower() 
    alphanum_key = lambda key: [ convert(c) for c in re.split('([0-9]+)', key) ] 
    return sorted(l, key = alphanum_key)

# fit_type is a globally defined tuple
def E_LJ_Fit_all(ind,e_0=0, s_0=0, e_1=0, s_1=0, e_2=0, s_2=0, e_3=0, s_3=0, e_4=0, s_4=0, e_5=0, s_5=0, e_6=0, s_6=0, e_7=0, s_7=0, e_8=0, s_8=0, e_9=0, s_9=0,e_10=0,s_10=0,\
                         e_11=0,s_11=0,e_12=0,s_12=0,e_13=0,s_13=0,e_14=0,s_14=0,e_15=0,s_15=0,e_16=0,s_16=0,e_17=0,s_17=0,e_18=0,s_18=0,e_19=0,s_19=0,e_20=0,s_20=0,\
                         e_21=0,s_21=0,e_22=0,s_22=0,e_23=0,s_23=0,e_24=0,s_24=0,e_25=0,s_25=0,e_26=0,s_26=0,e_27=0,s_27=0,e_28=0,s_28=0,e_29=0,s_29=0,e_30=0,s_30=0,\
                         e_31=0,s_31=0,e_32=0,s_32=0,e_33=0,s_33=0,e_34=0,s_34=0,e_35=0,s_35=0,e_36=0,s_36=0,e_37=0,s_37=0,e_38=0,s_38=0,e_39=0,s_39=0,e_40=0,s_40=0,\
                         e_41=0,s_41=0,e_42=0,s_42=0,e_43=0,s_43=0,e_44=0,s_44=0,e_45=0,s_45=0,e_46=0,s_46=0,e_47=0,s_47=0,e_48=0,s_48=0,e_49=0,s_49=0,e_50=0,s_50=0,\
                         e_51=0,s_51=0,e_52=0,s_52=0,e_53=0,s_53=0,e_54=0,s_54=0,e_55=0,s_55=0,e_56=0,s_56=0,e_57=0,s_57=0,e_58=0,s_58=0,e_59=0,s_59=0,e_60=0,s_60=0,\
                         e_61=0,s_61=0,e_62=0,s_62=0,e_63=0,s_63=0,e_64=0,s_64=0,e_65=0,s_65=0,e_66=0,s_66=0,e_67=0,s_67=0,e_68=0,s_68=0,e_69=0,s_69=0,e_70=0,s_70=0,\
                         e_71=0,s_71=0,e_72=0,s_72=0,e_73=0,s_73=0,e_74=0,s_74=0,e_75=0,s_75=0,e_76=0,s_76=0,e_77=0,s_77=0,e_78=0,s_78=0,e_79=0,s_79=0,e_80=0,s_80=0,\
                         e_81=0,s_81=0,e_82=0,s_82=0,e_83=0,s_83=0,e_84=0,s_84=0,e_85=0,s_85=0,e_86=0,s_86=0,e_87=0,s_87=0,e_88=0,s_88=0,e_89=0,s_89=0,e_90=0,s_90=0,\
                         e_91=0,s_91=0,e_92=0,s_92=0,e_93=0,s_93=0,e_94=0,s_94=0,e_95=0,s_95=0,e_96=0,s_96=0,e_97=0,s_97=0,e_98=0,s_98=0,e_99=0,s_99=0,e_100=0,s_100=0):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,Current_Fit_Pairs,Data

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()
    
    # Find defined e and s variables. *vars holds the variable names *vals holds the variable values
    e_vars = natural_sort([ i for i in local_vars if 'e_' in i and "__" not in i and local_vars[i] != 0 ])
    e_vals = [ locals()[i] for i in e_vars ]
    s_vars = natural_sort([ i for i in local_vars if 's_' in i and "__" not in i and local_vars[i] != 0 ])
    s_vals = [ locals()[i] for i in s_vars ]

    # Impose constraints on epsilon values
    for count_i,i in enumerate(e_vars):
        if e_vals[count_i] < 0.001 or e_vals[count_i] > 30.0: return ones(len(ind))*1.E20

    # Impose constraints on sigma values
    for count_i,i in enumerate(s_vars):
        if s_vals[count_i] > 10.0 or s_vals[count_i] < Pair_min_dict[Current_Fit_Pairs[count_i]]*2.**(-1./6.): return ones(len(ind))*1.E20

    # Initialize eps and sigma arrays for the vectorized calculations
    eps_array = array([ VDW_dict[i][0] for i in Data[ind[0]]["pair_type_vector"] ])
    sigma_array = array([ VDW_dict[i][1] for i in Data[ind[0]]["pair_type_vector"] ])
    for count_c,c in enumerate(Current_Fit_Pairs):
        for i in [ count_j for count_j,j in enumerate(Data[ind[0]]["pair_type_vector"]) if ( j == c or (j[1],j[0]) == c ) ]:        
            eps_array[i]   = e_vals[count_c]
            sigma_array[i] = s_vals[count_c]

    # Initialize E_LJ for the current configuration
    E_LJ = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        E_LJ[count_d] = sum( ( 4.0*eps_array*sigma_array**(12.0)*Data[d]["pair_vector-12"] - 4.0*eps_array*sigma_array**(6.0)*Data[d]["pair_vector-6"] ) )

    return E_LJ

# fit_type is a globally defined tuple
def E_BUCK_Fit_all(ind,a_0=0, b_0=0, c_0=0, a_1=0, b_1=0, c_1=0, a_2=0, b_2=0, c_2=0, a_3=0, b_3=0, c_3=0, a_4=0, b_4=0, c_4=0, a_5=0, b_5=0, c_5=0, a_6=0, b_6=0, c_6=0, a_7=0, b_7=0, c_7=0,\
                       a_8=0, b_8=0, c_8=0, a_9=0, b_9=0, c_9=0, a_10=0,b_10=0,c_10=0,a_11=0,b_11=0,c_11=0,a_12=0,b_12=0,c_12=0,a_13=0,b_13=0,c_13=0,a_14=0,b_14=0,c_14=0,a_15=0,b_15=0,c_15=0,\
                       a_16=0,b_16=0,c_16=0,a_17=0,b_17=0,c_17=0,a_18=0,b_18=0,c_18=0,a_19=0,b_19=0,c_19=0,a_20=0,b_20=0,c_20=0,a_21=0,b_21=0,c_21=0,a_22=0,b_22=0,c_22=0,a_23=0,b_23=0,c_23=0,\
                       a_24=0,b_24=0,c_24=0,a_25=0,b_25=0,c_25=0,a_26=0,b_26=0,c_26=0,a_27=0,b_27=0,c_27=0,a_28=0,b_28=0,c_28=0,a_29=0,b_29=0,c_29=0,a_30=0,b_30=0,c_30=0,\
                       a_31=0,b_31=0,c_31=0,a_32=0,b_32=0,c_32=0,a_33=0,b_33=0,c_33=0,a_34=0,b_34=0,c_34=0,a_35=0,b_35=0,c_35=0,a_36=0,b_36=0,c_36=0,a_37=0,b_37=0,c_37=0,\
                       a_38=0,b_38=0,c_38=0,a_39=0,b_39=0,c_39=0,a_40=0,b_40=0,c_40=0,a_41=0,b_41=0,c_41=0,a_42=0,b_42=0,c_42=0,a_43=0,b_43=0,c_43=0,a_44=0,b_44=0,c_44=0,\
                       a_45=0,b_45=0,c_45=0,a_46=0,b_46=0,c_46=0,a_47=0,b_47=0,c_47=0,a_48=0,b_48=0,c_48=0,a_49=0,b_49=0,c_49=0,a_50=0,b_50=0,c_50=0):

    # In python you need to expand the scope of global variables at every level
    global VDW_dict,Current_Fit_Pairs,Data

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()
    
    # Find defined e and s variables. *vars holds the variable names *vals holds the variable values
    a_vars = natural_sort([ i for i in local_vars if 'a_' in i and "__" not in i and local_vars[i] != 0 ])
    a_vals = [ locals()[i] for i in a_vars ]
    b_vars = natural_sort([ i for i in local_vars if 'b_' in i and "__" not in i and local_vars[i] != 0 ])
    b_vals = [ locals()[i] for i in b_vars ]
    c_vars = natural_sort([ i for i in local_vars if 'c_' in i and "__" not in i and local_vars[i] != 0 ])
    c_vals = [ locals()[i] for i in c_vars ]


    # Impose constraints on values
    for count_i,i in enumerate(a_vars):
        if a_vals[count_i] < 0.001 or b_vals[count_i] < 0.001 or c_vals[count_i] < 0.001\
        or True not in [ a_vals[count_i]*exp(-j/b_vals[count_i]) - c_vals[count_i]*j**(-6.0) > 10.0 for j in arange(0.1,5.1,0.01) ]:
            return ones(len(ind))*1.E20

    # Initialize E_BUCK for the current configuration
    E_BUCK = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            ran_flag = 0
            for count_j,j in enumerate(Current_Fit_Pairs):
                if i == j or (i[1],i[0]) == j:
                    E_BUCK[count_d] += sum( ( a_vals[count_j]*exp(-Data[d]["pairs"][i]/b_vals[count_j]) - c_vals[count_j]*Data[d]["pairs"][i]**(-6.0) ) )
                    ran_flag = 1
                    break
            if ran_flag == 0:
                E_BUCK[count_d] += sum( ( VDW_dict[i][0]*exp(-Data[d]["pairs"][i]/VDW_dict[i][1]) - VDW_dict[i][2]*Data[d]["pairs"][i]**(-6.0) ) )

    return E_BUCK

# fit_type is a globally defined tuple
def E_LJ_Fit_all_UA(ind,e_0=0, s_0=0, e_1=0, s_1=0, e_2=0, s_2=0, e_3=0, s_3=0, e_4=0, s_4=0, e_5=0, s_5=0, e_6=0, s_6=0, e_7=0, s_7=0, e_8=0, s_8=0, e_9=0, s_9=0,e_10=0,s_10=0,\
                        e_11=0,s_11=0,e_12=0,s_12=0,e_13=0,s_13=0,e_14=0,s_14=0,e_15=0,s_15=0,e_16=0,s_16=0,e_17=0,s_17=0,e_18=0,s_18=0,e_19=0,s_19=0,e_20=0,s_20=0,\
                        e_21=0,s_21=0,e_22=0,s_22=0,e_23=0,s_23=0,e_24=0,s_24=0,e_25=0,s_25=0,e_26=0,s_26=0,e_27=0,s_27=0,e_28=0,s_28=0,e_29=0,s_29=0,e_30=0,s_30=0,\
                        e_31=0,s_31=0,e_32=0,s_32=0,e_33=0,s_33=0,e_34=0,s_34=0,e_35=0,s_35=0,e_36=0,s_36=0,e_37=0,s_37=0,e_38=0,s_38=0,e_39=0,s_39=0,e_40=0,s_40=0,\
                        e_41=0,s_41=0,e_42=0,s_42=0,e_43=0,s_43=0,e_44=0,s_44=0,e_45=0,s_45=0,e_46=0,s_46=0,e_47=0,s_47=0,e_48=0,s_48=0,e_49=0,s_49=0,e_50=0,s_50=0,\
                        e_51=0,s_51=0,e_52=0,s_52=0,e_53=0,s_53=0,e_54=0,s_54=0,e_55=0,s_55=0,e_56=0,s_56=0,e_57=0,s_57=0,e_58=0,s_58=0,e_59=0,s_59=0,e_60=0,s_60=0,\
                        e_61=0,s_61=0,e_62=0,s_62=0,e_63=0,s_63=0,e_64=0,s_64=0,e_65=0,s_65=0,e_66=0,s_66=0,e_67=0,s_67=0,e_68=0,s_68=0,e_69=0,s_69=0,e_70=0,s_70=0,\
                        e_71=0,s_71=0,e_72=0,s_72=0,e_73=0,s_73=0,e_74=0,s_74=0,e_75=0,s_75=0,e_76=0,s_76=0,e_77=0,s_77=0,e_78=0,s_78=0,e_79=0,s_79=0,e_80=0,s_80=0,\
                        e_81=0,s_81=0,e_82=0,s_82=0,e_83=0,s_83=0,e_84=0,s_84=0,e_85=0,s_85=0,e_86=0,s_86=0,e_87=0,s_87=0,e_88=0,s_88=0,e_89=0,s_89=0,e_90=0,s_90=0,\
                        e_91=0,s_91=0,e_92=0,s_92=0,e_93=0,s_93=0,e_94=0,s_94=0,e_95=0,s_95=0,e_96=0,s_96=0,e_97=0,s_97=0,e_98=0,s_98=0,e_99=0,s_99=0,e_100=0,s_100=0):
        
    # In python you need to expand the scope of global variables at every level
    global VDW_dict,Current_Fit_Pairs,Data

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()
    
    # Find defined e and s variables. *vars holds the variable names *vals holds the variable values
    e_vars = natural_sort([ i for i in local_vars if 'e_' in i and "__" not in i and local_vars[i] != 0 ])
    e_vals = [ locals()[i] for i in e_vars ]

    # Impose constraints on epsilon values
    for count_i,i in enumerate(e_vars):
        if e_vals[count_i] < 0.001 or e_vals[count_i] > 30.0: return ones(len(ind))*1.E20

    # Initialize E_LJ for the current configuration
    E_LJ = zeros(len(ind))

    # Iterate over the configuration index
    for count_d,d in enumerate(ind):

        # Iterate over list of pair-seps for each pair-type (Data[d]["pairs"] is a dictionary with keys for each pair type linked to arrays of the instances of pair separations for that type)
        for i in list(Data[d]["pairs"].keys()):
            ran_flag = 0
            for count_j,j in enumerate(Current_Fit_Pairs):
                if i == j or (i[1],i[0]) == j:
                    E_LJ[count_d] += sum( ( 4.0*e_vals[count_j]*VDW_dict[j][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*e_vals[count_j]*VDW_dict[j][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )
                    ran_flag = 1
                    break
            if ran_flag == 0:
                E_LJ[count_d] += sum( ( 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(12.0)*Data[d]["pairs"][i]**(-12.0) - 4.0*VDW_dict[i][0]*VDW_dict[i][1]**(6.0)*Data[d]["pairs"][i]**(-6.0) ) )

    return E_LJ

# Calculate the effective sigma value for united atoms based on the volume of the 
# bonded atoms (sum of invidivual hard sphere volumes minus their overlap)
def calc_eff_sigma(seps,sigma_1=0,sigma_2=0,sigma_3=0,sigma_4=0,sigma_5=0,sigma_6=0,sigma_7=0,sigma_8=0,sigma_9=0,sigma_10=0):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()
    
    # Find defined sigmas
    sigma_vars = natural_sort([ i for i in local_vars if 'sigma_' in i and "__" not in i and local_vars[i] != 0 ])
    sigma_vals = [ locals()[i] for i in sigma_vars ]

    V_eff = 0.0

    # Add individual volumes to V_eff
    for i in sigma_vals:
        V_eff += 4.0/3.0*pi*i**(3.0)

    # Subtract out overlapping regions from V_eff
    for count_i,i in enumerate(sigma_vals):
        for count_j,j in enumerate(sigma_vals):

            # Avoid double counting pairs
            if count_j <= count_i: continue

            # If there is no overlap between the spheres subract nothing
            if seps[count_i,count_j] > i+j: continue
            
            # If one sphere completely contains the other then subtract off the smaller sphere volume
            if seps[count_i,count_j] < abs(i-j): 
                if i <= j:
                    V_eff = V_eff - 4.0/3.0*pi*i**(3.0)
                else:
                    V_eff = V_eff - 4.0/3.0*pi*j**(3.0)
                continue

            # Subtract the overlap
            else:
                V_eff = V_eff - ( pi * (sigma_1+sigma_2-seps[count_i,count_j])**(2.0) *\
                (seps[count_i,count_j]**(2.0) + 2.0*seps[count_i,count_j]*i + 2.0*seps[count_i,count_j]*j - 3.0*i**(2.0) - 3.0*j**(2.0) + 6.0*i*j ) ) / \
                (12.0*seps[count_i,count_j])

    # Calculate the effective sigma value based on the effective volume (V_1+V_2-V_overlap)
    sigma_eff = ( 3.0 / (4.0*pi) * V_eff )**(1.0/3.0)

    return sigma_eff

# Calculate the effective sigma value for united atoms based on the volume of the 
# bonded atoms (sum of invidivual hard sphere volumes minus their overlap)
def calc_eff_vol(seps,sigma_1=0,sigma_2=0,sigma_3=0,sigma_4=0,sigma_5=0,sigma_6=0,sigma_7=0,sigma_8=0,sigma_9=0,sigma_10=0):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()
    
    # Find defined sigmas
    sigma_vars = natural_sort([ i for i in local_vars if 'sigma_' in i and "__" not in i and local_vars[i] != 0 ])
    sigma_vals = [ locals()[i] for i in sigma_vars ]

    V_eff = 0.0

    # Add individual volumes to V_eff
    for i in sigma_vals:
        V_eff += 4.0/3.0*pi*i**(3.0)

    # Subtract out overlapping regions from V_eff
    for count_i,i in enumerate(sigma_vals):
        for count_j,j in enumerate(sigma_vals):

            # Avoid double counting pairs
            if count_j <= count_i: continue

            # If there is no overlap between the spheres subract nothing
            if seps[count_i,count_j] > i+j: continue
            
            # If one sphere completely contains the other then subtract off the smaller sphere volume
            if seps[count_i,count_j] < abs(i-j): 
                if i <= j:
                    V_eff = V_eff - 4.0/3.0*pi*i**(3.0)
                else:
                    V_eff = V_eff - 4.0/3.0*pi*j**(3.0)
                continue

            # Subtract the overlap
            else:
                V_eff = V_eff - ( pi * (i+j-seps[count_i,count_j])**(2.0) *\
                (seps[count_i,count_j]**(2.0) + 2.0*seps[count_i,count_j]*i + 2.0*seps[count_i,count_j]*j - 3.0*i**(2.0) - 3.0*j**(2.0) + 6.0*i*j ) ) / \
                (12.0*seps[count_i,count_j])

    return V_eff

# this function converts the AA-sigma values to UA values on the basis of effective volume
# and the AA-eps values on the basis of interchain LJ-interaction energies.
def calc_eff_UA(folder,VDW_dict,DB_Pairs,Data):

    # In python you need to expand the scope of global variables at every level
    global Pair_min_dict

    # Print diagnostic
    print("\nCalculating all-atom interaction potentials, averaged over all configurations...")
        
    # Collect types that will be converted into UA from all configurations
    UA_types = {}
    UA_inds_A  = {}
    UA_inds_B  = {}
    removed_H_types = []
    keys = list(Data.keys())

    # Iterate over the configurations and collect the UA_types and their corresponding atom indices in each configuration
    for count_i,i in enumerate(keys):
        
        # Iterate over the molecule a adj_mat and identify carbons with attached hydrogens.
        for count_j,j in enumerate(Data[i]["adj_mat_a"]):
            
            # If the current type is a carbon then add its type along with any attached hydrogens
            if int(Data[i]["types_a"][count_j].split('[')[1].split(']')[0]) == 6:

                # Grab the UA_type and indices as tuples (sort by type to make sure the UA_types are unique and independent of the basis ordering of the adj_mat)
                current = [ (Data[i]["types_a"][count_j],count_j) ] + \
                          sorted([ (Data[i]["types_a"][count_k],count_k) for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_a'][count_k].split('[')[1].split(']')[0]) == 1 ])

                # Add attached H_types to UA_H_list
                if len(current) > 1:
                    for k in current[1:]:
                        if k[0] not in removed_H_types: removed_H_types += [k[0]]

            # Else, simply grab the type and its index
            else:
                current = [ (Data[i]["types_a"][count_j],count_j) ]

            # Separate the types from the atom indices
            current_type = tuple([ k[0] for k in current ])
            current_inds = [ k[1] for k in current ]

            # If the type doesn't exist in UA_inds_A, initialize the dictionary list for this type
            if current_type not in list(UA_inds_A.keys()):
                UA_inds_A[current_type] = [ [] for k in range(len(keys)) ]

            # Add the current indices to the current_type's list in UA_inds_A dictionary 
            UA_inds_A[current_type][count_i] += [ current_inds ]

        # Iterate over the molecule b adj_mat and identify carbons with attached hydrogens.
        for count_j,j in enumerate(Data[i]["adj_mat_b"]):

            # If the current type is a carbon then add its type along with any attached hydrogens
            if int(Data[i]["types_b"][count_j].split('[')[1].split(']')[0]) == 6: 

                # Grab the UA_type and indices as tuples (sort by type to make sure the UA_types are unique and independent of the basis ordering of the adj_mat)
                current = [ (Data[i]["types_b"][count_j],count_j) ] + \
                          sorted([ (Data[i]["types_b"][count_k],count_k) for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_b'][count_k].split('[')[1].split(']')[0]) == 1 ])

                # Add attached H_types to UA_H_list
                if len(current) > 1:
                    for k in current[1:]:
                        if k[0] not in removed_H_types: removed_H_types += [k[0]]

            # Else, simply grab the type and its index
            else:
                current = [ (Data[i]["types_b"][count_j],count_j) ]

            # Separate the types from the atom indices
            current_type = tuple([ k[0] for k in current ])
            current_inds = [ k[1] for k in current ]

            # If the type doesn't exist in UA_inds_B, initialize the dictionary list for this type
            if current_type not in list(UA_inds_B.keys()):
                UA_inds_B[current_type] = [ [] for k in range(len(keys)) ]

            # Add the current indices to the current_type's list in UA_inds_B dictionary 
            UA_inds_B[current_type][count_i] += [ current_inds ]

    # Add reverse combinations of keys to VDW_dict (they are removed
    # by the write_params function to avoid saving redundant pairs)
    VDW_dict_keys = list(VDW_dict.keys())
    for i in VDW_dict_keys:
        if (i[1],i[0]) not in list(VDW_dict.keys()): VDW_dict[(i[1],i[0])]=VDW_dict[i]

    # Initialize the new VDW dictionary
    # (avoids the complications associated with updates to VDW_dict while fits are being performed)
    VDW_new = {}
    for i in list(VDW_dict.keys()):        
        VDW_new[i] = VDW_dict[i]

    # Print diagnostic
    print("Calculating united-atom interaction potentials based on fits to all-atom potentials:\n")

    # Iterate over all pairs and calculate the average effective AA-potentials 
    # Fit the UA-potentials, save plots, and assign the approriate elements in VDW_dict
    global charge_a,charge_b,fit_type
    disp_start = 10.0
    disp_step  = 0.05
    seps = arange(disp_step,disp_start+disp_step,disp_step)
    assigned_types = []
    assigned_flag = 0
    for i in list(UA_inds_A.keys()):
        for j in list(UA_inds_B.keys()):            

            # Skip if fit_type is not in the fit list
            if (i[0],j[0]) not in DB_Pairs and (j[0],i[0]) not in DB_Pairs:
                continue

            # Avoid redundant calculations
            if (i,j) in assigned_types: continue                

            # Designate the pair type being fit/assigned.
            fit_type = (i[0],j[0])

            # If either atom is a UA-H then the epsilon value is set to zero
            if i[0] in removed_H_types or j[0] in removed_H_types:
                VDW_new[(i[0],j[0])] = (0.0,VDW_new[(i[0],j[0])][1])
                VDW_new[(j[0],i[0])] = (0.0,VDW_new[(j[0],i[0])][1])
                assigned_flag = 1

            # If neither atom is a UA_atom then the parameters are copied as is.
            elif len(i) <= 1 and len(j) <= 1 and assigned_flag == 0: 
                VDW_new[(i[0],j[0])] = VDW_new[(i[0],j[0])]
                VDW_new[(j[0],i[0])] = VDW_new[(j[0],i[0])]
                assigned_flag = 1

            # If one or both are UA-C types
            if (len(i) > 1 or len(j) > 1) and assigned_flag == 0:

                # If i is in both UA_inds_A and UA_inds_B, and if j is in both UA_inds_B and UA_inds_A,
                # then the potentials generated from each are combined.
                if i in list(UA_inds_B.keys()) and j in list(UA_inds_A.keys()) and i != j:
                    Potential_a,counts_a = calc_eff_pot(i,j,UA_inds_A,UA_inds_B,Data,keys,VDW_dict,fit='lj',disp_start=disp_start,disp_step=disp_step)
                    Potential_b,counts_b = calc_eff_pot(j,i,UA_inds_A,UA_inds_B,Data,keys,VDW_dict,fit='lj',disp_start=disp_start,disp_step=disp_step)
                    AA_Potential = Potential_a * float(counts_a)/float(counts_a+counts_b) + Potential_b * float(counts_b)/float(counts_a+counts_b)

                # Else, the potential is calculated directly from a single loop over types i in mol A and types j in mol B in all configurations
                else:
                    AA_Potential,counts = calc_eff_pot(i,j,UA_inds_A,UA_inds_B,Data,keys,VDW_dict,fit='lj',disp_start=disp_start,disp_step=disp_step)
                    
                # Find the UA_charges on types a and b
                for count_k,k in enumerate(keys):
                    if UA_inds_A[i][count_k] != []:
                        charge_a = Data[k]["charges_a_UA"][UA_inds_A[i][count_k][0][0]]
                        break
                for count_k,k in enumerate(keys):
                    if UA_inds_B[j][count_k] != []:
                        charge_b = Data[k]["charges_b_UA"][UA_inds_B[j][count_k][0][0]]
                        break

                # Find fit region of the AA_Potential
                # If there is not a minimum, then region starting at min_sep is used
                # Note: this sometimes happens for non-binding groups like CH3-CH3 interactions)
                if where(AA_Potential == AA_Potential.min())[0][0] == len(AA_Potential) - 1:
                    print("Using asymmtote")
                    min_array = abs(seps-Pair_min_dict[fit_type]-0.5)   # 0.25 factor is just an empirical offset to help sample a bit of the repulsive wall
                    fit_inds = list(range(where(min_array == min_array.min())[0][0],len(AA_Potential)))

                # If a minimum exists, then the portion of the binding curve with energy less than the asymtotic energy is 
                # used for fitting the UA-parameters
                else:
#                    min_ind = [ count_k for count_k,k in enumerate(AA_Potential) if k == min(AA_Potential) ][0]
#                    fit_inds = range(min_ind-0,len(AA_Potential))
                    fit_inds = [ count_k for count_k,k in enumerate(AA_Potential) if k <= AA_Potential[-1] ]
                    # For very shallow minima (i.e. well-depth < 0.25 kcal/mol, additional points are added to sample the repulsive shoulder
                    if AA_Potential[-1] - min(AA_Potential) < 0.25:
                        print("Increasing repulsive sampling...")
                        fit_inds += list(range(min(fit_inds)-5,min(fit_inds)))  # Add 5 datapoints to increase sampling of the repulsive shoulder

                # Some dependence on initial guess for epsilon has been observed. To avoid that, a loop over
                # a range of reasonable initial guesses for epsilon is performed and the best fit is kept
                xhi2_best = 0.0
                params_best = []
                for f in arange(0.001,VDW_dict[fit_type][0]+1.0,0.1):
                    
                    # Set current initial guess and shift sigma if necessary (if the AA-parameters are being read, sometimes round-off errors put the sigma values outside of the allowed boundaries) 
                    initial_guess = (f,VDW_dict[fit_type][1])
                    if initial_guess[1] < Pair_min_dict[fit_type]*2.**(-1./6.): initial_guess = (initial_guess[0],Pair_min_dict[fit_type]*2.**(-1./6.) )

#                 # Shift the values (if the AA-parameters are being read, sometimes round-off errors put the sigma/eps values outside of the allowed boundaries)            
#                 initial_guess = VDW_dict[fit_type]
#                 if initial_guess[0] < 0.002: initial_guess = (0.1,VDW_dict[fit_type][1])   # I've found that if the AA-eps value is non-binding then it can throw off the fit
#                 if initial_guess[0] < 0.001: initial_guess = (0.001,initial_guess[1])
#                 if initial_guess[1] < Pair_min_dict[fit_type]*2.**(-1./6.): initial_guess = (initial_guess[0],Pair_min_dict[fit_type]*2.**(-1./6.) )

                    # Calculate the fit parameters
                    params,err           = curve_fit(E_UA_Fit,seps[fit_inds],AA_Potential[fit_inds],p0=initial_guess,maxfev=1000000)
                    
                    # Calculate xhi2 for this fit
                    xhi2 = mean((E_UA_Fit(seps[fit_inds],*params) - AA_Potential[fit_inds])**(2.0))
                    
                    # Check for best-fit update
                    if params_best == [] or xhi2 < xhi2_best:
                        params_best = params
                        xhi2_best = xhi2

                # Update the parameters in the VDW dictionary
                VDW_new[(i[0],j[0])] = params_best
                VDW_new[(j[0],i[0])] = params_best

                # Calculate the resulting UA_potential
                UA_Potential = E_UA_Fit(seps,*params_best)

                # Plot a comparison of the AA and UA potentials
                plot_UA_fits(folder,"{}_{}_AA-UA_fit_pot.png".format(i[0],j[0]),seps,AA_Potential,UA_Potential,fit_inds,fit='lj',r_max=disp_start)

                # Toggle the assigned flag 
                assigned_flag = 1
    
            # Add the types that have been evaluated
            assigned_types += [(i,j)]
            assigned_types += [(j,i)]

            # Print diagnostic
            print("\tpair {:80}: {} ->     {}"\
            .format(str(fit_type),", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[fit_type]]),", ".join([ "{:<10.4f}".format(j) for j in [VDW_new[fit_type][0],VDW_new[fit_type][1]]])))
            assigned_flag = 0

    return VDW_new

def E_UA_Fit(seps,eps,sigma):

    # In python you need to expand the scope of global variables at every level
    global fit_type,Pair_min_dict,charge_a,charge_b

    # If the attempted value for the fit type violates the sanity checks on the fit range then the energy is set to an astronomical value
    # NOTE: the restriction on sigma_min is to stop the minimum of the fit potential from lieing outside of the sampled configuration space
    if fit_type in list(Pair_min_dict.keys()) and ( eps < 0.001 or eps > 30.0 or sigma > 10.0 or sigma < Pair_min_dict[fit_type]*2.**(-1./6.) ): return ones(len(seps))*1.E20        

    Potential = zeros(len(seps))
    for count_i,i in enumerate(seps):

        # Add LJ terms (Note: eps and sigma are defined in units of kcal/mol and angstroms, respectively)
        Potential[count_i] += 4.0 * eps * ( (sigma/i)**(12.0) - (sigma/i)**(6.0) )

        # Add Coulombic terms (NOTE: Ecoul_Const makes the conversion to kcal/mol)
#        Potential[count_i] += charge_a*charge_b/i * Ecoul_Const

    return Potential
        
def plot_UA_fits(folder,name,seps,y_vals_1,y_vals_2,fit_inds,fit='lj',r_max=10.0):

    # Plot all pair potentials
    # NOTE: write_params call cleans up VDW_dict.keys to remove redundant pair_types 
    # Initialize figure and generate plot
    fig = plt.figure(figsize=(6,4))
    ax = plt.subplot(111)
    ax.plot(seps,y_vals_1,color=(0.05,0.35,0.75),linestyle='-',linewidth=2.0)
    ax.plot(seps,y_vals_2,color=(0.05,0.75,0.35),linestyle='-',linewidth=2.0)
    if fit_inds != []:
        ax.plot(seps[fit_inds],y_vals_2[fit_inds],color=(0.75,0.05,0.35),linestyle='None',linewidth=0.0,marker='o',markersize=5,markeredgecolor=(0.75,0.05,0.35))

    # Set limits based on the largest range
    y_min = 1000.0
    if fit == 'lj':
        while (1):
            if y_min > min(y_vals_1):# or y_min > min(y_vals_2):
                y_min -= 5
            else:
                break
    if y_min >= 5: y_max = y_min+15
    else: y_max = 5
    ax.set_xlim([0,r_max])
    ax.set_ylim([y_min,y_max])

    # Set Labels and Save Name
    ax.set_ylabel("$\mathrm{E_{LJ} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')
    ax.set_xlabel("$r\mathrm{_{ij} \, (kcal \cdot mol^{-1})}$",fontsize=32,labelpad=10,fontweight='bold')

    # Format ticks and plot box
    ax.tick_params(axis='both', which='major',labelsize=12,pad=10,direction='out',width=2,length=4)
    ax.tick_params(axis='both', which='minor',labelsize=12,pad=10,direction='out',width=2,length=3)
    [j.set_linewidth(2) for j in ax.spines.values()]

    # Save the figure
    savefig(folder+'/'+name, dpi=300, bbox_inches='tight')
    close(fig)

def calc_eff_pot(A_type,B_type,UA_inds_A,UA_inds_B,Data,keys,VDW_dict,fit='lj',disp_start=10.0,disp_step=0.01):

    # NOTE: UA_inds_*[*_type] is a list of lists. Each element in the list corresponds to a frame in Data[keys]
    #       and each element in the sublists is an atom index corresponding to an atom in the UA_group being parametrized
    #       e.g. for UA_inds_A[type] = [ [ [1,2,3],[4,5,6],[7,8,9] ] [ [1,2,3],[4,5,6],[7,8,9] ] ], the type might be a CH2
    #       united-atom, and [1,2,3] the atom indices of one instance of that type in molecule A in the first configuration (keys[0]),
    #       [4,5,6] the atom indices of the second instance of that type in molecula A in the first configuration (keys[0]), etc.
    count = 0
    steps = int(disp_start/disp_step)
    Potential = zeros(steps)

    # Iterate over the frames
    for count_k,k in enumerate(keys):

        # Iterate over the groups in UA_inds_A[A_type][count_k]
        for i in UA_inds_A[A_type][count_k]:

            # Iterate over the groups in UA_inds_B[B_type][count_k]
            for j in UA_inds_B[B_type][count_k]:

                # grab geometries, types, and charges for each group of atoms
                geo_a = Data[k]['geo_a'][i]
                geo_b = Data[k]['geo_b'][j]
                types_a = [ Data[k]['types_a'][m] for m in i ]
                types_b = [ Data[k]['types_b'][m] for m in j ]
                charges_a = [ Data[k]['charges_a'][m] for m in i ]
                charges_b = [ Data[k]['charges_b'][m] for m in j ]

                # Increment the configuration counter
                count += 1
                
                # Define initial separation and the separation vector
                sep   = norm(geo_b[0,:]-geo_a[0,:])
                sep_v = (geo_b[0,:]-geo_a[0,:])/sep
                
                # separate the two groups to the starting displacement (disp_start)
                geo_b += ( sep_v * (disp_start-sep) )
                if fit == 'lj':

                    # Calculate the LJ+Coulombic interaction potential for the two groups of atoms at each separation
                    for m in range(steps):
                        for count_n,n in enumerate(geo_a):
                            for count_p,p in enumerate(geo_b):

                                # Calculate pair separation
                                r = norm(n-p)

                                # Add LJ terms (Note: eps and sigma are defined in units of kcal/mol and angstroms, respectively)
                                Potential[steps-m-1] += 4.0 * VDW_dict[(types_a[count_n],types_b[count_p])][0] * \
                                                      ( (VDW_dict[(types_a[count_n],types_b[count_p])][1]/r)**(12.0) - (VDW_dict[(types_a[count_n],types_b[count_p])][1]/r)**(6.0) )

                                # Add Coulombic terms (NOTE: Ecoul_Const makes the conversion to kcal/mol)
#                                Potential[steps-m-1] += charges_a[count_n]*charges_b[count_p]/r * Ecoul_Const
                        
                        # Displace group B (bring closer by disp_step)
                        geo_b -= ( sep_v * disp_step )

    Potential = Potential / float(count)
    return Potential,count

# this function converts the AA-sigma values to UA values on the basis of effective volume
# and the AA-eps values on the basis of interchain LJ-interaction energies.
def calc_UA_parameters(VDW_dict,DB_Pairs,Data,vol_scale,vol_method):

    # In python you need to expand the scope of global variables at every level
    global Pair_min_dict

    # Find bonds and average bond lengths
    bond_lengths = {}
    bond_counts = {}
    for i in list(Data.keys()):

        # Find the bonds in molecule a
        for count_j,j in enumerate(Data[i]["adj_mat_a"]):
            for count_k,k in enumerate(j):
                if k == 1:
                    if (Data[i]["types_a"][count_j],Data[i]["types_a"][count_k]) not in list(bond_lengths.keys()): bond_lengths[(Data[i]["types_a"][count_j],Data[i]["types_a"][count_k])] = 0.0
                    if (Data[i]["types_a"][count_j],Data[i]["types_a"][count_k]) not in list(bond_counts.keys()): bond_counts[(Data[i]["types_a"][count_j],Data[i]["types_a"][count_k])] = 0
                    bond_lengths[(Data[i]["types_a"][count_j],Data[i]["types_a"][count_k])] += norm(Data[i]["geo_a"][count_j]-Data[i]["geo_a"][count_k])
                    bond_counts[(Data[i]["types_a"][count_j],Data[i]["types_a"][count_k])] += 1

        # Find the bonds in molecule b
        for count_j,j in enumerate(Data[i]["adj_mat_b"]):
            for count_k,k in enumerate(j):
                if k == 1:
                    if (Data[i]["types_b"][count_j],Data[i]["types_b"][count_k]) not in list(bond_lengths.keys()): bond_lengths[(Data[i]["types_b"][count_j],Data[i]["types_b"][count_k])] = 0.0
                    if (Data[i]["types_b"][count_j],Data[i]["types_b"][count_k]) not in list(bond_counts.keys()): bond_counts[(Data[i]["types_b"][count_j],Data[i]["types_b"][count_k])] = 0
                    bond_lengths[(Data[i]["types_b"][count_j],Data[i]["types_b"][count_k])] += norm(Data[i]["geo_b"][count_j]-Data[i]["geo_b"][count_k])
                    bond_counts[(Data[i]["types_b"][count_j],Data[i]["types_b"][count_k])] += 1

    # Calculate Average Bond Lengths
    for i in list(bond_lengths.keys()):
        bond_lengths[i] = bond_lengths[i]/bond_counts[i]
        
    # Collect types that will be converted into UA
    UA_types = {}
    removed_H_types = []
    for i in list(Data.keys()):
        
        # Iterate over the molecule a adj_mat and identify carbons with attached hydrogens.
        for count_j,j in enumerate(Data[i]["adj_mat_a"]):
            if int(Data[i]["types_a"][count_j].split('[')[1].split(']')[0]) != 6: continue
            UA_types[Data[i]["types_a"][count_j]] = [ Data[i]["types_a"][count_j] ] +\
                                                    [ Data[i]["types_a"][count_k] for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_a'][count_k].split('[')[1].split(']')[0]) == 1 ]

        # Iterate over the molecule b adj_mat and identify carbons with attached hydrogens.
        for count_j,j in enumerate(Data[i]["adj_mat_b"]):
            if int(Data[i]["types_b"][count_j].split('[')[1].split(']')[0]) != 6: continue
            UA_types[Data[i]["types_b"][count_j]] = [ Data[i]["types_b"][count_j] ] +\
                                                    [ Data[i]["types_b"][count_k] for count_k,k in enumerate(j) if k == 1 and int(Data[i]['types_b'][count_k].split('[')[1].split(']')[0]) == 1 ]

    # Add reverse combinations of keys to VDW_dict (they are removed
    # by the write_params function to avoid saving redundant pairs)
    VDW_dict_keys = list(VDW_dict.keys())
    for i in VDW_dict_keys:
        if (i[1],i[0]) not in list(VDW_dict.keys()): VDW_dict[(i[1],i[0])]=VDW_dict[i]

    # Initialize the new VDW dictionary to hold the parameters with updated sigma radii
    # (avoids the complications associated with updates to VDW_dict while fits are being performed)
    VDW_new = {}

    # Calculate sigmas based on the scaled volume of the largest sphere that encompasses the UA atom group
    if vol_method == 'radius':
        
        # Print diagnostic
        print("\nCalculating LJ parameters based on the (radial) group volume and the total configurational LJ interaction energy:\n")

        # Loop over the VDW dictionary and update UA sigma values based on membership in UA_types    
        for i in list(VDW_dict.keys()):

            # Initialize V_eff
            V_eff = 0.0

            # If neither atom type is in the UA_list then set the excluded volume based
            # on the existing sigma value
            if i[0] not in list(UA_types.keys()) and i[1] not in list(UA_types.keys()):
                V_eff = 4.0/3.0*pi*(VDW_dict[i][1])**(3.0)

            # Calculate effective radius for type_1 on the basis of bond length and the sigma values for H and C
            if i[0] in list(UA_types.keys()):
                sigmas = [VDW_dict[i][1]]
                for count_j,j in enumerate(UA_types[i[0]][1:]):

                    # Displacement for C-H is based on the bond length
                    sep = bond_lengths[(i[0],j)]

                    # Add effective sigma based on the bond length plus the attached H sigma when interating with type_2
                    sigmas += [VDW_dict[(j,i[1])][1]+sep]

                # Calculate the effective volume based on the largest radius
                max_rad = max(sigmas)
                if max_rad == sigmas[0]:
                    V_eff += vol_scale*4.0/3.0*pi*max_rad**(3.0)
                    #V_eff += 4.0/3.0*pi*max_rad**(3.0) # old, leaves the value unscaled if the AA-sigma of the core atom determined the outer volume.      
                else:
                    V_eff += vol_scale*4.0/3.0*pi*max_rad**(3.0)

            # Calculate the effective volume for type_2
            if i[1] in list(UA_types.keys()):
                sigmas = [VDW_dict[i][1]]
                for count_j,j in enumerate(UA_types[i[1]][1:]):

                    # Displacement for C-H is based on the bond length
                    sep = bond_lengths[(i[1],j)]

                    # Add effective sigma based on the bond length plus the attached H sigma when interating with type_2
                    sigmas += [VDW_dict[(i[0],j)][1]+sep]

                # Calculate the effective volume based on the largest radius
                max_rad = max(sigmas)
                if max_rad == sigmas[0]:
                    V_eff += vol_scale*4.0/3.0*pi*max_rad**(3.0)
                    #V_eff += 4.0/3.0*pi*max_rad**(3.0) # old, leaves the value unscaled if the AA-sigma of the core atom determined the outer volume.
                else:
                    V_eff += vol_scale*4.0/3.0*pi*max_rad**(3.0)

            # If both types are UA then subtract off the C-C interation because its volume contribution was double counted
            if i[0] in list(UA_types.keys()) and i[1] in list(UA_types.keys()):
                V_eff -= 4.0/3.0*pi*(VDW_dict[i][1])**(3.0)

            # Calculate sigma_eff based on the excluded volume for the pair-wise interaction
            sigma_eff = ( 3.0 / (4.0*pi) * V_eff )**(1.0/3.0)

            # Assign the value to the VDW_dict
            VDW_new[i] = (VDW_dict[i][0],sigma_eff)

    # VOLUME BASED SIGMA VALUES
    # Loop over the VDW dictionary and update UA sigma values based on membership in UA_types
    if vol_method == 'volume':

        # Print diagnostic
        print("\nCalculating LJ parameters based on the (MC-calculated) group volume and the total configurational LJ interaction energy:\n")

        # Initialize sobol quasi-random value list (used in the calc_eff_vol_MC for monte-carlo volume calculations)
        sobol_list = i4_sobol_generate( 3, 1000000, 10 )

        # Loop over the VDW dictionary and update UA sigma values based on membership in UA_types    
        for i in list(VDW_dict.keys()):

            # Initialize V_eff
            V_eff = 0.0

            # If neither atom type is in the UA_list then set the excluded volume based
            # on the existing sigma value
            if i[0] not in list(UA_types.keys()) and i[1] not in list(UA_types.keys()):
                V_eff = 4.0/3.0*pi*(VDW_dict[i][1])**(3.0)

            # Calculate the effective volume for type_1
            if i[0] in list(UA_types.keys()):
                seps = zeros([len(UA_types[i[0]]),len(UA_types[i[0]])])
                for count_j,j in enumerate(UA_types[i[0]]):
                    for count_k,k in enumerate(UA_types[i[0]]):

                        # The self distance is 0
                        if count_j == count_k:
                            seps[count_j,count_k] = 0.0

                        # Only explicitly set separations for the central carbon
                        elif count_j == 0 or count_k == 0:
                            seps[count_j,count_k] = bond_lengths[(j,k)]

                        # The separations between all attached hydrogens are set to an arbitrarily large value (enforces the assumption of no overlap)
                        else:
                            seps[count_j,count_k] = 1000.0

                # Collect the sigma values for the carbon and all attached hydrogens with the i[1] atom type
                sigmas = [ VDW_dict[(j,i[1])][1] for j in UA_types[i[0]] ]

                # Calculate the effective volume of the C with its joined atoms (w.r.t. the i[0] atom type)
                # V_eff += calc_eff_vol(seps,*sigmas) # old sphere based method
                V_eff += calc_eff_vol_MC(sobol_list,seps,*sigmas)

            # Calculate the effective volume for type_2
            if i[1] in list(UA_types.keys()):
                seps = zeros([len(UA_types[i[1]]),len(UA_types[i[1]])])
                for count_j,j in enumerate(UA_types[i[1]]):
                    for count_k,k in enumerate(UA_types[i[1]]):

                        # The self distance is 0
                        if count_j == count_k:
                            seps[count_j,count_k] = 0.0

                        # Only explicitly set separations for the central carbon
                        elif count_j == 0 or count_k == 0:
                            seps[count_j,count_k] = bond_lengths[(j,k)]

                        # The separations between all attached hydrogens are set to an arbitrarily large value (enforces the assumption of no overlap)
                        else:
                            seps[count_j,count_k] = 1000.0

                # Collect the sigma values for the carbon and all attached hydrogens with the i[1] atom type
                sigmas = [ VDW_dict[(i[0],j)][1] for j in UA_types[i[1]] ]

                # Calculate the effective volume of the C with its joined atoms (w.r.t. the i[0] atom type)
                # V_eff += calc_eff_vol(seps,*sigmas) # old sphere based method
                V_eff += calc_eff_vol_MC(sobol_list,seps,*sigmas)

            # If both types are UA then subtract off the C-C interation because its volume contribution was double counted
            if i[0] in list(UA_types.keys()) and i[1] in list(UA_types.keys()):
                V_eff -= 4.0/3.0*pi*(VDW_dict[i][1])**(3.0)

            # Calculate sigma_eff based on the excluded volume for the pair-wise interaction
            sigma_eff = ( 3.0 / (4.0*pi) * V_eff )**(1.0/3.0)

            # Assign the value to the VDW_dict
            VDW_new[i] = (VDW_dict[i][0],sigma_eff)

    # Loop over the VDW dictionary and update UA sigma values based on membership in UA_types    
    run_types = []
    for i in list(VDW_dict.keys()):

        # Avoid redundant runs
        if i in run_types: continue
        else: run_types += [(i[0],i[1]),(i[1],i[0])]

        # Initialize total energy and types variables
        # set_flag is used to check if the algorithm's assign a VDW values or if the parameter was unassigned by the following loops
        # tot_LJ holds the LJ contribution for the interaction between the UA-groups
        # types is used to avoid redundant additions
        set_flag = 0
        tot_LJ   = 0.0
        types    = []

        # If either atom is a UA-H then the epsilon value is set to zero
        if ( i[0] in [ k for j in list(UA_types.keys()) for k in UA_types[j][1:] ] ) or ( i[1] in [ k for j in list(UA_types.keys()) for k in UA_types[j][1:] ] ):
            set_flag = 1
            VDW_new[(i[0],i[1])] = (0.0,VDW_new[i][1])
            VDW_new[(i[1],i[0])] = (0.0,VDW_new[i][1])

        # If neither atom is a UA_atom then the parameters are copied as is.
        elif i[0] not in list(UA_types.keys()) and i[1] not in list(UA_types.keys()):
            set_flag = 1
            VDW_new[(i[0],i[1])] = VDW_new[i]
            VDW_new[(i[1],i[0])] = VDW_new[i]

        # If one or both of the atoms are UA-C then the epsilon value is adjusted based on 
        # the configurationally summed LJ interaction energy for the group A: group B interactions.
        else:

            # If the first type is a UA-C then the contribution of it's group of atoms interacting with i[1] are added to tot_LJ
            if i[0] in list(UA_types.keys()):

                # Iterate over the group 1:atom 2 interactions and add up their contribution to the total_LJ energy
                # Note: double-counting is avoided using the check against the types list.
                for count_j,j in enumerate(UA_types[i[0]]):
                    if (i[1],j) in types: continue
                    if count_j == 0:
                        tot_LJ += E_LJ_Type_w_con((i[1],j),[],[])  # No constraints for central carbon
                    else:
                        tot_LJ += E_LJ_Type_w_con((i[1],j),[],UA_types[i[0]][0]) # Constraints for hydrogens attached to central carbon (first element in UA_types).
#                    tot_LJ += E_LJ_Type((i[1],j))
                    types  += [(i[1],j),(j,i[1])]

            # If the second type is a UA-C then the contribution of it's group of atoms interacting with i[0] are added to tot_LJ
            if i[1] in list(UA_types.keys()):

                # Iterate over the group 2:atom 1 interactions and add up their contribution to the total_LJ energy
                # Note: double-counting is avoided using the check against the types list.
                for count_j,j in enumerate(UA_types[i[1]]):
                    if (i[0],j) in types: continue
                    if count_j == 0:
                        tot_LJ += E_LJ_Type_w_con((i[0],j),[],[])  # No constraints for central carbon
                    else:
                        tot_LJ += E_LJ_Type_w_con((i[0],j),[],UA_types[i[1]][0]) # Constraints for hydrogens attached to central carbon (first element in UA_types).
#                    tot_LJ += E_LJ_Type((i[0],j))                
                    types  += [(i[0],j),(j,i[0])]

            # If both types in the pair are UA-C then the UA-H:UA-H interactions need to be added to tot_LJ
            if i[0] in list(UA_types.keys()) and i[1] in list(UA_types.keys()):

                # Add the UA Hydrogen interactions to the total UA energy
                # Note: The UA-C:UA-C interaction has already been included in the previous loops, the comparison with the types list avoids the 
                #       redundant evaluation.
                for j in UA_types[i[0]]:
                    for k in UA_types[i[1]]:
                        if (j,k) in types: continue
                        tot_LJ += E_LJ_Type_w_con((j,k),UA_types[i[0]][0],UA_types[i[1]][0]) # Constraints for hydrogens attached to central carbon (first element in UA_types).
#                        tot_LJ += E_LJ_Type((j,k))
                        types += [(j,k),(k,j)]

            # Calculate the scaled epsilon value if either of the types involved in the current pair are UA_types.
            if i[0] in list(UA_types.keys()) or i[1] in list(UA_types.keys()):

                # Initialize variables for sigma_scale, the objective function, and the initial guess for the new eps and sigma values
                # Note: if eps_new and sigma_new aren't reassigned by the while loop, the program prints a warning that the algorithm failed. 
                sigma_scale     = 2.0
                f               = 1E6
                eps_new         = VDW_dict[i][0]
                sigma_new       = VDW_new[i][1]

                # Loop over different combinations of sigma_scale and eps_scale, using the objective function
                # to determine the optimal combination with respect to the original AA-eps value and volume based sigma.
                while 1:

                    # The E_LJ_Type function uses the values in VDW_dict to calculate E_LJ for a given interaction over all configurations. 
                    # Since we want the total LJ energy for the current interaction (i) but with the new sigma, we need to temporarily assign 
                    # the new sigma (VDW_new[i][1]) to VDW_dict. For consistency, the value is reassigned afterward.
                    old_sigma = VDW_dict[i][1]
                    VDW_dict[(i[0],i[1])] = (VDW_dict[(i[0],i[1])][0],sigma_scale*VDW_new[i][1])
                    VDW_dict[(i[1],i[0])] = (VDW_dict[(i[1],i[0])][0],sigma_scale*VDW_new[i][1])
#                    current = E_LJ_Type(i)
                    current = E_LJ_Type_w_con(i,[],[])
                    VDW_dict[(i[0],i[1])] = (VDW_dict[(i[0],i[1])][0],old_sigma)
                    VDW_dict[(i[1],i[0])] = (VDW_dict[(i[1],i[0])][0],old_sigma)

                    # Calculate eps_scale based on the ratio of LJ_energies, eps_current based on the AA-eps and the scale factor,
                    # and the scale factor, and sigma_current based on sigma_scale and the volume-based sigma value.
                    eps_scale     = (tot_LJ/current) 
                    eps_current   = eps_scale * VDW_dict[i][0]
                    sigma_current = sigma_scale * VDW_new[i][1]
                    
                    # Update objective value
                    f_current = (1.0 - eps_current/VDW_dict[i][0])**(2.0) + (1.0 - sigma_current/VDW_new[i][1])**(2.0)

                    # Only consider updates with positive epsilon values (and, commented out right now, sigma values that are at least as large as the original AA sigmas)
#                    if f_current < f and eps_current > 0.0:  # and sigma_current > VDW_dict[i][1]:
                    if f_current < f and eps_current > 0.0  and sigma_current > Pair_min_dict[i]:
                        set_flag  = 1            
                        f         = f_current
                        eps_new   = eps_current
                        sigma_new = sigma_current

                    # Update sigma_scale and check break condition (only positive sigma values are sampled)
                    sigma_scale -= 0.001
                    if sigma_scale < 0.0: break

                # Update the VDW parameters
                VDW_new[(i[0],i[1])] = (eps_new,sigma_new)
                VDW_new[(i[1],i[0])] = (eps_new,sigma_new)
                
        # Print diagnostic
        if set_flag == 1:
            print("\tpair {:80}: {} ->     {}".format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<10.4f}".format(j) for j in [VDW_new[i][0],VDW_new[i][1]]])))
        else:
            print("\tpair {:80}: {} ->     {} (WARNING: algorithm failed to assign parameters)".\
                  format(i,", ".join([ "{:<10.4f}".format(j) for j in VDW_dict[i]]),", ".join([ "{:<10.4f}".format(j) for j in [VDW_new[i][0],VDW_new[i][1]]])))

    return VDW_new

# Calculate the effective sigma value for united atoms based on the volume of the 
# bonded atoms (sum of invidivual hard sphere volumes minus their overlap)
def calc_eff_vol_MC(sobol_list,seps,sigma_1=0,sigma_2=0,sigma_3=0,sigma_4=0,sigma_5=0,sigma_6=0,sigma_7=0,sigma_8=0,sigma_9=0,sigma_10=0):

    # Initialize local variable dictionary (used for determining what has been defined)
    local_vars = locals()
    
    # Find defined sigmas
    sigma_vars = natural_sort([ i for i in local_vars if 'sigma_' in i and "__" not in i and local_vars[i] != 0 ])
    sigma_vals = [ locals()[i] for i in sigma_vars ]

    V_eff = 0.0

    # Add central atom volume to V_eff 
    V_eff += 4.0/3.0*pi*sigma_vals[0]**(3.0)
    
    # Subtract out overlapping regions from V_eff
    for count_i,i in enumerate(sigma_vals):
        for count_j,j in enumerate(sigma_vals):

            # Only sum over the central-atom UA_H bonds
            if count_i != 0 or count_j == 0: continue

            # Add bond volume (central-atom + hydrogen + cylinder the size of hydrogen connecting the two)
            V_eff += calc_pair_volume(i,j,seps[count_i,count_j],sobol_list)

            # Subtract off the central-atom volume because it was already included in V_eff
            V_eff -= 4.0/3.0*pi*sigma_vals[0]**(3.0)

    return V_eff

# Monte-carlo volume calculator for two spheres connected by a cylinder with the radius of the smaller sphere.
# r_1 and r_2 are the sphere radii, length is their separation, and sobol_list is a preinitialized list of 
# quasi-random triples for sampling the enclosing volume.
def calc_pair_volume(r_1,r_2,length,sobol_list):

    # Find larger r
    if r_1 >= r_2: r_big = r_1; r_small = r_2
    else: r_big = r_2; r_small = r_1

    # If one sphere completely contains the other then return the volume of the larger sphere
    if length < abs(r_1-r_2): return 4.0/3.0*pi*r_big**(3.0)
    
    V_box = (r_1+r_2+length) * (2.0*r_big) * (2.0*r_big)
    counts = 0

    # Scale the samples to fill the box (only one quadrant is sampled since the object is 4-fold symmetric)
    samples = array([ [sobol_list[0][i],sobol_list[1][i],sobol_list[2][i]] for i in range(len(sobol_list[0])) ])*array([r_1+length+r_2,r_big,r_big]) 

    # Monte-carlo sample the volume
    for count_p,p in enumerate(samples):

        # sphere one is located at r_1,r_big,r_big
        if (p[0]-r_1)**2.0 + (p[1]-r_big)**2.0 + (p[2]-r_big)**2.0 < r_1**2.0: counts += 1; continue

        # sphere two is located at r_1+length,r_big,r_big
        if (p[0]-(r_1+length))**2.0 + (p[1]-r_big)**2.0 + (p[2]-r_big)**2.0 < r_2**2.0: counts += 1; continue

        # the connecting cylinder is oriented along the x-axis from r_1 to r_1+length, centered at r_big,r_big in y and z, with a radius of r_small
        if (p[0] >= r_1 and p[0] < (r_1+length) ) and (p[1]-r_big)**2.0 + (p[2]-r_big)**2.0 < r_small**2.0: counts += 1; continue

    return V_box * float(counts)/float(len(samples))

def i4_bit_hi1 ( n ):

#*****************************************************************************80
#
## I4_BIT_HI1 returns the position of the high 1 bit base 2 in an integer.
#
#  Example:
#
#       N    Binary     BIT
#    ----    --------  ----
#       0           0     0
#       1           1     1
#       2          10     2
#       3          11     2 
#       4         100     3
#       5         101     3
#       6         110     3
#       7         111     3
#       8        1000     4
#       9        1001     4
#      10        1010     4
#      11        1011     4
#      12        1100     4
#      13        1101     4
#      14        1110     4
#      15        1111     4
#      16       10000     5
#      17       10001     5
#    1023  1111111111    10
#    1024 10000000000    11
#    1025 10000000001    11
#
#  Licensing:
#
#    This code is distributed under the MIT license.
#
#  Modified:
#
#    22 February 2011
#
#  Author:
#
#    Original MATLAB version by John Burkardt.
#    PYTHON version by Corrado Chisari
#
#  Parameters:
#
#    Input, integer N, the integer to be measured.
#    N should be nonnegative.  If N is nonpositive, the value will always be 0.
#
#    Output, integer BIT, the number of bits base 2.
#
  i = math.floor ( n )
  bit = 0
  while ( 1 ):
    if ( i <= 0 ):
      break
    bit += 1
    i = math.floor ( i / 2. )
  return bit

def i4_bit_lo0 ( n ):

#*****************************************************************************80
#
## I4_BIT_LO0 returns the position of the low 0 bit base 2 in an integer.
#
#  Example:
#
#       N    Binary     BIT
#    ----    --------  ----
#       0           0     1
#       1           1     2
#       2          10     1
#       3          11     3 
#       4         100     1
#       5         101     2
#       6         110     1
#       7         111     4
#       8        1000     1
#       9        1001     2
#      10        1010     1
#      11        1011     3
#      12        1100     1
#      13        1101     2
#      14        1110     1
#      15        1111     5
#      16       10000     1
#      17       10001     2
#    1023  1111111111     1
#    1024 10000000000     1
#    1025 10000000001     1
#
#  Licensing:
#
#    This code is distributed under the MIT license.
#
#  Modified:
#
#    22 February 2011
#
#  Author:
#
#    Original MATLAB version by John Burkardt.
#    PYTHON version by Corrado Chisari
#
#  Parameters:
#
#    Input, integer N, the integer to be measured.
#    N should be nonnegative.
#
#    Output, integer BIT, the position of the low 1 bit.
#
	bit = 0
	i = math.floor ( n )
	while ( 1 ):
		bit = bit + 1
		i2 = math.floor ( i / 2. )
		if ( i == 2 * i2 ):
			break

		i = i2
	return bit

def i4_sobol_generate ( m, n, skip ):

#*****************************************************************************80
#
## I4_SOBOL_GENERATE generates a Sobol dataset.
#
#  Licensing:
#
#    This code is distributed under the MIT license.
#
#  Modified:
#
#    22 February 2011
#
#  Author:
#
#    Original MATLAB version by John Burkardt.
#    PYTHON version by Corrado Chisari
#
#  Parameters:
#
#    Input, integer M, the spatial dimension.
#
#    Input, integer N, the number of points to generate.
#
#    Input, integer SKIP, the number of initial points to skip.
#
#    Output, real R(M,N), the points.
#
	r=zeros((m,n))
	for j in range (1, n+1):
		seed = skip + j - 2
		[ r[0:m,j-1], seed ] = i4_sobol ( m, seed )
	return r

def i4_sobol ( dim_num, seed ):

#*****************************************************************************80
#
## I4_SOBOL generates a new quasirandom Sobol vector with each call.
#
#  Discussion:
#
#    The routine adapts the ideas of Antonov and Saleev.
#
#  Licensing:
#
#    This code is distributed under the MIT license.
#
#  Modified:
#
#    22 February 2011
#
#  Author:
#
#    Original FORTRAN77 version by Bennett Fox.
#    MATLAB version by John Burkardt.
#    PYTHON version by Corrado Chisari
#
#  Reference:
#
#    Antonov, Saleev,
#    USSR Computational Mathematics and Mathematical Physics,
#    olume 19, 1980, pages 252 - 256.
#
#    Paul Bratley, Bennett Fox,
#    Algorithm 659:
#    Implementing Sobol's Quasirandom Sequence Generator,
#    ACM Transactions on Mathematical Software,
#    Volume 14, Number 1, pages 88-100, 1988.
#
#    Bennett Fox,
#    Algorithm 647:
#    Implementation and Relative Efficiency of Quasirandom 
#    Sequence Generators,
#    ACM Transactions on Mathematical Software,
#    Volume 12, Number 4, pages 362-376, 1986.
#
#    Ilya Sobol,
#    USSR Computational Mathematics and Mathematical Physics,
#    Volume 16, pages 236-242, 1977.
#
#    Ilya Sobol, Levitan, 
#    The Production of Points Uniformly Distributed in a Multidimensional 
#    Cube (in Russian),
#    Preprint IPM Akad. Nauk SSSR, 
#    Number 40, Moscow 1976.
#
#  Parameters:
#
#    Input, integer DIM_NUM, the number of spatial dimensions.
#    DIM_NUM must satisfy 1 <= DIM_NUM <= 40.
#
#    Input/output, integer SEED, the "seed" for the sequence.
#    This is essentially the index in the sequence of the quasirandom
#    value to be generated.	On output, SEED has been set to the
#    appropriate next value, usually simply SEED+1.
#    If SEED is less than 0 on input, it is treated as though it were 0.
#    An input value of 0 requests the first (0-th) element of the sequence.
#
#    Output, real QUASI(DIM_NUM), the next quasirandom vector.
#
	global atmost
	global dim_max
	global dim_num_save
	global initialized
	global lastq
	global log_max
	global maxcol
	global poly
	global recipd
	global seed_save
	global v

	if ( not 'initialized' in list(globals().keys()) ):
		initialized = 0
		dim_num_save = -1

	if ( not initialized or dim_num != dim_num_save ):
		initialized = 1
		dim_max = 40
		dim_num_save = -1
		log_max = 30
		seed_save = -1
#
#	Initialize (part of) V.
#
		v = zeros((dim_max,log_max))
		v[0:40,0] = transpose([ \
			1, 1, 1, 1, 1, 1, 1, 1, 1, 1, \
			1, 1, 1, 1, 1, 1, 1, 1, 1, 1, \
			1, 1, 1, 1, 1, 1, 1, 1, 1, 1, \
			1, 1, 1, 1, 1, 1, 1, 1, 1, 1 ])

		v[2:40,1] = transpose([ \
			1, 3, 1, 3, 1, 3, 3, 1, \
			3, 1, 3, 1, 3, 1, 1, 3, 1, 3, \
			1, 3, 1, 3, 3, 1, 3, 1, 3, 1, \
			3, 1, 1, 3, 1, 3, 1, 3, 1, 3 ])

		v[3:40,2] = transpose([ \
			7, 5, 1, 3, 3, 7, 5, \
			5, 7, 7, 1, 3, 3, 7, 5, 1, 1, \
			5, 3, 3, 1, 7, 5, 1, 3, 3, 7, \
			5, 1, 1, 5, 7, 7, 5, 1, 3, 3 ])

		v[5:40,3] = transpose([ \
			1, 7, 9,13,11, \
			1, 3, 7, 9, 5,13,13,11, 3,15, \
			5, 3,15, 7, 9,13, 9, 1,11, 7, \
			5,15, 1,15,11, 5, 3, 1, 7, 9 ])
	
		v[7:40,4] = transpose([ \
			9, 3,27, \
			15,29,21,23,19,11,25, 7,13,17, \
			1,25,29, 3,31,11, 5,23,27,19, \
			21, 5, 1,17,13, 7,15, 9,31, 9 ])

		v[13:40,5] = transpose([ \
							37,33, 7, 5,11,39,63, \
		 27,17,15,23,29, 3,21,13,31,25, \
			9,49,33,19,29,11,19,27,15,25 ])

		v[19:40,6] = transpose([ \
			13, \
			33,115, 41, 79, 17, 29,119, 75, 73,105, \
			7, 59, 65, 21,	3,113, 61, 89, 45,107 ])

		v[37:40,7] = transpose([ \
			7, 23, 39 ])
#
#	Set POLY.
#
		poly= [ \
			1,	 3,	 7,	11,	13,	19,	25,	37,	59,	47, \
			61,	55,	41,	67,	97,	91, 109, 103, 115, 131, \
			193, 137, 145, 143, 241, 157, 185, 167, 229, 171, \
			213, 191, 253, 203, 211, 239, 247, 285, 369, 299 ]

		atmost = 2**log_max - 1
#
#	Find the number of bits in ATMOST.
#
		maxcol = i4_bit_hi1 ( atmost )
#
#	Initialize row 1 of V.
#
		v[0,0:maxcol] = 1

#
#	Things to do only if the dimension changed.
#
	if ( dim_num != dim_num_save ):
#
#	Check parameters.
#
		if ( dim_num < 1 or dim_max < dim_num ):
			print('I4_SOBOL - Fatal error!') 
			print('	The spatial dimension DIM_NUM should satisfy:') 
			print('		1 <= DIM_NUM <= %d'%dim_max)
			print('	But this input value is DIM_NUM = %d'%dim_num)
			return

		dim_num_save = dim_num
#
#	Initialize the remaining rows of V.
#
		for i in range(2 , dim_num+1):
#
#	The bits of the integer POLY(I) gives the form of polynomial I.
#
#	Find the degree of polynomial I from binary encoding.
#
			j = poly[i-1]
			m = 0
			while ( 1 ):
				j = math.floor ( j / 2. )
				if ( j <= 0 ):
					break
				m = m + 1
#
#	Expand this bit pattern to separate components of the logical array INCLUD.
#
			j = poly[i-1]
			includ=zeros(m)
			for k in range(m, 0, -1):
				j2 = math.floor ( j / 2. )
				includ[k-1] =  (j != 2 * j2 )
				j = j2
#
#	Calculate the remaining elements of row I as explained
#	in Bratley and Fox, section 2.
#
			for j in range( m+1, maxcol+1 ):
				newv = v[i-1,j-m-1]
				l = 1
				for k in range(1, m+1):
					l = 2 * l
					if ( includ[k-1] ):
						newv = bitwise_xor ( int(newv), int(l * v[i-1,j-k-1]) )
				v[i-1,j-1] = newv
#
#	Multiply columns of V by appropriate power of 2.
#
		l = 1
		for j in range( maxcol-1, 0, -1):
			l = 2 * l
			v[0:dim_num,j-1] = v[0:dim_num,j-1] * l
#
#	RECIPD is 1/(common denominator of the elements in V).
#
		recipd = 1.0 / ( 2 * l )
		lastq=zeros(dim_num)

	seed = int(math.floor ( seed ))

	if ( seed < 0 ):
		seed = 0

	if ( seed == 0 ):
		l = 1
		lastq=zeros(dim_num)

	elif ( seed == seed_save + 1 ):
#
#	Find the position of the right-hand zero in SEED.
#
		l = i4_bit_lo0 ( seed )

	elif ( seed <= seed_save ):

		seed_save = 0
		l = 1
		lastq=zeros(dim_num)

		for seed_temp in range( int(seed_save), int(seed)):
			l = i4_bit_lo0 ( seed_temp )
			for i in range(1 , dim_num+1):
				lastq[i-1] = bitwise_xor ( int(lastq[i-1]), int(v[i-1,l-1]) )

		l = i4_bit_lo0 ( seed )

	elif ( seed_save + 1 < seed ):

		for seed_temp in range( int(seed_save + 1), int(seed) ):
			l = i4_bit_lo0 ( seed_temp )
			for i in range(1, dim_num+1):
				lastq[i-1] = bitwise_xor ( int(lastq[i-1]), int(v[i-1,l-1]) )

		l = i4_bit_lo0 ( seed )
#
#	Check that the user is not calling too many times!
#
	if ( maxcol < l ):
		print('I4_SOBOL - Fatal error!')
		print('	Too many calls!')
		print('	MAXCOL = %d\n'%maxcol)
		print('	L =			%d\n'%l)
		return
#
#	Calculate the new components of QUASI.
#
	quasi=zeros(dim_num)
	for i in range( 1, dim_num+1):
		quasi[i-1] = lastq[i-1] * recipd
		lastq[i-1] = bitwise_xor ( int(lastq[i-1]), int(v[i-1,l-1]) )

	seed_save = seed
	seed = seed + 1

	return [ quasi, seed ]

# Checks for job completion in the discovered orca output files
def keep_complete(Names):
    complete = []
    for i in Names:
        with open(i,'r') as f:
            for lines in f:
                if "SCF NOT CONVERGED AFTER" in lines:
                    break
                if "****ORCA TERMINATED NORMALLY****" in lines:
                    complete += [i]
    return complete

# Checks for UA carbons
def UA_carbon(type):
    local_adjmat,labels = type_adjmat(type)
    if labels[0] == "6" and len([ i for count_i,i in enumerate(local_adjmat[0]) if i == 1 and labels[count_i] == "1" ]) > 0:
        return True
    else:
        return False

# This function generates a local adjacency matrix from an atomtype label
def type_adjmat(type):

    # Initialize breaks (indices for brackets in "type"), atoms (start:end index tuples for the atom labels),
    # labels ( basis of atom labels indexed to the adjmat), and adj_mat (the local adjacency matrix with connectivity relationships)
    breaks = [ count_i for count_i,i in enumerate(type) if i in ['[',']'] ]  
    atoms  = [ (i+1,i+breaks[count_i+1]-i) for count_i,i in enumerate(breaks[:-1]) if breaks[count_i+1]-i > 1 ]    
    labels = [ type[i[0]:i[1]] for i in atoms ]   
    adj_mat = zeros([len(atoms),len(atoms)])      

    # Loop over atoms
    for count_i,i in enumerate(atoms):

        # Initialize variables
        starting_index = breaks.index(i[1])       # index of the nearest bracket forward from the current atom label
        break_count=0                             # counter for keeping track of who needs parsing and who is connected to who

        # Loop over brackets
        for count_j,j in enumerate(breaks[starting_index:-1]):

            # The atom has no new connections if the first bracket is closed
            if count_j == 0 and type[j] == "]":  
                break

            # Increment break_count + 1 for "open" brackets
            if type[j] == "[": break_count += 1

            # Increment break_count - 1 for "closed" brackets
            if type[j] == "]": break_count -= 1

            # When break_count == 1 and the parser resides at an open bracket, the next atom past the bracket is connected
            if break_count == 1 and type[j] == "[":
                idx = next( count_k for count_k,k in enumerate(atoms) if k[0] == j+1 )
                adj_mat[count_i,idx] = 1
                adj_mat[idx,count_i] = 1

    return adj_mat, labels
def fit_method(method,x_vals,y_vals,Pairs,weight,fit,E_config_tot,delta_xhi2_thresh=0.00001,min_cycles=100,max_cycles=1000,mixing_rule='none',L2_s=0.0,L2_e=0.0,VDW_0=None,outlier_option=False):
   method_list = ['lstsq','global','Boltzmann','linear','quadratic','Boltzmann_tot']
   if method not in method_list:
      print("Error: {} is not a supported fitting method, available option: {}".format(method," ".join(method_list)))
      quit()
   if method == 'global':
      fit_vals, xhi = global_fit(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule,L2_s,L2_e,VDW_0,outlier_option)
   if method == 'lstsq':
      fit_vals, xhi = lstsq_fit(x_vals,y_vals,Pairs,fit)
   if method == 'Boltzmann':
      fit_vals, xhi = global_fit_Boltz(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule,0,0,VDW_0,outlier_option)
   if method == 'linear':
      fit_vals, xhi = global_fit_Boltz(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule,0,0,VDW_0,outlier_option,weighting_function='linear')
   if method == 'quadratic':
      fit_vals, xhi = global_fit_Boltz(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule,0,0,VDW_0,outlier_option,weighting_function='quadratic')
   if method == 'Boltzmann_tot':
      fit_vals, xhi = global_fit_Boltz(x_vals,y_vals,Pairs,weight,fit,delta_xhi2_thresh,min_cycles,max_cycles,mixing_rule,0,0,VDW_0,outlier_option,E_config_tot)
   return fit_vals,xhi


def plot_method_convergence(folder,VDW_dict_multi,xhi_total,VDW_0=None):
   global VDW_dict
   # if input VDW_0(UFF value) then include UFF into plot
   # Collect data sets into plottable lists    
   keys = list(VDW_dict.keys())
   methods =  list(VDW_dict_multi.keys())
   for i in keys:
      if i[0] > i[1]:
         if (i[1],i[0]) in list(VDW_dict.keys()):             
            VDW_dict.pop((i[1],i[0]))
            for j in VDW_dict_multi:               
               VDW_dict_multi[j].pop((i[1],i[0]))
   sets = { i:{"eps":[],"sigma":[]} for i in list(VDW_dict.keys()) }
   for i in list(VDW_dict.keys()):
      for j in VDW_dict_multi:
            sets[i]["eps"] += [VDW_dict_multi[j][i][0]]
            sets[i]["sigma"] += [VDW_dict_multi[j][i][1]]

   ###### PLOT SIGMA CONVERGENCE ######
   #fig = plt.figure()
   #ax = plt.subplot(111)
   fig,ax = plt.subplots(1,2,figsize=(10,4))
   if len(methods) == 1: ax=[ax]
   plot_handles = []
   color_list = [(0.05,0.35,0.75),(0.05,0.8,0.6),(0.9,0.3,0.05),(0.35,0.7,0.9),(0.9,0.5,0.7),(0.9,0.6,0.05),(0.95,0.9,0.25),(0.05,0.05,0.05)]   # Supposedly more CB-friendly
   color_list += [ (i[0]*0.8,i[1]*0.8,i[2]*0.8) for i in color_list ] # Extend 
   color_list += [ (i[0]*0.6,i[1]*0.6,i[2]*0.6) for i in color_list ] # Extend 
   color_list += [ (i[0]*0.4,i[1]*0.4,i[2]*0.4) for i in color_list ]*100 # Extend 

   for count_i,i in enumerate(sets.keys()):            
      plot_handle, = ax[0].plot(arange(1,len(methods)+1),array(sets[i]["sigma"]),\
                            linestyle="-",linewidth=2.0,\
                            marker=".",markersize=20,color=color_list[count_i],label=i)
      plot_handle, = ax[1].plot(arange(1,len(methods)+1),array(sets[i]["eps"]),\
                            linestyle="-",linewidth=2.0,\
                            marker=".",markersize=20,color=color_list[count_i],label=i)
      # if dashes_list[count_i] != []:
      #     plot_handle.set_dashes(dashes_list[count_i])
      # plot_handle.set_dash_capstyle('round')
      # plot_handle.set_dash_joinstyle('round')
   #ax.set_ylabel("sigma (ang)".decode('utf-8'),fontsize=32,labelpad=10,fontweight='bold')
   ax[0].set_ylabel("sigma (ang)",fontsize=20,labelpad=10,fontweight='bold')
   ax[1].set_ylabel("eps (kcal/mol)",fontsize=20,labelpad=10,fontweight='bold')
   labels = [ '{} \nxhi={:2.4f}'.format(i,xhi_total[i]) for i in methods]
   for i in range(2):
      ax[i].set_xlabel("method",fontsize=32,labelpad=10,fontweight='bold')
      ax[i].set_xticks(arange(1,len(methods)+1))
      ax[i].set_xticklabels(labels,{'fontsize': 10})
      #ax[i].tick_params(axis='both', which='major',labelsize=24,pad=10,direction='out',width=3,length=6)
      #ax[i].tick_params(axis='both', which='minor',labelsize=24,pad=10,direction='out',width=2,length=4)
      ax[i].tick_params(axis='both', which='major',pad=10,direction='out',width=3,length=6)
      ax[i].tick_params(axis='both', which='minor',pad=10,direction='out',width=2,length=4)
      [j.set_linewidth(3) for j in ax[i].spines.values()]

   # # Set limits
   # y_min,y_max = ax.get_ylim()
   # x_min,x_max = ax.get_xlim()
   # if args.x_min != []: x_min = args.x_min
   # if args.x_max != []: x_max = args.x_max
   # if args.y_min != []: y_min = args.y_min
   # if args.y_max != []: y_max = args.y_max
   # ax.set_xlim([x_min,x_max])
   # ax.set_ylim([y_min,y_max])

   # if args.log_y:
   #     ax.set_yscale('log')
   # if args.log_x:
   #     ax.set_xscale('log')

   plt.tight_layout()
   # Put a legend to the right of the current axis and use modified save call to account for the legend artist
   ax[1].legend(loc='best',frameon=False)
   handles, labels = ax[1].get_legend_handles_labels()
   lgd = ax[1].legend(handles, labels, loc='center left', bbox_to_anchor=(1,0.5))
   savefig(folder+"/method.png", dpi=300, bbox_extra_artists=(lgd,), bbox_inches='tight')            
   close(fig)


# Create logger to save stdout to logfile
class Logger(object):

    def __init__(self,folder):
        self.terminal = sys.stdout
        self.log = open(folder+"/extract_vdw.log", "a")

    def write(self, message):
        self.terminal.write(message)
        self.log.write(message)  
    def flush(self):
        pass

if __name__ == "__main__":
   main(sys.argv[1:])
